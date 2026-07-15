"""
HEP (Heartbeat Evoked Potential) Comparison: Systole vs Diastole

This script loads pickle files from pickles_sleep_stage/EDF/{sleep_stage},
extracts systole and diastole segments, computes HEP for each phase,
and compares them statistically and visually.
"""

import os
import pickle
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from scipy import stats
from scipy.signal import find_peaks, hilbert, butter, filtfilt
from scipy.stats import entropy
from scipy.ndimage import median_filter
import mne
from mne.stats import permutation_cluster_1samp_test
import neurokit2 as nk
try:
    import pywt
    PYWT_AVAILABLE = True
except ImportError:
    PYWT_AVAILABLE = False
import warnings
import streamlit as st
from io import BytesIO
try:
    import pynapple as nap
    PYNAPPLE_AVAILABLE = True
except ImportError:
    PYNAPPLE_AVAILABLE = False
    # Don't call st.warning() here - it must be called after st.set_page_config()

# Try to import PDF generation libraries
try:
    from reportlab.lib.pagesizes import letter, A4  # noqa: F401
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, Table, TableStyle  # noqa: F401
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT  # noqa: F401
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

# Try to import PPTX generation libraries
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN  # noqa: F401
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

warnings.filterwarnings('ignore')

# Fix for MNE pickle loading
import sys  # noqa: E402
import mne.io.array  # noqa: E402
sys.modules['mne.io.array.array'] = mne.io.array

# Import _plot_metric_vs_hrv from utils
from utils.eeg_utils import _plot_metric_vs_hrv  # noqa: E402


def download_plot_as_svg(fig, filename, cache_key=None):
    """
    Create a download button for a matplotlib figure as SVG.
    Uses session state caching to prevent page reloads.
    
    Parameters
    ----------
    fig : matplotlib.figure.Figure
        The figure to download
    filename : str
        Base filename for the download (without extension)
    cache_key : str, optional
        Unique cache key for session state. If None, uses filename.
    """
    if 'st' not in globals() or fig is None:
        return
    
    if cache_key is None:
        cache_key = f"svg_{filename}"
    
    # Convert figure to SVG and cache in session state
    if cache_key not in st.session_state:
        svg_buffer = BytesIO()
        fig.savefig(svg_buffer, format='svg', bbox_inches='tight')
        svg_buffer.seek(0)
        st.session_state[cache_key] = svg_buffer.getvalue()
    
    # Create download button
    st.download_button(
        label="📥 Download Plot as SVG",
        data=st.session_state[cache_key],
        file_name=f"{filename}.svg",
        mime="image/svg+xml",
        key=f"download_svg_{cache_key}",
        use_container_width=False
    )


def display_plot_with_download(fig, filename, cache_key=None):
    """
    Display a matplotlib figure in Streamlit and add SVG download button.
    Wrapper function to reduce code repetition.
    
    Parameters
    ----------
    fig : matplotlib.figure.Figure
        The figure to display and download
    filename : str
        Base filename for the download (without extension)
    cache_key : str, optional
        Unique cache key for session state. If None, uses filename.
    """
    if 'st' in globals():
        st.pyplot(fig)
        download_plot_as_svg(fig, filename, cache_key)
    else:
        plt.show()


def load_pickle_files(sleep_stage='N1', base_dir='pickles_sleep_stage/EDF'):
    """
    Load all pickle files from the specified sleep stage directory.
    
    Parameters
    ----------
    sleep_stage : str
        Sleep stage to load (N1, N2, N3, R, W, or 'All' to load from all stages)
    base_dir : str
        Base directory containing sleep stage subdirectories
    
    Returns
    -------
    list : List of tuples (file_path, patient_id, raw_object, sleep_stage)
    """
    loaded_files = []
    
    # If "All" is selected, load from all sleep stages
    if sleep_stage == 'All':
        sleep_stages = ['N1', 'N2', 'N3', 'R', 'W']
        total_files = 0
        
        for stage in sleep_stages:
            pickle_dir = os.path.join(base_dir, stage)
            
            if not os.path.exists(pickle_dir):
                continue
            
            pickle_files = [f for f in os.listdir(pickle_dir) if f.endswith('.pkl')]
            
            if not pickle_files:
                continue
            
            for pickle_file in pickle_files:
                file_path = os.path.join(pickle_dir, pickle_file)
                
                # Extract patient ID from filename
                patient_id = pickle_file.replace('.pkl', '').split('.edf')[0]
                
                try:
                    with open(file_path, 'rb') as f:
                        raw = pickle.load(f)
                    
                    # Add sleep stage to the tuple for tracking
                    loaded_files.append((file_path, patient_id, raw, stage))
                    total_files += 1
                except Exception as e:
                    if 'st' in globals():
                        st.warning(f"Error loading {pickle_file} from {stage}: {e}")
                    else:
                        print(f"Error loading {pickle_file} from {stage}: {e}")
        
        if 'st' in globals():
            st.info(f"Found {total_files} pickle files across all sleep stages")
        else:
            print(f"Found {total_files} pickle files across all sleep stages")
        
        return loaded_files
    
    # Single sleep stage loading (original behavior)
    pickle_dir = os.path.join(base_dir, sleep_stage)
    
    if not os.path.exists(pickle_dir):
        if 'st' in globals():
            st.error(f"Directory not found: {pickle_dir}")
        else:
            print(f"Directory not found: {pickle_dir}")
        return []
    
    pickle_files = [f for f in os.listdir(pickle_dir) if f.endswith('.pkl')]
    
    if not pickle_files:
        if 'st' in globals():
            st.warning(f"No pickle files found in {pickle_dir}")
        else:
            print(f"No pickle files found in {pickle_dir}")
        return []
    
    if 'st' in globals():
        st.info(f"Found {len(pickle_files)} pickle files in {pickle_dir}")
    else:
        print(f"Found {len(pickle_files)} pickle files in {pickle_dir}")
    
    for pickle_file in pickle_files:
        file_path = os.path.join(pickle_dir, pickle_file)
        
        # Extract patient ID from filename
        patient_id = pickle_file.replace('.pkl', '').split('.edf')[0]
        
        try:
            with open(file_path, 'rb') as f:
                raw = pickle.load(f)
            loaded_files.append((file_path, patient_id, raw))
        except Exception as e:
            if 'st' in globals():
                st.warning(f"Error loading {pickle_file}: {e}")
            else:
                print(f"Error loading {pickle_file}: {e}")
            continue
    
    if 'st' in globals():
        st.success(f"Successfully loaded {len(loaded_files)} files")
    else:
        print(f"Successfully loaded {len(loaded_files)} files")
    return loaded_files


def calculate_total_avg_bpm(loaded_files):
    """
    Calculate the total average BPM (beats per minute) across all loaded files.
    
    Parameters
    ----------
    loaded_files : list
        List of tuples (file_path, patient_id, raw_object)
    
    Returns
    -------
    float : Total average BPM across all files
    dict : Dictionary with per-file BPM values and summary statistics
    """
    all_bpms = []
    file_bpms = {}
    
    for file_path, patient_id, raw in loaded_files:
        try:
            # Get sampling frequency
            sfreq = raw.info['sfreq']
            
            # Get channel names and data
            ch_names = raw.ch_names
            data = raw.get_data()
            
            # Find ECG channel
            ch_lower = [ch.lower() for ch in ch_names]
            ecg_indices = [i for i, ch in enumerate(ch_lower) if 'ecg' in ch or 'ekg' in ch]
            
            if not ecg_indices:
                continue
            
            ecg_ch_idx = ecg_indices[0]
            ecg_signal = data[ecg_ch_idx, :]
            
            # Clean ECG signal
            try:
                # bandpass 0.5 - 5
                # ecg_signal = nk.ecg_bandpass(ecg_signal, sampling_rate=sfreq, low_cut=0.5, high_cut=20)
                # 
                ecg_clean = nk.ecg_clean(ecg_signal, sampling_rate=sfreq)
            except Exception:
                ecg_clean = ecg_signal
            
            # Detect R-peaks
            _, rpk = nk.ecg_peaks(ecg_clean, sampling_rate=sfreq)
            rpeaks = rpk['ECG_R_Peaks']
            
            if len(rpeaks) < 2:
                continue
            # Calculate BPM from R-R intervals by number of beats / total time
            hr_bpm = (len(rpeaks) / (len(ecg_signal) / sfreq)) * 60 # beats per minute
            # Calculate average BPM for this file
            all_bpms.append(hr_bpm)
            file_bpms[patient_id] = {
                'avg_bpm': hr_bpm,
                'num_beats': len(rpeaks),
                'duration_sec': len(ecg_signal) / sfreq
            }
            
        except Exception as e:
            if 'st' in globals():
                st.warning(f"Error calculating BPM for {patient_id}: {e}")
            else:
                print(f"Error calculating BPM for {patient_id}: {e}")
            continue
    
    # Calculate total average BPM
    if all_bpms:
        total_avg_bpm = np.mean(all_bpms)
        total_std_bpm = np.std(all_bpms)
        
        results = {
            'total_avg_bpm': total_avg_bpm,
            'total_std_bpm': total_std_bpm,
            'total_min_bpm': np.min(all_bpms),
            'total_max_bpm': np.max(all_bpms),
            'num_files': len(file_bpms),
            'file_bpms': file_bpms
        }
        return total_avg_bpm, results
    else:
        return None, {'total_avg_bpm': None, 'num_files': 0, 'file_bpms': {}}


def extract_spike_trigger_events(ecg_data, sfreq, before_event_sec, after_event_sec):
    """
    Extract spike-triggered averaging event windows from ECG data.
    
    For each R-peak detected in the ECG data, extracts a time window
    from [R-peak - before_event_sec] to [R-peak + after_event_sec],
    and returns a 2D array where each row represents one event.
    
    Parameters
    ----------
    ecg_data : np.ndarray
        1D array of ECG signal data
    sfreq : float
        Sampling frequency in Hz
    before_event_sec : float
        Time in seconds to sample before the R-peak event
    after_event_sec : float
        Time in seconds to sample after the R-peak event
    
    Returns
    -------
    np.ndarray
        2D array of shape (n_events, n_samples_per_event) where:
        - n_events: number of detected R-peaks with valid windows
        - n_samples_per_event: number of samples in each window
        Each row contains the ECG data for one event window
    np.ndarray
        Array of R-peak indices (sample indices) for each event
    dict
        Metadata including:
        - n_events: number of events extracted
        - n_samples_per_event: number of samples per event
        - valid_rpeaks: array of valid R-peak indices
        - skipped_events: number of events skipped (out of bounds)
    """
    # Convert time windows to sample indices
    before_event_samples = int(np.round(before_event_sec * sfreq))
    after_event_samples = int(np.round(after_event_sec * sfreq))
    n_samples_per_event = before_event_samples + after_event_samples + 1  # +1 for the R-peak itself
    
    _, rpk = nk.ecg_peaks(ecg_data, sampling_rate=sfreq)
    rpeaks = rpk['ECG_R_Peaks']
    
    if len(rpeaks) == 0:
        return np.array([]), np.array([]), {
            'n_events': 0,
            'n_samples_per_event': n_samples_per_event,
            'valid_rpeaks': np.array([]),
            'skipped_events': 0
        }
    
    # Filter R-peaks to ensure they're within valid range
    valid_rpeaks = []
    events = []
    skipped = 0
    
    for rpeak_idx in rpeaks:
        # Calculate window boundaries
        start_idx = rpeak_idx - before_event_samples
        end_idx = rpeak_idx + after_event_samples + 1  # +1 to include the end sample
        
        # Check if window is within data bounds
        if start_idx >= 0 and end_idx <= len(ecg_data):
            # Extract event window
            event_window = ecg_data[start_idx:end_idx]
            
            # Ensure all events have the same length (handle edge cases)
            if len(event_window) == n_samples_per_event:
                events.append(event_window)
                valid_rpeaks.append(rpeak_idx)
            else:
                skipped += 1
        else:
            skipped += 1
    
    if len(events) == 0:
        return np.array([]), np.array([]), {
            'n_events': 0,
            'n_samples_per_event': n_samples_per_event,
            'valid_rpeaks': np.array([]),
            'skipped_events': skipped
        }
    
    # Convert to 2D numpy array
    events_array = np.array(events)
    valid_rpeaks_array = np.array(valid_rpeaks)
    
    metadata = {
        'n_events': len(events_array),
        'n_samples_per_event': n_samples_per_event,
        'valid_rpeaks': valid_rpeaks_array,
        'skipped_events': skipped,
        'before_event_samples': before_event_samples,
        'after_event_samples': after_event_samples,
        'before_event_sec': before_event_sec,
        'after_event_sec': after_event_sec
    }
    
    return events_array, valid_rpeaks_array, metadata


def extract_systole_diastole_segments(raw, sfreq):
    """
    Extract systole (RST) and diastole (PQR) segments from ECG data.
    
    Parameters
    ----------
    raw : mne.io.Raw
        MNE Raw object with EEG/ECG data
    sfreq : float
        Sampling frequency
    
    Returns
    -------
    dict : Dictionary with 'systole' and 'diastole' keys, each containing:
        - 'eeg_segments': List of EEG data arrays during that phase
        - 'ecg_segments': List of ECG data arrays during that phase
        - 'indices': List of (start, end) sample index tuples
    """
    # Get channel names and data
    ch_names = raw.ch_names
    data = raw.get_data()
    
    # Find ECG channel
    ch_lower = [ch.lower() for ch in ch_names]
    ecg_indices = [i for i, ch in enumerate(ch_lower) if 'ecg' in ch or 'ekg' in ch]
    
    if not ecg_indices:
        return None
    
    ecg_ch_idx = ecg_indices[0]
    ecg_signal = data[ecg_ch_idx, :]
    
    # Detect R-peaks
    _, rpk = nk.ecg_peaks(ecg_signal, sampling_rate=sfreq)
    rpeaks = rpk['ECG_R_Peaks']
    
    if len(rpeaks) < 2:
        return None
    
    # Find EEG channels
    eeg_ch_idx = [i for i, ch in enumerate(ch_names) 
                  if 'eeg' in ch.lower() or any(elec in ch.upper() for elec in ['FP', 'F', 'C', 'P', 'O', 'T', 'A'])]
    
    if not eeg_ch_idx:
        return None
    
    eeg_data = data[eeg_ch_idx, :]
    
    # Extract segments
    systole_segments_eeg = []
    systole_segments_ecg = []
    systole_indices = []
    
    diastole_segments_eeg = []
    diastole_segments_ecg = []
    diastole_indices = []
    
    # For each R-R interval
    for i in range(len(rpeaks) - 1):
        r_start = rpeaks[i]
        r_end = rpeaks[i + 1]
        rr_length = r_end - r_start
        
        if rr_length < 10:  # Skip very short intervals
            continue
        
        # Systole: R to T (RST) - approximately first 40% of R-R interval (R peak to T wave)
        systole_end = r_start + int(0.4 * rr_length)
        if systole_end > r_start and systole_end < len(ecg_signal):
            systole_segments_eeg.append(eeg_data[:, r_start:systole_end])
            systole_segments_ecg.append(ecg_signal[r_start:systole_end])
            systole_indices.append((r_start, systole_end))
        
        # Diastole: P to R (PQR) - approximately last 30% of previous R-R interval (P wave to R peak)
        if i > 0:
            prev_r_start = rpeaks[i - 1]
            prev_rr_length = r_start - prev_r_start
            # P wave typically starts around 70-80% into the previous R-R interval
            diastole_start = prev_r_start + int(0.7 * prev_rr_length)
            diastole_end = r_start
            if diastole_start < diastole_end and diastole_end < len(ecg_signal) and diastole_start >= 0:
                diastole_segments_eeg.append(eeg_data[:, diastole_start:diastole_end])
                diastole_segments_ecg.append(ecg_signal[diastole_start:diastole_end])
                diastole_indices.append((diastole_start, diastole_end))
    
    # Return segments
    if systole_segments_eeg and diastole_segments_eeg:
        return {
            'systole': {
                'eeg_segments': systole_segments_eeg,
                'ecg_segments': systole_segments_ecg,
                'indices': systole_indices
            },
            'diastole': {
                'eeg_segments': diastole_segments_eeg,
                'ecg_segments': diastole_segments_ecg,
                'indices': diastole_indices
            }
        }
    
    return None


def compute_hep_from_segments(segments_dict, ch_names, sfreq):
    """
    Compute HEP (Heartbeat Evoked Potential) by averaging EEG segments.
    
    Parameters
    ----------
    segments_dict : dict
        Dictionary with 'systole' and 'diastole' keys containing segments
    ch_names : list
        List of EEG channel names
    sfreq : float
        Sampling frequency
    
    Returns
    -------
    dict : Dictionary with averaged HEP for systole and diastole
    """
    hep_results = {}
    
    for phase in ['systole', 'diastole']:
        if phase not in segments_dict:
            continue
        
        eeg_segments = segments_dict[phase]['eeg_segments']
        
        if not eeg_segments:
            continue
        
        # Find maximum length for alignment
        max_len = max(seg.shape[1] for seg in eeg_segments)
        
        # Align and average segments
        aligned_segments = []
        for seg in eeg_segments:
            # Pad shorter segments with NaN
            if seg.shape[1] < max_len:
                padded = np.full((seg.shape[0], max_len), np.nan)
                padded[:, :seg.shape[1]] = seg
                aligned_segments.append(padded)
            else:
                aligned_segments.append(seg)
        
        # Stack and compute mean (ignoring NaNs)
        stacked = np.stack(aligned_segments, axis=0)  # (n_segments, n_channels, n_timepoints)
        hep_mean = np.nanmean(stacked, axis=0)  # (n_channels, n_timepoints)
        hep_std = np.nanstd(stacked, axis=0)  # (n_channels, n_timepoints)
        
        # Create time vector (in seconds)
        time_vector = np.arange(hep_mean.shape[1]) / sfreq
        
        hep_results[phase] = {
            'mean': hep_mean,
            'std': hep_std,
            'time': time_vector,
            'n_segments': len(eeg_segments),
            'ch_names': ch_names[:hep_mean.shape[0]]
        }
    
    return hep_results


def plot_hr_with_segments(segments, raw, patient_id, sfreq, save_path=None, max_duration_sec=60):
    """
    Plot heart rate vs time with systole and diastole segments highlighted.
    
    Parameters
    ----------
    segments : dict
        Dictionary with 'systole' and 'diastole' segments containing indices
    raw : mne.io.Raw
        MNE Raw object with ECG data
    patient_id : str
        Patient identifier
    sfreq : float
        Sampling frequency
    save_path : str, optional
        Path to save the figure
    max_duration_sec : float
        Maximum duration to plot in seconds (default: 60)
    """
    # Get ECG channel
    ch_names = raw.ch_names
    ch_lower = [ch.lower() for ch in ch_names]
    ecg_indices = [i for i, ch in enumerate(ch_lower) if 'ecg' in ch or 'ekg' in ch]
    
    if not ecg_indices:
        print(f"No ECG channel found for {patient_id}")
        return
    
    ecg_ch_idx = ecg_indices[0]
    data = raw.get_data()
    ecg_signal = data[ecg_ch_idx, :]
    
    # Limit to max_duration_sec
    max_samples = int(max_duration_sec * sfreq)
    if len(ecg_signal) > max_samples:
        ecg_signal = ecg_signal[:max_samples]
    
    # Create time vector
    time_vector = np.arange(len(ecg_signal)) / sfreq
    
    # Calculate heart rate from R-peaks
    try:
        ecg_clean = nk.ecg_clean(ecg_signal, sampling_rate=sfreq)
    except Exception:
        ecg_clean = ecg_signal
    
    try:
        _, rpk = nk.ecg_peaks(ecg_clean, sampling_rate=sfreq)
        rpeaks = rpk['ECG_R_Peaks']
    except Exception:
        peaks, _ = find_peaks(ecg_clean, distance=int(sfreq * 0.3))
        rpeaks = peaks
    
    # Filter R-peaks within the time window
    rpeaks = rpeaks[rpeaks < len(ecg_signal)]
    r_times = rpeaks / sfreq
    
    # Calculate heart rate (beats per minute) from R-R intervals
    if len(rpeaks) > 1:
        rr_intervals = np.diff(rpeaks) / sfreq  # in seconds
        hr = 60.0 / rr_intervals  # beats per minute
        
        # Create HR time series (assign HR to midpoint between R-peaks)
        hr_times = (r_times[:-1] + r_times[1:]) / 2
    else:
        hr = np.array([])
        hr_times = np.array([])
    
    # Create figure
    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(14, 8), sharex=True)
    fig.suptitle(f'Heart Rate and Cardiac Phases: {patient_id}', fontsize=14, fontweight='bold')
    
    # Plot ECG signal
    ax1.plot(time_vector, ecg_signal, 'k-', linewidth=0.5, alpha=0.7, label='ECG Signal')
    ax1.scatter(r_times, ecg_signal[rpeaks], color='red', s=30, zorder=5, label='R-peaks')
    ax1.set_ylabel('ECG Amplitude (V)', fontsize=11)
    ax1.set_title('ECG Signal with R-peaks', fontsize=12)
    ax1.legend(loc='upper right')
    ax1.grid(True, alpha=0.3)
    
    # Plot heart rate
    if len(hr) > 0:
        ax2.plot(hr_times, hr, 'b-', linewidth=1.5, alpha=0.8, label='Heart Rate')
        ax2.set_ylabel('Heart Rate (BPM)', fontsize=11)
        ax2.set_xlabel('Time (s)', fontsize=11)
        ax2.set_title('Heart Rate vs Time with Systole/Diastole Phases', fontsize=12)
        ax2.legend(loc='upper right')
        ax2.grid(True, alpha=0.3)
        
        # Highlight systole segments
        systole_indices = segments.get('systole', {}).get('indices', [])
        systole_plotted = False
        for start_idx, end_idx in systole_indices:
            if end_idx < len(ecg_signal):
                start_time = start_idx / sfreq
                end_time = end_idx / sfreq
                if end_time <= max_duration_sec:
                    ax2.axvspan(start_time, end_time, alpha=0.3, color='red', 
                              label='Systole' if not systole_plotted else '')
                    systole_plotted = True
        
        # Highlight diastole segments
        diastole_indices = segments.get('diastole', {}).get('indices', [])
        diastole_plotted = False
        for start_idx, end_idx in diastole_indices:
            if end_idx < len(ecg_signal):
                start_time = start_idx / sfreq
                end_time = end_idx / sfreq
                if end_time <= max_duration_sec:
                    ax2.axvspan(start_time, end_time, alpha=0.3, color='green', 
                              label='Diastole' if not diastole_plotted else '')
                    diastole_plotted = True
        
        # Update legend to include phases
        ax2.legend(loc='upper right')
        
        # Set x-axis limit
        ax2.set_xlim([0, min(max_duration_sec, time_vector[-1])])
    else:
        ax2.text(0.5, 0.5, 'No R-peaks detected', ha='center', va='center', 
                transform=ax2.transAxes, fontsize=12)
        ax2.set_ylabel('Heart Rate (BPM)', fontsize=11)
        ax2.set_xlabel('Time (s)', fontsize=11)
    
    plt.tight_layout()
    
    if save_path:
        plt.savefig(save_path, dpi=300, bbox_inches='tight')
        if 'st' in globals():
            st.info(f"HR plot saved to {save_path}")
        else:
            print(f"HR plot saved to {save_path}")
    
    # Display in Streamlit if available
    display_plot_with_download(fig, f"hr_with_segments_{patient_id}", f"hr_segments_{patient_id}")
    plt.close(fig)


def process_all_files(loaded_files):
    """
    Process all loaded files to extract HEP for systole and diastole.
    
    Parameters
    ----------
    loaded_files : list
        List of tuples (file_path, patient_id, raw_object)
    
    Returns
    -------
    pd.DataFrame : DataFrame with HEP results for each patient
    """
    results = []
    
    for file_path, patient_id, raw in loaded_files:
        try:
            sfreq = raw.info['sfreq']
            
            # Extract systole and diastole segments
            segments = extract_systole_diastole_segments(raw, sfreq)
            
            if segments is None:
                if 'st' in globals():
                    st.warning(f"Could not extract segments for {patient_id}")
                else:
                    print(f"Could not extract segments for {patient_id}")
                continue
            
            # Plot HR with segments for the first patient as a sample
            if len(results) == 0:
                plot_hr_with_segments(segments, raw, patient_id, sfreq, 
                                     save_path=None)  # Don't save, display in Streamlit
            
            # Get EEG channel names
            ch_names = raw.ch_names
            eeg_ch_idx = [i for i, ch in enumerate(ch_names) 
                         if 'eeg' in ch.lower() or any(elec in ch.upper() for elec in ['FP', 'F', 'C', 'P', 'O', 'T', 'A'])]
            eeg_ch_names = [ch_names[i] for i in eeg_ch_idx]
            
            # Compute HEP
            hep_results = compute_hep_from_segments(segments, eeg_ch_names, sfreq)
            
            if not hep_results:
                continue
            
            # Extract features for comparison
            for phase in ['systole', 'diastole']:
                if phase not in hep_results:
                    continue
                
                hep_data = hep_results[phase]
                mean_hep = hep_data['mean']
                
                # Compute summary statistics across channels and time
                # Mean amplitude
                mean_amplitude = np.nanmean(mean_hep)
                
                # Peak amplitude (max absolute value)
                peak_amplitude = np.nanmax(np.abs(mean_hep))
                
                # RMS (root mean square)
                rms = np.sqrt(np.nanmean(mean_hep**2))
                
                # Standard deviation
                std_amplitude = np.nanstd(mean_hep)
                
                # Number of segments
                n_segments = hep_data['n_segments']
                
                results.append({
                    'patient_id': patient_id,
                    'phase': phase,
                    'mean_amplitude': mean_amplitude,
                    'peak_amplitude': peak_amplitude,
                    'rms': rms,
                    'std_amplitude': std_amplitude,
                    'n_segments': n_segments,
                    'n_channels': mean_hep.shape[0]
                })
        
        except Exception as e:
            print(f"Error processing {patient_id}: {e}")
            continue
    
    return pd.DataFrame(results)


def calculate_modulation_index(segments, raw, sfreq, n_phase_bins=2):
    """
    Calculate correlation and modulation index (MI) between systole/diastole phases 
    and EEG amplitude time series using Tort et al. (2010) method.
    
    Parameters
    ----------
    segments : dict
        Dictionary with 'systole' and 'diastole' segments containing indices
    raw : mne.io.Raw
        MNE Raw object with EEG data
    sfreq : float
        Sampling frequency
    n_phase_bins : int
        Number of phase bins (default: 2 for systole/diastole)
    
    Returns
    -------
    dict : Dictionary containing:
        - 'correlation': Pearson correlation between phase and amplitude
        - 'modulation_index': MI value (Tort et al. 2010)
        - 'mean_amplitude_systole': Mean EEG amplitude during systole
        - 'mean_amplitude_diastole': Mean EEG amplitude during diastole
        - 'amplitude_distribution': Distribution of amplitudes across phases
        - 'entropy': Entropy of amplitude distribution
        - 'max_entropy': Maximum possible entropy
    """
    # Get EEG channels
    ch_names = raw.ch_names
    eeg_ch_idx = [i for i, ch in enumerate(ch_names) 
                 if 'eeg' in ch.lower() or any(elec in ch.upper() for elec in ['FP', 'F', 'C', 'P', 'O', 'T', 'A'])]
    
    if not eeg_ch_idx:
        return None
    
    data = raw.get_data()
    eeg_data = data[eeg_ch_idx, :]  # (n_channels, n_timepoints)
    
    # Calculate amplitude (envelope) using Hilbert transform
    from scipy.signal import hilbert
    amplitude_time_series = []
    
    for ch_idx in range(eeg_data.shape[0]):
        # Get analytic signal using Hilbert transform
        analytic_signal = hilbert(eeg_data[ch_idx, :])
        # Get amplitude envelope
        amplitude = np.abs(analytic_signal)
        amplitude_time_series.append(amplitude)
    
    # Average across channels
    amplitude_ts = np.mean(amplitude_time_series, axis=0)
    
    # Create phase assignment array (0 = systole, 1 = diastole, -1 = neither)
    phase_assignment = np.full(len(amplitude_ts), -1, dtype=int)
    
    # Assign systole segments
    systole_indices = segments.get('systole', {}).get('indices', [])
    for start_idx, end_idx in systole_indices:
        if end_idx < len(phase_assignment):
            phase_assignment[start_idx:end_idx] = 0  # Systole
    
    # Assign diastole segments
    diastole_indices = segments.get('diastole', {}).get('indices', [])
    for start_idx, end_idx in diastole_indices:
        if end_idx < len(phase_assignment):
            # Only assign if not already assigned (systole takes precedence if overlap)
            mask = (phase_assignment[start_idx:end_idx] == -1)
            phase_assignment[start_idx:end_idx][mask] = 1  # Diastole
    
    # Extract amplitudes for each phase
    systole_mask = phase_assignment == 0
    diastole_mask = phase_assignment == 1
    
    systole_amplitudes = amplitude_ts[systole_mask]
    diastole_amplitudes = amplitude_ts[diastole_mask]
    
    if len(systole_amplitudes) == 0 or len(diastole_amplitudes) == 0:
        return None
    
    # Calculate mean amplitudes
    mean_amplitude_systole = np.mean(systole_amplitudes)
    mean_amplitude_diastole = np.mean(diastole_amplitudes)
    
    # Calculate correlation between phase and amplitude
    # Create arrays: phase (0 or 1) and corresponding amplitude
    phase_values = np.concatenate([
        np.zeros(len(systole_amplitudes)),
        np.ones(len(diastole_amplitudes))
    ])
    amplitude_values = np.concatenate([systole_amplitudes, diastole_amplitudes])
    
    if len(phase_values) > 1 and len(amplitude_values) > 1:
        correlation, p_value = stats.pearsonr(phase_values, amplitude_values)
    else:
        correlation, p_value = np.nan, np.nan
    
    # Calculate Modulation Index (MI) using Tort et al. (2010) method
    # MI = (H_max - H) / H_max
    # where H is the entropy of the amplitude distribution across phase bins
    # and H_max is the maximum entropy (uniform distribution)
    
    # Create amplitude distribution across phase bins
    # For systole/diastole, we have 2 bins
    amplitude_distribution = np.array([
        np.mean(systole_amplitudes),
        np.mean(diastole_amplitudes)
    ])
    
    # Normalize to get probability distribution
    amplitude_distribution_sum = np.sum(amplitude_distribution)
    if amplitude_distribution_sum > 0:
        p_distribution = amplitude_distribution / amplitude_distribution_sum
    else:
        p_distribution = np.array([0.5, 0.5])  # Uniform if no signal
    
    # Calculate entropy
    # Remove zeros to avoid log(0)
    p_distribution_clean = p_distribution[p_distribution > 0]
    if len(p_distribution_clean) > 0:
        H = entropy(p_distribution_clean, base=2)  # Base 2 for bits
    else:
        H = 0
    
    # Maximum entropy (uniform distribution)
    H_max = np.log2(n_phase_bins)
    
    # Modulation Index
    if H_max > 0:
        modulation_index = (H_max - H) / H_max
    else:
        modulation_index = 0
    
    return {
        'correlation': correlation,
        'correlation_p_value': p_value,
        'modulation_index': modulation_index,
        'mean_amplitude_systole': mean_amplitude_systole,
        'mean_amplitude_diastole': mean_amplitude_diastole,
        'amplitude_distribution': amplitude_distribution,
        'entropy': H,
        'max_entropy': H_max,
        'n_systole_samples': len(systole_amplitudes),
        'n_diastole_samples': len(diastole_amplitudes)
    }


def process_all_files_with_mi(loaded_files):
    """
    Process all loaded files to extract HEP and calculate modulation index.
    
    Parameters
    ----------
    loaded_files : list
        List of tuples (file_path, patient_id, raw_object)
    
    Returns
    -------
    tuple : (pd.DataFrame, pd.DataFrame)
        First DataFrame: HEP results
        Second DataFrame: Modulation index results
    """
    results = []
    mi_results = []
    
    for file_path, patient_id, raw in loaded_files:
        try:
            sfreq = raw.info['sfreq']
            
            # Extract systole and diastole segments
            segments = extract_systole_diastole_segments(raw, sfreq)
            
            if segments is None:
                print(f"Could not extract segments for {patient_id}")
                continue
            
            # Calculate modulation index
            mi_data = calculate_modulation_index(segments, raw, sfreq)
            
            if mi_data:
                mi_results.append({
                    'patient_id': patient_id,
                    **mi_data
                })
            
            # Get EEG channel names
            ch_names = raw.ch_names
            eeg_ch_idx = [i for i, ch in enumerate(ch_names) 
                         if 'eeg' in ch.lower() or any(elec in ch.upper() for elec in ['FP', 'F', 'C', 'P', 'O', 'T', 'A'])]
            eeg_ch_names = [ch_names[i] for i in eeg_ch_idx]
            
            # Compute HEP
            hep_results = compute_hep_from_segments(segments, eeg_ch_names, sfreq)
            
            if not hep_results:
                continue
            
            # Extract features for comparison
            for phase in ['systole', 'diastole']:
                if phase not in hep_results:
                    continue
                
                hep_data = hep_results[phase]
                mean_hep = hep_data['mean']
                
                # Compute summary statistics across channels and time
                mean_amplitude = np.nanmean(mean_hep)
                peak_amplitude = np.nanmax(np.abs(mean_hep))
                rms = np.sqrt(np.nanmean(mean_hep**2))
                std_amplitude = np.nanstd(mean_hep)
                n_segments = hep_data['n_segments']
                
                results.append({
                    'patient_id': patient_id,
                    'phase': phase,
                    'mean_amplitude': mean_amplitude,
                    'peak_amplitude': peak_amplitude,
                    'rms': rms,
                    'std_amplitude': std_amplitude,
                    'n_segments': n_segments,
                    'n_channels': mean_hep.shape[0]
                })
        
        except Exception as e:
            print(f"Error processing {patient_id}: {e}")
            continue
    
    return pd.DataFrame(results), pd.DataFrame(mi_results)


def compare_systole_diastole(df):
    """
    Perform statistical comparison between systole and diastole HEP.
    
    Parameters
    ----------
    df : pd.DataFrame
        DataFrame with HEP results
    
    Returns
    -------
    dict : Dictionary with statistical test results
    """
    if df.empty:
        return None
    
    # Separate systole and diastole data
    systole_df = df[df['phase'] == 'systole']
    diastole_df = df[df['phase'] == 'diastole']
    
    # Get common patients (paired comparison)
    common_patients = set(systole_df['patient_id']).intersection(set(diastole_df['patient_id']))
    
    if len(common_patients) < 2:
        if 'st' in globals():
            st.warning("Not enough paired data for comparison")
        else:
            print("Not enough paired data for comparison")
        return None
    
    # Prepare paired data
    systole_paired = systole_df[systole_df['patient_id'].isin(common_patients)].sort_values('patient_id')
    diastole_paired = diastole_df[diastole_df['patient_id'].isin(common_patients)].sort_values('patient_id')
    
    results = {}
    
    # Compare each metric
    metrics = ['mean_amplitude', 'peak_amplitude', 'rms', 'std_amplitude']
    
    for metric in metrics:
        systole_vals = systole_paired[metric].values
        diastole_vals = diastole_paired[metric].values
        
        # Remove NaN values
        valid_mask = ~(np.isnan(systole_vals) | np.isnan(diastole_vals))
        systole_vals = systole_vals[valid_mask]
        diastole_vals = diastole_vals[valid_mask]
        
        if len(systole_vals) < 2:
            continue
        
        # Paired t-test
        t_stat, p_value_ttest = stats.ttest_rel(diastole_vals, systole_vals)
        
        # Wilcoxon signed-rank test (non-parametric)
        try:
            w_stat, p_value_wilcoxon = stats.wilcoxon(diastole_vals, systole_vals)
        except Exception:
            p_value_wilcoxon = np.nan
        
        # Effect size (Cohen's d for paired samples)
        differences = diastole_vals - systole_vals
        cohens_d = np.mean(differences) / np.std(differences, ddof=1) if np.std(differences, ddof=1) > 0 else 0
        
        results[metric] = {
            'systole_mean': np.mean(systole_vals),
            'systole_std': np.std(systole_vals),
            'diastole_mean': np.mean(diastole_vals),
            'diastole_std': np.std(diastole_vals),
            'mean_difference': np.mean(differences),
            'p_value_ttest': p_value_ttest,
            'p_value_wilcoxon': p_value_wilcoxon,
            'cohens_d': cohens_d,
            'n_patients': len(systole_vals)
        }
    
    return results


def plot_comparison(df, stats_results, save_path=None):
    """
    Create visualization comparing systole and diastole HEP.
    
    Parameters
    ----------
    df : pd.DataFrame
        DataFrame with HEP results
    stats_results : dict
        Statistical test results
    save_path : str, optional
        Path to save the figure
    """
    if df.empty:
        print("No data to plot")
        return
    
    # Set style
    sns.set_style('whitegrid')
    plt.rcParams['figure.figsize'] = (15, 10)
    
    # Create figure with subplots
    fig, axes = plt.subplots(2, 2, figsize=(15, 10))
    # unique patients present in both phases (paired comparison)
    common_patients = set(df[df['phase'] == 'systole']['patient_id']).intersection(
        set(df[df['phase'] == 'diastole']['patient_id'])
    )
    fig.suptitle(f'HEP Comparison: Systole vs Diastole (n={len(common_patients)})', fontsize=16, fontweight='bold')
    
    metrics = ['mean_amplitude', 'peak_amplitude', 'rms', 'std_amplitude']
    metric_labels = ['Mean Amplitude (V)', 'Peak Amplitude (V)', 'RMS (V)', 'Std Amplitude (V)']
    
    for idx, (metric, label) in enumerate(zip(metrics, metric_labels)):
        ax = axes[idx // 2, idx % 2]
        
        # Prepare data for this metric
        systole_vals = df[df['phase'] == 'systole'][metric].dropna()
        diastole_vals = df[df['phase'] == 'diastole'][metric].dropna()
        
        # Box plot
        data_to_plot = [systole_vals, diastole_vals]
        bp = ax.boxplot(data_to_plot, labels=['Systole', 'Diastole'], patch_artist=True)
        
        # Color the boxes
        colors = ['lightblue', 'lightcoral']
        for patch, color in zip(bp['boxes'], colors):
            patch.set_facecolor(color)
        
        ax.set_ylabel(label)
        ax.set_title(f'{label.replace(" (V)", "")}')
        ax.grid(True, alpha=0.3)
        
        # Add statistical annotation if available
        if stats_results and metric in stats_results:
            stats_data = stats_results[metric]
            p_val = stats_data['p_value_wilcoxon']
            if not np.isnan(p_val):
                if p_val < 0.001:
                    sig_text = '***'
                elif p_val < 0.01:
                    sig_text = '**'
                elif p_val < 0.05:
                    sig_text = '*'
                else:
                    sig_text = 'ns'
                
                # Add significance line
                y_max = max(systole_vals.max(), diastole_vals.max())
                y_min = min(systole_vals.min(), diastole_vals.min())
                y_range = y_max - y_min
                
                ax.plot([1, 2], [y_max + 0.05 * y_range, y_max + 0.05 * y_range], 
                       'k-', linewidth=1.5)
                ax.text(1.5, y_max + 0.08 * y_range, sig_text, 
                       ha='center', fontsize=12, fontweight='bold')
                
                # Add p-value text
                ax.text(1.5, y_max + 0.12 * y_range, f'p = {p_val:.3e}', 
                       ha='center', fontsize=9)
    
    plt.tight_layout()
    
    if save_path:
        plt.savefig(save_path, dpi=300, bbox_inches='tight')
        if 'st' in globals():
            st.info(f"Figure saved to {save_path}")
        else:
            print(f"Figure saved to {save_path}")
    
    # Display in Streamlit if available
    display_plot_with_download(fig, "hep_comparison_systole_diastole", "hep_comparison_plot")
    plt.close(fig)


def print_statistics(stats_results):
    """
    Print statistical comparison results.
    
    Parameters
    ----------
    stats_results : dict
        Statistical test results
    """
    if not stats_results:
        if 'st' in globals():
            st.warning("No statistical results to display")
        else:
            print("No statistical results to display")
        return
    
    if 'st' in globals():
        st.markdown("## 📊 Statistical Comparison: Systole vs Diastole HEP")
        st.markdown("---")
    else:
        print("\n" + "="*80)
        print("STATISTICAL COMPARISON: Systole vs Diastole HEP")
        print("="*80)
    
    for metric, results in stats_results.items():
        metric_name = metric.upper().replace('_', ' ')
        
        if 'st' in globals():
            st.markdown(f"### {metric_name}")
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Systole", f"{results['systole_mean']:.4f} ± {results['systole_std']:.4f} V")
            with col2:
                st.metric("Diastole", f"{results['diastole_mean']:.4f} ± {results['diastole_std']:.4f} V")
            with col3:
                st.metric("Difference", f"{results['mean_difference']:.4f} V")
            
            col4, col5, col6 = st.columns(3)
            with col4:
                st.metric("Cohen's d", f"{results['cohens_d']:.4f}")
            with col5:
                p_val = results['p_value_wilcoxon']
                st.metric("Wilcoxon p-value", f"{p_val:.3e}")
            with col6:
                st.metric("N patients", f"{results['n_patients']}")
            
            # Significance interpretation
            p_val = results['p_value_wilcoxon']
            if not np.isnan(p_val):
                if p_val < 0.001:
                    sig = "*** (highly significant)"
                    st.success(f"Significance: {sig}")
                elif p_val < 0.01:
                    sig = "** (very significant)"
                    st.success(f"Significance: {sig}")
                elif p_val < 0.05:
                    sig = "* (significant)"
                    st.success(f"Significance: {sig}")
                else:
                    sig = "ns (not significant)"
                    st.info(f"Significance: {sig}")
            st.markdown("---")
        else:
            print(f"\n{metric_name}:")
            print(f"  Systole:   {results['systole_mean']:.4f} ± {results['systole_std']:.4f} V")
            print(f"  Diastole:  {results['diastole_mean']:.4f} ± {results['diastole_std']:.4f} V")
            print(f"  Difference: {results['mean_difference']:.4f} V")
            print(f"  Cohen's d:  {results['cohens_d']:.4f}")
            print(f"  Paired t-test p-value:     {results['p_value_ttest']:.3e}")
            print(f"  Wilcoxon test p-value:     {results['p_value_wilcoxon']:.3e}")
            print(f"  N patients: {results['n_patients']}")
            
            # Significance interpretation
            p_val = results['p_value_wilcoxon']
            if not np.isnan(p_val):
                if p_val < 0.001:
                    sig = "*** (highly significant)"
                elif p_val < 0.01:
                    sig = "** (very significant)"
                elif p_val < 0.05:
                    sig = "* (significant)"
                else:
                    sig = "ns (not significant)"
                print(f"  Significance: {sig}")
    
    if 'st' not in globals():
        print("\n" + "="*80)


def generate_pdf_report(loaded_files, sleep_stage, output_dir, save_results):
    """
    Generate a PDF report of the HEP analysis.
    
    Parameters
    ----------
    loaded_files : list
        List of loaded files
    sleep_stage : str
        Sleep stage being analyzed
    output_dir : str
        Output directory for saving PDF
    save_results : bool
        Whether to save results
    
    Returns
    -------
    bytes or None
        PDF file as bytes if successful, None otherwise
    """
    if not REPORTLAB_AVAILABLE:
        st.error("reportlab is not available. Install with: `pip install reportlab`")
        return None
    
    try:
        # Create PDF in memory
        buffer = BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4, 
                               rightMargin=72, leftMargin=72,
                               topMargin=72, bottomMargin=18)
        
        # Container for the 'Flowable' objects
        elements = []
        
        # Define styles
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1f77b4'),
            spaceAfter=30,
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor('#2c3e50'),
            spaceAfter=12,
            spaceBefore=12
        )
        
        # Title
        elements.append(Paragraph("HEP Comparison: Systole vs Diastole", title_style))
        elements.append(Spacer(1, 0.2*inch))
        
        # Summary information
        elements.append(Paragraph("Analysis Summary", heading_style))
        elements.append(Paragraph(f"<b>Sleep Stage:</b> {sleep_stage}", styles['Normal']))
        num_patients = len(loaded_files) if loaded_files else 0
        elements.append(Paragraph(f"<b>Number of Patients:</b> {num_patients}", styles['Normal']))
        elements.append(Paragraph(f"<b>Output Directory:</b> {output_dir}", styles['Normal']))
        elements.append(Spacer(1, 0.3*inch))
        
        # Add information about available data
        if loaded_files and len(loaded_files) > 0:
            elements.append(Paragraph("Patient IDs:", heading_style))
            patient_ids = [f[1] if len(f) > 1 else f"Patient_{i}" for i, f in enumerate(loaded_files)]
            for pid in patient_ids[:10]:  # Limit to first 10
                elements.append(Paragraph(f"• {pid}", styles['Normal']))
            if len(patient_ids) > 10:
                elements.append(Paragraph(f"... and {len(patient_ids) - 10} more", styles['Normal']))
            elements.append(Spacer(1, 0.2*inch))
        
        # Add note about plots
        elements.append(Paragraph("Note:", heading_style))
        elements.append(Paragraph(
            "This PDF contains a summary of the analysis. For detailed plots and visualizations, "
            "please refer to the interactive Streamlit interface or check the saved plot files in the output directory.",
            styles['Normal']
        ))
        elements.append(Spacer(1, 0.2*inch))
        
        # Add information about saved files if save_results is True
        if save_results:
            elements.append(Paragraph("Saved Files:", heading_style))
            elements.append(Paragraph(
                f"All analysis results and plots have been saved to: <b>{output_dir}</b>",
                styles['Normal']
            ))
            elements.append(Spacer(1, 0.2*inch))
        
        # Add timestamp
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        elements.append(Spacer(1, 0.3*inch))
        elements.append(Paragraph(f"<i>Report generated on: {timestamp}</i>", 
                                 ParagraphStyle('Timestamp', parent=styles['Normal'], 
                                               fontSize=9, textColor=colors.grey, alignment=TA_CENTER)))
        
        # Build PDF
        doc.build(elements)
        
        # Get PDF bytes
        pdf_bytes = buffer.getvalue()
        buffer.close()
        
        return pdf_bytes
        
    except Exception as e:
        st.error(f"Error generating PDF: {e}")
        import traceback
        traceback.print_exc()
        return None


def generate_pptx_report(loaded_files, sleep_stage, output_dir, save_results):
    """
    Generate a PPTX report of the HEP analysis.
    
    Parameters
    ----------
    loaded_files : list
        List of loaded files
    sleep_stage : str
        Sleep stage being analyzed
    output_dir : str
        Output directory for saving PPTX
    save_results : bool
        Whether to save results
    
    Returns
    -------
    bytes or None
        PPTX file as bytes if successful, None otherwise
    """
    if not PPTX_AVAILABLE:
        st.error("python-pptx is not available. Install with: `pip install python-pptx`")
        return None
    
    try:
        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "HEP Comparison: Systole vs Diastole"
        subtitle.text = f"Sleep Stage: {sleep_stage}\nAnalysis Report"
        
        # Summary slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title text box
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.8)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = "Analysis Summary"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(31, 119, 180)  # Blue color
        
        # Add content text box
        content_top = Inches(1.5)
        content_height = Inches(5)
        content_box = slide.shapes.add_textbox(left, content_top, width, content_height)
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # Add summary information
        num_patients = len(loaded_files) if loaded_files else 0
        summary_text = f"Sleep Stage: {sleep_stage}\n"
        summary_text += f"Number of Patients: {num_patients}\n"
        summary_text += f"Output Directory: {output_dir}\n\n"
        
        # Add patient IDs
        if loaded_files and len(loaded_files) > 0:
            summary_text += "Patient IDs:\n"
            patient_ids = [f[1] if len(f) > 1 else f"Patient_{i}" for i, f in enumerate(loaded_files)]
            for pid in patient_ids[:10]:  # Limit to first 10
                summary_text += f"• {pid}\n"
            if len(patient_ids) > 10:
                summary_text += f"... and {len(patient_ids) - 10} more\n"
            summary_text += "\n"
        
        # Add note about plots
        summary_text += "Note:\n"
        summary_text += "This presentation contains a summary of the analysis. "
        summary_text += "For detailed plots and visualizations, please refer to the interactive Streamlit interface "
        summary_text += "or check the saved plot files in the output directory.\n\n"
        
        # Add information about saved files
        if save_results:
            summary_text += "Saved Files:\n"
            summary_text += f"All analysis results and plots have been saved to: {output_dir}\n\n"
        
        # Add timestamp
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        summary_text += f"Report generated on: {timestamp}"
        
        # Set text content
        para = content_frame.paragraphs[0]
        para.text = summary_text
        para.font.size = Pt(12)
        para.font.color.rgb = RGBColor(0, 0, 0)
        
        # Save to BytesIO
        buffer = BytesIO()
        prs.save(buffer)
        pptx_bytes = buffer.getvalue()
        buffer.close()
        
        return pptx_bytes
        
    except Exception as e:
        st.error(f"Error generating PPTX: {e}")
        import traceback
        traceback.print_exc()
        return None


def main():
    """
    Main function to run HEP comparison analysis with Streamlit.
    """
    st.set_page_config(page_title="HEP Systole/Diastole Comparison", layout="wide")
    
    # Check and warn about pynapple availability after set_page_config()
    if not PYNAPPLE_AVAILABLE:
        st.warning("⚠️ pynapple is not available. PETH raster plot functionality will be disabled. Install with: `pip install pynapple`")
    
    st.title("🧠 HEP Comparison: Systole vs Diastole")
    st.markdown("Heartbeat Evoked Potential (HEP) analysis comparing systole and diastole phases")
    
    # Sidebar configuration
    st.sidebar.header("Configuration")
    sleep_stage = st.sidebar.selectbox(
        "Sleep Stage",
        options=['N1', 'N2', 'N3', 'R', 'W', 'All'],
        index=0,
        help="Select the sleep stage to analyze (or 'All' to analyze all stages combined)"
    )
    
    pickles_root = 'pickles_sleep_stage'
    if os.path.exists(pickles_root):
        folder_options = sorted([d for d in os.listdir(pickles_root) if os.path.isdir(os.path.join(pickles_root, d))])
    else:
        folder_options = ['EDF']
        
    selected_folder = st.sidebar.selectbox(
        "Base Directory Folder",
        options=folder_options,
        index=folder_options.index('EDF') if 'EDF' in folder_options else 0,
        help="Select the folder within pickles_sleep_stage"
    )
    base_dir = os.path.join(pickles_root, selected_folder)
    
    output_dir = st.sidebar.text_input(
        "Output Directory",
        value='hep_comparison_results',
        help="Directory to save results"
    )
    
    save_results = st.sidebar.checkbox("Save Results to Files", value=True)
    
    # PETH binsize configuration
    binsize = st.sidebar.number_input(
        "PETH Bin Size (seconds)",
        min_value=0.01,
        max_value=0.50,
        value=0.05,
        step=0.01,
        format="%.2f",
        help="Bin size for PETH rate calculation in seconds (default: 0.05s = 50ms)"
    )
    
    # Aggregation method configuration
    aggregation_method = st.sidebar.selectbox(
        "Aggregation Method",
        options=['median', 'mean'],
        index=0,  # Default to median
        help="Method for aggregating PETH rates and event densities across patients (default: median)"
    )
    
    # Test Type configuration
    test_type = st.sidebar.selectbox(
        "Significance Test Type",
        options=["Cluster-Based Permutation", "Circular Statistics (Rayleigh)"],
        index=0,
        help="Choose the statistical test for significance."
    )

    # Jitter configuration
    jitter_amount = st.sidebar.number_input(
        "Jitter Amount (seconds)",
        min_value=0.0,
        max_value=1.0,
        value=0.25,
        step=0.01,
        help="Amount of jitter to apply to R-peak timestamps for null hypothesis testing (default: 0.25)"
    )

    # Processing mode configuration
    processing_mode = st.sidebar.selectbox(
        "Processing Mode",
        options=['single', 'multithread'],
        index=1,  # Default to single
        help="Process patients sequentially (single) or in parallel using multithreading (multithread)"
    )
    
    # Create output directory
    os.makedirs(output_dir, exist_ok=True)
    
    # Display sleep stage info
    if sleep_stage == 'All':
        st.markdown("**Sleep Stage:** All (N1, N2, N3, R, W combined)")
    else:
        st.markdown(f"**Sleep Stage:** {sleep_stage}")
    st.markdown(f"**Base Directory:** {base_dir}")
    st.markdown("---")
    
    # Load pickle files
    with st.spinner("Loading pickle files..."):
        loaded_files_raw = load_pickle_files(sleep_stage=sleep_stage, base_dir=base_dir)
    
    if not loaded_files_raw:
        st.error("No files loaded. Please check the directory path.")
        return
    
    # Handle tuple format - preserve sleep_stage information for plotting
    # Check if files include sleep_stage (4 elements) or old format (3 elements)
    sleep_stage_mapping = {}  # {patient_id: sleep_stage} for tracking
    if loaded_files_raw and len(loaded_files_raw[0]) == 4:
        # New format: (file_path, patient_id, raw, sleep_stage)
        loaded_files = [(f[0], f[1], f[2]) for f in loaded_files_raw]
        # Create mapping of patient_id to sleep_stage
        for f in loaded_files_raw:
            sleep_stage_mapping[f[1]] = f[3]
        # Get unique sleep stages for display
        sleep_stages_in_data = list(set([f[3] for f in loaded_files_raw]))
        if sleep_stage == 'All' and len(sleep_stages_in_data) > 1:
            st.info(f"📊 Loaded {len(loaded_files)} files from {len(sleep_stages_in_data)} sleep stages: {', '.join(sorted(sleep_stages_in_data))}")
    else:
        # Old format: (file_path, patient_id, raw) - keep as is
        loaded_files = loaded_files_raw
        # If sleep_stage is not 'All', assign all files to that sleep_stage
        if sleep_stage != 'All':
            for f in loaded_files:
                sleep_stage_mapping[f[1]] = sleep_stage
    
    # PDF Download button in sidebar (after files are loaded)
    st.sidebar.markdown("---")
    st.sidebar.header("Export")
    
    # Generate PDF when button is clicked
    if st.sidebar.button("📥 Generate PDF Report", help="Generate and download a PDF report of the analysis"):
        with st.sidebar:
            with st.spinner("Generating PDF report..."):
                pdf_bytes = generate_pdf_report(loaded_files=loaded_files, sleep_stage=sleep_stage, 
                                                output_dir=output_dir, save_results=save_results)
                if pdf_bytes:
                    # Store PDF in session state
                    st.session_state['pdf_bytes'] = pdf_bytes
                    st.session_state['pdf_generated'] = True
                    st.success("PDF report generated!")
                else:
                    st.error("Failed to generate PDF report.")
    
    # Show download button if PDF is generated
    if 'pdf_generated' in st.session_state and st.session_state.get('pdf_generated', False):
        if 'pdf_bytes' in st.session_state:
            st.sidebar.download_button(
                label="⬇️ Download PDF",
                data=st.session_state['pdf_bytes'],
                file_name=f"HEP_Analysis_{sleep_stage}_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.pdf",
                mime="application/pdf",
                key="pdf_download_sidebar"
            )
    
    # Generate PPTX when button is clicked
    if st.sidebar.button("📊 Generate PPTX Report", help="Generate and download a PPTX report of the analysis"):
        with st.sidebar:
            with st.spinner("Generating PPTX report..."):
                pptx_bytes = generate_pptx_report(loaded_files=loaded_files, sleep_stage=sleep_stage, 
                                                  output_dir=output_dir, save_results=save_results)
                if pptx_bytes:
                    # Store PPTX in session state
                    st.session_state['pptx_bytes'] = pptx_bytes
                    st.session_state['pptx_generated'] = True
                    st.success("PPTX report generated!")
                else:
                    st.error("Failed to generate PPTX report.")
    
    # Show download button if PPTX is generated
    if 'pptx_generated' in st.session_state and st.session_state.get('pptx_generated', False):
        if 'pptx_bytes' in st.session_state:
            st.sidebar.download_button(
                label="⬇️ Download PPTX",
                data=st.session_state['pptx_bytes'],
                file_name=f"HEP_Analysis_{sleep_stage}_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key="pptx_download_sidebar"
            )
        
    st.markdown("## 📊 First Patient EEG-ECG Alignment")
    if st.button("Generate First Patient EEG-ECG Alignment Plot"):
        plot_first_patient_eeg_ecg_aligned(loaded_files, sleep_stage, output_dir, save_results, binsize=binsize, aggregation_method=aggregation_method, processing_mode=processing_mode, sleep_stage_mapping=sleep_stage_mapping)
    
    # Run BPM peri-event analysis (PETH raster plots)
    st.markdown("---")
    st.markdown("## 📊 BPM Peri-Event Analysis")
    BPM_peri_event_analysis(loaded_files, sleep_stage, output_dir, save_results, binsize=binsize, aggregation_method=aggregation_method, processing_mode=processing_mode, sleep_stage_mapping=sleep_stage_mapping, jitter_amount=jitter_amount, test_type=test_type)

    # Run HEP systole/diastole comparison analysis
    st.markdown("---")
    st.markdown("## 🧠 HEP Systole/Diastole Comparison Analysis")
    MI_systole_diastole_comparison(loaded_files, sleep_stage, output_dir, save_results)


def BPM_peri_event_analysis(loaded_files, sleep_stage, output_dir, save_results, binsize=0.05, aggregation_method='median', processing_mode='single', sleep_stage_mapping=None, jitter_amount=0.25, test_type="Cluster-Based Permutation"):
    """
    Run BPM peri-event analysis and create PETH raster plots.
    
    Parameters
    ----------
    binsize : float, optional
        Bin size for PETH rate calculation in seconds (default: 0.05)
    aggregation_method : str, optional
        Method for aggregating across patients: 'median' or 'mean' (default: 'median')
    processing_mode : str, optional
        Processing mode: 'single' for sequential or 'multithread' for parallel (default: 'single')
    sleep_stage_mapping : dict, optional
        Dictionary mapping patient_id to sleep_stage (default: None)
    jitter_amount : float, optional
         Amount of jitter to apply (default: 0.0)
    test_type: str, optional
        Test type (default: "Cluster-Based Permutation")
    """
    # Create PETH raster plot for the first file
    if loaded_files:
        significance_results = plot_peth_raster(loaded_files, file_index=0, minmax=(-0.5, 1.0), 
                        binsize=binsize, save_plot=save_results, save_dir=output_dir, 
                        aggregation_method=aggregation_method, processing_mode=processing_mode,
                        sleep_stage_mapping=sleep_stage_mapping, sleep_stage=sleep_stage, jitter_amount=jitter_amount, test_type=test_type)
        
        # Add download button for significance results
        if significance_results:
            # Convert significance_results to DataFrame
            significance_data = []
            for patient_id, patient_results in significance_results.items():
                for result in patient_results:
                    significance_data.append({
                        'patient_id': patient_id,
                        'electrode': result['electrode'],
                        'p_value': result['p_value'],
                        'mean_diff': result['mean_diff'],
                        'significant': result['significant']
                    })
            
            if significance_data:
                significance_df = pd.DataFrame(significance_data)
                
                # Cache CSV data in session state to prevent reruns
                cache_key = f"significance_csv_{sleep_stage}"
                if cache_key not in st.session_state:
                    st.session_state[cache_key] = significance_df.to_csv(index=False)
                
                # Create download button
                st.download_button(
                    label="📥 Download Electrode Significance Data (CSV)",
                    data=st.session_state[cache_key],
                    file_name=f"electrode_significance_{sleep_stage}_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.csv",
                    mime="text/csv",
                    key=f"download_significance_data_{sleep_stage}",
                    use_container_width=False
                )
        
        return significance_results


def clean_ecg_advanced(ecg_signal, sampling_rate, median_window_ms=300, 
                       wavelet='db4', wavelet_levels=6, zscore_threshold=3.0):
    """
    Advanced ECG signal cleaning using multiple techniques:
    1. Median filter to suppress spikes
    2. Wavelet denoising
    3. Z-score thresholding to remove extreme values
    
    Parameters
    ----------
    ecg_signal : np.ndarray
        1D array of ECG signal data
    sampling_rate : float
        Sampling frequency in Hz
    median_window_ms : float, optional
        Median filter window size in milliseconds (default: 300, range: 200-400)
    wavelet : str, optional
        Wavelet type for denoising: 'db4' or 'sym4' (default: 'db4')
    wavelet_levels : int, optional
        Number of decomposition levels for wavelet (default: 6, range: 5-8)
    zscore_threshold : float, optional
        Z-score threshold for removing extreme values (default: 3.0)
    
    Returns
    -------
    np.ndarray
        Cleaned ECG signal
    dict
        Dictionary with cleaning information including:
        - 'median_filtered': signal after median filter
        - 'wavelet_denoised': signal after wavelet denoising
        - 'n_extreme_samples': number of samples removed by z-score thresholding
        - 'methods_applied': list of methods successfully applied
    """
    if len(ecg_signal) == 0:
        return ecg_signal, {'methods_applied': [], 'n_extreme_samples': 0}
    
    methods_applied = []
    info = {}
    cleaned_signal = ecg_signal.copy()
    
    # Step 1: Median filter to suppress spikes
    # Convert window size from ms to samples
    median_window_samples = int(np.round(median_window_ms * sampling_rate / 1000.0))
    # Ensure window size is odd (required for median filter)
    if median_window_samples % 2 == 0:
        median_window_samples += 1
    # Ensure minimum window size of 3
    median_window_samples = max(3, median_window_samples)
    # Use scipy.ndimage.median_filter for better performance
    median_filtered = median_filter(cleaned_signal, size=median_window_samples)
    cleaned_signal = median_filtered
    methods_applied.append('median_filter')
    info['median_filtered'] = median_filtered
    info['median_window_samples'] = median_window_samples
    
    # Step 2: Wavelet denoising
    if PYWT_AVAILABLE:
        try:
            # Validate wavelet
            if wavelet not in ['db4', 'sym4']:
                wavelet = 'db4'  # Default to db4 if invalid
            
            # Validate levels
            max_level = pywt.dwt_max_level(len(cleaned_signal), wavelet)
            wavelet_levels = min(max(wavelet_levels, 5), min(8, max_level))
            
            # Perform wavelet decomposition
            coeffs = pywt.wavedec(cleaned_signal, wavelet, level=wavelet_levels)
            
            # Estimate noise standard deviation using median absolute deviation (MAD)
            # of the finest detail coefficients
            detail_coeffs = coeffs[-1]
            if len(detail_coeffs) > 0:
                sigma = np.median(np.abs(detail_coeffs)) / 0.6745  # MAD estimator
            else:
                sigma = np.std(detail_coeffs) if len(detail_coeffs) > 0 else 1.0
            
            # Apply soft thresholding to detail coefficients
            threshold = sigma * np.sqrt(2 * np.log(len(cleaned_signal)))
            coeffs_thresh = [coeffs[0]]  # Keep approximation coefficients
            for i in range(1, len(coeffs)):
                coeffs_thresh.append(pywt.threshold(coeffs[i], threshold, mode='soft'))
            
            # Reconstruct signal
            wavelet_denoised = pywt.waverec(coeffs_thresh, wavelet)
            
            # Ensure same length (wavelet reconstruction might add samples)
            if len(wavelet_denoised) > len(cleaned_signal):
                wavelet_denoised = wavelet_denoised[:len(cleaned_signal)]
            elif len(wavelet_denoised) < len(cleaned_signal):
                # Pad with last value if shorter
                padding = np.full(len(cleaned_signal) - len(wavelet_denoised), 
                                wavelet_denoised[-1])
                wavelet_denoised = np.concatenate([wavelet_denoised, padding])
            
            cleaned_signal = wavelet_denoised
            methods_applied.append('wavelet_denoising')
            info['wavelet_denoised'] = wavelet_denoised
            info['wavelet_used'] = wavelet
            info['wavelet_levels'] = wavelet_levels
            info['threshold'] = threshold
        except Exception as e:
            if 'st' in globals():
                st.warning(f"Wavelet denoising failed: {e}")
            else:
                print(f"Wavelet denoising failed: {e}")
    else:
        if 'st' in globals():
            st.warning("PyWavelets not available. Skipping wavelet denoising.")
        else:
            print("PyWavelets not available. Skipping wavelet denoising.")
    
    # Step 3: Z-score thresholding to remove extreme values
    try:
        # Calculate z-scores
        mean_signal = np.mean(cleaned_signal)
        std_signal = np.std(cleaned_signal)
        
        if std_signal > 0:
            z_scores = np.abs((cleaned_signal - mean_signal) / std_signal)
            
            # Find extreme values
            extreme_mask = z_scores > zscore_threshold
            n_extreme = np.sum(extreme_mask)
            
            if n_extreme > 0:
                # Replace extreme values with median of surrounding samples
                # or interpolate
                extreme_indices = np.where(extreme_mask)[0]
                
                # Create a copy for interpolation
                signal_interp = cleaned_signal.copy()
                
                # For each extreme value, replace with median of neighbors
                for idx in extreme_indices:
                    # Get neighborhood (avoid edges)
                    start_idx = max(0, idx - 5)
                    end_idx = min(len(cleaned_signal), idx + 6)
                    neighborhood = cleaned_signal[start_idx:end_idx]
                    # Exclude the extreme value itself
                    neighborhood = neighborhood[neighborhood != cleaned_signal[idx]]
                    if len(neighborhood) > 0:
                        signal_interp[idx] = np.median(neighborhood)
                    else:
                        # Fallback: use mean
                        signal_interp[idx] = mean_signal
                
                cleaned_signal = signal_interp
                methods_applied.append('zscore_thresholding')
                info['n_extreme_samples'] = n_extreme
                info['zscore_threshold'] = zscore_threshold
            else:
                info['n_extreme_samples'] = 0
        else:
            info['n_extreme_samples'] = 0
    except Exception as e:
        if 'st' in globals():
            st.warning(f"Z-score thresholding failed: {e}")
        else:
            print(f"Z-score thresholding failed: {e}")
        info['n_extreme_samples'] = 0
    
    info['methods_applied'] = methods_applied
    
    return cleaned_signal, info


def _process_patient_for_averaged_peth(file_data, minmax, binsize):
    """
    Helper function to process a single patient for averaged PETH.
    
    Parameters
    ----------
    file_data : tuple
        (file_path, patient_id, raw_object)
    minmax : tuple
        Time window around events in seconds
    binsize : float
        Bin size for rate calculation in seconds
    
    Returns
    -------
    tuple : (patient_id, channel_data, eeg_ch_names) or None if processing fails
    """
    file_path, patient_id, raw = file_data
    try:
        # Get sampling frequency
        sfreq = raw.info['sfreq']
        
        # Get channel names and data
        ch_names = raw.ch_names
        data = raw.get_data()
        
        # Find ECG channel
        ch_lower = [ch.lower() for ch in ch_names]
        ecg_indices = [i for i, ch in enumerate(ch_lower) if 'ecg' in ch or 'ekg' in ch]
        
        if not ecg_indices:
            return None
        
        ecg_ch_idx = ecg_indices[0]
        ecg_signal = data[ecg_ch_idx, :]
        
        # Clean ECG and detect R-peaks
        try:
            ecg_signal_clean, _ = clean_ecg_advanced(
                ecg_signal, 
                sampling_rate=sfreq,
                median_window_ms=300,
                wavelet='db4',
                wavelet_levels=5,
                zscore_threshold=3.0
            )
            _, rpk = nk.ecg_peaks(ecg_signal_clean, sampling_rate=sfreq)
            rpeaks = rpk['ECG_R_Peaks']
        except Exception:
            return None
        
        if len(rpeaks) < 2:
            return None
        
        # Convert R-peak indices to timestamps
        rpeak_times = rpeaks / sfreq
        rpeak_ts = nap.Ts(t=rpeak_times, time_units="s")
        
        # Find EEG channels
        eeg_ch_indices = [i for i, ch in enumerate(ch_names) 
                        if 'eeg' in ch.lower() or any(elec in ch.upper() for elec in ['FP', 'F', 'C', 'P', 'O', 'T', 'A'])]
        
        if not eeg_ch_indices:
            return None
        
        # Limit to first 10 channels for consistency
        eeg_ch_indices = eeg_ch_indices[:10]
        eeg_ch_names = [ch_names[i] for i in eeg_ch_indices]
        
        # Get EEG data
        eeg_data = data[eeg_ch_indices, :]
        
        # Process each EEG channel
        channel_data = {}  # {channel_name: {'rates_list': [], 'events_list': []}}
        
        for ch_idx, ch_name in enumerate(eeg_ch_names):
            try:
                eeg_signal = eeg_data[ch_idx, :]
                
                # Convert continuous EEG signal to pynapple Tsd object
                time_index = np.arange(len(eeg_signal)) / sfreq
                eeg_signal_tsd = nap.Tsd(t=time_index, d=eeg_signal, time_units="s")
                
                # Compute perievent alignment around R-peaks using continuous signal
                peth_eeg = nap.compute_perievent_continuous(
                    timeseries=eeg_signal_tsd,
                    tref=rpeak_ts,
                    minmax=minmax,
                    time_unit="s"
                )
                
                if len(peth_eeg) > 0:
                        # Initialize channel data structure
                        if ch_name not in channel_data:
                            channel_data[ch_name] = {'rates_list': [], 'events_list': []}
                        
                        # Calculate rates for this patient
                        try:
                            peth_eeg_count = peth_eeg.count(binsize)
                            if hasattr(peth_eeg_count, 'values'):
                                count_vals = peth_eeg_count.values
                                if count_vals.ndim == 2:
                                    rates = np.mean(count_vals, axis=1) / binsize
                                else:
                                    rates = count_vals / binsize
                                
                                if hasattr(peth_eeg_count, 'index'):
                                    time_axis = peth_eeg_count.index.values
                                else:
                                    time_axis = np.arange(len(rates)) * binsize + minmax[0]
                                
                                channel_data[ch_name]['rates_list'].append((time_axis, rates))
                        except Exception:
                            pass
                        
                        # Collect events for raster plot
                        try:
                            peth_eeg_tsd = peth_eeg.to_tsd()
                            channel_data[ch_name]['events_list'].append({
                                'times': peth_eeg_tsd.index.values,
                                'events': peth_eeg_tsd.values
                            })
                        except Exception:
                            pass
            
            except Exception:
                continue
        
        return (patient_id, channel_data, eeg_ch_names)
    
    except Exception:
        return None

def plot_regional_averaged_eeg_peth_raster(loaded_files, minmax=(-0.5, 1.0), binsize=0.005, save_plot=False, save_dir=None, aggregation_method='median', processing_mode='single', sleep_stage_mapping=None, sleep_stage='N1'):
    """
    Create averaged PETH plots for each electrode across ALL patients.
    Uses continuous EEG signal averaging.
    """
    if not PYNAPPLE_AVAILABLE:
        if 'st' in globals():
            st.error("pynapple is not available. Please install it to use regional PETH plots.")
        else:
            print("pynapple is not available. Please install it to use regional PETH plots.")
        return []

    if not loaded_files:
        return []

    # Map for all possible electrodes we care about to collect data
    regional_electrodes = {
        'F': ['Fp1', 'Fpz', 'Fp2', 'F7', 'F3', 'Fz', 'F4', 'F8'],
        'P': ['P3', 'Pz', 'P4'],
        'T': ['T3', 'T4', 'T5', 'T6'],
        'O': ['O1', 'Oz', 'O2'],
        'C': ['C3', 'Cz', 'C4']
    }
    # Flatten list of electrodes of interest
    electrodes_of_interest = set()
    for region_list in regional_electrodes.values():
        for elec in region_list:
            electrodes_of_interest.add(elec.upper())

    # Data structure: {electrode: [(time_axis, mean_signal), ...]}
    all_electrode_data = {}
    for elec in electrodes_of_interest:
        all_electrode_data[elec] = []

    if 'st' in globals():
        progress_bar = st.progress(0)
        status_text = st.empty()

    total_files = len(loaded_files)
    
    # Process all patients
    for idx, (file_path, patient_id, raw) in enumerate(loaded_files):
        if 'st' in globals():
            status_text.text(f"Processing patient {idx+1}/{total_files}: {patient_id}...")
            progress_bar.progress((idx) / total_files)

        try:
            # Get sampling frequency
            sfreq = raw.info['sfreq']
            ch_names = raw.ch_names
            data = raw.get_data()

            # Find ECG channel
            ch_lower = [ch.lower() for ch in ch_names]
            ecg_indices = [i for i, ch in enumerate(ch_lower) if 'ecg' in ch or 'ekg' in ch]

            if not ecg_indices:
                continue

            ecg_ch_idx = ecg_indices[0]
            ecg_signal = data[ecg_ch_idx, :]

            # Clean ECG and detect R-peaks
            try:
                ecg_signal_clean, _ = clean_ecg_advanced(
                    ecg_signal, 
                    sampling_rate=sfreq,
                    median_window_ms=300,
                    wavelet='db4',
                    wavelet_levels=5,
                    zscore_threshold=3.0
                )
                _, rpk = nk.ecg_peaks(ecg_signal_clean, sampling_rate=sfreq)
                rpeaks = rpk['ECG_R_Peaks']
            except Exception:
                continue

            if len(rpeaks) < 2:
                continue

            rpeak_times = rpeaks / sfreq
            rpeak_ts = nap.Ts(t=rpeak_times, time_units="s")

            # Channel mapping
            ch_name_to_idx = {ch.upper(): i for i, ch in enumerate(ch_names)}

            # Process each electrode of interest if present
            for elec_upper in electrodes_of_interest:
                if elec_upper not in ch_name_to_idx:
                    continue
                
                # finding original name
                _ = next((ch for ch in ch_names if ch.upper() == elec_upper), elec_upper)
                
                ch_idx = ch_name_to_idx[elec_upper]
                
                try:
                    eeg_signal = data[ch_idx, :]
                    time_index = np.arange(len(eeg_signal)) / sfreq
                    eeg_signal_tsd = nap.Tsd(t=time_index, d=eeg_signal, time_units="s")

                    peth_eeg = nap.compute_perievent_continuous(
                        timeseries=eeg_signal_tsd,
                        tref=rpeak_ts,
                        minmax=minmax,
                        time_unit="s"
                    )

                    if len(peth_eeg) > 0:
                        if hasattr(peth_eeg, 'values'):
                            values = peth_eeg.values
                            if values.ndim == 2:
                                mean_signal = np.nanmean(values, axis=1)
                            else:
                                mean_signal = values
                            
                            if hasattr(peth_eeg, 'index'):
                                time_axis = peth_eeg.index.values
                            else:
                                time_axis = np.linspace(minmax[0], minmax[1], len(mean_signal))
                            
                            # Find display name in regional_electrodes values
                            display_name = elec_upper
                            for r_list in regional_electrodes.values():
                                for r_elec in r_list:
                                    if r_elec.upper() == elec_upper:
                                        display_name = r_elec
                                        break
                            
                            if display_name in all_electrode_data:
                                all_electrode_data[display_name].append((time_axis, mean_signal))

                except Exception:
                    continue

        except Exception:
            continue

    if 'st' in globals():
        progress_bar.empty()
        status_text.empty()

    # Create plots
    figures = []
    
    # Filter electrodes with data
    electrodes_with_data = [e for e in all_electrode_data.keys() if all_electrode_data[e]]
    
    # Group electrodes by region
    electrodes_by_region = {}
    for region, region_elecs in regional_electrodes.items():
        electrodes_by_region[region] = [e for e in electrodes_with_data if e.upper() in [re.upper() for re in region_elecs]]

    for region, region_elecs in electrodes_by_region.items():
        if not region_elecs:
            continue

        n_electrodes = len(region_elecs)
        # Single column layout
        n_cols = 1
        n_rows = n_electrodes

        fig_region, axes = plt.subplots(n_rows, n_cols, figsize=(8, 4 * n_rows), constrained_layout=True)
        
        if n_electrodes > 1:
            axes_flat = axes.flatten()
        else:
            axes_flat = [axes]
        
        # Hide unused
        for i in range(n_electrodes, len(axes_flat)):
            axes_flat[i].set_visible(False)

        for elec_idx, electrode in enumerate(region_elecs):
            ax = axes_flat[elec_idx]
            electrode_data_list = all_electrode_data.get(electrode, [])
            
            if electrode_data_list:
                try:
                    all_time_axes = [t for t, _ in electrode_data_list]
                    all_signals = [s for _, s in electrode_data_list]
                    
                    min_time = max([t[0] for t in all_time_axes])
                    max_time = min([t[-1] for t in all_time_axes])
                    common_time = np.arange(min_time, max_time + binsize, binsize)
                    
                    interpolated_signals = []
                    for time_axis, signal in zip(all_time_axes, all_signals):
                        if len(time_axis) > 1 and len(signal) > 1:
                            interp_sig = np.interp(common_time, time_axis, signal)
                            interpolated_signals.append(interp_sig)
                    
                    if interpolated_signals:
                        n_patients = len(interpolated_signals)
                        interpolated_signals_array = np.array(interpolated_signals)
                        
                        if aggregation_method == 'median':
                            agg_signal = np.median(interpolated_signals_array, axis=0)
                            spread_lower, spread_upper = bootstrap_median_ci(interpolated_signals_array, n_bootstrap=1000)
                            label_text = 'Median'
                            spread_label = '95% CI'
                        else:
                            agg_signal = np.mean(interpolated_signals_array, axis=0)
                            std_signal = np.std(interpolated_signals_array, axis=0)
                            if n_patients > 1:
                                se_signal = std_signal / np.sqrt(n_patients)
                                t_critical = stats.t.ppf(0.975, df=n_patients - 1)
                                spread_lower = agg_signal - t_critical * se_signal
                                spread_upper = agg_signal + t_critical * se_signal
                                spread_label = '95% CI'
                            else:
                                spread_lower = agg_signal - std_signal
                                spread_upper = agg_signal + std_signal
                                spread_label = '±1 SD'
                            label_text = 'Mean'
                        
                        ax.plot(common_time, agg_signal, linewidth=2, color='blue', label=label_text)
                        ax.fill_between(common_time, spread_lower, spread_upper, alpha=0.3, color='blue', label=spread_label)
                        
                        ax.set_title(f"{electrode} (n={n_patients})", fontsize=11, fontweight='bold')
                        ax.set_ylabel("Amplitude (V)", fontsize=10)
                        ax.set_xlabel("Time from R-peak (s)", fontsize=10)
                        ax.axvline(0.0, color='red', linestyle='--', linewidth=1, alpha=0.5)
                        ax.grid(True, alpha=0.3)
                        ax.set_xlim(minmax)
                        ax.legend(fontsize=8, loc='upper right')
                    else:
                        ax.text(0.5, 0.5, 'No valid interpolated data', ha='center', va='center', transform=ax.transAxes)
                except Exception as e:
                    ax.text(0.5, 0.5, f'Error: {str(e)[:20]}', ha='center', va='center', transform=ax.transAxes, color='red')
            else:
                ax.text(0.5, 0.5, 'No data', ha='center', va='center', transform=ax.transAxes)
        
        # Add overarching title
        fig_region.suptitle(f"{region} Region - Averaged HEP Waveforms (All Patients)", fontsize=14, fontweight='bold')
        
        figures.append(fig_region)
        
        # Save if requested
        if save_plot and save_dir:
            save_path = os.path.join(save_dir, f'regional_all_patients_{region}.png')
            fig_region.savefig(save_path, bbox_inches='tight')
            if 'st' in globals():
                st.success(f"Saved {save_path}")
            
        display_plot_with_download(fig_region, f"all_patients_{region}", f"all_patients_{region}")
        plt.close(fig_region)

    return figures
 

                    




def bootstrap_median_ci(rates_array, n_bootstrap=1000, confidence=0.95):
    """
    Calculate 95% confidence interval for median using bootstrap percentile method.
    
    Parameters
    ----------
    rates_array : np.ndarray
        Array of shape (n_patients, n_timepoints) containing rates for each patient
    n_bootstrap : int
        Number of bootstrap iterations (default: 1000)
    confidence : float
        Confidence level (default: 0.95 for 95% CI)
    
    Returns
    -------
    lower_bound : np.ndarray
        Lower bound of CI for each timepoint
    upper_bound : np.ndarray
        Upper bound of CI for each timepoint
    """
    n_patients = len(rates_array)
    
    if n_patients == 1:
        # Single patient: use MAD as fallback
        median_rates = np.median(rates_array, axis=0)
        mad_rates = np.median(np.abs(rates_array - median_rates), axis=0)
        return median_rates - mad_rates, median_rates + mad_rates
    
    # Bootstrap resampling
    bootstrap_medians = []
    for _ in range(n_bootstrap):
        # Resample patients with replacement
        bootstrap_sample = rates_array[np.random.choice(n_patients, size=n_patients, replace=True)]
        bootstrap_median = np.median(bootstrap_sample, axis=0)
        bootstrap_medians.append(bootstrap_median)
    
    bootstrap_medians = np.array(bootstrap_medians)
    
    # Calculate percentiles for confidence interval
    alpha = 1 - confidence
    lower_percentile = (alpha / 2) * 100
    upper_percentile = (1 - alpha / 2) * 100
    
    lower_bound = np.percentile(bootstrap_medians, lower_percentile, axis=0)
    upper_bound = np.percentile(bootstrap_medians, upper_percentile, axis=0)
    
    return lower_bound, upper_bound


def bootstrap_pvalue_test(rates_list, test_window=(-0.1, 0.1), baseline_window=(-0.5, -0.5), n_bootstrap=100):
    """
    Perform bootstrapping test to compare event rates in test window vs baseline window.
    
    Parameters
    ----------
    rates_list : list
        List of (time_axis, rates) tuples from different patients
    test_window : tuple
        Time window around R-peak to test (default: (-0.1, 0.1))
    baseline_window : tuple
        Baseline time window before R-peak (default: (-0.5, -0.2))
    n_bootstrap : int
        Number of bootstrap iterations (default: 10000)
    
    Returns
    -------
    float : p-value from bootstrapping test
    float : observed mean difference (test - baseline)
    """
    if not rates_list or len(rates_list) < 2:
        return np.nan, np.nan
    
    # Collect differences for each patient
    differences = []
    
    for time_axis, rates in rates_list:
        if len(time_axis) < 2 or len(rates) < 2:
            continue
        
        # Find indices for test and baseline windows
        test_mask = (time_axis >= test_window[0]) & (time_axis <= test_window[1])
        baseline_mask = (time_axis >= baseline_window[0]) & (time_axis <= baseline_window[1])
        
        if np.sum(test_mask) == 0 or np.sum(baseline_mask) == 0:
            continue
        
        # Calculate mean rates in each window
        test_rate = np.mean(rates[test_mask])
        baseline_rate = np.mean(rates[baseline_mask])
        
        # Calculate difference
        diff = test_rate - baseline_rate
        differences.append(diff)
    
    if len(differences) < 2:
        return np.nan, np.nan
    
    differences = np.array(differences)
    observed_mean_diff = np.mean(differences)
    
    # Bootstrapping test: test if mean difference is significantly different from zero
    # Null hypothesis: mean difference = 0
    n_samples = len(differences)
    
    # Center the differences to create null distribution (mean = 0)
    centered_differences = differences - observed_mean_diff
    
    # Bootstrap: resample from centered differences and calculate means using pynapple Randomization
    if not PYNAPPLE_AVAILABLE:
        # Fallback to numpy if pynapple is not available
        bootstrap_means = []
        for _ in range(n_bootstrap):
            bootstrap_sample = np.random.choice(centered_differences, size=n_samples, replace=True)
            bootstrap_mean = np.mean(bootstrap_sample)
            bootstrap_means.append(bootstrap_mean)
        bootstrap_means = np.array(bootstrap_means)
    else:
        # Use pynapple Randomization for bootstrap resampling
        # Create a Ts object with indices as timestamps for resampling
        indices_ts = nap.Ts(t=np.arange(n_samples, dtype=float), time_units="s")
        bootstrap_means = []
        for _ in range(n_bootstrap):
            # Resample timestamps (indices) using pynapple's resample_timestamps
            resampled_indices_ts = nap.resample_timestamps(indices_ts)
            # Extract timestamps and convert to integer indices
            # resample_timestamps returns timestamps within [0, n_samples) range
            resampled_timestamps = resampled_indices_ts.index.values
            # Convert to integer indices, ensuring they're within valid range
            resampled_indices = np.clip(np.round(resampled_timestamps).astype(int), 0, n_samples - 1)
            # Use indices to sample from centered_differences
            bootstrap_sample = centered_differences[resampled_indices]
            bootstrap_mean = np.mean(bootstrap_sample)
            bootstrap_means.append(bootstrap_mean)
        bootstrap_means = np.array(bootstrap_means)
    
    # Calculate p-value: proportion of bootstrap means that are as extreme or more extreme
    # Two-tailed test: how often do we get a mean as extreme as observed under null hypothesis?
    p_value = np.mean(np.abs(bootstrap_means) >= np.abs(observed_mean_diff))
    
    # Ensure p-value is at least 1/n_bootstrap (minimum resolution)
    p_value = max(p_value, 1.0 / n_bootstrap)
    
    return p_value, observed_mean_diff


def max_statistic_permutation_test(interpolated_rates, common_time, baseline_window=(-0.5, -0.2), 
                                   n_permutations=1000, alpha_levels=[0.05, 0.01], aggregation_method='median'):
    """
    Perform max-statistic permutation test to identify significant bins in PETH plots.
    Controls for multiple comparisons using the maximum statistic method.
    
    Parameters
    ----------
    interpolated_rates : np.ndarray
        Array of shape (n_patients, n_bins) with interpolated rates per patient per bin
    common_time : np.ndarray
        Time axis array corresponding to bins
    baseline_window : tuple, optional
        Time window defining baseline period (default: (-0.5, -0.2))
    n_permutations : int, optional
        Number of permutations to perform (default: 1000)
    alpha_levels : list, optional
        Significance levels to test (default: [0.05, 0.01])
    aggregation_method : str, optional
        Method for aggregation: 'median' or 'mean' (default: 'median')
    
    Returns
    -------
    dict : Dictionary containing:
        - 'p_values': Array of p-values per bin
        - 'significant_bins': Dictionary with alpha level as key, boolean array as value
        - 'observed_deviations': Array of observed median deviations per bin
        - 'baseline_median': Baseline median value
        - 'max_deviation_distribution': Array of maximum deviations from permutations
    """
    if interpolated_rates is None or len(interpolated_rates) == 0:
        return {
            'p_values': None,
            'significant_bins': {},
            'observed_deviations': None,
            'baseline_median': None,
            'max_deviation_distribution': None
        }
    
    # Convert to numpy array if it's a list
    if isinstance(interpolated_rates, list):
        # Handle list of tuples (rates, patient_id) format
        if len(interpolated_rates) > 0 and isinstance(interpolated_rates[0], tuple):
            rates_array = np.array([r for r, _ in interpolated_rates])
        else:
            rates_array = np.array(interpolated_rates)
    else:
        rates_array = interpolated_rates
    
    if rates_array.ndim != 2:
        return {
            'p_values': None,
            'significant_bins': {},
            'observed_deviations': None,
            'baseline_median': None,
            'max_deviation_distribution': None
        }
    
    n_patients, n_bins = rates_array.shape
    
    if n_patients < 2 or n_bins < 2:
        return {
            'p_values': None,
            'significant_bins': {},
            'observed_deviations': None,
            'baseline_median': None,
            'max_deviation_distribution': None
        }
    
    # Find baseline bins
    baseline_mask = (common_time >= baseline_window[0]) & (common_time <= baseline_window[1])
    baseline_indices = np.where(baseline_mask)[0]
    
    if len(baseline_indices) == 0:
        return {
            'p_values': None,
            'significant_bins': {},
            'observed_deviations': None,
            'baseline_median': None,
            'max_deviation_distribution': None
        }
    
    # Calculate baseline median across all patients and baseline bins
    baseline_rates = rates_array[:, baseline_indices]
    if aggregation_method == 'median':
        baseline_median = np.median(baseline_rates)
    else:  # mean
        baseline_median = np.mean(baseline_rates)
    
    # Calculate observed deviations for each bin
    if aggregation_method == 'median':
        observed_aggregated = np.median(rates_array, axis=0)
    else:  # mean
        observed_aggregated = np.mean(rates_array, axis=0)
    
    observed_deviations = observed_aggregated - baseline_median
    observed_abs_deviations = np.abs(observed_deviations)
    
    # Perform permutations
    max_deviations = []
    
    for perm_idx in range(n_permutations):
        # Resample patients with replacement (bootstrap approach)
        permuted_indices = np.random.choice(n_patients, size=n_patients, replace=True)
        permuted_rates = rates_array[permuted_indices, :]
        
        # Calculate aggregated rates for permuted data
        if aggregation_method == 'median':
            permuted_aggregated = np.median(permuted_rates, axis=0)
        else:  # mean
            permuted_aggregated = np.mean(permuted_rates, axis=0)
        
        # Calculate deviations from baseline for each bin
        permuted_deviations = permuted_aggregated - baseline_median
        permuted_abs_deviations = np.abs(permuted_deviations)
        
        # Record maximum absolute deviation across all bins
        max_deviation = np.max(permuted_abs_deviations)
        max_deviations.append(max_deviation)
    
    max_deviations = np.array(max_deviations)
    
    # Calculate p-values for each bin
    # P-value = proportion of permutations where max deviation >= observed deviation for this bin
    p_values = np.zeros(n_bins)
    for bin_idx in range(n_bins):
        observed_abs_dev = observed_abs_deviations[bin_idx]
        # Count how many permutations had max deviation >= this bin's observed deviation
        p_value = np.mean(max_deviations >= observed_abs_dev)
        # Ensure minimum resolution
        p_values[bin_idx] = max(p_value, 1.0 / n_permutations)
    
    # Determine significant bins for each alpha level
    significant_bins = {}
    for alpha in alpha_levels:
        significant_bins[alpha] = p_values < alpha
    
    return {
        'p_values': p_values,
        'significant_bins': significant_bins,
        'observed_deviations': observed_deviations,
        'baseline_median': baseline_median,
        'max_deviation_distribution': max_deviations
    }



# -----------------------------------------------------------------------------
# Circular Statistics Helpers
# -----------------------------------------------------------------------------

def calculate_kuiper_statistic(alpha1, alpha2):
    """
    Calculate Kuiper's test statistic V for two samples of circular data.
    alpha1, alpha2: arrays of angles in radians.
    """
    n1 = len(alpha1)
    n2 = len(alpha2)
    
    if n1 == 0 or n2 == 0:
        return 0.0
    
    # Sort and wrap to [0, 2pi)
    alpha1 = np.sort(np.mod(alpha1, 2*np.pi))
    alpha2 = np.sort(np.mod(alpha2, 2*np.pi))
    
    # Combine and sort to find all evaluation points
    all_alphas = np.concatenate([alpha1, alpha2])
    
    # Label the points: 0 for alpha1, 1 for alpha2
    labels = np.concatenate([np.zeros(n1), np.ones(n2)])
    
    # Sort labels based on angles
    sort_idx = np.argsort(all_alphas)
    sorted_labels = labels[sort_idx]
    
    # Calculate CDFs
    cdf1 = np.cumsum(sorted_labels == 0) / n1
    cdf2 = np.cumsum(sorted_labels == 1) / n2
    
    # Difference
    diff = cdf1 - cdf2
    
    D_plus = np.max(diff)
    D_minus = np.max(-diff)
    
    return D_plus + D_minus

def permutation_kuiper_test(alpha1, alpha2, n_permutations=1000):
    """
    Perform a permutation test using Kuiper's statistic.
    Compares two circular distributions.
    """
    # Remove NaNs
    alpha1 = alpha1[~np.isnan(alpha1)]
    alpha2 = alpha2[~np.isnan(alpha2)]
    
    if len(alpha1) < 5 or len(alpha2) < 5:
        return 1.0
    
    # Calculate observed statistic
    V_obs = calculate_kuiper_statistic(alpha1, alpha2)
    
    # Combine
    combined = np.concatenate([alpha1, alpha2])
    n1 = len(alpha1)
    
    interactions_greater = 0
    
    for _ in range(n_permutations):
        # Shuffle
        np.random.shuffle(combined)
        perm_samp1 = combined[:n1]
        perm_samp2 = combined[n1:]
        
        V_perm = calculate_kuiper_statistic(perm_samp1, perm_samp2)
        
        if V_perm >= V_obs:
            interactions_greater += 1
            
    return (interactions_greater + 1) / (n_permutations + 1)


def test_electrode_significance_circular(loaded_files, band=(8, 12), n_permutations=1000, 
                                        significance_threshold=0.05, processing_mode='single'):
    """
    Test if heart modulates EEG electrodes by comparing phase distributions at R-peaks 
    vs P, Q, S, T waves using circular statistics (Kuiper's Test via permutation).
    
    Parameters
    ----------
    loaded_files : list
        List of tuples (file_path, patient_id, raw_object)
    band : tuple, optional
        Frequency band for phase extraction (low, high) in Hz (default: (8, 12) Alpha)
        If None, uses broadband (not recommended for phase).
    n_permutations : int, optional
        Number of permutations for significance testing (default: 1000)
    significance_threshold : float, optional
        Threshold for significance (default: 0.05)
    processing_mode : str
        'single' or 'multithread' (default: 'single')
        
    Returns
    -------
    dict : Results per patient
    """
    if not loaded_files:
        return {}
    
    results = {}
    
    if 'st' in globals():
        progress_bar = st.progress(0)
        status_text = st.empty()
        
    for file_idx, (file_path, patient_id, raw) in enumerate(loaded_files):
        try:
            if 'st' in globals():
                status_text.text(f"Circular Stats: Processing {patient_id} ({file_idx + 1}/{len(loaded_files)})...")
                progress_bar.progress((file_idx + 1) / len(loaded_files))
            
            sfreq = raw.info['sfreq']
            ch_names = raw.ch_names
            data = raw.get_data()
            
            # 1. ECG Processing
            ecg_indices = [i for i, ch in enumerate(ch_names) if 'ecg' in ch.lower() or 'ekg' in ch.lower()]
            if not ecg_indices:
                continue
            ecg_signal = data[ecg_indices[0], :]
            
            # Clean and Peaks
            try:
                ecg_clean = nk.ecg_clean(ecg_signal, sampling_rate=sfreq)
                _, rpk = nk.ecg_peaks(ecg_clean, sampling_rate=sfreq)
                r_peaks = rpk['ECG_R_Peaks']
                
                # Delineate P, Q, S, T
                waves, _ = nk.ecg_delineate(ecg_clean, r_peaks, sampling_rate=sfreq, method='dwt')
                
                # Extract indices (some might be NaNs)
                indices = {
                    'R': r_peaks,
                    'P': np.array(waves['ECG_P_Peaks']),
                    'Q': np.array(waves['ECG_Q_Peaks']),
                    'S': np.array(waves['ECG_S_Peaks']),
                    'T': np.array(waves['ECG_T_Peaks'])
                }
                
                # Remove NaNs from indices
                for key in indices:
                    indices[key] = indices[key][~np.isnan(indices[key])].astype(int)
                    # Clip to data bounds
                    indices[key] = indices[key][(indices[key] >= 0) & (indices[key] < data.shape[1])]
                    
            except Exception as e:
                print(f"ECG processing failed for {patient_id}: {e}")
                continue
                
            # 2. EEG Processing
            eeg_indices = [i for i, ch in enumerate(ch_names) 
                          if 'eeg' in ch.lower() or any(elec in ch.upper() for elec in ['FP', 'F', 'C', 'P', 'O', 'T', 'A'])]
            
            if not eeg_indices:
                continue
            
            patient_results = []
            
            # Pre-filter EEG if band is provided
            eeg_data = data[eeg_indices, :]
            eeg_ch_names = [ch_names[i] for i in eeg_indices]
            
            if band is not None:
                # Use MNE filtering if possible, or scipy
                # raw.exclude... we already have numpy array 'eeg_data'
                # Let's use scipy butterworth
                nyq = 0.5 * sfreq
                low = band[0] / nyq
                high = band[1] / nyq
                b, a = butter(4, [low, high], btype='band')
                eeg_data_filtered = filtfilt(b, a, eeg_data, axis=1)
            else:
                eeg_data_filtered = eeg_data
                
            # Compute Hilbert Phase
            # analytic_signal = hilbert(eeg_data_filtered) # This might be memory intensive across all channels
            # Do per channel
            
            for ch_idx, ch_name in enumerate(eeg_ch_names):
                try:
                    sig = eeg_data_filtered[ch_idx, :]
                    analytic = hilbert(sig)
                    phase = np.angle(analytic) # Radians [-pi, pi]
                    
                    # Extract phases at events
                    phases = {}
                    for wave_name, wave_idxs in indices.items():
                        if len(wave_idxs) > 0:
                            phases[wave_name] = phase[wave_idxs]
                        else:
                            phases[wave_name] = np.array([])
                    
                    if len(phases['R']) < 10: # Minimum trials
                        continue
                        
                    # Compare R to P, Q, S, T
                    comparisons = {}
                    is_significant_any = False
                    
                    for wave in ['P', 'Q', 'S', 'T']:
                        if len(phases[wave]) < 5:
                            comparisons[wave] = np.nan
                            continue
                            
                        # Run Permutation Test
                        p_val = permutation_kuiper_test(phases['R'], phases[wave], n_permutations=n_permutations)
                        comparisons[wave] = p_val
                        
                        if p_val < significance_threshold: # Simple threshold, maybe should correct?
                            is_significant_any = True
                    
                    patient_results.append({
                        'electrode': ch_name,
                        'p_values': comparisons,
                        'significant': is_significant_any
                    })
                    
                except Exception as e:
                    print(f"Error channel {ch_name}: {e}")
                    continue
            
            results[patient_id] = patient_results
            
        except Exception as e:
            print(f"Error processing {patient_id}: {e}")
            continue
            
    return results


def test_electrode_significance_per_patient(loaded_files, minmax=(-0.5, 1.0), binsize=0.1, 
                                            test_window=(-0.1, 0.1), baseline_window=(-0.5, -0.2), 
                                            n_bootstrap=100, significance_threshold=0.05, processing_mode='single', jitter_amount=0.0):
    """
    Test significance of RR to channel peri-event analysis for each patient individually.
    For each patient, tests which electrodes show significant differences between test and baseline windows.
    
    Parameters
    ----------
    loaded_files : list
        List of tuples (file_path, patient_id, raw_object)
    minmax : tuple, optional
        Time window around events in seconds (default: (-0.5, 1.0))
    binsize : float, optional
        Bin size for rate calculation in seconds (default: 0.05)
    test_window : tuple, optional
        Time window around R-peak to test (default: (-0.1, 0.1))
    baseline_window : tuple, optional
        Baseline time window before R-peak (default: (-0.5, -0.2))
    n_bootstrap : int, optional
        Number of bootstrap iterations (default: 10000)
    significance_threshold : float, optional
        P-value threshold for significance (default: 0.05)
    processing_mode : str, optional
        Processing mode: 'single' for sequential or 'multithread' for parallel (default: 'single')
    jitter_amount : float, optional
        Amount of jitter to apply to timestamps in seconds (default: 0.0)
    
    Returns
    -------
    dict : Dictionary with patient_id as keys and lists of significant electrodes as values
    """
    if not PYNAPPLE_AVAILABLE:
        if 'st' in globals():
            st.error("pynapple is not available. Please install it to use electrode significance testing.")
        else:
            print("pynapple is not available. Please install it to use electrode significance testing.")
        return {}
    
    if not loaded_files:
        return {}
    
    results = {}  # {patient_id: [(electrode, p_value, mean_diff), ...]}
    results_full_time = {}
    if 'st' in globals():
        progress_bar = st.progress(0)
        status_text = st.empty()
    
    for file_idx, (file_path, patient_id, raw) in enumerate(loaded_files):
        try:
            if 'st' in globals():
                status_text.text(f"Testing significance for {patient_id} ({file_idx + 1}/{len(loaded_files)})...")
                progress_bar.progress((file_idx + 1) / len(loaded_files))
            
            # Get sampling frequency
            sfreq = raw.info['sfreq']
            
            # Get channel names and data
            ch_names = raw.ch_names
            data = raw.get_data()
            
            # Find ECG channel
            ch_lower = [ch.lower() for ch in ch_names]
            ecg_indices = [i for i, ch in enumerate(ch_lower) if 'ecg' in ch or 'ekg' in ch]
            
            if not ecg_indices:
                continue
            
            ecg_ch_idx = ecg_indices[0]
            ecg_signal = data[ecg_ch_idx, :]
            
            # Clean ECG and detect R-peaks
            try:
                ecg_signal_clean, _ = clean_ecg_advanced(
                    ecg_signal, 
                    sampling_rate=sfreq,
                    median_window_ms=300,
                    wavelet='db4',
                    wavelet_levels=5,
                    zscore_threshold=3.0
                )
                _, rpk = nk.ecg_peaks(ecg_signal_clean, sampling_rate=sfreq)
                rpeaks = rpk['ECG_R_Peaks']
            except Exception as e:
                raise e
            
            if len(rpeaks) < 2:
                continue
            
            # Convert R-peak indices to timestamps
            rpeak_times = rpeaks / sfreq
            rpeak_ts = nap.Ts(t=rpeak_times, time_units="s")

            # Find EEG channels
            eeg_ch_indices = [i for i, ch in enumerate(ch_names) 
                            if 'eeg' in ch.lower() or any(elec in ch.upper() for elec in ['FP', 'F', 'C', 'P', 'O', 'T', 'A'])]

            if eeg_ch_indices:
                eeg_data = data[eeg_ch_indices, :]
                eeg_ch_names = [ch_names[i] for i in eeg_ch_indices]
                patient_results = []
                patient_results_list = []
                
                # Process each EEG channel
                for ch_idx, ch_name in enumerate(eeg_ch_names):
                    try:
                        eeg_signal = eeg_data[ch_idx, :]
                        time_index = np.arange(len(eeg_signal)) / sfreq
                        eeg_signal_tsd = nap.Tsd(t=time_index, d=eeg_signal, time_units="s")

                        # 1. OBSERVED HEP (Average across real R-peaks)
                        peth_obs = nap.compute_perievent_continuous(eeg_signal_tsd, rpeak_ts, minmax=minmax)
                        mask = (peth_obs.index.values >= test_window[0]) & (peth_obs.index.values <= test_window[1])
                        
                        # X_obs shape: (Time_points,) -> average across trials
                        X_obs = np.mean(peth_obs.values[mask, :], axis=1)

                        # 2. NULL DISTRIBUTIONS (Jittering)
                        # We generate many "fake" HEPs to see what happens by chance
                        null_heps = []
                        for _ in range(n_bootstrap):
                            # Jitter timestamps using pynapple
                            jittered_ts = nap.process.jitter_timestamps(rpeak_ts, max_jitter=jitter_amount)
                            
                            # Compute HEP with jittered peaks
                            peth_null = nap.compute_perievent_continuous(eeg_signal_tsd, jittered_ts, minmax=minmax)
                            
                            # Average across trials for this permutation and store
                            null_heps.append(np.mean(peth_null.values[mask, :], axis=1))

                        null_matrix = np.array(null_heps) # Shape: (n_permutations, n_time_points)

                        # 3. PREPARE FOR MNE CLUSTER TEST
                        # We subtract the mean of the null from our observed data
                        # Format for 1samp: (n_observations, n_times). 
                        # Here, observations are the differences between observed HEP and each null HEP
                        X_diff = X_obs - null_matrix 

                        # Define threshold based on t-distribution
                        threshold_value = stats.t.ppf(1 - 0.05 / 2, df=n_bootstrap - 1)

                        T_obs, clusters, cluster_pv, H0 = permutation_cluster_1samp_test(
                            X_diff,
                            n_permutations=n_bootstrap,
                            threshold=threshold_value,
                            tail=0,
                            out_type='indices',
                            verbose=False
                        )

                        # 4. STORE RESULTS
                        significant = False
                        p_value = 1.0
                        if len(cluster_pv) > 0:
                            p_value = np.min(cluster_pv)
                            significant = p_value < significance_threshold

                        patient_results.append({
                            'electrode': ch_name,
                            'p_value': p_value,
                            'significant': significant,
                            'mean_amplitude': np.mean(X_obs) # Observed HEP mean in window
                        })
                        # The index of peth_obs contains the time relative to the R-peak (e.g., -0.2 to 0.8s)
                        relative_times = peth_obs.index.values[mask]

                        significant_windows = []

                        for i, clu_idx in enumerate(clusters):
                            if cluster_pv[i] < 0.05:
                                # Extract the start and end relative to the R-peak
                                # clu_idx is a tuple of arrays; for 1D, we take the first element
                                time_start = relative_times[clu_idx[0][0]]
                                time_end = relative_times[clu_idx[0][-1]]
                                                                
                                significant_windows.append({
                                    'start': time_start,
                                    'end': time_end,
                                    'p_value': cluster_pv[i]
                                })
                        patient_results_list.append(
                            {
                                'patient_id': patient_id,
                                'electrode': ch_name,
                                'significant_windows': significant_windows

                            }
                        )
                    
                    except Exception:
                        continue
            
            # Store results for this patient
            if patient_results: 
                results[patient_id] = patient_results
                results_full_time[patient_id] = patient_results_list
        
        except Exception:
            continue
    
    if 'st' in globals():
        progress_bar.empty()
        status_text.empty()
    
    return results,results_full_time


def plot_significance_percentage_topomap(significance_results, save_plot=False, save_dir=None):
    """
    Plot an MNE topomap showing the percentage of patients with significant electrodes.
    
    Parameters
    ----------
    significance_results : dict
        Dictionary returned from test_electrode_significance_per_patient
        Format: {patient_id: [{'electrode': str, 'p_value': float, 'mean_diff': float, 'significant': bool}, ...]}
    save_plot : bool, optional
        Whether to save the plot (default: False)
    save_dir : str, optional
        Directory to save the plot (default: None)
    
    Returns
    -------
    matplotlib.figure.Figure or None
        Figure object if plot was created, None otherwise
    """
    if not significance_results:
        if 'st' in globals():
            st.warning("No significance results available for topomap.")
        else:
            print("No significance results available for topomap.")
        return None
    
    # Count significant and total occurrences for each electrode
    electrode_counts = {}  # {electrode: {'significant': count, 'total': count}}
    
    for patient_id, patient_results in significance_results.items():
        for result in patient_results:
            electrode = result['electrode']
            if electrode not in electrode_counts:
                electrode_counts[electrode] = {'significant': 0, 'total': 0}
            
            electrode_counts[electrode]['total'] += 1
            if result['significant']:
                electrode_counts[electrode]['significant'] += 1
    
    # Calculate percentages
    electrode_percentages = {}
    for electrode, counts in electrode_counts.items():
        if counts['total'] > 0:
            percentage = (counts['significant'] / counts['total']) * 100
            electrode_percentages[electrode] = percentage
    
    if not electrode_percentages:
        if 'st' in globals():
            st.warning("No electrode data available for topomap.")
        else:
            print("No electrode data available for topomap.")
        return None
    
    # Need at least 3 electrodes for a meaningful topomap
    if len(electrode_percentages) < 3:
        if 'st' in globals():
            st.warning(f"Insufficient electrodes ({len(electrode_percentages)}) for topomap. Need at least 3.")
        else:
            print(f"Insufficient electrodes ({len(electrode_percentages)}) for topomap. Need at least 3.")
        return None
    
    try:
        # Create montage (standard 10-20 system)
        montage = mne.channels.make_standard_montage('standard_1020')
        
        # Get available channels and values
        available_channels = []
        values = []
        used_electrodes = set()  # Track which electrodes we've matched
        
        # Normalize montage channel names (uppercase)
        montage_ch_names_upper = {ch.upper(): ch for ch in montage.ch_names}
        
        for electrode, percentage in electrode_percentages.items():
            if electrode in used_electrodes:
                continue
                
            electrode_upper = electrode.upper()
            matched = False
            
            # Strategy 1: Exact match (case-insensitive)
            if electrode_upper in montage_ch_names_upper:
                available_channels.append(montage_ch_names_upper[electrode_upper])
                values.append(percentage)
                used_electrodes.add(electrode)
                matched = True
                continue
            
            # Strategy 2: Try common variations
            # Remove common prefixes/suffixes and try again
            electrode_clean = electrode_upper
            # Remove common suffixes like '_EEG', '-EEG', etc.
            for suffix in ['_EEG', '-EEG', '_E', '-E', ' EEG', ' E']:
                if electrode_clean.endswith(suffix):
                    electrode_clean = electrode_clean[:-len(suffix)]
                    break
            
            if electrode_clean in montage_ch_names_upper:
                available_channels.append(montage_ch_names_upper[electrode_clean])
                values.append(percentage)
                used_electrodes.add(electrode)
                matched = True
                continue
            
            # Strategy 3: Try matching by first 2-3 characters (for regional matching)
            # This handles cases like 'Fp1' matching 'FP1', 'F3' matching 'F3', etc.
            if len(electrode_clean) >= 2:
                prefix = electrode_clean[:2]
                # Try to find montage channels that start with the same prefix
                for montage_ch_upper, montage_ch_orig in montage_ch_names_upper.items():
                    if montage_ch_upper.startswith(prefix) and montage_ch_orig not in available_channels:
                        # Check if the numbers match (e.g., Fp1 should match FP1, not FP2)
                        # Extract numbers from both
                        elec_num = ''.join([c for c in electrode_clean if c.isdigit()])
                        montage_num = ''.join([c for c in montage_ch_upper if c.isdigit()])
                        
                        # If numbers match or both have no numbers, it's a good match
                        if elec_num == montage_num or (not elec_num and not montage_num):
                            available_channels.append(montage_ch_orig)
                            values.append(percentage)
                            used_electrodes.add(electrode)
                            matched = True
                            break
                    
                    if matched:
                        break
        
        if len(available_channels) < 3:
            if 'st' in globals():
                st.warning(f"Could not match enough electrodes to montage. Found {len(available_channels)} matches.")
            else:
                print(f"Could not match enough electrodes to montage. Found {len(available_channels)} matches.")
            return None
        
        # Create MNE info object
        info = mne.create_info(ch_names=available_channels, sfreq=256, ch_types='eeg')
        info.set_montage(montage)
        
        # Create EvokedArray with percentage values
        values_array = np.array(values).reshape(-1, 1)
        evoked = mne.EvokedArray(values_array, info)
        
        # Calculate min and max for colorbar limits
        vmin = np.min(values) if len(values) > 0 else 0
        vmax = np.max(values) if len(values) > 0 else 100
        
        # Plot topomap
        fig, ax = plt.subplots(figsize=(10, 8))
        im, _ = mne.viz.plot_topomap(
            evoked.data[:, 0], 
            evoked.info, 
            axes=ax, 
            show=False,
            cmap='Reds',
            vlim=(vmin, vmax)  # Use actual min and max of percentage values
        )
        fig.colorbar(im, ax=ax, label='% of Patients with Significant Electrode')
        ax.set_title('Percentage of Patients with Significant RR-to-Channel Coupling\n(Per Electrode)', 
                    fontsize=14, fontweight='bold', pad=20)
        plt.tight_layout()
        
        # Save plot if requested
        if save_plot and save_dir:
            os.makedirs(save_dir, exist_ok=True)
            save_path = os.path.join(save_dir, 'electrode_significance_percentage_topomap.png')
            plt.savefig(save_path, dpi=300, bbox_inches='tight')
            if 'st' in globals():
                st.success(f"Topomap saved to {save_path}")
            else:
                print(f"Topomap saved to {save_path}")
        
        # Display in streamlit if available
        display_plot_with_download(fig, "electrode_significance_percentage_topomap", "topomap_significance")
        return fig
        
    except Exception as e:
        if 'st' in globals():
            st.error(f"Error creating topomap: {e}")
        else:
            print(f"Error creating topomap: {e}")
        import traceback
        traceback.print_exc()
        return None


def plot_regional_electrode_averaged_eeg(loaded_files, minmax=(-0.5, 1.0), binsize=0.005, save_plot=False, save_dir=None, aggregation_method='median', processing_mode='single', sleep_stage_mapping=None, sleep_stage='N1'):
    """
    Create averaged PETH and raster plots for each electrode within each region (F, P, T, O, C).
    Each region gets its own figure with subplots for each electrode, averaged across all patients.
    
    Parameters
    ----------
    loaded_files : list
        List of tuples (file_path, patient_id, raw_object)
    minmax : tuple, optional
        Time window around events in seconds (default: (-0.5, 1.0))
    binsize : float, optional
        Bin size for rate calculation in seconds (default: 0.05)
    save_plot : bool, optional
        Whether to save the plot (default: False)
    save_dir : str, optional
        Directory to save the plot (default: None)
    aggregation_method : str, optional
        Method for aggregating across patients: 'median' or 'mean' (default: 'median')
    processing_mode : str, optional
        Processing mode: 'single' for sequential or 'multithread' for parallel (default: 'single')
    
    Returns
    -------
    list : List of matplotlib.figure.Figure objects (one per region)
    """
    if not PYNAPPLE_AVAILABLE:
        if 'st' in globals():
            st.error("pynapple is not available. Please install it to use regional electrode plots.")
        else:
            print("pynapple is not available. Please install it to use regional electrode plots.")
        return []
    
    if not loaded_files:
        return []
    
    # Define electrode groupings by region
    regional_electrodes = {
        'F': ['Fp1', 'Fpz', 'Fp2', 'F7', 'F3', 'Fz', 'F4', 'F8'],
        'P': ['P3', 'Pz', 'P4'],
        'T': ['T3', 'T4', 'T5', 'T6'],
        'O': ['O1', 'Oz', 'O2'],
        'C': ['C3', 'Cz', 'C4']
    }
    
    # Collect data from all patients, grouped by region and electrode
    # Structure: {region: {electrode: {'rates_list': [], 'events_list': []}}}
    all_electrode_data = {}
    for region in regional_electrodes.keys():
        all_electrode_data[region] = {}
        for electrode in regional_electrodes[region]:
            all_electrode_data[region][electrode] = {'rates_list': [], 'events_list': []}
    
    if 'st' in globals():
        progress_bar = st.progress(0)
        status_text = st.empty()
    
    for file_idx, (file_path, patient_id, raw) in enumerate(loaded_files):
        try:
            if 'st' in globals():
                status_text.text(f"Processing {patient_id} for regional electrode plots ({file_idx + 1}/{len(loaded_files)})...")
                progress_bar.progress((file_idx + 1) / len(loaded_files))
            
            # Get sampling frequency
            sfreq = raw.info['sfreq']
            
            # Get channel names and data
            ch_names = raw.ch_names
            data = raw.get_data()
            
            # Find ECG channel
            ch_lower = [ch.lower() for ch in ch_names]
            ecg_indices = [i for i, ch in enumerate(ch_lower) if 'ecg' in ch or 'ekg' in ch]
            
            if not ecg_indices:
                continue
            
            ecg_ch_idx = ecg_indices[0]
            ecg_signal = data[ecg_ch_idx, :]
            
            # Clean ECG and detect R-peaks
            try:
                ecg_signal_clean, _ = clean_ecg_advanced(
                    ecg_signal, 
                    sampling_rate=sfreq,
                    median_window_ms=300,
                    wavelet='db4',
                    wavelet_levels=5,
                    zscore_threshold=3.0
                )
                _, rpk = nk.ecg_peaks(ecg_signal_clean, sampling_rate=sfreq)
                rpeaks = rpk['ECG_R_Peaks']
            except Exception:
                continue
            
            if len(rpeaks) < 2:
                continue
            
            # Convert R-peak indices to timestamps
            rpeak_times = rpeaks / sfreq
            rpeak_ts = nap.Ts(t=rpeak_times, time_units="s")
            
            # Create mapping of channel names to indices
            ch_name_to_idx = {ch.upper(): i for i, ch in enumerate(ch_names)}
            
            # Process each region and electrode
            for region, electrodes in regional_electrodes.items():
                for electrode in electrodes:
                    electrode_upper = electrode.upper()
                    
                    # Find channel index
                    if electrode_upper not in ch_name_to_idx:
                        continue
                    
                    ch_idx = ch_name_to_idx[electrode_upper]
                    
                    try:
                        eeg_signal = data[ch_idx, :]
                        
                        # Convert continuous EEG signal to pynapple Tsd object
                        time_index = np.arange(len(eeg_signal)) / sfreq
                        eeg_signal_tsd = nap.Tsd(t=time_index, d=eeg_signal, time_units="s")
                        
                        # Compute perievent alignment around R-peaks using continuous signal
                        peth_eeg = nap.compute_perievent_continuous(
                            timeseries=eeg_signal_tsd,
                            tref=rpeak_ts,
                            minmax=minmax,
                            time_unit="s"
                        )
                        
                        if len(peth_eeg) > 0:
                                # Calculate rates for this patient and electrode
                                try:
                                    peth_eeg_count = peth_eeg.count(binsize)
                                    if hasattr(peth_eeg_count, 'values'):
                                        count_vals = peth_eeg_count.values
                                        if count_vals.ndim == 2:
                                            rates = np.mean(count_vals, axis=1) / binsize
                                        else:
                                            rates = count_vals / binsize
                                        
                                        if hasattr(peth_eeg_count, 'index'):
                                            time_axis = peth_eeg_count.index.values
                                        else:
                                            time_axis = np.arange(len(rates)) * binsize + minmax[0]
                                        
                                        all_electrode_data[region][electrode]['rates_list'].append((time_axis, rates))
                                except Exception:
                                    pass
                                
                                # Collect events for raster plot
                                try:
                                    peth_eeg_tsd = peth_eeg.to_tsd()
                                    all_electrode_data[region][electrode]['events_list'].append({
                                        'times': peth_eeg_tsd.index.values,
                                        'events': peth_eeg_tsd.values
                                    })
                                except Exception:
                                    pass
                    
                    except Exception:
                        continue
        
        except Exception as e:
            if 'st' in globals():
                st.warning(f"Error processing {patient_id} for regional electrode plots: {e}")
            else:
                print(f"Error processing {patient_id} for regional electrode plots: {e}")
            continue
    
    if 'st' in globals():
        progress_bar.empty()
        status_text.empty()
    
    # Create plots for each region
    figures = []
    
    for region, electrodes in regional_electrodes.items():
        # Filter electrodes that have data
        electrodes_with_data = [e for e in electrodes if all_electrode_data[region][e]['rates_list'] or all_electrode_data[region][e]['events_list']]
        
        if not electrodes_with_data:
            continue
        
        # Create figure for this region
        n_electrodes = len(electrodes_with_data)
        n_cols = 2  # PETH and Raster side by side
        n_rows = n_electrodes
        
        fig_region, axes = plt.subplots(n_rows, n_cols, figsize=(14, 3 * n_rows), sharex='col')
        if n_electrodes == 1:
            axes = axes.reshape(1, -1)
        
        for elec_idx, electrode in enumerate(electrodes_with_data):
            ax_peth = axes[elec_idx, 0]
            ax_raster = axes[elec_idx, 1]
            
            electrode_data = all_electrode_data[region][electrode]
            
            # Plot 1: Average PETH
            if electrode_data['rates_list']:
                try:
                    # Interpolate all rates to common time axis
                    all_time_axes = [t for t, _ in electrode_data['rates_list']]
                    all_rates = [r for _, r in electrode_data['rates_list']]
                    
                    # Find common time range
                    min_time = max([t[0] for t in all_time_axes])
                    max_time = min([t[-1] for t in all_time_axes])
                    common_time = np.arange(min_time, max_time + binsize, binsize)
                    
                    # Interpolate each patient's rates to common time axis
                    interpolated_rates = []
                    for time_axis, rates in zip(all_time_axes, all_rates):
                        if len(time_axis) > 1 and len(rates) > 1:
                            interp_rates = np.interp(common_time, time_axis, rates)
                            interpolated_rates.append(interp_rates)
                    
                    if interpolated_rates:
                        # Aggregate across patients using selected method
                        if aggregation_method == 'median':
                            agg_rates = np.median(interpolated_rates, axis=0)
                            # Calculate 95% confidence interval using bootstrap percentile method
                            interpolated_rates_array = np.array(interpolated_rates)
                            spread_lower, spread_upper = bootstrap_median_ci(interpolated_rates_array, n_bootstrap=1000)
                            spread_label = '95% CI'
                            label_text = 'Median'
                        else:  # mean
                            agg_rates = np.mean(interpolated_rates, axis=0)
                            std_rates = np.std(interpolated_rates, axis=0)
                            spread_lower = agg_rates - std_rates
                            spread_upper = agg_rates + std_rates
                            label_text = 'Mean'
                            spread_label = '±1 SD'
                        
                        # Perform bootstrapping test
                        p_value, mean_diff = bootstrap_pvalue_test(
                            electrode_data['rates_list'],
                            test_window=(-0.1, 0.1),  # Window around R-peak
                            baseline_window=(-0.5, -0.2),  # Baseline window
                            n_bootstrap=10000
                        )
                        
                        # Format p-value for display
                        if not np.isnan(p_value):
                            if p_value < 0.001:
                                p_str = "p < 0.001"
                            elif p_value < 0.01:
                                p_str = f"p = {p_value:.3f}"
                            else:
                                p_str = f"p = {p_value:.3f}"
                        else:
                            p_str = "p = N/A"
                        
                        ax_peth.plot(common_time, agg_rates, linewidth=2, color='blue', label=label_text)
                        ax_peth.fill_between(common_time, spread_lower, spread_upper, 
                                            alpha=0.3, color='blue', label=spread_label)
                        ax_peth.set_ylabel("Rate (events/sec)", fontsize=10)
                        ax_peth.set_title(f"{label_text} PETH - {electrode} (n={len(interpolated_rates)} patients, {p_str})", 
                                         fontsize=11, fontweight='bold')
                        
                        # Perform max-statistic permutation test for bin-wise significance
                        try:
                            interpolated_rates_array = np.array(interpolated_rates)
                            sig_test_results = max_statistic_permutation_test(
                                interpolated_rates_array, common_time, 
                                baseline_window=(-0.5, -0.2),
                                n_permutations=1000,
                                alpha_levels=[0.05, 0.01],
                                aggregation_method=aggregation_method
                            )
                            
                            if sig_test_results['p_values'] is not None:
                                # Calculate y-axis range for marker placement
                                y_min, y_max = ax_peth.get_ylim()
                                y_range = y_max - y_min
                                marker_offset = y_range * 0.05  # 5% above the line
                                
                                # Plot markers for significant bins
                                sig_05 = sig_test_results['significant_bins'].get(0.05, np.zeros(len(common_time), dtype=bool))
                                sig_01 = sig_test_results['significant_bins'].get(0.01, np.zeros(len(common_time), dtype=bool))
                                
                                # Markers for p < 0.01 (higher significance)
                                if np.any(sig_01):
                                    sig_01_times = common_time[sig_01]
                                    sig_01_rates = agg_rates[sig_01] + marker_offset
                                    ax_peth.scatter(sig_01_times, sig_01_rates, marker='*', s=100, 
                                                   color='red', zorder=5, alpha=0.8, 
                                                   label='p < 0.01' if not any([np.any(sig_01), np.any(sig_05)]) else '')
                                
                                # Markers for p < 0.05 (but not < 0.01)
                                sig_05_only = sig_05 & ~sig_01
                                if np.any(sig_05_only):
                                    sig_05_times = common_time[sig_05_only]
                                    sig_05_rates = agg_rates[sig_05_only] + marker_offset * 0.7
                                    ax_peth.scatter(sig_05_times, sig_05_rates, marker='*', s=60, 
                                                   color='orange', zorder=5, alpha=0.7,
                                                   label='p < 0.05' if not np.any(sig_01) else '')
                        except Exception as e:
                            # Silently fail if test cannot be performed
                            if 'st' in globals():
                                pass
                            else:
                                print(f"Warning: Could not perform significance test for {electrode}: {e}")
                        
                        ax_peth.axvline(0.0, color='red', linestyle='--', linewidth=1, alpha=0.5)
                        ax_peth.legend(fontsize=8)
                        ax_peth.grid(True, alpha=0.3)
                        ax_peth.set_xlim(minmax)
                        ax_peth.tick_params(axis='x', labelsize=8, bottom=True)
                except Exception as e:
                    ax_peth.text(0.5, 0.5, f'PETH Error: {str(e)[:40]}', 
                               ha='center', va='center', transform=ax_peth.transAxes, fontsize=8, color='red')
            else:
                ax_peth.text(0.5, 0.5, 'No rate data', 
                           ha='center', va='center', transform=ax_peth.transAxes, fontsize=9)
            
            # Plot 2: Average/Median Event Density
            if electrode_data['events_list']:
                try:
                    # Calculate histogram for each patient separately
                    time_bins = np.arange(minmax[0], minmax[1] + binsize, binsize)
                    patient_histograms = []
                    
                    for event_data in electrode_data['events_list']:
                        times = event_data['times']
                        if len(times) > 0:
                            hist_counts, _ = np.histogram(times, bins=time_bins)
                            patient_histograms.append(hist_counts)
                    
                    if patient_histograms and len(patient_histograms) > 0:
                        # Stack histograms and aggregate using selected method
                        histograms_array = np.array(patient_histograms)
                        
                        if aggregation_method == 'median':
                            agg_counts = np.median(histograms_array, axis=0)
                            label_text = 'Median'
                        else:  # mean
                            agg_counts = np.mean(histograms_array, axis=0)
                            label_text = 'Mean'
                        
                        # Create time axis for histogram (bin centers)
                        bin_centers = (time_bins[:-1] + time_bins[1:]) / 2
                        
                        # Plot as bar plot showing aggregated event density
                        ax_raster.bar(bin_centers, agg_counts, width=binsize*0.8, 
                                     color='#0C7BDC', alpha=0.6, edgecolor='#0A5FA8', linewidth=0.5)
                        ax_raster.set_ylabel(f"{label_text} Events/Bin", fontsize=10)
                        n_patients = len(patient_histograms)
                        ax_raster.set_title(f"{label_text} Event Density - {electrode} (n={n_patients} patients)", 
                                           fontsize=11, fontweight='bold')
                        ax_raster.axvline(0.0, color='red', linestyle='--', linewidth=1, alpha=0.5)
                        ax_raster.set_xlim(minmax)
                        ax_raster.grid(True, alpha=0.3, axis='y')
                        ax_raster.tick_params(axis='x', labelsize=8, bottom=True)
                        ax_raster.tick_params(axis='y', labelsize=8)
                except Exception as e:
                    ax_raster.text(0.5, 0.5, f'Raster Error: {str(e)[:40]}', 
                                 ha='center', va='center', transform=ax_raster.transAxes, fontsize=8, color='red')
            else:
                ax_raster.text(0.5, 0.5, 'No event data', 
                             ha='center', va='center', transform=ax_raster.transAxes, fontsize=9)
            
            # Set x-label and ticks for all subplots
            ax_peth.set_xlabel("Time from R-peak (s)", fontsize=10)
            ax_raster.set_xlabel("Time from R-peak (s)", fontsize=10)
            
            # Ensure x-axis ticks are visible
            ax_peth.tick_params(labelsize=8, axis='both', which='major')
            ax_raster.tick_params(labelsize=8, axis='both', which='major')
            ax_peth.tick_params(axis='x', labelsize=8, bottom=True)
            ax_raster.tick_params(axis='x', labelsize=8, bottom=True)
        
        # Determine aggregation label for title
        agg_label = "Median" if aggregation_method == 'median' else "Average"
        plt.suptitle(f"{region} Region - {agg_label} PETH & Raster Plots (All Patients)", 
                    fontsize=14, fontweight='bold', y=0.995)
        plt.tight_layout(rect=[0, 0, 1, 0.99])
        
        # Save plot if requested
        if save_plot and save_dir:
            save_path_region = os.path.join(save_dir, f'eeg_peth_raster_{region}_region_averaged.png')
            plt.savefig(save_path_region, dpi=300, bbox_inches='tight')
            if 'st' in globals():
                st.success(f"{region} region averaged plots saved to {save_path_region}")
            else:
                print(f"{region} region averaged plots saved to {save_path_region}")
        
        # Display in streamlit if available
        display_plot_with_download(fig_region, f"eeg_peth_raster_{region}_region_averaged", f"regional_electrode_{region}")
        figures.append(fig_region)
        plt.close(fig_region)
    
    return figures


def plot_significant_patients_electrode_peth(loaded_files, significance_results, minmax=(-0.5, 1.0), binsize=0.05, save_plot=False, save_dir=None, aggregation_method='median', processing_mode='single', sleep_stage_mapping=None, sleep_stage='N1'):
    """
    Create averaged PETH plots for each electrode, showing only patients where that electrode is significant.
    
    Parameters
    ----------
    loaded_files : list
        List of tuples (file_path, patient_id, raw_object)
    significance_results : dict
        Dictionary returned from test_electrode_significance_per_patient
        Format: {patient_id: [{'electrode': str, 'p_value': float, 'mean_diff': float, 'significant': bool}, ...]}
    minmax : tuple, optional
        Time window around events in seconds (default: (-0.5, 1.0))
    binsize : float, optional
        Bin size for rate calculation in seconds (default: 0.05)
    save_plot : bool, optional
        Whether to save the plot (default: False)
    save_dir : str, optional
        Directory to save the plot (default: None)
    aggregation_method : str, optional
        Method for aggregating across patients: 'median' or 'mean' (default: 'median')
    processing_mode : str, optional
        Processing mode: 'single' for sequential or 'multithread' for parallel (default: 'single')
    
    Returns
    -------
    list : List of matplotlib.figure.Figure objects
    """
    if not PYNAPPLE_AVAILABLE:
        if 'st' in globals():
            st.error("pynapple is not available. Please install it to use significant patients PETH plots.")
        else:
            print("pynapple is not available. Please install it to use significant patients PETH plots.")
        return []
    
    if not loaded_files or not significance_results:
        return []
    
    # Build mapping: {electrode: [list of patient_ids that are significant]}
    electrode_significant_patients = {}  # {electrode: [patient_id, ...]}
    
    for patient_id, patient_results in significance_results.items():
        for result in patient_results:
            if result.get('significant', False):
                electrode = result['electrode']
                if electrode not in electrode_significant_patients:
                    electrode_significant_patients[electrode] = []
                electrode_significant_patients[electrode].append(patient_id)
    
    if not electrode_significant_patients:
        if 'st' in globals():
            st.warning("No significant electrodes found for any patient.")
        else:
            print("No significant electrodes found for any patient.")
        return []
    
    # Create mapping from patient_id to file index
    patient_to_file_idx = {}
    for idx, (file_path, patient_id, raw) in enumerate(loaded_files):
        patient_to_file_idx[patient_id] = idx
    
    # Collect data from significant patients only, grouped by electrode
    # Structure: {electrode: {'rates_list': [], 'events_list': []}}
    all_electrode_data = {}
    for electrode in electrode_significant_patients.keys():
        all_electrode_data[electrode] = []
    
    if 'st' in globals():
        progress_bar = st.progress(0)
        status_text = st.empty()
    
    total_patients = sum(len(patients) for patients in electrode_significant_patients.values())
    processed = 0
    
    # Process each electrode's significant patients
    for electrode, significant_patient_ids in electrode_significant_patients.items():
        for patient_id in significant_patient_ids:
            try:
                if patient_id not in patient_to_file_idx:
                    continue
                
                file_idx = patient_to_file_idx[patient_id]
                file_path, patient_id_check, raw = loaded_files[file_idx]
                
                if 'st' in globals():
                    processed += 1
                    status_text.text(f"Processing {electrode} - {patient_id} ({processed}/{total_patients})...")
                    progress_bar.progress(processed / total_patients)
                
                # Get sampling frequency
                sfreq = raw.info['sfreq']
                
                # Get channel names and data
                ch_names = raw.ch_names
                data = raw.get_data()
                
                # Find ECG channel
                ch_lower = [ch.lower() for ch in ch_names]
                ecg_indices = [i for i, ch in enumerate(ch_lower) if 'ecg' in ch or 'ekg' in ch]
                
                if not ecg_indices:
                    continue
                
                ecg_ch_idx = ecg_indices[0]
                ecg_signal = data[ecg_ch_idx, :]
                
                # Clean ECG and detect R-peaks
                try:
                    ecg_signal_clean, _ = clean_ecg_advanced(
                        ecg_signal, 
                        sampling_rate=sfreq,
                        median_window_ms=300,
                        wavelet='db4',
                        wavelet_levels=5,
                        zscore_threshold=3.0
                    )
                    _, rpk = nk.ecg_peaks(ecg_signal_clean, sampling_rate=sfreq)
                    rpeaks = rpk['ECG_R_Peaks']
                except Exception:
                    continue
                
                if len(rpeaks) < 2:
                    continue
                
                # Convert R-peak indices to timestamps
                rpeak_times = rpeaks / sfreq
                rpeak_ts = nap.Ts(t=rpeak_times, time_units="s")
                
                # Find electrode channel
                ch_name_to_idx = {ch.upper(): i for i, ch in enumerate(ch_names)}
                electrode_upper = electrode.upper()
                
                if electrode_upper not in ch_name_to_idx:
                    continue
                
                ch_idx = ch_name_to_idx[electrode_upper]
                
                try:
                    eeg_signal = data[ch_idx, :]
                    
                    # Convert continuous EEG signal to pynapple Tsd object
                    time_index = np.arange(len(eeg_signal)) / sfreq
                    eeg_signal_tsd = nap.Tsd(t=time_index, d=eeg_signal, time_units="s")
                    
                    # Compute perievent alignment around R-peaks using continuous signal
                    peth_eeg = nap.compute_perievent_continuous(
                        timeseries=eeg_signal_tsd,
                        tref=rpeak_ts,
                        minmax=minmax,
                        time_unit="s"
                    )
                    
                    if len(peth_eeg) > 0:
                        # Calculate continuous mean signal for this patient and electrode
                        # peth_eeg is a TsdFrame (Time x Trials)
                        # We want to average across trials (axis 1)
                        if hasattr(peth_eeg, 'values'):
                            values = peth_eeg.values
                            if values.ndim == 2:
                                    mean_signal = np.nanmean(values, axis=1)
                            else:
                                mean_signal = values
                            
                            if hasattr(peth_eeg, 'index'):
                                time_axis = peth_eeg.index.values
                            else:
                                time_axis = np.linspace(minmax[0], minmax[1], len(mean_signal))
                            
                            all_electrode_data[electrode].append((time_axis, mean_signal))                    
                        # Note: Raster plot (events_list) is not populated for continuous data 
                        # as it requires different visualization than event histograms.
                
                except Exception:
                    continue
            
            except Exception:
                continue
    
    if 'st' in globals():
        progress_bar.empty()
        status_text.empty()
    
    # Filter electrodes that have data
    electrodes_with_data = [e for e in all_electrode_data.keys() if all_electrode_data[e]]
    
    # Create plots - group by region for better organization
    regional_electrodes = {
        'F': ['Fp1', 'Fpz', 'Fp2', 'F7', 'F3', 'Fz', 'F4', 'F8'],
        'P': ['P3', 'Pz', 'P4'],
        'T': ['T3', 'T4', 'T5', 'T6'],
        'O': ['O1', 'Oz', 'O2'],
        'C': ['C3', 'Cz', 'C4']
    }
    
    # Group electrodes by region
    electrodes_by_region = {}
    for region, region_elecs in regional_electrodes.items():
        electrodes_by_region[region] = [e for e in electrodes_with_data if e.upper() in [re.upper() for re in region_elecs]]
    
    figures = []
    
    # Create plots for each region that has significant electrodes
    for region, region_elecs in electrodes_by_region.items():
        if not region_elecs:
            continue
        
        # Create figure for this region
        n_electrodes = len(region_elecs)
        
        # Grid settings - using single column as requested
        n_cols = 1
        n_rows = n_electrodes
            
        fig_region, axes = plt.subplots(n_rows, n_cols, figsize=(8, 4 * n_rows), constrained_layout=True)
        # Flatten axes for easy iteration
        if n_electrodes > 1:
            axes_flat = axes.flatten()
        else:
            axes_flat = [axes]
            
        # Hide unused subplots
        for i in range(n_electrodes, len(axes_flat)):
            axes_flat[i].set_visible(False)
            
        for elec_idx, electrode in enumerate(region_elecs):
            ax = axes_flat[elec_idx]
            
            # Get data for this electrode
            electrode_data_list = all_electrode_data.get(electrode, [])
            _ = len(electrode_significant_patients.get(electrode, []))
            
            # Calculate stats for this electrode from significance_results
            electrode_slopes = []
            total_patients_with_electrode = 0
            
            for pid, res_list in significance_results.items():
                for res in res_list:
                    if res['electrode'] == electrode:
                        total_patients_with_electrode += 1
                        if res['significant']:
                            electrode_slopes.append(res['mean_diff'])
            
            sig_percentage = 0
            if total_patients_with_electrode > 0:
                sig_percentage = (len(electrode_slopes) / total_patients_with_electrode) * 100
                
            group_p_value = np.nan
            # Filter out NaNs and Infs
            electrode_slopes = [s for s in electrode_slopes if np.isfinite(s)]
            
            if len(electrode_slopes) >= 2:
                # One-sample t-test of slopes against 0
                try:
                    _, group_p_value = stats.ttest_1samp(electrode_slopes, 0)
                except Exception:
                    group_p_value = np.nan
            
            p_str = "p=N/A"
            if not np.isnan(group_p_value):
                if group_p_value < 0.05:
                    p_str = "p<0.001" if group_p_value < 0.001 else f"p={group_p_value:.3f}"
        
            if electrode_data_list:
                try:
                    # electrode_data_list contains tuples of (time_axis, mean_signal)
                    # Interpolate all rates to common time axis
                    all_time_axes = [t for t, _ in electrode_data_list]
                    all_signals = [s for _, s in electrode_data_list]
                    
                    # Find common time range
                    min_time = max([t[0] for t in all_time_axes])
                    max_time = min([t[-1] for t in all_time_axes])
                    common_time = np.arange(min_time, max_time + binsize, binsize)
                    
                    # Interpolate each patient's signal to common time axis
                    interpolated_signals = []
                    for time_axis, signal in zip(all_time_axes, all_signals):
                        if len(time_axis) > 1 and len(signal) > 1:
                            interp_sig = np.interp(common_time, time_axis, signal)
                            interpolated_signals.append(interp_sig)
                    
                    if interpolated_signals:
                        # Aggregate across significant patients using selected method
                        n_patients = len(interpolated_signals)
                        interpolated_signals_array = np.array(interpolated_signals)
                        
                        if aggregation_method == 'median':
                            agg_signal = np.median(interpolated_signals_array, axis=0)
                            # Calculate 95% confidence interval using bootstrap percentile method
                            spread_lower, spread_upper = bootstrap_median_ci(interpolated_signals_array, n_bootstrap=1000)
                            label_text = 'Median'
                            spread_label = '95% CI'
                        else:  # mean
                            agg_signal = np.mean(interpolated_signals_array, axis=0)
                            std_signal = np.std(interpolated_signals_array, axis=0)
                            # Calculate 95% confidence interval
                            if n_patients > 1:
                                # Standard error
                                se_signal = std_signal / np.sqrt(n_patients)
                                # t-critical value for 95% CI (two-tailed)
                                t_critical = stats.t.ppf(0.975, df=n_patients - 1)
                                # Confidence interval bounds
                                spread_lower = agg_signal - t_critical * se_signal
                                spread_upper = agg_signal + t_critical * se_signal
                                spread_label = '95% CI'
                            else:
                                # If only one patient, use SD as fallback
                                spread_lower = agg_signal - std_signal
                                spread_upper = agg_signal + std_signal
                                spread_label = '±1 SD'
                            label_text = 'Mean'
                        
                        # Plot aggregated signal
                        ax.plot(common_time, agg_signal, linewidth=2, color='blue', label=label_text)
                        
                        # Update title with stats
                        if p_str != "p=N/A":
                            ax.set_title(f"{electrode} (n={n_patients}, {sig_percentage:.0f}% sig, {p_str})", fontsize=11, fontweight='bold')   
                        else:
                            ax.set_title(f"{electrode} (n={n_patients}, {sig_percentage:.0f}% sig)", fontsize=11, fontweight='bold')    
                        ax.fill_between(common_time, spread_lower, spread_upper, 
                                            alpha=0.3, color='blue', label=spread_label)
                        
                        ax.set_ylabel("Amplitude (V)", fontsize=10)
                        ax.set_xlabel("Time from R-peak (s)", fontsize=10)
                        
                        ax.axvline(0.0, color='red', linestyle='--', linewidth=1, alpha=0.5)
                        ax.grid(True, alpha=0.3)
                        ax.set_xlim(minmax)
                        ax.legend(fontsize=8, loc='upper right')
                        
                    else:
                        ax.text(0.5, 0.5, 'No valid interpolated data', 
                               ha='center', va='center', transform=ax.transAxes, fontsize=9)
                except Exception as e:
                    ax.text(0.5, 0.5, f'Error: {str(e)[:40]}', 
                               ha='center', va='center', transform=ax.transAxes, fontsize=8, color='red')
            else:
                ax.text(0.5, 0.5, 'No data', 
                           ha='center', va='center', transform=ax.transAxes, fontsize=9)
        
        # Determine aggregation label for title
        agg_label = "Median" if aggregation_method == 'median' else "Average"
        fig_region.suptitle(f"{region} Region - {agg_label} HEP Waveforms (Significant Electrodes)", 
                    fontsize=14, fontweight='bold', y=1.02)
        
        # Save plot if requested
        if save_plot and save_dir:
            save_path_region = os.path.join(save_dir, f'hep_waveforms_{region}_region_significant_electrodes.png')
            fig_region.savefig(save_path_region, dpi=300, bbox_inches='tight')
            if 'st' in globals():
                st.success(f"{region} region plots saved to {save_path_region}")
            else:
                print(f"{region} region plots saved to {save_path_region}")
        
        # Display in streamlit if available
        display_plot_with_download(fig_region, f"hep_waveforms_{region}_region", f"hep_waveforms_{region}")
        figures.append(fig_region)
        plt.close(fig_region)


def plot_peth_raster_per_patient(loaded_files, significance_results, minmax=(-0.5, 1.0), binsize=0.05, save_plot=False, save_dir=None, aggregation_method='median', processing_mode='single', sleep_stage_mapping=None, sleep_stage='N1'):
    """
    Create PETH and raster plots with one line per patient for each electrode.
    Shows individual patient traces instead of averaging.
    
    Parameters
    ----------
    loaded_files : list
        List of tuples (file_path, patient_id, raw_object)
    significance_results : dict
        Dictionary returned from test_electrode_significance_per_patient
        Format: {patient_id: [{'electrode': str, 'p_value': float, 'mean_diff': float, 'significant': bool}, ...]}
    minmax : tuple, optional
        Time window around events in seconds (default: (-0.5, 1.0))
    binsize : float, optional
        Bin size for rate calculation in seconds (default: 0.05)
    save_plot : bool, optional
        Whether to save the plot (default: False)
    save_dir : str, optional
        Directory to save the plot (default: None)
    aggregation_method : str, optional
        Method for aggregating across patients: 'median' or 'mean' (default: 'median')
        Note: This function shows individual traces, so this parameter is not used but kept for API consistency.
    processing_mode : str, optional
        Processing mode: 'single' for sequential or 'multithread' for parallel (default: 'single')
    
    Returns
    -------
    list : List of matplotlib.figure.Figure objects
    """
    if not PYNAPPLE_AVAILABLE:
        if 'st' in globals():
            st.error("pynapple is not available. Please install it to use per-patient PETH plots.")
        else:
            print("pynapple is not available. Please install it to use per-patient PETH plots.")
        return []
    
    if not loaded_files or not significance_results:
        return []
    
    # Build mapping: {electrode: [(patient_id, time_axis, rates, events), ...]}
    electrode_patient_data = {}  # {electrode: [(patient_id, time_axis, rates, events_dict), ...]}
    
    # Create mapping from patient_id to file index
    patient_to_file_idx = {}
    for idx, (file_path, patient_id, raw) in enumerate(loaded_files):
        patient_to_file_idx[patient_id] = idx
    
    # Collect data from significant patients, grouped by electrode
    for patient_id, patient_results in significance_results.items():
        for result in patient_results:
            if result.get('significant', False):
                electrode = result['electrode']
                if electrode not in electrode_patient_data:
                    electrode_patient_data[electrode] = []
    
    if 'st' in globals():
        progress_bar = st.progress(0)
        status_text = st.empty()
    
    total_combinations = sum(len([r for r in results if r.get('significant', False)]) 
                            for results in significance_results.values())
    processed = 0
    
    # Process each significant patient-electrode combination
    for patient_id, patient_results in significance_results.items():
        if patient_id not in patient_to_file_idx:
            continue
        
        file_idx = patient_to_file_idx[patient_id]
        file_path, patient_id_check, raw = loaded_files[file_idx]
        
        for result in patient_results:
            if not result.get('significant', False):
                continue
            
            electrode = result['electrode']
            
            try:
                if 'st' in globals():
                    processed += 1
                    status_text.text(f"Processing {electrode} - {patient_id} ({processed}/{total_combinations})...")
                    progress_bar.progress(processed / total_combinations)
                
                # Get sampling frequency
                sfreq = raw.info['sfreq']
                
                # Get channel names and data
                ch_names = raw.ch_names
                data = raw.get_data()
                
                # Find ECG channel
                ch_lower = [ch.lower() for ch in ch_names]
                ecg_indices = [i for i, ch in enumerate(ch_lower) if 'ecg' in ch or 'ekg' in ch]
                
                if not ecg_indices:
                    continue
                
                ecg_ch_idx = ecg_indices[0]
                ecg_signal = data[ecg_ch_idx, :]
                
                # Clean ECG and detect R-peaks
                try:
                    ecg_signal_clean, _ = clean_ecg_advanced(
                        ecg_signal, 
                        sampling_rate=sfreq,
                        median_window_ms=300,
                        wavelet='db4',
                        wavelet_levels=5,
                        zscore_threshold=3.0
                    )
                    _, rpk = nk.ecg_peaks(ecg_signal_clean, sampling_rate=sfreq)
                    rpeaks = rpk['ECG_R_Peaks']
                except Exception:
                    continue
                
                if len(rpeaks) < 2:
                    continue
                
                # Convert R-peak indices to timestamps
                rpeak_times = rpeaks / sfreq
                rpeak_ts = nap.Ts(t=rpeak_times, time_units="s")
                
                # Find electrode channel
                ch_name_to_idx = {ch.upper(): i for i, ch in enumerate(ch_names)}
                electrode_upper = electrode.upper()
                
                if electrode_upper not in ch_name_to_idx:
                    continue
                
                ch_idx = ch_name_to_idx[electrode_upper]
                
                try:
                    eeg_signal = data[ch_idx, :]
                    
                    # Convert continuous EEG signal to pynapple Tsd object
                    time_index = np.arange(len(eeg_signal)) / sfreq
                    eeg_signal_tsd = nap.Tsd(t=time_index, d=eeg_signal, time_units="s")
                    
                    # Compute perievent alignment around R-peaks using continuous signal
                    peth_eeg = nap.compute_perievent_continuous(
                        timeseries=eeg_signal_tsd,
                        tref=rpeak_ts,
                        minmax=minmax,
                        time_unit="s"
                    )
                    
                    if len(peth_eeg) > 0:
                        # Calculate continuous mean signal for this patient (average across trials)
                        if hasattr(peth_eeg, 'values'):
                            values = peth_eeg.values
                            if values.ndim == 2:
                                mean_signal = np.nanmean(values, axis=1)
                            else:
                                mean_signal = values
                            
                            if hasattr(peth_eeg, 'index'):
                                time_axis = peth_eeg.index.values
                            else:
                                time_axis = np.linspace(minmax[0], minmax[1], len(mean_signal))
                            
                            # Store with None for events_dict as this is continuous data
                            electrode_patient_data[electrode].append((patient_id, time_axis, mean_signal, None))
                
                except Exception:
                    continue
            
            except Exception:
                continue
    
    if 'st' in globals():
        progress_bar.empty()
        status_text.empty()
    
    # Filter electrodes that have data
    electrodes_with_data = [e for e in electrode_patient_data.keys() 
                           if len(electrode_patient_data[e]) > 0]
    
    # Create plots - group by region for better organization
    regional_electrodes = {
        'F': ['Fp1', 'Fpz', 'Fp2', 'F7', 'F3', 'Fz', 'F4', 'F8'],
        'P': ['P3', 'Pz', 'P4'],
        'T': ['T3', 'T4', 'T5', 'T6'],
        'O': ['O1', 'Oz', 'O2'],
        'C': ['C3', 'Cz', 'C4']
    }
    
    # Group electrodes by region
    electrodes_by_region = {}
    for region, region_elecs in regional_electrodes.items():
        electrodes_by_region[region] = [e for e in electrodes_with_data if e.upper() in [re.upper() for re in region_elecs]]
    
    figures = []
    
    # Create colormap for patients
    import matplotlib.cm as cm
    unique_patients = list(significance_results.keys())
    colors = cm.rainbow(np.linspace(0, 1, len(unique_patients)))
    patient_colors = {pat: color for pat, color in zip(unique_patients, colors)}
    
    # Create plots for each region that has significant electrodes
    for region, region_elecs in electrodes_by_region.items():
        if not region_elecs:
            continue
        
        # Create figure for this region
        n_electrodes = len(region_elecs)
        n_cols = 1  # Only HEP
        n_rows = n_electrodes
        
        fig_region, axes = plt.subplots(n_rows, n_cols, figsize=(10, 4 * n_rows), sharex=True)
        if n_electrodes == 1:
            axes = np.array([axes])
        
        for elec_idx, electrode in enumerate(region_elecs):
            ax_peth = axes[elec_idx]
            
            patient_data_list = electrode_patient_data[electrode]
            n_patients = len(patient_data_list)
            
            # Plot 1: Continuous PETH with one line per patient
            if patient_data_list:
                try:
                    # Find common time range for interpolation
                    all_time_axes = [time_axis for _, time_axis, _, _ in patient_data_list]
                    min_time = max([t[0] for t in all_time_axes])
                    max_time = min([t[-1] for t in all_time_axes])
                    common_time = np.arange(min_time, max_time + binsize, binsize)
                    
                    # Plot each patient's PETH
                    for patient_id, time_axis, mean_signal, _ in patient_data_list:
                        if len(time_axis) > 1 and len(mean_signal) > 1:
                            # Interpolate to common time axis
                            interp_signal = np.interp(common_time, time_axis, mean_signal)
                            color = patient_colors.get(patient_id, 'gray')
                            ax_peth.plot(common_time, interp_signal, linewidth=1.5, 
                                       alpha=0.7, color=color, label=patient_id)
                    
                    ax_peth.set_ylabel("Amplitude (V)", fontsize=10)
                    ax_peth.set_title(f"HEP per Patient - {electrode} (n={n_patients})", 
                                     fontsize=11, fontweight='bold')
                    ax_peth.axvline(0.0, color='red', linestyle='--', linewidth=1, alpha=0.5)
                    
                    # Legend only if few patients to avoid crowding
                    if n_patients <= 10:
                        ax_peth.legend(fontsize=7, ncol=2, loc='upper right')
                        
                    ax_peth.grid(True, alpha=0.3)
                    ax_peth.set_xlim(minmax)
                    ax_peth.tick_params(axis='x', labelsize=8, bottom=True)
                    ax_peth.tick_params(axis='y', labelsize=8)
                except Exception as e:
                    ax_peth.text(0.5, 0.5, f'Plot Error: {str(e)[:40]}', 
                               ha='center', va='center', transform=ax_peth.transAxes, fontsize=8, color='red')
            else:
                ax_peth.text(0.5, 0.5, 'No data', 
                           ha='center', va='center', transform=ax_peth.transAxes, fontsize=9) 
        
        plt.suptitle(f"{region} Region - Individual Patient HEPs", 
                    fontsize=14, fontweight='bold', y=0.995)
        plt.tight_layout(rect=[0, 0, 1, 0.99])
        
        # Save plot if requested
        if save_plot and save_dir:
            save_path_region = os.path.join(save_dir, f'hep_per_patient_{region}_region.png')
            plt.savefig(save_path_region, dpi=300, bbox_inches='tight')
            if 'st' in globals():
                st.success(f"{region} region per-patient plots saved to {save_path_region}")
            else:
                print(f"{region} region per-patient plots saved to {save_path_region}")
        
        # Display in streamlit if available
        display_plot_with_download(fig_region, f"hep_per_patient_{region}_region", f"regional_per_patient_{region}")
        figures.append(fig_region)
        plt.close(fig_region)

    return figures


def calc_rayleigh(angles):
    """
    Calculate Rayleigh test for uniformity.
    angles: array of angles in radians
    Returns: p-value
    """
    n = len(angles)
    if n == 0:
        return 1.0
    
    # Mean resultant length
    r_vector = np.sum(np.exp(1j * angles))
    r = np.abs(r_vector) / n
    
    # Rayleigh statistic Z
    z = n * (r ** 2)
    
    # P-value approximation (valid for large n, typically n > 10)
    # p = exp(-z) * [1 + (2z - z^2) / (4n) - ...]
    # Simple approximation:
    p_value = np.exp(-z)
    
    return p_value, r

def test_circular_statistics_r_to_r(loaded_files):
    """
    Compute circular statistics on Heartbeat Evoked Potentials.
    Iterates through heartbeats (R-to-R intervals), finds the phase of the
    maximum EEG amplitude, and performs circular statistics.
    
    Returns
    -------
    sig_results_rayleigh : dict
        Results of Rayleigh test for uniformity (Phase Locking).
    sig_results_ww : dict
        Results of Watson-Williams test (Mean phase difference).
    sig_results_kuiper : dict
        Results of Kuiper's test (Distribution shape diff).
    """
    if not PYNAPPLE_AVAILABLE:
        print("pynapple not available. Skipping circular stats.")
        return {}, {}, {}

    sig_results_rayleigh = {}
    _ = {}
    _ = {}
    
    # Helper to clean and get R-peaks (reused logic)
    def get_rpeaks(raw):
        sfreq = raw.info['sfreq']
        ch_names = raw.ch_names
        data = raw.get_data()
        
        # Find ECG
        ch_lower = [ch.lower() for ch in ch_names]
        ecg_indices = [i for i, ch in enumerate(ch_lower) if 'ecg' in ch or 'ekg' in ch]
        if not ecg_indices:
            return None, sfreq, data, ch_names
        
        ecg_ch_idx = ecg_indices[0]
        ecg_signal = data[ecg_ch_idx, :]
        
        try:
            # Basic cleaning if advanced not available or just use raw if clean fails
            # Assuming clean_ecg_advanced is available in scope as seen in other funcs
            if 'clean_ecg_advanced' in globals():
                ecg_clean, _ = clean_ecg_advanced(
                    ecg_signal, sampling_rate=sfreq,
                    median_window_ms=300, wavelet='db4',
                    wavelet_levels=5, zscore_threshold=3.0
                )
            else:
                ecg_clean = ecg_signal # Fallback
            
            _, rpk = nk.ecg_peaks(ecg_clean, sampling_rate=sfreq)
            rpeaks = rpk['ECG_R_Peaks']
            return rpeaks, sfreq, data, ch_names
        except Exception:
            return None, sfreq, data, ch_names

    if 'st' in globals():
        progress_bar = st.progress(0)
        status_text = st.empty()

    for file_idx, (file_path, patient_id, raw) in enumerate(loaded_files):
        if 'st' in globals():
            status_text.text(f"Circular Stats: Processing {patient_id}...")
            progress_bar.progress((file_idx) / len(loaded_files))

        rpeaks, sfreq, data, ch_names = get_rpeaks(raw)
        if rpeaks is None or len(rpeaks) < 2:
            continue
            
        # Filter for EEG channels
        eeg_indices = [i for i, ch in enumerate(ch_names) 
                       if 'eeg' in ch.lower() or any(elec in ch.upper() for elec in ['FP', 'F', 'C', 'P', 'O', 'T', 'A'])]
        
        patient_results = []
        
        for ch_idx in eeg_indices:
            electrode = ch_names[ch_idx]
            eeg_signal = data[ch_idx, :]
            
            # Extract phases of max amplitude in each R-R interval
            peak_phases = []
            
            for i in range(len(rpeaks) - 1):
                start = rpeaks[i]
                end = rpeaks[i+1]
                length = end - start
                if length <= 0:
                    continue
                
                # Segment
                segment = eeg_signal[start:end]
                
                # Find index of max absolute amplitude
                if len(segment) == 0:
                    continue
                
                # Option: Max amplitude
                peak_idx = np.argmax(np.abs(segment))
                
                # Convert to phase (0 to 2pi)
                # 0 = R[i], 2pi = R[i+1]
                phase = 2 * np.pi * (peak_idx / length)
                peak_phases.append(phase)
            
            if len(peak_phases) > 10: # Minimum trials
                peak_phases = np.array(peak_phases)
                p_val, r_len = calc_rayleigh(peak_phases)
                
                # Calculate mean angle
                mean_angle = stats.circmean(peak_phases)
                
                is_sig = p_val < 0.05
                
                patient_results.append({
                    'electrode': electrode,
                    'p_value': p_val,
                    'mean_diff': r_len, # Using R-length as proxy for strength for now
                    'mean_angle': mean_angle,
                    'significant': is_sig,
                    'phases': peak_phases # Store for potential WW/Kuiper later
                })
        
        sig_results_rayleigh[patient_id] = patient_results
        
        # NOTE: Watson-Williams and Kuiper require Group comparisons (e.g. Healthy vs Patient)
        # or Condition comparisons (Systole vs Diastole tasks? Hits vs Misses?).
        # Currently we only have one list of files. 
        # So we leave sig_results_ww and sig_results_kuiper empty or as placeholders.
        
    if 'st' in globals():
        progress_bar.empty()
        status_text.empty()
        
    return sig_results_rayleigh

def plot_peth_raster(loaded_files, file_index=0, minmax=(-0.5, 1.0), binsize=0.05, save_plot=False, save_dir=None, aggregation_method='median', processing_mode='single', sleep_stage_mapping=None, sleep_stage='N1', jitter_amount=0.25, test_type="Cluster-Based Permutation"):
    """
    Create a Peri-Event Time Histogram (PETH) raster plot using pynapple.
    
    This function loads a file from loaded_files, extracts R-peaks from ECG,
    and creates a raster plot showing the timing of events around R-peaks.
    
    Parameters
    ----------
    loaded_files : list
        List of tuples (file_path, patient_id, raw_object)
    file_index : int, optional
        Index of the file to process (default: 0)
    minmax : tuple, optional
        Time window around events in seconds (default: (-0.5, 1.0))
    binsize : float, optional
        Bin size for rate calculation in seconds (default: 0.05)
    save_plot : bool, optional
        Whether to save the plot (default: False)
    save_dir : str, optional
        Directory to save the plot (default: None)
    aggregation_method : str, optional
        Method for aggregating across patients: 'median' or 'mean' (default: 'median')
    processing_mode : str, optional
        Processing mode: 'single' for sequential or 'multithread' for parallel (default: 'single')
    sleep_stage_mapping : dict, optional
        Dictionary mapping patient_id to sleep_stage (default: None)
    sleep_stage : str, optional
        Current sleep stage being analyzed (default: 'N1')
    jitter_amount : float, optional
         Amount of jitter to apply (default: 0.0)
    test_type: str, optional
        Test type (default: "Cluster-Based Permutation")
    
    Returns
    -------
    matplotlib.figure.Figure or None
        Figure object if plot was created, None otherwise
    """
    
    significance_results = {}
    test_name = "Unknown"
    
    if test_type == "Cluster-Based Permutation":
        test_name = "Cluster-Based Permutation"
        # Run permutation test with jitter support
        significance_results, significance_results_all_t = test_electrode_significance_per_patient(
            loaded_files, 
            minmax=minmax, 
            binsize=binsize,
            test_window=(-0.1, 0.1), 
            baseline_window=(-0.5, -0.2),
            n_bootstrap=100, 
            significance_threshold=0.05, 
            processing_mode=processing_mode,
            jitter_amount=jitter_amount
        )
    else: 
        # Default to Circular Statistics
        test_name = "Rayleigh"
        significance_results = test_circular_statistics_r_to_r(loaded_files)

    # title for the stat test type
    if 'st' in globals():
        st.title(f"Statistics Results: {test_name}")
        if jitter_amount > 0 and test_type == "Cluster-Based Permutation":
             st.info(f"Using Jittered Data (Control): Jitter = {jitter_amount}s")

        electrode_counts = {}
        for patient_id, patient_results in significance_results.items():
            for result in patient_results:
                electrode = result['electrode']
                if electrode not in electrode_counts:
                    electrode_counts[electrode] = {'significant': 0, 'total': 0}
                    electrode_counts[electrode]['total'] += 1
                    if result['significant']:
                        electrode_counts[electrode]['significant'] += 1
            
        topomap_data = []
        for electrode, counts in electrode_counts.items():
            percentage = (counts['significant'] / counts['total']) * 100 if counts['total'] > 0 else 0
            topomap_data.append({
                'electrode': electrode,
                'percentage_significant': percentage,
                'n_significant': counts['significant'],
                'n_total': counts['total']
            })
            
        if topomap_data:
            topomap_df = pd.DataFrame(topomap_data)
            # Cache CSV data in session state to prevent reruns
            cache_key = "topomap_percentage_csv"
            if cache_key not in st.session_state:
                st.session_state[cache_key] = topomap_df.to_csv(index=False)
                
            st.download_button(
                label="📥 Download Topomap Percentage Data (CSV)",
                data=st.session_state[cache_key],
                file_name=f"topomap_percentage_data_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.csv",
                mime="text/csv",
                key="download_topomap_data",
                use_container_width=False
            )
        
        # Create averaged PETH plots for significant patients only
        if len(loaded_files) > 1:
            if 'st' in globals():
                st.markdown("---")
                # Determine aggregation label for title
                agg_label = "Median" if aggregation_method == 'median' else "Average"
                st.markdown(f"## 📊 {agg_label} PETH Plots (Significant Patients Only)")
            
            plot_significant_patients_electrode_peth(loaded_files, significance_results, 
                                                    minmax=minmax, binsize=binsize,
                                                    save_plot=save_plot, save_dir=save_dir,
                                                    aggregation_method=aggregation_method, processing_mode=processing_mode,
                                                    sleep_stage_mapping=sleep_stage_mapping, sleep_stage=sleep_stage)
            
            # Create PETH and raster plots with one line per patient
            if 'st' in globals():
                st.markdown("---")
                st.markdown("## 📊 PETH & Raster Plots (One Line per Patient)")
            
            plot_peth_raster_per_patient(loaded_files, significance_results, 
                                        minmax=minmax, binsize=binsize,
                                        save_plot=save_plot, save_dir=save_dir,
                                        aggregation_method=aggregation_method, processing_mode=processing_mode,
                                        sleep_stage_mapping=sleep_stage_mapping, sleep_stage=sleep_stage)
            
            if 'st' in globals():
                significance_data = []
                for patient_id, patient_results in significance_results.items():
                    for result in patient_results:
                        if result.get('significant', False):
                            significance_data.append({
                                'patient_id': patient_id,
                                'electrode': result['electrode'],
                                'p_value': result['p_value'],
                                'mean_diff': result['mean_diff'],
                                'significant': result['significant']
                            })
                
                if significance_data:
                    sig_patients_df = pd.DataFrame(significance_data)
                    # Cache CSV data in session state to prevent reruns
                    cache_key = "significant_patients_peth_csv"
                    if cache_key not in st.session_state:
                        st.session_state[cache_key] = sig_patients_df.to_csv(index=False)
                    
                    if 'st' in globals():
                        st.markdown("---")
                        st.markdown("## 🗺️ Electrode Significance Topomap (Percentage of Patients)")
                    
                    plot_significance_percentage_topomap(significance_results, save_plot=save_plot, save_dir=save_dir)

                    st.download_button(
                        label="📥 Download Significant Patients PETH Data (CSV)",
                        data=st.session_state[cache_key],
                        file_name=f"significant_patients_peth_data_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.csv",
                        mime="text/csv",
                        key="download_significant_patients_peth",
                        use_container_width=False
                    )
                
    
    # Prepare electrode data for download buttons (used by multiple plot sections)
    all_electrode_data = []
    if 'st' in globals() and significance_results:
        for patient_id, patient_results in significance_results.items():
            for result in patient_results:
                all_electrode_data.append({
                    'patient_id': patient_id,
                    'electrode': result['electrode'],
                    'p_value': result['p_value'],
                    'mean_diff': result['mean_diff'],
                    'significant': result['significant']
                })
    
        # Add download button for averaged plots data (using significance results)
        if 'st' in globals() and all_electrode_data:
            averaged_df = pd.DataFrame(all_electrode_data)
            # Cache CSV data in session state to prevent reruns
            cache_key = "averaged_eeg_peth_csv"
            if cache_key not in st.session_state:
                st.session_state[cache_key] = averaged_df.to_csv(index=False)
            
            st.download_button(
                label="📥 Download Averaged EEG Channel PETH Data (CSV)",
                data=st.session_state[cache_key],
                file_name=f"averaged_eeg_peth_data_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.csv",
                mime="text/csv",
                key="download_averaged_eeg_peth",
                use_container_width=False
            )
        
        # Create regional averaged plots
        if 'st' in globals():
            st.markdown("---")
            # Determine aggregation label for title
            agg_label = "Median" if aggregation_method == 'median' else "Average"
            st.markdown(f"## 📊 Regional {agg_label} EEG PETH & Raster Plots (F, C, T, P, O - First 21 Electrodes)")
        
        plot_regional_averaged_eeg_peth_raster(loaded_files, minmax=minmax, binsize=binsize,
                                                save_plot=save_plot, save_dir=save_dir,
                                                aggregation_method=aggregation_method, processing_mode=processing_mode,
                                                sleep_stage_mapping=sleep_stage_mapping, sleep_stage=sleep_stage)
        
        # Add download button for regional averaged plots data
        if 'st' in globals() and all_electrode_data:
            regional_df = pd.DataFrame(all_electrode_data)
            # Cache CSV data in session state to prevent reruns
            cache_key = "regional_averaged_peth_csv"
            if cache_key not in st.session_state:
                st.session_state[cache_key] = regional_df.to_csv(index=False)
            
            st.download_button(
                label="📥 Download Regional Averaged PETH Data (CSV)",
                data=st.session_state[cache_key],
                file_name=f"regional_averaged_peth_data_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.csv",
                mime="text/csv",
                key="download_regional_averaged_peth",
                use_container_width=False
            )
        
        # Create regional electrode plots (one plot per region with subplots for each electrode)
        if 'st' in globals():
            st.markdown("---")
            st.markdown("## 📊 Regional Electrode Plots (F, P, T, O, C - One Plot Per Region)")
        
        plot_regional_electrode_averaged_eeg(loaded_files, minmax=minmax, binsize=binsize,
                                            save_plot=save_plot, save_dir=save_dir,
                                            aggregation_method=aggregation_method, processing_mode=processing_mode,
                                            sleep_stage_mapping=sleep_stage_mapping, sleep_stage=sleep_stage)
        
        # Add download button for regional electrode plots data
        if 'st' in globals() and all_electrode_data:
            regional_elec_df = pd.DataFrame(all_electrode_data)
            # Cache CSV data in session state to prevent reruns
            cache_key = "regional_electrode_peth_csv"
            if cache_key not in st.session_state:
                st.session_state[cache_key] = regional_elec_df.to_csv(index=False)
            
            st.download_button(
                label="📥 Download Regional Electrode PETH Data (CSV)",
                data=st.session_state[cache_key],
                file_name=f"regional_electrode_peth_data_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.csv",
                mime="text/csv",
                key="download_regional_electrode_peth",
                use_container_width=False
            )
    
    return significance_results
        
    

def MI_systole_diastole_comparison(loaded_files, sleep_stage, output_dir, save_results):

    # Process all files (with MI calculation)
    with st.spinner("Processing files to extract HEP and calculate modulation index..."):
        df, mi_df = process_all_files_with_mi(loaded_files)
    
    if df.empty:
        st.error("No data extracted. Exiting.")
        return
    
    # Save results to CSV
    if save_results:
        csv_path = os.path.join(output_dir, f'hep_results_{sleep_stage}.csv')
        df.to_csv(csv_path, index=False)
        st.success(f"Results saved to {csv_path}")
    
    # Add download button for HEP results
    # Cache CSV data in session state to prevent reruns
    cache_key = f"hep_results_csv_{sleep_stage}"
    if cache_key not in st.session_state:
        st.session_state[cache_key] = df.to_csv(index=False)
    
    st.download_button(
        label="📥 Download HEP Results Data (CSV)",
        data=st.session_state[cache_key],
        file_name=f"hep_results_{sleep_stage}_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.csv",
        mime="text/csv",
        key=f"download_hep_results_{sleep_stage}",
        use_container_width=False
    )
    
    # Display data summary
    st.markdown("## 📋 Data Summary")
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("Total Patients", len(df['patient_id'].unique()))
    with col2:
        st.metric("Systole Records", len(df[df['phase'] == 'systole']))
    with col3:
        st.metric("Diastole Records", len(df[df['phase'] == 'diastole']))
    
    # Save modulation index results
    if not mi_df.empty:
        if save_results:
            mi_csv_path = os.path.join(output_dir, f'modulation_index_{sleep_stage}.csv')
            mi_df.to_csv(mi_csv_path, index=False)
            st.success(f"Modulation index results saved to {mi_csv_path}")
        
        # Display summary statistics for MI
        st.markdown("## 📊 Modulation Index (Tort et al. 2010) Summary")
        st.markdown("---")
        
        col1, col2, col3, col4 = st.columns(4)
        with col1:
            st.metric("Mean MI", f"{mi_df['modulation_index'].mean():.4f} ± {mi_df['modulation_index'].std():.4f}")
        with col2:
            st.metric("Mean Correlation", f"{mi_df['correlation'].mean():.4f} ± {mi_df['correlation'].std():.4f}")
        with col3:
            st.metric("Mean Amplitude Systole", f"{mi_df['mean_amplitude_systole'].mean():.4f} ± {mi_df['mean_amplitude_systole'].std():.4f}")
        with col4:
            st.metric("Mean Amplitude Diastole", f"{mi_df['mean_amplitude_diastole'].mean():.4f} ± {mi_df['mean_amplitude_diastole'].std():.4f}")
        
        # Display MI dataframe
        st.markdown("### Modulation Index Data")
        st.dataframe(mi_df, use_container_width=True)
        
        # Add download button for modulation index data
        # Cache CSV data in session state to prevent reruns
        cache_key = f"modulation_index_csv_{sleep_stage}"
        if cache_key not in st.session_state:
            st.session_state[cache_key] = mi_df.to_csv(index=False)
        
        st.download_button(
            label="📥 Download Modulation Index Data (CSV)",
            data=st.session_state[cache_key],
            file_name=f"modulation_index_{sleep_stage}_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.csv",
            mime="text/csv",
            key=f"download_modulation_index_{sleep_stage}",
            use_container_width=False
        )
    else:
        st.warning("No modulation index data calculated.")
    
    # Perform statistical comparison
    st.markdown("---")
    with st.spinner("Performing statistical comparison..."):
        stats_results = compare_systole_diastole(df)
    
    if stats_results:
        # Display statistics
        print_statistics(stats_results)
        
        # Save statistics to CSV
        if save_results:
            stats_df = pd.DataFrame(stats_results).T
            stats_csv_path = os.path.join(output_dir, f'hep_statistics_{sleep_stage}.csv')
            stats_df.to_csv(stats_csv_path)
            st.success(f"Statistics saved to {stats_csv_path}")
        
        # Add download button for statistics
        # Cache CSV data in session state to prevent reruns
        cache_key = f"statistics_csv_{sleep_stage}"
        if cache_key not in st.session_state:
            stats_df = pd.DataFrame(stats_results).T
            st.session_state[cache_key] = stats_df.to_csv()
        
        st.download_button(
            label="📥 Download Statistical Comparison Data (CSV)",
            data=st.session_state[cache_key],
            file_name=f"hep_statistics_{sleep_stage}_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.csv",
            mime="text/csv",
            key=f"download_statistics_{sleep_stage}",
            use_container_width=False
        )
        
        # Create plots
        st.markdown("---")
        st.markdown("## 📈 Comparison Plots")
        with st.spinner("Creating comparison plots..."):
            plot_comparison(df, stats_results, save_path=None)
        
        # Add download button next to comparison plots
        # Cache CSV data in session state to prevent reruns
        cache_key = f"comparison_plot_csv_{sleep_stage}"
        if cache_key not in st.session_state:
            st.session_state[cache_key] = df.to_csv(index=False)
        
        st.download_button(
            label="📥 Download Comparison Plot Data (CSV)",
            data=st.session_state[cache_key],
            file_name=f"comparison_plot_data_{sleep_stage}_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.csv",
            mime="text/csv",
            key=f"download_comparison_plot_data_{sleep_stage}",
            use_container_width=False
        )
    else:
        st.warning("Could not perform statistical comparison")
        # Still create plots without statistics
        st.markdown("## 📈 Comparison Plots")
        plot_comparison(df, None, save_path=None)
    
    # Create PairGrid plot with systole/diastole as hue
    st.markdown("---")
    st.markdown("## 🔄 PairGrid Plot")
    with st.spinner("Creating PairGrid plot..."):
        plot_hep_pairgrid(df, sleep_stage=sleep_stage, save_plot=save_results, save_dir=output_dir)
    
    # Add download button for PairGrid plot data
    # Cache CSV data in session state to prevent reruns
    cache_key = f"pairgrid_csv_{sleep_stage}"
    if cache_key not in st.session_state:
        st.session_state[cache_key] = df.to_csv(index=False)
    
    st.download_button(
        label="📥 Download PairGrid Plot Data (CSV)",
        data=st.session_state[cache_key],
        file_name=f"pairgrid_plot_data_{sleep_stage}_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.csv",
        mime="text/csv",
        key=f"download_pairgrid_data_{sleep_stage}",
        use_container_width=False
    )
    
    st.success("✅ Analysis complete!")


def plot_hep_pairgrid(df, sleep_stage='N1', save_plot=True, save_dir='hep_comparison_results'):
    """
    Plot HEP metrics using PairGrid with systole/diastole as hue.
    
    Parameters
    ----------
    df : pd.DataFrame
        DataFrame with HEP results containing 'phase', 'mean_amplitude', 
        'peak_amplitude', 'rms', 'std_amplitude' columns
    sleep_stage : str
        Sleep stage identifier for naming
    save_plot : bool
        Whether to save the plot
    save_dir : str
        Directory to save the plot
    
    Returns
    -------
    None
    """
    if df.empty:
        if 'st' in globals():
            st.warning("No data to plot in PairGrid")
        else:
            print("No data to plot in PairGrid")
        return
    
    # Check if we have both systole and diastole data
    phases = df['phase'].unique()
    if 'systole' not in phases or 'diastole' not in phases:
        if 'st' in globals():
            st.warning("Need both systole and diastole data for PairGrid plot")
        else:
            print("Need both systole and diastole data for PairGrid plot")
        return
    
    # Get common patients (paired data)
    systole_df = df[df['phase'] == 'systole'].sort_values('patient_id')
    diastole_df = df[df['phase'] == 'diastole'].sort_values('patient_id')
    
    common_patients = set(systole_df['patient_id']).intersection(set(diastole_df['patient_id']))
    
    if len(common_patients) < 2:
        if 'st' in globals():
            st.warning("Not enough paired data for PairGrid plot")
        else:
            print("Not enough paired data for PairGrid plot")
        return
    
    # Filter to common patients and ensure same order
    systole_paired = systole_df[systole_df['patient_id'].isin(common_patients)].sort_values('patient_id')
    diastole_paired = diastole_df[diastole_df['patient_id'].isin(common_patients)].sort_values('patient_id')
    
    # Verify patient order matches
    if not np.array_equal(systole_paired['patient_id'].values, diastole_paired['patient_id'].values):
        if 'st' in globals():
            st.info("Patient order mismatch, sorting...")
        else:
            print("Patient order mismatch, sorting...")
        systole_paired = systole_paired.sort_values('patient_id')
        diastole_paired = diastole_paired.sort_values('patient_id')
    
    # Define metrics to plot
    metrics = ['mean_amplitude', 'peak_amplitude', 'rms', 'std_amplitude']
    metric_labels = ['Mean Amplitude', 'Peak Amplitude', 'RMS', 'Std Amplitude']
    
    # Find common valid indices across all metrics
    n_patients = len(systole_paired)
    
    # Start with all indices as potentially valid
    valid_mask = np.ones(n_patients, dtype=bool)
    
    # Check each metric and find common valid indices
    for metric in metrics:
        if metric in systole_paired.columns and metric in diastole_paired.columns:
            systole_vals = systole_paired[metric].values
            diastole_vals = diastole_paired[metric].values
            metric_valid = ~(np.isnan(systole_vals) | np.isnan(diastole_vals))
            valid_mask = valid_mask & metric_valid
    
    valid_indices = np.where(valid_mask)[0]
    
    if len(valid_indices) < 2:
        if 'st' in globals():
            st.warning("Not enough valid data across all metrics for PairGrid plot")
        else:
            print("Not enough valid data across all metrics for PairGrid plot")
        return
    
    # Extract arrays for each metric using common valid indices
    arrs = []
    labels = []
    for metric, label in zip(metrics, metric_labels):
        if metric in systole_paired.columns and metric in diastole_paired.columns:
            systole_vals = systole_paired[metric].values[valid_indices]
            diastole_vals = diastole_paired[metric].values[valid_indices]
            
            # Combine systole and diastole values
            combined_vals = np.concatenate([systole_vals, diastole_vals])
            arrs.append(combined_vals)
            labels.append(label)
    
    if len(arrs) < 2:
        if 'st' in globals():
            st.warning("Not enough valid metrics for PairGrid plot")
        else:
            print("Not enough valid metrics for PairGrid plot")
        return
    
    # Create hue array (systole or diastole) matching the combined arrays
    n_valid = len(valid_indices)
    hue = pd.Series(['Systole'] * n_valid + ['Diastole'] * n_valid)
    
    # Create time array (patient index)
    t = np.arange(len(arrs[0]))
    
    # Create name for the plot
    plot_name = f'hep_pairgrid_{sleep_stage}'
    
    # Call _plot_metric_vs_hrv
    if 'st' in globals():
        st.info("Creating PairGrid plot for HEP metrics (hue: Systole/Diastole)...")
    else:
        print("\nCreating PairGrid plot for HEP metrics (hue: Systole/Diastole)...")
    
    _plot_metric_vs_hrv(
        t=t,
        arrs=arrs,
        labels=labels,
        save_plot=save_plot,
        name=plot_name,
        save_dir=save_dir,
        is_streamlit=('st' in globals()),
        hue=hue,
        size_values=None
    )
    
    if save_plot and 'st' not in globals():
        print(f"PairGrid plot saved to {save_dir}/pairgrid_all_{plot_name}.png")



def plot_first_patient_eeg_ecg_aligned(loaded_files, sleep_stage, output_dir, save_results, binsize=0.05, aggregation_method='median', processing_mode='single', sleep_stage_mapping=None):
    """
    Plot EEG data aligned to ECG R-peaks for the first patient.
    Plots y-axis in microVolts vs x-axis time (-0.5, 1.0) s.
    One plot per electrode (grouped by region).
    """
    if not PYNAPPLE_AVAILABLE:
        if 'st' in globals():
            st.error("pynapple is not available.")
        return

    if not loaded_files:
        if 'st' in globals():
            st.warning("No files loaded.")
        return

    # Use first patient
    file_path, patient_id, raw = loaded_files[0]
    
    if 'st' in globals():
        st.info(f"Plotting aligned EEG for first patient: {patient_id}")
    
    # Parameters
    minmax = (-0.5, 1.0)
    
    # Get sampling frequency and data
    sfreq = raw.info['sfreq']
    ch_names = raw.ch_names
    data = raw.get_data()
    
    # Process ECG
    ch_lower = [ch.lower() for ch in ch_names]
    ecg_indices = [i for i, ch in enumerate(ch_lower) if 'ecg' in ch or 'ekg' in ch]
    
    if not ecg_indices:
        if 'st' in globals():
            st.warning(f"No ECG channel found in {patient_id}")
        return
        
    ecg_ch_idx = ecg_indices[0]
    ecg_signal = data[ecg_ch_idx, :]
    
    # Clean ECG
    try:
        ecg_signal_clean, _ = clean_ecg_advanced(
            ecg_signal, 
            sampling_rate=sfreq,
            median_window_ms=300,
            wavelet='db4',
            wavelet_levels=5,
            zscore_threshold=3.0
        )
        _, rpk = nk.ecg_peaks(ecg_signal_clean, sampling_rate=sfreq)
        rpeaks = rpk['ECG_R_Peaks']
    except Exception as e:
        if 'st' in globals():
            st.error(f"Error processing ECG for {patient_id}: {e}")
        return

    if len(rpeaks) < 2:
        if 'st' in globals():
            st.warning("Not enough R-peaks found.")
        return

    rpeak_times = rpeaks / sfreq
    rpeak_ts = nap.Ts(t=rpeak_times, time_units="s")
    
    # Process EEG Channels
    eeg_ch_indices = [i for i, ch in enumerate(ch_names) 
                    if 'eeg' in ch.lower() or any(elec in ch.upper() for elec in ['FP', 'F', 'C', 'P', 'O', 'T', 'A'])]
    
    if not eeg_ch_indices:
        if 'st' in globals():
            st.warning("No EEG channels found.")
        return

    # Regional grouping
    regional_electrodes = {
        'F': ['Fp1', 'Fpz', 'Fp2', 'F7', 'F3', 'Fz', 'F4', 'F8'],
        'P': ['P3', 'Pz', 'P4'],
        'T': ['T3', 'T4', 'T5', 'T6'],
        'O': ['O1', 'Oz', 'O2'],
        'C': ['C3', 'Cz', 'C4']
    }
    
    # Map channels to regions
    channels_by_region = {region: [] for region in regional_electrodes}
    other_channels = []
    
    for idx in eeg_ch_indices:
        ch_name = ch_names[idx]
        matched = False
        for region, elecs in regional_electrodes.items():
            if ch_name.upper() in [e.upper() for e in elecs]:
                channels_by_region[region].append((idx, ch_name))
                matched = True
                break
        if not matched:
            other_channels.append((idx, ch_name))
            
    # Remove empty regions
    channels_by_region = {k: v for k, v in channels_by_region.items() if v}
    if other_channels:
        channels_by_region['Other'] = other_channels

    # Plotting
    for region, channels in channels_by_region.items():
        if not channels:
            continue
            
        n_electrodes = len(channels)
        n_cols = 1
        n_rows = n_electrodes
        
        fig, axes = plt.subplots(n_rows, n_cols, figsize=(10, 3 * n_rows), sharex=True)
        if n_electrodes == 1:
            axes = np.array([axes])
        else:
            axes = axes.flatten()
            
        for i, (ch_idx, ch_name) in enumerate(channels):
            ax = axes[i]
            eeg_signal = data[ch_idx, :]
            
            # PETH
            time_index = np.arange(len(eeg_signal)) / sfreq
            eeg_tsd = nap.Tsd(t=time_index, d=eeg_signal, time_units="s")
            
            peth = nap.compute_perievent_continuous(
                timeseries=eeg_tsd,
                tref=rpeak_ts,
                minmax=minmax,
                time_unit="s"
            )
            
            if len(peth) > 0:
                # Mean across trials
                if hasattr(peth, 'values') and peth.values.ndim == 2:
                    mean_signal = np.nanmean(peth.values, axis=1)
                    time_axis = peth.index.values
                    
                    # Plot all trials
                    # ax.plot(time_axis, peth.values, color='gray', alpha=0.2, linewidth=0.5)
                    
                else:
                     ax.text(0.5, 0.5, "No data dimensions", ha='center')
                     continue
                
                ax.plot(time_axis, mean_signal, color='blue', linewidth=1.5, label='Mean')
                
                # Add +/- 1 SD shading
                std_signal = np.nanstd(peth.values, axis=1)
                ax.fill_between(time_axis, mean_signal - std_signal, mean_signal + std_signal, 
                                color='blue', alpha=0.2, label='±1 SD')
                                
                ax.set_title(f"{ch_name} (aligned to ECG)", fontsize=10, fontweight='bold')
                ax.set_ylabel("Amplitude (V)", fontsize=9)
                ax.grid(True, alpha=0.3)
                ax.axvline(0, color='red', linestyle='--', alpha=0.5)
                # add legend
                ax.legend()
                
        axes[-1].set_xlabel("Time (s)", fontsize=10)
        plt.suptitle(f"{patient_id} - Region {region} - Aligned EEG", y=0.995)
        plt.tight_layout()
        
        # Display
        display_plot_with_download(fig, f"aligned_eeg_{patient_id}_{region}", f"aligned_eeg_{patient_id}_{region}")
        plt.close(fig)

if __name__ == "__main__":
    main()
import os
import sys
import argparse

import streamlit as st

import pickle

import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
import mne
import neurokit2 as nk
import numpy as np
import pandas as pd
from scipy import signal
from scipy.ndimage import median_filter
from scipy.signal import butter, filtfilt, find_peaks
from scipy import stats
from scipy.stats import kurtosis
from scipy.ndimage import label
from typing import Tuple
import re
from collections import Counter
from concurrent.futures import ProcessPoolExecutor, as_completed
import math
import threading
import time
from mne_icalabel import label_components

DEFAULT_PATIENT_WORKERS = max(1, min(8, os.cpu_count() or 1))
MIN_COMPLETE_PATIENT_PICKLE_BYTES = 64 * 1024
STANDARD_1020_EEG_CHANNELS = (
    "Fp1", "Fp2", "F7", "F3", "Fz", "F4", "F8",
    "T3", "T7", "C3", "Cz", "C4", "T4", "T8",
    "T5", "P7", "P3", "Pz", "P4", "T6", "P8",
    "O1", "Oz", "O2",
)
STANDARD_1020_EEG_CHANNEL_KEYS = {
    channel.lower() for channel in STANDARD_1020_EEG_CHANNELS
}
DEFAULT_HEP_GROUPS = [
    'The_Human_Sleep_Project',
    'Young_The_Human_Sleep_Project',
    'Harvard_Electroencephalography',
]
MAX_SPECTRAL_POWER_RATIO = 0.4
HEP_WINDOW_MIN_PTP_UV = 0.05
HEP_WINDOW_MIN_STD_UV = 0.01
HEP_WINDOW_MAX_ABS_UV = 5000.0
HEP_WINDOW_MAX_ROUGHNESS = 2.2
HEP_WINDOW_MAX_BAD_CHANNEL_FRACTION = 0.35

def default_hep_groups(available_groups, fallback_count=2):
    groups = [g for g in DEFAULT_HEP_GROUPS if g in available_groups]
    return groups or available_groups[:fallback_count]

def _parse_mode_arg():
    """Parse --mode argument from sys.argv (Streamlit strips the -- separator)."""
    try:
        parser = argparse.ArgumentParser(add_help=False)
        parser.add_argument('--mode', type=str, default=None)
        parsed, _ = parser.parse_known_args(sys.argv[1:])
        return parsed.mode
    except SystemExit:
        return None

def get_peaks(signal_data, use_smart_threshold=True, mad_multiplier=10, **kwargs):
    """
    Centralized function to find peaks in a signal.
    By default, calculates a robust smart threshold (different for each patient)
    based on the signal's Median Absolute Deviation (MAD), unless 'height' 
    is explicitly provided in kwargs.
    Wraps scipy.signal.find_peaks.
    """
    if use_smart_threshold and 'height' not in kwargs:
        # Calculate dynamic threshold for the specific patient's signal 
        mad = np.median(np.abs(signal_data - np.median(signal_data))) + 1e-12
        smart_height = np.median(signal_data) + mad_multiplier * mad
        kwargs['height'] = smart_height

    return find_peaks(signal_data, **kwargs)

try:
    import pywt
    PYWT_AVAILABLE = True
except ImportError:
    PYWT_AVAILABLE = False

try:
    import torch
    import torch.nn as nn
    TORCH_AVAILABLE = True
except ImportError:
    TORCH_AVAILABLE = False

try:
    import pynapple as nap
    PYNAPPLE_AVAILABLE = True
except ImportError:
    PYNAPPLE_AVAILABLE = False
    st.error("Pynapple (nap) not installed. Please install it.")


try:
    import heartpy as hp
    HEARTPY_AVAILABLE = True
except ImportError:
    HEARTPY_AVAILABLE = False

import io
import traceback

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

st.set_page_config(layout="wide")
if not hasattr(st, "subtitle"):
    st.subtitle = st.subheader
if not hasattr(st, "divide"):
    st.divide = st.divider if hasattr(st, "divider") else lambda: st.markdown("---")

_original_st_pyplot = st.pyplot
_original_st_plotly_chart = getattr(st, "plotly_chart", None)

ICA_CLEANED_TITLE_SUFFIX = " (ICA cleaned)"
ICA_CLEANED_FIGURE_NOTE = "This plot uses ICA-cleaned data."
REPETITIVE_RPEAK_EXCLUDE_PERC = 20.0
HEP_FLIP_CACHE_SUFFIX = "_rpeak_shape_v1"

def _append_ica_cleaned_to_title(title):
    if not title or ICA_CLEANED_TITLE_SUFFIX in title:
        return title
    return f"{title}{ICA_CLEANED_TITLE_SUFFIX}"

def _append_variability_label_to_title(title, labels):
    if not title:
        return title
    title_upper = str(title).upper()
    label_text = " ".join(str(label) for label in (labels or [])).upper()
    if "MAD" in label_text and "MAD" not in title_upper:
        return f"{title} | Median +/- MAD"
    if ("SEM" in label_text or "STANDARD ERROR" in label_text) and "SEM" not in title_upper:
        return f"{title} | Mean +/- SEM"
    if (
        "STD" in label_text
        or " SD" in label_text
        or "± 1 SD" in label_text
        or "+/- 1 SD" in label_text
        or "STANDARD DEVIATION" in label_text
    ) and not any(token in title_upper for token in ("STD", " SD", "STANDARD DEVIATION")):
        return f"{title} | Mean +/- STD"
    return title

def _mark_matplotlib_titles_with_variability_label(fig):
    if fig is None:
        return
    try:
        all_labels = []
        for ax in fig.axes:
            _, labels = ax.get_legend_handles_labels()
            all_labels.extend(labels)
            title = ax.get_title()
            if title:
                ax.set_title(_append_variability_label_to_title(title, labels))
        if fig._suptitle:
            fig._suptitle.set_text(
                _append_variability_label_to_title(fig._suptitle.get_text(), all_labels)
            )
    except Exception:
        pass

def _mark_matplotlib_titles_ica_cleaned(fig):
    if fig is None:
        return
    try:
        _mark_matplotlib_titles_with_variability_label(fig)
        if getattr(fig, "_skip_ica_cleaned_title", False):
            return
        if fig._suptitle:
            fig._suptitle.set_text(_append_ica_cleaned_to_title(fig._suptitle.get_text()))
        for ax in fig.axes:
            title = ax.get_title()
            if title:
                ax.set_title(_append_ica_cleaned_to_title(title))
    except Exception:
        pass

def _mark_plotly_titles_ica_cleaned(fig):
    if fig is None:
        return
    try:
        labels = [getattr(trace, "name", "") for trace in getattr(fig, "data", []) or []]
        title = getattr(getattr(fig, "layout", None), "title", None)
        if title is not None and getattr(title, "text", None):
            labeled_title = _append_variability_label_to_title(title.text, labels)
            fig.update_layout(title_text=_append_ica_cleaned_to_title(labeled_title))
        for annotation in getattr(getattr(fig, "layout", None), "annotations", []) or []:
            if getattr(annotation, "text", None):
                annotation.text = _append_ica_cleaned_to_title(annotation.text)
    except Exception:
        pass

def _get_matplotlib_title(fig):
    title = "Figure"
    try:
        if fig._suptitle:
            title = fig._suptitle.get_text()
        else:
            for ax in fig.axes:
                if ax.get_title():
                    title = ax.get_title()
                    break
    except Exception:
        pass
    return title

def _get_matplotlib_description(_fig, _title):
    if getattr(_fig, "_skip_ica_cleaned_title", False):
        return ""
    return ICA_CLEANED_FIGURE_NOTE

def _get_plotly_title(fig):
    title = "Interactive Figure"
    try:
        layout_title = getattr(getattr(fig, "layout", None), "title", None)
        if layout_title is not None and getattr(layout_title, "text", None):
            title = layout_title.text
        else:
            for annotation in getattr(getattr(fig, "layout", None), "annotations", []) or []:
                if getattr(annotation, "text", None):
                    title = annotation.text
                    break
    except Exception:
        pass
    return title

def _get_plotly_description(fig, title):
    description = ICA_CLEANED_FIGURE_NOTE
    try:
        subplot_titles = [
            annotation.text
            for annotation in getattr(getattr(fig, "layout", None), "annotations", []) or []
            if getattr(annotation, "text", None)
        ]
        if subplot_titles:
            description += f" Subplots include: {', '.join(subplot_titles)}."
        trace_count = len(getattr(fig, "data", []) or [])
        if trace_count:
            description += f" It contains {trace_count} plotted trace{'s' if trace_count != 1 else ''}."
    except Exception:
        pass
    return description

def _is_generic_figure_title(title):
    return not title or str(title).strip() in {"Figure", "Interactive Figure"}

def _clear_matplotlib_rendered_title(fig, title):
    try:
        if fig is None:
            return
        if fig._suptitle:
            fig._suptitle.set_visible(False)
            fig._suptitle.set_text("")
        titled_axes = [ax for ax in fig.axes if ax.get_title()]
        if len(titled_axes) == 1 and titled_axes[0].get_title() == title:
            titled_axes[0].set_title("")
    except Exception:
        pass

def _clear_plotly_rendered_title(fig):
    try:
        if fig is not None:
            fig.update_layout(title_text="")
    except Exception:
        pass

def _write_plot_description(description):
    if not description.strip():
        return
    note_key = "_ica_cleaned_figure_note_written"
    if description.strip() == ICA_CLEANED_FIGURE_NOTE:
        if st.session_state.get(note_key, False):
            return
        st.session_state[note_key] = True
    st.caption(description)

def custom_st_pyplot(fig=None, clear_figure=None, **kwargs):
    if fig is None:
        fig = plt.gcf()

    _mark_matplotlib_titles_ica_cleaned(fig)

    title = _get_matplotlib_title(fig)
    description = _get_matplotlib_description(fig, title)
        

    if 'pptx_figures_data' not in st.session_state:
        st.session_state.pptx_figures_data = []

    if not _is_generic_figure_title(title):
        st.subtitle(title)
    _original_st_pyplot(fig, clear_figure=clear_figure, **kwargs)
    _write_plot_description(description)

    if PPTX_AVAILABLE:
        buf = io.BytesIO()
        fig.savefig(buf, format='png', bbox_inches='tight')
        buf.seek(0)

        st.session_state.pptx_figures_data.append({
            'title': title,
            'description': description,
            'image': buf
        })

st.pyplot = custom_st_pyplot

def custom_st_plotly_chart(fig=None, **kwargs):
    _mark_plotly_titles_ica_cleaned(fig)
    title = _get_plotly_title(fig)
    description = _get_plotly_description(fig, title)
    if not _is_generic_figure_title(title):
        st.subtitle(title)
    chart = _original_st_plotly_chart(fig, **kwargs)
    _write_plot_description(description)
    return chart

if _original_st_plotly_chart is not None:
    st.plotly_chart = custom_st_plotly_chart

def generate_pptx():
    if not PPTX_AVAILABLE:
        return None

    prs = Presentation()

    try:
        title_only_slide_layout = prs.slide_layouts[5]
    except IndexError:
        title_only_slide_layout = prs.slide_layouts[0]

    for item in st.session_state.get('pptx_figures_data', []):
        slide = prs.slides.add_slide(title_only_slide_layout)

        if slide.shapes.title:
            slide.shapes.title.text = item['title']

        pic_left = Inches(0.5)
        pic_top = Inches(1.5)
        pic_width = Inches(5.5)

        item['image'].seek(0)
        slide.shapes.add_picture(item['image'], pic_left, pic_top, width=pic_width)

        txBox_left = Inches(6.2)
        txBox_top = Inches(1.5)
        txBox_width = Inches(3.3)
        txBox_height = Inches(5.0)

        txBox = slide.shapes.add_textbox(txBox_left, txBox_top, txBox_width, txBox_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        # Add paragraphs
        p = tf.paragraphs[0]
        p.text = item['description']
        p.font.size = Pt(14)

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

def bandpass_ecg(x, fs, lo=3.0, hi=40.0, order=4):
    nyq = 0.5 * fs
    b, a = signal.butter(order, [lo/nyq, hi/nyq], btype="bandpass")
    return signal.filtfilt(b, a, x)

def detect_rpeaks_abs(x_filt, fs):
    # Make it robust to inversion: detect on absolute signal
    x = np.abs(x_filt)

    # Smooth a bit to emphasize QRS energy
    win = int(0.08 * fs)  # ~80 ms
    win = max(win, 3) | 1  # odd >= 3
    x_smooth = signal.savgol_filter(x, win, polyorder=2)

    # Adaptive-ish threshold using MAD
    mad = np.median(np.abs(x_smooth - np.median(x_smooth))) + 1e-12
    thr = np.median(x_smooth) + 3.5 * mad

    # Refractory period to avoid double detections
    distance = int(0.25 * fs)  # 250 ms
    peaks, _ = get_peaks(x_smooth, height=thr, distance=distance)
    return peaks

def decide_inversion_from_template(x_filt, rpeaks, fs, qrs_ms=120):
    # Build an average beat template centered at R
    half = int((qrs_ms / 1000.0) * fs)  # +- qrs_ms
    if len(rpeaks) < 5:
        return False, None  # not enough info

    # Keep only peaks that allow full window extraction
    rpeaks = rpeaks[(rpeaks > half) & (rpeaks < len(x_filt) - half)]
    if len(rpeaks) < 5:
        return False, None

    beats = np.stack([x_filt[p-half:p+half+1] for p in rpeaks], axis=0)

    # Robust average: median across beats resists outliers
    template = np.median(beats, axis=0)

    # Polarity decision: is the dominant QRS deflection negative?
    pos_peak = np.max(template)
    neg_peak = np.min(template)  # negative value
    
    # FLIP ONLY the scans were 90% of max is smaller than abs(min) of the r peaks.
    inverted = (abs(neg_peak) > 0.9 * abs(pos_peak))

    return inverted, template

def fix_inverted_ecg(ecg, fs, lo=3.0, hi=40.0):
    ecg = np.asarray(ecg).astype(float)

    ecg_f = bandpass_ecg(ecg, fs, lo=lo, hi=hi)
    rpeaks = detect_rpeaks_abs(ecg_f, fs)

    inverted, template = decide_inversion_from_template(ecg_f, rpeaks, fs)

    ecg_fixed = -ecg if inverted else ecg
    info = {
        "inverted_detected": bool(inverted),
        "n_rpeaks_used": int(len(rpeaks)),
        "rpeaks": rpeaks,
        "template_filtered": template,  # median beat on filtered signal
    }
    return ecg_fixed, info

def clean_ecg_high_fidelity(ecg_signal, sampling_rate, wavelet='db1', wavelet_levels=3, zscore_threshold=5.0):
    """
    Improved ECG cleaning to preserve R-peak amplitude and HEP components.
    """
    if len(ecg_signal) == 0:
        return ecg_signal, {'methods_applied': [], 'n_extreme_samples': 0}
    
    methods_applied = []
    info = {}
    cleaned_signal = ecg_signal.copy()
    
    # --- Step 1: Gentle Median Filter ---
    # REDUCED: 300ms was killing the R-peak. 
    # 10-20ms is enough to kill 'pops' without touching the R-wave.
    median_window_ms = 20 
    median_window_samples = int(np.round(median_window_ms * sampling_rate / 1000.0))
    if median_window_samples % 2 == 0:
        median_window_samples += 1
    
    median_filtered = median_filter(cleaned_signal, size=max(3, median_window_samples))
    cleaned_signal = median_filtered
    methods_applied.append('gentle_median_filter')
    
    # --- Step 2: Bandpass Filter (3-40 Hz) ---
    # Replaces Wavelet Denoising as per user request.
    try:
        lowcut = .5
        highcut = 40.0
        nyquist = 0.5 * sampling_rate
        low = lowcut / nyquist
        high = highcut / nyquist
        b, a = butter(2, [low, high], btype='band') # Order 2 is usually sufficient for ECG
        
        cleaned_signal = filtfilt(b, a, cleaned_signal)
        methods_applied.append('bandpass_3_40Hz')
    except Exception as e:
        if 'st' in globals():
            st.warning(f"Bandpass filter failed: {e}")
        # Fallback to original signal if filter fails
        cleaned_signal = cleaned_signal

    # --- Step 3: Outlier Handling (Z-score) ---
    # INCREASED: 3.0 often catches the R-wave itself. 5.0-6.0 catches actual artifacts.
    mean_val = np.mean(cleaned_signal)
    std_val = np.std(cleaned_signal)
    if std_val > 0:
        z_scores = (cleaned_signal - mean_val) / std_val
        extreme_mask = np.abs(z_scores) > zscore_threshold
        # Instead of replacing with median, consider clipping or ignoring 
        # to avoid 'flat-top' peaks.
        cleaned_signal[extreme_mask] = np.sign(cleaned_signal[extreme_mask]) * (zscore_threshold * std_val)
        methods_applied.append('zscore_clipping')

    info['methods_applied'] = methods_applied
    return cleaned_signal, info

def detect_rpeaks_robust(ecg_signal_clean, sfreq, return_log=False):
    """
    Detects R-peaks using WFDB xqrs_detect, refines them to local maxima within 0.05s window,
    and filters them to ensure a minimum distance of 501ms between consecutive peaks.
    """
    try:
        import wfdb.processing as wp
        # 1. WFDB Detection
        rpeaks_wfdb = wp.xqrs_detect(ecg_signal_clean, fs=sfreq, verbose=False)
        
        # 2. Refine to local maxima
        if len(rpeaks_wfdb) > 0:
            rpeaks = wp.correct_peaks(ecg_signal_clean, rpeaks_wfdb, 
                                      search_radius=int(0.05*sfreq),
                                      smooth_window_size=int(0.1*sfreq))
            rpeaks = np.unique(rpeaks)
            rpeaks = np.sort(rpeaks)

            rpeaks_sec = rpeaks / sfreq
            rr_intervals = np.diff(rpeaks_sec)
            unique_rr, counts = np.unique(rr_intervals, return_counts=True)
            
            mask = np.ones(len(rpeaks), dtype=bool)
            repetitive_info = []
            for rr, count in zip(unique_rr, counts):
                if count > len(rr_intervals) * 0.1:
                    perc = count / len(rr_intervals) * 100
                    msg = f"RR interval {rr:.3f}s appears {count} times ({perc:.1f}%)"
                    repetitive_info.append(msg)
                    bad_idx = np.where(rr_intervals == rr)[0] + 1
                    mask[bad_idx] = False
            
            removed_perc = np.sum(~mask) / len(mask) * 100
            log_msg = None
            if repetitive_info or removed_perc > 0:
                log_msg = {
                    'total': len(mask),
                    'removed': np.sum(~mask),
                    'perc': removed_perc,
                    'info': repetitive_info,
                    'skipped': removed_perc > REPETITIVE_RPEAK_EXCLUDE_PERC
                }
                
            rpeaks = rpeaks[mask]
            
            # 3. Filter peaks: if < 550ms between each other, keep first, remove second
            min_dist = int(0.4 * sfreq)
            if len(rpeaks) > 0:
                filtered_rpeaks = [rpeaks[0]]
                for i in range(1, len(rpeaks)):
                    if rpeaks[i] - filtered_rpeaks[-1] >= min_dist:
                        filtered_rpeaks.append(rpeaks[i])
                rpeaks = np.array(filtered_rpeaks)
                
        if return_log:
            return rpeaks, log_msg
        return rpeaks
    except Exception:
        if return_log:
            return np.array([], dtype=int), None
        return np.array([], dtype=int)


def detect_eeg_inversion(ch_data, sfreq, threshold=0.8):
    """
    Detects if an EEG channel is inverted by comparing local maxima and minima.
    If the absolute value is higher in the local minimum at more than threshold (80%),
    the channel is considered inverted.
    """
    # Bandpass filter 1-40Hz to remove DC offset and high frequency noise
    nyq = 0.5 * sfreq
    b, a = butter(2, [1.0/nyq, 40.0/nyq], btype='bandpass')
    filt_data = filtfilt(b, a, ch_data.astype(float))
    
    # Find all local maxima and minima with a small prominence to ignore noise
    prom = np.std(filt_data) * 0.1
    peaks, _ = find_peaks(filt_data, prominence=prom)
    troughs, _ = find_peaks(-filt_data, prominence=prom)
    
    if len(peaks) < 5 or len(troughs) < 5:
        return False
        
    is_inverted_count = 0
    total_comparisons = 0
    
    # For each trough, find the nearest peaks before and after and compare magnitude
    for t_idx in troughs:
        before = peaks[peaks < t_idx]
        after = peaks[peaks > t_idx]
        if len(before) == 0 or len(after) == 0:
            continue
        
        # Determine the "local" peak magnitude (max of the adjacent peaks)
        p_near = max(filt_data[before[-1]], filt_data[after[0]])
        t_val = np.abs(filt_data[t_idx])
        
        if t_val > p_near:
            is_inverted_count += 1
        total_comparisons += 1
        
    if total_comparisons == 0:
        return False
        
    return (is_inverted_count / total_comparisons) > threshold

def drop_non_eeg_channels(raw):
    """
    Remove non-EEG electrodes from the raw object using a single drop_channels call.
    """
    to_drop = [ch for ch in ['LOC', 'SpO2','eog'] if ch in raw.ch_names]
    prefixes = ('el', 'trx', 'png', 'beat','pul')
    
    for ch_name in raw.ch_names:
        if ch_name.lower().startswith(prefixes) and ch_name not in to_drop:
            to_drop.append(ch_name)
            
    if to_drop:
        raw.drop_channels(to_drop)
    return raw

def process_file_data(raw, patient_id):
    """
    Cleans ECG and extracts R-peaks using the provided logic.
    Returns the raw data, sfreq, and R-peak times if successful.
    """
    # Parameters
    minmax = (-0.3, .4)
    
    # Get sampling frequency and data
    sfreq = raw.info['sfreq']
    ch_names = raw.ch_names
    data = raw.get_data() * 1e6
    
    # Process ECG
    ch_lower = [ch.lower() for ch in ch_names]
    ecg_indices = [i for i, ch in enumerate(ch_lower) if 'ecg' in ch or 'ekg' in ch]
    
    if not ecg_indices:
        try:
            st.warning(f"No ECG channel found in {patient_id}")
        except Exception:
            print(f"WARNING: No ECG channel found in {patient_id}")
        return None
        
    ecg_ch_idx = ecg_indices[0]
    ecg_signal = data[ecg_ch_idx, :]
    
    # EEG inversion is now handled in get_group_individuals after HEP extraction
    # Fix inverted ECG
    ecg_signal, _ = fix_inverted_ecg(ecg_signal, sfreq)
    
    # Clean ECG
    ecg_signal_clean, _ = clean_ecg_high_fidelity(
        ecg_signal, 
        sampling_rate=sfreq
    )
    # Detect R-peaks using robust method
    rpeaks, log_msg = detect_rpeaks_robust(ecg_signal_clean, sfreq, return_log=True)

    if hasattr(raw, '_data'):
        raw._data[ecg_ch_idx, :] = ecg_signal_clean / 1e6

    if len(rpeaks) < 2:
        try:
            st.warning(
                f"Not enough R-peaks found for patient **{patient_id}** "
                f"(detected: {len(rpeaks)}). "
                f"Reason: {log_msg}"
            )
        except Exception:
            print(f"WARNING: Not enough R-peaks for {patient_id} (detected: {len(rpeaks)}). {log_msg}")
        return None

    rpeak_times = rpeaks / sfreq
    rpeak_ts = nap.Ts(t=rpeak_times, time_units="s")
    return raw, sfreq, rpeak_ts, rpeaks, minmax, log_msg


def _epoch_spectral_power_ratio_hf_lf(epochs, sfreq, low_band=(1.0, 30.0), high_band=(30.0, 80.0)):
    """Return high/low spectral power ratio for epochs shaped (n_epochs, n_channels, n_times)."""
    epochs = np.asarray(epochs, dtype=float)
    if epochs.ndim != 3 or epochs.size == 0:
        return np.empty((0, 0), dtype=float)
    n_times = epochs.shape[-1]
    if n_times < 4:
        return np.full(epochs.shape[:2], np.nan, dtype=float)
    nperseg = int(min(max(8, round(float(sfreq) * 0.5)), n_times))
    freqs, psd = signal.welch(epochs, fs=float(sfreq), axis=-1, nperseg=nperseg)
    low_mask = (freqs >= low_band[0]) & (freqs <= min(low_band[1], float(sfreq) / 2.0))
    high_mask = (freqs >= high_band[0]) & (freqs <= min(high_band[1], float(sfreq) / 2.0))
    ratios = np.full(epochs.shape[:2], np.nan, dtype=float)
    if not np.any(low_mask) or not np.any(high_mask):
        return ratios
    low_power = np.trapz(psd[..., low_mask], freqs[low_mask], axis=-1)
    high_power = np.trapz(psd[..., high_mask], freqs[high_mask], axis=-1)
    valid_low = np.isfinite(low_power) & (low_power > 0)
    ratios[valid_low] = high_power[valid_low] / (low_power[valid_low] + 1e-18)
    return ratios


def _valid_hep_window_mask(epochs, sfreq):
    """Return per-window keep mask after flat-line, artifact, and spectral checks."""
    epochs = np.asarray(epochs, dtype=float)
    if epochs.ndim != 3 or epochs.shape[0] == 0:
        return np.zeros(0, dtype=bool)

    ptp_uv = np.ptp(epochs, axis=2) * 1e6
    std_uv = np.nanstd(epochs, axis=2) * 1e6
    max_abs_uv = np.nanmax(np.abs(epochs), axis=2) * 1e6
    diff_std = np.nanstd(np.diff(epochs, axis=2), axis=2)
    trace_std = np.nanstd(epochs, axis=2)
    roughness = diff_std / (trace_std + 1e-12)

    bad_by_ch = (
        ~np.isfinite(ptp_uv)
        | ~np.isfinite(std_uv)
        | (ptp_uv < HEP_WINDOW_MIN_PTP_UV)
        | (std_uv < HEP_WINDOW_MIN_STD_UV)
        | (max_abs_uv > HEP_WINDOW_MAX_ABS_UV)
        | (np.isfinite(roughness) & (roughness > HEP_WINDOW_MAX_ROUGHNESS))
    )

    spectral_ratio = _epoch_spectral_power_ratio_hf_lf(epochs, sfreq)
    bad_by_ch |= ~np.isfinite(spectral_ratio) | (spectral_ratio >= MAX_SPECTRAL_POWER_RATIO)

    bad_fraction = np.mean(bad_by_ch, axis=1)
    return bad_fraction <= HEP_WINDOW_MAX_BAD_CHANNEL_FRACTION


def compute_hep_avg(raw, rpeaks, sfreq, minmax=(-0.5, 1.0), rpeak_ts=None):
    """
    Computes HEP (Heartbeat Evoked Potential) for all EEG channels.
    Only complete, non-flat, not-too-noisy heartbeat windows are averaged.
    """
    tmin, tmax = minmax
    
    # Identify EEG channels
    ch_names = raw.ch_names
    eeg_indices = [i for i, ch in enumerate(ch_names) 
                  if 'eeg' in ch.lower() or any(elec in ch.upper() for elec in ['FP', 'F', 'C', 'P', 'O', 'T', 'A'])]
    
    if not eeg_indices:
        return None, None, None
        
    # Pick EEG channels
    eeg_ch_names = [ch_names[i] for i in eeg_indices]

    pre = int(round(abs(tmin) * sfreq))
    post = int(round(tmax * sfreq))
    offsets = np.arange(-pre, post + 1)
    valid_rpeaks = np.asarray(rpeaks, dtype=int)
    valid_rpeaks = valid_rpeaks[
        (valid_rpeaks - pre >= 0) & (valid_rpeaks + post < raw.n_times)
    ]
    if len(valid_rpeaks) == 0:
        return None, offsets / sfreq, eeg_ch_names

    data = raw.get_data(picks=eeg_indices)
    epochs = np.asarray([data[:, peak + offsets] for peak in valid_rpeaks], dtype=float)
    keep_mask = _valid_hep_window_mask(epochs, sfreq)
    if int(np.sum(keep_mask)) == 0:
        try:
            st.warning(
                "No valid HEP windows remained after flat/noise and "
                f"spectral_power_ratio_hf_lf < {MAX_SPECTRAL_POWER_RATIO} filtering."
            )
        except Exception:
            pass
        return None, offsets / sfreq, eeg_ch_names

    mean_data = np.nanmean(epochs[keep_mask], axis=0)
    return mean_data, offsets / sfreq, eeg_ch_names



def compute_ecg_hep_avg(raw, rpeaks, sfreq, minmax=(-0.5, 1.0), rpeak_ts=None):
    """
    Computes HEP (Heartbeat Evoked Potential) for the ECG channel.
    """
    tmin, tmax = minmax
    ch_names = raw.ch_names
    ch_lower = [ch.lower() for ch in ch_names]
    ecg_indices = [i for i, ch in enumerate(ch_lower) if 'ecg' in ch or 'ekg' in ch]
    
    if not ecg_indices:
        return None, None, None
        
    ecg_ch_name = ch_names[ecg_indices[0]]

    # Pynapple implementation
    data = raw.get_data(picks=[ecg_indices[0]]).T # (n_times, 1)
    tsd_frame = nap.TsdFrame(t=raw.times, d=data, columns=[ecg_ch_name])
    perievent = nap.compute_perievent_continuous(tsd_frame, rpeak_ts, minmax=minmax)
    # Average across trials (axis 1)
    mean_data = perievent.nanmean(axis=1).values.T # (1, n_times)
    return mean_data, perievent.t, [ecg_ch_name]


def _is_ecg_channel(ch_name):
    ch = ch_name.lower()
    return 'ecg' in ch or 'ekg' in ch


def _is_eeg_channel(ch_name):
    ch = ch_name.strip()
    ch_upper = ch.upper()
    ch_lower = ch.lower()
    return (
        ch_lower.startswith('eeg')
        or re.match(r'^[A-Za-z]{1,3}[0-9]+$', ch) is not None
        or re.match(r'^[A-Za-z]{1,2}z$', ch, re.IGNORECASE) is not None
    )


def _is_non_eeg_feature_channel(ch_name):
    ch = ch_name.strip()
    ch_lower = ch.lower()
    excluded_exact = {'status', 'mkr+', 'marker', 'event', 'annotations'}
    excluded_prefixes = ('stim', 'trigger', 'trig', 'dc', 'sync')
    if ch_lower in excluded_exact:
        return False
    if ch_lower.startswith(excluded_prefixes):
        return False
    if _is_ecg_channel(ch) or _is_eeg_channel(ch):
        return False
    return True


def compute_non_eeg_aligned_avg(raw, minmax=(-0.5, 1.0), rpeak_ts=None):
    """
    Computes ECG-aligned averages for all non-EEG / non-ECG channels.
    """
    if rpeak_ts is None:
        return None, None, None

    candidate_indices = [
        idx for idx, ch in enumerate(raw.ch_names)
        if _is_non_eeg_feature_channel(ch)
    ]
    if not candidate_indices:
        return None, None, None

    ch_names = [raw.ch_names[i] for i in candidate_indices]
    data = raw.get_data(picks=candidate_indices).T
    tsd_frame = nap.TsdFrame(t=raw.times, d=data, columns=ch_names)
    perievent = nap.compute_perievent_continuous(tsd_frame, rpeak_ts, minmax=minmax)
    mean_data = perievent.nanmean(axis=1).values.T
    return mean_data, perievent.t, ch_names


def find_ecg_t_peak_time(ecg_hep_data, times, t_window=(0.15, 0.50)):
    """
    Find the ECG T-wave peak time in the averaged ECG HEP epoch.

    The ECG signal is already polarity-corrected earlier in the pipeline, so
    the T wave should be positive. This intentionally uses the maximum inside
    the expected T-wave window rather than the largest absolute value near the
    R peak.
    """
    if ecg_hep_data is None or times is None:
        return None

    times = np.asarray(times, dtype=float)
    ecg_trace = np.asarray(ecg_hep_data, dtype=float).squeeze()
    if ecg_trace.ndim != 1 or len(ecg_trace) != len(times):
        return None

    mask = (times >= t_window[0]) & (times <= t_window[1])
    if not np.any(mask):
        return None

    window_values = ecg_trace[mask]
    finite_mask = np.isfinite(window_values)
    if not np.any(finite_mask):
        return None

    window_times = times[mask][finite_mask]
    window_values = window_values[finite_mask]
    return float(window_times[int(np.nanargmax(window_values))])


def score_negative_dip_multimethod(
    eeg_trace,
    times,
    center_time,
    pre_window=0.10,
    post_window=0.0,
    baseline_pre_window=0.50,
    swing_threshold=50.0,
    z_threshold=3.0,
    prominence_threshold=25.0,
    min_votes=2,
):
    """
    Score whether a negative deflection around ``center_time`` is large enough
    to be considered an inverted EEG response.

    Three complementary criteria are used:
    1. Local peak-to-trough swing percentage.
    2. Z-score depth relative to the preceding baseline window.
    3. Prominence as a percentage of the trace's global range.

    A flip is recommended when the dip is negative and at least ``min_votes`` of
    the three criteria pass. This distinguishes a true deep downward deflection
    from a small wiggle around zero.
    """
    if eeg_trace is None or times is None or center_time is None:
        return False, {}

    times = np.asarray(times, dtype=float)
    eeg_trace = np.asarray(eeg_trace, dtype=float).squeeze()
    if eeg_trace.ndim != 1 or len(eeg_trace) != len(times):
        return False, {}

    target_mask = (times >= center_time - pre_window) & (times <= center_time + post_window)
    if not np.any(target_mask):
        return False, {}

    target_times = times[target_mask]
    target_values = eeg_trace[target_mask]
    finite_mask = np.isfinite(target_values)
    if not np.any(finite_mask):
        return False, {}

    target_times = target_times[finite_mask]
    target_values = target_values[finite_mask]
    if len(target_values) < 3:
        return False, {}

    dip_idx = int(np.nanargmin(target_values))
    peak_idx = int(np.nanargmax(target_values))
    dip_amp = float(target_values[dip_idx])
    peak_amp = float(target_values[peak_idx])
    dip_time = float(target_times[dip_idx])

    preceding_values = target_values[:dip_idx + 1]
    if len(preceding_values):
        shoulder_peak = float(np.nanmax(preceding_values))
        shoulder_peak_idx = int(np.nanargmax(preceding_values))
        shoulder_peak_time = float(target_times[shoulder_peak_idx])
    else:
        shoulder_peak = peak_amp
        shoulder_peak_time = float(target_times[peak_idx])

    if shoulder_peak > 0:
        swing_pct = ((shoulder_peak - dip_amp) / (abs(shoulder_peak) + 1e-12)) * 100.0
    else:
        swing_pct = 0.0
    swing_pass = bool(dip_amp < 0 and shoulder_peak > 0 and swing_pct >= swing_threshold)

    baseline_mask = (times >= center_time - baseline_pre_window) & (times < center_time - pre_window)
    baseline_values = eeg_trace[baseline_mask]
    baseline_values = baseline_values[np.isfinite(baseline_values)]
    if len(baseline_values) >= 5:
        baseline_mean = float(np.nanmean(baseline_values))
        baseline_std = float(np.nanstd(baseline_values, ddof=1))
    else:
        baseline_mean = float(np.nanmean(target_values))
        baseline_std = float(np.nanstd(target_values, ddof=1)) if len(target_values) > 1 else 0.0
    z_depth = (baseline_mean - dip_amp) / (baseline_std + 1e-12)
    z_pass = bool(dip_amp < baseline_mean - z_threshold * (baseline_std + 1e-12))

    global_range = float(np.nanmax(eeg_trace) - np.nanmin(eeg_trace))
    if np.isfinite(global_range) and global_range > 0:
        try:
            troughs, props = find_peaks(-target_values, prominence=0)
            if len(troughs):
                nearest = int(np.argmin(np.abs(troughs - dip_idx)))
                prominence = float(props["prominences"][nearest])
            else:
                left_shoulder = float(np.nanmax(target_values[:dip_idx + 1])) if dip_idx > 0 else dip_amp
                right_shoulder = float(np.nanmax(target_values[dip_idx:])) if dip_idx < len(target_values) - 1 else dip_amp
                prominence = float(min(left_shoulder, right_shoulder) - dip_amp)
            prominence_pct = max(0.0, prominence / (global_range + 1e-12) * 100.0)
        except Exception:
            prominence = np.nan
            prominence_pct = 0.0
    else:
        prominence = np.nan
        prominence_pct = 0.0
    prominence_pass = bool(dip_amp < 0 and prominence_pct >= prominence_threshold)

    votes = int(swing_pass) + int(z_pass) + int(prominence_pass)
    invert = bool(dip_amp < 0 and votes >= min_votes)

    return invert, {
        "method": "negative_dip_multiscore",
        "center_time": float(center_time),
        "dip_time": dip_time,
        "dip_amp": dip_amp,
        "window_max_amp": peak_amp,
        "shoulder_peak_amp": shoulder_peak,
        "shoulder_peak_time": shoulder_peak_time,
        "swing_pct": float(swing_pct),
        "swing_threshold": float(swing_threshold),
        "swing_pass": bool(swing_pass),
        "baseline_mean": baseline_mean,
        "baseline_std": baseline_std,
        "z_depth": float(z_depth),
        "z_threshold": float(z_threshold),
        "z_pass": bool(z_pass),
        "prominence": float(prominence) if np.isfinite(prominence) else np.nan,
        "prominence_pct": float(prominence_pct),
        "prominence_threshold": float(prominence_threshold),
        "prominence_pass": bool(prominence_pass),
        "vote_count": int(votes),
        "min_votes": int(min_votes),
        "flip_reason": "multi_score_deep_dip" if invert else "multi_score_small_dip",
        "pre_window": float(pre_window),
        "post_window": float(post_window),
        "baseline_pre_window": float(baseline_pre_window),
    }


def should_invert_eeg_from_t_wave(eeg_trace, times, ecg_t_peak_time, pre_window=0.10, post_window=0.0):
    """
    Decide whether an EEG HEP trace is inverted around the ECG T-wave peak.

    The check is constrained to the ``pre_window`` immediately before the ECG
    T peak. The EEG is considered inverted only when the EEG T-wave candidate
    is negative, is the local minimum of that pre-T window, and is not also the
    local maximum. This avoids the old behavior of flipping from an unrelated
    absolute deflection near the R peak.
    """
    invert, info = score_negative_dip_multimethod(
        eeg_trace,
        times,
        center_time=ecg_t_peak_time,
        pre_window=pre_window,
        post_window=post_window,
        baseline_pre_window=0.50,
        swing_threshold=50.0,
        z_threshold=3.0,
        prominence_threshold=25.0,
        min_votes=2,
    )

    # Preserve legacy keys used elsewhere in the dashboard.
    if info:
        info.update({
            "ecg_t_peak_time": float(ecg_t_peak_time),
            "eeg_t_peak_time": float(info.get("dip_time", np.nan)),
            "eeg_t_peak_amp": float(info.get("dip_amp", np.nan)),
        })
    return invert, info


def _smooth_rpeak_window(values, times, smooth_ms=25.0):
    """Smooth the R-peak window enough to ignore single-sample derivative zeros."""
    values = np.asarray(values, dtype=float)
    times = np.asarray(times, dtype=float)
    if len(values) < 5:
        return values.copy()

    dt = float(np.nanmedian(np.diff(times))) if len(times) > 1 else 0.0
    if not np.isfinite(dt) or dt <= 0:
        return median_filter(values, size=min(3, len(values)))

    win = int(round((smooth_ms / 1000.0) / dt))
    win = max(5, win)
    if win % 2 == 0:
        win += 1
    if win >= len(values):
        win = len(values) - 1 if len(values) % 2 == 0 else len(values)
    if win < 5:
        return median_filter(values, size=min(3, len(values)))

    try:
        return signal.savgol_filter(values, window_length=win, polyorder=2, mode="interp")
    except Exception:
        return median_filter(values, size=min(3, len(values)))


def _linear_slope(x, y):
    """Least-squares slope for a short trend segment."""
    x = np.asarray(x, dtype=float)
    y = np.asarray(y, dtype=float)
    finite = np.isfinite(x) & np.isfinite(y)
    if np.sum(finite) < 2:
        return np.nan
    x = x[finite]
    y = y[finite]
    x = x - np.nanmean(x)
    denom = float(np.dot(x, x))
    if denom <= 1e-18:
        return np.nan
    return float(np.dot(x, y - np.nanmean(y)) / denom)


def _evaluate_rpeak_extremum_shape(times, values, idx, kind, flank_ms=35.0):
    """
    Score whether idx is a clear valley or peak with a trend on both sides.
    A real valley has negative left slope and positive right slope; a real peak
    has positive left slope and negative right slope.
    """
    n = len(values)
    if idx <= 1 or idx >= n - 2:
        return {
            "is_clear": False,
            "reason": "edge_candidate",
            "prominence": 0.0,
            "left_slope": np.nan,
            "right_slope": np.nan,
        }

    center_t = float(times[idx])
    flank_sec = flank_ms / 1000.0
    left_mask = (times >= center_t - flank_sec) & (times < center_t)
    right_mask = (times > center_t) & (times <= center_t + flank_sec)

    if np.sum(left_mask) < 2:
        left_mask = np.zeros(n, dtype=bool)
        left_mask[max(0, idx - 3):idx] = True
    if np.sum(right_mask) < 2:
        right_mask = np.zeros(n, dtype=bool)
        right_mask[idx + 1:min(n, idx + 4)] = True

    if np.sum(left_mask) < 2 or np.sum(right_mask) < 2:
        return {
            "is_clear": False,
            "reason": "insufficient_flank_points",
            "prominence": 0.0,
            "left_slope": np.nan,
            "right_slope": np.nan,
        }

    left_slope = _linear_slope(times[left_mask], values[left_mask])
    right_slope = _linear_slope(times[right_mask], values[right_mask])
    center_amp = float(values[idx])
    left_level = float(np.nanmedian(values[left_mask]))
    right_level = float(np.nanmedian(values[right_mask]))
    local_range = float(np.nanmax(values) - np.nanmin(values))
    noise = float(np.nanmedian(np.abs(np.diff(values) - np.nanmedian(np.diff(values))))) if len(values) > 3 else 0.0
    prominence_threshold = max(2.5 * noise, 0.12 * local_range, 1e-12)
    slope_threshold = prominence_threshold / max(flank_sec, 1e-6) * 0.20

    if kind == "valley":
        prominence = min(left_level - center_amp, right_level - center_amp)
        has_shape = (
            np.isfinite(left_slope) and np.isfinite(right_slope)
            and left_slope < -slope_threshold
            and right_slope > slope_threshold
        )
    else:
        prominence = min(center_amp - left_level, center_amp - right_level)
        has_shape = (
            np.isfinite(left_slope) and np.isfinite(right_slope)
            and left_slope > slope_threshold
            and right_slope < -slope_threshold
        )

    is_clear = bool(has_shape and prominence >= prominence_threshold)
    reason = "clear_shape" if is_clear else "weak_or_noisy_shape"
    return {
        "is_clear": is_clear,
        "reason": reason,
        "prominence": float(max(prominence, 0.0)),
        "prominence_threshold": float(prominence_threshold),
        "slope_threshold": float(slope_threshold),
        "left_slope": float(left_slope) if np.isfinite(left_slope) else np.nan,
        "right_slope": float(right_slope) if np.isfinite(right_slope) else np.nan,
        "left_level": left_level,
        "right_level": right_level,
        "center_amp": center_amp,
        "center_time": center_t,
    }


def should_flip_eeg_around_r_peak(eeg_trace, times, pre_window=0.10, post_window=0.10):
    """
    Decide whether an EEG HEP channel is inverted around the R-peak.

    The decision uses only the ±100 ms R-peak window. A channel is flipped only
    when the smoothed global minimum forms a clear valley: values trend downward
    into the minimum and upward away from it. A clear peak/parabola shape means
    the polarity is kept. This avoids flipping on single noisy samples where the
    derivative is briefly zero.
    """
    if eeg_trace is None or times is None:
        return False, {}

    times = np.asarray(times, dtype=float)
    eeg_trace = np.asarray(eeg_trace, dtype=float).squeeze()
    if eeg_trace.ndim != 1 or len(eeg_trace) != len(times):
        return False, {}

    mask = (times >= -pre_window) & (times <= post_window)
    if not np.any(mask):
        return False, {}

    window_times = times[mask]
    window_values = eeg_trace[mask]
    finite_mask = np.isfinite(window_values)
    if not np.any(finite_mask):
        return False, {}

    window_times = window_times[finite_mask]
    window_values = window_values[finite_mask]
    if len(window_values) < 5:
        return False, {}

    smooth_values = _smooth_rpeak_window(window_values, window_times)
    min_idx = int(np.nanargmin(smooth_values))
    max_idx = int(np.nanargmax(smooth_values))
    min_amp = float(smooth_values[min_idx])
    max_amp = float(smooth_values[max_idx])

    valley = _evaluate_rpeak_extremum_shape(window_times, smooth_values, min_idx, "valley")
    peak = _evaluate_rpeak_extremum_shape(window_times, smooth_values, max_idx, "peak")

    if valley["is_clear"] and not peak["is_clear"]:
        flip = True
        flip_reason = "clear_r_peak_valley"
    elif peak["is_clear"] and not valley["is_clear"]:
        flip = False
        flip_reason = "clear_r_peak_peak"
    elif valley["is_clear"] and peak["is_clear"]:
        flip = bool(valley["prominence"] > peak["prominence"])
        flip_reason = "valley_stronger_than_peak" if flip else "peak_stronger_than_valley"
    else:
        flip = False
        flip_reason = "no_clear_r_peak_extremum"

    t0_idx = int(np.argmin(np.abs(window_times)))
    t0_amp = float(smooth_values[t0_idx])

    return flip, {
        "method": "r_peak_smoothed_extremum_shape",
        "t0_amp": t0_amp,
        "t0_time": float(window_times[t0_idx]),
        "window_min_amp": min_amp,
        "window_min_time": float(window_times[min_idx]),
        "window_max_amp": max_amp,
        "window_max_time": float(window_times[max_idx]),
        "valley": valley,
        "peak": peak,
        "flip_reason": flip_reason,
        "pre_window": float(pre_window),
        "post_window": float(post_window),
    }


def should_flip_eeg_swing_percentage(eeg_trace, times, pre_window=0.10, post_window=0.10, threshold=50.0):
    """
    Detect an inverted R-peak dip in the EEG using a multi-score depth rule.

    The first score remains the Local Peak-to-Trough Ratio:
    Drop% = (V_peak - V_dip) / V_peak * 100.

    It is combined with:
    - Z-score depth relative to the preceding 500 ms.
    - Prominence/global-range percentage.

    Returns True only when the dip is negative and at least two of the three
    depth scores pass, which makes the flip decision more robust to small
    wiggles around x=0.
    """
    flip, info = score_negative_dip_multimethod(
        eeg_trace,
        times,
        center_time=0.0,
        pre_window=pre_window,
        post_window=post_window,
        baseline_pre_window=0.50,
        swing_threshold=threshold,
        z_threshold=3.0,
        prominence_threshold=25.0,
        min_votes=2,
    )
    if not info:
        return False, {}

    # Preserve legacy field names expected by logging code.
    info.update({
        "method": "r_peak_multiscore",
        "v_peak": float(info.get("shoulder_peak_amp", np.nan)),
        "v_dip": float(info.get("dip_amp", np.nan)),
        "drop_pct": float(info.get("swing_pct", np.nan)),
        "threshold": float(threshold),
        "flip_reason": "swing_pct_dip" if flip else info.get("flip_reason", "multi_score_small_dip"),
    })
    return flip, info


def flip_eeg_channels_around_r_peak(raw, hep_data, times, ch_names, eeg_indices):
    """
    Flip EEG channels whose R-peak window contains a clear negative valley.
    """
    valid_eeg_indices = [
        idx for idx in eeg_indices
        if idx < len(hep_data) and len(hep_data[idx]) > 0
    ]
    if not valid_eeg_indices:
        return [], {}

    flipped_channels = []
    per_channel_info = {}
    raw_ch_names = raw.ch_names
    for idx in valid_eeg_indices:
        ch_name = ch_names[idx]
        should_flip, flip_info = should_flip_eeg_around_r_peak(
            hep_data[idx], times, pre_window=0.10, post_window=0.10
        )
        per_channel_info[ch_name] = flip_info
        if should_flip and ch_name in raw_ch_names:
            raw_idx = raw_ch_names.index(ch_name)
            raw._data[raw_idx, :] = -raw._data[raw_idx, :]
            flipped_channels.append(ch_name)

    return flipped_channels, {"per_channel": per_channel_info}


def flip_eeg_channels_if_t_wave_inverted(raw, hep_data, times, ch_names, eeg_indices, ecg_hep_data):
    """
    Flip individual raw EEG channels when their EEG T wave is inverted.

    Each standard EEG channel is checked separately. If that channel has a
    negative local minimum in the 100 ms before the ECG T-wave peak, only that
    channel is flipped in ``raw`` and the caller can recompute HEP from the
    corrected raw data.
    """
    valid_eeg_indices = [
        idx for idx in eeg_indices
        if idx < len(hep_data) and len(hep_data[idx]) > 0
    ]
    if not valid_eeg_indices:
        return [], {}

    ecg_t_peak_time = find_ecg_t_peak_time(ecg_hep_data, times)

    flipped_channels = []
    per_channel_info = {}
    raw_ch_names = raw.ch_names
    for idx in valid_eeg_indices:
        ch_name = ch_names[idx]
        should_flip, flip_info = should_invert_eeg_from_t_wave(
            hep_data[idx], times, ecg_t_peak_time, pre_window=0.10, post_window=0.0
        )
        per_channel_info[ch_name] = flip_info
        if should_flip and ch_name in raw_ch_names:
            raw_idx = raw_ch_names.index(ch_name)
            raw._data[raw_idx, :] = -raw._data[raw_idx, :]
            flipped_channels.append(ch_name)

    return flipped_channels, {
        "ecg_t_peak_time": ecg_t_peak_time,
        "per_channel": per_channel_info,
    }


def flip_eeg_channels_swing_pct_then_t_wave(raw, hep_data, times, ch_names, eeg_indices, ecg_hep_data):
    """
    Flip EEG channels for non-Berkeley data.

    Each channel is first tested with the Swing Percentage formula in the
    ±100 ms R-peak window.  If the dip fails the ≥50 % threshold, the channel
    falls back to the EEG T-wave inversion check.
    """
    valid_eeg_indices = [
        idx for idx in eeg_indices
        if idx < len(hep_data) and len(hep_data[idx]) > 0
    ]
    if not valid_eeg_indices:
        return [], {}

    ecg_t_peak_time = find_ecg_t_peak_time(ecg_hep_data, times)

    flipped_channels = []
    per_channel_info = {}
    raw_ch_names = raw.ch_names

    for idx in valid_eeg_indices:
        ch_name = ch_names[idx]

        flip, flip_info = should_flip_eeg_swing_percentage(
            hep_data[idx], times, pre_window=0.10, post_window=0.10, threshold=50.0
        )

        if not flip:
            flip_tw, flip_info_tw = should_invert_eeg_from_t_wave(
                hep_data[idx], times, ecg_t_peak_time, pre_window=0.10, post_window=0.0
            )
            flip_info = {**flip_info, "t_wave_fallback": flip_info_tw}
            flip = flip_tw

        per_channel_info[ch_name] = {**flip_info, "used_t_wave_fallback": not flip_info.get("flip_reason") == "swing_pct_dip" and flip}
        if flip and ch_name in raw_ch_names:
            raw_idx = raw_ch_names.index(ch_name)
            raw._data[raw_idx, :] = -raw._data[raw_idx, :]
            flipped_channels.append(ch_name)

    return flipped_channels, {
        "ecg_t_peak_time": ecg_t_peak_time,
        "per_channel": per_channel_info,
    }


def process_and_invert_hep(_raw, rpeaks, sfreq, minmax, _rpeak_ts, patient_id, group_name=''):
    """
    Computes HEP and ECG HEP, then fixes EEG polarity when needed.

    Per-channel flip is based only on the R-peak window: invert when the
    smoothed ±100 ms window has a clear negative valley, and keep polarity
    when it has a clear positive peak/parabola.
    """
    hep_data, times, ch_names = compute_hep_avg(_raw, rpeaks, sfreq, minmax, rpeak_ts=_rpeak_ts)
    ecg_hep_data, _, ecg_ch_names = compute_ecg_hep_avg(_raw, rpeaks, sfreq, minmax, rpeak_ts=_rpeak_ts)

    # Extract raw ECG signal (1D) for quality checks
    ch_lower = [ch.lower() for ch in _raw.ch_names]
    ecg_indices = [i for i, ch in enumerate(ch_lower) if 'ecg' in ch or 'ekg' in ch]
    ecg_data = _raw.get_data(picks=[ecg_indices[0]]).squeeze() if ecg_indices else None

    if hep_data is None:
        return hep_data, times, ch_names, ecg_hep_data, ecg_ch_names, ecg_data, []

    # Match standard 10-20 EEG channel names: letter(s)+digit(s) OR letter+z (midline like Fz, Cz, Pz, Oz)
    eeg_indices = [i for i, ch in enumerate(ch_names)
                  if re.match(r'^[A-Za-z]{1,3}[0-9]+$', ch) or re.match(r'^[A-Za-z]{1,2}z$', ch, re.IGNORECASE)]

    flipped_channels, flip_info = flip_eeg_channels_around_r_peak(
        _raw, hep_data, times, ch_names, eeg_indices
    )
    if flipped_channels:
        flipped_details = []
        for ch in flipped_channels:
            ch_info = flip_info.get("per_channel", {}).get(ch, {})
            min_amp_uv = ch_info.get("window_min_amp", np.nan) * 1e6
            max_amp_uv = ch_info.get("window_max_amp", np.nan) * 1e6
            flipped_details.append((ch, min_amp_uv, max_amp_uv))
        detail_text = ", ".join(
            f"{ch} (min {mn:.2f} uV, max {mx:.2f} uV)"
            for ch, mn, mx in flipped_details[:8]
        )
        if len(flipped_details) > 8:
            detail_text += f", +{len(flipped_details) - 8} more"
        flip_msg = (
            f"Flipped {len(flipped_channels)} inverted EEG channel(s) for {patient_id} "
            f"(R-peak ±100 ms clear valley shape): {detail_text}"
        )
        try:
            st.info(flip_msg)
        except Exception:
            print(flip_msg)

    if flipped_channels:
        hep_data, times, ch_names = compute_hep_avg(_raw, rpeaks, sfreq, minmax, rpeak_ts=_rpeak_ts)
        ecg_hep_data, _, ecg_ch_names = compute_ecg_hep_avg(_raw, rpeaks, sfreq, minmax, rpeak_ts=_rpeak_ts)

    return hep_data, times, ch_names, ecg_hep_data, ecg_ch_names, ecg_data, flipped_channels


def _apply_ica_ecg_removal(raw, patient_id):
    """
    Applies MNE ICA to remove ECG artifact components from EEG channels.
    Finds components correlated with the ECG channel and excludes them.
    Returns the cleaned raw object (in-place modification).
    """
    try:
        ch_names = raw.ch_names
        ch_lower = [ch.lower() for ch in ch_names]
        ecg_indices = [i for i, ch in enumerate(ch_lower) if 'ecg' in ch or 'ekg' in ch]
        if not ecg_indices:
            print(f"[ICA] No ECG channel found for {patient_id}, skipping ICA.")
            return raw

        ecg_ch_name = ch_names[ecg_indices[0]]
        eeg_picks = mne.pick_types(raw.info, eeg=True, ecg=False, stim=False, exclude='bads')
        if len(eeg_picks) < 2:
            print(f"[ICA] Not enough EEG channels for ICA in {patient_id}, skipping.")
            return raw

        n_components = min(15, len(eeg_picks) - 1)
        # ICLabel requirements: extended infomax, 1–100 Hz bandpass, CAR.
        # Prefer picard (extended infomax equivalent); fall back to infomax if unavailable.
        try:
            import picard  # noqa: F401
            _ica_method = 'picard'
            _ica_fit_params = dict(ortho=False, extended=True)
        except ImportError:
            _ica_method = 'infomax'
            _ica_fit_params = dict(extended=True)
        ica = mne.preprocessing.ICA(
            n_components=n_components,
            method=_ica_method,
            fit_params=_ica_fit_params,
            random_state=42,
            max_iter=500,
        )
        # Fit on 1–100 Hz bandpass + CAR — satisfies all three ICLabel preconditions.
        raw_filt = raw.copy().filter(1.0, 100.0, picks=eeg_picks, verbose=False)
        # Add this before calling label_components
        raw_filt.set_montage('standard_1020') # Or your specific montage
        raw_filt.set_eeg_reference('average', projection=False, verbose=False)
        ica.fit(raw_filt, picks=eeg_picks, verbose=False)

        # --- ECG-based detection ---
        bad_ecg_method = 'correlation'
        ecg_inds, ecg_scores = ica.find_bads_ecg(raw, ch_name=ecg_ch_name, verbose=False, method=bad_ecg_method)
        print(f"[ICA] {patient_id}: ECG-based components: {ecg_inds}")

        # --- ICLabel-based detection ---
        iclabel_inds = []
        try:
            # raw_filt already satisfies ICLabel: 1–100 Hz + CAR + extended infomax.
            ic_labels = label_components(raw_filt, ica, method='iclabel')
            labels = ic_labels['labels']
            probs  = ic_labels['y_pred_proba']
            exclude_labels = {'muscle', 'eye', 'heart', 'line_noise', 'channel_noise'}
            for idx, (lbl, prob) in enumerate(zip(labels, probs)):
                if lbl in exclude_labels and max(prob) > 0.5:
                    iclabel_inds.append(idx)
            print(f"[ICA] {patient_id}: ICLabel-based components: {iclabel_inds} "
                  f"(labels: {[labels[i] for i in iclabel_inds]})")
        except Exception as e_iclabel:
            print(f"[ICA] {patient_id}: ICLabel failed, using ECG-only: {e_iclabel}")

        # --- Combine and apply ---
        # Fallback: if bad_ecg_method found no components, always remove the top ECG-scoring one
        if not ecg_inds and ecg_scores is not None and len(ecg_scores) > 0:
            top_ecg = int(np.argmax(np.abs(ecg_scores)))
            print(f"[ICA] {patient_id}: {bad_ecg_method} found no components, using top-scoring fallback: {top_ecg} (score={ecg_scores[top_ecg]:.4f})")
            ecg_inds = [top_ecg]
        exclude_inds = sorted(set(ecg_inds) | set(iclabel_inds))
        if exclude_inds:
            ica.exclude = exclude_inds
            ica.apply(raw, verbose=False)
            print(f"[ICA] {patient_id}: removed {len(exclude_inds)} component(s) total: {exclude_inds}")
        else:
            print(f"[ICA] {patient_id}: no artifact components identified.")
    except Exception as e:
        print(f"[ICA] Error during ICA for {patient_id}: {e}")
    return raw


def _process_patient_worker(args):
    """Module-level worker for ProcessPoolExecutor. Returns result tuple or None on failure."""
    if len(args) == 4:
        f_path, patient_id, apply_ica, group_name = args
    elif len(args) == 3:
        f_path, patient_id, apply_ica = args
        group_name = ''
    else:
        f_path, patient_id = args
        apply_ica = False
        group_name = ''
    try:
        with open(f_path, 'rb') as f:
            raw = pickle.load(f)
        raw = drop_non_eeg_channels(raw)
        results = process_file_data(raw, patient_id)
        if results is None:
            return None
        raw, sfreq, rpeak_ts, rpeaks, minmax, log_msg = results
        if apply_ica:
            raw = _apply_ica_ecg_removal(raw, patient_id)
        hep_data, times, ch_names, ecg_hep_data, ecg_ch_names, ecg_data, flipped_channels = process_and_invert_hep(
            raw, rpeaks, sfreq, minmax, rpeak_ts, patient_id, group_name=group_name
        )

        return (patient_id, hep_data, times, ch_names, rpeaks, ecg_hep_data, ecg_ch_names, log_msg, flipped_channels)
    except Exception as e:
        print(f"[Worker] Error processing {patient_id}: {e}")
        return None


def _process_non_eeg_patient_worker(args):
    """Worker that computes ECG-aligned averages for non-EEG channels."""
    if len(args) == 3:
        f_path, patient_id, apply_ica = args
    else:
        f_path, patient_id = args
        apply_ica = False

    try:
        with open(f_path, 'rb') as f:
            raw = pickle.load(f)

        results = process_file_data(raw, patient_id)
        if results is None:
            return None

        raw, sfreq, rpeak_ts, rpeaks, minmax, log_msg = results
        if apply_ica:
            raw = _apply_ica_ecg_removal(raw, patient_id)

        non_eeg_data, times, ch_names = compute_non_eeg_aligned_avg(
            raw, minmax=minmax, rpeak_ts=rpeak_ts
        )
        if non_eeg_data is None or times is None or ch_names is None:
            return None

        return (patient_id, non_eeg_data, times, ch_names, rpeaks, log_msg)
    except Exception as e:
        print(f"[Worker] Error processing non-EEG channels for {patient_id}: {e}")
        traceback.print_exc()
        return None


def _individuals_have_flip_metadata(individuals):
    """Return True when cached individual tuples include flipped-channel metadata."""
    if not individuals:
        return False
    return all(len(ind) > 8 for ind in individuals)


def _patient_id_from_pickle_name(filename):
    return filename.replace('.pkl', '').replace('.edf', '')


def _cache_filename_for_individuals(test_run=False, apply_ica=False, min_eeg_channels=None):
    ica_suffix = '_ica' if apply_ica else ''
    channel_suffix = (
        f'_min{int(min_eeg_channels)}_1020eeg'
        if min_eeg_channels else ''
    )
    cache_filename = (
        f'individuals_cache_test{ica_suffix}{channel_suffix}.pkl'
        if test_run else f'individuals_cache{ica_suffix}{channel_suffix}.pkl'
    )
    return cache_filename.replace('.pkl', f'{HEP_FLIP_CACHE_SUFFIX}.pkl')


def _individuals_cache_meta_path(cache_path):
    return cache_path.replace('.pkl', '_meta.pkl')


def _standard_1020_channel_count(channel_names):
    """Count distinct channels from the dashboard's canonical 10-20 list."""
    return len({
        str(channel).strip().lower()
        for channel in (channel_names or [])
        if str(channel).strip().lower() in STANDARD_1020_EEG_CHANNEL_KEYS
    })


def _dump_pickle_atomic_with_heartbeat(payload, destination):
    """Atomically write a cache while reporting disk-I/O progress every 10 seconds."""
    temp_path = f"{destination}.tmp.{os.getpid()}"
    stop_event = threading.Event()
    started = time.monotonic()

    def report_progress():
        while not stop_event.wait(10):
            try:
                size_mb = os.path.getsize(temp_path) / (1024 * 1024)
            except OSError:
                size_mb = 0.0
            elapsed = int(time.monotonic() - started)
            print(
                f"[cache write] {os.path.basename(destination)}: "
                f"{size_mb:.1f} MiB written, {elapsed}s elapsed",
                flush=True,
            )

    reporter = threading.Thread(target=report_progress, daemon=True)
    reporter.start()
    try:
        with open(temp_path, 'wb') as f:
            pickle.dump(payload, f)
            f.flush()
            os.fsync(f.fileno())
        os.replace(temp_path, destination)
        elapsed = int(time.monotonic() - started)
        size_mb = os.path.getsize(destination) / (1024 * 1024)
        print(
            f"[cache write] {os.path.basename(destination)} complete: "
            f"{size_mb:.1f} MiB in {elapsed}s",
            flush=True,
        )
    finally:
        stop_event.set()
        reporter.join(timeout=1)
        try:
            os.remove(temp_path)
        except FileNotFoundError:
            pass


def _load_fresh_individuals_cache(cache_path, patient_file_paths, expected_pids, group_name, sleep_stage):
    """Return cached individuals when source pickle files and patient IDs match."""
    if not os.path.exists(cache_path):
        return None, "missing"

    cache_mtime = os.path.getmtime(cache_path)
    newer_pickles = [
        os.path.basename(path)
        for path in patient_file_paths
        if os.path.getmtime(path) > cache_mtime
    ]
    if newer_pickles:
        return None, f"{len(newer_pickles)} source pickle(s) newer than cache"

    try:
        with open(cache_path, 'rb') as f:
            individuals = pickle.load(f)
    except Exception as e:
        return None, f"failed to load cache: {e}"

    cached_pids = {ind[0] for ind in individuals if ind is not None and len(ind) > 0}
    expected_pids = set(expected_pids)
    comparison_pids = cached_pids
    cache_version_valid = False
    meta_path = _individuals_cache_meta_path(cache_path)
    if os.path.exists(meta_path):
        try:
            with open(meta_path, 'rb') as f:
                meta = pickle.load(f)
            comparison_pids = set(meta.get('attempted_pids', cached_pids))
            cache_version_valid = meta.get('cache_version') == HEP_FLIP_CACHE_SUFFIX
        except Exception:
            comparison_pids = cached_pids

    if comparison_pids != expected_pids:
        added = sorted(expected_pids - comparison_pids)
        removed = sorted(comparison_pids - expected_pids)
        details = []
        if added:
            details.append(f"{len(added)} new patient(s)")
        if removed:
            details.append(f"{len(removed)} removed/excluded patient(s)")
        return None, ", ".join(details) or "patient set changed"

    if individuals and not _individuals_have_flip_metadata(individuals):
        return None, "cache version missing R-peak flip metadata"
    if not individuals and not cache_version_valid:
        return None, "empty cache version could not be verified"

    print(f"Using cache for {group_name}/{sleep_stage}: {os.path.basename(cache_path)}")
    return individuals, "fresh"


def _load_incremental_individuals_cache(cache_path, args_list):
    """Load reusable cached results and return only patient jobs needing work."""
    if not os.path.exists(cache_path):
        return [], args_list, set(), True

    try:
        with open(cache_path, 'rb') as f:
            cached_individuals = pickle.load(f)
    except Exception as e:
        print(f"Could not reuse incremental cache {os.path.basename(cache_path)}: {e}")
        return [], args_list, set(), True

    # Old tuple formats cannot be mixed with results produced by the current
    # R-peak/channel-flipping implementation.
    if not _individuals_have_flip_metadata(cached_individuals):
        return [], args_list, set(), True

    cached_by_pid = {
        ind[0]: ind for ind in cached_individuals
        if ind is not None and len(ind) > 0
    }
    attempted_pids = set(cached_by_pid)
    recorded_sources = {}
    meta_path = _individuals_cache_meta_path(cache_path)
    if os.path.exists(meta_path):
        try:
            with open(meta_path, 'rb') as f:
                meta = pickle.load(f)
            if meta.get('cache_version') == HEP_FLIP_CACHE_SUFFIX:
                attempted_pids = set(meta.get('attempted_pids', attempted_pids))
                recorded_sources = meta.get('source_files', {}) or {}
        except Exception as e:
            print(f"Could not read incremental cache metadata: {e}")

    expected_pids = {args[1] for args in args_list}
    reusable_by_pid = {
        pid: ind for pid, ind in cached_by_pid.items() if pid in expected_pids
    }
    cache_data_changed = set(reusable_by_pid) != set(cached_by_pid)
    pending_args = []
    changed_pids = set()
    for args in args_list:
        path, pid = args[0], args[1]
        source_name = os.path.basename(path)
        recorded_mtime = recorded_sources.get(source_name)
        source_changed = (
            recorded_mtime is not None
            and os.path.getmtime(path) != recorded_mtime
        )
        if pid not in attempted_pids or source_changed:
            pending_args.append(args)
            if source_changed:
                changed_pids.add(pid)
                reusable_by_pid.pop(pid, None)
                cache_data_changed = True

    if pending_args:
        cache_data_changed = True
    return list(reusable_by_pid.values()), pending_args, changed_pids, cache_data_changed


def get_group_individuals(
    group_name, sleep_stage, base_path, test_run=False, recompute_cache=False,
    apply_ica=False, parallel_workers=DEFAULT_PATIENT_WORKERS,
    min_eeg_channels=None,
):
    """
    Loads all files for a group/sleep_stage and returns individual HEPs.
    Returns: list of (patient_id, hep_data, times, ch_names)
    """
    group_dir = os.path.join(base_path, group_name, sleep_stage)
    if not os.path.exists(group_dir):
        return []

    # Exclude cache files when looking for patient pkl files.
    all_patient_files = sorted([
        f for f in os.listdir(group_dir)
        if f.endswith('.pkl')
        and not f.startswith('individuals_cache')
        and not f.startswith('non_eeg_individuals_cache')
    ])
    patient_files = [
        f for f in all_patient_files
        if os.path.getsize(os.path.join(group_dir, f)) >= MIN_COMPLETE_PATIENT_PICKLE_BYTES
    ]
    incomplete_count = len(all_patient_files) - len(patient_files)
    if incomplete_count:
        print(
            f"Skipping {incomplete_count} incomplete patient pickle(s) smaller than "
            f"{MIN_COMPLETE_PATIENT_PICKLE_BYTES // 1024} KiB in {group_name}/{sleep_stage}"
        )
    if not patient_files:
        return []

    if test_run:
        patient_files = patient_files[:10]

    excluded_pids = set()
    excluded_csv = os.path.join(base_path, "excluded_patients.csv")
    if os.path.exists(excluded_csv):
        try:
            exc_df = pd.read_csv(excluded_csv)
            excluded_pids = set(exc_df['patient_id'].dropna().str.strip().tolist())
        except Exception:
            pass

    args_list = []
    for f in patient_files:
        pid = _patient_id_from_pickle_name(f)
        base_pid = pid.split('_')[0]
        if pid not in excluded_pids and base_pid not in excluded_pids:
            args_list.append((os.path.join(group_dir, f), pid, apply_ica, group_name))

    n_jobs = len(args_list)
    if n_jobs == 0:
        return []

    cache_filename = _cache_filename_for_individuals(
        test_run=test_run,
        apply_ica=apply_ica,
        min_eeg_channels=min_eeg_channels,
    )
    cache_path = os.path.join(group_dir, cache_filename)
    patient_file_paths = [args[0] for args in args_list]
    expected_pids = [args[1] for args in args_list]
    if not recompute_cache:
        individuals, cache_reason = _load_fresh_individuals_cache(
            cache_path,
            patient_file_paths,
            expected_pids,
            group_name,
            sleep_stage,
        )
        if individuals is not None:
            return individuals
        if min_eeg_channels:
            print(
                f"Filtered cache (at least {int(min_eeg_channels)} canonical "
                "10-20 EEG channels) "
                f"for {group_name}/{sleep_stage} is {cache_reason}; "
                "refreshing it from the full cache"
            )
        else:
            print(f"Cache for {group_name}/{sleep_stage} needs recompute: {cache_reason}")

    if min_eeg_channels:
        # Build the channel-filtered cache from the normal cache. That cache is
        # itself incremental, so only newly added/changed raw patients need
        # expensive processing before this inexpensive filtering step.
        all_individuals = get_group_individuals(
            group_name,
            sleep_stage,
            base_path,
            test_run=test_run,
            recompute_cache=False,
            apply_ica=apply_ica,
            parallel_workers=parallel_workers,
            min_eeg_channels=None,
        )
        individuals = [
            ind for ind in all_individuals
            if ind is not None
            and len(ind) > 3
            and ind[3] is not None
            and _standard_1020_channel_count(ind[3]) >= int(min_eeg_channels)
        ]
        filtered_data_needs_write = True
        if os.path.exists(cache_path):
            try:
                with open(cache_path, 'rb') as f:
                    previous_filtered = pickle.load(f)
                previous_pids = {ind[0] for ind in previous_filtered if ind is not None}
                current_pids = {ind[0] for ind in individuals if ind is not None}
                with open(_individuals_cache_meta_path(cache_path), 'rb') as f:
                    previous_meta = pickle.load(f)
                recorded_sources = previous_meta.get('source_files', {}) or {}
                current_sources_unchanged = all(
                    recorded_sources.get(os.path.basename(path)) == os.path.getmtime(path)
                    for path in patient_file_paths
                )
                filtered_data_needs_write = not (
                    previous_pids == current_pids and current_sources_unchanged
                )
            except Exception:
                filtered_data_needs_write = True
        try:
            if filtered_data_needs_write:
                print(
                    f"Writing {len(individuals)} filtered patient result(s) to "
                    f"{os.path.basename(cache_path)}..."
                )
                _dump_pickle_atomic_with_heartbeat(individuals, cache_path)
            else:
                print(
                    f"Filtered patient data is unchanged; keeping existing "
                    f"{os.path.basename(cache_path)}"
                )
            with open(_individuals_cache_meta_path(cache_path), 'wb') as f:
                pickle.dump({
                    'attempted_pids': expected_pids,
                    'successful_pids': [ind[0] for ind in individuals],
                    'source_files': {
                        os.path.basename(path): os.path.getmtime(path)
                        for path in patient_file_paths
                    },
                    'apply_ica': apply_ica,
                    'test_run': test_run,
                    'min_eeg_channels': int(min_eeg_channels),
                    'channel_filter': 'canonical_1020_v1',
                    'cache_version': HEP_FLIP_CACHE_SUFFIX,
                }, f)
            print(
                f"Saved {len(individuals)} patient(s) with at least "
                f"{int(min_eeg_channels)} canonical 10-20 EEG channels to "
                f"{os.path.basename(cache_path)}"
            )
        except Exception as e:
            if 'st' in globals():
                st.warning(
                    f"Failed to save channel-filtered cache for "
                    f"{group_name} / {sleep_stage}: {e}"
                )
        return individuals

    individuals = []
    jobs_to_run = args_list
    cache_data_changed = True
    if not recompute_cache:
        individuals, jobs_to_run, changed_pids, cache_data_changed = _load_incremental_individuals_cache(
            cache_path, args_list
        )
        if individuals or len(jobs_to_run) < n_jobs:
            new_count = len(jobs_to_run) - len(changed_pids)
            print(
                f"Reusing {len(individuals)} cached patient(s) for "
                f"{group_name}/{sleep_stage}; processing {new_count} new and "
                f"{len(changed_pids)} changed patient(s)"
            )

    n_pending = len(jobs_to_run)
    if n_pending == 0:
        # The patient set may only have shrunk, so persist the pruned cache and
        # refreshed metadata below without launching a worker pool.
        print(f"No new or changed patients to process for {group_name}/{sleep_stage}")

    progress_bar = st.progress(0, text=f"Loading {group_name} / {sleep_stage} patients...")
    status_text = st.empty()

    completed = 0
    if n_pending:
        max_workers = max(1, min(int(parallel_workers), n_pending, os.cpu_count() or 1))
        status_text.text(f"Starting {n_pending} patient jobs with {max_workers} worker(s)")
        with ProcessPoolExecutor(max_workers=max_workers) as executor:
            futures = {executor.submit(_process_patient_worker, a): a[1] for a in jobs_to_run}
            for future in as_completed(futures, timeout=30000):
                patient_id = futures[future]
                completed += 1
                status_text.text(f"Finished {completed}/{n_pending}: {patient_id}")
                progress_bar.progress(completed / n_pending, text=f"Loading new patients ({completed}/{n_pending})")
                try:
                    result = future.result(timeout=12000)
                except Exception as e:
                    print(f"[Worker] Timeout or error for {patient_id}: {e}")
                    result = None
                if result is not None:
                    individuals.append(result)

    expected_order = {pid: index for index, pid in enumerate(expected_pids)}
    individuals.sort(key=lambda ind: expected_order.get(ind[0], len(expected_order)))

    progress_bar.empty()
    status_text.empty()
    
    # Save to local cache
    try:
        if cache_data_changed:
            print(
                f"Writing {len(individuals)} patient result(s) to "
                f"{os.path.basename(cache_path)}..."
            )
            _dump_pickle_atomic_with_heartbeat(individuals, cache_path)
        else:
            print(f"Patient data unchanged; keeping existing {os.path.basename(cache_path)}")
        with open(_individuals_cache_meta_path(cache_path), 'wb') as f:
            pickle.dump({
                'attempted_pids': expected_pids,
                'successful_pids': [ind[0] for ind in individuals if ind is not None and len(ind) > 0],
                'source_files': {
                    os.path.basename(path): os.path.getmtime(path)
                    for path in patient_file_paths
                },
                'apply_ica': apply_ica,
                'test_run': test_run,
                'cache_version': HEP_FLIP_CACHE_SUFFIX,
            }, f)
    except Exception as e:
        if 'st' in globals():
            st.warning(f"Failed to save cache for {group_name} / {sleep_stage}: {e}")

    return individuals


def get_group_non_eeg_individuals(
    group_name, sleep_stage, base_path, test_run=False, recompute_cache=False,
    apply_ica=False, parallel_workers=DEFAULT_PATIENT_WORKERS,
):
    """
    Loads all files for a group/sleep_stage and returns ECG-aligned averages
    for non-EEG channels.
    Returns: list of (patient_id, non_eeg_data, times, ch_names, rpeaks, log_msg)
    """
    group_dir = os.path.join(base_path, group_name, sleep_stage)
    if not os.path.exists(group_dir):
        return []

    patient_files = [f for f in os.listdir(group_dir) if f.endswith('.pkl') and not f.startswith('non_eeg_individuals_cache') and not f.startswith('individuals_cache')]
    if not patient_files:
        return []

    ica_suffix = '_ica' if apply_ica else ''
    cache_filename = (
        f'non_eeg_individuals_cache_test{ica_suffix}.pkl'
        if test_run else f'non_eeg_individuals_cache{ica_suffix}.pkl'
    )
    cache_path = os.path.join(group_dir, cache_filename)

    if os.path.exists(cache_path) and not test_run and not recompute_cache:
        cache_mtime = os.path.getmtime(cache_path)
        is_cache_valid = True
        for f in patient_files:
            if os.path.getmtime(os.path.join(group_dir, f)) > cache_mtime:
                is_cache_valid = False
                break

        if is_cache_valid:
            try:
                with open(cache_path, 'rb') as f:
                    individuals = pickle.load(f)

                current_pids = set(f.replace('.pkl', '').replace('.edf', '') for f in patient_files)
                cache_invalid = any(ind[0] not in current_pids for ind in individuals)
                if not cache_invalid:
                    return individuals
            except Exception as e:
                if 'st' in globals():
                    st.warning(f"Failed to load non-EEG cache: {e}. Recomputing...")

    if test_run:
        patient_files = patient_files[:10]

    individuals = []
    progress_bar = st.progress(0, text=f"Loading non-EEG channels for {group_name} / {sleep_stage} patients...")
    status_text = st.empty()
    n_files = len(patient_files)

    excluded_pids = set()
    excluded_csv = os.path.join(base_path, "excluded_patients.csv")
    if os.path.exists(excluded_csv):
        try:
            exc_df = pd.read_csv(excluded_csv)
            excluded_pids = set(exc_df['patient_id'].dropna().str.strip().tolist())
        except Exception:
            pass

    args_list = []
    for f in patient_files:
        pid = f.replace('.pkl', '').replace('.edf', '')
        base_pid = pid.split('_')[0]
        if pid not in excluded_pids and base_pid not in excluded_pids:
            args_list.append((os.path.join(group_dir, f), pid, apply_ica))

    n_jobs = len(args_list)
    if n_jobs == 0:
        progress_bar.empty()
        status_text.empty()
        return []

    completed = 0
    max_workers = max(1, min(int(parallel_workers), n_jobs, os.cpu_count() or 1))
    status_text.text(f"Starting {n_jobs} non-EEG patient jobs with {max_workers} worker(s)")
    with ProcessPoolExecutor(max_workers=max_workers) as executor:
        futures = {executor.submit(_process_non_eeg_patient_worker, a): a[1] for a in args_list}
        for future in as_completed(futures, timeout=30000):
            patient_id = futures[future]
            completed += 1
            status_text.text(f"Finished {completed}/{n_jobs}: {patient_id}")
            progress_bar.progress(completed / n_jobs, text=f"Loading patients ({completed}/{n_jobs})")
            try:
                result = future.result(timeout=12000)
            except Exception as e:
                print(f"[Worker] Timeout or error for non-EEG patient {patient_id}: {e}")
                result = None
            if result is not None:
                individuals.append(result)

    progress_bar.empty()
    status_text.empty()

    try:
        with open(cache_path, 'wb') as f:
            pickle.dump(individuals, f)
    except Exception as e:
        if 'st' in globals():
            st.warning(f"Failed to save non-EEG cache for {group_name} / {sleep_stage}: {e}")

    return individuals

def _get_pynapple_jitter_shift(max_jitter_sec, sfreq):
    """
    Returns a random circular shift (in samples) using pynapple's jitter_timestamps
    if available, otherwise falls back to numpy uniform sampling.

    pynapple.randomize.jitter_timestamps expects a nap.Ts object (discrete event times)
    and shifts each timestamp by an independent uniform draw in [-max_jitter, +max_jitter].
    We create a dummy single-event Ts at t=0 and read off the jittered position as our shift.
    """
    if PYNAPPLE_AVAILABLE:
        dummy_ts = nap.Ts(t=np.array([0.0]))
        jittered = nap.randomize.jitter_timestamps(dummy_ts, max_jitter=max_jitter_sec)
        shift_sec = float(jittered.t[0])   # in [-max_jitter_sec, +max_jitter_sec]
    else:
        shift_sec = np.random.uniform(-max_jitter_sec, max_jitter_sec)
    return int(shift_sec * sfreq)


def _patient_cluster_stats(patient_trace, null_traces, p_threshold, n_permutations, times=None):
    """
    Single-patient cluster-mass permutation test.

    For a single subject we cannot compute a t-statistic across subjects, so we
    use the raw signal amplitude as the "statistic" and define the cluster threshold
    as the |amplitude| > mean + 2*std of the null distribution of peak amplitudes.

    Parameters
    ----------
    patient_trace : np.ndarray  shape (n_times,)
        The patient's observed HEP trace.
    null_traces : np.ndarray  shape (n_permutations, n_times)
        Jitter-null distribution for this patient.
    p_threshold : float
    n_permutations : int

    Returns
    -------
    dict
        {
            'p_value': float,
            'n_windows': int,
            'total_duration_sec': float,
        }
    """
    # Threshold based on null distribution of absolute amplitudes
    null_abs = np.abs(null_traces)          # (n_perm, n_times)
    amp_thresh = np.percentile(null_abs, (1 - p_threshold) * 100)

    # Observed cluster mass
    obs_mask = np.abs(patient_trace) > amp_thresh
    obs_labels, n_obs_clusters = label(obs_mask)
    if n_obs_clusters == 0:
        return {
            'p_value': 1.0,
            'n_windows': 0,
            'total_duration_sec': 0.0,
        }

    obs_mass = max(np.sum(np.abs(patient_trace[obs_labels == i + 1])) for i in range(n_obs_clusters))

    # Null cluster masses
    null_masses = np.zeros(n_permutations)
    for p in range(n_permutations):
        null_trace = null_traces[p]
        null_mask = np.abs(null_trace) > amp_thresh
        null_labels, n_null_c = label(null_mask)
        if n_null_c > 0:
            null_masses[p] = max(np.sum(np.abs(null_trace[null_labels == j + 1])) for j in range(n_null_c))

    p_val = (np.sum(null_masses >= obs_mass) + 1) / (n_permutations + 1)
    sig_window_count = 0
    total_duration_sec = 0.0
    if p_val < p_threshold:
        for i in range(n_obs_clusters):
            indices = np.where(obs_labels == i + 1)[0]
            if len(indices) == 0:
                continue
            sig_window_count += 1
            if times is not None and len(indices) > 1:
                total_duration_sec += float(times[indices[-1]] - times[indices[0]])
            elif times is not None and len(times) > 1:
                total_duration_sec += float(np.mean(np.diff(times)))

    return {
        'p_value': float(p_val),
        'n_windows': int(sig_window_count),
        'total_duration_sec': float(max(total_duration_sec, 0.0)),
    }


def _patient_cluster_p_value(patient_trace, null_traces, p_threshold, n_permutations):
    """Backward-compatible wrapper returning only the single-patient p-value."""
    return _patient_cluster_stats(
        patient_trace,
        null_traces,
        p_threshold,
        n_permutations,
    )['p_value']


def permutation_cluster_jitter_test(avg_hep, times, n_permutations=100, p_threshold=0.05, jitter_sec=None,
                                     qrs_exclude_sec=0.05):
    """
    Cluster-based permutation test with pynapple jitter, per-patient significance,
    and Fisher-combined group p-value.

    Parameters
    ----------
    avg_hep : np.ndarray  shape (n_subjects, n_times)
        HEP data matrix — one row per subject.
    times : np.ndarray  shape (n_times,)
        Time axis (seconds).
    n_permutations : int
        Permutations for the null distribution (default 100).
    p_threshold : float
        Cluster-level alpha (default 0.05).
    jitter_sec : float or None
        Maximum circular-shift jitter in seconds.  Uses pynapple's
        ``jitter_timestamps`` to draw the shift amount for each subject.
        If None, shifts span the full epoch length.

    Returns
    -------
    significant_windows : list of dict
        Each dict has 'start', 'end', 'p_value' for the GROUP-level test.
        The 'p_value' is the Fisher-combined p-value across patients when
        all per-patient p-values are available; otherwise it is the standard
        permutation p-value.
    t_obs : np.ndarray  shape (n_times,)
        Observed t-statistic trace (group t-test vs 0).
    per_patient_info : dict
        'p_values'       : list[float]  – per-patient p-value
        'significant'    : list[bool]   – True if patient is sig. at p_threshold
        'n_significant'  : int          – count of significant patients
        'fisher_p'       : float        – Fisher-combined p-value across patients
    """
    n_subjects, n_times = avg_hep.shape

    # ── Sampling rate & max jitter in samples ────────────────────────────────
    sfreq = 1.0 / np.mean(np.diff(times))
    if jitter_sec is not None:
        max_jitter_sec = float(jitter_sec)
        max_shift      = int(max_jitter_sec * sfreq)
    else:
        max_jitter_sec = float(n_times) / sfreq   # full epoch
        max_shift      = n_times

    # ── Per-patient null traces (shared design) ───────────────────────────────
    # Build one null-trace matrix per patient: (n_permutations, n_times)
    # Jitter amounts are drawn via pynapple jitter_timestamps.
    patient_null = [np.zeros((n_permutations, n_times)) for _ in range(n_subjects)]
    for perm in range(n_permutations):
        for s in range(n_subjects):
            shift = _get_pynapple_jitter_shift(max_jitter_sec, sfreq)
            shift = np.clip(shift, -max_shift, max_shift)
            patient_null[s][perm] = np.roll(avg_hep[s], shift)

    # ── Per-patient cluster permutation p-values ─────────────────────────────
    patient_pvals = []
    patient_n_windows = []
    patient_total_duration_sec = []
    for s in range(n_subjects):
        patient_stats = _patient_cluster_stats(
            avg_hep[s],
            patient_null[s],
            p_threshold,
            n_permutations,
            times=times,
        )
        patient_pvals.append(patient_stats['p_value'])
        patient_n_windows.append(patient_stats['n_windows'])
        patient_total_duration_sec.append(patient_stats['total_duration_sec'])

    patient_significant = [p < p_threshold for p in patient_pvals]
    n_significant       = int(sum(patient_significant))

    # Fisher's combined probability test across patients
    # X² = -2 * Σ ln(p_i),  df = 2*k
    eps = 1e-15
    log_sum = sum(np.log(max(p, eps)) for p in patient_pvals)
    chi2_stat  = -2.0 * log_sum
    fisher_p   = float(1.0 - stats.chi2.cdf(chi2_stat, df=2 * n_subjects))

    per_patient_info = {
        'p_values'    : patient_pvals,
        'n_windows'   : patient_n_windows,
        'total_duration_sec': patient_total_duration_sec,
        'significant' : patient_significant,
        'n_significant': n_significant,
        'fisher_p'    : fisher_p,
    }

    # ── Group-level cluster permutation test ─────────────────────────────────
    # t-thresh for group test (n_subjects - 1 df)
    t_thresh = stats.t.ppf(1 - p_threshold / 2, df=max(n_subjects - 1, 1))

    # Observed t-stats
    t_obs, _ = stats.ttest_1samp(avg_hep, 0)

    # QRS complex dominates ±qrs_exclude_sec around the R-peak (cardiac field
    # artifact, not neural signal) — exclude it from cluster detection.
    qrs_mask = np.abs(times) <= qrs_exclude_sec

    obs_mask = (np.abs(t_obs) > t_thresh) & ~qrs_mask
    obs_labels, n_clusters = label(obs_mask)

    if n_clusters == 0:
        return [], t_obs, per_patient_info

    obs_cluster_masses = [
        np.sum(np.abs(t_obs[obs_labels == i + 1])) for i in range(n_clusters)
    ]

    # Build group null distribution from the already-computed patient_null matrices
    null_dist = np.zeros(n_permutations)
    for perm in range(n_permutations):
        # Stack jittered traces for all subjects at this permutation index
        hep_jittered = np.vstack([patient_null[s][perm] for s in range(n_subjects)])
        t_null, _    = stats.ttest_1samp(hep_jittered, 0)
        null_mask    = (np.abs(t_null) > t_thresh) & ~qrs_mask
        null_labels, n_null_c = label(null_mask)
        if n_null_c > 0:
            null_dist[perm] = max(
                np.sum(np.abs(t_null[null_labels == j + 1])) for j in range(n_null_c)
            )

    # Compile significant windows – substitute Fisher p-value where informative
    significant_windows = []
    for i in range(n_clusters):
        perm_p = (np.sum(null_dist >= obs_cluster_masses[i]) + 1) / (n_permutations + 1)
        # Use the more conservative of the two
        combined_p = max(perm_p, fisher_p) if fisher_p < p_threshold else perm_p
        if perm_p < p_threshold:
            indices = np.where(obs_labels == i + 1)[0]
            significant_windows.append({
                'start'      : times[indices[0]],
                'end'        : times[indices[-1]],
                'p_value'    : perm_p,
                'fisher_p'   : fisher_p,
                'n_sig_patients': n_significant,
                'n_patients' : n_subjects,
            })

    return significant_windows, t_obs, per_patient_info


_perm_test_call_count = 0


def permutation_two_group_cluster_test(hep_a, hep_b, times, n_permutations=100, p_threshold=0.05, jitter_sec=None,
                                        label_a="Group A", label_b="Group B", channel_label=None, button_key=None,
                                        show_explanation=True):
    """
    Cluster-based permutation test using Pynapple's independent subject jitter.

    hep_a, hep_b: np.ndarray (n_subjects, n_times)
    times: np.ndarray (n_times)
    label_a, label_b: display names shown in the explanation UI
    button_key: unique Streamlit widget key for the explain button (auto-generated if None)
    show_explanation: set False for batch calls that should not render Streamlit explanation buttons
    """
    global _perm_test_call_count
    _perm_test_call_count += 1
    if button_key is None:
        button_key = f"explain_perm_test_{_perm_test_call_count}"

    n_a, n_times = hep_a.shape
    n_b = hep_b.shape[0]
    n_total = n_a + n_b
    
    # We vstack to create the "Pool" for label swapping
    pooled = np.vstack([hep_a, hep_b]) 
    
    # Calculate Observed T-stat (Welch's for unequal n)
    t_obs, _ = stats.ttest_ind(hep_a, hep_b, axis=0, equal_var=False)
    
    # Observed Cohen's d (Pooled SD)
    d_num = np.mean(hep_a, 0) - np.mean(hep_b, 0)
    d_den = np.sqrt((np.var(hep_a, 0) + np.var(hep_b, 0)) / 2)
    cohens_d = d_num / np.where(d_den == 0, 1e-12, d_den)

    # Threshold for clustering
    df = n_a + n_b - 2
    t_thresh = stats.t.ppf(1 - p_threshold / 2, df=df)

    # Find observed clusters
    obs_mask = np.abs(t_obs) > t_thresh
    obs_labels, n_clusters = label(obs_mask)
    
    if n_clusters == 0:
        if show_explanation:
            explain_permutation_test(
                hep_a, hep_b, times, [], t_obs, cohens_d,
                label_a=label_a, label_b=label_b,
                p_threshold=p_threshold, jitter_sec=jitter_sec,
                null_dist=np.zeros(n_permutations), channel_label=channel_label, button_key=button_key,
            )
        return [], t_obs, cohens_d

    obs_cluster_masses = [np.sum(np.abs(t_obs[obs_labels == i + 1])) for i in range(n_clusters)]

    # --- Permutation Loop ---
    null_dist = np.zeros(n_permutations)
    
    for p in range(n_permutations):
        # 1. Shuffle group labels
        perm_idx = np.random.permutation(n_total)
        perm_data = pooled[perm_idx]
        
        group_a = perm_data[:n_a]
        group_b = perm_data[n_a:]
        
        # 2. Apply Pynapple Jitter to each subject independently
        if jitter_sec is not None:
            group_a = _nap_jitter(group_a, times, jitter_sec)
            group_b = _nap_jitter(group_b, times, jitter_sec)

        # 3. Calculate Null T-stat
        t_null, _ = stats.ttest_ind(group_a, group_b, axis=0, equal_var=False)
        
        # 4. Max Cluster Mass for Null Distribution
        null_mask = np.abs(t_null) > t_thresh
        null_labels, n_null = label(null_mask)
        if n_null > 0:
            null_dist[p] = max(np.sum(np.abs(t_null[null_labels == j + 1])) for j in range(n_null))

    # --- Final Significance Compilation ---
    significant_windows = []
    for i in range(n_clusters):
        p_val = (np.sum(null_dist >= obs_cluster_masses[i]) + 1) / (n_permutations + 1)
        if p_val < p_threshold:
            indices = np.where(obs_labels == i + 1)[0]
            significant_windows.append({
                'start': times[indices[0]],
                'end': times[indices[-1]],
                'p_value': p_val,
                'direction': 'A>B' if np.mean(t_obs[indices]) > 0 else 'B>A'
            })

    if show_explanation:
        explain_permutation_test(
            hep_a, hep_b, times, significant_windows, t_obs, cohens_d,
            label_a=label_a, label_b=label_b,
            p_threshold=p_threshold, jitter_sec=jitter_sec,
            null_dist=null_dist, channel_label=channel_label, button_key=button_key,
        )
    return significant_windows, t_obs, cohens_d


def explain_permutation_test(
    hep_a, hep_b, times, significant_windows, t_obs, cohens_d,
    label_a="Group A", label_b="Group B",
    p_threshold=0.05, n_demo_perms=300, jitter_sec=None,
    null_dist=None, channel_label=None, button_key="explain_perm_test"
):
    """
    Streamlit UI: a button that opens a step-by-step visual explanation of
    permutation_two_group_cluster_test. Pass the inputs and outputs of that
    function directly.

    Parameters
    ----------
    hep_a, hep_b  : np.ndarray (n_subjects, n_times)
    times         : np.ndarray (n_times)
    significant_windows, t_obs, cohens_d : returned by permutation_two_group_cluster_test
    label_a, label_b : display names for the two groups
    p_threshold   : same threshold used in the test
    null_dist     : pre-computed null distribution from the test (skips re-running permutations)
    n_demo_perms  : permutations to run only when null_dist is not provided
    jitter_sec    : same value used in the test
    button_key    : unique Streamlit widget key (change if calling multiple times)
    """
    ch_suffix = f" — {channel_label}" if channel_label else ""
    if not st.button(f"Explain Permutation Cluster Test{ch_suffix}", key=button_key):
        return

    st.markdown(f"#### Group HEP Analysis{ch_suffix}")

    sig_windows_a, _, _ = permutation_cluster_jitter_test(
        hep_a, times,
        n_permutations=n_demo_perms,
        p_threshold=p_threshold,
        jitter_sec=jitter_sec,
    )
    sig_windows_b, _, _ = permutation_cluster_jitter_test(
        hep_b, times,
        n_permutations=n_demo_perms,
        p_threshold=p_threshold,
        jitter_sec=jitter_sec,
    )

    # ── Figure 1: Group A median ± MAD ───────────────────────────────────────
    st.markdown(f"**{label_a}** — median ± MAD")
    fig_a, ax_a = plt.subplots(figsize=(14, 5))
    _render_hep_spaghetti(ax_a, hep_a, times, label_a, "#1f77b4")
    _annotate_sig_windows(ax_a, sig_windows_a)
    ax_a.set_title(f"{label_a} - HEP Median +/- MAD (n={hep_a.shape[0]}){ch_suffix}")
    ax_a.legend(loc='upper right', fontsize=7, ncol=3)
    fig_a.tight_layout()
    st.pyplot(fig_a, use_container_width=False)
    plt.close(fig_a)

    # ── Figure 2: Group B median ± MAD ───────────────────────────────────────
    st.markdown(f"**{label_b}** — median ± MAD")
    fig_b, ax_b = plt.subplots(figsize=(14, 5))
    _render_hep_spaghetti(ax_b, hep_b, times, label_b, "#d62728")
    _annotate_sig_windows(ax_b, sig_windows_b)
    ax_b.set_title(f"{label_b} - HEP Median +/- MAD (n={hep_b.shape[0]}){ch_suffix}")
    ax_b.legend(loc='upper right', fontsize=7, ncol=3)
    fig_b.tight_layout()
    st.pyplot(fig_b, use_container_width=False)
    plt.close(fig_b)

    # ── Figure 3: Combined — median ± MAD for both groups ────────────────────
    st.markdown("**Both groups combined** — median ± MAD")
    fig_c, ax_c = plt.subplots(figsize=(14, 5))
    for _data, _lbl, _col in [(hep_a, label_a, "#1f77b4"), (hep_b, label_b, "#d62728")]:
        _med, _mad = median_and_mad(_data * 1e6, axis=0)
        ax_c.plot(times, _med, color=_col, linewidth=2.5, label=f"{_lbl} (n={_data.shape[0]})")
        ax_c.fill_between(times, _med - _mad, _med + _mad, color=_col, alpha=0.2)
    for _win in (significant_windows or []):
        ax_c.axvspan(_win['start'], _win['end'], color='orange', alpha=0.28,
                     label=f"p={_win['p_value']:.3f}")
    ax_c.axhline(0, color='black', linewidth=0.6, linestyle='--')
    ax_c.axvline(0, color='gray', linewidth=0.6, linestyle=':')
    ax_c.set_ylabel('Amplitude (µV)')
    ax_c.set_xlabel('Time (s)')
    ax_c.set_title(f"Both groups - HEP Median +/- MAD{ch_suffix}")
    ax_c.legend(loc='upper right', fontsize=9)
    fig_c.tight_layout()
    st.pyplot(fig_c, use_container_width=False)
    plt.close(fig_c)

def _nap_jitter(data_matrix, times, jitter_sec):
    """Helper to apply independent circular shifts to rows."""
    dt = times[1] - times[0]
    out = np.empty_like(data_matrix)
    for i in range(data_matrix.shape[0]):
        shift_val = np.random.uniform(-jitter_sec, jitter_sec)
        shift_samples = int(round(shift_val / dt))
        out[i] = np.roll(data_matrix[i], shift_samples)
    return out

def save_hep_to_downloads(hep_a: np.ndarray, hep_b: np.ndarray, times: np.ndarray,
                          label_a: str = "group_a", label_b: str = "group_b",
                          downloads_dir: str = "/storage/pblab_shared_data2/Nir/Cobrad/.downloads") -> str:
    """Save hep_a and hep_b matrices (and times) as .npz to .downloads/.

    Returns the path of the saved file.
    """
    os.makedirs(downloads_dir, exist_ok=True)
    filename = f"hep_{label_a}_vs_{label_b}.npz"
    filepath = os.path.join(downloads_dir, filename)
    np.savez(filepath, hep_a=hep_a, hep_b=hep_b, times=times,
             label_a=np.array(label_a), label_b=np.array(label_b))
    return filepath


def finalize_plot(fig, ax, title, avg_hep=None, times=None, n_subjects=None, significant_windows=None, all_heps=None, mad_hep=None, show_legend=True):
    """
    Applies common styling to the plot, optionally plots the average/median and significance, and displays it.
    Appends N, min p-value, and Cohen's d to the title if data is available.
    If mad_hep is provided, treats avg_hep as the median and draws a ±MAD ribbon.
    """
    if avg_hep is not None and times is not None:
        if mad_hep is not None:
            label = f"Group Median (n={n_subjects})" if n_subjects is not None else "Group Median"
            ax.fill_between(times, (avg_hep - mad_hep) * 1e6, (avg_hep + mad_hep) * 1e6,
                            color='blue', alpha=0.15, label='± MAD')
        else:
            label = f"Group Average (n={n_subjects})" if n_subjects is not None else "Group Average"
        ax.plot(times, avg_hep * 1e6, color='blue', linewidth=2, label=label)

    if significant_windows:
        for window in significant_windows:
            start, end = window['start'], window['end']
            p_val = window['p_value']
            
            # Plot rectangle
            ax.axvspan(start, end, color='orange', alpha=0.2)
            
            # Plot asterisk
            mid_point = (start + end) / 2
            
            # Determine asterisks
            if p_val < 0.001:
                asterisks = "***"
            elif p_val < 0.01:
                asterisks = "**"
            else:
                asterisks = "*"
            current_ymax = ax.get_ylim()[1]
            ax.text(mid_point, current_ymax-.4, asterisks, ha='center', va='bottom', fontsize=16, color='orange', fontweight='bold')

    # Build stats suffix for title
    stats_parts = []
    if n_subjects is not None:
        stats_parts.append(f"N={n_subjects}")
    if significant_windows:
        min_p = min(w['p_value'] for w in significant_windows)
        min_p_win = min(significant_windows, key=lambda w: w['p_value'])
        win_str = f"[{min_p_win['start']:.3f}s, {min_p_win['end']:.3f}s]"
        stats_parts.append(f"p={min_p:.2e} {win_str}")
        # Per-patient count and Fisher p (present when using pynapple-jitter test)
        n_sig_pt = min_p_win.get('n_sig_patients')
        n_tot_pt = min_p_win.get('n_patients')
        fish_p   = min_p_win.get('fisher_p')
        if n_sig_pt is not None and n_tot_pt is not None:
            stats_parts.append(f"{n_sig_pt}/{n_tot_pt} pts sig")
        if fish_p is not None:
            stats_parts.append(f"Fisher p={fish_p:.3f}")
    else:
        stats_parts.append("p=n.s.")
    if all_heps is not None and len(all_heps) > 1:
        try:
            heps_arr = np.array(all_heps)  # (n_subjects, n_times)
            # Cohen's d at the peak absolute amplitude
            mean_arr = np.mean(heps_arr, axis=0)
            std_arr = np.std(heps_arr, axis=0, ddof=1)
            peak_idx = np.argmax(np.abs(mean_arr))
            d = abs(mean_arr[peak_idx]) / (std_arr[peak_idx] + 1e-12)
            stats_parts.append(f"d={d:.2f}")
        except Exception:
            pass

    full_title = title
    if mad_hep is not None and "MAD" not in full_title.upper():
        full_title = f"{full_title} | Median +/- MAD"
    if stats_parts:
        full_title = f"{full_title}\n({', '.join(stats_parts)})"

    ax.set_xlabel("Time (s)")
    ax.set_ylabel("Amplitude (μV)")
    ax.set_title(full_title)
    if show_legend:
        handles, labels = ax.get_legend_handles_labels()
        if len(handles) > 12:
            ax.legend(handles, labels, fontsize=7, ncol=2, loc='upper right',
                      bbox_to_anchor=(1.0, 1.0), framealpha=0.6)
        else:
            ax.legend(handles, labels)
    ax.grid(True)
    ax.axvline(0, color='r', linestyle='--', alpha=0.5)
    if times is not None:
        ax.set_xlim(times[0], times[-1])
    fig.tight_layout()
    st.pyplot(fig, use_container_width=True)

def apply_ecg_regression(hep_matrix, ecg_vector):
    """Remove ECG artifact from HEP epochs via linear regression."""
    out = hep_matrix.copy().astype(float)
    if ecg_vector is None:
        return out
    ecg_v = np.asarray(ecg_vector).squeeze()
    if ecg_v.ndim != 1 or len(ecg_v) != out.shape[1]:
        return out
    ecg_den = np.dot(ecg_v, ecg_v)
    if ecg_den < 1e-20:
        return out
    eig = np.dot(out, ecg_v) / ecg_den
    return out - eig[:, None] * ecg_v[None, :]


def zscore_per_subject(matrix):
    """Z-score each row (subject) of a (n_subj, n_times) matrix."""
    mu = np.nanmean(matrix, axis=1, keepdims=True)
    sigma = np.nanstd(matrix, axis=1, ddof=1, keepdims=True)
    sigma = np.where(sigma == 0, 1e-12, sigma)
    return (matrix - mu) / sigma


def scale_matrix(matrix, use_zscore):
    """Z-score (per subject) or convert V→µV."""
    if use_zscore:
        return zscore_per_subject(matrix)
    return matrix * 1e6


def median_and_mad(matrix, axis=0):
    """Return nan-median and median absolute deviation along axis."""
    arr = np.asarray(matrix, dtype=float)
    med = np.nanmedian(arr, axis=axis)
    mad = np.nanmedian(np.abs(arr - np.expand_dims(med, axis=axis)), axis=axis)
    return med, mad


def summarize_channels_without_reference_cancellation(channel_matrix, cancellation_ratio=0.05):
    """
    Average channels unless a common-average reference makes the channel mean
    collapse toward zero; in that case use the channel median as a robust
    representative trace.
    """
    channel_matrix = np.asarray(channel_matrix, dtype=float)
    if channel_matrix.ndim != 2 or channel_matrix.shape[0] == 0:
        return np.array([])

    mean_trace = np.nanmean(channel_matrix, axis=0)
    if channel_matrix.shape[0] < 3:
        return mean_trace

    mean_scale = np.nanstd(mean_trace)
    channel_scale = np.nanmedian(np.nanstd(channel_matrix, axis=1))
    if np.isfinite(channel_scale) and channel_scale > 0:
        if mean_scale / channel_scale < cancellation_ratio:
            return np.nanmedian(channel_matrix, axis=0)
    return mean_trace


def identify_common_eeg_channels(individuals, min_fraction=0.5):
    """Return EEG channel names present in >= min_fraction of individuals."""
    all_ch_sets = [set(ind[3]) for ind in individuals]
    counts = Counter([ch for s in all_ch_sets for ch in s])
    return [
        ch for ch, count in counts.items()
        if count >= len(individuals) * min_fraction
        and (re.match(r'^[a-zA-Z]{1,2}[0-9]*$', ch) or re.match(r'^[a-zA-Z]z$', ch))
    ]


def load_excluded_pids(base_path):
    """Load globally excluded patient IDs from base_path/excluded_patients.csv."""
    csv_path = os.path.join(base_path, "excluded_patients.csv")
    if not os.path.exists(csv_path):
        return []
    try:
        df = pd.read_csv(csv_path)
        if 'patient_id' in df.columns:
            return [str(pid) for pid in df['patient_id'].tolist()]
    except Exception as e:
        if 'st' in globals():
            st.warning(f"Failed to load excluded patients CSV: {e}")
    return []


def filter_excluded(individuals, excluded_pids):
    """Remove individuals whose base pid or full pid is in excluded_pids."""
    return [
        ind for ind in individuals
        if (str(ind[0]).split('_')[0] if '_' in str(ind[0]) else str(ind[0])) not in excluded_pids
        and str(ind[0]) not in excluded_pids
    ]


def get_hemisphere_channels(channels):
    """Split channel list into (left, right, mid) by electrode naming convention."""
    left = [ch for ch in channels if re.search(r'[13579]$', ch)]
    right = [ch for ch in channels if re.search(r'[2468]$', ch)]
    mid = [ch for ch in channels if re.search(r'z$', ch, re.IGNORECASE)]
    return left, right, mid


_TOPO_STD19    = ['Fp1','Fp2','F7','F3','Fz','F4','F8','C3','Cz','C4','P3','Pz','P4','O1','O2']
_TOPO_ALIASES  = {'T7':'T3','T8':'T4','P7':'T5','P8':'T6'}

def _render_topomap(ch_values, ax, title, cmap='RdBu_r', vlim=None, colorbar_label=''):
    """
    Plot a topomap from a {channel_name: scalar} dict using the standard 10-20 montage.
    Uses only the 19 standard channels + T3/T4/T5/T6 aliases (same set as Single Group
    Analysis) so MNE never sees more channels than it has sphere positions for.
    Returns the AxesImage or None.
    """
    mont = mne.channels.make_standard_montage('standard_1020')
    mont_upper = [c.upper() for c in mont.ch_names]
    pch, pd = [], []
    for ch, val in ch_values.items():
        cu = ch.upper()
        if cu in mont_upper:
            pch.append(mont.ch_names[mont_upper.index(cu)])
            pd.append(val)
    # pad missing standard-19
    for bc in _TOPO_STD19:
        if not any(bc.upper() == x.upper() for x in pch):
            pch.append(bc)
            pd.append(0.0 if cmap == 'RdBu_r' else 0.05)
    # pad missing T3/T4/T5/T6 aliases
    for nn, on in _TOPO_ALIASES.items():
        if not any(c2.upper() in [nn.upper(), on.upper()] for c2 in pch):
            pch.append(nn)
            pd.append(0.0 if cmap == 'RdBu_r' else 0.05)
    if not pch:
        return None
    info2 = mne.create_info(ch_names=pch, sfreq=250., ch_types='eeg')
    info2.set_montage(mont, on_missing='ignore')
    _valid = np.array([
        not np.any(np.isnan(ch['loc'][:3])) and np.any(ch['loc'][:3] != 0)
        for ch in info2['chs']
    ])
    if not np.any(_valid):
        return None
    pd_arr = np.array(pd)
    da = pd_arr[_valid]
    info2 = mne.pick_info(info2, np.where(_valid)[0])
    if vlim is None:
        vmax = np.max(np.abs(da)) or 1.0
        vlim = (-vmax, vmax)
    res = mne.viz.plot_topomap(da, info2, axes=ax, cmap=cmap, vlim=vlim, extrapolate='head', show=False)
    im = res[0] if isinstance(res, tuple) else res
    cb = plt.colorbar(im, ax=ax)
    if colorbar_label:
        cb.set_label(colorbar_label)
    ax.set_title(title)
    return im


def _render_pval_topomap(p_vals, ax_t, title):
    """Plot a p-value topomap into ax_t. Returns the AxesImage or None."""
    _mont = mne.channels.make_standard_montage('standard_1020')
    _mont_upper = [c.upper() for c in _mont.ch_names]
    _pch, _pd = [], []
    for _c, _p in p_vals.items():
        _cu = _c.upper()
        if _cu in _mont_upper:
            _pch.append(_mont.ch_names[_mont_upper.index(_cu)])
            _pd.append(_p)
    _std19 = ['Fp1', 'Fp2', 'F7', 'F3', 'Fz', 'F4', 'F8',
              'C3', 'Cz', 'C4', 'P3', 'Pz', 'P4', 'O1', 'O2']
    _aliases = {'T7': 'T3', 'T8': 'T4', 'P7': 'T5', 'P8': 'T6'}
    for _bc in _std19:
        if not any(_bc.upper() == _x.upper() for _x in _pch):
            _pch.append(_bc); _pd.append(0.05)
    for _nn, _on in _aliases.items():
        if not any(_c2.upper() in [_nn.upper(), _on.upper()] for _c2 in _pch):
            _pch.append(_nn); _pd.append(0.05)
    if not _pch:
        return None
    _info2 = mne.create_info(ch_names=_pch, sfreq=250., ch_types='eeg')
    _info2.set_montage(_mont, on_missing='ignore')
    _valid2 = np.array([
        not np.any(np.isnan(_ch['loc'][:3])) and np.any(_ch['loc'][:3] != 0)
        for _ch in _info2['chs']
    ])
    if not np.any(_valid2):
        return None
    _pch_v = [_pch[i] for i in np.where(_valid2)[0]]
    _da = np.clip(np.array(_pd)[_valid2], 0, 0.05)
    _info2 = mne.pick_info(_info2, np.where(_valid2)[0])
    _res = mne.viz.plot_topomap(
        _da, _info2, axes=ax_t,
        cmap='Reds_r', names=_pch_v, vlim=(0, 0.05), extrapolate='head'
    )
    _im = _res[0] if isinstance(_res, tuple) else _res
    _cb = plt.colorbar(_im, ax=ax_t)
    _cb.set_label("p-value")
    ax_t.set_title(title)
    return _im


def _to_montage_channel_name(ch_name, montage):
    """Return the montage spelling for a channel name, including old temporal aliases."""
    if ch_name is None:
        return None

    aliases = {
        'T3': 'T7',
        'T4': 'T8',
        'T5': 'P7',
        'T6': 'P8',
    }
    mont_upper = {ch.upper(): ch for ch in montage.ch_names}
    ch_upper = str(ch_name).upper()
    ch_upper = aliases.get(ch_upper, ch_upper)
    return mont_upper.get(ch_upper)


def _render_existing_channel_pval_topomap(
    p_vals, ax_t, title, p_threshold=0.05, show_names=True
):
    """
    Plot a p-value topomap using only channels supplied in p_vals.
    Unlike _render_pval_topomap, this does not pad missing electrodes, so labels
    are shown only for electrodes that actually exist in that group.
    """
    montage = mne.channels.make_standard_montage('standard_1020')
    plotted_names = []
    display_names = []
    plotted_pvals = []

    for ch, p_val in p_vals.items():
        mont_name = _to_montage_channel_name(ch, montage)
        if mont_name is None or mont_name in plotted_names:
            continue
        plotted_names.append(mont_name)
        display_names.append(str(ch))
        plotted_pvals.append(float(p_val) if np.isfinite(p_val) else 1.0)

    if len(plotted_names) < 4:
        ax_t.set_title(title)
        ax_t.text(
            0.5, 0.5, "Too few montage electrodes",
            ha='center', va='center', transform=ax_t.transAxes, color='gray'
        )
        ax_t.axis('off')
        return None

    info = mne.create_info(ch_names=plotted_names, sfreq=250., ch_types='eeg')
    info.set_montage(montage, on_missing='ignore')
    valid = np.array([
        not np.any(np.isnan(ch['loc'][:3])) and np.any(ch['loc'][:3] != 0)
        for ch in info['chs']
    ])
    if np.sum(valid) < 4:
        ax_t.set_title(title)
        ax_t.text(
            0.5, 0.5, "Too few valid montage positions",
            ha='center', va='center', transform=ax_t.transAxes, color='gray'
        )
        ax_t.axis('off')
        return None

    names_valid = [display_names[i] for i in np.where(valid)[0]]
    data = np.clip(np.asarray(plotted_pvals, dtype=float)[valid], 0, p_threshold)
    info = mne.pick_info(info, np.where(valid)[0])
    names_arg = names_valid if show_names else None

    res = mne.viz.plot_topomap(
        data,
        info,
        axes=ax_t,
        cmap='Reds_r',
        names=names_arg,
        vlim=(0, p_threshold),
        extrapolate='head',
        show=False,
    )
    im = res[0] if isinstance(res, tuple) else res
    cb = plt.colorbar(im, ax=ax_t)
    cb.set_label("p-value")
    ax_t.set_title(title)
    return im


def _get_common_channels(individuals):
    """Return the intersection of EEG channel names across a list of individual tuples."""
    common_channels = None
    for ind in individuals or []:
        ch_names = ind[3] if len(ind) > 3 else None
        if ch_names is None:
            continue
        if common_channels is None:
            common_channels = set(ch_names)
        else:
            common_channels &= set(ch_names)
    return sorted(common_channels) if common_channels else []


def _stack_traces_with_common_length(traces_a, traces_b=None, times=None):
    """Trim trace lists to a common length and return stacked arrays."""
    arrays_a = [np.asarray(t, dtype=float).squeeze() for t in traces_a if t is not None]
    arrays_b = [np.asarray(t, dtype=float).squeeze() for t in (traces_b or []) if t is not None]
    all_arrays = arrays_a + arrays_b
    if not all_arrays:
        return None, None, None

    min_len = min(len(arr) for arr in all_arrays)
    if min_len <= 1:
        return None, None, None

    stack_a = np.array([arr[:min_len] for arr in arrays_a], dtype=float) if arrays_a else None
    stack_b = np.array([arr[:min_len] for arr in arrays_b], dtype=float) if arrays_b else None
    times_use = np.asarray(times[:min_len], dtype=float) if times is not None else None
    return stack_a, stack_b, times_use


def _compute_patient_csd_map(hep_data, times, ch_names, candidate_channels):
    """Convert one patient's ICA-cleaned HEP to CSD space and return {channel: trace}."""
    valid_channels = [ch for ch in candidate_channels if ch in ch_names]
    if len(valid_channels) < 4 or len(times) < 2:
        return None

    idx = [ch_names.index(ch) for ch in valid_channels]
    data = np.asarray(hep_data[idx, :], dtype=float)
    sfreq = 1.0 / max(float(times[1] - times[0]), 1e-9)
    montage = mne.channels.make_standard_montage('standard_1020')

    try:
        info = mne.create_info(ch_names=valid_channels, sfreq=sfreq, ch_types='eeg')
        info.set_montage(montage, on_missing='ignore')
        evoked = mne.EvokedArray(data, info, tmin=float(times[0]), verbose=False)
        evoked_csd = mne.preprocessing.compute_current_source_density(evoked)
        return {
            ch: evoked_csd.data[i].copy()
            for i, ch in enumerate(evoked_csd.ch_names)
        }
    except Exception:
        return None


def _compute_stage_or_group_pvals(individuals, n_permutations=200, jitter_sec=0.1, use_csd=False):
    """Compute per-channel minimum cluster p-values against baseline for one set of individuals."""
    if not individuals:
        return {}

    common_channels = _get_common_channels(individuals)
    if not common_channels:
        return {}

    channel_traces = {ch: [] for ch in common_channels}
    times_ref = None

    for ind in individuals:
        hep_data, times, ch_names = ind[1], ind[2], ind[3]
        if hep_data is None or times is None or ch_names is None:
            continue
        if times_ref is None:
            times_ref = np.asarray(times, dtype=float)

        if use_csd:
            csd_map = _compute_patient_csd_map(hep_data, times, list(ch_names), common_channels)
            if not csd_map:
                continue
            for ch in common_channels:
                if ch in csd_map:
                    channel_traces[ch].append(csd_map[ch])
        else:
            for ch in common_channels:
                if ch in ch_names:
                    channel_traces[ch].append(np.asarray(hep_data[ch_names.index(ch)], dtype=float))

    channel_pvals = {}
    for ch, traces in channel_traces.items():
        if len(traces) < 3:
            continue
        stacked, _, times_use = _stack_traces_with_common_length(traces, times=times_ref)
        if stacked is None or times_use is None or stacked.shape[0] < 3:
            continue
        try:
            sig_windows, _, _ = permutation_cluster_jitter_test(
                stacked,
                times_use,
                n_permutations=n_permutations,
                p_threshold=0.05,
                jitter_sec=jitter_sec,
            )
            channel_pvals[ch] = min((w['p_value'] for w in sig_windows), default=1.0)
        except Exception:
            channel_pvals[ch] = 1.0

    return channel_pvals


def _compute_two_group_channel_pvals(individuals_a, individuals_b, n_permutations=200, jitter_sec=0.1, use_csd=False):
    """Compute per-channel minimum cluster p-values for the difference between two groups."""
    if not individuals_a or not individuals_b:
        return {}

    common_a = set(_get_common_channels(individuals_a))
    common_b = set(_get_common_channels(individuals_b))
    common_channels = sorted(common_a & common_b)
    if not common_channels:
        return {}

    channel_traces_a = {ch: [] for ch in common_channels}
    channel_traces_b = {ch: [] for ch in common_channels}
    times_ref = None

    for inds, out_map in ((individuals_a, channel_traces_a), (individuals_b, channel_traces_b)):
        for ind in inds:
            hep_data, times, ch_names = ind[1], ind[2], ind[3]
            if hep_data is None or times is None or ch_names is None:
                continue
            if times_ref is None:
                times_ref = np.asarray(times, dtype=float)

            if use_csd:
                csd_map = _compute_patient_csd_map(hep_data, times, list(ch_names), common_channels)
                if not csd_map:
                    continue
                for ch in common_channels:
                    if ch in csd_map:
                        out_map[ch].append(csd_map[ch])
            else:
                for ch in common_channels:
                    if ch in ch_names:
                        out_map[ch].append(np.asarray(hep_data[ch_names.index(ch)], dtype=float))

    channel_pvals = {}
    for ch in common_channels:
        traces_a = channel_traces_a[ch]
        traces_b = channel_traces_b[ch]
        if len(traces_a) < 2 or len(traces_b) < 2:
            continue
        stack_a, stack_b, times_use = _stack_traces_with_common_length(traces_a, traces_b, times_ref)
        if stack_a is None or stack_b is None or times_use is None:
            continue
        if stack_a.shape[0] < 2 or stack_b.shape[0] < 2:
            continue
        try:
            sig_windows, _, _ = permutation_two_group_cluster_test(
                stack_a,
                stack_b,
                times_use,
                n_permutations=n_permutations,
                p_threshold=0.05,
                jitter_sec=jitter_sec,
            )
            channel_pvals[ch] = min((w['p_value'] for w in sig_windows), default=1.0)
        except Exception:
            channel_pvals[ch] = 1.0

    return channel_pvals


def _render_hep_spaghetti(ax, hep_matrix, times, label, color, pids=None):
    """Plot group median ± MAD HEP trace on ax.

    Parameters
    ----------
    ax         : matplotlib Axes
    hep_matrix : np.ndarray (n_subj, n_times) in Volts — plotted as µV
    times      : np.ndarray (n_times)
    label      : str — group label shown in the mean legend entry
    color      : str — color for the mean line
    pids       : list[str] | None — optional per-subject labels (used for legends)
    """
    n_subj = hep_matrix.shape[0]
    med_hep, mad_hep = median_and_mad(hep_matrix * 1e6, axis=0)
    ax.plot(times, med_hep, color=color, linewidth=2.5,
            label=f'Median {label} (n={n_subj})', zorder=5)
    ax.fill_between(times, med_hep - mad_hep, med_hep + mad_hep,
                    color=color, alpha=0.20, label='± MAD')
    ax.axhline(0, color='black', linewidth=0.6, linestyle='--')
    ax.axvline(0, color='gray', linewidth=0.6, linestyle=':')
    ax.set_ylabel('Amplitude (µV)')
    ax.set_xlabel('Time (s)')


def _annotate_sig_windows(ax, significant_windows, y_pad_frac=0.06):
    """Overlay significant windows with p-value labels on a waveform axis."""
    if not significant_windows:
        return

    y_min, y_max = ax.get_ylim()
    y_span = max(y_max - y_min, 1e-9)
    text_y = y_max - y_span * y_pad_frac
    seen = set()
    for win in significant_windows:
        start = win['start']
        end = win['end']
        p_val = win.get('p_value', np.nan)
        ax.axvspan(start, end, color='orange', alpha=0.18)
        p_text = "<0.001" if np.isfinite(p_val) and p_val < 0.001 else f"{p_val:.3f}"
        key = (round(start, 6), round(end, 6), p_text)
        if key in seen:
            continue
        seen.add(key)
        ax.text(
            (start + end) / 2,
            text_y,
            f"p={p_text}",
            ha='center',
            va='top',
            fontsize=8,
            color='darkorange',
            fontweight='bold',
            bbox=dict(boxstyle='round,pad=0.15', facecolor='white', alpha=0.75, edgecolor='none'),
        )


def _summarize_channel_group_difference(sig_windows, cohens_d, times):
    """Collapse significant per-channel results into compact circular-plot metrics."""
    if cohens_d is None or times is None or len(times) == 0:
        return {
            'effect_size': 0.0,
            'signed_effect': 0.0,
            'total_sig_duration': 0.0,
            'min_p': 1.0,
            'n_windows': 0,
        }

    sig_windows = sig_windows or []
    if not sig_windows:
        return {
            'effect_size': 0.0,
            'signed_effect': 0.0,
            'total_sig_duration': 0.0,
            'min_p': 1.0,
            'n_windows': 0,
        }

    total_sig_duration = 0.0
    effect_values = []
    signed_values = []
    min_p = 1.0

    cohens_d = np.asarray(cohens_d)
    for win in sig_windows:
        start = float(win.get('start', times[0]))
        end = float(win.get('end', times[-1]))
        min_p = min(min_p, float(win.get('p_value', 1.0)))
        total_sig_duration += max(0.0, end - start)

        idx = np.where((times >= start) & (times <= end))[0]
        if idx.size == 0:
            continue
        win_cd = cohens_d[idx]
        if win_cd.size == 0:
            continue
        effect_values.append(float(np.nanmax(np.abs(win_cd))))
        signed_values.append(float(np.nanmean(win_cd)))

    if not effect_values:
        return {
            'effect_size': 0.0,
            'signed_effect': 0.0,
            'total_sig_duration': total_sig_duration,
            'min_p': min_p,
            'n_windows': len(sig_windows),
        }

    return {
        'effect_size': float(np.nanmax(effect_values)),
        'signed_effect': float(np.nanmean(signed_values)),
        'total_sig_duration': float(total_sig_duration),
        'min_p': float(min_p),
        'n_windows': len(sig_windows),
    }


def _rayleigh_test(angles):
    """Rayleigh test of uniformity for circular data.

    Returns (R_bar, p_value) where R_bar is the mean resultant length
    (0 = uniform, 1 = perfectly concentrated at one phase).
    Uses the standard Rayleigh Z statistic and Zar (1999) p-value approximation.
    """
    n = len(angles)
    if n < 2:
        return 0.0, 1.0
    angles = np.asarray(angles, dtype=float)
    C = float(np.sum(np.cos(angles)))
    S = float(np.sum(np.sin(angles)))
    R = np.sqrt(C ** 2 + S ** 2)
    R_bar = R / n
    Z = n * R_bar ** 2
    # Zar (1999) p-value approximation; good for n >= 2
    p = np.exp(-Z)
    if n < 50:
        p = p * (
            1
            + (2 * Z - Z ** 2) / (4 * n)
            - (24 * Z - 132 * Z ** 2 + 76 * Z ** 3 - 9 * Z ** 4) / (288 * n ** 2)
        )
    return float(R_bar), float(np.clip(p, 0.0, 1.0))


def _plot_per_electrode_circular_summary(
    channel_stats, group_a, group_b, selected_stage,
    channel_data=None, sig_windows_per_ch=None
):
    """Circular HEP summary: ONE polar subplot per electrode.

    The circle maps the HEP time axis:
      - R-peak (t = 0) -> 0 degrees (right / 3-o-clock)
      - t = +0.4 s     -> 180 degrees (left / 9-o-clock)
      - Positive time (0 to 0.4 s) is the top semicircle.
      - Negative time (-0.4 to 0 s) is the bottom semicircle.

    Radial fill shows the mean group difference (group_a minus group_b).
    Orange transparent bands mark statistically significant windows.
    The Rayleigh test checks whether significant time points cluster at
    a particular HEP phase (not uniformly distributed around the circle).
    """
    # Gather valid channels
    if channel_data:
        valid_chs = [
            ch for ch, (ca, cb, t_el) in channel_data.items()
            if ca.shape[0] > 0 and cb.shape[0] > 0 and len(t_el) > 0
        ]
    else:
        valid_chs = [ch for ch, s in channel_stats.items() if s.get('effect_size', 0) > 0]

    if not valid_chs:
        return None

    # Layout
    n_ch = len(valid_chs)
    n_cols = min(4, n_ch)
    n_rows = int(np.ceil(n_ch / n_cols))

    fig, axes = plt.subplots(
        n_rows, n_cols,
        figsize=(4.5 * n_cols, 4.5 * n_rows),
        subplot_kw={'projection': 'polar'},
    )
    axes_flat = np.array(axes).flatten() if n_ch > 1 else np.array([axes])

    dir_colors = {'pos': '#1f77b4', 'neg': '#d95f02'}
    base_r = 1.0
    theta_full = np.linspace(-np.pi, np.pi, 360)

    for i, ch in enumerate(valid_chs):
        ax = axes_flat[i]

        # Data
        diff = None
        t_el = np.array([])
        if channel_data and ch in channel_data:
            ca, cb, t_el = channel_data[ch]
            med_a = np.nanmedian(ca, axis=0)
            med_b = np.nanmedian(cb, axis=0)
            diff = med_a - med_b

        sig_wins = (sig_windows_per_ch or {}).get(ch, [])

        # Polar axes setup
        ax.set_theta_zero_location('N')   # 0 deg = top = R-peak (t=0)
        ax.set_theta_direction(1)          # counter-clockwise: top = positive time
        ax.set_xticks([])
        ax.set_yticks([])
        ax.grid(False)
        ax.spines['polar'].set_visible(False)

        # Base circle
        ax.plot(theta_full, np.full(360, base_r), color='lightgray', linewidth=1.2, zorder=1)

        # Radial group-difference fill
        ray_tag = ''
        if diff is not None and len(t_el) > 0:
            theta_arr = t_el * np.pi / 0.4   # t=0->0 deg, t=0.4->180 deg, t=-0.4->-180 deg

            max_amp = float(np.nanmax(np.abs(diff))) if np.any(np.isfinite(diff)) else 1.0
            if max_amp < 1e-12:
                max_amp = 1.0
            scale = 0.6 / max_amp   # max deviation = 0.6 radius units

            r_diff = base_r + diff * scale

            # Significant-window highlight (drawn first, behind difference fill)
            for win in sig_wins:
                t_s = float(win.get('start', 0))
                t_e = float(win.get('end', 0))
                if t_e <= t_s:
                    continue
                n_pts = max(20, int((t_e - t_s) * 200))
                t_range = np.linspace(t_s, t_e, n_pts)
                theta_range = t_range * np.pi / 0.4
                ax.fill_between(
                    theta_range,
                    0,
                    base_r + 0.65,
                    color='orange', alpha=0.18, linewidth=0, zorder=2,
                )

            # Positive difference (group_a > group_b) - blue outward fill
            pos = diff >= 0
            if np.any(pos):
                ax.fill_between(theta_arr, base_r, r_diff,
                                where=pos, color=dir_colors['pos'],
                                alpha=0.60, linewidth=0, zorder=3)

            # Negative difference (group_b > group_a) - orange/red inward fill
            neg = diff < 0
            if np.any(neg):
                ax.fill_between(theta_arr, r_diff, base_r,
                                where=neg, color=dir_colors['neg'],
                                alpha=0.60, linewidth=0, zorder=3)

            # R-peak marker
            ax.plot([0], [base_r], 'rv', markersize=9, zorder=10, clip_on=False)

            # Rayleigh test on significant time points
            sig_angles = []
            for win in sig_wins:
                mask = (t_el >= float(win.get('start', -999))) & (t_el <= float(win.get('end', 999)))
                sig_angles.extend((t_el[mask] * np.pi / 0.4).tolist())

            if len(sig_angles) >= 2:
                R_bar, ray_p = _rayleigh_test(sig_angles)
                p_str = '<0.001' if ray_p < 0.001 else f'{ray_p:.3f}'
                ray_tag = f'Rayleigh p={p_str}  R={R_bar:.2f}'
            else:
                ray_tag = 'no sig. windows'

        # Cardinal time labels
        for t_val, t_label in [(0.0, '0 s\n(R-peak)'), (0.2, '0.2 s'),
                               (-0.2, '-0.2 s'), (0.4, '+/-0.4 s')]:
            ang = t_val * np.pi / 0.4
            ax.text(ang, base_r + 0.82, t_label,
                    ha='center', va='center', fontsize=6, color='#555555')

        # Axes limits and title
        ax.set_ylim(0, base_r + 0.9)
        title_line = ch if not ray_tag else f'{ch}\n{ray_tag}'
        ax.set_title(title_line, fontsize=8.5, fontweight='bold', pad=5)

    # Hide unused subplots
    for j in range(len(valid_chs), len(axes_flat)):
        axes_flat[j].set_visible(False)

    # Figure-level legend
    legend_handles = [
        plt.Line2D([0], [0], color=dir_colors['pos'], linewidth=4,
                   label=f'{group_a} > {group_b} (positive diff)'),
        plt.Line2D([0], [0], color=dir_colors['neg'], linewidth=4,
                   label=f'{group_b} > {group_a} (negative diff)'),
        plt.Line2D([0], [0], color='orange', linewidth=8, alpha=0.5,
                   label='Significant window'),
        plt.Line2D([0], [0], color='red', marker='v', markersize=9,
                   linestyle='none', label='R-peak (t = 0)'),
    ]
    fig.legend(handles=legend_handles, loc='lower center',
               bbox_to_anchor=(0.5, -0.01), ncol=2, fontsize=9, frameon=False)

    fig.suptitle(
        f'Circular HEP Summary per Electrode  -  {group_a} vs {group_b}  -  Stage: {selected_stage}\n'
        'Circle = HEP time axis.  Top (0\u00b0) = R-peak.  Bottom = \u00b10.4 s.  '
        'Right half = positive time.  Left half = negative time.',
        fontsize=10, fontweight='bold', y=1.01,
    )
    fig.tight_layout()
    return fig


def run_compare_groups_analysis(base_path, selected_stage):
    """
    Logic for Comparing Groups mode.
    Plots a publication-ready, statistically annotated comparison of the HEP
    (Heartbeat-Evoked Potential) between all available groups, showing:
      1. Group median waveforms + MAD ribbons per group
      2. Difference waveform + cluster-permutation significance + Cohen's d (if 2 groups)
      3. Channel x time heatmap of mean group difference
      4. Summary statistics table
    """
    available_groups = sorted([g for g in os.listdir(base_path) if os.path.isdir(os.path.join(base_path, g))])
    if not available_groups:
        st.error("No groups found in the data directory.")
        return

    # ── Group selection ──────────────────────────────────────────────────────
    default_groups = default_hep_groups(available_groups)
    selected_groups = st.multiselect(
        "Select Groups to Compare",
        options=available_groups,
        default=default_groups,
        key="cmp_selected_groups",
    )
    if not selected_groups:
        st.warning("Please select at least one group.")
        return

    # ── UI controls ─────────────────────────────────────────────────────────
    with st.expander("⚙️ Group Comparison Settings", expanded=True):
        col1, col2, col3, col4, col5, col6, col7 = st.columns(7)
        with col1:
            if st.runtime.exists():
                test_run = st.checkbox("Test Run (5 files/group)", value=False, key="cmp_test_run")
            else:
                test_run = True
        with col2:
            n_permutations = st.slider("Permutations", 50, 500, 200, 50, key="cmp_n_perm")
        with col3:
            jitter_sec = st.number_input("Jitter (s)", 0.01, 0.5, 0.1, 0.05, key="cmp_jitter")
        with col4:
            use_zscore = st.checkbox("Z-score subjects", value=True, key="cmp_zscore",
                                     help="Normalise each subject's trace to zero-mean unit-variance before averaging. Turn off to keep raw µV values.")
        with col5:
            apply_ica = st.checkbox("Apply ICA", value=True, key="cmp_apply_ica",
                                    help="Apply ICA to remove ECG artifact components from EEG channels.")
        with col6:
            parallel_workers = st.number_input(
                "Patient workers",
                min_value=1,
                max_value=max(1, os.cpu_count() or 1),
                value=DEFAULT_PATIENT_WORKERS,
                step=1,
                key="cmp_patient_workers",
                help="Number of patients processed in parallel. Lower this if memory runs out.",
            )
        with col7:
            recompute_cache = st.button("Recompute Cache", key="cmp_recompute_cache",
                                        help="Force reprocessing of all patient data, ignoring disk cache.")
    amp_ylabel = "Amplitude (Z-score)" if use_zscore else "Amplitude (µV)"

    # ── Load globally excluded patients ─────────────────────────────────────
    global_excluded_pids = load_excluded_pids(base_path)

    # ── Load individual HEPs per group ──────────────────────────────────────
    if recompute_cache and hasattr(get_group_individuals, "clear"):
        get_group_individuals.clear()
    group_individuals = {}   # group -> list of individual tuples
    for group in selected_groups:
        with st.spinner(f"Loading {group}…"):
            inds = get_group_individuals(
                group, selected_stage, base_path,
                test_run=test_run, apply_ica=apply_ica,
                recompute_cache=recompute_cache,
                parallel_workers=parallel_workers,
            )
            
        if inds:
            # Filter out globally excluded patients by checking if the base patient ID is in the excluded list
            inds_filtered = filter_excluded(inds, global_excluded_pids)
            
            if inds_filtered:
                group_individuals[group] = inds_filtered
            else:
                st.warning(f"All loaded data for group **{group}** in stage {selected_stage} was globally excluded.")
        else:
            st.warning(f"No valid data for group **{group}** in stage {selected_stage}.")

    if not group_individuals:
        st.error("No valid data found for any group.")
        return

    groups_with_data = list(group_individuals.keys())
    n_groups = len(groups_with_data)
    # unique patient ids per group (each individual tuple is (patient_id, ...))
    _group_n_labels = {
        g: f"{g} (n={len(set(ind[0] for ind in inds))})"
        for g, inds in group_individuals.items()
    }
    groups_label = " vs ".join(_group_n_labels[g] for g in groups_with_data) if len(groups_with_data) == 2 \
        else ", ".join(_group_n_labels[g] for g in groups_with_data)

    def _comparison_title(plot_name, extra=None):
        parts = [plot_name, f"Stage: {selected_stage}", f"Groups: {groups_label}"]
        if extra:
            parts.append(extra)
        return " | ".join(parts)

    def _safe_csv_filename(name):
        safe = re.sub(r"[^A-Za-z0-9_.-]+", "_", str(name)).strip("_")
        return f"{safe}.csv"

    def _append_median_trace_row(rows, plot_name, scope, group, matrix, times_vec):
        if matrix is None or times_vec is None:
            return
        arr = np.asarray(matrix, dtype=float)
        t = np.asarray(times_vec, dtype=float)
        if arr.ndim != 2 or arr.shape[0] == 0 or len(t) == 0:
            return
        min_len = min(arr.shape[1], len(t))
        if min_len == 0:
            return
        median_trace, _ = median_and_mad(arr[:, :min_len], axis=0)
        row = {
            "plot": plot_name,
            "scope": scope,
            "stage": selected_stage,
            "groups": groups_label,
            "group": group,
            "statistic": "median",
            "shaded_band": "MAD",
            "n": arr.shape[0],
            "unit": "Z-score" if use_zscore else "uV",
        }
        for time_val, amp_val in zip(t[:min_len], median_trace):
            row[f"t={time_val:.6f}s"] = amp_val
        rows.append(row)

    def _median_trace_rows_to_csv(rows):
        return pd.DataFrame(rows).to_csv(index=False).encode("utf-8")

    # ── Identify common EEG channels across all groups ───────────────────────
    # Same logic as Per-Channel Analysis: ≥50% of subjects, 1-2 letter prefix
    all_inds_flat = [ind for inds in group_individuals.values() for ind in inds]
    common_channels = identify_common_eeg_channels(all_inds_flat)

    # Build per-group arrays:  (n_subjects, n_times)  averaged across channels
    group_hep_matrix = {}   # group -> np.ndarray (n_subj, n_times)
    group_times = {}
    group_hep_per_channel = {}   # group -> dict{ch: (n_subj, n_times)}

    for group, inds in group_individuals.items():
        subj_mean_heps = []
        subj_ch_heps = {ch: [] for ch in common_channels}
        times_ref = None
        for ind in inds:
            pid, hep_data, times, ch_names, rpeaks, ecg_hep, ecg_ch_names = ind[:7]
            if hep_data is None or times is None:
                continue
            times_ref = times
            # Summarize common channels, guarding against average-reference cancellation.
            valid_ch_indices = [ch_names.index(ch) for ch in common_channels if ch in ch_names]
            if valid_ch_indices:
                subj_mean_heps.append(
                    summarize_channels_without_reference_cancellation(
                        hep_data[valid_ch_indices, :]
                    )
                )
            # Per-channel
            for ch in common_channels:
                if ch in ch_names:
                    ch_idx = ch_names.index(ch)
                    subj_ch_heps[ch].append(hep_data[ch_idx])

        if subj_mean_heps:
            group_hep_matrix[group] = np.array(subj_mean_heps)     # (n_subj, n_times)
            group_times[group] = times_ref
            group_hep_per_channel[group] = {
                ch: np.array(v) for ch, v in subj_ch_heps.items() if v
            }

    if not group_hep_matrix:
        st.error("No valid HEP data could be computed.")
        return

    # ── Optional: Z-score each subject's trace ───────────────────────────────
    for group in list(group_hep_matrix.keys()):
        group_hep_matrix[group] = scale_matrix(group_hep_matrix[group], use_zscore)

        for ch in group_hep_per_channel.get(group, {}):
            ch_mat = group_hep_per_channel[group][ch]  # (n_subj, n_times)
            if ch_mat.ndim == 2 and ch_mat.shape[0] > 0:
                group_hep_per_channel[group][ch] = scale_matrix(ch_mat, use_zscore)

    # Use times from first available group
    times = group_times[groups_with_data[0]]

    # ── Colour palette ───────────────────────────────────────────────────────
    PALETTE = ['#2196F3', '#FF5722', '#4CAF50', '#9C27B0', '#FF9800']
    group_color = {g: PALETTE[i % len(PALETTE)] for i, g in enumerate(groups_with_data)}

    st.markdown("---")
    # ═══════════════════════════════════════════════════════════════════════
    # PLOT 0 — ECG Comparison across groups (average + individual traces)
    # ═══════════════════════════════════════════════════════════════════════
    st.subheader("🫀 ECG Comparison across Groups")

    # Collect per-group ECG traces
    group_ecg_matrix = {}   # group -> list of 1-D arrays (µV)
    ecg_times_ref = None
    for group, inds in group_individuals.items():
        ecg_traces = []
        for ind in inds:
            ecg_hep = ind[5]
            ind_times = ind[2]
            if ecg_hep is not None:
                trace = np.asarray(ecg_hep).squeeze()
                if trace.ndim == 1 and len(trace) > 0:
                    ecg_traces.append(trace * 1e6)   # convert V → µV
                    if ecg_times_ref is None and ind_times is not None:
                        ecg_times_ref = ind_times
        if ecg_traces:
            group_ecg_matrix[group] = ecg_traces

    if group_ecg_matrix and ecg_times_ref is not None:
        fig_ecg_cmp, ax_ecg_cmp = plt.subplots(figsize=(14, 5))
        ax_ecg_cmp.axvline(0, color='red', linestyle='--', alpha=0.6, label='R-peak (t=0)')
        ax_ecg_cmp.axhline(0, color='black', linewidth=0.5, alpha=0.3)

        for group in groups_with_data:
            traces = group_ecg_matrix.get(group)
            if not traces:
                continue
            color = group_color[group]
            # Align all traces to the shortest length
            min_len = min(len(tr) for tr in traces)
            t_ecg = ecg_times_ref[:min_len]
            traces_arr = np.array([tr[:min_len] for tr in traces])   # (n_subj, n_times)
            med_ecg, mad_ecg = median_and_mad(traces_arr, axis=0)

            ax_ecg_cmp.plot(t_ecg, med_ecg, color=color, linewidth=2.5,
                            label=f"{group}  (n={len(traces)})")
            ax_ecg_cmp.fill_between(t_ecg, med_ecg - mad_ecg, med_ecg + mad_ecg,
                                    color=color, alpha=0.22)

        ax_ecg_cmp.set_xlabel("Time relative to R-peak (s)", fontsize=12)
        ax_ecg_cmp.set_ylabel("Amplitude (µV)", fontsize=12)
        ax_ecg_cmp.set_title(
            _comparison_title("ECG Median +/- MAD Comparison"),
            fontsize=14, fontweight='bold'
        )
        ax_ecg_cmp.legend(fontsize=11)
        ax_ecg_cmp.grid(True, alpha=0.25)
        ax_ecg_cmp.set_xlim(ecg_times_ref[0], ecg_times_ref[-1])
        fig_ecg_cmp.tight_layout()
        st.pyplot(fig_ecg_cmp, use_container_width=True)
        plt.close(fig_ecg_cmp)
    else:
        st.info("No ECG data available for cross-group comparison.")

    st.markdown("---")
    # ═══════════════════════════════════════════════════════════════════════
    # PLOT 1b — Group Median + Median Absolute Deviation
    # Shows two figures: before and after ECG regression (matching the
    # "Group HEP Analysis > Per-Channel Analysis > Average" method).
    # ═══════════════════════════════════════════════════════════════════════
    st.subheader("📈 Median HEP per Group")

    # Always use non-ICA individuals for these plots; ECG regression (the
    # same method as Per-Channel Analysis) is applied post-hoc on epochs.
    if not apply_ica:
        group_individuals_for_1b = group_individuals
    else:
        group_individuals_for_1b = {}
        for group in groups_with_data:
            with st.spinner(f"Loading raw (no ICA) data for {group}…"):
                inds_raw = get_group_individuals(
                    group, selected_stage, base_path,
                    test_run=test_run, apply_ica=False, recompute_cache=recompute_cache,
                    parallel_workers=parallel_workers,
                )
            filtered = filter_excluded(inds_raw, global_excluded_pids) if inds_raw else []
            if filtered:
                group_individuals_for_1b[group] = filtered
            elif group in group_individuals:
                # Fall back to the already loaded group data so downstream
                # comparison plots do not crash when raw/no-ICA cache is missing.
                group_individuals_for_1b[group] = group_individuals[group]
                st.info(
                    f"Using loaded data for {group} in the post-cleaned comparison plots "
                    f"because raw/no-ICA individuals were unavailable for {selected_stage}."
                )

    # ── Build (n_subj, n_times) matrices and apply scaling ───────────────
    # apply_ecg_regression=True applies the same ECG linear regression used
    # in Per-Channel Analysis (projects out the ECG epoch from EEG epochs).
    def _build_scaled_matrix(group_inds_dict, clean_ecg=False):
        mat_dict, times_dict = {}, {}
        for grp, inds in group_inds_dict.items():
            subj_heps, times_ref = [], None
            for ind in inds:
                hep_data, ind_times, ch_names = ind[1], ind[2], ind[3]
                ecg_hep = ind[5] if len(ind) > 5 else None
                if hep_data is None or ind_times is None:
                    continue
                times_ref = ind_times
                hep_clean = apply_ecg_regression(hep_data, ecg_hep) if clean_ecg else hep_data.copy().astype(float)
                valid_idx = [ch_names.index(ch) for ch in common_channels if ch in ch_names]
                if valid_idx:
                    subj_heps.append(
                        summarize_channels_without_reference_cancellation(
                            hep_clean[valid_idx, :]
                        )
                    )
            if subj_heps:
                mat_dict[grp] = np.array(subj_heps)
                times_dict[grp] = times_ref
        for grp in list(mat_dict):
            mat_dict[grp] = scale_matrix(mat_dict[grp], use_zscore)
        return mat_dict, times_dict

    group_hep_matrix_raw, group_times_raw = _build_scaled_matrix(
        group_individuals_for_1b, clean_ecg=False
    )
    group_hep_matrix_ica, group_times_ica = _build_scaled_matrix(
        group_individuals_for_1b, clean_ecg=True
    )

    def _build_scaled_side_matrix(group_inds_dict, side_chs, clean_ecg=False):
        mat_dict, times_dict = {}, {}
        for grp, inds in group_inds_dict.items():
            subj_heps, times_ref = [], None
            for ind in inds:
                hep_data, ind_times, ch_names = ind[1], ind[2], ind[3]
                ecg_hep = ind[5] if len(ind) > 5 else None
                if hep_data is None or ind_times is None:
                    continue
                times_ref = ind_times
                hep_clean = apply_ecg_regression(hep_data, ecg_hep) if clean_ecg else hep_data.copy().astype(float)
                valid_idx = [ch_names.index(ch) for ch in side_chs if ch in ch_names]
                if valid_idx:
                    subj_heps.append(
                        summarize_channels_without_reference_cancellation(
                            hep_clean[valid_idx, :]
                        )
                    )
            if subj_heps:
                mat_dict[grp] = scale_matrix(np.array(subj_heps), use_zscore)
                times_dict[grp] = times_ref
        return mat_dict, times_dict

    def _build_scaled_per_channel_matrix(group_inds_dict, clean_ecg=False):
        ch_dict, times_dict = {}, {}
        for grp, inds in group_inds_dict.items():
            subj_ch_heps = {ch: [] for ch in common_channels}
            times_ref = None
            for ind in inds:
                hep_data, ind_times, ch_names = ind[1], ind[2], ind[3]
                ecg_hep = ind[5] if len(ind) > 5 else None
                if hep_data is None or ind_times is None:
                    continue
                times_ref = ind_times
                hep_clean = apply_ecg_regression(hep_data, ecg_hep) if clean_ecg else hep_data.copy().astype(float)
                for ch in common_channels:
                    if ch in ch_names:
                        subj_ch_heps[ch].append(hep_clean[ch_names.index(ch)])
            ch_dict[grp] = {}
            for ch, rows in subj_ch_heps.items():
                if rows:
                    ch_dict[grp][ch] = scale_matrix(np.array(rows), use_zscore)
            if times_ref is not None:
                times_dict[grp] = times_ref
        return ch_dict, times_dict

    group_hep_per_channel_ica, group_times_per_channel_ica = _build_scaled_per_channel_matrix(
        group_individuals_for_1b, clean_ecg=True
    )
    active_group_hep_matrix = group_hep_matrix
    active_group_times = group_times
    active_group_hep_per_channel = group_hep_per_channel
    active_group_times_per_channel = group_times
    groups_with_active_data = [
        g for g in groups_with_data
        if g in active_group_hep_matrix and g in active_group_times
    ]

    # ── Per-selected-group significant topomaps ──────────────────────────
    # These maps intentionally use each group's own electrode set, not the
    # cross-group common_channels list, so groups with different montages keep
    # only their real electrode labels.
    st.markdown("---")
    st.subheader("Significant Brain Maps per Selected Group")
    st.caption(
        "Each figure is computed against zero for one selected group using the same ECG-regressed "
        "per-channel traces and full-epoch minimum cluster p-value logic as the Significant Channels Topomap. "
        "Only electrodes that exist in that group are labeled."
    )

    def _build_group_specific_channel_matrix(inds):
        candidate_channels = identify_common_eeg_channels(inds or [])
        montage = mne.channels.make_standard_montage('standard_1020')
        candidate_channels = [
            ch for ch in candidate_channels
            if _to_montage_channel_name(ch, montage) is not None
            and not any(x in str(ch).lower() for x in ['ecg', 'ekg', 'eog', 'emg', 'resp'])
        ]
        channel_rows = {ch: [] for ch in candidate_channels}
        times_ref = None
        for ind in inds or []:
            hep_data, ind_times, ch_names = ind[1], ind[2], ind[3]
            ecg_hep = ind[5] if len(ind) > 5 else None
            if hep_data is None or ind_times is None or ch_names is None:
                continue
            if times_ref is None:
                times_ref = np.asarray(ind_times, dtype=float)
            hep_clean = apply_ecg_regression(np.asarray(hep_data, dtype=float), ecg_hep)
            for ch in candidate_channels:
                if ch in ch_names:
                    channel_rows[ch].append(np.asarray(hep_clean[ch_names.index(ch)], dtype=float))

        channel_matrix = {}
        for ch, rows in channel_rows.items():
            if len(rows) < 3:
                continue
            min_len = min(len(row) for row in rows)
            if min_len <= 1:
                continue
            channel_matrix[ch] = scale_matrix(
                np.asarray([row[:min_len] for row in rows], dtype=float),
                use_zscore,
            )
        return channel_matrix, times_ref

    with st.spinner("Computing per-group significant brain maps..."):
        group_topomap_results = {}
        for group in groups_with_data:
            group_inds_for_map = group_individuals_for_1b.get(group, group_individuals.get(group, []))
            n_group_patients = len(group_inds_for_map)
            ch_matrix, t_ref = _build_group_specific_channel_matrix(group_inds_for_map)
            if not ch_matrix or t_ref is None:
                group_topomap_results[group] = None
                continue

            pvals = {}
            n_by_ch = {}
            for ch, mat in ch_matrix.items():
                min_len = min(mat.shape[1], len(t_ref))
                t_use = t_ref[:min_len]
                mat_use = mat[:, :min_len]
                if mat_use.shape[0] < 3 or mat_use.shape[1] < 2:
                    continue
                try:
                    sig_windows, _, _ = permutation_cluster_jitter_test(
                        mat_use,
                        t_use,
                        n_permutations=100,
                        p_threshold=0.05,
                        jitter_sec=jitter_sec,
                    )
                    pvals[ch] = min((w['p_value'] for w in sig_windows), default=1.0)
                    n_by_ch[ch] = mat_use.shape[0]
                except Exception:
                    pvals[ch] = 1.0
                    n_by_ch[ch] = mat_use.shape[0]
            group_topomap_results[group] = (pvals, n_by_ch, n_group_patients)

    for group in groups_with_data:
        result = group_topomap_results.get(group)
        if not result:
            st.info(f"{group}: no montage-compatible common electrodes with at least 3 subjects.")
            continue
        pvals, n_by_ch, n_group_patients = result
        if not pvals:
            st.info(f"{group}: no electrodes available for full-epoch cluster testing.")
            continue

        sig_channels = [ch for ch, p_val in pvals.items() if p_val < 0.05]
        title = (
            f"{group} | {selected_stage} | full epoch minimum cluster p-value | "
            f"{len(sig_channels)}/{len(pvals)} sig. electrodes | N={n_group_patients}"
        )
        fig_group_topo, ax_group_topo = plt.subplots(figsize=(5.5, 4.5))
        ok_group_topo = _render_existing_channel_pval_topomap(
            pvals,
            ax_group_topo,
            title,
            p_threshold=0.05,
            show_names=True,
        )
        if ok_group_topo is not None:
            fig_group_topo.tight_layout()
            st.pyplot(fig_group_topo, use_container_width=False)
        else:
            st.warning(f"{group}: could not match channels to standard 10-20 montage.")
        plt.close(fig_group_topo)

        if sig_channels:
            sig_rows = [
                {
                    "Electrode": ch,
                    "N": n_by_ch.get(ch, ""),
                    "p-value": "<0.001" if pvals[ch] < 0.001 else f"{pvals[ch]:.3f}",
                }
                for ch in sorted(sig_channels, key=lambda c: pvals[c])
            ]
            with st.expander(f"{group}: significant electrodes", expanded=False):
                st.dataframe(pd.DataFrame(sig_rows), use_container_width=True)
        else:
            st.info(f"{group}: no significant electrodes across the full epoch (p < 0.05).")

    # ── Render helper ────────────────────────────────────────────────────
    def _render_spaghetti(mat_dict, times_dict, title_suffix, mark_ica_cleaned=True):
        fig, ax = plt.subplots(figsize=(14, 5))
        fig._skip_ica_cleaned_title = not mark_ica_cleaned
        ax.axvline(0, color='red', linestyle='--', alpha=0.6, label='R-peak (t=0)')
        ax.axhline(0, color='black', linewidth=0.5, alpha=0.3)
        for grp in groups_with_data:
            if grp not in mat_dict:
                continue
            mat = mat_dict[grp]
            t = times_dict[grp]
            n_subj = mat.shape[0]
            median_hep, mad_hep = median_and_mad(mat, axis=0)
            color = group_color[grp]
            ax.plot(t, median_hep, color=color, linewidth=2.5, label=f"{grp}  (n={n_subj})")
            ax.fill_between(t, median_hep - mad_hep, median_hep + mad_hep,
                            color=color, alpha=0.20)
        ax.set_xlabel("Time relative to R-peak (s)", fontsize=12)
        ax.set_ylabel(amp_ylabel, fontsize=12)
        ax.set_title(
            _comparison_title("HEP Median +/- MAD", title_suffix),
            fontsize=14, fontweight='bold',
        )
        ax.legend(fontsize=11)
        ax.grid(True, alpha=0.25)
        if times_dict:
            first_grp = next((g for g in groups_with_data if g in times_dict), None)
            if first_grp is not None:
                t0 = times_dict.get(first_grp)
                if t0 is not None:
                    ax.set_xlim(t0[0], t0[-1])
        fig.tight_layout()
        st.pyplot(fig, use_container_width=True)
        plt.close(fig)

    _render_spaghetti(group_hep_matrix_raw, group_times_raw, "(Raw / loaded data)", mark_ica_cleaned=False)
    _render_spaghetti(group_hep_matrix_ica, group_times_ica, "(ECG regression comparison)", mark_ica_cleaned=True)

    # ═══════════════════════════════════════════════════════════════════════
    # PLOT 1c — Hemisphere Comparison: Left vs Right electrodes
    # Left  = channels ending in odd digit  (e.g. F3, C3, P3, O1)
    # Right = channels ending in even digit (e.g. F4, C4, P4, O2)
    # ═══════════════════════════════════════════════════════════════════════
    if group_hep_per_channel and common_channels:
        st.subheader("🧠 Hemisphere Comparison: Left vs Right Electrodes")

        left_chs, right_chs, _ = get_hemisphere_channels(common_channels)

        if left_chs and right_chs:
            fig1c, (ax_L, ax_R) = plt.subplots(1, 2, figsize=(16, 5), sharey=True)

            for ax, side_chs, side_label in [
                (ax_L, left_chs,  "Left hemisphere"),
                (ax_R, right_chs, "Right hemisphere"),
            ]:
                side_mat_dict, side_times_dict = _build_scaled_side_matrix(
                    group_individuals_for_1b, side_chs, clean_ecg=apply_ica
                )
                ax.axvline(0, color='red', linestyle='--', alpha=0.6)
                ax.axhline(0, color='black', linewidth=0.5, alpha=0.3)

                for group in groups_with_data:
                    if group not in side_mat_dict:
                        continue
                    side_avg = side_mat_dict[group]
                    t = side_times_dict[group]
                    n_subj = side_avg.shape[0]
                    grand, mad = median_and_mad(side_avg, axis=0)
                    color = group_color[group]
                    ax.plot(t, grand, color=color, linewidth=2.5,
                            label=f"{group}  (n={n_subj})")
                    ax.fill_between(t, grand - mad, grand + mad,
                                    color=color, alpha=0.18)

                ax.set_title(f"{side_label} - Median +/- MAD\n({', '.join(side_chs)})",
                             fontsize=12, fontweight='bold')
                ax.set_xlabel("Time relative to R-peak (s)", fontsize=11)
                ax.grid(True, alpha=0.25)
                ax.legend(fontsize=10)
                if active_group_times:
                    first_grp = next((g for g in groups_with_active_data if g in active_group_times), None)
                    if first_grp is not None:
                        t0 = active_group_times.get(first_grp)
                        if t0 is not None:
                            ax.set_xlim(t0[0], t0[-1])

            ax_L.set_ylabel(amp_ylabel, fontsize=11)
            fig1c.suptitle(
                _comparison_title("Left vs Right Hemisphere HEP Median +/- MAD"),
                fontsize=14, fontweight='bold'
            )
            fig1c.tight_layout()
            st.pyplot(fig1c, use_container_width=True)
            plt.close(fig1c)
        else:
            st.info("Not enough lateralised channels to plot a hemisphere comparison "
                    f"(left: {left_chs}, right: {right_chs}).")

    # ═══════════════════════════════════════════════════════════════════════
    # PLOT 2 — Difference + permutation significance + Cohen's d
    # (only when exactly 2 groups are available)
    # ═══════════════════════════════════════════════════════════════════════
    sig_windows_global = []
    if n_groups == 2 and len(groups_with_active_data) == 2:
        g_a, g_b = groups_with_active_data[0], groups_with_active_data[1]
        hep_a = active_group_hep_matrix[g_a]
        hep_b = active_group_hep_matrix[g_b]
        t_a = active_group_times[g_a]
        t_b = active_group_times[g_b]

        # Align lengths
        min_len = min(hep_a.shape[1], hep_b.shape[1])
        hep_a = hep_a[:, :min_len]
        hep_b = hep_b[:, :min_len]
        t_common = t_a[:min_len]

        st.subheader(f"📊 Statistical Comparison: {g_a}  vs  {g_b}")
        with st.spinner(f"Running cluster permutation test ({n_permutations} permutations)…"):
            sig_windows_global, t_stat, cohens_d = permutation_two_group_cluster_test(
                hep_a, hep_b, t_common,
                n_permutations=n_permutations,
                p_threshold=0.05,
                jitter_sec=jitter_sec,
                channel_label="Average",
            )

        # ── Build figure with 2 stacked panels ──────────────────────────────
        fig2, (ax_diff, ax_d) = plt.subplots(
            2, 1, figsize=(14, 8),
            gridspec_kw={'height_ratios': [3, 1.5]},
            sharex=True
        )

        # -- Top panel: group medians + difference curve --------------------
        med_a, mad_a = median_and_mad(hep_a, axis=0)
        med_b, mad_b = median_and_mad(hep_b, axis=0)
        diff = med_a - med_b

        ax_diff.plot(t_common, med_a, color=group_color[g_a], linewidth=2, label=f"{g_a} median")
        ax_diff.fill_between(t_common, med_a - mad_a, med_a + mad_a, color=group_color[g_a], alpha=0.15)
        ax_diff.plot(t_common, med_b, color=group_color[g_b], linewidth=2, label=f"{g_b} median")
        ax_diff.fill_between(t_common, med_b - mad_b, med_b + mad_b, color=group_color[g_b], alpha=0.15)
        ax_diff.plot(t_common, diff, color='black', linewidth=1.5, linestyle='--', label=f"{g_a}−{g_b} diff")
        ax_diff.axhline(0, color='black', linewidth=0.4, alpha=0.3)
        ax_diff.axvline(0, color='red', linestyle='--', alpha=0.5)

        # Significance spans + asterisks
        for win in sig_windows_global:
            ax_diff.axvspan(win['start'], win['end'], color='orange', alpha=0.25)
            ax_d.axvspan(win['start'], win['end'], color='orange', alpha=0.25)
            mid = (win['start'] + win['end']) / 2
            p_val = win['p_value']
            stars = '***' if p_val < 0.001 else ('**' if p_val < 0.01 else '*')
            ymax = ax_diff.get_ylim()[1]
            # p-value label
            p_txt = f"p={'<0.001' if p_val < 0.001 else f'{p_val:.3f}'}"
            ax_diff.text(mid, ymax * 0.95, stars, ha='center', va='top',
                         fontsize=16, color='darkorange', fontweight='bold')
            ax_diff.text(mid, ymax * 0.82, p_txt, ha='center', va='top',
                         fontsize=9, color='darkorange')

        ax_diff.set_ylabel(amp_ylabel, fontsize=11)
        if sig_windows_global:
            min_p = min(w['p_value'] for w in sig_windows_global)
            p_tag = 'p<0.001' if min_p < 0.001 else f'p={min_p:.3f}'
        else:
            p_tag = 'p=n.s.'
        ax_diff.set_title(
            _comparison_title(
                "Average HEP Group Difference - Median +/- MAD",
                f"n={hep_a.shape[0]}+{hep_b.shape[0]} | {p_tag} | permutations={n_permutations}",
            ),
            fontsize=13, fontweight='bold'
        )
        ax_diff.legend(fontsize=10)
        ax_diff.grid(True, alpha=0.2)

        # -- Bottom panel: Cohen's d ----------------------------------------
        ax_d.plot(t_common, cohens_d, color='#5c35cc', linewidth=2, label="Cohen's d (signed)")
        ax_d.axhline(0,    color='black', linewidth=0.5, alpha=0.3)
        ax_d.axhline(0.5,  color='gray',  linewidth=0.8, linestyle=':', alpha=0.6, label='d=0.5 (medium)')
        ax_d.axhline(-0.5, color='gray',  linewidth=0.8, linestyle=':', alpha=0.6)
        ax_d.axhline(0.8,  color='gray',  linewidth=0.8, linestyle='--', alpha=0.5, label='d=0.8 (large)')
        ax_d.axhline(-0.8, color='gray',  linewidth=0.8, linestyle='--', alpha=0.5)
        ax_d.axvline(0, color='red', linestyle='--', alpha=0.5)
        ax_d.fill_between(t_common, 0, cohens_d,
                          where=(cohens_d > 0), color='#2196F3', alpha=0.2, interpolate=True)
        ax_d.fill_between(t_common, 0, cohens_d,
                          where=(cohens_d < 0), color='#FF5722', alpha=0.2, interpolate=True)
        ax_d.set_xlabel("Time relative to R-peak (s)", fontsize=11)
        ax_d.set_ylabel("Cohen's d", fontsize=11)
        ax_d.legend(fontsize=9, loc='upper right')
        ax_d.grid(True, alpha=0.2)
        if t_common is not None:
            ax_d.set_xlim(t_common[0], t_common[-1])

        fig2.tight_layout()
        st.pyplot(fig2, use_container_width=True)
        plt.close(fig2)

        # ── Significant windows table ────────────────────────────────────────
        if sig_windows_global:
            st.markdown("**Significant time windows (cluster permutation test):**")
            rows = []
            for win in sig_windows_global:
                p_val = win['p_value']
                rows.append({
                    "Start (s)": f"{win['start']:.3f}",
                    "End (s)": f"{win['end']:.3f}",
                    "Duration (ms)": f"{(win['end']-win['start'])*1000:.1f}",
                    "p-value": "<0.001" if p_val < 0.001 else f"{p_val:.3f}",
                    "Direction": win.get('direction', ''),
                })
            st.dataframe(pd.DataFrame(rows), use_container_width=True)
        else:
            st.info("No significant time windows found (cluster permutation test, p > 0.05).")


        # ── Hemisphere statistical comparison ───────────────────────────────
        st.markdown("---")
        st.subheader(f"🧠 Hemisphere Statistical Comparison: {g_a} vs {g_b}")
        with st.expander("ℹ️ About Hemisphere Comparison", expanded=False):
            st.markdown(
                "Cluster-permutation test run separately on **left** (odd-numbered) "
                "and **right** (even-numbered) electrodes, averaged within each hemisphere."
            )

        if active_group_hep_per_channel and common_channels:
            left_chs_stat, right_chs_stat, _ = get_hemisphere_channels(common_channels)

            def _build_hemisphere_matrix(group, side_chs):
                """Return (n_subj, n_times) averaged over side_chs for a group."""
                side_mat_dict, _ = _build_scaled_side_matrix(
                    group_individuals_for_1b, side_chs, clean_ecg=apply_ica
                )
                if group not in side_mat_dict:
                    return None
                return side_mat_dict[group]

            for side_label, side_chs in [("Left hemisphere", left_chs_stat), ("Right hemisphere", right_chs_stat)]:
                if not side_chs:
                    st.info(f"No {side_label} channels found.")
                    continue

                hem_a = _build_hemisphere_matrix(g_a, side_chs)
                hem_b = _build_hemisphere_matrix(g_b, side_chs)

                if hem_a is None or hem_b is None:
                    st.warning(f"Not enough data for {side_label} comparison.")
                    continue

                # Align time lengths
                min_len_h = min(hem_a.shape[1], hem_b.shape[1], len(t_common))
                hem_a = hem_a[:, :min_len_h]
                hem_b = hem_b[:, :min_len_h]
                t_h = t_common[:min_len_h]

                st.markdown(f"**{side_label}** ({', '.join(side_chs)})")
                with st.spinner(f"Running permutation test for {side_label}…"):
                    sig_hem, t_hem_stat, cd_hem = permutation_two_group_cluster_test(
                        hem_a, hem_b, t_h,
                        n_permutations=n_permutations,
                        p_threshold=0.05,
                        jitter_sec=jitter_sec,
                        channel_label=side_label,
                    )

                fig_hem, (ax_hem_diff, ax_hem_d) = plt.subplots(
                    2, 1, figsize=(14, 7),
                    gridspec_kw={'height_ratios': [3, 1.5]},
                    sharex=True
                )

                med_ha, mad_ha = median_and_mad(hem_a, axis=0)
                med_hb, mad_hb = median_and_mad(hem_b, axis=0)
                diff_h = med_ha - med_hb

                ax_hem_diff.axvline(0, color='red', linestyle='--', alpha=0.5)
                ax_hem_diff.axhline(0, color='black', linewidth=0.4, alpha=0.3)
                ax_hem_diff.plot(t_h, med_ha, color=group_color[g_a], linewidth=2, label=f"{g_a} median")
                ax_hem_diff.fill_between(t_h, med_ha - mad_ha, med_ha + mad_ha, color=group_color[g_a], alpha=0.15)
                ax_hem_diff.plot(t_h, med_hb, color=group_color[g_b], linewidth=2, label=f"{g_b} median")
                ax_hem_diff.fill_between(t_h, med_hb - mad_hb, med_hb + mad_hb, color=group_color[g_b], alpha=0.15)
                ax_hem_diff.plot(t_h, diff_h, color='black', linewidth=1.5, linestyle='--', label=f"{g_a}−{g_b} diff")

                for win in sig_hem:
                    ax_hem_diff.axvspan(win['start'], win['end'], color='orange', alpha=0.25)
                    ax_hem_d.axvspan(win['start'], win['end'], color='orange', alpha=0.25)
                    mid = (win['start'] + win['end']) / 2
                    p_v = win['p_value']
                    stars = '***' if p_v < 0.001 else ('**' if p_v < 0.01 else '*')
                    ymax_h = ax_hem_diff.get_ylim()[1]
                    p_txt_h = f"p={'<0.001' if p_v < 0.001 else f'{p_v:.3f}'}"
                    ax_hem_diff.text(mid, ymax_h * 0.95, stars, ha='center', va='top',
                                     fontsize=16, color='darkorange', fontweight='bold')
                    ax_hem_diff.text(mid, ymax_h * 0.82, p_txt_h, ha='center', va='top',
                                     fontsize=9, color='darkorange')

                p_tag_h = ('p<0.001' if min(w['p_value'] for w in sig_hem) < 0.001
                           else f"p={min(w['p_value'] for w in sig_hem):.3f}") if sig_hem else 'p=n.s.'
                ax_hem_diff.set_title(
                    _comparison_title(
                        f"{side_label} Group Difference - Median +/- MAD",
                        f"n={hem_a.shape[0]}+{hem_b.shape[0]} | {p_tag_h}",
                    ),
                    fontsize=13, fontweight='bold'
                )
                ax_hem_diff.set_ylabel(amp_ylabel, fontsize=11)
                ax_hem_diff.legend(fontsize=10)
                ax_hem_diff.grid(True, alpha=0.2)

                ax_hem_d.plot(t_h, cd_hem, color='#5c35cc', linewidth=2, label="Cohen's d")
                ax_hem_d.axhline(0,    color='black', linewidth=0.5, alpha=0.3)
                ax_hem_d.axhline( 0.5, color='gray', linewidth=0.8, linestyle=':', alpha=0.6, label='d=0.5')
                ax_hem_d.axhline(-0.5, color='gray', linewidth=0.8, linestyle=':', alpha=0.6)
                ax_hem_d.axhline( 0.8, color='gray', linewidth=0.8, linestyle='--', alpha=0.5, label='d=0.8')
                ax_hem_d.axhline(-0.8, color='gray', linewidth=0.8, linestyle='--', alpha=0.5)
                ax_hem_d.axvline(0, color='red', linestyle='--', alpha=0.5)
                ax_hem_d.fill_between(t_h, 0, cd_hem, where=(cd_hem > 0), color='#2196F3', alpha=0.2, interpolate=True)
                ax_hem_d.fill_between(t_h, 0, cd_hem, where=(cd_hem < 0), color='#FF5722', alpha=0.2, interpolate=True)
                ax_hem_d.set_xlabel("Time relative to R-peak (s)", fontsize=11)
                ax_hem_d.set_ylabel("Cohen's d", fontsize=11)
                ax_hem_d.legend(fontsize=9, loc='upper right')
                ax_hem_d.grid(True, alpha=0.2)
                ax_hem_d.set_xlim(t_h[0], t_h[-1])

                fig_hem.tight_layout()
                st.pyplot(fig_hem, use_container_width=True)
                plt.close(fig_hem)

                side_key = side_label.lower().replace(" ", "_")
                hem_csv_rows = []
                _append_median_trace_row(
                    hem_csv_rows,
                    "Hemisphere Group Difference",
                    side_label,
                    g_a,
                    hem_a,
                    t_h,
                )
                _append_median_trace_row(
                    hem_csv_rows,
                    "Hemisphere Group Difference",
                    side_label,
                    g_b,
                    hem_b,
                    t_h,
                )
                if hem_csv_rows:
                    st.download_button(
                        label=f"Download {side_label} median traces CSV",
                        data=_median_trace_rows_to_csv(hem_csv_rows),
                        file_name=_safe_csv_filename(
                            f"hemisphere_group_difference_{selected_stage}_{side_key}_{g_a}_vs_{g_b}"
                        ),
                        mime="text/csv",
                        key=f"download_hemisphere_median_{selected_stage}_{side_key}_{g_a}_{g_b}",
                    )

                if sig_hem:
                    rows_h = []
                    for win in sig_hem:
                        p_v = win['p_value']
                        rows_h.append({
                            "Start (s)": f"{win['start']:.3f}",
                            "End (s)": f"{win['end']:.3f}",
                            "Duration (ms)": f"{(win['end']-win['start'])*1000:.1f}",
                            "p-value": "<0.001" if p_v < 0.001 else f"{p_v:.3f}",
                            "Direction": win.get('direction', ''),
                        })
                    st.dataframe(pd.DataFrame(rows_h), use_container_width=True)
                else:
                    st.info(f"No significant windows in {side_label} (p > 0.05).")
        else:
            st.info("Per-hemisphere analysis requires per-channel data.")

        # ── Per-electrode statistical comparison ────────────────────────────
        st.markdown("---")
        st.subheader(f"📡 Per-Electrode Statistical Comparison: {g_a} vs {g_b}")
        with st.expander("ℹ️ About Per-Electrode Comparison", expanded=False):
            st.markdown(
                "Each subplot shows group medians ± MAD for one EEG channel. "
                "Orange spans mark cluster-permutation significant windows."
            )

        if active_group_hep_per_channel and common_channels:
            n_cols_el = 4
            n_rows_el = int(np.ceil(len(common_channels) / n_cols_el))
            fig_el, axes_el = plt.subplots(
                n_rows_el, n_cols_el,
                figsize=(5 * n_cols_el, 3.5 * n_rows_el),
                sharex=True
            )
            axes_el_flat = np.array(axes_el).flatten()

            per_ch_save = {}  # collect ca/cb per channel for download
            circular_channel_stats = {}
            sig_windows_per_ch = {}

            for idx_ch, ch in enumerate(common_channels):
                ax_el = axes_el_flat[idx_ch]
                ch_a = active_group_hep_per_channel.get(g_a, {}).get(ch)
                ch_b = active_group_hep_per_channel.get(g_b, {}).get(ch)

                if ch_a is None or ch_b is None or ch_a.shape[0] == 0 or ch_b.shape[0] == 0:
                    ax_el.set_title(ch, fontsize=9)
                    ax_el.text(0.5, 0.5, 'No data', ha='center', va='center',
                               transform=ax_el.transAxes, fontsize=8, color='gray')
                    ax_el.axis('off')
                    continue

                t_a_el = active_group_times_per_channel.get(g_a)
                t_b_el = active_group_times_per_channel.get(g_b)
                if t_a_el is None or t_b_el is None:
                    ax_el.set_title(ch, fontsize=9)
                    ax_el.text(0.5, 0.5, 'No timebase', ha='center', va='center',
                               transform=ax_el.transAxes, fontsize=8, color='gray')
                    ax_el.axis('off')
                    continue

                min_t_el = min(ch_a.shape[1], ch_b.shape[1], len(t_a_el), len(t_b_el))
                ca = ch_a[:, :min_t_el]
                cb = ch_b[:, :min_t_el]
                t_el = t_a_el[:min_t_el]

                per_ch_save[ch] = (ca, cb, t_el)

                med_ca, mad_ca = median_and_mad(ca, axis=0)
                med_cb, mad_cb = median_and_mad(cb, axis=0)

                ax_el.axvline(0, color='red', linestyle='--', alpha=0.5, linewidth=0.8)
                ax_el.axhline(0, color='black', linewidth=0.3, alpha=0.3)
                ax_el.plot(t_el, med_ca, color=group_color[g_a], linewidth=1.5,
                           label=f"{g_a} (n={ca.shape[0]})")
                ax_el.fill_between(t_el, med_ca - mad_ca, med_ca + mad_ca,
                                   color=group_color[g_a], alpha=0.18)
                ax_el.plot(t_el, med_cb, color=group_color[g_b], linewidth=1.5,
                           label=f"{g_b} (n={cb.shape[0]})")
                ax_el.fill_between(t_el, med_cb - mad_cb, med_cb + mad_cb,
                                   color=group_color[g_b], alpha=0.18)

                # Quick permutation test for this channel (fewer perms for speed)
                try:
                    sig_el, _, cd_el = permutation_two_group_cluster_test(
                        ca, cb, t_el,
                        n_permutations=max(50, n_permutations // 4),
                        p_threshold=0.05,
                        jitter_sec=jitter_sec,
                        channel_label=ch,
                    )
                    for win in sig_el:
                        ax_el.axvspan(win['start'], win['end'], color='orange', alpha=0.28)
                    if sig_el:

                        min_p_el = min(w['p_value'] for w in sig_el)
                        p_tag_el = 'p<0.001' if min_p_el < 0.001 else f'p={min_p_el:.3f}'
                    else:
                        p_tag_el = 'n.s.'
                    circular_channel_stats[ch] = _summarize_channel_group_difference(sig_el, cd_el, t_el)
                    sig_windows_per_ch[ch] = sig_el
                except Exception:
                    p_tag_el = ''
                    circular_channel_stats[ch] = {
                        'effect_size': 0.0,
                        'signed_effect': 0.0,
                        'total_sig_duration': 0.0,
                        'min_p': 1.0,
                        'n_windows': 0,
                    }
                    sig_windows_per_ch[ch] = []

                ax_el.set_title(f"{ch}  [{p_tag_el}]", fontsize=9, fontweight='bold')
                ax_el.set_xlim(t_el[0], t_el[-1])
                ax_el.grid(True, alpha=0.2)
                ax_el.tick_params(labelsize=7)
                if idx_ch % n_cols_el == 0:
                    ax_el.set_ylabel(amp_ylabel, fontsize=8)

            # Hide unused axes
            for j in range(len(common_channels), len(axes_el_flat)):
                axes_el_flat[j].axis('off')

            # Determine valid N for each group
            n_a = len(group_individuals.get(g_a, []))
            n_b = len(group_individuals.get(g_b, []))

            # Single legend outside
            handles_el = [
                plt.Line2D([0], [0], color=group_color[g_a], linewidth=2, label=f"{g_a} median (N={n_a})"),
                plt.Line2D([0], [0], color=group_color[g_b], linewidth=2, label=f"{g_b} median (N={n_b})"),
            ]
            fig_el.legend(handles=handles_el, loc='upper right', fontsize=10,
                          bbox_to_anchor=(1.0, 1.0))
            fig_el.suptitle(
                _comparison_title("Per-Electrode HEP Comparison | Statistic: Median | Shaded band: MAD"),
                fontsize=14, fontweight='bold', y=1.01
            )
            fig_el.tight_layout()
            st.pyplot(fig_el, use_container_width=True)
            plt.close(fig_el)

            per_electrode_csv_rows = []
            for ch, (ca, cb, t_el) in per_ch_save.items():
                _append_median_trace_row(
                    per_electrode_csv_rows,
                    "Per-Electrode HEP Comparison",
                    ch,
                    g_a,
                    ca,
                    t_el,
                )
                _append_median_trace_row(
                    per_electrode_csv_rows,
                    "Per-Electrode HEP Comparison",
                    ch,
                    g_b,
                    cb,
                    t_el,
                )
            if per_electrode_csv_rows:
                st.download_button(
                    label="Download per-electrode median traces CSV",
                    data=_median_trace_rows_to_csv(per_electrode_csv_rows),
                    file_name=_safe_csv_filename(
                        f"per_electrode_hep_comparison_{selected_stage}_{g_a}_vs_{g_b}"
                    ),
                    mime="text/csv",
                    key=f"download_per_electrode_median_{selected_stage}_{g_a}_{g_b}",
                )

            fig_circular = _plot_per_electrode_circular_summary(
                circular_channel_stats, g_a, g_b, selected_stage,
                channel_data=per_ch_save,
                sig_windows_per_ch=sig_windows_per_ch,
            )
            if fig_circular is not None:
                st.pyplot(fig_circular, use_container_width=True)
                with st.expander("ℹ️ About Circular HEP Summary", expanded=False):
                    st.markdown(
                        "**Circular HEP summary:** Each subplot is one electrode. "
                        "The circle maps HEP time — R-peak at right (0°), ±0.4 s at left (180°). "
                        "Top half = positive time, bottom half = negative time. "
                        "Radial fill shows mean group difference. "
                        "Orange bands = significant windows (cluster-permutation test). "
                        "The Rayleigh test checks whether significant windows cluster at a particular HEP phase."
                    )
                plt.close(fig_circular)
            else:
                st.info("No significant per-electrode differences were available for the circular summary plot.")

        else:
            st.info("Per-electrode plots require per-channel data.")

    elif n_groups == 2:
        missing_clean = [g for g in groups_with_data if g not in groups_with_active_data]
        st.warning(
            "Skipping statistical comparison because comparison data was not "
            f"available for: {', '.join(missing_clean)}."
        )
    elif n_groups > 2:
        st.info(
            f"Significance testing is shown for exactly 2 groups. "
            f"You have {n_groups} groups loaded. Select a subset or compare pairs manually."
        )

    # ═══════════════════════════════════════════════════════════════════════
    # PLOT 3 — Channel × Time heatmap of mean difference (2 groups only)
    # ═══════════════════════════════════════════════════════════════════════
    if n_groups == 2 and common_channels:
        g_a, g_b = groups_with_data[0], groups_with_data[1]
        st.subheader(f"🗺️ Topographic Difference Heatmap: {g_a} − {g_b}  (per EEG channel)")

        diff_matrix = []   # (n_channels, n_times)
        valid_chs = []
        for ch in common_channels:
            ch_a = active_group_hep_per_channel.get(g_a, {}).get(ch)
            ch_b = active_group_hep_per_channel.get(g_b, {}).get(ch)
            if ch_a is not None and ch_b is not None and len(ch_a) > 0 and len(ch_b) > 0:
                min_t = min(ch_a.shape[1], ch_b.shape[1])
                diff_row = np.nanmedian(ch_a[:, :min_t], 0) - np.nanmedian(ch_b[:, :min_t], 0)
                diff_matrix.append(diff_row)
                valid_chs.append(ch)

        if diff_matrix:
            diff_arr = np.array(diff_matrix)   # (n_ch, n_times)
            t_hm = times[:diff_arr.shape[1]]

            fig3, ax3 = plt.subplots(figsize=(14, max(4, len(valid_chs) * 0.35 + 2)))
            vmax = np.nanpercentile(np.abs(diff_arr), 98)

            im = ax3.imshow(
                diff_arr,
                aspect='auto',
                origin='lower',
                cmap='RdBu_r',
                vmin=-vmax, vmax=vmax,
                extent=[t_hm[0], t_hm[-1], -0.5, len(valid_chs) - 0.5]
            )
            ax3.set_yticks(range(len(valid_chs)))
            ax3.set_yticklabels(valid_chs, fontsize=8)
            ax3.axvline(0, color='black', linestyle='--', alpha=0.7, linewidth=1.5)

            # Mark significant windows on heatmap
            for win in sig_windows_global:
                ax3.axvspan(win['start'], win['end'], color='yellow', alpha=0.18)

            cbar = fig3.colorbar(im, ax=ax3, pad=0.01, fraction=0.025)
            cbar.set_label(f"Median diff {'Z-score' if use_zscore else 'µV'} ({g_a}−{g_b})", fontsize=10)

            ax3.set_xlabel("Time relative to R-peak (s)", fontsize=11)
            ax3.set_ylabel("EEG Channel", fontsize=11)
            ax3.set_title(
                _comparison_title("Median HEP Difference per Channel - Median +/- MAD"),
                fontsize=13, fontweight='bold'
            )
            fig3.tight_layout()
            st.pyplot(fig3, use_container_width=True)
            plt.close(fig3)
        else:
            st.warning("No common channels with data in both groups for the heatmap.")

    # ═══════════════════════════════════════════════════════════════════════
    # PLOT 4 — Spatial Topomaps of Difference & P-value
    # ═══════════════════════════════════════════════════════════════════════
    if n_groups == 2 and common_channels:
        st.divider()
        st.subheader("Spatial Topomaps - Group Difference & Significance")
        st.caption(f"Select a time window to visualize the {g_a} − {g_b} average difference distribution and spatial significance across the scalp.")
        
        col_t1, col_t2 = st.columns(2)
        with col_t1:
            topo_diff_tmin = st.number_input("Start Time (ms)", min_value=-500, max_value=1000, value=200, step=10, key="topo_diff_tmin_compare")
        with col_t2:
            topo_diff_tmax = st.number_input("End Time (ms)", min_value=-500, max_value=1000, value=400, step=10, key="topo_diff_tmax_compare")
            
        if topo_diff_tmin >= topo_diff_tmax:
            st.warning("Start time must be less than end time.")
        else:
            try:
                montage = mne.channels.make_standard_montage('standard_1020')
                montage_ch_names_upper = [ch.upper() for ch in montage.ch_names]
                
                times_ms = times * 1000
                t_mask = (times_ms >= topo_diff_tmin) & (times_ms <= topo_diff_tmax)
                
                if not np.any(t_mask):
                    st.warning("Invalid time window selected (no data points).")
                else:
                    g_a, g_b = groups_with_data[0], groups_with_data[1]
                    hep_a = active_group_hep_matrix[g_a]
                    hep_b = active_group_hep_matrix[g_b]
                    
                    min_len = min(hep_a.shape[1], hep_b.shape[1])
                    t_mask_min = t_mask[:min_len]
                    
                    if not np.any(t_mask_min):
                        st.warning("Invalid time window after alignment.")
                    else:
                        
                        plot_ch_names = []
                        plot_data_diff = []
                        plot_data_pval_diff = []

                        plot_data_amp_a = []
                        plot_data_amp_b = []
                        plot_data_pval_a = []
                        plot_data_pval_b = []

                        t_win = times[:min_len][t_mask_min]

                        n_a_subj = active_group_hep_matrix[g_a].shape[0]
                        n_b_subj = active_group_hep_matrix[g_b].shape[0]
                        pmod_scores_a = np.zeros(n_a_subj)
                        pmod_scores_b = np.zeros(n_b_subj)
                        pmod_ch_count = 0
                        _pmod_per_ch_a = []
                        _pmod_per_ch_b = []

                        for i, ch in enumerate(common_channels):
                            ch_upper = ch.upper()
                            if ch_upper in montage_ch_names_upper:
                                m_idx = montage_ch_names_upper.index(ch_upper)
                                standard_name = montage.ch_names[m_idx]
                                
                                # Extract channel data
                                ch_a_data = active_group_hep_per_channel.get(g_a, {}).get(ch)
                                ch_b_data = active_group_hep_per_channel.get(g_b, {}).get(ch)
                                
                                if ch_a_data is not None and ch_b_data is not None:
                                    ch_a_full = ch_a_data[:, :min_len]
                                    ch_b_full = ch_b_data[:, :min_len]
                                    t_full = times[:min_len]
                                    ch_a_win = ch_a_full[:, t_mask_min]
                                    ch_b_win = ch_b_full[:, t_mask_min]
                                    
                                    # Median amplitudes in window. Data is already scaled
                                    # to either Z-score or microvolts upstream.
                                    amp_a = np.nanmedian(ch_a_win)
                                    amp_b = np.nanmedian(ch_b_win)
                                    diff_median = amp_a - amp_b
                                    
                                    # P-value calculation for Difference (A vs B)
                                    sig_windows_diff, _, _ = permutation_two_group_cluster_test(
                                        ch_a_win, ch_b_win, t_win,
                                        n_permutations=n_permutations,
                                        p_threshold=0.05,
                                        jitter_sec=jitter_sec,
                                        channel_label=ch,
                                    )
                                    p_val_diff = min([w['p_value'] for w in sig_windows_diff]) if sig_windows_diff else 1.0

                                    # P-value calculation for Group A (vs 0)
                                    sig_windows_a, _, per_pt_a = permutation_cluster_jitter_test(
                                        ch_a_win, t_win,
                                        n_permutations=n_permutations,
                                        p_threshold=0.05,
                                        jitter_sec=jitter_sec,
                                    )
                                    p_val_a = min([w['p_value'] for w in sig_windows_a]) if sig_windows_a else 1.0

                                    # P-value calculation for Group B (vs 0)
                                    sig_windows_b, _, per_pt_b = permutation_cluster_jitter_test(
                                        ch_b_win, t_win,
                                        n_permutations=n_permutations,
                                        p_threshold=0.05,
                                        jitter_sec=jitter_sec,
                                    )
                                    p_val_b = min([w['p_value'] for w in sig_windows_b]) if sig_windows_b else 1.0

                                    # Accumulate per-patient modulation from patient p-values, the
                                    # number of significant windows, and their total duration.
                                    _eps = 1e-15
                                    _, _, per_pt_a_full = permutation_cluster_jitter_test(
                                        ch_a_full, t_full,
                                        n_permutations=n_permutations,
                                        p_threshold=0.05,
                                        jitter_sec=jitter_sec,
                                    )
                                    _, _, per_pt_b_full = permutation_cluster_jitter_test(
                                        ch_b_full, t_full,
                                        n_permutations=n_permutations,
                                        p_threshold=0.05,
                                        jitter_sec=jitter_sec,
                                    )
                                    for _si, _pv in enumerate(per_pt_a_full.get('p_values', [])):
                                        if _si < n_a_subj:
                                            _n_win = per_pt_a_full.get('n_windows', [0] * n_a_subj)[_si]
                                            _dur = per_pt_a_full.get('total_duration_sec', [0.0] * n_a_subj)[_si]
                                            _weight = _n_win * _dur
                                            pmod_scores_a[_si] += -np.log10(max(_pv, _eps)) * _weight
                                    for _si, _pv in enumerate(per_pt_b_full.get('p_values', [])):
                                        if _si < n_b_subj:
                                            _n_win = per_pt_b_full.get('n_windows', [0] * n_b_subj)[_si]
                                            _dur = per_pt_b_full.get('total_duration_sec', [0.0] * n_b_subj)[_si]
                                            _weight = _n_win * _dur
                                            pmod_scores_b[_si] += -np.log10(max(_pv, _eps)) * _weight
                                    pmod_ch_count += 1
                                    _pmod_per_ch_a.append(per_pt_a_full)
                                    _pmod_per_ch_b.append(per_pt_b_full)

                                    plot_ch_names.append(standard_name)
                                    plot_data_diff.append(diff_median)
                                    plot_data_pval_diff.append(p_val_diff)
                                    
                                    plot_data_amp_a.append(amp_a)
                                    plot_data_amp_b.append(amp_b)
                                    plot_data_pval_a.append(p_val_a)
                                    plot_data_pval_b.append(p_val_b)

                        # Pad missing 10-20 standard channels
                        standard_19_base = ['Fp1', 'Fp2', 'F7', 'F3', 'Fz', 'F4', 'F8', 'C3', 'Cz', 'C4', 'P3', 'Pz', 'P4', 'O1', 'O2']
                        aliases = {'T7': 'T3', 'T8': 'T4', 'P7': 'T5', 'P8': 'T6'}
                        
                        def pad_channel_if_missing(ch_name, pad_amp, pad_pval):
                            if not any(ch_name.upper() == p_ch.upper() for p_ch in plot_ch_names):
                                plot_ch_names.append(ch_name)
                                plot_data_diff.append(0.0)
                                plot_data_pval_diff.append(pad_pval)
                                plot_data_amp_a.append(pad_amp)
                                plot_data_amp_b.append(pad_amp)
                                plot_data_pval_a.append(pad_pval)
                                plot_data_pval_b.append(pad_pval)
                        
                        for base_ch in standard_19_base:
                            pad_channel_if_missing(base_ch, 0.0, 0.05)
                                
                        for new_name, old_name in aliases.items():
                            if not any(ch.upper() in [new_name.upper(), old_name.upper()] for ch in plot_ch_names):
                                pad_channel_if_missing(new_name, 0.0, 0.05)

                        if not plot_ch_names:
                            st.warning("No channels matched the standard montage.")
                        else:
                            # Build channel->value dicts for _render_topomap
                            ch_amp_a   = dict(zip(plot_ch_names, plot_data_amp_a))
                            ch_amp_b   = dict(zip(plot_ch_names, plot_data_amp_b))
                            ch_diff    = dict(zip(plot_ch_names, plot_data_diff))
                            ch_pval_a  = dict(zip(plot_ch_names, plot_data_pval_a))
                            ch_pval_b  = dict(zip(plot_ch_names, plot_data_pval_b))
                            ch_pval_d  = dict(zip(plot_ch_names, plot_data_pval_diff))

                            amp_label = f"Amplitude ({'Z-score' if use_zscore else 'µV'})"
                            diff_label = f"Median Diff ({'Z-score' if use_zscore else 'µV'})"

                            st.markdown("#### 1. Median Amplitude in Time Window")
                            fig_amp, axes_amp = plt.subplots(1, 2, figsize=(10, 4))
                            _render_topomap(ch_amp_a, axes_amp[0], f"{g_a} Amplitude",
                                            cmap='RdBu_r', colorbar_label=amp_label)
                            _render_topomap(ch_amp_b, axes_amp[1], f"{g_b} Amplitude",
                                            cmap='RdBu_r', colorbar_label=amp_label)
                            fig_amp.suptitle(
                                _comparison_title(
                                    "Spatial Median Amplitude",
                                    f"{topo_diff_tmin}-{topo_diff_tmax} ms",
                                ),
                                fontsize=13, fontweight='bold'
                            )
                            fig_amp.tight_layout()
                            st.pyplot(fig_amp, use_container_width=False)
                            plt.close(fig_amp)

                            st.markdown("#### 2. Significance vs Baseline (p-value)")
                            fig_pval_ind, axes_pval_ind = plt.subplots(1, 2, figsize=(10, 4))
                            _render_pval_topomap(ch_pval_a, axes_pval_ind[0], f"{g_a} P-value (vs 0)")
                            _render_pval_topomap(ch_pval_b, axes_pval_ind[1], f"{g_b} P-value (vs 0)")
                            fig_pval_ind.suptitle(
                                _comparison_title(
                                    "Spatial Significance vs Baseline",
                                    f"{topo_diff_tmin}-{topo_diff_tmax} ms",
                                ),
                                fontsize=13, fontweight='bold'
                            )
                            fig_pval_ind.tight_layout()
                            st.pyplot(fig_pval_ind, use_container_width=False)
                            plt.close(fig_pval_ind)

                            st.markdown("#### 3. Group Difference & Significance")
                            fig_diff_topo, axes_diff_topo = plt.subplots(1, 2, figsize=(10, 4))
                            _render_topomap(ch_diff, axes_diff_topo[0], f"Difference ({g_a} - {g_b})",
                                            cmap='RdBu_r', colorbar_label=diff_label)
                            _render_topomap(ch_pval_d, axes_diff_topo[1], "P-value (Diff)",
                                            cmap='Reds_r', vlim=(0, 0.05), colorbar_label='p-value')
                            fig_diff_topo.suptitle(
                                _comparison_title(
                                    "Spatial Group Difference and Significance",
                                    f"{topo_diff_tmin}-{topo_diff_tmax} ms",
                                ),
                                fontsize=13, fontweight='bold'
                            )
                            fig_diff_topo.tight_layout()
                            st.pyplot(fig_diff_topo, use_container_width=False)
                            plt.close(fig_diff_topo)

                            st.markdown("#### 4. Per-Patient HBCI (Heart-Brain Coupling Index)")
                            _eps_hbci = 1e-15
                            p_threshold = 0.05
                            hbci_a = np.zeros(n_a_subj)
                            hbci_b = np.zeros(n_b_subj)
                            _n_sig_ch_a = np.zeros(n_a_subj)
                            _n_sig_ch_b = np.zeros(n_b_subj)
                            if _pmod_per_ch_a:
                                for _ch_ppa, _ch_ppb in zip(_pmod_per_ch_a, _pmod_per_ch_b):
                                    for _si in range(n_a_subj):
                                        _pv = _ch_ppa['p_values'][_si] if _si < len(_ch_ppa['p_values']) else 1.0
                                        _dur = _ch_ppa['total_duration_sec'][_si] if _si < len(_ch_ppa.get('total_duration_sec', [])) else 0.0
                                        if _pv < p_threshold:
                                            hbci_a[_si] += _dur * (-np.log10(max(_pv, _eps_hbci)))
                                            _n_sig_ch_a[_si] += 1
                                    for _si in range(n_b_subj):
                                        _pv = _ch_ppb['p_values'][_si] if _si < len(_ch_ppb['p_values']) else 1.0
                                        _dur = _ch_ppb['total_duration_sec'][_si] if _si < len(_ch_ppb.get('total_duration_sec', [])) else 0.0
                                        if _pv < p_threshold:
                                            hbci_b[_si] += _dur * (-np.log10(max(_pv, _eps_hbci)))
                                            _n_sig_ch_b[_si] += 1
                                _n_total_ch = max(pmod_ch_count, 1)
                                hbci_a = (_n_sig_ch_a / _n_total_ch) * hbci_a
                                hbci_b = (_n_sig_ch_b / _n_total_ch) * hbci_b
                                from scipy.stats import mannwhitneyu
                                _u_stat_h, p_hbci = mannwhitneyu(hbci_a, hbci_b, alternative='two-sided')
                                fig_hbci, ax_hbci = plt.subplots(figsize=(5, 5))
                                bp_h = ax_hbci.boxplot(
                                    [hbci_a, hbci_b], tick_labels=[g_a, g_b],
                                    patch_artist=True, widths=0.4,
                                    medianprops=dict(color='black', linewidth=2),
                                )
                                _hbci_colors = ['#2196F3', '#FF5722']
                                for patch, c in zip(bp_h['boxes'], _hbci_colors):
                                    patch.set_facecolor(c)
                                    patch.set_alpha(0.55)
                                _rng_h = np.random.default_rng(42)
                                for _i, (scores, c) in enumerate(zip([hbci_a, hbci_b], _hbci_colors), start=1):
                                    _jit = _rng_h.uniform(-0.12, 0.12, size=len(scores))
                                    ax_hbci.scatter(np.full(len(scores), _i) + _jit, scores,
                                                   color=c, alpha=0.8, s=35, zorder=3,
                                                   edgecolors='white', linewidths=0.5)
                                _all_h = np.concatenate([hbci_a, hbci_b])
                                _y_max_h = max(float(np.max(_all_h)) if len(_all_h) > 0 else 1.0, 1e-6)
                                _bar_top_h = _y_max_h * 1.08
                                _tick_h2 = _y_max_h * 0.03
                                ax_hbci.plot([1, 1, 2, 2],
                                            [_bar_top_h - _tick_h2, _bar_top_h, _bar_top_h, _bar_top_h - _tick_h2],
                                            'k-', linewidth=1)
                                if p_hbci < 0.001:
                                    _p_lbl = "p < 0.001 ***"
                                elif p_hbci < 0.01:
                                    _p_lbl = f"p = {p_hbci:.3f} **"
                                elif p_hbci < 0.05:
                                    _p_lbl = f"p = {p_hbci:.3f} *"
                                else:
                                    _p_lbl = f"p = {p_hbci:.3f} (n.s.)"
                                ax_hbci.text(1.5, _bar_top_h + _tick_h2 * 0.5, _p_lbl,
                                            ha='center', va='bottom', fontsize=10)
                                ax_hbci.set_ylabel("HBCI score")
                                ax_hbci.set_title(
                                    f"Per-Patient HBCI | Stage: {selected_stage} | Groups: {g_a} vs {g_b}\n"
                                    "(channel_fraction x duration x -log10(p))"
                                )
                                ax_hbci.set_ylim(bottom=0, top=_y_max_h * 1.18)
                                fig_hbci.tight_layout()
                                st.pyplot(fig_hbci, use_container_width=False)
                                plt.close(fig_hbci)
                                col_a, col_b = st.columns(2)
                                col_a.metric(f"{g_a} Median HBCI", f"{np.median(hbci_a):.4f}", delta=f"n={n_a_subj}")
                                col_b.metric(f"{g_b} Median HBCI", f"{np.median(hbci_b):.4f}", delta=f"n={n_b_subj}")
                                st.caption(
                                    "HBCI per patient = channel_fraction x sum_channels(duration_sec x -log10(p)). "
                                    "Same formula as the single-patient HBCI index. "
                                    "Higher values = more channels with longer, more significant HEP modulation."
                                )
                            else:
                                st.info("No per-channel permutation data available to compute HBCI.")

            except Exception as e:
                import traceback as _tb
                _details = _tb.format_exc()
                st.error(
                    f"**Error generating topomaps** ({type(e).__name__}): {e}\n\n"
                    f"Likely cause: channel count mismatch between data and MNE montage positions."
                )
                with st.expander("Traceback — Topomap Error"):
                    st.code(_details)

    # ═══════════════════════════════════════════════════════════════════════
    # SUMMARY STATISTICS TABLE
    # ═══════════════════════════════════════════════════════════════════════
    st.subheader("📋 Summary Statistics")
    summary_rows = []
    for group in groups_with_data:
        mat = active_group_hep_matrix[group]  # already normalised (Z-scored or µV)
        grand, grand_mad = median_and_mad(mat, axis=0)
        n_subj = mat.shape[0]
        peak_idx = np.argmax(np.abs(grand))
        peak_amp = grand[peak_idx]
        peak_t = active_group_times[group][peak_idx]
        median_amp = np.nanmedian(grand)
        mad_amp = np.nanmedian(np.abs(grand - median_amp))
        unit_lbl = "Z" if use_zscore else "µV"
        row = {
            "Group": group,
            "N": n_subj,
            f"Median amplitude ({unit_lbl})": f"{median_amp:.4f}",
            f"MAD amplitude ({unit_lbl})": f"{mad_amp:.4f}",
            f"Peak amplitude ({unit_lbl})": f"{peak_amp:.4f}",
            "Peak latency (s)": f"{peak_t:.3f}",
        }
        # Add between-group stats at peak time (2-group case)
        if n_groups == 2:
            g_a, g_b = groups_with_data[0], groups_with_data[1]
            hep_at_peak_a = active_group_hep_matrix[g_a][:, peak_idx]
            hep_at_peak_b = active_group_hep_matrix[g_b][:, peak_idx]
            t_stat_peak, p_val_peak = stats.ttest_ind(hep_at_peak_a, hep_at_peak_b)
            na2, nb2 = len(hep_at_peak_a), len(hep_at_peak_b)
            sd_pool = np.sqrt(
                ((na2 - 1) * np.var(hep_at_peak_a, ddof=1) + (nb2 - 1) * np.var(hep_at_peak_b, ddof=1))
                / (na2 + nb2 - 2)
            )
            d_peak = abs(np.mean(hep_at_peak_a) - np.mean(hep_at_peak_b)) / (sd_pool + 1e-12)
            if group == g_a:
                row["t-stat at peak"] = f"{t_stat_peak:.3f}"
                row["p-value at peak"] = "<0.001" if p_val_peak < 0.001 else f"{p_val_peak:.3f}"
                row["Cohen's d at peak"] = f"{d_peak:.3f}"
        summary_rows.append(row)

    st.dataframe(pd.DataFrame(summary_rows), use_container_width=True)




if TORCH_AVAILABLE:
    class ConvBlock(nn.Module):
        def __init__(self, in_c, out_c):
            super().__init__()
            self.conv1 = nn.Conv1d(in_c, out_c, kernel_size=3, padding=1)
            self.bn1 = nn.BatchNorm1d(out_c)
            self.conv2 = nn.Conv1d(out_c, out_c, kernel_size=3, padding=1)
            self.bn2 = nn.BatchNorm1d(out_c)
            self.relu = nn.ReLU()
        
        def forward(self, x):
            return self.relu(self.bn2(self.conv2(self.relu(self.bn1(self.conv1(x))))))

    class UNet1D(nn.Module):
        def __init__(self):
            super().__init__()
            self.enc1 = ConvBlock(1, 64)
            self.enc2 = ConvBlock(64, 128)
            self.enc3 = ConvBlock(128, 256)
            self.pool = nn.MaxPool1d(2)
            self.up2 = nn.ConvTranspose1d(256, 128, 2, stride=2)
            self.dec2 = ConvBlock(256, 128)
            self.up1 = nn.ConvTranspose1d(128, 64, 2, stride=2)
            self.dec1 = ConvBlock(128, 64)
            self.out = nn.Conv1d(64, 1, 1)
        
        def forward(self, x):
            e1 = self.enc1(x)
            e2 = self.enc2(self.pool(e1))
            b  = self.enc3(self.pool(e2))
            d2 = self.dec2(torch.cat([self.up2(b), e2], dim=1))
            d1 = self.dec1(torch.cat([self.up1(d2), e1], dim=1))
            return torch.sigmoid(self.out(d1))

    class ECG_CNN_RPeak(nn.Module):
        def __init__(self):
            super().__init__()
            self.encoder = nn.Sequential(
                nn.Conv1d(1, 32, kernel_size=7, padding=3),
                nn.BatchNorm1d(32), nn.ReLU(),
                nn.Conv1d(32, 64, kernel_size=5, padding=2),
                nn.BatchNorm1d(64), nn.ReLU(),
                nn.Conv1d(64, 128, kernel_size=3, padding=1),
                nn.BatchNorm1d(128), nn.ReLU(),
            )
            self.classifier = nn.Conv1d(128, 1, kernel_size=1)
        
        def forward(self, x):
            return torch.sigmoid(self.classifier(self.encoder(x)))

    class BiLSTM_RPeak(nn.Module):
        def __init__(self):
            super().__init__()
            self.bilstm = nn.LSTM(input_size=1, hidden_size=64, num_layers=2, bidirectional=True, batch_first=True)
            self.fc = nn.Linear(128, 1)
        
        def forward(self, x):
            out, _ = self.bilstm(x)
            return torch.sigmoid(self.fc(out))

def swt_preprocess(ecg, level=4, wavelet='db4'):
    if not PYWT_AVAILABLE:
        return np.zeros((2, len(ecg)))
    coeffs = pywt.swt(ecg, wavelet, level=level)
    detail_l4 = coeffs[level-1][1]       # Detail at level 4
    derivative = np.diff(ecg, prepend=ecg[0])
    return np.stack([detail_l4, derivative], axis=0)  # 2-channel DL input

def plot_ecg_cleaning_comparison(raw_segment, clean_segment, sfreq, patient_id):
    """
    Plots a 10 second segment of raw vs cleaned ECG and compares R-peak detection methods.
    """
    times = np.arange(len(raw_segment)) / sfreq

    # --- Compute 3 Methods ---
    methods_peaks = {}
    
    progress_bar = st.progress(0)
    status_text = st.empty()

    # 1. NeuroKit2
    status_text.text("Running NeuroKit2 (1/3)...")
    try:
        import neurokit2 as nk
        _, rpk_nk = nk.ecg_peaks(clean_segment, sampling_rate=sfreq,method='promac')
        methods_peaks['NeuroKit2'] = rpk_nk['ECG_R_Peaks']
    except Exception as e:
        methods_peaks[f'NeuroKit2 (Error: {e})'] = []
    progress_bar.progress(33)
    
    # 2. WFDB
    status_text.text("Running WFDB (2/3)...")
    try:
        import wfdb.processing as wp
        rpeaks_wfdb = wp.xqrs_detect(clean_segment, fs=sfreq, verbose=False)
        # Optional: correct peak locations to local maxima
        r_peaks_corrected = wp.correct_peaks(clean_segment, rpeaks_wfdb, 
                                            search_radius=int(0.05*sfreq),
                                            smooth_window_size=int(0.1*sfreq))
        methods_peaks['WFDB XQRS'] = r_peaks_corrected
    except Exception as e:
        methods_peaks[f'WFDB (Error: {e})'] = []
    progress_bar.progress(66)

    # 3. WFDB Robust
    status_text.text("Running WFDB part 2 (3/3)...")
    try:
        methods_peaks['WFDB Robust'] = detect_rpeaks_robust(clean_segment, sfreq)
    except Exception as e:
        methods_peaks[f'WFDB Robust (Error: {e})'] = []

    progress_bar.progress(100)
    status_text.text("Processing Complete!")

    # --- Plotting ---
    method_names = list(methods_peaks.keys())
    fig, axes = plt.subplots(len(method_names), 1, figsize=(16, 3 * len(method_names)), sharex=True)
    
    if len(method_names) == 1:
        axes = [axes]

    for idx, (method_name, rpeaks) in enumerate(methods_peaks.items()):
        ax = axes[idx]
        ax.plot(times, raw_segment, color='gray', alpha=0.3, label='Raw' if idx == 0 else None)
        ax.plot(times, clean_segment, color='red', alpha=0.7, label='Cleaned' if idx == 0 else None)
        
        for i, peak in enumerate(rpeaks):
            peak_time = peak / sfreq
            ax.axvline(peak_time, color='green', linestyle='--', alpha=0.6, 
                       label='R-peak' if (idx == 0 and i == 0) else None)
        
        ax.set_ylabel("Amp (μV)")
        ax.set_title(f"Method: {method_name}")
        if idx == 0:
            ax.legend(loc='upper right')
        ax.grid(True, alpha=0.3)

    axes[-1].set_xlabel("Time (s)")
    plt.tight_layout()
    st.pyplot(fig, use_container_width=True)

def plot_ecg_hep_individuals(individuals, selected_group, selected_stage):
    """
    Plots ECG HEP for all individuals and their average.
    """
    fig, ax = plt.subplots(figsize=(16, 9))
    all_ecg_heps = []
    # allow user to select how many individuals to plot
    n_individuals = st.slider("Number of individuals to plot", min_value=1, max_value=len(individuals), value=len(individuals))
    for ind in individuals[:n_individuals]:
        # ind structure: (patient_id, hep_data, times, ch_names, rpeaks, ecg_hep_data, ecg_ch_names)
        ecg_hep = ind[5]
        if ecg_hep is not None:
            hep = ecg_hep[0]
            ax.plot(ind[2], hep * 1e6, color='gray', alpha=0.3, linewidth=1)
            all_ecg_heps.append(hep)
    
    if all_ecg_heps:
        avg_ecg = np.nanmean(all_ecg_heps, axis=0)
        finalize_plot(
            fig, ax, 
            f"ECG - Group: {selected_group} - Stage: {selected_stage}",
            avg_hep=avg_ecg, 
            times=individuals[0][2], 
            n_subjects=len(all_ecg_heps)
        )
        
        # Calculate epochs stats
        epochs_data = [{"Patient ID": ind[0], "Epochs": len(ind[4])} for ind in individuals]
        if epochs_data:
            df_epochs = pd.DataFrame(epochs_data)
            avg_epochs = df_epochs["Epochs"].mean()
            st.write(f"**Average Epochs per Patient:** {avg_epochs:.2f}")
            
            with st.expander("See Epochs per Patient"):
                st.dataframe(df_epochs)
                
    else:
        st.warning("No ECG data found for this group.")

def analyze_ecg_repetitive_pattern(ecg_signal, sfreq, patient_id=""):
    """
    Detects repetitive waveform templates/motifs in ECG signal.
    This identifies identical wave patterns that repeat throughout the recording.
    
    Returns:
    --------
    pattern_metrics : dict
        Dictionary containing:
        - 'patient_id': Patient identifier
        - 'template': Extracted repetitive template waveform
        - 'template_times': Time axis for template
        - 'match_locations': Sample indices where template repeats
        - 'match_correlations': Correlation strength at each match location
        - 'avg_correlation': Average correlation strength
        - 'repetition_rate': Matches per second
        - 'has_repetitive_pattern': Boolean indicating significant repetition
    """
    from scipy.signal import correlate, find_peaks

    # Use a manageable segment for analysis (first 60 seconds)
    max_samples = min(len(ecg_signal), int(60 * sfreq))
    signal_segment = ecg_signal[:max_samples]
    
    # Template extraction parameters
    # Try different template durations to find repeating patterns
    template_durations = [0.05, 0.1, 0.2, 0.3, 0.5]  # seconds
    
    best_template = None
    best_matches = []
    best_correlations = []
    best_avg_corr = 0
    best_duration = 0
    
    for template_dur in template_durations:
        template_len = int(template_dur * sfreq)
        
        if template_len >= len(signal_segment) // 4:
            continue
        
        # Strategy: Find the most repetitive segment by testing multiple candidates
        # Sample multiple starting points
        n_candidates = min(20, len(signal_segment) // template_len)
        candidate_starts = np.linspace(0, len(signal_segment) - template_len, n_candidates, dtype=int)
        
        for start_idx in candidate_starts:
            # Extract candidate template
            template = signal_segment[start_idx:start_idx + template_len]
            
            # Normalize template
            template_norm = (template - np.mean(template)) / (np.std(template) + 1e-10)
            
            # Normalize signal segment
            signal_norm = (signal_segment - np.mean(signal_segment)) / (np.std(signal_segment) + 1e-10)
            
            # Cross-correlate to find matches
            correlation = correlate(signal_norm, template_norm, mode='valid')
            
            # Find peaks in correlation (high correlation = good match)
            # Require minimum separation between matches
            min_separation = template_len // 2
            threshold = 0.7  # Correlation threshold for a "good match"
            
            peak_indices, properties = get_peaks(correlation, 
                                                   height=threshold, 
                                                   distance=min_separation)
            
            if len(peak_indices) > 0:
                avg_corr = np.mean(correlation[peak_indices])
                
                # Keep the template with the most repetitions and highest avg correlation
                if avg_corr > best_avg_corr and len(peak_indices) >= 3:
                    best_avg_corr = avg_corr
                    best_template = template
                    best_matches = peak_indices
                    best_correlations = correlation[peak_indices]
                    best_duration = template_dur
    
    # Prepare return values
    if best_template is not None:
        template_times = np.arange(len(best_template)) / sfreq
        repetition_rate = len(best_matches) / (len(signal_segment) / sfreq)
        has_repetitive_pattern = len(best_matches) >= 5 and best_avg_corr >= 0.75
    else:
        # No strong repetitive pattern found
        template_times = np.array([])
        repetition_rate = 0
        has_repetitive_pattern = False
        best_template = np.array([])
        best_matches = []
        best_correlations = []
    
    return {
        'patient_id': patient_id,
        'template': best_template,
        'template_times': template_times,
        'template_duration': best_duration,
        'match_locations': best_matches,
        'match_correlations': best_correlations,
        'avg_correlation': best_avg_corr,
        'repetition_rate': repetition_rate,
        'n_matches': len(best_matches),
        'has_repetitive_pattern': has_repetitive_pattern,
        'signal_segment': signal_segment,  # For visualization
        'sfreq': sfreq
    }

def handle_ecg_noise_detection(base_path, selected_group, selected_stage):
    """
    Analyzes ECG for repetitive waveform patterns in all patients and displays results.
    """
    st.subheader(f"ECG Repetitive Pattern Analysis - {selected_group} {selected_stage}")
    
    st.info("This analysis detects identical waveform patterns that repeat throughout each patient's ECG recording.")
    
    group_dir = os.path.join(base_path, selected_group, selected_stage)
    files = [f for f in os.listdir(group_dir) if f.endswith('.pkl')]
    
    if not files:
        st.warning("No files found for pattern analysis.")
        return
    
    all_metrics = []
    
    progress_bar = st.progress(0)
    
    for idx, f_name in enumerate(files):
        progress_bar.progress((idx + 1) / len(files))
        
        f_path = os.path.join(group_dir, f_name)
        patient_id = f_name.replace('.pkl', '').replace('.edf', '')
        
        try:
            with open(f_path, 'rb') as f:
                raw = pickle.load(f)
            
            # Extract ECG channel
            sfreq = raw.info['sfreq']
            ch_names = raw.ch_names
            ch_lower = [ch.lower() for ch in ch_names]
            ecg_indices = [i for i, ch in enumerate(ch_lower) if 'ecg' in ch or 'ekg' in ch]
            
            if not ecg_indices:
                continue
            
            ecg_idx = ecg_indices[0]
            ecg_signal = raw.get_data(picks=[ecg_idx])[0] * 1e6  # Convert to uV
            
            # Analyze for repetitive patterns
            metrics = analyze_ecg_repetitive_pattern(ecg_signal, sfreq, patient_id)
            all_metrics.append(metrics)
            
        except Exception as e:
            st.warning(f"Error processing {patient_id}: {e}")
            continue
    
    progress_bar.empty()
    
    if not all_metrics:
        st.error("No valid ECG data found for pattern analysis.")
        return
    
    # --- Summary Statistics ---
    st.markdown("### Group Summary")
    
    n_patients = len(all_metrics)
    n_with_pattern = sum(1 for m in all_metrics if m['has_repetitive_pattern'])
    
    # Calculate averages only for patients with detected patterns
    valid_metrics = [m for m in all_metrics if len(m['template']) > 0]
    if valid_metrics:
        avg_repetition_rate = np.mean([m['repetition_rate'] for m in valid_metrics])
        avg_correlation = np.mean([m['avg_correlation'] for m in valid_metrics])
        avg_n_matches = np.mean([m['n_matches'] for m in valid_metrics])
    else:
        avg_repetition_rate = 0
        avg_correlation = 0
        avg_n_matches = 0
    
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric("Total Patients", n_patients)
    with col2:
        st.metric("With Repetitive Pattern", f"{n_with_pattern} ({100*n_with_pattern/n_patients:.1f}%)")
    with col3:
        st.metric("Avg Repetition Rate", f"{avg_repetition_rate:.2f} /sec")
    with col4:
        st.metric("Avg Correlation", f"{avg_correlation:.3f}")
    
    # --- Per-Patient Table ---
    st.markdown("### Per-Patient Pattern Metrics")
    
    df_data = []
    for m in all_metrics:
        if len(m['template']) > 0:
            df_data.append({
                'Patient ID': m['patient_id'],
                'Pattern Duration (ms)': f"{m['template_duration']*1000:.1f}",
                'Num Repetitions': m['n_matches'],
                'Repetition Rate (/sec)': f"{m['repetition_rate']:.2f}",
                'Avg Correlation': f"{m['avg_correlation']:.3f}",
                'Has Strong Pattern': '✓' if m['has_repetitive_pattern'] else '✗'
            })
        else:
            df_data.append({
                'Patient ID': m['patient_id'],
                'Pattern Duration (ms)': 'N/A',
                'Num Repetitions': 0,
                'Repetition Rate (/sec)': '0.00',
                'Avg Correlation': 'N/A',
                'Has Strong Pattern': '✗'
            })
    
    df = pd.DataFrame(df_data)
    st.dataframe(df, use_container_width=True)
    
    # --- Visualization: Individual Patient Templates ---
    st.markdown("### Individual Patient Templates")
    st.markdown("Shows the extracted repetitive waveform pattern for each patient")
    
    # Filter patients with valid templates
    patients_with_templates = [m for m in all_metrics if len(m['template']) > 0]
    
    if not patients_with_templates:
        st.warning("No repetitive patterns detected in any patient.")
        return
    
    n_show = st.slider("Number of patients to show", min_value=1, 
                       max_value=min(12, len(patients_with_templates)), 
                       value=min(6, len(patients_with_templates)), 
                       key="pattern_patient_slider")
    
    # Create grid of template plots
    n_cols = 3
    n_rows = int(np.ceil(n_show / n_cols))
    
    fig_templates, axes = plt.subplots(n_rows, n_cols, figsize=(15, 4*n_rows))
    if n_rows == 1 and n_cols == 1:
        axes = np.array([axes])
    axes = axes.flatten()
    
    for i in range(n_show):
        ax = axes[i]
        m = patients_with_templates[i]
        
        # Plot the template
        ax.plot(m['template_times'] * 1000, m['template'], color='red', linewidth=2)
        
        pattern_status = "STRONG" if m['has_repetitive_pattern'] else "WEAK"
        ax.set_title(f"{m['patient_id']} - {pattern_status}\n"
                     f"Repeats: {m['n_matches']}, Corr: {m['avg_correlation']:.2f}", 
                     fontsize=9)
        ax.set_xlabel('Time (ms)', fontsize=8)
        ax.set_ylabel('Amplitude (μV)', fontsize=8)
        ax.grid(True, alpha=0.3)
        ax.tick_params(labelsize=8)
    
    # Hide unused subplots
    for j in range(n_show, len(axes)):
        axes[j].axis('off')
    
    plt.tight_layout()
    st.pyplot(fig_templates, use_container_width=True)
    
    # --- Detailed View: Template in Context ---
    st.markdown("### Detailed Pattern View")
    st.markdown("View the repetitive pattern in context of the full ECG signal")
    
    selected_patient_idx = st.selectbox(
        "Select patient for detailed view",
        range(len(patients_with_templates)),
        format_func=lambda i: patients_with_templates[i]['patient_id']
    )
    
    m = patients_with_templates[selected_patient_idx]
    
    # Show template
    st.markdown(f"**Patient: {m['patient_id']}**")
    st.markdown(f"- Pattern duration: {m['template_duration']*1000:.1f} ms")
    st.markdown(f"- Found {m['n_matches']} repetitions (rate: {m['repetition_rate']:.2f} /sec)")
    st.markdown(f"- Average correlation: {m['avg_correlation']:.3f}")
    
    col_a, col_b = st.columns(2)
    
    with col_a:
        # Plot template
        fig_temp, ax_temp = plt.subplots(figsize=(7, 4))
        ax_temp.plot(m['template_times'] * 1000, m['template'], color='red', linewidth=2)
        ax_temp.set_title("Extracted Repetitive Template", fontsize=12, fontweight='bold')
        ax_temp.set_xlabel('Time (ms)')
        ax_temp.set_ylabel('Amplitude (μV)')
        ax_temp.grid(True, alpha=0.3)
        st.pyplot(fig_temp, use_container_width=True)
    
    with col_b:
        # Plot correlation values at match locations
        fig_corr, ax_corr = plt.subplots(figsize=(7, 4))
        ax_corr.bar(range(len(m['match_correlations'])), m['match_correlations'], 
                    color='steelblue', alpha=0.7)
        ax_corr.axhline(0.75, color='red', linestyle='--', label='Strong Match Threshold')
        ax_corr.set_title("Correlation Strength per Repetition", fontsize=12, fontweight='bold')
        ax_corr.set_xlabel('Repetition Number')
        ax_corr.set_ylabel('Correlation Coefficient')
        ax_corr.set_ylim(0, 1)
        ax_corr.legend()
        ax_corr.grid(True, alpha=0.3, axis='y')
        st.pyplot(fig_corr, use_container_width=True)
    
    # Show signal with marked repetitions
    st.markdown("**Signal with Pattern Locations Marked**")
    
    # User can select time window
    signal_duration = len(m['signal_segment']) / m['sfreq']
    view_window = st.slider("View window (seconds)", 
                            min_value=2, 
                            max_value=min(30, int(signal_duration)), 
                            value=10,
                            key="pattern_window_slider")
    
    start_time = st.slider("Start time (seconds)", 
                          min_value=0.0, 
                          max_value=max(0.0, signal_duration - view_window),
                          value=0.0,
                          step=0.5,
                          key="pattern_start_slider")
    
    start_sample = int(start_time * m['sfreq'])
    end_sample = min(len(m['signal_segment']), int((start_time + view_window) * m['sfreq']))
    
    signal_view = m['signal_segment'][start_sample:end_sample]
    times_view = np.arange(len(signal_view)) / m['sfreq'] + start_time
    
    fig_signal, ax_signal = plt.subplots(figsize=(15, 5))
    ax_signal.plot(times_view, signal_view, color='black', linewidth=1, alpha=0.7, label='ECG Signal')
    
    # Mark pattern locations in the view window
    for match_loc in m['match_locations']:
        match_time = match_loc / m['sfreq']
        if start_time <= match_time <= start_time + view_window:
            ax_signal.axvline(match_time, color='red', linestyle='--', alpha=0.5, linewidth=1.5)
    
    ax_signal.set_title("ECG Signal with Repetitive Pattern Locations (Red lines)", fontsize=12)
    ax_signal.set_xlabel('Time (s)')
    ax_signal.set_ylabel('Amplitude (μV)')
    ax_signal.grid(True, alpha=0.3)
    ax_signal.legend()
    st.pyplot(fig_signal, use_container_width=True)
    
    # --- Average Template Across All Patients ---
    st.markdown("### Average Template Across All Patients")
    
    # Find common duration (use median)
    template_lengths = [len(m['template']) for m in patients_with_templates]
    median_length = int(np.median(template_lengths))
    
    # Interpolate all templates to median length and average
    from scipy.interpolate import interp1d
    
    aligned_templates = []
    for m in patients_with_templates:
        if len(m['template']) > 0:
            # Interpolate to median length
            x_old = np.linspace(0, 1, len(m['template']))
            x_new = np.linspace(0, 1, median_length)
            f = interp1d(x_old, m['template'], kind='linear')
            template_aligned = f(x_new)
            aligned_templates.append(template_aligned)
    
    if aligned_templates:
        avg_template = np.mean(aligned_templates, axis=0)
        std_template = np.std(aligned_templates, axis=0)
        
        # Use median duration for time axis
        median_duration = np.median([m['template_duration'] for m in patients_with_templates])
        avg_times = np.linspace(0, median_duration * 1000, median_length)
        
        fig_avg, ax_avg = plt.subplots(figsize=(10, 5))
        ax_avg.plot(avg_times, avg_template, color='blue', linewidth=3, label=f'Average (n={len(aligned_templates)})')
        ax_avg.fill_between(avg_times, avg_template - std_template, avg_template + std_template,
                           color='blue', alpha=0.2, label='± 1 SD')
        
        ax_avg.set_title(f"Average Repetitive Pattern - Mean +/- STD - {selected_group} {selected_stage}",
                        fontsize=14, fontweight='bold')
        ax_avg.set_xlabel('Time (ms)', fontsize=11)
        ax_avg.set_ylabel('Amplitude (μV)', fontsize=11)
        ax_avg.grid(True, alpha=0.3)
    ax_avg.legend()
    st.pyplot(fig_avg, use_container_width=True)

    # --- Cleaning Methods Comparison ---
    st.markdown("### Cleaning Methods Comparison")
    st.markdown("Visualizing the effect of different artifact removal methods on the selected patient's signal.")
    
    if 'm' in locals() and len(m['template']) > 0:
        # Use the same patient selected in "Detailed Pattern View"
        
        # 1. Define Cleaning Functions locally or call them if defined outside
        # Defining them here for clarity in context, or better move to global scope. 
        # I'll define them here to access 'm' easily if needed, but standard practice is outside.
        # Let's call the helper functions defined above (we will add them to the file scope in a moment).
        # Since I can't add them to file scope in this single Replace block easily without replacing the whole file,
        # I will define them as nested functions or just implement the logic inline for this block.
        # Actually, I can add the helper functions *before* this function in a previous edit, or just include them here.
        # Given the constraint of one Replace block, I will include the logic here.
        
        # ... Wait, I should add the helper functions properly at module level. 
        # But this tool call is for `replace_file_content`. I can only do one contiguous block.
        # I will implement the logic inside `handle_ecg_noise_detection` or helper functions *inside* this block 
        # if I want to do it in one go. However, to keep code clean, I should probably use `multi_replace` 
        # if I want to insert functions elsewhere. 
        # BUT, for now, let's implement the cleaning logic directly here for the specific signal view.
        
        signal_seg = m['signal_segment']
        template = m['template']
        match_locs = m['match_locations']
        template_len = len(template)
        
        # Method 1: Global Template Subtraction
        clean_global = signal_seg.copy()
        for start_idx in match_locs:
            if start_idx + template_len <= len(clean_global):
                clean_global[start_idx:start_idx+template_len] -= template
                
        # Method 2: Scaled Template Subtraction
        clean_scaled = signal_seg.copy()
        template_norm_sq = np.sum(template**2)
        if template_norm_sq > 0:
            for start_idx in match_locs:
                if start_idx + template_len <= len(clean_scaled):
                    segment = clean_scaled[start_idx:start_idx+template_len]
                    scale = np.dot(segment, template) / template_norm_sq
                    # Limit scale to reasonable bounds to avoid subtracting noise that looks like template
                    # scale = np.clip(scale, 0.5, 2.0) # Optional
                    clean_scaled[start_idx:start_idx+template_len] -= scale * template
        
        # Method 3: Local Moving Average
        clean_local = signal_seg.copy()
        window_size = 10
        n_matches = len(match_locs)
        
        for i, start_idx in enumerate(match_locs):
            # Find local window indices
            start_w = max(0, i - window_size // 2)
            end_w = min(n_matches, start_w + window_size)
            if end_w - start_w < window_size and start_w > 0:
                start_w = max(0, end_w - window_size)
            
            local_indices = match_locs[start_w:end_w]
            
            # Compute local template
            local_segments = []
            for loc in local_indices:
                if loc + template_len <= len(signal_seg):
                    local_segments.append(signal_seg[loc:loc+template_len])
            
            if local_segments:
                local_temp = np.mean(local_segments, axis=0)
                if start_idx + template_len <= len(clean_local):
                    clean_local[start_idx:start_idx+template_len] -= local_temp

        # --- Plotting ---
        # Reuse the zoomed window from earlier
        # start_sample, end_sample, times_view defined above in "Signal with Pattern Locations Marked"
        
        # We need to slice the cleaned signals to match the view window
        # Ensure we have the view variables
        if 'start_sample' in locals() and 'end_sample' in locals() and 'times_view' in locals():
            sig_view_raw = signal_seg[start_sample:end_sample]
            sig_view_global = clean_global[start_sample:end_sample]
            sig_view_scaled = clean_scaled[start_sample:end_sample]
            sig_view_local = clean_local[start_sample:end_sample]
            
            fig_clean, axes_clean = plt.subplots(4, 1, figsize=(15, 10), sharex=True, sharey=True)
            
            # 1. Original
            axes_clean[0].plot(times_view, sig_view_raw, color='black', alpha=0.8, linewidth=1)
            axes_clean[0].set_title("Original Signal (Red lines = detected patterns)", fontsize=10, fontweight='bold')
            # Add markers
            for match_loc in match_locs:
                match_time = match_loc / m['sfreq']
                if times_view[0] <= match_time <= times_view[-1]:
                    axes_clean[0].axvline(match_time, color='red', linestyle='--', alpha=0.5)
            
            # 2. Global
            axes_clean[1].plot(times_view, sig_view_global, color='#1f77b4', linewidth=1)
            axes_clean[1].set_title("Method 1: Global Template Subtraction", fontsize=10, fontweight='bold')
            
            # 3. Scaled
            axes_clean[2].plot(times_view, sig_view_scaled, color='#ff7f0e', linewidth=1)
            axes_clean[2].set_title("Method 2: Scaled Template Subtraction (Best Fit)", fontsize=10, fontweight='bold')
            
            # 4. Local
            axes_clean[3].plot(times_view, sig_view_local, color='#2ca02c', linewidth=1)
            axes_clean[3].set_title(f"Method 3: Local Moving Average (Window={window_size})", fontsize=10, fontweight='bold')
            
            for ax in axes_clean:
                ax.grid(True, alpha=0.3)
                ax.set_ylabel("Amplitude (μV)")
            
            axes_clean[-1].set_xlabel("Time (s)")
            
            plt.tight_layout()
            st.pyplot(fig_clean, use_container_width=True)
        else:
            st.warning("Please interact with the Signal View slider above to initialize the plot.")

    st.markdown("""
    **Interpretation Guide:**
    - **Template**: The extracted waveform pattern that repeats throughout the recording
    - **Repetition Rate**: How many times per second the pattern appears
    - **Correlation**: How similar each repetition is to the template (1.0 = perfect match)
    - **Strong Pattern**: Pattern repeats ≥5 times with correlation ≥0.75
    - **Red vertical lines**: Locations where the repetitive pattern was detected
    """)

def render_hbci(patient_channel_sig, ch_names, times, selected_pid):
    """Compute and display the Heart-Brain Coupling Index for one patient."""
    st.divider()

    if not patient_channel_sig:
        st.info("HBCI: no significant channels found for this patient.")
        return

    channel_scores = {}
    channel_min_p = {}
    for ch, ch_data in patient_channel_sig.items():
        sig_windows = [w for w in ch_data.get('windows', []) if w['p_value'] < 0.05]
        if not sig_windows:
            continue
        score = sum(
            (w['end'] - w['start']) * (-np.log10(max(w['p_value'], 1e-10)))
            for w in sig_windows
        )
        channel_scores[ch] = score
        channel_min_p[ch] = min(max(w['p_value'], 1e-10) for w in sig_windows)

    if not channel_scores:
        st.info("HBCI: no significant channels found for this patient.")
        return

    n_sig_channels = len(channel_scores)
    n_total_channels = max(len(ch_names), 1)
    channel_fraction = n_sig_channels / n_total_channels
    hbci_raw = sum(channel_scores.values())
    hbci = channel_fraction * hbci_raw

    sig_ch_list = list(channel_scores.keys())
    scores_arr = np.array([channel_scores[c] for c in sig_ch_list])
    min_p_arr = np.array([channel_min_p[c] for c in sig_ch_list])

    # Plot 1: Per-channel contribution bar chart
    st.subheader("HBCI – Channel Contributions")
    fig1, ax1 = plt.subplots(figsize=(max(6, len(sig_ch_list) * 0.6 + 2), 5))
    norm1 = plt.Normalize(vmin=0, vmax=0.05)
    cmap1 = plt.cm.RdYlGn_r
    bar_colors = [cmap1(norm1(p)) for p in min_p_arr]
    ax1.bar(sig_ch_list, scores_arr, color=bar_colors)
    sm1 = plt.cm.ScalarMappable(cmap=cmap1, norm=norm1)
    sm1.set_array([])
    fig1.colorbar(sm1, ax=ax1).set_label("Min p-value")
    ax1.set_title(f"HBCI Channel Contributions – {selected_pid}")
    ax1.set_xlabel("Channel")
    ax1.set_ylabel("Score (duration × -log₁₀(p))")
    plt.xticks(rotation=45, ha='right')
    plt.tight_layout()
    st.pyplot(fig1, use_container_width=True)
    plt.close(fig1)

    # Plot 2: Window heatmap: channels × time bins
    st.subheader("HBCI – Significant Window Heatmap")
    n_bins = 100
    t0, t1 = float(times[0]), float(times[-1])
    bin_edges = np.linspace(t0, t1, n_bins + 1)
    heatmap = np.zeros((len(sig_ch_list), n_bins))
    for ci, ch in enumerate(sig_ch_list):
        for w in patient_channel_sig[ch].get('windows', []):
            if w['p_value'] >= 0.05:
                continue
            log_p = -np.log10(max(w['p_value'], 1e-10))
            for bi in range(n_bins):
                if w['start'] < bin_edges[bi + 1] and w['end'] > bin_edges[bi]:
                    heatmap[ci, bi] = max(heatmap[ci, bi], log_p)
    fig2, ax2 = plt.subplots(figsize=(12, max(3, len(sig_ch_list) * 0.4 + 1.5)))
    im2 = ax2.imshow(
        heatmap, aspect='auto', origin='lower',
        cmap='hot_r', vmin=0, vmax=3,
        extent=[t0, t1, -0.5, len(sig_ch_list) - 0.5]
    )
    fig2.colorbar(im2, ax=ax2).set_label("-log₁₀(p)")
    ax2.set_yticks(range(len(sig_ch_list)))
    ax2.set_yticklabels(sig_ch_list)
    ax2.set_xlabel("Time (s)")
    ax2.set_ylabel("Channel")
    ax2.set_title(f"HBCI Significant Windows – {selected_pid}")
    plt.tight_layout()
    st.pyplot(fig2, use_container_width=True)
    plt.close(fig2)

    # Plot 3: HBCI summary breakdown
    st.subheader("Heart-Brain Coupling Index (HBCI)")
    fig3, (ax3l, ax3r) = plt.subplots(1, 2, figsize=(12, 4))
    tab20 = plt.cm.tab20
    left_offset = 0.0
    for ci, (ch, sc) in enumerate(channel_scores.items()):
        ax3l.barh("HBCI", sc, left=left_offset, color=tab20(ci % 20), label=ch)
        left_offset += sc
    ax3l.set_title("HBCI Score Breakdown")
    ax3l.set_xlabel("Score")
    ax3l.legend(loc='lower right', fontsize=7, ncol=2)
    ax3r.axis('off')
    ax3r.set_facecolor('#f0f8ff')
    fig3.patch.set_facecolor('#f0f8ff')
    ax3r.text(0.5, 0.62, f"HBCI = {hbci:.3f}",
              transform=ax3r.transAxes, ha='center', va='center',
              fontsize=36, fontweight='bold')
    ax3r.text(0.5, 0.25,
              f"Sig. channels: {n_sig_channels} / {n_total_channels}\n"
              f"Channel fraction: {channel_fraction:.3f}\n"
              f"Raw score: {hbci_raw:.3f}",
              transform=ax3r.transAxes, ha='center', va='center',
              fontsize=13, color='#333333')
    plt.tight_layout()
    st.pyplot(fig3, use_container_width=True)
    plt.close(fig3)

    st.metric("HBCI Score", f"{hbci:.4f}")


def handle_single_patient_view(individuals, selected_group, selected_stage, base_path,
                               n_permutations=100, p_threshold=0.05, jitter_sec=0.1):
    """
    Handles the selection and plotting of all channels for a single patient,
    plus a raw ECG viewer with time controls.
    """
    patient_ids = [ind[0] for ind in individuals]
    selected_pid = st.selectbox("Select Patient for All Channels", patient_ids)
    
    # --- Raw ECG Viewer ---
    st.markdown("### Raw ECG Viewer")
    
    # Locate and load the file for the selected patient
    # We need to find the file that starts with selected_pid in the group folder
    group_dir = os.path.join(base_path, selected_group, selected_stage)
    ecg_raw_data = None
    ecg_sfreq = None
    
    raw_obj = None
    if os.path.exists(group_dir):
        # Find file starting with patient_id
        # Note: ind[0] might be just the base ID, but filenames might be longer.
        # Let's assume filenames start with the ID.
        potential_files = [f for f in os.listdir(group_dir) if f.startswith(selected_pid) and f.endswith('.pkl')]
        if potential_files:
            file_path = os.path.join(group_dir, potential_files[0])
            try:
                with open(file_path, 'rb') as f:
                    raw_obj = pickle.load(f)
                    
                ecg_sfreq = raw_obj.info['sfreq']
                ch_names = raw_obj.ch_names
                ch_lower = [ch.lower() for ch in ch_names]
                ecg_indices = [i for i, ch in enumerate(ch_lower) if 'ecg' in ch or 'ekg' in ch]
                
                if ecg_indices:
                    ecg_idx = ecg_indices[0]
                    # Get all data for this channel
                    ecg_raw_data = raw_obj.get_data(picks=[ecg_idx])[0] * 1e6 # Convert to uV
                else:
                    st.warning("No ECG channel found in the raw file.")
                    
            except Exception as e:
                st.error(f"Error loading raw file for ECG viewer: {e}")
        else:
             st.warning(f"Could not find source file for {selected_pid} to display raw ECG.")
             
    if ecg_raw_data is not None and ecg_sfreq is not None:
        duration_sec = len(ecg_raw_data) / ecg_sfreq
        
        col1, col2 = st.columns(2)
        with col1:
            start_time = st.number_input("Start Time (s)", min_value=0.0, max_value=duration_sec, value=0.0, step=1.0)
        with col2:
            view_duration = st.slider("View Duration (s)", min_value=1, max_value=60, value=10)
            
        # Plot
        start_sample = int(start_time * ecg_sfreq)
        end_sample = min(len(ecg_raw_data), start_sample + int(view_duration * ecg_sfreq))
        
        # Add plot mode selector
        plot_mode = st.radio("ECG Plot Mode", ["Overlay", "Separate Subplots"], horizontal=True)

        segment = ecg_raw_data[start_sample:end_sample]
        segment_times = np.linspace(start_time, start_time + len(segment)/ecg_sfreq, len(segment))
        
        # Buffer and Clean logic
        clean_segment = None
        # Initialize fixed_segment to match raw segment length initially (or None)
        fixed_segment = None 
        
        try:
             # Add buffer
            buffer_samples = int(1.0 * ecg_sfreq) # 1 sec buffer
            start_buf = max(0, start_sample - buffer_samples)
            end_buf = min(len(ecg_raw_data), end_sample + buffer_samples)
            
            segment_buf = ecg_raw_data[start_buf:end_buf]
            
            # 1. Fix Inverted
            fixed_segment_buf, fit_info = fix_inverted_ecg(segment_buf, fs=ecg_sfreq)

            # 2. Clean based on fixed signal
            clean_segment_buf, _ = clean_ecg_high_fidelity(fixed_segment_buf, sampling_rate=ecg_sfreq)
            
            # Crop back to original window
            crop_start = start_sample - start_buf
            crop_end = crop_start + (end_sample - start_sample)
            
            clean_segment = clean_segment_buf[crop_start:crop_end]
            fixed_segment = fixed_segment_buf[crop_start:crop_end]
            
        except Exception as e:
            st.warning(f"Could not clean ECG segment: {e}")

        # --- FFT Frequency Spectrum ---
        st.markdown("#### FFT Frequency Spectrum")
        try:
            from scipy.fft import rfft, rfftfreq

            nyquist = ecg_sfreq / 2
            fft_xlim = st.slider(
                "FFT Frequency Range (Hz)",
                min_value=0.0,
                max_value=float(min(500, nyquist)),
                value=(0.0, 10.0),
                step=0.5,
                key="fft_xlim_slider",
            )

            n = len(segment)
            freqs = rfftfreq(n, d=1.0 / ecg_sfreq)

            # Raw FFT
            fft_raw = np.abs(rfft(segment)) / n * 2  # one-sided amplitude

            fig_fft, ax_fft = plt.subplots(figsize=(12, 4))
            ax_fft.semilogy(freqs, fft_raw, color='black', linewidth=1, alpha=0.8, label='Raw ECG')

            if clean_segment is not None and len(clean_segment) == len(segment):
                fft_clean = np.abs(rfft(clean_segment)) / n * 2
                ax_fft.semilogy(freqs, fft_clean, color='red', linewidth=1, alpha=0.8, label='Cleaned ECG')

            ax_fft.set_xlim(fft_xlim[0], fft_xlim[1])
            ax_fft.set_title(f"FFT Spectrum - {selected_pid} (window: {start_time:.1f}–{start_time + view_duration:.1f} s)")
            ax_fft.set_xlabel("Frequency (Hz)")
            ax_fft.set_ylabel("Amplitude (μV, log scale)")
            ax_fft.legend(loc='upper right')
            ax_fft.grid(True, which='both', alpha=0.3)
            st.pyplot(fig_fft, use_container_width=True)
        except Exception as e:
            st.warning(f"Could not compute FFT: {e}")

        # Plotting based on mode
        if plot_mode == "Overlay":
            fig_ecg_raw, ax_ecg = plt.subplots(figsize=(12, 4))
            ax_ecg.plot(segment_times, segment, color='black', linewidth=1, label='Raw ECG', alpha=0.5)
            
            if fixed_segment is not None and len(fixed_segment) == len(segment):
                 # Plot fixed signal if valid
                 ax_ecg.plot(segment_times, fixed_segment, color='blue', linestyle='--', linewidth=1, alpha=0.7, label='Inverted (Fixed)')

            if clean_segment is not None and len(clean_segment) == len(segment):
                 ax_ecg.plot(segment_times, clean_segment, color='red', linewidth=1, alpha=0.8, label='Cleaned ECG')
            
            ax_ecg.set_title(f"Raw vs Fixed vs Cleaned ECG - {selected_pid} (Start: {start_time}s)")
            ax_ecg.set_xlabel("Time (s)")
            ax_ecg.set_ylabel("Amplitude (μV)")
            ax_ecg.grid(True, alpha=0.3)
            ax_ecg.legend(loc='upper right')
            st.pyplot(fig_ecg_raw, use_container_width=True)

            # --- Cleaned ECG with WFDB Robust R-peaks ---
            if clean_segment is not None:
                st.markdown("#### Cleaned ECG with WFDB Robust R-peaks")
                try:
                    # Detect R-peaks on the cleaned segment
                    rpeaks_loc = detect_rpeaks_robust(clean_segment, ecg_sfreq)
                    
                    fig_rpeaks, ax_rpeaks = plt.subplots(figsize=(12, 4))
                    ax_rpeaks.plot(segment_times, clean_segment, color='red', linewidth=1, label='Cleaned ECG')
                    
                    if len(rpeaks_loc) > 0:
                        # rpeaks_loc are indices relative to start of clean_segment
                        rpeak_times = segment_times[rpeaks_loc]
                        rpeak_amps = clean_segment[rpeaks_loc]
                        
                        ax_rpeaks.plot(rpeak_times, rpeak_amps, 'bo', label='WFDB Robust R-peaks')
                        
                        for rt in rpeak_times:
                             ax_rpeaks.axvline(rt, color='blue', linestyle='--', alpha=0.3)
                             
                        # Highlight exactly 500ms intervals
                        if len(rpeak_times) > 1:
                            rpeak_diffs = np.diff(rpeak_times)
                            # Using 3 decimal points for 500ms (0.500s) precision
                            exact_500_mask = np.round(rpeak_diffs, 3) == 0.500
                            
                            # Find all peaks that are part of a 500ms interval
                            peak_is_500 = np.zeros(len(rpeak_times), dtype=bool)
                            peak_is_500[:-1] |= exact_500_mask
                            peak_is_500[1:] |= exact_500_mask
                            
                            if np.any(peak_is_500):
                                ax_rpeaks.plot(rpeak_times[peak_is_500], rpeak_amps[peak_is_500], 
                                               'ro', markersize=12, fillstyle='none', markeredgewidth=2, 
                                               label='Exactly 500ms Apart')
                                               
                        # Calculate Median, MAD and BPM for the title
                        med = np.median(clean_segment)
                        mad = np.median(np.abs(clean_segment - med))
                        if len(rpeak_times) > 1:
                            avg_rr = np.mean(np.diff(rpeak_times))
                            bpm = 60.0 / avg_rr if avg_rr > 0 else 0
                            title_str = f"Cleaned ECG with WFDB Robust R-peaks - {selected_pid} (Median: {med:.2f}μV, MAD: {mad:.2f}μV, ~{bpm:.0f} BPM)"
                        else:
                            title_str = f"Cleaned ECG with WFDB Robust R-peaks - {selected_pid} (Median: {med:.2f}μV, MAD: {mad:.2f}μV)"
                    else:
                        st.info("No R-peaks detected by WFDB in this segment.")
                        if clean_segment is not None:
                            med = np.median(clean_segment)
                            mad = np.median(np.abs(clean_segment - med))
                        else:
                            med, mad = 0, 0
                        title_str = f"Cleaned ECG with WFDB Robust R-peaks - {selected_pid} (Median: {med:.2f}μV, MAD: {mad:.2f}μV)"
                        
                    ax_rpeaks.set_title(title_str)
                    ax_rpeaks.set_xlabel("Time (s)")
                    ax_rpeaks.set_ylabel("Amplitude (μV)")
                    ax_rpeaks.grid(True, alpha=0.3)
                    ax_rpeaks.legend(loc='upper right')
                    st.pyplot(fig_rpeaks, use_container_width=True)
                except Exception as e:
                    st.warning(f"Error detecting R-peaks with WFDB: {e}")


        else: # Separate Subplots
            fig_ecg_raw, axes = plt.subplots(3, 1, figsize=(12, 9), sharex=True)
            
            # Raw
            axes[0].plot(segment_times, segment, color='black', linewidth=1)
            axes[0].set_title("Raw ECG")
            axes[0].set_ylabel("Amp (μV)")
            axes[0].grid(True, alpha=0.3)
            
            # Fixed
            if fixed_segment is not None and len(fixed_segment) == len(segment):
                axes[1].plot(segment_times, fixed_segment, color='blue', linewidth=1)
            else:
                axes[1].text(0.5, 0.5, "Processing Failed", ha='center', va='center')
            axes[1].set_title("Inverted (Fixed) ECG")
            axes[1].set_ylabel("Amp (μV)")
            axes[1].grid(True, alpha=0.3)

            # Cleaned
            if clean_segment is not None and len(clean_segment) == len(segment):
                axes[2].plot(segment_times, clean_segment, color='red', linewidth=1)
            else:
                axes[2].text(0.5, 0.5, "Cleaning Failed", ha='center', va='center')
                
            axes[2].set_title("Cleaned ECG")
            axes[2].set_ylabel("Amp (μV)")
            axes[2].set_xlabel("Time (s)")
            axes[2].grid(True, alpha=0.3)
            
            plt.tight_layout()
            st.pyplot(fig_ecg_raw, use_container_width=True)

        # --- ECG Cleaning Comparison ---
        st.markdown(f"### ECG Cleaning Methods Comparison ({selected_pid})")
        if st.checkbox("Show ECG Cleaning Comparison", value=False):
            if clean_segment is not None:
                # Plot the comparison using the pre-cleaned segment and raw segment
                plot_ecg_cleaning_comparison(segment, clean_segment, ecg_sfreq, selected_pid)
            else:
                 st.warning(f"Could not load source file for {selected_pid} to display cleaning comparison.")
                 
        # --- Raw EEG Viewer ---
        st.markdown(f"**Raw EEG Viewer ({selected_pid})**")
        # Filter for EEG channels only (exclude ECG, EOG, etc.)
        eeg_channels = [ch for ch in ch_names if not any(x in ch.lower() for x in ['ecg', 'ekg', 'eog', 'emg', 'resp'])]
        selected_eeg_ch = st.selectbox("Select EEG Electrode", eeg_channels)
        
        if selected_eeg_ch:
             try:
                 eeg_ch_idx = ch_names.index(selected_eeg_ch)
                 eeg_raw_data = raw_obj.get_data(picks=[eeg_ch_idx])[0] * 1e6
                 
                 eeg_segment = eeg_raw_data[start_sample:end_sample]
                 
                 fig_eeg_raw, ax_eeg_raw = plt.subplots(figsize=(12, 4))
                 ax_eeg_raw.plot(segment_times, eeg_segment, color='blue', linewidth=1)
                 ax_eeg_raw.set_title(f"Raw EEG - {selected_eeg_ch} (Time Aligned)")
                 ax_eeg_raw.set_xlabel("Time (s)")
                 ax_eeg_raw.set_ylabel("Amplitude (μV)")
                 ax_eeg_raw.grid(True, alpha=0.3)
                 st.pyplot(fig_eeg_raw, use_container_width=True)
                 
             except Exception as e:
                 st.error(f"Error extracting EEG channel {selected_eeg_ch}: {e}")



        # --- HRV Analysis (RR Intervals) ---
        st.markdown(f"**HRV Analysis (RR Intervals) - {selected_pid}**")
        
        # We need to get the rpeaks for this patient.
        # Check individuals list for the selected patient
        rpeaks_current = None
        for ind_chk in individuals:
            if ind_chk[0] == selected_pid:
                # ind structure: (patient_id, hep_data, times, ch_names, rpeaks, ecg_hep_data, ecg_ch_names)
                rpeaks_current = ind_chk[4]
                break
        
        if rpeaks_current is not None and len(rpeaks_current) > 1:
            # Calculate RR intervals in ms
            # fs is needed. Use ecg_sfreq from raw file above, or assume it's same for group.
            if ecg_sfreq is None:
                 pass

            if ecg_sfreq:
                rr_intervals_samples = np.diff(rpeaks_current)
                rr_intervals_ms = (rr_intervals_samples / ecg_sfreq) * 1000
                
                # Filter outliers (550ms to 1300ms)
                valid_rr = rr_intervals_ms[(rr_intervals_ms >= 0) & (rr_intervals_ms <= 1300)]
                n_excluded = len(rr_intervals_ms) - len(valid_rr)
                
                if len(valid_rr) > 0:
                    rr_mean = np.mean(valid_rr)
                    rr_std = np.std(valid_rr)
                    rr_median = np.median(valid_rr)
                    hr_mean = 60000 / rr_mean
                    
                    st.write(f"**Mean RR:** {rr_mean:.1f} ms | **Median RR:** {rr_median:.1f} ms | **SDNN:** {rr_std:.1f} ms | **Mean HR:** {hr_mean:.1f} bpm")
                    if n_excluded > 0:
                        st.caption(f"(Excluded {n_excluded} intervals outside 550-1300ms range)")
                    
                    fig_hrv, ax_hrv = plt.subplots(figsize=(10, 4))
                    ax_hrv.hist(valid_rr, bins=30, color='purple', alpha=0.7, edgecolor='black')
                    ax_hrv.set_title(f"RR Interval Distribution - {selected_pid} (550-1300ms) - {hr_mean:.1f} BPM")
                    ax_hrv.set_xlabel("RR Interval (ms)")
                    ax_hrv.set_ylabel("Count")
                    ax_hrv.set_xlim(400, 1300)
                    ax_hrv.grid(True, alpha=0.3)
                    
                    # Add vertical line for mean
                    ax_hrv.axvline(rr_mean, color='red', linestyle='dashed', linewidth=2, label=f'Mean: {rr_mean:.1f}ms')
                    ax_hrv.legend()
                    
                    st.pyplot(fig_hrv, use_container_width=True)
                else:
                    st.warning("No valid RR intervals found (all outside 550-1300ms).")
            else:
                 st.warning("Sampling rate unknown (file load failed), cannot calculate RR intervals in ms.")
        else:
            st.warning("Not enough R-peaks to calculate HRV.")

    for ind in individuals:
        if ind[0] == selected_pid:
            # Determine grid size
            ch_names = ind[3]
            n_channels = len(ch_names)
            n_cols = 5
            n_rows = int(np.ceil(n_channels / n_cols))
            
            fig, axes = plt.subplots(n_rows, n_cols, figsize=(20, 4 * n_rows), sharex=True, sharey=True)
            axes = axes.flatten()
            
            times = ind[2]
            full_hep = ind[1] # (n_ch, n_times)
            
            # Get ECG HEP data if available
            ecg_hep = ind[5] # (1, n_times)
            
            # remove LOC ROC channels
            full_hep = full_hep[:-2]
            ch_names = ch_names[:-2]
            # remove ECG channels
            non_ecg_mask = [i for i, ch in enumerate(ch_names) if not any(x in ch.lower() for x in ['ecg', 'ekg'])]
            full_hep = full_hep[non_ecg_mask]
            ch_names = [ch_names[i] for i in non_ecg_mask]

            patient_channel_sig = {}
            patient_window_count = None
            if raw_obj is not None and PYNAPPLE_AVAILABLE:
                try:
                    raw_sig, sfreq_sig, rpeak_ts_sig, _, minmax_sig, _ = process_file_data(
                        drop_non_eeg_channels(raw_obj.copy()),
                        selected_pid
                    )
                    target_channels = [ch for ch in ch_names if ch in raw_sig.ch_names]
                    if target_channels:
                        target_indices = [raw_sig.ch_names.index(ch) for ch in target_channels]
                        tsd_sig = nap.TsdFrame(
                            t=raw_sig.times,
                            d=raw_sig.get_data(picks=target_indices).T,
                            columns=target_channels
                        )
                        perievent_sig = nap.compute_perievent_continuous(
                            tsd_sig, rpeak_ts_sig, minmax=minmax_sig
                        )
                        perievent_vals = np.asarray(perievent_sig.values)
                        times_sig = np.asarray(perievent_sig.t)
                        if perievent_vals.ndim == 2:
                            perievent_vals = perievent_vals[:, :, np.newaxis]

                        _sig_progress = st.progress(0, text=f"Computing significance: channel 0/{len(target_channels)}")
                        for sig_idx, ch_name_sig in enumerate(target_channels):
                            _sig_progress.progress(
                                (sig_idx + 1) / len(target_channels),
                                text=f"Computing significance: {ch_name_sig} ({sig_idx + 1}/{len(target_channels)})"
                            )
                            ch_epochs = np.asarray(perievent_vals[:, :, sig_idx]).T
                            if ch_epochs.ndim != 2 or ch_epochs.shape[0] < 2 or ch_epochs.shape[1] < 2:
                                continue
                            valid_epoch_mask = ~np.all(np.isnan(ch_epochs), axis=1)
                            ch_epochs = ch_epochs[valid_epoch_mask]
                            if ch_epochs.shape[0] < 2:
                                continue
                            sig_windows_ch, _, per_pt_info_ch = permutation_cluster_jitter_test(
                                ch_epochs,
                                times_sig,
                                n_permutations=n_permutations,
                                p_threshold=p_threshold,
                                jitter_sec=jitter_sec,
                            )
                            patient_channel_sig[ch_name_sig] = {
                                'windows': sig_windows_ch,
                                'fisher_p': per_pt_info_ch.get('fisher_p', 1.0),
                                'n_sig_epochs': per_pt_info_ch.get('n_significant', 0),
                                'n_epochs': ch_epochs.shape[0],
                            }
                        _sig_progress.empty()
                        if patient_channel_sig:
                            patient_window_count = max(v.get('n_epochs', 0) for v in patient_channel_sig.values())
                except Exception as e:
                    st.warning(f"Could not compute single-patient per-channel significance: {e}")
            # Determine symmetric limits for alignment
            # EEG Limits
            max_eeg = np.nanmax(np.abs(full_hep * 1e6)) * 1.1 # 10% margin
            ylim_eeg = (-max_eeg, max_eeg)

            # ECG Limits
            if ecg_hep is not None:
                max_ecg = np.nanmax(np.abs(ecg_hep[0] * 1e6)) * 1.1 # 10% margin
                ylim_ecg = (-max_ecg, max_ecg)
            
            for i, (ch_name, ch_data) in enumerate(zip(ch_names, full_hep)):
                ax = axes[i]
                ax.plot(times, ch_data * 1e6)
                ch_skew = stats.skew(ch_data)
                sig_info = patient_channel_sig.get(ch_name, {})
                sig_windows = sig_info.get('windows', [])
                if sig_windows:
                    _annotate_sig_windows(ax, sig_windows, y_pad_frac=0.08)
                    min_p_ch = min(w['p_value'] for w in sig_windows)
                    p_tag = 'p<0.001' if min_p_ch < 0.001 else f"p={min_p_ch:.3f}"
                else:
                    p_tag = "n.s."
                ax.set_title(f"{ch_name}\nskew={ch_skew:.2f} | {p_tag}")
                ax.grid(True)
                ax.axvline(0, color='r', linestyle='--', alpha=0.5)
                ax.set_ylim(ylim_eeg)
                
                # Add ECG HEP on secondary axis
                if ecg_hep is not None:
                    ax2 = ax.twinx()
                    # ecg_hep is (1, n_times), so take [0]
                    ax2.plot(times, ecg_hep[0] * 1e6, color='gray', linestyle='--', alpha=0.7,linewidth=1)
                    ax2.set_ylim(ylim_ecg)
                    
                    # Only show right y-axis labels for the rightmost column to reduce clutter
                    if (i + 1) % n_cols == 0:
                        ax2.set_ylabel("ECG (μV)", color='gray', fontsize=8)
                    else:
                        ax2.set_yticklabels([])
                    
                    ax2.tick_params(axis='y', labelcolor='gray', labelsize=8)

                if i % n_cols == 0:
                    ax.set_ylabel("Amp (μV)")
                if i >= n_channels - n_cols:
                    ax.set_xlabel("Time (s)")
                    
            # Hide unused subplots
            for j in range(i + 1, len(axes)):
                axes[j].axis('off')
            
            st.divider()
            st.title(f"All Channels for {selected_pid} - {selected_group} {selected_stage}")
            if patient_window_count is not None:
                n_sig_channels = sum(1 for v in patient_channel_sig.values() if v.get('windows'))
                st.caption(
                    f"Orange spans mark single-patient significant windows computed across that patient's "
                    f"{patient_window_count} heartbeat windows with permutation jitter "
                    f"(channels significant: {n_sig_channels}/{len(ch_names)})."
                )
            fig.suptitle(f"All Channels - {selected_pid} - {selected_group} {selected_stage}", fontsize=16)
            plt.tight_layout()
            st.pyplot(fig, use_container_width=True)



            # --- Butterfly Plot (Z-Scored) ---
            st.subheader(f"All Channels Z-Scored ({selected_pid})")
            fig_z, ax_z = plt.subplots(figsize=(14, 6)) # Increased width for legend

            # Compute Common Average Reference (CAR)
            # Use only EEG channels for CAR (exclude potential artifact channels)
            # Filter: strictly exclude ECG/EOG/EMG and require digit or 'z' suffix
            eeg_indices_car = [
                idx for idx, name in enumerate(ch_names)
                if not any(excl in name.upper() for excl in ['ECG', 'EKG', 'EOG', 'EMG', 'RESP', 'PULSE'])
                and (name[-1].isdigit() or name.lower().endswith('z'))
            ]
            
            if eeg_indices_car:
                car_signal = np.mean(full_hep[eeg_indices_car], axis=0)
            else:
                # Fallback if no specific EEG channels found
                car_signal = np.mean(full_hep, axis=0)

            for i, (ch_name, ch_data) in enumerate(zip(ch_names, full_hep)):
                # Filter: last char must be digit
                if not ch_name[-1].isdigit():
                    continue
                    
                last_digit = int(ch_name[-1])
                
                # Color logic: Even=Blue, Odd=Red
                if last_digit % 2 == 0:
                    color = 'blue'
                else:
                    color = 'red'

                # Subtract CAR
                ch_data_car = ch_data - car_signal

                # Calculate Z-score per channel on CAR-subtracted data
                if np.std(ch_data_car) > 0:
                    z_data = (ch_data_car - np.mean(ch_data_car)) / np.std(ch_data_car)
                else:
                    z_data = np.zeros_like(ch_data_car)
                    
                ax_z.plot(times, z_data, color=color, alpha=0.3, linewidth=1, label=ch_name)
            
            # Plot CAR signal (Z-scored)
            if np.std(car_signal) > 0:
                z_car = (car_signal - np.mean(car_signal)) / np.std(car_signal)
                ax_z.plot(times, z_car, color='black', alpha=0.8, linewidth=2, linestyle='--', label='CAR (EEG only)')

            ax_z.set_title(f"Butterfly Plot - Z-Scored (CAR Subtracted, Even=Blue, Odd=Red) - {selected_pid}")
            ax_z.set_xlabel("Time (s)")
            ax_z.set_ylabel("Z-Score")
            ax_z.axvline(0, color='r', linestyle='--', alpha=0.8)
            ax_z.grid(True, alpha=0.3)
            
            # Add legend outside
            ax_z.legend(bbox_to_anchor=(1.05, 1), loc="upper left", borderaxespad=0., fontsize='small', ncol=1)
            
            st.pyplot(fig_z, use_container_width=True)

            # --- Butterfly Plot (Raw uV) with ECG ---
            st.subheader(f"Butterfly Plot - Raw uV ({selected_pid})")
            fig_raw, ax_raw = plt.subplots(figsize=(14, 6))

            # Plot EEG channels
            for i, (ch_name, ch_data) in enumerate(zip(ch_names, full_hep)):
                 # Filter: last char must be digit
                if not ch_name[-1].isdigit():
                    continue
                last_digit = int(ch_name[-1])
                color = 'blue' if last_digit % 2 == 0 else 'red'
                
                ax_raw.plot(times, ch_data * 1e6, color=color, alpha=0.3, linewidth=1)

            ax_raw.set_xlabel("Time (s)")
            ax_raw.set_ylabel("EEG Amplitude (μV)")
            ax_raw.axvline(0, color='r', linestyle='--', alpha=0.8)
            ax_raw.grid(True, alpha=0.3)

            # Secondary Axis for ECG
            if ecg_hep is not None:
                ax_ecg = ax_raw.twinx()
                ax_ecg.plot(times, ecg_hep[0] * 1e6, color='green', alpha=0.6, linewidth=2, linestyle='--', label='ECG')
                ax_ecg.set_ylabel("ECG Amplitude (μV)", color='green')
                ax_ecg.tick_params(axis='y', labelcolor='green')
                # Add legend for ECG
                ax_ecg.legend(loc='upper right')

            ax_raw.set_title(f"Butterfly Plot - Raw uV (Even=Blue, Odd=Red) - {selected_pid}")
            
            # Add legend for Even/Odd EEG (optional, purely visual)
            from matplotlib.lines import Line2D
            custom_lines = [Line2D([0], [0], color='blue', lw=2),
                            Line2D([0], [0], color='red', lw=2)]
            ax_raw.legend(custom_lines, ['Even Channels', 'Odd Channels'], loc='upper left')

            st.pyplot(fig_raw, use_container_width=True)

            # --- Butterfly Plot - All Channels Colored by Region ---
            st.subheader(f"Butterfly Plot - Region Colored ({selected_pid})")
            fig_reg_all, ax_reg_all = plt.subplots(figsize=(14, 6))

            region_color_map = {
                'F': 'blue',
                'C': 'green',
                'T': 'red',
                'P': 'purple',
                'O': 'orange'
            }
            
            # Helper to determine region
            def get_region_color(ch_name):
                name_upper = ch_name.upper()
                # Check regions in specific order to handle overlaps if any (e.g. FC -> Frontal-Central)
                # Typically F first checks F, FC, Fp. C checks C, CP.
                # Simplest heuristic: check letter presence.
                for reg, col in region_color_map.items():
                   if reg in name_upper:
                       return col, reg
                return 'gray', 'Other'

            # Plot EEG channels
            used_regions = set()
            for i, (ch_name, ch_data) in enumerate(zip(ch_names, full_hep)):
                color, region = get_region_color(ch_name)
                ax_reg_all.plot(times, ch_data * 1e6, color=color, alpha=0.3, linewidth=1)
                used_regions.add(region)

            ax_reg_all.set_xlabel("Time (s)")
            ax_reg_all.set_ylabel("EEG Amplitude (μV)")
            ax_reg_all.axvline(0, color='r', linestyle='--', alpha=0.8)
            ax_reg_all.grid(True, alpha=0.3)

            # Secondary Axis for ECG
            if ecg_hep is not None:
                ax_ecg_reg = ax_reg_all.twinx()
                ax_ecg_reg.plot(times, ecg_hep[0] * 1e6, color='black', alpha=0.6, linewidth=2, linestyle='--', label='ECG')
                ax_ecg_reg.set_ylabel("ECG Amplitude (μV)", color='black')
                ax_ecg_reg.tick_params(axis='y', labelcolor='black')
                ax_ecg_reg.legend(loc='upper right')

            ax_reg_all.set_title(f"Butterfly Plot - Region Colored - {selected_pid}")
            
            # Custom legend for regions
            custom_lines_reg = [Line2D([0], [0], color=region_color_map[reg], lw=2) for reg in region_color_map if reg in used_regions]
            custom_labels_reg = [f"Region {reg}" for reg in region_color_map if reg in used_regions]
            if 'Other' in used_regions:
                custom_lines_reg.append(Line2D([0], [0], color='gray', lw=2))
                custom_labels_reg.append("Region Other")
            
            ax_reg_all.legend(custom_lines_reg, custom_labels_reg, loc='upper left')

            st.pyplot(fig_reg_all, use_container_width=True)

            # --- Butterfly Plot - Odd Channels (Region Colored) ---
            st.subheader(f"Butterfly Plot - Odd Channels ({selected_pid})")
            fig_odd, ax_odd = plt.subplots(figsize=(14, 6))
            
            used_regions_odd = set()
            has_odd = False
            for i, (ch_name, ch_data) in enumerate(zip(ch_names, full_hep)):
                if not ch_name[-1].isdigit():
                    continue
                if int(ch_name[-1]) % 2 == 0:
                    continue # Skip even

                color, region = get_region_color(ch_name)
                ax_odd.plot(times, ch_data * 1e6, color=color, alpha=0.3, linewidth=1)
                used_regions_odd.add(region)
                has_odd = True

            if has_odd:
                ax_odd.set_xlabel("Time (s)")
                ax_odd.set_ylabel("EEG Amplitude (μV)")
                ax_odd.axvline(0, color='r', linestyle='--', alpha=0.8)
                ax_odd.grid(True, alpha=0.3)

                if ecg_hep is not None:
                    ax_ecg_odd = ax_odd.twinx()
                    ax_ecg_odd.plot(times, ecg_hep[0] * 1e6, color='black', alpha=0.6, linewidth=2, linestyle='--', label='ECG')
                    ax_ecg_odd.set_ylabel("ECG Amplitude (μV)", color='black')
                    ax_ecg_odd.tick_params(axis='y', labelcolor='black')
                    ax_ecg_odd.legend(loc='upper right')

                ax_odd.set_title(f"Butterfly Plot - Odd Channels - {selected_pid}")
                
                 # Legend
                custom_lines_odd = [Line2D([0], [0], color=region_color_map[reg], lw=2) for reg in region_color_map if reg in used_regions_odd]
                custom_labels_odd = [f"Region {reg}" for reg in region_color_map if reg in used_regions_odd]
                if 'Other' in used_regions_odd:
                    custom_lines_odd.append(Line2D([0], [0], color='gray', lw=2))
                    custom_labels_odd.append("Region Other")
                
                ax_odd.legend(custom_lines_odd, custom_labels_odd, loc='upper left')
                
                st.pyplot(fig_odd, use_container_width=True)
            else:
                plt.close(fig_odd)

            # --- Butterfly Plot - Even Channels (Region Colored) ---
            st.subheader(f"Butterfly Plot - Even Channels ({selected_pid})")
            fig_even, ax_even = plt.subplots(figsize=(14, 6))
            
            used_regions_even = set()
            has_even = False
            for i, (ch_name, ch_data) in enumerate(zip(ch_names, full_hep)):
                if not ch_name[-1].isdigit():
                    continue
                if int(ch_name[-1]) % 2 != 0:
                    continue # Skip odd

                color, region = get_region_color(ch_name)
                ax_even.plot(times, ch_data * 1e6, color=color, alpha=0.3, linewidth=1)
                used_regions_even.add(region)
                has_even = True

            if has_even:
                ax_even.set_xlabel("Time (s)")
                ax_even.set_ylabel("EEG Amplitude (μV)")
                ax_even.axvline(0, color='r', linestyle='--', alpha=0.8)
                ax_even.grid(True, alpha=0.3)

                if ecg_hep is not None:
                    ax_ecg_even = ax_even.twinx()
                    ax_ecg_even.plot(times, ecg_hep[0] * 1e6, color='black', alpha=0.6, linewidth=2, linestyle='--', label='ECG')
                    ax_ecg_even.set_ylabel("ECG Amplitude (μV)", color='black')
                    ax_ecg_even.tick_params(axis='y', labelcolor='black')
                    ax_ecg_even.legend(loc='upper right')

                ax_even.set_title(f"Butterfly Plot - Even Channels - {selected_pid}")

                # Legend
                custom_lines_even = [Line2D([0], [0], color=region_color_map[reg], lw=2) for reg in region_color_map if reg in used_regions_even]
                custom_labels_even = [f"Region {reg}" for reg in region_color_map if reg in used_regions_even]
                if 'Other' in used_regions_even:
                    custom_lines_even.append(Line2D([0], [0], color='gray', lw=2))
                    custom_labels_even.append("Region Other")
                
                ax_even.legend(custom_lines_even, custom_labels_even, loc='upper left')

                st.pyplot(fig_even, use_container_width=True)
            else:
                plt.close(fig_even)

            # --- Regional Butterfly Plots ---
            regions = ['F', 'T', 'P', 'O']
            for region in regions:
                # Find channels containing the region letter (e.g., 'F' in 'Fp1', 'F3', 'Fz')
                region_indices = [idx for idx, name in enumerate(ch_names) if region in name.upper()]
                
                if not region_indices:
                    continue

                st.subheader(f"Butterfly Plot - Region {region} ({selected_pid})")
                fig_reg, ax_reg = plt.subplots(figsize=(14, 6))
                
                # Setup colormap
                cmap = plt.get_cmap('tab10')
                has_channels = False
                
                # Collect handles for legend
                region_handles = []
                region_labels = []

                for k, idx in enumerate(region_indices):
                    ch_name = ch_names[idx]
                    ch_data = full_hep[idx]
                    
                    # Use a unique color for each channel
                    color = cmap(k % 10)
                    
                    line, = ax_reg.plot(times, ch_data * 1e6, color=color, alpha=0.6, linewidth=1.5, label=ch_name)
                    has_channels = True
                    region_handles.append(line)
                    region_labels.append(ch_name)

                if not has_channels:
                    plt.close(fig_reg)
                    continue

                ax_reg.set_xlabel("Time (s)")
                ax_reg.set_ylabel("EEG Amplitude (μV)")
                ax_reg.axvline(0, color='r', linestyle='--', alpha=0.8)
                ax_reg.grid(True, alpha=0.3)
                
                # Secondary Axis for ECG
                if ecg_hep is not None:
                    ax_ecg_reg = ax_reg.twinx()
                    ecg_line, = ax_ecg_reg.plot(times, ecg_hep[0] * 1e6, color='black', alpha=0.4, linewidth=2, linestyle='--', label='ECG')
                    ax_ecg_reg.set_ylabel("ECG Amplitude (μV)", color='black')
                    ax_ecg_reg.tick_params(axis='y', labelcolor='black')
                    
                    # Add ECG to legend
                    region_handles.append(ecg_line)
                    region_labels.append('ECG')

                ax_reg.set_title(f"Butterfly Plot - Region {region} - {selected_pid}")
                
                # Add legend
                ax_reg.legend(region_handles, region_labels, loc='upper left', bbox_to_anchor=(1.05, 1), borderaxespad=0.)
                
                st.pyplot(fig_reg, use_container_width=True)

            # --- Single Electrode Window Explorer ---
            st.subheader(f"Electrode Window Explorer ({selected_pid})")

            eeg_ch_options = [ch for ch in ch_names if not any(x in ch.lower() for x in ['ecg', 'ekg', 'eog', 'emg', 'resp'])]
            col_elec, col_nwin = st.columns(2)
            with col_elec:
                selected_elec = st.selectbox("Select Electrode", eeg_ch_options, key="window_explorer_elec")
            with col_nwin:
                n_windows_show = st.slider("Number of Windows to Show", min_value=1, max_value=100, value=20, key="window_explorer_n")

            raw_obj_available = False
            try:
                _ = raw_obj
                raw_obj_available = True
            except NameError:
                pass

            if raw_obj_available and ecg_sfreq is not None and selected_elec:
                try:
                    rpeaks_ind = ind[4]
                    rpeak_times_ind = rpeaks_ind / ecg_sfreq
                    rpeak_ts_ind = nap.Ts(t=rpeak_times_ind, time_units="s")
                    tmin_win = float(times[0])
                    tmax_win = float(times[-1])

                    raw_ch_names_list = raw_obj.ch_names
                    if selected_elec in raw_ch_names_list:
                        elec_idx = raw_ch_names_list.index(selected_elec)
                        elec_data = raw_obj.get_data(picks=[elec_idx]).T  # (n_times, 1)
                        tsd_elec = nap.TsdFrame(t=raw_obj.times, d=elec_data, columns=[selected_elec])
                        perievent_elec = nap.compute_perievent_continuous(tsd_elec, rpeak_ts_ind, minmax=(tmin_win, tmax_win))

                        pe_values = perievent_elec.values  # (n_times, n_events, 1) or (n_times, n_events)
                        pe_times = perievent_elec.t
                        if pe_values.ndim == 3:
                            pe_values = pe_values[:, :, 0]  # (n_times, n_events)

                        n_trials_total = pe_values.shape[1]
                        n_show = min(n_windows_show, n_trials_total)

                        fig_explorer, ax_explorer = plt.subplots(figsize=(14, 6))

                        for trial_idx in range(n_show):
                            ax_explorer.plot(pe_times, pe_values[:, trial_idx] * 1e6,
                                             color='steelblue', alpha=0.2, linewidth=0.8)

                        avg_signal = np.nanmean(pe_values[:, :n_show], axis=1)
                        ax_explorer.plot(pe_times, avg_signal * 1e6, color='navy', linewidth=2.5,
                                         label=f'Average (n={n_show})')

                        ax_explorer.axvline(0, color='r', linestyle='--', alpha=0.8)
                        ax_explorer.set_xlabel("Time (s)")
                        ax_explorer.set_ylabel("EEG Amplitude (μV)")
                        ax_explorer.set_title(f"Electrode {selected_elec} - {n_show} Windows - {selected_pid}")
                        ax_explorer.legend()
                        ax_explorer.grid(True, alpha=0.3)

                        st.caption(f"Showing {n_show} of {n_trials_total} available windows")
                        st.pyplot(fig_explorer, use_container_width=True)
                        plt.close(fig_explorer)
                    else:
                        st.warning(f"Electrode {selected_elec} not found in raw data.")
                except Exception as e:
                    st.error(f"Error computing per-window data: {e}")
            else:
                st.warning("Raw data not available for window-level analysis.")

            # ── Heart-Brain Coupling Index (HBCI) ──────────────────────────
            render_hbci(patient_channel_sig, ch_names, times, selected_pid)


def plot_patients_butterfly_comparison(individuals, selected_group, selected_stage):
    """
    Plots a grid of butterfly plots, one method per patient.
    """
    n_pats = len(individuals)
    n_cols = 2
    n_rows = int(np.ceil(n_pats / n_cols))
    
    fig, axes = plt.subplots(n_rows, n_cols, figsize=(16, 5 * n_rows), sharex=True, sharey=True)
    if n_pats == 1:
        axes = [axes]
    else:
        axes = axes.flatten()
        
    for idx, ind in enumerate(individuals):
        ax = axes[idx]
        patient_id = ind[0]
        full_hep = ind[1]
        times = ind[2]
        ch_names = ind[3]
        ecg_hep = ind[5]
        
        # Plot EEG channels
        for i, (ch_name, ch_data) in enumerate(zip(ch_names, full_hep)):
            # Filter: last char must be digit
            if not ch_name[-1].isdigit():
                continue
            last_digit = int(ch_name[-1])
            color = 'blue' if last_digit % 2 == 0 else 'red'
            
            ax.plot(times, ch_data * 1e6, color=color, alpha=0.3, linewidth=1)

        ax.grid(True, alpha=0.3)
        ax.set_ylim(-30, 30)
        ax.axvline(0, color='r', linestyle='--', alpha=0.8)
        
        # Secondary Axis for ECG
        if ecg_hep is not None:
            ax_ecg = ax.twinx()
            ax_ecg.plot(times, ecg_hep[0] * 1e6, color='green', alpha=0.6, linewidth=1.5, linestyle='--')
            
            ax_ecg.set_ylim(-150, 150)

            # Only label rightmost plots
            if (idx + 1) % n_cols == 0:
                 ax_ecg.set_ylabel("ECG (μV)", color='green')
            else:
                 ax_ecg.set_yticklabels([])
                 
            ax_ecg.tick_params(axis='y', labelcolor='green')

        ax.set_title(f"{patient_id}")
        
        if idx % n_cols == 0:
            ax.set_ylabel("EEG (μV)")
        if idx >= n_pats - n_cols:
            ax.set_xlabel("Time (s)")

    # Hide unused subplots
    for j in range(len(individuals), len(axes)):
        axes[j].axis('off')
        
    fig.suptitle(f"Patients Comparison (Butterfly EEG + ECG) - {selected_group} {selected_stage}", fontsize=16)
    st.pyplot(fig, use_container_width=True)

def _granger_causality_index(cause, effect, lag=8):
    """
    Compute a normalized Granger causality index (cause → effect).
    Compares RSS of restricted AR model vs full ARX model.
    Returns a value in [0, 1]; higher = stronger directional coupling.
    """
    n = len(effect)
    if n <= lag * 3:
        return 0.0
    Y = effect[lag:]
    X_r = np.column_stack([effect[lag - k - 1: n - k - 1] for k in range(lag)])
    X_f = np.column_stack(
        [effect[lag - k - 1: n - k - 1] for k in range(lag)] +
        [cause[lag - k - 1: n - k - 1] for k in range(lag)]
    )
    try:
        beta_r, _, _, _ = np.linalg.lstsq(X_r, Y, rcond=None)
        rss_r = np.sum((Y - X_r @ beta_r) ** 2)
        beta_f, _, _, _ = np.linalg.lstsq(X_f, Y, rcond=None)
        rss_f = np.sum((Y - X_f @ beta_f) ** 2)
        if rss_r < 1e-12 or rss_f >= rss_r:
            return 0.0
        gc = np.log(rss_r / rss_f)
        return float(np.clip(gc / (gc + 1.0), 0.0, 1.0))
    except Exception:
        return 0.0


def _compute_eeg_band_power_series(eeg_data, sfreq, fs_out=4.0, bands=None, window_sec=8.0):
    """
    Compute EEG band power time series using a sliding Hann-windowed FFT.

    Returns:
        band_powers : dict {band_name: ndarray (n_channels, n_steps)}
        t_power     : ndarray of time points in seconds (relative to start of eeg_data)
    """
    if bands is None:
        bands = {'delta': (0.5, 4), 'theta': (4, 8), 'alpha': (8, 12), 'beta': (12, 30)}
    n_ch, n_samples = eeg_data.shape
    window_samples = int(window_sec * sfreq)
    step_samples = max(1, int(sfreq / fs_out))
    if window_samples > n_samples:
        return None, None
    n_steps = (n_samples - window_samples) // step_samples + 1
    freqs = np.fft.rfftfreq(window_samples, d=1.0 / sfreq)
    hann = np.hanning(window_samples)
    band_powers = {b: np.zeros((n_ch, n_steps)) for b in bands}
    for step_i in range(n_steps):
        start = step_i * step_samples
        segment = eeg_data[:, start: start + window_samples]
        fft_power = np.abs(np.fft.rfft(segment * hann, axis=1)) ** 2
        for band, (fmin, fmax) in bands.items():
            mask = (freqs >= fmin) & (freqs < fmax)
            if mask.any():
                band_powers[band][:, step_i] = np.mean(fft_power[:, mask], axis=1)
    t_power = np.arange(n_steps) * step_samples / sfreq
    return band_powers, t_power


def _compute_bhi_patient(raw, rpeaks, sfreq, bands=None, fs_hrv=4.0, lag=8):
    """
    Compute BHI (Brain-Heart Interaction) indexes for one patient using
    a Granger-causality approach analogous to Catrambone et al. (2019).

    Returns dict {'HtB': {band: value}, 'BtH': {band: value}} or None.
    """
    if bands is None:
        bands = {'delta': (0.5, 4), 'theta': (4, 8), 'alpha': (8, 12), 'beta': (12, 30)}

    # --- HRV series ---
    rr_intervals = np.diff(rpeaks) / sfreq
    valid_mask = (rr_intervals > 0.3) & (rr_intervals < 2.0)
    rpeaks_trimmed = rpeaks[:-1][valid_mask]
    rr_valid = rr_intervals[valid_mask]
    if len(rr_valid) < 20:
        return None
    t_rpeaks = rpeaks_trimmed / sfreq
    t_hrv = np.arange(t_rpeaks[0], t_rpeaks[-1], 1.0 / fs_hrv)
    hrv_series = np.interp(t_hrv, t_rpeaks, rr_valid)

    # --- EEG channels ---
    ch_names = raw.ch_names
    eeg_idx = [i for i, ch in enumerate(ch_names)
               if re.match(r'^[A-Za-z]{1,3}[0-9]+$', ch) or re.match(r'^[A-Za-z]{1,2}z$', ch, re.IGNORECASE)]
    if not eeg_idx:
        return None
    eeg_data = raw.get_data(picks=eeg_idx)  # (n_channels, n_samples)

    # --- Band power time series ---
    band_powers, t_power = _compute_eeg_band_power_series(eeg_data, sfreq, fs_out=fs_hrv, bands=bands)
    if band_powers is None:
        return None

    # EEG power time axis is relative to start of recording
    t_eeg_abs = rpeaks[0] / sfreq + t_power  # anchor to first R-peak region

    # --- Align common time window ---
    t_start = max(t_hrv[0], t_eeg_abs[0])
    t_end = min(t_hrv[-1], t_eeg_abs[-1])
    if t_end <= t_start:
        return None
    t_common = np.arange(t_start, t_end, 1.0 / fs_hrv)
    if len(t_common) < lag * 3:
        return None

    hrv_aligned = np.interp(t_common, t_hrv, hrv_series)
    hrv_z = (hrv_aligned - np.mean(hrv_aligned)) / (np.std(hrv_aligned) + 1e-10)

    # --- Granger causality per channel per band ---
    results = {'HtB': {}, 'BtH': {}}
    for band in bands:
        htb_vals, bth_vals = [], []
        for ch_i in range(len(eeg_idx)):
            eeg_aligned = np.interp(t_common, t_eeg_abs, band_powers[band][ch_i])
            eeg_z = (eeg_aligned - np.mean(eeg_aligned)) / (np.std(eeg_aligned) + 1e-10)
            htb_vals.append(_granger_causality_index(hrv_z, eeg_z, lag=lag))
            bth_vals.append(_granger_causality_index(eeg_z, hrv_z, lag=lag))
        results['HtB'][band] = float(np.mean(htb_vals))
        results['BtH'][band] = float(np.mean(bth_vals))
    return results


def plot_bhi_analysis(filtered_individuals, selected_group, selected_stage, base_path):
    """
    Loads raw EEG data per patient, computes BHI indexes, and displays group results.

    Brain-Heart Interplay is quantified via a Granger-causality ARX approach
    inspired by Catrambone et al. (2019), Annals of Biomedical Engineering.
    """
    st.subheader(f"Brain-Heart Interplay (BHI) — {selected_group} / {selected_stage}")
    with st.expander("ℹ️ About BHI Analysis", expanded=False):
        st.markdown(
            "Directional coupling between heartbeat (HRV) and EEG band power, using a "
            "Granger-causality index analogous to the ARX model of "
            "[Catrambone et al. (2019)](https://doi.org/10.1007/s10439-019-02249-y).\n\n"
            "- **HtB** (Heart → Brain): HRV fluctuations predicting EEG band power\n"
            "- **BtH** (Brain → Heart): EEG band power predicting HRV"
        )

    bands = {'delta': (0.5, 4), 'theta': (4, 8), 'alpha': (8, 12), 'beta': (12, 30)}
    lag_order = st.slider(
        "ARX lag order (samples at 4 Hz)", min_value=2, max_value=16, value=8,
        key="bhi_lag", help="Number of past samples used in the ARX model. 8 = 2 s memory at 4 Hz."
    )

    group_dir = os.path.join(base_path, selected_group, selected_stage)
    patient_results = {}
    progress = st.progress(0, text="Computing BHI...")
    n = len(filtered_individuals)

    for i, ind in enumerate(filtered_individuals):
        pid = ind[0]
        progress.progress((i + 1) / n, text=f"BHI: {pid} ({i + 1}/{n})")
        potential_files = [f for f in os.listdir(group_dir) if f.startswith(pid) and f.endswith('.pkl')]
        if not potential_files:
            continue
        file_path = os.path.join(group_dir, potential_files[0])
        try:
            with open(file_path, 'rb') as fh:
                raw_obj = pickle.load(fh)
            raw_obj = drop_non_eeg_channels(raw_obj)
            result = process_file_data(raw_obj, pid)
            if result is None:
                continue
            raw_obj, sfreq, _rpeak_ts, rpeaks, _minmax, _ = result
            bhi = _compute_bhi_patient(raw_obj, rpeaks, sfreq, bands=bands, lag=lag_order)
            if bhi is not None:
                patient_results[pid] = bhi
        except Exception as e:
            st.warning(f"BHI: skipped {pid} — {e}")

    progress.empty()

    if not patient_results:
        st.error("Could not compute BHI for any patient.")
        return

    band_names = list(bands.keys())
    htb_group = {b: [r['HtB'][b] for r in patient_results.values()] for b in band_names}
    bth_group = {b: [r['BtH'][b] for r in patient_results.values()] for b in band_names}

    colors = ['#4C72B0', '#DD8452', '#55A868', '#C44E52']
    fig, axes = plt.subplots(1, 2, figsize=(14, 6))
    x = np.arange(len(band_names))
    rng = np.random.default_rng(42)

    for ax, data, direction in [
        (axes[0], htb_group, 'Heart \u2192 Brain (HtB)'),
        (axes[1], bth_group, 'Brain \u2192 Heart (BtH)'),
    ]:
        means = [np.mean(data[b]) for b in band_names]
        sems = [np.std(data[b]) / np.sqrt(max(len(data[b]), 1)) for b in band_names]
        ax.bar(x, means, 0.55, yerr=sems, capsize=5, color=colors, alpha=0.8)
        for j, band in enumerate(band_names):
            ys = data[band]
            xs = np.full(len(ys), j) + rng.uniform(-0.12, 0.12, len(ys))
            ax.scatter(xs, ys, color='k', s=20, alpha=0.6, zorder=5)
        ax.set_xticks(x)
        ax.set_xticklabels(band_names, fontsize=11)
        ax.set_xlabel("EEG Frequency Band")
        ax.set_ylabel("Granger Causality Index")
        ax.set_title(f"{direction} - Mean +/- SEM", fontsize=13)
        ax.grid(True, alpha=0.3, axis='y')
        ax.set_ylim(bottom=0)

    fig.suptitle(
        f"BHI - Mean +/- SEM — {selected_group} / {selected_stage}  (n={len(patient_results)} patients)",
        fontsize=14
    )
    plt.tight_layout()
    st.pyplot(fig, use_container_width=True)

    # Summary table
    rows = [
        {
            'Band': b,
            'HtB mean': f"{np.mean(htb_group[b]):.4f}",
            'HtB SEM': f"{np.std(htb_group[b]) / np.sqrt(len(htb_group[b])):.4f}",
            'BtH mean': f"{np.mean(bth_group[b]):.4f}",
            'BtH SEM': f"{np.std(bth_group[b]) / np.sqrt(len(bth_group[b])):.4f}",
        }
        for b in band_names
    ]
    st.dataframe(pd.DataFrame(rows), use_container_width=True)


def plot_group_ecg_analysis(individuals, selected_group, selected_stage):
    """
    Plots the average ECG of the whole group and individual patient ECGs.
    """
    st.subheader(f"Group ECG Analysis - {selected_group} {selected_stage}")

    # 1. Group Average ECG
    all_ecg_heps = []
    for ind in individuals:
        ecg_hep = ind[5]
        if ecg_hep is not None:
            all_ecg_heps.append(ecg_hep[0])
    
    if all_ecg_heps:
        avg_ecg = np.nanmean(all_ecg_heps, axis=0)
        times = individuals[0][2]
        
        fig_avg, ax_avg = plt.subplots(figsize=(10, 5))
        ax_avg.plot(times, avg_ecg * 1e6, color='black', linewidth=2, label='Group Average ECG')
        
        # Optional: Add standard error shading
        std_ecg = np.nanstd(all_ecg_heps, axis=0)
        sem_ecg = std_ecg / np.sqrt(len(all_ecg_heps))
        ax_avg.fill_between(times, (avg_ecg - sem_ecg) * 1e6, (avg_ecg + sem_ecg) * 1e6, color='gray', alpha=0.3, label='SEM')

        ax_avg.set_title(f"Group Average ECG - Mean +/- SEM (n={len(all_ecg_heps)})")
        ax_avg.set_xlabel("Time (s)")
        ax_avg.set_ylabel("Amplitude (μV)")
        ax_avg.axvline(0, color='r', linestyle='--', alpha=0.5)
        ax_avg.legend()
        ax_avg.grid(True, alpha=0.3)
        st.pyplot(fig_avg, use_container_width=True)
    else:
        st.warning("No ECG data found for this group to calculate average.")

    # 2. RR Interval Analysis
    st.markdown("#### RR Interval Analysis")
    
    # Calculate RR intervals for all individuals
    all_rr_intervals = []
    patient_rr_intervals = {} # {pid: rr_intervals}
    
    if individuals:
        # Estimate sampling rate from time vector of first individual
        # times = individuals[0][2]
        # fs = 1 / np.mean(np.diff(times))
        # Better: use the median diff to be robust
        times = individuals[0][2]
        dt = np.median(np.diff(times))
        fs = 1.0 / dt
        
        for ind in individuals:
            pid = ind[0]
            rpeaks = ind[4]
            
            if rpeaks is not None and len(rpeaks) > 1:
                # Calculate RR intervals in seconds
                rr_samples = np.diff(rpeaks)
                rr_sec = rr_samples / fs
                
                # Filter physiological range? (e.g., 0.3s to 2.0s corresponds to 30-200 BPM)
                # For now, let's keep it raw but maybe exclude obvious artifacts if needed
                # valid_rr = rr_sec[(rr_sec > 0.3) & (rr_sec < 2.0)]
                valid_rr = rr_sec 
                
                all_rr_intervals.extend(valid_rr)
                patient_rr_intervals[pid] = valid_rr

    # 2a. Group RR Histogram
    if all_rr_intervals:
        fig_rr_group, ax_rr_group = plt.subplots(figsize=(10, 5))
        ax_rr_group.hist(all_rr_intervals, bins=50, color='purple', alpha=0.7, edgecolor='black')
        
        avg_rr = np.mean(all_rr_intervals)
        std_rr = np.std(all_rr_intervals)
        
        ax_rr_group.set_title(f"Group RR Interval Distribution (n={len(all_rr_intervals)} beats)\nMean={avg_rr:.3f}s ({60/avg_rr:.1f} BPM), SD={std_rr:.3f}s")
        ax_rr_group.set_xlabel("RR Interval (s)")
        ax_rr_group.set_ylabel("Count")
        ax_rr_group.set_xlim(right=1.3)
        ax_rr_group.grid(True, alpha=0.3)
        
        st.pyplot(fig_rr_group, use_container_width=True)
    else:
        st.info("No RR intervals could be calculated.")

    # 2b. Individual RR Histograms
    if patient_rr_intervals:
        st.markdown("#### Individual RR Interval Histograms")
        
        # Use existing slider or adding a new one? Use existing slider for consistency if possible?
        # The existing slider `n_individuals` is already defined below for ECG plots. 
        # But we haven't reached that part of the code yet (it's in original lines 1944+).
        # We are inserting BEFORE line 1944. 
        # So we can define a slider here or reuse the one for ECG plots if we move code.
        # Let's add a separate section for this to be clean.
        
        # Grid layout
        pids = list(patient_rr_intervals.keys())
        n_pats_rr = len(pids)
        
        # Slider to limit number of plots if too many
        n_show_rr = st.slider("Number of patients to show for RR Histograms", 1, n_pats_rr, min(4, n_pats_rr), key="rr_hist_slider")
        
        n_cols_rr = 2
        n_rows_rr = int(np.ceil(n_show_rr / n_cols_rr))
        
        fig_rr_ind, axes_rr_ind = plt.subplots(n_rows_rr, n_cols_rr, figsize=(12, 4 * n_rows_rr))
        if n_show_rr == 1:
            axes_rr_ind = [axes_rr_ind]
        else:
            axes_rr_ind = axes_rr_ind.flatten()
            
        for idx in range(n_show_rr):
            pid = pids[idx]
            rr_data = patient_rr_intervals[pid]
            ax = axes_rr_ind[idx]
            
            ax.hist(rr_data, bins=30, color='teal', alpha=0.6, edgecolor='black')
            
            mean_rr = np.mean(rr_data)
            std_rr = np.std(rr_data)
            
            ax.set_title(f"{pid}\nMean={mean_rr:.3f}s, SD={std_rr:.3f}s")
            ax.set_xlabel("RR (s)")
            ax.set_xlim(right=1.3)
            ax.grid(True, alpha=0.3)
            
        # Hide unused
        for j in range(n_show_rr, len(axes_rr_ind)):
            axes_rr_ind[j].axis('off')
            
        fig_rr_ind.tight_layout()
        st.pyplot(fig_rr_ind, use_container_width=True)


    # 3. Individual Patients
    st.markdown("#### Individual Patient ECGs")
    n_individuals = st.slider("Number of individual patients to show", min_value=1, max_value=len(individuals), value=len(individuals), key="group_ecg_slider")
    
    # Grid layout for individuals
    n_cols = 2
    n_rows = int(np.ceil(n_individuals / n_cols))
    
    fig_ind, axes = plt.subplots(n_rows, n_cols, figsize=(14, 4 * n_rows), sharex=True)
    if n_individuals == 1:
        axes = [axes]
    else:
        axes = axes.flatten()
        
    for i in range(n_individuals):
        ind = individuals[i]
        ax = axes[i]
        patient_id = ind[0]
        ecg_hep = ind[5]
        times = ind[2]
        
        if ecg_hep is not None:
             ax.plot(times, ecg_hep[0] * 1e6, color='green', linewidth=1.5)
        else:
             ax.text(0.5, 0.5, "No ECG Data", ha='center', va='center')
             
        ax.set_title(f"{patient_id}")
        ax.grid(True, alpha=0.3)
        ax.axvline(0, color='r', linestyle='--', alpha=0.5)
        
        if i % n_cols == 0:
            ax.set_ylabel("ECG (μV)")
        if i >= n_individuals - n_cols:
            ax.set_xlabel("Time (s)")
            
    # Hide unused
    if hasattr(axes, '__len__'):    
        for j in range(n_individuals, len(axes)):
            axes[j].axis('off')
            
    st.pyplot(fig_ind, use_container_width=True)


def handle_ecg_reduction(filtered_individuals, selected_group, selected_stage):
    """
    For each patient: z-score normalises the ECG HEP and each EEG channel HEP independently
    (bringing both to the same amplitude scale), then subtracts the normalised ECG from the
    normalised EEG. Plots per-patient and group-average original, ECG overlay, and residual.
    """
    st.subheader("ECG Signal Reduction (Normalised Subtraction)")
    st.caption(
        "Both the ECG HEP and each EEG channel HEP are z-score normalised independently "
        "so they share the same amplitude scale. "
        "Residual = z-scored EEG − z-scored ECG."
    )

    # Determine common EEG channels
    common_channels = identify_common_eeg_channels(filtered_individuals)
    if not common_channels:
        st.warning("No common EEG channels found.")
        return

    display_channels = ['Average'] + common_channels

    ch_data = {ch: {'orig_norm': [], 'ecg_norm': [], 'residual': [], 'pids': []} for ch in display_channels}
    times = None

    for ind in filtered_individuals:
        pid, hep_full, t, ch_names, rpeaks, ecg_hep, ecg_ch = ind[:7]
        if ecg_hep is None:
            continue
        ecg_signal = np.asarray(ecg_hep).squeeze()
        if ecg_signal.ndim != 1:
            continue
        if times is None:
            times = t
        if len(ecg_signal) != len(t):
            continue

        ecg_std = ecg_signal.std()
        if ecg_std < 1e-20:
            continue
        ecg_z = (ecg_signal - ecg_signal.mean()) / ecg_std

        valid_indices = []
        for ch in common_channels:
            if ch not in ch_names:
                continue
            idx = ch_names.index(ch)
            eeg_trace = hep_full[idx].copy()
            eeg_std = eeg_trace.std()
            if eeg_std < 1e-20:
                continue
            eeg_z = (eeg_trace - eeg_trace.mean()) / eeg_std
            residual = eeg_z - ecg_z
            ch_data[ch]['orig_norm'].append(eeg_z)
            ch_data[ch]['ecg_norm'].append(ecg_z)
            ch_data[ch]['residual'].append(residual)
            ch_data[ch]['pids'].append(pid)
            valid_indices.append(idx)

        if valid_indices:
            avg_trace = summarize_channels_without_reference_cancellation(
                hep_full[valid_indices]
            )
            avg_std = avg_trace.std()
            if avg_std < 1e-20:
                continue
            avg_z = (avg_trace - avg_trace.mean()) / avg_std
            residual_avg = avg_z - ecg_z
            ch_data['Average']['orig_norm'].append(avg_z)
            ch_data['Average']['ecg_norm'].append(ecg_z)
            ch_data['Average']['residual'].append(residual_avg)
            ch_data['Average']['pids'].append(pid)

    if times is None:
        st.warning("No valid ECG HEP data found.")
        return

    for ch_name in display_channels:
        d = ch_data[ch_name]
        if not d['orig_norm']:
            continue
        orig_arr = np.array(d['orig_norm'])
        ecg_arr = np.array(d['ecg_norm'])
        resid_arr = np.array(d['residual'])
        pids = d['pids']
        n_subj = len(pids)

        with st.expander(f"ECG Reduction — Channel: {ch_name}  (n={n_subj})", expanded=(ch_name == 'Average')):
            fig, axes = plt.subplots(1, 2, figsize=(14, 5))
            cmap = plt.get_cmap('tab20' if n_subj <= 20 else 'hsv', max(n_subj, 1))
            colors = [cmap(i / max(n_subj - 1, 1)) for i in range(n_subj)]

            # --- subplot 1: normalised EEG + ECG overlay per patient
            ax = axes[0]
            for i, (eeg_z, ecg_z, pid) in enumerate(zip(orig_arr, ecg_arr, pids)):
                ax.plot(times, eeg_z, color=colors[i], alpha=0.35, linewidth=0.8)
            if n_subj:
                ax.plot(times, orig_arr.mean(axis=0), color='steelblue', linewidth=2, label='EEG avg (z)')
                ax.plot(times, ecg_arr.mean(axis=0), color='crimson', linewidth=2,
                        linestyle='--', label='ECG avg (z)')
            ax.set_title("Normalised EEG + ECG (z-score)", fontsize=10)
            ax.set_xlabel("Time (s)")
            ax.set_ylabel("Amplitude (z-score)")
            ax.axvline(0, color='gray', linewidth=0.8, linestyle='--')
            ax.axhline(0, color='gray', linewidth=0.5)
            ax.legend(fontsize=8)
            ax.grid(alpha=0.2)

            # --- subplot 2: residual per patient
            ax = axes[1]
            for i, (res, pid) in enumerate(zip(resid_arr, pids)):
                ax.plot(times, res, color=colors[i], alpha=0.35, linewidth=0.8)
            if n_subj:
                ax.plot(times, resid_arr.mean(axis=0), color='black', linewidth=2, label='Residual avg')
            ax.set_title("Residual HEP (EEG_z − ECG_z)", fontsize=10)
            ax.set_xlabel("Time (s)")
            ax.set_ylabel("Amplitude (z-score)")
            ax.axvline(0, color='gray', linewidth=0.8, linestyle='--')
            ax.axhline(0, color='gray', linewidth=0.5)
            ax.legend(fontsize=8)
            ax.grid(alpha=0.2)

            fig.suptitle(
                f"ECG Reduction — Channel: {ch_name} | {selected_group} / {selected_stage}",
                fontsize=11, fontweight='bold'
            )
            fig.tight_layout()
            st.pyplot(fig, use_container_width=True)
            plt.close(fig)


def _shift_trace_with_nan(trace, shift_samples):
    """Shift a 1D trace by integer samples while padding exposed edges with NaN."""
    trace = np.asarray(trace, dtype=float)
    shifted = np.full_like(trace, np.nan, dtype=float)
    if len(trace) == 0:
        return shifted
    if shift_samples == 0:
        shifted[:] = trace
    elif shift_samples > 0:
        shifted[shift_samples:] = trace[:-shift_samples]
    else:
        shifted[:shift_samples] = trace[-shift_samples:]
    return shifted


def _safe_nanargmax(signal_1d):
    """Return argmax for finite values only, or None if the trace is all non-finite."""
    arr = np.asarray(signal_1d, dtype=float)
    finite_mask = np.isfinite(arr)
    if not finite_mask.any():
        return None
    finite_idx = np.flatnonzero(finite_mask)
    return int(finite_idx[np.argmax(arr[finite_mask])])


def _safe_nanargabsmax(signal_1d):
    """Return index of the largest absolute finite value, or None if all values are non-finite."""
    arr = np.asarray(signal_1d, dtype=float)
    finite_mask = np.isfinite(arr)
    if not finite_mask.any():
        return None
    finite_idx = np.flatnonzero(finite_mask)
    return int(finite_idx[np.argmax(np.abs(arr[finite_mask]))])


def _compute_ecg_guided_attenuation(eeg_trace, ecg_trace_aligned):
    """
    Reduce EEG amplitude with a time-varying gain guided by the aligned ECG.
    Stronger ECG magnitude yields stronger attenuation, while the per-channel
    attenuation strength is controlled by the EEG/ECG correlation.
    """
    eeg_trace = np.asarray(eeg_trace, dtype=float)
    ecg_trace_aligned = np.asarray(ecg_trace_aligned, dtype=float)
    valid = np.isfinite(eeg_trace) & np.isfinite(ecg_trace_aligned)
    if not valid.any():
        nan_arr = np.full_like(eeg_trace, np.nan, dtype=float)
        return np.nan, nan_arr, eeg_trace.copy(), nan_arr

    eeg_valid = eeg_trace[valid]
    ecg_valid = ecg_trace_aligned[valid]

    ecg_centered = ecg_valid - np.mean(ecg_valid)
    ecg_abs = np.abs(ecg_centered)
    ecg_abs_max = float(np.max(ecg_abs))
    if ecg_abs_max < 1e-20:
        nan_arr = np.full_like(eeg_trace, np.nan, dtype=float)
        return 0.0, nan_arr, eeg_trace.copy(), nan_arr

    ecg_weight = ecg_abs / ecg_abs_max
    ecg_weight = np.power(ecg_weight, 0.75)
    eeg_std = float(np.std(eeg_valid))
    ecg_std = float(np.std(ecg_valid))
    if eeg_std < 1e-20 or ecg_std < 1e-20:
        corr_strength = 0.0
    else:
        corr_strength = float(np.corrcoef(eeg_valid, ecg_valid)[0, 1])
        if not np.isfinite(corr_strength):
            corr_strength = 0.0

    corr_abs = np.abs(corr_strength)
    attenuation_strength = float(np.clip(0.15 + 1.35 * corr_abs, 0.0, 0.98))
    gain_valid = 1.0 - attenuation_strength * ecg_weight
    gain_valid = np.power(np.clip(gain_valid, 0.02, 1.0), 1.35)
    gain_valid = np.clip(gain_valid, 0.02, 1.0)

    gain = np.full_like(eeg_trace, np.nan, dtype=float)
    gain[valid] = gain_valid

    cleaned = eeg_trace.copy()
    cleaned[valid] = eeg_valid * gain_valid

    removed = np.full_like(eeg_trace, np.nan, dtype=float)
    removed[valid] = eeg_valid - cleaned[valid]
    return attenuation_strength, gain, cleaned, removed


def handle_nc_hep_cleaning(filtered_individuals, selected_group, selected_stage):
    """
    Peak-aligned ECG artifact reduction for HEP windows.
    For each patient:
      1. Average across EEG channels to get a representative EEG HEP.
      2. Find the EEG global maximum and ECG global maximum inside the HEP window.
      3. Shift the ECG HEP so both peaks align.
      4. Estimate how strongly each EEG channel follows the aligned ECG.
      5. Reduce EEG amplitude with a time-varying gain where the ECG is strongest.
    """
    st.subheader("NC_HEP_CLEANING")
    st.caption(
        "The ECG HEP template is time-aligned to the patient EEG HEP peak, then each EEG channel is attenuated "
        "according to how strongly it correlates with the aligned ECG. Where the ECG is high, the EEG is reduced more."
    )

    if not filtered_individuals:
        st.warning("No patient data available.")
        return

    common_channels = identify_common_eeg_channels(filtered_individuals)
    if not common_channels:
        st.warning("No common EEG channels found.")
        return

    patient_options = [ind[0] for ind in filtered_individuals if len(ind) >= 7 and ind[5] is not None]
    if not patient_options:
        st.warning("No patients with valid ECG HEP data were found.")
        return

    selected_pid = st.selectbox(
        "Select Patient for NC_HEP_CLEANING",
        patient_options,
        key="nc_hep_cleaning_patient",
    )
    ind = next((x for x in filtered_individuals if x[0] == selected_pid), None)
    if ind is None:
        st.warning("Selected patient could not be loaded.")
        return

    pid, hep_full, times, ch_names, _, ecg_hep, _ = ind[:7]
    ecg_trace = np.asarray(ecg_hep).squeeze() if ecg_hep is not None else None
    if ecg_trace is None or ecg_trace.ndim != 1:
        st.warning(f"No valid ECG HEP trace found for {pid}.")
        return

    available_channels = [ch for ch in common_channels if ch in ch_names]
    if not available_channels:
        st.warning(f"No common EEG channels available for {pid}.")
        return

    eeg_idx = [ch_names.index(ch) for ch in available_channels]
    eeg_mat = np.asarray(hep_full[eeg_idx], dtype=float)
    if eeg_mat.ndim != 2 or eeg_mat.shape[1] != len(times):
        st.warning("EEG HEP data shape does not match the time vector.")
        return

    eeg_global = np.nanmean(eeg_mat, axis=0)
    eeg_peak_idx = _safe_nanargabsmax(eeg_global)
    ecg_peak_idx = _safe_nanargabsmax(ecg_trace)
    if eeg_peak_idx is None or ecg_peak_idx is None:
        st.warning("Could not detect valid global maxima for EEG and ECG.")
        return

    channel_strengths = []
    channel_shift_samples = []
    cleaned_mat = np.full_like(eeg_mat, np.nan, dtype=float)
    gain_mat = np.full_like(eeg_mat, np.nan, dtype=float)
    removed_mat = np.full_like(eeg_mat, np.nan, dtype=float)
    aligned_ecg_mat = np.full_like(eeg_mat, np.nan, dtype=float)
    for i_ch, eeg_trace in enumerate(eeg_mat):
        eeg_ch_peak_idx = _safe_nanargabsmax(eeg_trace)
        if eeg_ch_peak_idx is None:
            aligned_ecg = np.full_like(ecg_trace, np.nan, dtype=float)
            shift_samples = np.nan
        else:
            shift_samples = int(eeg_ch_peak_idx - ecg_peak_idx)
            aligned_ecg = _shift_trace_with_nan(ecg_trace, shift_samples)
        strength, gain, cleaned, removed = _compute_ecg_guided_attenuation(eeg_trace, aligned_ecg)
        channel_strengths.append(strength)
        channel_shift_samples.append(shift_samples)
        aligned_ecg_mat[i_ch] = aligned_ecg
        gain_mat[i_ch] = gain
        cleaned_mat[i_ch] = cleaned
        removed_mat[i_ch] = removed

    channel_strengths = np.asarray(channel_strengths, dtype=float)
    channel_shift_samples = np.asarray(channel_shift_samples, dtype=float)
    avg_aligned_ecg = np.nanmean(aligned_ecg_mat, axis=0)
    avg_gain = np.nanmean(gain_mat, axis=0)
    avg_cleaned = np.nanmean(cleaned_mat, axis=0)
    avg_removed = np.nanmean(removed_mat, axis=0)
    sfreq = 1.0 / np.mean(np.diff(times)) if len(times) > 1 else np.nan
    peak_shift_ms = np.nan
    finite_shift = channel_shift_samples[np.isfinite(channel_shift_samples)]
    if finite_shift.size and np.isfinite(sfreq) and sfreq > 0:
        peak_shift_ms = float(np.mean(finite_shift) / sfreq * 1000.0)

    with st.expander("Alignment Summary", expanded=True):
        st.write(f"Patient: **{pid}**")
        st.write(f"EEG channels used: **{len(available_channels)}**")
        st.write(f"ECG dominant peak: **{times[ecg_peak_idx] * 1000:.1f} ms**")
        st.write(f"Global EEG dominant peak: **{times[eeg_peak_idx] * 1000:.1f} ms**")
        if finite_shift.size:
            st.write(
                f"Electrode-specific ECG shift: **{np.mean(finite_shift):.1f} ± {np.std(finite_shift):.1f} samples** "
                f"({peak_shift_ms:.1f} ms mean)"
            )
        finite_strengths = channel_strengths[np.isfinite(channel_strengths)]
        if finite_strengths.size:
            st.write(
                f"Per-channel attenuation strength: mean **{np.mean(finite_strengths):.3f}**, "
                f"median **{np.median(finite_strengths):.3f}**, std **{np.std(finite_strengths):.3f}**"
            )

    fig_main, axes_main = plt.subplots(2, 2, figsize=(16, 10))

    ax = axes_main[0, 0]
    ax.plot(times * 1000, eeg_global * 1e6, color="steelblue", linewidth=2.5, label="Global EEG HEP")
    ax.plot(times * 1000, ecg_trace * 1e6, color="crimson", linewidth=2, alpha=0.85, label="Original ECG HEP")
    ax.axvline(times[eeg_peak_idx] * 1000, color="steelblue", linestyle="--", alpha=0.7)
    ax.axvline(times[ecg_peak_idx] * 1000, color="crimson", linestyle="--", alpha=0.7)
    ax.set_title("Before Alignment")
    ax.set_xlabel("Time (ms)")
    ax.set_ylabel("Amplitude (uV)")
    ax.grid(alpha=0.25)
    ax.legend(fontsize=8)

    ax = axes_main[0, 1]
    ax.plot(times * 1000, eeg_global * 1e6, color="steelblue", linewidth=2.5, label="Global EEG HEP")
    ax.plot(times * 1000, avg_aligned_ecg * 1e6, color="darkorange", linewidth=2,
            alpha=0.9, label="Mean electrode-aligned ECG HEP")
    ax.axvline(times[eeg_peak_idx] * 1000, color="black", linestyle="--", alpha=0.7, label="Global EEG peak")
    ax.set_title("After Electrode-Specific Peak Alignment")
    ax.set_xlabel("Time (ms)")
    ax.set_ylabel("Amplitude (uV)")
    ax.grid(alpha=0.25)
    ax.legend(fontsize=8)

    ax = axes_main[1, 0]
    ax.plot(times * 1000, eeg_global * 1e6, color="steelblue", linewidth=2.5, label="Global EEG HEP")
    ax.plot(times * 1000, avg_removed * 1e6, color="purple", linewidth=2,
            linestyle="--", label="Mean attenuation amount")
    ax.plot(times * 1000, avg_cleaned * 1e6, color="black", linewidth=2.5, label="Cleaned global EEG HEP")
    ax.set_title("Global EEG Before vs Reduced Part vs Cleaned")
    ax.set_xlabel("Time (ms)")
    ax.set_ylabel("Amplitude (uV)")
    ax.grid(alpha=0.25)
    ax.legend(fontsize=8)

    ax = axes_main[1, 1]
    valid_order = np.argsort(np.nan_to_num(channel_strengths, nan=np.inf))
    ordered_strengths = channel_strengths[valid_order]
    ordered_names = [available_channels[i] for i in valid_order]
    ax.barh(ordered_names, ordered_strengths, color="teal", alpha=0.8, edgecolor="black", linewidth=0.4)
    ax.axvline(0, color="black", linestyle="--", linewidth=0.8)
    ax.set_title("Per-Channel ECG Attenuation Strength")
    ax.set_xlabel("Attenuation strength")
    ax.set_ylabel("EEG channel")
    ax.grid(axis="x", alpha=0.25)

    ax_gain = ax.twinx()
    ax_gain.plot(times * 1000, avg_gain, color="darkgreen", linewidth=1.8, alpha=0.8)
    ax_gain.set_ylabel("Mean gain", color="darkgreen")
    ax_gain.tick_params(axis='y', colors="darkgreen")
    ax_gain.set_ylim(0, 1.05)

    fig_main.suptitle(
        f"NC_HEP_CLEANING — {pid} | {selected_group} / {selected_stage}",
        fontsize=12,
        fontweight="bold",
    )
    fig_main.tight_layout()
    st.pyplot(fig_main, use_container_width=True)
    plt.close(fig_main)

    n_show = min(6, len(available_channels))
    default_channels = available_channels[:n_show]
    selected_channels = st.multiselect(
        "Channels to inspect after NC_HEP_CLEANING",
        options=available_channels,
        default=default_channels,
        key="nc_hep_cleaning_channels",
    )

    if selected_channels:
        n_rows = len(selected_channels)
        fig_ch, axes_ch = plt.subplots(n_rows, 1, figsize=(14, 3.2 * n_rows), sharex=True)
        if n_rows == 1:
            axes_ch = [axes_ch]

        for ax, ch in zip(axes_ch, selected_channels):
            ch_idx = available_channels.index(ch)
            eeg_orig = eeg_mat[ch_idx]
            aligned_ecg = aligned_ecg_mat[ch_idx]
            eeg_removed = removed_mat[ch_idx]
            eeg_clean = cleaned_mat[ch_idx]
            gain_trace = gain_mat[ch_idx]
            valid_overlay = np.isfinite(eeg_orig) & np.isfinite(aligned_ecg)
            if valid_overlay.any():
                ecg_overlay = np.full_like(aligned_ecg, np.nan, dtype=float)
                ecg_overlay_valid = aligned_ecg[valid_overlay]
                eeg_overlay_valid = eeg_orig[valid_overlay]
                ecg_overlay_centered = ecg_overlay_valid - np.mean(ecg_overlay_valid)
                eeg_overlay_centered = eeg_overlay_valid - np.mean(eeg_overlay_valid)
                denom_overlay = float(np.dot(ecg_overlay_centered, ecg_overlay_centered))
                if denom_overlay > 1e-20:
                    overlay_scale = float(np.dot(eeg_overlay_centered, ecg_overlay_centered) / denom_overlay)
                    ecg_overlay[valid_overlay] = np.mean(eeg_overlay_valid) + overlay_scale * ecg_overlay_centered
                else:
                    ecg_overlay[:] = np.nan
            else:
                ecg_overlay = np.full_like(aligned_ecg, np.nan, dtype=float)
            ax.plot(times * 1000, eeg_orig * 1e6, color="steelblue", linewidth=1.8, label=f"{ch} original")
            ax.plot(times * 1000, ecg_overlay * 1e6, color="crimson", linewidth=1.4,
                    linestyle=":", label="Aligned ECG (EEG-scaled)")
            ax.plot(times * 1000, eeg_removed * 1e6, color="darkorange", linewidth=1.4,
                    linestyle="--", label=f"{ch} reduced part")
            ax.plot(times * 1000, eeg_clean * 1e6, color="black", linewidth=1.9, label=f"{ch} cleaned")
            ax.axvline(0, color="gray", linestyle="--", linewidth=0.8, alpha=0.6)
            ax.grid(alpha=0.25)
            ax.set_ylabel("uV")
            strength_txt = channel_strengths[ch_idx]
            shift_txt = channel_shift_samples[ch_idx]
            gain_min = np.nanmin(gain_trace) if np.isfinite(gain_trace).any() else np.nan
            if np.isfinite(strength_txt):
                shift_ms = shift_txt / sfreq * 1000.0 if np.isfinite(shift_txt) and np.isfinite(sfreq) and sfreq > 0 else np.nan
                ax.set_title(
                    f"{ch} | attenuation = {strength_txt:.3f}, min gain = {gain_min:.3f}, "
                    f"shift = {shift_txt:.0f} samples ({shift_ms:.1f} ms)"
                )
            else:
                ax.set_title(f"{ch} | attenuation = NaN")
            ax.legend(fontsize=8, loc="upper right")

        axes_ch[-1].set_xlabel("Time (ms)")
        fig_ch.suptitle(f"NC_HEP_CLEANING Channel Inspection — {pid}", fontsize=12, fontweight="bold")
        fig_ch.tight_layout()
        st.pyplot(fig_ch, use_container_width=True)
        plt.close(fig_ch)

    summary_df = pd.DataFrame({
        "channel": available_channels,
        "attenuation_strength": channel_strengths,
        "shift_samples": channel_shift_samples,
        "orig_peak_uV": np.nanmax(eeg_mat, axis=1) * 1e6,
        "cleaned_peak_uV": np.nanmax(cleaned_mat, axis=1) * 1e6,
        "reduced_peak_uV": np.nanmax(removed_mat, axis=1) * 1e6,
        "min_gain": np.nanmin(gain_mat, axis=1),
    })
    st.dataframe(summary_df.style.format({
        "attenuation_strength": "{:.4f}",
        "shift_samples": "{:.1f}",
        "orig_peak_uV": "{:.3f}",
        "cleaned_peak_uV": "{:.3f}",
        "reduced_peak_uV": "{:.3f}",
        "min_gain": "{:.4f}",
    }), use_container_width=True)


def handle_ica_ecg_cleaning(filtered_individuals, selected_group, selected_stage, base_path=None, snr_min=1.5, snr_max=6.0):
    """
    Visualises ECG-component removal from HEP data using a regression approach:
      - ECG HEP is treated as the single ICA component.
      - Per-channel eigenvalue = regression coefficient of that channel's HEP onto the ECG HEP.
      - Cleaned HEP = original HEP - eigenvalue × ECG HEP.
    Shows per-patient and group-average before/after plots, plus an eigenvalue bar chart.
    """
    st.subheader("ICA ECG Cleaning of HEP")
    st.caption(
        "ECG HEP is used as the sole ICA component. "
        "The eigenvalue (regression weight) measures how much cardiac signal contaminates each EEG channel. "
        "Cleaned = Original − eigenvalue × ECG HEP."
    )

    # ------------------------------------------------------------------ collect data
    # Determine common EEG channels
    common_channels = identify_common_eeg_channels(filtered_individuals)

    if not common_channels:
        st.warning("No common EEG channels found.")
        return

    display_channels = ['Average'] + common_channels

    # Per-channel storage: list of (orig, cleaned, eigenvalue) per patient
    ch_data = {ch: {'orig': [], 'clean': [], 'eig': [], 'pids': [], 'ecg': []} for ch in display_channels}
    times = None

    for ind in filtered_individuals:
        pid, hep_full, t, ch_names, rpeaks, ecg_hep, ecg_ch = ind[:7]
        if ecg_hep is None:
            continue
        ecg_signal = np.asarray(ecg_hep).squeeze()
        if ecg_signal.ndim != 1:
            continue
        if times is None:
            times = t
        if len(ecg_signal) != len(t):
            continue

        ecg_denom = np.dot(ecg_signal, ecg_signal)
        if ecg_denom < 1e-20:
            continue

        # Per common channel
        valid_indices = []
        for ch in common_channels:
            if ch not in ch_names:
                continue
            idx = ch_names.index(ch)
            eeg_trace = hep_full[idx].copy()
            eig = np.dot(eeg_trace, ecg_signal) / ecg_denom
            cleaned = eeg_trace - eig * ecg_signal
            ch_data[ch]['orig'].append(eeg_trace)
            ch_data[ch]['clean'].append(cleaned)
            ch_data[ch]['eig'].append(eig)
            ch_data[ch]['pids'].append(pid)
            ch_data[ch]['ecg'].append(ecg_signal)
            valid_indices.append(idx)

        # Average across common channels
        if valid_indices:
            avg_trace = summarize_channels_without_reference_cancellation(
                hep_full[valid_indices]
            )
            eig_avg = np.dot(avg_trace, ecg_signal) / ecg_denom
            cleaned_avg = avg_trace - eig_avg * ecg_signal
            ch_data['Average']['orig'].append(avg_trace)
            ch_data['Average']['clean'].append(cleaned_avg)
            ch_data['Average']['eig'].append(eig_avg)
            ch_data['Average']['pids'].append(pid)
            ch_data['Average']['ecg'].append(ecg_signal)

    if times is None:
        st.warning("No valid ECG HEP data found.")
        return

    # ------------------------------------------------------------------ eigenvalue summary
    st.markdown("#### Eigenvalue (ECG Regression Weight) per Channel")
    chan_labels, eig_means, eig_stds = [], [], []
    for ch in display_channels:
        eigs = ch_data[ch]['eig']
        if eigs:
            chan_labels.append(ch)
            eig_means.append(float(np.mean(eigs)))
            eig_stds.append(float(np.std(eigs)))

    if chan_labels:
        fig_eig, ax_eig = plt.subplots(figsize=(max(8, len(chan_labels) * 0.45), 4))
        x = np.arange(len(chan_labels))
        bars = ax_eig.bar(x, eig_means, yerr=eig_stds, capsize=3,
                          color=['steelblue' if c != 'Average' else 'darkorange' for c in chan_labels],
                          alpha=0.8, edgecolor='black', linewidth=0.5)
        ax_eig.axhline(0, color='black', linewidth=0.8, linestyle='--')
        ax_eig.set_xticks(x)
        ax_eig.set_xticklabels(chan_labels, rotation=45, ha='right', fontsize=8)
        ax_eig.set_ylabel("Eigenvalue (regression weight)", fontsize=9)
        ax_eig.set_title(
            f"ECG Component Weight per EEG Channel — {selected_group} / {selected_stage}\n"
            f"(mean ± SD across {len(filtered_individuals)} patients)",
            fontsize=10
        )
        ax_eig.grid(axis='y', alpha=0.3)
        fig_eig.tight_layout()
        st.pyplot(fig_eig, use_container_width=True)
        plt.close(fig_eig)

    # ------------------------------------------------------------------ IC identification panels
    st.markdown("#### Identification of Cardiac Independent Components (ICs)")
    st.caption(
        "IC #1 = ECG regression component (weights = regression eigenvalues per channel). "
        "IC #2 = first principal component of the ICA residuals (SVD-derived). "
        "IC #1 typically shows a strong lateralised gradient; IC #2 shows a more central/frontal distribution."
    )

    # Derive IC #2 timeseries via SVD of residuals (using Average channel patient data)
    ic2_timeseries = None
    ic2_weights = {}
    avg_d = ch_data.get('Average', {})
    if avg_d.get('orig') and avg_d.get('eig') and times is not None:
        # Build residual matrix: rows=patients, cols=time
        residual_mat = []
        ecg_list_ic2 = []
        for ind2 in filtered_individuals:
            _pid2, _hep2, _t2, _ch2, _rp2, _ecg2, _ecg_ch2 = ind2[:7]
            if _ecg2 is None:
                continue
            _ecg2_sq = np.asarray(_ecg2).squeeze()
            if _ecg2_sq.ndim != 1 or len(_ecg2_sq) != len(times):
                continue
            ecg_list_ic2.append(_ecg2_sq)
        # Use Average channel orig and eig
        for orig_tr, eig_v, ecg_tr in zip(avg_d['orig'], avg_d['eig'], ecg_list_ic2[:len(avg_d['orig'])]):
            residual_mat.append(orig_tr - eig_v * ecg_tr)
        if len(residual_mat) >= 2:
            R = np.array(residual_mat)
            R -= R.mean(axis=0, keepdims=True)
            try:
                _U, _s, _Vt = np.linalg.svd(R, full_matrices=False)
                ic2_timeseries = _Vt[0]
            except Exception:
                ic2_timeseries = None

    # Compute IC #2 weight per channel
    if ic2_timeseries is not None:
        ic2_denom = np.dot(ic2_timeseries, ic2_timeseries)
        if ic2_denom > 1e-20:
            for ch_c in common_channels:
                if ch_data[ch_c]['orig']:
                    ch_mean = np.mean(ch_data[ch_c]['orig'], axis=0)
                    ic2_weights[ch_c] = np.dot(ch_mean, ic2_timeseries) / ic2_denom
                else:
                    ic2_weights[ch_c] = 0.0

    topo_channels = []
    ic1_w = []
    # IC Topographies: horizontal bar charts
    if chan_labels:
        topo_channels = [c for c in chan_labels if c != 'Average']
        ic1_w = [eig_means[chan_labels.index(c)] for c in topo_channels if c in chan_labels]
        ic2_w = [ic2_weights.get(c, 0.0) for c in topo_channels]
        if topo_channels and ic1_w:
            fig_topo, axes_topo = plt.subplots(1, 2, figsize=(14, max(4, len(topo_channels) * 0.3)))
            cmap_div = plt.get_cmap('RdBu_r')
            for ax_t, weights_t, ic_label in zip(axes_topo, [ic1_w, ic2_w],
                                                  ['IC #1 — ECG Regression Component',
                                                   'IC #2 — Residual Principal Component']):
                w_arr = np.array(weights_t)
                abs_max = np.abs(w_arr).max() + 1e-12
                colors_t = [cmap_div(0.5 + 0.5 * w / abs_max) for w in w_arr]
                bars_t = ax_t.barh(topo_channels, w_arr, color=colors_t, edgecolor='black', linewidth=0.4)
                ax_t.axvline(0, color='black', linewidth=0.8)
                ax_t.set_title(ic_label, fontsize=10)
                ax_t.set_xlabel("Regression weight", fontsize=9)
                ax_t.tick_params(axis='y', labelsize=7)
                ax_t.grid(axis='x', alpha=0.3)
            fig_topo.suptitle(
                f"IC Topographies — {selected_group} / {selected_stage}",
                fontsize=11, fontweight='bold'
            )
            fig_topo.tight_layout()
            st.pyplot(fig_topo, use_container_width=True)
            plt.close(fig_topo)

    # IC Weight × Time colormaps: electrode (y) vs time (x), color = weight × IC timeseries
    # Shows which electrode at which time carries the most cardiac regression weight
    if topo_channels and ic1_w and times is not None:
        # Build group-mean ECG HEP for IC#1 timeseries
        ecg_traces_wt = []
        for ind_wt in filtered_individuals:
            _pid_wt, _hep_wt, _t_wt, _ch_wt, _rp_wt, _ecg_wt, _ech_wt = ind_wt[:7]
            if _ecg_wt is None:
                continue
            _ecg_wt_sq = np.asarray(_ecg_wt).squeeze()
            if _ecg_wt_sq.ndim == 1 and len(_ecg_wt_sq) == len(times):
                ecg_traces_wt.append(_ecg_wt_sq)
        if ecg_traces_wt:
            ic1_ts_mean = np.mean(ecg_traces_wt, axis=0)  # (n_times,)
            # IC#1 weight-time matrix: channels × time = eig_weight[ch] * ic1_timeseries[t]
            ic1_wt_matrix = np.array(ic1_w)[:, np.newaxis] * ic1_ts_mean[np.newaxis, :]  # (n_ch, n_times)
            ic1_wt_matrix *= 1e6  # → µV

            n_panels_wt = 2 if (ic2_timeseries is not None and ic2_weights) else 1
            fig_wt, axes_wt = plt.subplots(n_panels_wt, 1,
                                           figsize=(14, 3.5 * n_panels_wt + 1),
                                           squeeze=False)
            times_ms_wt = times * 1000

            # Panel 1 — IC #1
            ax_wt = axes_wt[0, 0]
            vmax_wt = np.abs(ic1_wt_matrix).max() + 1e-12
            im1 = ax_wt.imshow(ic1_wt_matrix, aspect='auto', origin='upper',
                               extent=[times_ms_wt[0], times_ms_wt[-1],
                                       len(topo_channels), 0],
                               cmap='RdBu_r', vmin=-vmax_wt, vmax=vmax_wt)
            ax_wt.set_yticks(np.arange(len(topo_channels)) + 0.5)
            ax_wt.set_yticklabels(topo_channels, fontsize=7)
            ax_wt.axvline(0, color='black', linewidth=1.2, linestyle='--', label='R-peak (t=0)')
            ax_wt.set_xlabel("Time (ms)", fontsize=9)
            ax_wt.set_ylabel("Electrode", fontsize=9)
            ax_wt.set_title("IC #1 — Regression Weight × IC Timeseries  (electrode × time)", fontsize=10)
            plt.colorbar(im1, ax=ax_wt, shrink=0.9, label='Weight × amplitude (µV)')

            # Panel 2 — IC #2 (if available)
            if n_panels_wt == 2:
                ic2_w_ordered = [ic2_weights.get(c, 0.0) for c in topo_channels]
                ic2_wt_matrix = np.array(ic2_w_ordered)[:, np.newaxis] * ic2_timeseries[np.newaxis, :] * 1e6
                ax_wt2 = axes_wt[1, 0]
                vmax_wt2 = np.abs(ic2_wt_matrix).max() + 1e-12
                im2 = ax_wt2.imshow(ic2_wt_matrix, aspect='auto', origin='upper',
                                    extent=[times_ms_wt[0], times_ms_wt[-1],
                                            len(topo_channels), 0],
                                    cmap='RdBu_r', vmin=-vmax_wt2, vmax=vmax_wt2)
                ax_wt2.set_yticks(np.arange(len(topo_channels)) + 0.5)
                ax_wt2.set_yticklabels(topo_channels, fontsize=7)
                ax_wt2.axvline(0, color='black', linewidth=1.2, linestyle='--')
                ax_wt2.set_xlabel("Time (ms)", fontsize=9)
                ax_wt2.set_ylabel("Electrode", fontsize=9)
                ax_wt2.set_title("IC #2 — Regression Weight × IC Timeseries  (electrode × time)", fontsize=10)
                plt.colorbar(im2, ax=ax_wt2, shrink=0.9, label='Weight × amplitude (µV)')

            fig_wt.suptitle(
                f"IC Regression Weight × Time — {selected_group} / {selected_stage}\n"
                "Red = positive cardiac contribution, Blue = negative. "
                "Bright bands = electrodes with strongest weight at that time.",
                fontsize=10, fontweight='bold'
            )
            fig_wt.tight_layout()
            st.pyplot(fig_wt, use_container_width=True)
            plt.close(fig_wt)
            st.caption(
                "Each row = one electrode. Each column = one time point (ms, aligned to R-peak at 0). "
                "Colour = regression weight × IC timeseries amplitude — shows which electrode carries "
                "the most cardiac-field contribution at each moment."
            )


    # ------------------------------------------------------------------ back-projection panel
    st.markdown("#### Back-Projection of Cardiac Field Artefact (CFA) ICs")
    st.caption(
        "Shows how IC #1 and IC #2 reconstruct the cardiac-field artefact. "
        "Orange = IC #1 contribution, green = IC #2 contribution, black dashed = combined reconstruction. "
        "Right panel: variance explained per channel."
    )

    avg_orig_list = ch_data['Average']['orig']
    avg_eig_list  = ch_data['Average']['eig']
    ecg_list_bp   = []
    for ind_bp in filtered_individuals:
        _pid_bp, _hep_bp, _t_bp, _ch_bp, _rp_bp, _ecg_bp, _ech_bp = ind_bp[:7]
        if _ecg_bp is None:
            continue
        _ecg_bp_sq = np.asarray(_ecg_bp).squeeze()
        if _ecg_bp_sq.ndim != 1 or len(_ecg_bp_sq) != len(times):
            continue
        ecg_list_bp.append(_ecg_bp_sq)

    if avg_orig_list and ecg_list_bp and times is not None:
        n_bp = min(len(avg_orig_list), len(avg_eig_list), len(ecg_list_bp))
        ic1_backprojs, ic2_backprojs, orig_traces_bp = [], [], []
        for i_bp in range(n_bp):
            eig_bp = avg_eig_list[i_bp]
            ecg_bp = ecg_list_bp[i_bp]
            orig_bp = avg_orig_list[i_bp]
            ic1_bp = eig_bp * ecg_bp
            ic1_backprojs.append(ic1_bp)
            orig_traces_bp.append(orig_bp)
            if ic2_timeseries is not None:
                res_bp = orig_bp - ic1_bp
                ic2_w_bp = np.dot(res_bp, ic2_timeseries) / (np.dot(ic2_timeseries, ic2_timeseries) + 1e-20)
                ic2_backprojs.append(ic2_w_bp * ic2_timeseries)

        ic1_mean_bp = np.mean(ic1_backprojs, axis=0) * 1e6  # µV
        orig_mean_bp = np.mean(orig_traces_bp, axis=0) * 1e6
        ic2_mean_bp = np.mean(ic2_backprojs, axis=0) * 1e6 if ic2_backprojs else np.zeros_like(ic1_mean_bp)
        combined_bp = ic1_mean_bp + ic2_mean_bp

        # Variance explained per channel
        var_labels_bp, var_ic1_bp, var_ic2_bp = [], [], []
        for ch_vp in common_channels:
            if not ch_data[ch_vp]['orig'] or not ch_data[ch_vp]['eig']:
                continue
            n_vp = min(len(ch_data[ch_vp]['orig']), len(ch_data[ch_vp]['eig']), len(ecg_list_bp))
            if n_vp == 0:
                continue
            var_orig_vp = np.mean([np.var(ch_data[ch_vp]['orig'][j]) for j in range(n_vp)]) + 1e-30
            var_ic1_vp  = np.mean([np.var(ch_data[ch_vp]['eig'][j] * ecg_list_bp[j]) for j in range(n_vp)])
            var_labels_bp.append(ch_vp)
            var_ic1_bp.append(min(100.0, var_ic1_vp / var_orig_vp * 100))
            if ic2_timeseries is not None:
                res_vp_list = [ch_data[ch_vp]['orig'][j] - ch_data[ch_vp]['eig'][j]*ecg_list_bp[j] for j in range(n_vp)]
                ic2_w_vp = [np.dot(r, ic2_timeseries)/(np.dot(ic2_timeseries, ic2_timeseries)+1e-20) for r in res_vp_list]
                var_ic2_vp = np.mean([np.var(w*ic2_timeseries) for w in ic2_w_vp])
                var_ic2_bp.append(min(100.0 - var_ic1_bp[-1], var_ic2_vp / var_orig_vp * 100))
            else:
                var_ic2_bp.append(0.0)

        total_var_mean = np.mean([v1 + v2 for v1, v2 in zip(var_ic1_bp, var_ic2_bp)]) if var_ic1_bp else 0.0

        fig_bp, axes_bp = plt.subplots(1, 2, figsize=(14, 5))
        # Left: waveforms
        ax_left = axes_bp[0]
        ax_left.plot(times, ic1_mean_bp, color='darkorange', linewidth=2, label='IC #1 (ECG component)')
        if ic2_backprojs:
            ax_left.plot(times, ic2_mean_bp, color='green', linewidth=2, label='IC #2 (Residual PC)')
        ax_left.plot(times, combined_bp, color='black', linewidth=2, linestyle='--', label='Combined (IC1+IC2)')
        ax_left.axvline(0, color='gray', linewidth=0.8, linestyle='--')
        ax_left.axhline(0, color='gray', linewidth=0.5)
        ax_left.set_title("Back-Projection of Cardiac ICs", fontsize=10)
        ax_left.set_xlabel("Time (s)")
        ax_left.set_ylabel("Amplitude (µV)")
        ax_left.legend(fontsize=8)
        ax_left.grid(alpha=0.2)
        ax_left.annotate(
            f"IC1+IC2 explain {total_var_mean:.1f}% of cardiac artefact variance (group mean)",
            xy=(0.02, 0.97), xycoords='axes fraction', fontsize=8,
            va='top', ha='left', color='darkred',
            bbox=dict(boxstyle='round,pad=0.3', facecolor='lightyellow', alpha=0.8)
        )
        # Right: variance explained bar chart
        ax_right = axes_bp[1]
        if var_labels_bp:
            x_vp = np.arange(len(var_labels_bp))
            ax_right.bar(x_vp, var_ic1_bp, color='steelblue', alpha=0.8, label='IC #1', edgecolor='black', linewidth=0.4)
            ax_right.bar(x_vp, var_ic2_bp, bottom=var_ic1_bp, color='green', alpha=0.7, label='IC #2', edgecolor='black', linewidth=0.4)
            ax_right.set_xticks(x_vp)
            ax_right.set_xticklabels(var_labels_bp, rotation=45, ha='right', fontsize=7)
            ax_right.set_ylabel("Variance explained (%)", fontsize=9)
            ax_right.set_title("Cardiac Artefact Variance Explained per Channel", fontsize=10)
            ax_right.legend(fontsize=8)
            ax_right.grid(axis='y', alpha=0.3)
            ax_right.set_ylim(0, 100)
        fig_bp.suptitle(
            f"CFA Back-Projection — {selected_group} / {selected_stage}",
            fontsize=11, fontweight='bold'
        )
        fig_bp.tight_layout()
        st.pyplot(fig_bp, use_container_width=True)
        plt.close(fig_bp)

    # ------------------------------------------------------------------ ICA+CSD per-patient lookup
    # Build _csd_ica_lookup[pid][ch] = CSD(ICA-cleaned) trace (V/m²)
    # Must run before the per-channel plot loop so it's available in subplots.
    _csd_ica_lookup: dict = {}
    _sfq_pre = float(1.0 / np.mean(np.diff(times))) if len(times) > 1 else 500.0
    _tmin_pre = float(times[0])

    def _apply_csd_pre(mc_data, ch_names_list):
        try:
            _info = mne.create_info(ch_names=list(ch_names_list), sfreq=_sfq_pre, ch_types='eeg')
            _ev = mne.EvokedArray(mc_data.astype(float), _info, tmin=_tmin_pre)
            _mont = mne.channels.make_standard_montage('standard_1005')
            _ev.set_montage(_mont, match_case=False, on_missing='ignore', verbose=False)
            _kept = [c for c in _ev.ch_names if c in _mont.ch_names]
            if len(_kept) < 3:
                return None, None
            _ev.pick_channels(_kept)
            _ev_csd = mne.preprocessing.compute_current_source_density(_ev)
            return _ev_csd.data.copy(), list(_ev_csd.ch_names)
        except Exception:
            return None, None

    for _ind_pre in filtered_individuals:
        _pid_pre, _hep_pre, _t_pre, _ch_pre, _, _ecg_pre, _ = _ind_pre[:7]
        if _ecg_pre is None:
            continue
        _ecg_pre_v = np.asarray(_ecg_pre).squeeze()
        if _ecg_pre_v.ndim != 1 or len(_ecg_pre_v) != len(times):
            continue
        _den_pre = np.dot(_ecg_pre_v, _ecg_pre_v)
        if _den_pre < 1e-20:
            continue
        _vchs = [c for c in common_channels if c in _ch_pre]
        _vidx = [_ch_pre.index(c) for c in _vchs]
        if len(_vchs) < 3:
            continue
        _raw_pre = _hep_pre[_vidx].astype(float)
        _eigs_pre = np.dot(_raw_pre, _ecg_pre_v) / _den_pre
        _ica_pre = _raw_pre - _eigs_pre[:, None] * _ecg_pre_v[None, :]
        _csd_pre_d, _csd_pre_chs = _apply_csd_pre(_ica_pre, _vchs)
        if _csd_pre_d is None:
            continue
        _csd_ica_lookup[_pid_pre] = {ch: _csd_pre_d[i] for i, ch in enumerate(_csd_pre_chs)}
        _csd_ica_lookup[_pid_pre]['Average'] = _csd_pre_d.mean(axis=0)

    # ------------------------------------------------------------------ per-channel plots
    st.markdown("#### Per-Channel Before / After ICA Cleaning")
    with st.expander("ℹ️ Cleaning Method Comparison — about this section", expanded=False):
        st.markdown("#### Individual Patient Averages: Cleaning Method Comparison")
        st.caption(
            "Group-average and individual patient HEP traces for 4 cleaning methods. "
            "CLEEGN: deep learning CNN from CECNL/CLEEGN (https://github.com/CECNL/CLEEGN). "
            "Thin lines = individual patients; thick lines = group average."
        )
    for ch_name in display_channels:
        d = ch_data[ch_name]
        if not d['orig']:
            continue
        orig_arr = np.array(d['orig']) * 1e6   # → µV
        clean_arr = np.array(d['clean']) * 1e6
        eigs = np.array(d['eig'])
        pids = d['pids']
        n_subj = len(pids)

        with st.expander(f"ICA Cleaning — Channel: {ch_name}  (n={n_subj})", expanded=(ch_name == 'Average')):
            _ck = f"ica_ch_{ch_name}_{selected_group}_{selected_stage}"
            _do_render = st.session_state.get(_ck, ch_name == 'Average')
            if not _do_render:
                st.button(
                    f"▶ Compute analysis for {ch_name}",
                    key=f"btn_{_ck}",
                    on_click=lambda k=_ck: st.session_state.update({k: True}),
                )
            if _do_render:
                fig, axes = plt.subplots(1, 2, figsize=(14, 5))

                # --- subplot 1: per-patient original
                ax = axes[0]
                cmap = plt.get_cmap('tab20' if n_subj <= 20 else 'hsv', max(n_subj, 1))
                colors = [cmap(i / max(n_subj - 1, 1)) for i in range(n_subj)]
                for i, (trace, pid) in enumerate(zip(orig_arr, pids)):
                    ax.plot(times, trace, color=colors[i], alpha=0.4, linewidth=0.8, label=pid)
                if n_subj:
                    ax.plot(times, orig_arr.mean(axis=0), color='black', linewidth=2, label='Group avg')
                ax.set_title("Original HEP", fontsize=10)
                ax.set_xlabel("Time (s)")
                ax.set_ylabel("Amplitude (µV)")
                ax.axvline(0, color='gray', linewidth=0.8, linestyle='--')
                ax.axhline(0, color='gray', linewidth=0.5)
                ax.grid(alpha=0.2)

                # --- subplot 2: per-patient cleaned
                ax = axes[1]
                for i, (trace, pid) in enumerate(zip(clean_arr, pids)):
                    ax.plot(times, trace, color=colors[i], alpha=0.4, linewidth=0.8)
                if n_subj:
                    ax.plot(times, clean_arr.mean(axis=0), color='black', linewidth=2)
                ax.set_title("Cleaned HEP (ECG removed)", fontsize=10)
                ax.set_xlabel("Time (s)")
                ax.axvline(0, color='gray', linewidth=0.8, linestyle='--')
                ax.axhline(0, color='gray', linewidth=0.5)
                ax.grid(alpha=0.2)

                fig.suptitle(
                    f"ICA ECG Cleaning — Channel: {ch_name} | {selected_group} / {selected_stage}",
                    fontsize=11, fontweight='bold'
                )
                fig.tight_layout()
                st.pyplot(fig, use_container_width=True)
                plt.close(fig)

                # ---- CLEEGN comparison: Raw vs ICA vs ICA+CLEEGN vs CLEEGN-only ----

                def _apply_cleegn_to_arr(traces_uv):
                    """Apply CLEEGN CNN to each row of traces_uv (n_subj, n_times) in µV.
                    Returns cleaned array same shape, or None if CLEEGN unavailable."""
                    try:
                        import subprocess, sys, os, tempfile
                        import torch

                        _cleegn_root = '/tmp/CLEEGN'
                        if not os.path.isdir(_cleegn_root):
                            subprocess.run(
                                ['git', 'clone', '--depth', '1',
                                 'https://github.com/CECNL/CLEEGN.git', _cleegn_root],
                                check=True, capture_output=True, timeout=120
                            )

                        if _cleegn_root not in sys.path:
                            sys.path.insert(0, _cleegn_root)

                        from utils.cleegn import CLEEGN as _CLEEGNModel  # noqa

                        _device = torch.device('cpu')
                        _n_times = traces_uv.shape[1]
                        _fs_hep = 1.0 / float(np.mean(np.diff(times))) if len(times) > 1 else 500.0

                        # Instantiate with n_chan=1 for single-channel HEP traces
                        _model = _CLEEGNModel(n_chan=1, fs=_fs_hep, N_F=1).to(_device)
                        _model.eval()

                        _win = min(_n_times, max(4, int(4.0 * _fs_hep)))
                        _stride = max(1, _win // 8)

                        cleaned = np.zeros_like(traces_uv)
                        for _si in range(traces_uv.shape[0]):
                            _x = traces_uv[_si].copy().astype(np.float32)
                            _std = _x.std() + 1e-10
                            _xn = (_x / _std).reshape(1, -1)  # (1, n_times)

                            _out = np.zeros(_n_times, dtype=np.float32)
                            _wsum = np.zeros(_n_times, dtype=np.float32)
                            from scipy.signal import windows as _wins
                            _hwin = _wins.hann(_win).astype(np.float32) + 1e-9

                            for _i0 in range(0, _n_times, _stride):
                                _i1 = _i0 + _win
                                if _i1 > _n_times:
                                    _i0 = _n_times - _win
                                    _i1 = _n_times
                                _seg = _xn[:, _i0:_i1]  # (1, win)
                                with torch.no_grad():
                                    _inp = torch.from_numpy(_seg[np.newaxis, np.newaxis]).to(_device)
                                    _res = _model(_inp).detach().cpu().squeeze().numpy()
                                if _res.ndim == 0:
                                    _res = np.full(_win, float(_res))
                                _res = np.asarray(_res).flatten()[:_win]
                                _out[_i0:_i1] += _res * _hwin
                                _wsum[_i0:_i1] += _hwin
                                if _i0 == _n_times - _win:
                                    break

                            _wsum = np.where(_wsum < 1e-9, 1.0, _wsum)
                            cleaned[_si] = (_out / _wsum) * _std

                        return cleaned

                    except Exception as _ce:
                        return None, str(_ce)

                _cleegn_available = False
                _cleegn_ica_arr = None
                _cleegn_raw_arr = None
                _cleegn_err = None

                _cleegn_result_ica = _apply_cleegn_to_arr(clean_arr)
                if isinstance(_cleegn_result_ica, tuple):
                    _cleegn_err = _cleegn_result_ica[1]
                else:
                    _cleegn_ica_arr = _cleegn_result_ica
                    _cleegn_raw_result = _apply_cleegn_to_arr(orig_arr)
                    if isinstance(_cleegn_raw_result, tuple):
                        _cleegn_err = _cleegn_raw_result[1]
                    else:
                        _cleegn_raw_arr = _cleegn_raw_result
                        _cleegn_available = True

                _fig_cmp, _ax_cmp = plt.subplots(figsize=(14, 6))

                # Individual patient thin lines
                _cmp_cmap = plt.get_cmap('tab20' if n_subj <= 20 else 'hsv', max(n_subj, 1))
                _cmp_colors = [_cmp_cmap(i / max(n_subj - 1, 1)) for i in range(n_subj)]
                for _pi in range(n_subj):
                    _ax_cmp.plot(times, orig_arr[_pi], color='black', alpha=0.1, linewidth=0.6)
                    _ax_cmp.plot(times, clean_arr[_pi], color='steelblue', alpha=0.1, linewidth=0.6)
                    if _cleegn_available:
                        _ax_cmp.plot(times, _cleegn_ica_arr[_pi], color='green', alpha=0.1, linewidth=0.6)
                        _ax_cmp.plot(times, _cleegn_raw_arr[_pi], color='red', alpha=0.1, linewidth=0.6)

                # Group average thick lines
                _ax_cmp.plot(times, orig_arr.mean(axis=0), color='black', linewidth=2.5, label='Raw')
                _ax_cmp.plot(times, clean_arr.mean(axis=0), color='steelblue', linewidth=2.5, label='ICA')
                if _cleegn_available:
                    _ax_cmp.plot(times, _cleegn_ica_arr.mean(axis=0), color='green', linewidth=2.5, label='ICA + CLEEGN')
                    _ax_cmp.plot(times, _cleegn_raw_arr.mean(axis=0), color='red', linewidth=2.5, label='CLEEGN only')
                else:
                    _ax_cmp.text(
                        0.5, 0.5,
                        f'CLEEGN not available\n{_cleegn_err or ""}',
                        transform=_ax_cmp.transAxes,
                        ha='center', va='center', fontsize=9, color='gray',
                        bbox=dict(boxstyle='round', facecolor='lightyellow', alpha=0.7)
                    )

                _ax_cmp.axvline(0, color='gray', linewidth=0.8, linestyle='--')
                _ax_cmp.axhline(0, color='gray', linewidth=0.5)
                _ax_cmp.set_xlabel("Time (s)")
                _ax_cmp.set_ylabel("Amplitude (µV)")
                _ax_cmp.legend(fontsize=10)
                _ax_cmp.grid(alpha=0.2)
                _fig_cmp.suptitle(
                    f"Cleaning Method Comparison — Channel: {ch_name} | {selected_group} / {selected_stage}\n"
                    f"Raw vs ICA vs ICA+CLEEGN vs CLEEGN-only  (n={n_subj} patients)",
                    fontsize=11, fontweight='bold'
                )
                _fig_cmp.tight_layout()
                st.pyplot(_fig_cmp, use_container_width=True)
                plt.close(_fig_cmp)

                if _cleegn_err and not _cleegn_available:
                    st.warning(f"CLEEGN could not be applied: {_cleegn_err}")
                # ---- end CLEEGN comparison ----

                # Individual patient subplots
                ncols = 4
                nrows = math.ceil(n_subj / ncols)
                fig2, axes2 = plt.subplots(nrows, ncols, figsize=(14, 3*nrows), sharex=True, sharey=True)
                if nrows == 1:
                    axes2 = axes2.reshape(1, -1)  # ensure 2D
                axes2 = axes2.flatten()

                # Calculate variance explained by IC1 and IC2 for each patient
                var_ic1_per_patient = []
                var_ic2_per_patient = []
                for i, pid in enumerate(pids):
                    orig_var = np.var(orig_arr[i]) + 1e-20

                    # IC1 contribution = original - cleaned
                    ic1_contrib = orig_arr[i] - clean_arr[i]
                    ic1_var = np.var(ic1_contrib)
                    ic1_pct = min(100.0, (ic1_var / orig_var) * 100)
                    var_ic1_per_patient.append(ic1_pct)

                    # IC2 contribution = residual after removing IC1
                    ic2_pct = 0.0
                    if ic2_timeseries is not None and ch_name == 'Average':
                        try:
                            for j_match, ind_match in enumerate(filtered_individuals):
                                if ind_match[0] == pid:
                                    if j_match < len(ecg_list_ic2):
                                        ecg_for_pid = ecg_list_ic2[j_match]
                                        eig_for_pid = eigs[i]
                                        residual = orig_arr[i]/1e6 - eig_for_pid * ecg_for_pid
                                        ic2_w = np.dot(residual, ic2_timeseries) / (np.dot(ic2_timeseries, ic2_timeseries) + 1e-20)
                                        ic2_contrib_var = np.var(ic2_w * ic2_timeseries)
                                        ic2_pct = min(100.0 - ic1_pct, (ic2_contrib_var / orig_var) * 100)
                                    break
                        except Exception:
                            ic2_pct = 0.0

                    var_ic2_per_patient.append(ic2_pct)

                # --- z-score each patient trace (normalise by original stats) ---
                _sig_mask  = (times >= 0.05) & (times <= 0.50)
                _noise_mask = (times >= -0.40) & (times <= -0.05)
                z_orig_list      = []
                z_clean_list     = []
                z_ecg_list       = []
                z_csd_ica_list   = []
                snr_orig_list    = []
                snr_clean_list   = []
                snr_csd_ica_list = []
                _ecg_raw = d.get('ecg', [])
                for i in range(len(pids)):
                    _mu  = orig_arr[i].mean()
                    _std = orig_arr[i].std() + 1e-12
                    _zo  = (orig_arr[i] - _mu) / _std
                    _zc  = (clean_arr[i] - _mu) / _std
                    z_orig_list.append(_zo)
                    z_clean_list.append(_zc)
                    if i < len(_ecg_raw):
                        _ecg_i = np.asarray(_ecg_raw[i], dtype=float)
                        _ecg_mu = _ecg_i.mean()
                        _ecg_std = _ecg_i.std() + 1e-12
                        z_ecg_list.append((_ecg_i - _ecg_mu) / _ecg_std)
                    else:
                        z_ecg_list.append(None)
                    _csd_tr = _csd_ica_lookup.get(pids[i], {}).get(ch_name)
                    if _csd_tr is not None and len(_csd_tr) == len(times):
                        _csd_mu  = _csd_tr.mean()
                        _csd_std = _csd_tr.std() + 1e-12
                        _z_csd   = (_csd_tr - _csd_mu) / _csd_std
                    else:
                        _z_csd = None
                    z_csd_ica_list.append(_z_csd)
                    _rms = lambda v: np.sqrt(np.mean(v ** 2))
                    _snr_o = 20.0 * np.log10(_rms(_zo[_sig_mask]) / (_rms(_zo[_noise_mask]) + 1e-12)) if _noise_mask.any() and _sig_mask.any() else 0.0
                    _snr_c = 20.0 * np.log10(_rms(_zc[_sig_mask]) / (_rms(_zc[_noise_mask]) + 1e-12)) if _noise_mask.any() and _sig_mask.any() else 0.0
                    snr_orig_list.append(float(_snr_o))
                    snr_clean_list.append(float(_snr_c))
                    if _z_csd is not None:
                        _snr_csd = 20.0 * np.log10(_rms(_z_csd[_sig_mask]) / (_rms(_z_csd[_noise_mask]) + 1e-12)) if _noise_mask.any() and _sig_mask.any() else 0.0
                        snr_csd_ica_list.append(float(_snr_csd))
                    else:
                        snr_csd_ica_list.append(float('nan'))

                # Filter to patients within the user-selected SNR range (on cleaned trace)
                _keep = [i for i in range(len(pids)) if snr_min <= snr_clean_list[i] <= snr_max]
                _n_kept = len(_keep)

                if _n_kept == 0:
                    st.caption(f"No patients with SNR in [{snr_min:.1f}, {snr_max:.1f}] dB on channel {ch_name}.")
                else:
                    st.subheader(f"Individual Patient Averages (SNR {snr_min:.1f}–{snr_max:.1f} dB) — {ch_name} | {selected_group} / {selected_stage}")
                    _ncols2 = 4
                    _nrows2 = math.ceil(_n_kept / _ncols2)
                    fig2, axes2 = plt.subplots(_nrows2, _ncols2, figsize=(14, 3 * _nrows2), sharex=True, sharey=True)
                    if _nrows2 == 1:
                        axes2 = np.array(axes2).reshape(1, -1)
                    axes2 = axes2.flatten()

                    for _plot_idx, i in enumerate(_keep):
                        ax = axes2[_plot_idx]
                        ax.plot(times, z_orig_list[i], color='red', alpha=0.7, linewidth=1, label='Original')
                        ax.plot(times, z_clean_list[i], color='blue', alpha=0.7, linewidth=1, label='Cleaned')
                        _ze = z_ecg_list[i]
                        _r2_str = ""
                        if _ze is not None and len(_ze) == len(times):
                            _y_lo = min(z_orig_list[i].min(), z_clean_list[i].min())
                            _y_hi = max(z_orig_list[i].max(), z_clean_list[i].max())
                            _ze_min, _ze_max = _ze.min(), _ze.max()
                            _ze_range = _ze_max - _ze_min if (_ze_max - _ze_min) > 1e-12 else 1.0
                            _ze_scaled = (_ze - _ze_min) / _ze_range * (_y_hi - _y_lo) + _y_lo
                            ax.plot(times, _ze_scaled, color='green', alpha=0.5, linewidth=0.8, linestyle='--', label='ECG')
                            _r = np.corrcoef(z_orig_list[i], _ze)[0, 1]
                            _r2_str = f"\nR²={_r**2 * 100:.1f}%"
                        ax.set_title(
                            f"{pids[i]}\nSNR: {snr_orig_list[i]:.1f}→{snr_clean_list[i]:.1f} dB{_r2_str}",
                            fontsize=7,
                        )
                        ax.axvline(0, color='gray', linewidth=0.8, linestyle='--')
                        ax.axhline(0, color='gray', linewidth=0.5)
                        ax.grid(alpha=0.2)
                        if _plot_idx == 0:
                            ax.legend(fontsize=6)
                    for j in range(_n_kept, len(axes2)):
                        axes2[j].set_visible(False)
                    fig2.suptitle(
                        f"Individual Patient Averages (SNR {snr_min:.1f}–{snr_max:.1f} dB) — {ch_name} | {selected_group} / {selected_stage}\n"
                        f"({_n_kept}/{len(pids)} patients shown)",
                        fontsize=11,
                    )
                    fig2.text(0.5, 0.02, 'Time (s)', ha='center', fontsize=10)
                    fig2.text(0.02, 0.5, 'Amplitude (z-score)', va='center', rotation='vertical', fontsize=10)
                    fig2.tight_layout(rect=[0.03, 0.03, 1, 0.95])
                    st.pyplot(fig2, use_container_width=True)
                    plt.close(fig2)

    # ------------------------------------------------------------------ interactive plotly hover chart
    try:
        from plotly.subplots import make_subplots
        import plotly.graph_objects as go

        _QUALITATIVE = [
            "#1f77b4", "#ff7f0e", "#2ca02c", "#d62728", "#9467bd",
            "#8c564b", "#e377c2", "#7f7f7f", "#bcbd22", "#17becf",
            "#aec7e8", "#ffbb78", "#98df8a", "#ff9896", "#c5b0d5",
            "#c49c94", "#f7b6d2", "#c7c7c7", "#dbdb8d", "#9edae5",
        ]

        ch_sel = st.selectbox(
            "Select channel for interactive view",
            options=display_channels,
            index=display_channels.index("Average") if "Average" in display_channels else 0,
            key="ica_hover_channel_select",
        )

        _d = ch_data.get(ch_sel, {})
        _orig_list = _d.get("orig", [])
        _clean_list = _d.get("clean", [])
        _pids = _d.get("pids", [])

        fig_hover = make_subplots(
            rows=1, cols=2,
            shared_yaxes=True,
            subplot_titles=("Original HEP", "Cleaned HEP"),
        )

        for _i, (_pid, _orig_trace, _clean_trace) in enumerate(
            zip(_pids, _orig_list, _clean_list)
        ):
            _color = _QUALITATIVE[_i % len(_QUALITATIVE)]
            _orig_uv = np.asarray(_orig_trace) * 1e6
            _clean_uv = np.asarray(_clean_trace) * 1e6

            fig_hover.add_trace(
                go.Scatter(
                    x=times,
                    y=_orig_uv,
                    mode="lines",
                    name=_pid,
                    line=dict(color=_color, width=1),
                    opacity=0.25,
                    showlegend=False,
                    hovertemplate="%{fullData.name}<br>t=%{x:.3f}s<br>%{y:.2f} µV",
                ),
                row=1, col=1,
            )
            fig_hover.add_trace(
                go.Scatter(
                    x=times,
                    y=_clean_uv,
                    mode="lines",
                    name=_pid,
                    line=dict(color=_color, width=1),
                    opacity=0.25,
                    showlegend=False,
                    hovertemplate="%{fullData.name}<br>t=%{x:.3f}s<br>%{y:.2f} µV",
                ),
                row=1, col=2,
            )

        if _orig_list:
            _orig_mean = np.array(_orig_list).mean(axis=0) * 1e6
            _clean_mean = np.array(_clean_list).mean(axis=0) * 1e6
            fig_hover.add_trace(
                go.Scatter(
                    x=times,
                    y=_orig_mean,
                    mode="lines",
                    name="Group average",
                    line=dict(color="black", width=3),
                    showlegend=True,
                    hovertemplate="Group average<br>t=%{x:.3f}s<br>%{y:.2f} µV",
                ),
                row=1, col=1,
            )
            fig_hover.add_trace(
                go.Scatter(
                    x=times,
                    y=_clean_mean,
                    mode="lines",
                    name="Group average",
                    line=dict(color="black", width=3),
                    showlegend=False,
                    hovertemplate="Group average<br>t=%{x:.3f}s<br>%{y:.2f} µV",
                ),
                row=1, col=2,
            )

        fig_hover.add_vline(
            x=0,
            line=dict(color="gray", width=1, dash="dash"),
            row=1, col=1,
        )
        fig_hover.add_vline(
            x=0,
            line=dict(color="gray", width=1, dash="dash"),
            row=1, col=2,
        )

        fig_hover.update_layout(
            height=520,
            title=dict(
                text=f"ICA ECG Cleaning — Interactive | {ch_sel} | {selected_group} / {selected_stage}",
                font=dict(size=13),
            ),
            hovermode="closest",
        )
        fig_hover.update_xaxes(title_text="Time (s)")
        fig_hover.update_yaxes(title_text="Amplitude (µV)", col=1)

        st.plotly_chart(fig_hover, use_container_width=True)

    except ImportError:
        st.info("Install plotly for the interactive chart: pip install plotly")

    # ------------------------------------------------------------------ CSD + SNR table
    st.markdown("---")
    st.markdown("#### Current Source Density (CSD) Analysis")
    st.caption(
        "CSD (surface Laplacian) sharpens spatial resolution by attenuating volume-conducted "
        "signals. Applied here to both raw and ICA-cleaned HEP data using MNE's standard "
        "10-05 montage. Channels without 3-D positions are dropped automatically."
    )

    # --- SNR helper (time-domain RMS, frequency-domain Welch, √N rule) ---
    def _snr_metrics(trace, times_arr, n_epochs=None, sfreq_hz=None):
        sig_m   = (times_arr >= 0.05)  & (times_arr <= 0.50)
        noise_m = (times_arr >= -0.40) & (times_arr <= -0.05)
        _rms = lambda v: float(np.sqrt(np.mean(v ** 2))) if len(v) > 0 else 1e-12

        # Time-domain: z-scored RMS ratio
        mu_t, sd_t = trace.mean(), trace.std() + 1e-12
        z_t = (trace - mu_t) / sd_t
        td_snr = (20.0 * np.log10(_rms(z_t[sig_m]) / (_rms(z_t[noise_m]) + 1e-12))
                  if sig_m.any() and noise_m.any() else 0.0)

        # Frequency-domain: Welch PSD ratio (1–30 Hz signal / 30–100 Hz noise floor)
        fd_snr = 0.0
        if sfreq_hz and sfreq_hz > 0:
            try:
                from scipy.signal import welch as _welch
                _nperseg = min(len(trace), max(32, int(sfreq_hz)))
                _f, _psd = _welch(trace, fs=sfreq_hz, nperseg=_nperseg)
                _sb = (_f >= 1.0)  & (_f <= 30.0)
                _nb = (_f > 30.0) & (_f <= 100.0)
                if _sb.any() and _nb.any():
                    fd_snr = 10.0 * np.log10(
                        np.mean(_psd[_sb]) / (np.mean(_psd[_nb]) + 1e-20))
            except Exception:
                fd_snr = 0.0

        # √N rule: theoretical SNR gain from averaging N epochs
        sqrt_n_db = 10.0 * np.log10(max(1, n_epochs)) if n_epochs is not None else float('nan')
        return float(td_snr), float(fd_snr), float(sqrt_n_db)

    # --- Infer sampling frequency from times vector ---
    _sfreq_ica = None
    if times is not None and len(times) > 1:
        _sfreq_ica = float(1.0 / np.mean(np.diff(times)))

    # --- MNE CSD helper ---
    def _apply_csd_mne(mc_data, ch_names_list, sfreq_v, tmin_v):
        """Apply MNE surface-Laplacian CSD to mc_data (n_ch × n_times).
        Returns (csd_array, kept_ch_names) or (None, None) on failure."""
        try:
            info_c = mne.create_info(ch_names=list(ch_names_list), sfreq=sfreq_v,
                                     ch_types='eeg')
            ev_c = mne.EvokedArray(mc_data.astype(float), info_c, tmin=tmin_v)
            mont = mne.channels.make_standard_montage('standard_1005')
            ev_c.set_montage(mont, match_case=False, on_missing='ignore',
                             verbose=False)
            kept = [c for c in ev_c.ch_names if c in mont.ch_names]
            if len(kept) < 3:
                return None, None
            ev_c.pick_channels(kept)
            ev_csd = mne.preprocessing.compute_current_source_density(ev_c)
            return ev_csd.data.copy(), list(ev_csd.ch_names)
        except Exception:
            return None, None

    # --- Per-patient CSD computation ---
    _csd_ch = {ch: {'pids': [], 'raw': [], 'ica': [], 'csd_raw': [], 'csd_ica': []}
               for ch in common_channels}
    _csd_avg_d = {'pids': [], 'raw': [], 'ica': [], 'csd_raw': [], 'csd_ica': []}

    for _ind_csd in filtered_individuals:
        _pid_c, _hep_c, _t_c, _ch_c, _rp_c, _ecg_c, _ = _ind_csd[:7]
        if _ecg_c is None or times is None:
            continue

        _valid_chs = [ch for ch in common_channels if ch in _ch_c]
        _valid_idx  = [_ch_c.index(ch) for ch in _valid_chs]
        if len(_valid_chs) < 3:
            continue

        _raw_mc  = _hep_c[_valid_idx].astype(float)              # (n_ch, n_t) in V
        _ica_mc  = apply_ecg_regression(_raw_mc, _ecg_c)         # (n_ch, n_t)

        _sfq = _sfreq_ica if _sfreq_ica else 500.0
        _tmin = float(times[0])

        _csd_raw_d, _csd_chs = _apply_csd_mne(_raw_mc, _valid_chs, _sfq, _tmin)
        _csd_ica_d, _        = _apply_csd_mne(_ica_mc, _valid_chs, _sfq, _tmin)
        if _csd_raw_d is None or _csd_ica_d is None:
            print(f"Skipping CSD for patient {_pid_c} due to montage issues.")
            continue

        # Per-channel storage (only channels that survived montage matching)
        for _ic, _chn in enumerate(_csd_chs):
            if _chn in _csd_ch and _chn in _valid_chs:
                _ri = _valid_chs.index(_chn)
                _csd_ch[_chn]['pids'].append(_pid_c)
                _csd_ch[_chn]['raw'].append(_raw_mc[_ri])
                _csd_ch[_chn]['ica'].append(_ica_mc[_ri])
                _csd_ch[_chn]['csd_raw'].append(_csd_raw_d[_ic])
                _csd_ch[_chn]['csd_ica'].append(_csd_ica_d[_ic])

        # Average across channels
        _csd_avg_d['pids'].append(_pid_c)
        _csd_avg_d['raw'].append(_raw_mc.mean(axis=0))
        _csd_avg_d['ica'].append(_ica_mc.mean(axis=0))
        _csd_avg_d['csd_raw'].append(_csd_raw_d.mean(axis=0))
        _csd_avg_d['csd_ica'].append(_csd_ica_d.mean(axis=0))

    _csd_ok = len(_csd_avg_d['pids']) > 0

    if _csd_ok:
        # --- 2×2 comparison plot ---
        _csd_ch_opts = ['Average'] + [ch for ch in common_channels
                                      if len(_csd_ch.get(ch, {}).get('pids', [])) > 0]
        _ch_csd_sel = st.selectbox(
            "Select channel for CSD comparison plot",
            options=_csd_ch_opts,
            index=0,
            key="csd_channel_select",
        )
        _cd_sel = _csd_avg_d if _ch_csd_sel == 'Average' else _csd_ch.get(_ch_csd_sel, _csd_avg_d)
        _pids_csd_sel = _cd_sel.get('pids', [])

        if _pids_csd_sel:
            _sc = 1e6
            _r_arr  = np.array(_cd_sel['raw'])     * _sc
            _i_arr  = np.array(_cd_sel['ica'])     * _sc
            _cr_arr = np.array(_cd_sel['csd_raw']) * _sc
            _ci_arr = np.array(_cd_sel['csd_ica']) * _sc

            fig_csd, axes_csd = plt.subplots(2, 2, figsize=(14, 7), sharex=True)
            _csd_panels = [
                ('Raw',       _r_arr,  'red'),
                ('ICA',       _i_arr,  'steelblue'),
                ('CSD',       _cr_arr, 'darkorchid'),
                ('ICA + CSD', _ci_arr, 'seagreen'),
            ]
            for _ax_c, (_lbl_c, _arr_c, _col_c) in zip(axes_csd.flatten(), _csd_panels):
                for _tr_c in _arr_c:
                    _ax_c.plot(times, _tr_c, color=_col_c, alpha=0.15, linewidth=0.8)
                _ax_c.plot(times, _arr_c.mean(axis=0), color=_col_c, linewidth=2.5,
                           label='Group mean')
                _ax_c.axvline(0, color='gray', linewidth=0.8, linestyle='--')
                _ax_c.axhline(0, color='gray', linewidth=0.5)
                _ax_c.set_title(_lbl_c, fontsize=11, fontweight='bold')
                _ax_c.set_xlabel("Time (s)", fontsize=9)
                _ax_c.set_ylabel("Amplitude (µV)", fontsize=9)
                _ax_c.grid(alpha=0.2)
                _ax_c.legend(fontsize=8)

            fig_csd.suptitle(
                f"Raw / ICA / CSD / ICA+CSD — {_ch_csd_sel} | {selected_group} / {selected_stage}",
                fontsize=12, fontweight='bold',
            )
            fig_csd.tight_layout()
            st.pyplot(fig_csd, use_container_width=True)
            plt.close(fig_csd)

            # --- Z-scored overlay: all four signals on shared axes ---
            def _zscore_arr(arr2d):
                """Z-score each row independently (per patient)."""
                mu = arr2d.mean(axis=1, keepdims=True)
                sd = arr2d.std(axis=1, keepdims=True) + 1e-12
                return (arr2d - mu) / sd

            _zr  = _zscore_arr(_r_arr)
            _zi  = _zscore_arr(_i_arr)
            _zcr = _zscore_arr(_cr_arr)
            _zci = _zscore_arr(_ci_arr)

            fig_z, axes_z = plt.subplots(2, 2, figsize=(14, 7), sharex=True, sharey=True)
            _z_panels = [
                ('Raw',       _zr,  'red'),
                ('ICA',       _zi,  'steelblue'),
                ('CSD',       _zcr, 'darkorchid'),
                ('ICA + CSD', _zci, 'seagreen'),
            ]
            for _ax_z, (_lbl_z, _arr_z, _col_z) in zip(axes_z.flatten(), _z_panels):
                for _tr_z in _arr_z:
                    _ax_z.plot(times, _tr_z, color=_col_z, alpha=0.15, linewidth=0.8)
                _ax_z.plot(times, _arr_z.mean(axis=0), color=_col_z, linewidth=2.5,
                           label='Group mean')
                _ax_z.axvline(0, color='gray', linewidth=0.8, linestyle='--')
                _ax_z.axhline(0, color='gray', linewidth=0.5)
                _ax_z.set_title(_lbl_z, fontsize=11, fontweight='bold')
                _ax_z.set_xlabel("Time (s)", fontsize=9)
                _ax_z.set_ylabel("Amplitude (z-score)", fontsize=9)
                _ax_z.grid(alpha=0.2)
                _ax_z.legend(fontsize=8)

            fig_z.suptitle(
                f"Z-scored — Raw / ICA / CSD / ICA+CSD — {_ch_csd_sel} | {selected_group} / {selected_stage}",
                fontsize=12, fontweight='bold',
            )
            fig_z.tight_layout()
            st.pyplot(fig_z, use_container_width=True)
            plt.close(fig_z)

        # --- SNR dataframe (Average channel across all patients) ---
        st.markdown("#### SNR Comparison: Raw / ICA / CSD / ICA+CSD")
        st.caption(
            "**TD SNR (dB)** — Time-domain RMS method: 20 × log₁₀(RMS[0.05–0.5 s] / RMS[−0.4–−0.05 s]) "
            "on z-scored trace.  "
            "**FD SNR (dB)** — Frequency-domain PSD method: 10 × log₁₀(PSD[1–30 Hz] / PSD[30–100 Hz]) "
            "via Welch.  "
            "**√N SNR (dB)** — Ensemble-averaging gain: 10 × log₁₀(N_epochs)."
        )

        _n_ep_map = {ind_ep[0]: len(ind_ep[4])
                     for ind_ep in filtered_individuals if ind_ep[4] is not None}
        _snr_rows = []
        for _is, _pid_s in enumerate(_csd_avg_d['pids']):
            _n_ep = _n_ep_map.get(_pid_s)
            _row = {'Patient': _pid_s, 'N epochs': _n_ep if _n_ep else '—'}
            _, _, _sqn = _snr_metrics(
                np.asarray(_csd_avg_d['raw'][_is]), times,
                n_epochs=_n_ep, sfreq_hz=_sfreq_ica)
            _row['√N SNR (dB)'] = round(_sqn, 2)
            for _sk, _sa in [
                ('Raw',     _csd_avg_d['raw']),
                ('ICA',     _csd_avg_d['ica']),
                ('CSD',     _csd_avg_d['csd_raw']),
                ('ICA+CSD', _csd_avg_d['csd_ica']),
            ]:
                if _is < len(_sa):
                    _td, _fd, _ = _snr_metrics(
                        np.asarray(_sa[_is]), times,
                        n_epochs=_n_ep, sfreq_hz=_sfreq_ica)
                    _row[f'{_sk} TD SNR (dB)'] = round(_td, 2)
                    _row[f'{_sk} FD SNR (dB)'] = round(_fd, 2)
            _snr_rows.append(_row)

        if _snr_rows:
            _df_snr = pd.DataFrame(_snr_rows)
            _col_ord = ['Patient', 'N epochs', '√N SNR (dB)']
            for _sig in ['Raw', 'ICA', 'CSD', 'ICA+CSD']:
                _col_ord += [f'{_sig} TD SNR (dB)', f'{_sig} FD SNR (dB)']
            _col_ord = [c for c in _col_ord if c in _df_snr.columns]
            st.dataframe(_df_snr[_col_ord], use_container_width=True)

        # ------------------------------------------------------------------ interactive Plotly median±MAD
        st.markdown("---")
        st.markdown("#### Interactive HEP: Median ± MAD by Method & Electrode")
        _col_left, _col_right = st.columns(2)
        with _col_left:
            _iplot_elec = st.selectbox(
                "Electrode",
                options=_csd_ch_opts,
                index=0,
                key="iplot_electrode",
            )
        with _col_right:
            _iplot_method = st.selectbox(
                "Method",
                options=["Raw", "ICA", "CSD", "ICA + CSD"],
                index=0,
                key="iplot_method",
            )
        _mkey_map = {"Raw": "raw", "ICA": "ica", "CSD": "csd_raw", "ICA + CSD": "csd_ica"}
        _mkey = _mkey_map[_iplot_method]
        _iplot_cd = _csd_avg_d if _iplot_elec == "Average" else _csd_ch.get(_iplot_elec, _csd_avg_d)
        _iplot_arr = np.array(_iplot_cd.get(_mkey, [])) * 1e6  # (n_patients, n_times) µV
        if _iplot_arr.size == 0 or _iplot_arr.ndim < 2:
            st.info("No data available for the selected electrode / method combination.")
        else:
            _iplot_med = np.median(_iplot_arr, axis=0)
            _iplot_mad = np.median(np.abs(_iplot_arr - _iplot_med), axis=0)
            try:
                import plotly.graph_objects as _go
                _color_map_iplot = {
                    "Raw":       "rgba(220,50,50,1)",
                    "ICA":       "rgba(70,130,180,1)",
                    "CSD":       "rgba(153,50,204,1)",
                    "ICA + CSD": "rgba(46,139,87,1)",
                }
                _fill_map_iplot = {
                    "Raw":       "rgba(220,50,50,0.18)",
                    "ICA":       "rgba(70,130,180,0.18)",
                    "CSD":       "rgba(153,50,204,0.18)",
                    "ICA + CSD": "rgba(46,139,87,0.18)",
                }
                _col_iplot = _color_map_iplot.get(_iplot_method, "rgba(70,130,180,1)")
                _fill_iplot = _fill_map_iplot.get(_iplot_method, "rgba(70,130,180,0.18)")
                _times_iplot = list(times)
                _upper_iplot = (_iplot_med + _iplot_mad).tolist()
                _lower_iplot = (_iplot_med - _iplot_mad).tolist()
                _ribbon_iplot = _go.Scatter(
                    x=_times_iplot + _times_iplot[::-1],
                    y=_upper_iplot + _lower_iplot[::-1],
                    fill="toself",
                    fillcolor=_fill_iplot,
                    line=dict(color="rgba(0,0,0,0)"),
                    name="± MAD",
                    showlegend=True,
                    hoverinfo="skip",
                )
                _median_iplot = _go.Scatter(
                    x=_times_iplot,
                    y=_iplot_med.tolist(),
                    mode="lines",
                    line=dict(color=_col_iplot, width=3),
                    name="Median",
                )
                fig_iplot = _go.Figure(data=[_ribbon_iplot, _median_iplot])
                fig_iplot.add_vline(
                    x=0,
                    line_width=1.5,
                    line_dash="dash",
                    line_color="gray",
                    annotation_text="R-peak",
                    annotation_position="top right",
                )
                fig_iplot.update_layout(
                    title=dict(
                        text=f"Median \u00b1 MAD \u2014 {_iplot_method} | {_iplot_elec} | {selected_group} / {selected_stage}",
                        font=dict(size=14),
                    ),
                    xaxis_title="Time (s)",
                    yaxis_title="Amplitude (\u00b5V)",
                    height=450,
                    showlegend=True,
                    legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
                    hovermode="x unified",
                )
                st.plotly_chart(fig_iplot, use_container_width=True)
            except ImportError:
                st.info("Install plotly for the interactive chart: pip install plotly")
    else:
        st.info(
            "CSD requires ≥ 3 common EEG channels matched to MNE's standard 10-05 montage. "
            "No CSD data could be computed for this group/stage."
        )


def handle_snr_optimization(filtered_individuals, selected_group, selected_stage, snr_min=1.5, snr_max=5.0):
    """
    Compares multiple HEP signal cleaning strategies by SNR for each patient.
    Methods: Raw, ECG Regression, Bandpass 1-30 Hz, Bandpass 1-15 Hz.
    SNR = 20 * log10(RMS[0.05-0.50 s] / RMS[-0.40 to -0.05 s]) on z-scored trace.
    """
    st.subheader("SNR Optimization — Compare Cleaning Methods")
    st.caption(
        "Each method is applied to the average HEP across common EEG channels per patient. "
        "SNR = 20 × log10(RMS[0.05–0.5 s] / RMS[−0.4 to −0.05 s]) on the z-scored trace."
    )

    METHOD_OPTIONS = [
        "Raw (no cleaning)",
        "ECG Regression",
        "Bandpass 1–30 Hz",
        "Bandpass 1–15 Hz",
    ]
    selected_methods = st.multiselect(
        "Select cleaning methods to compare",
        options=METHOD_OPTIONS,
        default=["Raw (no cleaning)", "ECG Regression"],
        key="snr_opt_methods",
    )
    if not selected_methods:
        st.warning("Select at least one method.")
        return

    # Gather common channels
    all_ch_sets = [set(ind[3]) for ind in filtered_individuals]
    counts_snr = Counter([ch for s in all_ch_sets for ch in s])
    common_channels_snr = [
        ch for ch, cnt in counts_snr.items()
        if cnt >= len(filtered_individuals) * 0.5
        and (re.match(r'^[a-zA-Z]{1,2}[0-9]+$', ch) or re.match(r'^[a-zA-Z]z$', ch))
    ]
    if not common_channels_snr:
        st.warning("No common EEG channels found.")
        return

    def _snr_opt_compute(trace, times):
        sig_mask  = (times >= 0.05) & (times <= 0.50)
        noise_mask = (times >= -0.40) & (times <= -0.05)
        if not sig_mask.any() or not noise_mask.any():
            return 0.0
        mu, sd = trace.mean(), trace.std() + 1e-12
        z = (trace - mu) / sd
        rms = lambda v: np.sqrt(np.mean(v ** 2))
        return float(20.0 * np.log10(rms(z[sig_mask]) / (rms(z[noise_mask]) + 1e-12)))

    def _bandpass_snr(trace, sfreq, lo, hi):
        nyq = sfreq / 2.0
        lo_n, hi_n = lo / nyq, min(hi / nyq, 0.999)
        if lo_n >= hi_n:
            return trace
        try:
            b, a = butter(4, [lo_n, hi_n], btype='band')
            return filtfilt(b, a, trace)
        except Exception:
            return trace

    results = {m: {} for m in selected_methods}
    times_ref = None

    for ind in filtered_individuals:
        pid, hep_full, t, ch_names, rpeaks, ecg_hep, ecg_ch = ind[:7]
        if times_ref is None:
            times_ref = t
        sfreq_est = 1.0 / (t[1] - t[0]) if len(t) > 1 else 256.0
        valid_idx = [ch_names.index(ch) for ch in common_channels_snr if ch in ch_names]
        if not valid_idx:
            continue
        avg_orig = summarize_channels_without_reference_cancellation(
            hep_full[valid_idx]
        ) * 1e6  # µV

        ecg_signal = None
        if ecg_hep is not None:
            ecg_s = np.asarray(ecg_hep).squeeze()
            if ecg_s.ndim == 1 and len(ecg_s) == len(t):
                ecg_signal = ecg_s

        for method in selected_methods:
            if method == "Raw (no cleaning)":
                trace = avg_orig.copy()
            elif method == "ECG Regression":
                if ecg_signal is not None:
                    denom = np.dot(ecg_signal, ecg_signal)
                    if denom > 1e-20:
                        eig = np.dot(avg_orig, ecg_signal) / denom
                        trace = avg_orig - eig * ecg_signal
                    else:
                        trace = avg_orig.copy()
                else:
                    trace = avg_orig.copy()
            elif method == "Bandpass 1–30 Hz":
                trace = _bandpass_snr(avg_orig, sfreq_est, 1.0, 30.0)
            elif method == "Bandpass 1–15 Hz":
                trace = _bandpass_snr(avg_orig, sfreq_est, 1.0, 15.0)
            else:
                trace = avg_orig.copy()
            results[method][pid] = _snr_opt_compute(trace, t)

    if times_ref is None:
        st.warning("No valid data found.")
        return

    all_pids = sorted({pid for m in results.values() for pid in m.keys()})
    if not all_pids:
        st.warning("No patient data to display.")
        return

    # ── Grouped bar chart ────────────────────────────────────────────────
    n_methods = len(selected_methods)
    n_pids = len(all_pids)
    x = np.arange(n_pids)
    bar_width = 0.8 / max(n_methods, 1)
    method_colors = ['#2196F3', '#FF5722', '#4CAF50', '#9C27B0', '#FF9800']

    fig_snr, ax_snr = plt.subplots(figsize=(max(10, n_pids * 0.5), 5))
    for m_idx, method in enumerate(selected_methods):
        snr_vals = [results[method].get(p, 0.0) for p in all_pids]
        offset = (m_idx - n_methods / 2 + 0.5) * bar_width
        ax_snr.bar(x + offset, snr_vals, width=bar_width * 0.9,
                   label=method, color=method_colors[m_idx % len(method_colors)],
                   alpha=0.8, edgecolor='black', linewidth=0.4)

    ax_snr.axhline(snr_min, color='red', linestyle='--', linewidth=1.2, label=f'Min filter ({snr_min:.1f} dB)')
    ax_snr.axhline(snr_max, color='orange', linestyle='--', linewidth=1.2, label=f'Max filter ({snr_max:.1f} dB)')
    ax_snr.set_xticks(x)
    ax_snr.set_xticklabels(all_pids, rotation=45, ha='right', fontsize=7)
    ax_snr.set_ylabel("SNR (dB)", fontsize=10)
    ax_snr.set_title(
        f"SNR by Cleaning Method — {selected_group} / {selected_stage}\n"
        f"(avg across common EEG channels, n={n_pids} patients)",
        fontsize=11, fontweight='bold'
    )
    ax_snr.legend(fontsize=8)
    ax_snr.grid(axis='y', alpha=0.3)
    fig_snr.tight_layout()
    st.pyplot(fig_snr, use_container_width=True)
    plt.close(fig_snr)

    # ── Summary table ────────────────────────────────────────────────────
    table_data = {'Patient': all_pids}
    for method in selected_methods:
        table_data[f'{method} SNR (dB)'] = [round(results[method].get(p, 0.0), 2) for p in all_pids]
    st.dataframe(pd.DataFrame(table_data).set_index('Patient'), use_container_width=True)


def handle_hep_t_peak_windows(filtered_individuals, selected_group, selected_stage):
    """
    Show sample HEP windows (patient averages) aligned to ECG, marking the
    ECG T-peak and the corresponding EEG T-wave peak per channel.
    Also shows per-channel and group averages.
    """
    st.subheader("HEP T-peak Windows — EEG aligned to ECG T-wave")
    st.caption(
        "Each panel shows the averaged HEP window for a sample of patients. "
        "A vertical dashed line marks the ECG T-peak (maximum of the averaged ECG HEP "
        "between 150–500 ms post R-peak). Dots mark the dominant EEG peak nearest "
        "to the ECG T-peak for each channel. The bottom panel shows the per-channel "
        "and group-average EEG HEP with peak markers."
    )

    if not filtered_individuals:
        st.warning("No individuals loaded.")
        return

    # Controls
    n_samples = st.slider(
        "Number of sample patients to show", min_value=1,
        max_value=min(10, len(filtered_individuals)), value=min(4, len(filtered_individuals)),
        key="hep_tpeak_n_samples"
    )
    t_search_min = st.number_input("T-peak search window — start (ms post R-peak)",
                                   min_value=50, max_value=400, value=150, step=10,
                                   key="hep_tpeak_tmin") / 1000.0
    t_search_max = st.number_input("T-peak search window — end (ms post R-peak)",
                                   min_value=100, max_value=700, value=500, step=10,
                                   key="hep_tpeak_tmax") / 1000.0
    eeg_peak_half_win = st.number_input(
        "EEG peak search half-window around ECG T-peak (ms)",
        min_value=10, max_value=150, value=60, step=10,
        key="hep_tpeak_eeg_hw"
    ) / 1000.0

    def find_t_peak_time(ecg_hep, times, t_min, t_max):
        mask = (times >= t_min) & (times <= t_max)
        if not np.any(mask):
            return None
        ecg_1d = np.asarray(ecg_hep).squeeze()
        if ecg_1d.ndim != 1 or len(ecg_1d) != len(times):
            return None
        sub = ecg_1d[mask]
        idx_local = np.argmax(np.abs(sub))
        return float(times[mask][idx_local])

    def find_eeg_peak_near(hep_ch, times, t_center, half_win):
        # EEG T-wave peak must occur BEFORE the ECG T-wave peak, not after
        mask = (times >= t_center - half_win) & (times < t_center)
        if not np.any(mask):
            return None, None
        sub = hep_ch[mask]
        idx_local = np.argmax(np.abs(sub))
        t_pk = float(times[mask][idx_local])
        amp = float(sub[idx_local]) * 1e6
        return t_pk, amp

    common_channels = identify_common_eeg_channels(filtered_individuals)
    if not common_channels:
        st.warning("No common EEG channels across individuals.")
        return

    sample_inds = filtered_individuals[:n_samples]

    st.markdown("### Sample Patient HEP Windows")
    for ind in sample_inds:
        pid = ind[0]
        hep_data = ind[1]
        times = np.asarray(ind[2])
        ch_names = ind[3]
        ecg_hep = ind[5] if len(ind) > 5 else None
        ecg_ch = ind[6] if len(ind) > 6 else []
        flip_metadata_available = len(ind) > 8
        flipped_channels = set(ind[8]) if flip_metadata_available and ind[8] else set()

        if hep_data is None or times is None:
            continue

        t_peak = None
        if ecg_hep is not None:
            t_peak = find_t_peak_time(ecg_hep, times, t_search_min, t_search_max)

        plot_chs = [ch for ch in common_channels if ch in ch_names]
        ecg_1d_check = np.asarray(ecg_hep).squeeze() if ecg_hep is not None else None
        has_ecg = (ecg_1d_check is not None and ecg_1d_check.ndim == 1
                   and len(ecg_1d_check) == len(times))
        n_rows = len(plot_chs) + (1 if has_ecg else 0)

        if n_rows == 0:
            continue

        with st.expander(f"Patient: {pid}", expanded=(pid == sample_inds[0][0])):
            fig, axes = plt.subplots(n_rows, 1, figsize=(14, 2.2 * n_rows), sharex=True)
            if n_rows == 1:
                axes = [axes]
            flip_note = (
                f" | {len(flipped_channels)}/{len(plot_chs)} EEG ch flipped"
                if flip_metadata_available
                else " | flip status unavailable"
            )
            fig.suptitle(
                f"HEP Window — {pid} ({selected_group} / {selected_stage}){flip_note}",
                fontsize=12, fontweight='bold'
            )
            times_ms = times * 1000
            ax_idx = 0

            if has_ecg:
                ax = axes[ax_idx]
                ecg_1d = ecg_1d_check * 1e6
                ax.plot(times_ms, ecg_1d, color='crimson', linewidth=1.5, label='ECG HEP')
                if t_peak is not None:
                    ax.axvline(t_peak * 1000, color='navy', linestyle='--', linewidth=1.2,
                               label=f'T-peak ({t_peak * 1000:.0f} ms)')
                    t_amp = float(np.interp(t_peak, times, ecg_1d_check)) * 1e6
                    ax.scatter([t_peak * 1000], [t_amp], color='navy', s=60, zorder=5)
                ax.set_ylabel(f"{ecg_ch[0] if ecg_ch else 'ECG'} (μV)")
                ax.legend(loc='upper right', fontsize=7)
                ax.axhline(0, color='gray', linewidth=0.5)
                ax.grid(True, alpha=0.3)
                ax_idx += 1

            flip_status_rows = []
            for ch in plot_chs:
                if ch not in ch_names:
                    continue
                ch_idx = ch_names.index(ch)
                hep_ch = hep_data[ch_idx]
                ax = axes[ax_idx]
                was_flipped = ch in flipped_channels
                flip_status = (
                    "flipped" if was_flipped
                    else "not flipped" if flip_metadata_available
                    else "unknown"
                )
                line_color = 'darkorange' if was_flipped else 'steelblue'
                ax.plot(times_ms, hep_ch * 1e6, color=line_color, linewidth=1.2, label=f"{ch} ({flip_status})")
                if t_peak is not None:
                    ax.axvline(t_peak * 1000, color='navy', linestyle='--', linewidth=0.8, alpha=0.6)
                    ep_t, ep_amp = find_eeg_peak_near(hep_ch, times, t_peak, eeg_peak_half_win)
                    if ep_t is not None:
                        ax.scatter([ep_t * 1000], [ep_amp], color='darkorange', s=60, zorder=5,
                                   label=f'EEG peak ({ep_t * 1000:.0f} ms, {ep_amp:.1f} μV)')
                ax.set_ylabel(f"{ch} [{flip_status}] (μV)")
                ax.legend(loc='upper right', fontsize=7)
                ax.axhline(0, color='gray', linewidth=0.5)
                ax.grid(True, alpha=0.3)
                ax_idx += 1
                flip_status_rows.append({
                    "Channel": ch,
                    "Flip status": flip_status,
                })

            axes[-1].set_xlabel("Time post R-peak (ms)")
            fig.tight_layout()
            st.pyplot(fig, use_container_width=True)
            plt.close(fig)
            if flip_status_rows:
                st.dataframe(pd.DataFrame(flip_status_rows), hide_index=True, use_container_width=True)

    st.markdown("### Per-Channel Group Average with T-peak Markers")

    all_t_peaks = []
    for ind in filtered_individuals:
        ecg_hep = ind[5] if len(ind) > 5 else None
        times = np.asarray(ind[2])
        if ecg_hep is not None:
            tp = find_t_peak_time(ecg_hep, times, t_search_min, t_search_max)
            if tp is not None:
                all_t_peaks.append(tp)
    group_t_peak = float(np.mean(all_t_peaks)) if all_t_peaks else None

    if group_t_peak is not None:
        st.info(
            f"Group mean ECG T-peak: **{group_t_peak * 1000:.1f} ms** post R-peak "
            f"(n={len(all_t_peaks)} individuals with ECG data)"
        )

    ref_times = None
    for ind in filtered_individuals:
        if ind[2] is not None:
            ref_times = np.asarray(ind[2])
            break

    if ref_times is None:
        st.warning("No valid time axis found.")
        return

    times_ms_ref = ref_times * 1000

    ch_group_heps = {}
    for ch in common_channels:
        ch_vals = []
        for ind in filtered_individuals:
            hd = ind[1]
            cns = ind[3]
            t = np.asarray(ind[2])
            if hd is None or ch not in cns:
                continue
            c_idx = cns.index(ch)
            if len(t) == hd.shape[1]:
                ch_vals.append(hd[c_idx])
        if ch_vals:
            ch_group_heps[ch] = np.nanmean(ch_vals, axis=0)

    if not ch_group_heps:
        st.warning("Could not compute group channel averages.")
        return

    group_avg = np.nanmean(list(ch_group_heps.values()), axis=0)
    n_chs = len(ch_group_heps)
    n_cols = 3
    n_rows_grid = (n_chs + n_cols - 1) // n_cols + 1

    fig, axes = plt.subplots(n_rows_grid, n_cols, figsize=(16, 2.8 * n_rows_grid))
    axes_flat = axes.flatten()

    for ax_i, (ch, avg_hep) in enumerate(ch_group_heps.items()):
        ax = axes_flat[ax_i]
        ax.plot(times_ms_ref, avg_hep * 1e6, color='steelblue', linewidth=1.5)
        ch_delta_str = ''
        if group_t_peak is not None:
            ax.axvline(group_t_peak * 1000, color='navy', linestyle='--', linewidth=1.0,
                       label=f'T-peak ({group_t_peak * 1000:.0f} ms)')
            ep_t, ep_amp = find_eeg_peak_near(avg_hep, ref_times, group_t_peak, eeg_peak_half_win)
            if ep_t is not None:
                ax.scatter([ep_t * 1000], [ep_amp], color='darkorange', s=60, zorder=5)
                ax.annotate(f'{ep_amp:.1f}μV', xy=(ep_t * 1000, ep_amp),
                            xytext=(5, 5), textcoords='offset points', fontsize=7, color='darkorange')
                delta_ms = (ep_t - group_t_peak) * 1000
                ch_delta_str = f'  Δ{delta_ms:+.0f} ms'
        ax.axhline(0, color='gray', linewidth=0.5)
        ax.set_title(f'{ch}{ch_delta_str}', fontsize=9, fontweight='bold')
        ax.set_xlabel('ms', fontsize=7)
        ax.set_ylabel('μV', fontsize=7)
        ax.tick_params(labelsize=7)
        ax.grid(True, alpha=0.3)

    if n_chs < len(axes_flat):
        ax = axes_flat[n_chs]
        ax.plot(times_ms_ref, group_avg * 1e6, color='darkgreen', linewidth=2.0,
                label='Group avg (all ch)')
        if group_t_peak is not None:
            ax.axvline(group_t_peak * 1000, color='navy', linestyle='--', linewidth=1.2,
                       label=f'T-peak ({group_t_peak * 1000:.0f} ms)')
            ep_t, ep_amp = find_eeg_peak_near(group_avg, ref_times, group_t_peak, eeg_peak_half_win)
            if ep_t is not None:
                ax.scatter([ep_t * 1000], [ep_amp], color='darkorange', s=80, zorder=5,
                           label=f'EEG T-peak ({ep_t * 1000:.0f} ms)')
        ax.axhline(0, color='gray', linewidth=0.5)
        ax.set_title('Group Average (all channels)', fontsize=9, fontweight='bold', color='darkgreen')
        ax.set_xlabel('ms', fontsize=7)
        ax.set_ylabel('μV', fontsize=7)
        ax.tick_params(labelsize=7)
        ax.legend(fontsize=7)
        ax.grid(True, alpha=0.3)

    for ax_i in range(n_chs + 1, len(axes_flat)):
        axes_flat[ax_i].set_visible(False)

    fig.suptitle(
        f"Group HEP — Per-Channel Averages with T-peak Markers\n"
        f"{selected_group} / {selected_stage} (n={len(filtered_individuals)} patients)",
        fontsize=12, fontweight='bold'
    )
    fig.tight_layout()
    st.pyplot(fig, use_container_width=True)
    plt.close(fig)

    # ── T-peak timing delta statistics (boxplot) ──────────────────
    st.markdown("### T-wave Peak Timing Δ — EEG vs ECG (per channel)")
    st.caption(
        "For each patient, the delta = EEG T-wave peak time − ECG T-peak time (ms). "
        "Positive = EEG lags ECG; negative = EEG leads. "
        "P-values: Wilcoxon signed-rank test vs zero (H₀: median Δ = 0). "
        "Orange line = median; box = IQR; whiskers = 1.5×IQR; dots = outliers."
    )

    from scipy.stats import wilcoxon, median_abs_deviation

    ch_deltas = {}
    for ch in common_channels:
        deltas = []
        for ind in filtered_individuals:
            hd = ind[1]
            cns = ind[3]
            ecg_hep = ind[5] if len(ind) > 5 else None
            t = np.asarray(ind[2])
            if hd is None or ch not in cns or ecg_hep is None:
                continue
            ecg_tp = find_t_peak_time(ecg_hep, t, t_search_min, t_search_max)
            if ecg_tp is None:
                continue
            c_idx = cns.index(ch)
            eeg_tp, _ = find_eeg_peak_near(hd[c_idx], t, ecg_tp, eeg_peak_half_win)
            if eeg_tp is not None:
                deltas.append((eeg_tp - ecg_tp) * 1000)
        if len(deltas) >= 2:
            ch_deltas[ch] = np.array(deltas)

    if not ch_deltas:
        st.info("Not enough data to compute per-channel timing statistics.")
    else:
        channels_sorted = list(ch_deltas.keys())
        data_list = [ch_deltas[ch] for ch in channels_sorted]

        fig_box, ax_box = plt.subplots(figsize=(max(8, len(channels_sorted) * 0.7 + 2), 5))

        ax_box.boxplot(
            data_list,
            labels=channels_sorted,
            patch_artist=True,
            medianprops=dict(color='darkorange', linewidth=2),
            boxprops=dict(facecolor='steelblue', alpha=0.4),
            whiskerprops=dict(color='steelblue'),
            capprops=dict(color='steelblue'),
            flierprops=dict(marker='o', color='gray', alpha=0.5, markersize=4),
        )

        ax_box.axhline(0, color='crimson', linewidth=1.2, linestyle='--', label='Zero (ECG T-peak)')
        ax_box.set_xlabel('EEG Channel', fontsize=10)
        ax_box.set_ylabel('Δ time (ms)  [EEG − ECG T-peak]', fontsize=10)
        ax_box.tick_params(axis='x', rotation=45, labelsize=8)
        ax_box.grid(True, axis='y', alpha=0.3)

        for i, ch in enumerate(channels_sorted):
            d = ch_deltas[ch]
            med = float(np.median(d))
            mad = float(median_abs_deviation(d))
            try:
                _, pval = wilcoxon(d)
            except Exception:
                pval = float('nan')
            if np.isnan(pval):
                stars = 'n/a'
            elif pval < 0.001:
                stars = '***'
            elif pval < 0.01:
                stars = '**'
            elif pval < 0.05:
                stars = '*'
            else:
                stars = 'ns'
            annotation = f'med={med:+.1f}\nMAD={mad:.1f}\np={pval:.3f} {stars}'
            ax_box.annotate(
                annotation,
                xy=(i + 1, med),
                xytext=(0, 18),
                textcoords='offset points',
                ha='center', va='bottom',
                fontsize=6,
                color='navy',
                bbox=dict(boxstyle='round,pad=0.2', fc='white', alpha=0.7, ec='none'),
            )

        ax_box.set_title(
            f"EEG T-wave vs ECG T-peak Timing Δ per Channel\n"
            f"{selected_group} / {selected_stage}  (n patients per channel shown in box)",
            fontsize=11, fontweight='bold'
        )
        ax_box.legend(fontsize=8)
        fig_box.tight_layout()
        st.pyplot(fig_box, use_container_width=True)
        plt.close(fig_box)

        rows = []
        for ch in channels_sorted:
            d = ch_deltas[ch]
            med = float(np.median(d))
            mad = float(median_abs_deviation(d))
            try:
                _, pval = wilcoxon(d)
            except Exception:
                pval = float('nan')
            stars = ('***' if pval < 0.001 else '**' if pval < 0.01
                     else '*' if pval < 0.05 else 'ns') if not np.isnan(pval) else 'n/a'
            rows.append({
                'Channel': ch,
                'n': len(d),
                'Median Δ (ms)': round(med, 2),
                'MAD (ms)': round(mad, 2),
                'p-value': round(pval, 4) if not np.isnan(pval) else None,
                'Significance': stars,
            })
        st.dataframe(pd.DataFrame(rows).set_index('Channel'), use_container_width=True)


def handle_ecg_free_ica(filtered_individuals, selected_group, selected_stage):
    """
    Automated ECG-free ICA component identification pipeline.
    Uses IC features (kurtosis, variance explained, spectral flatness) without ECG reference.
    Reference: ECG-free automatic cardiac IC classification using explicit feature flowcharts.
    """
    try:
        from sklearn.decomposition import FastICA
        FASTICA_AVAILABLE = True
    except ImportError:
        FASTICA_AVAILABLE = False

    st.subheader("Automated ECG-Free ICA Component Identification")
    st.caption(
        "Identifies cardiac ICA components purely from their features — no ECG reference required. "
        "Features: kurtosis (supra-Gaussian spikes → high kurtosis), variance explained, "
        "and spectral flatness (cardiac ICs have low spectral flatness / tonal structure). "
        "Top-K cardiac components are removed; the cleaned HEP is shown."
    )

    if not FASTICA_AVAILABLE:
        st.warning("scikit-learn is not installed; falling back to PCA-based decomposition.")

    n_components = st.slider(
        "Number of ICA/PCA components", min_value=2, max_value=20, value=5,
        key="ecg_free_ica_ncomp"
    )
    top_k = st.slider(
        "Top-K cardiac components to remove", min_value=1, max_value=n_components, value=1,
        key="ecg_free_ica_topk"
    )

    # Pipeline description
    st.info(
        "**Automated ECG-Free IC Classification Flowchart**\n\n"
        "1. Stack patient HEP matrices → shape (channels × time-points)\n"
        "2. Apply FastICA (or PCA fallback) → extract N independent components\n"
        "3. For each IC compute:\n"
        "   - **Kurtosis**: measures peakedness; cardiac ICs tend to have high kurtosis\n"
        "   - **Variance explained**: proportion of total signal variance\n"
        "   - **Spectral flatness**: low value → tonal/repetitive structure (cardiac-like)\n"
        "4. Rank ICs by composite score = kurtosis × (1 − spectral flatness)\n"
        "5. Remove top-K ranked ICs from the signal\n"
        "6. Reconstruct cleaned HEP and display"
    )

    # Collect common channels
    common_channels = identify_common_eeg_channels(filtered_individuals)

    if not common_channels:
        st.warning("No common EEG channels found.")
        return

    # Build data matrix: shape (n_subjects * n_channels, n_times)
    hep_rows = []
    row_labels = []
    times = None

    for ind in filtered_individuals:
        pid, hep_full, t, ch_names = ind[:4]
        if hep_full is None:
            continue
        if times is None:
            times = t
        for ch in common_channels:
            if ch not in ch_names:
                continue
            idx = ch_names.index(ch)
            hep_rows.append(hep_full[idx].copy())
            row_labels.append((pid, ch))

    if not hep_rows or times is None:
        st.warning("No valid HEP data found.")
        return

    X = np.array(hep_rows)  # shape: (n_rows, n_times)
    n_comp_actual = min(n_components, X.shape[0], X.shape[1])

    # Decomposition
    if FASTICA_AVAILABLE:
        try:
            ica_model = FastICA(n_components=n_comp_actual, random_state=42, max_iter=500)
            sources = ica_model.fit_transform(X)  # (n_rows, n_comp_actual)
            mixing = ica_model.mixing_            # (n_times, n_comp_actual)
            method_label = "FastICA"
        except Exception:
            FASTICA_AVAILABLE = False

    if not FASTICA_AVAILABLE:
        U, s, Vt = np.linalg.svd(X - X.mean(axis=0), full_matrices=False)
        sources = U[:, :n_comp_actual] * s[:n_comp_actual]
        mixing = Vt[:n_comp_actual, :].T
        method_label = "PCA (fallback)"

    st.markdown(f"*Decomposition method: **{method_label}**, {n_comp_actual} components*")

    # Compute IC features
    def spectral_flatness(signal):
        """Wiener entropy: geometric mean / arithmetic mean of power spectrum."""
        f_power = np.abs(np.fft.rfft(signal)) ** 2 + 1e-30
        log_mean = np.mean(np.log(f_power))
        arith_mean = np.mean(f_power)
        return float(np.exp(log_mean) / arith_mean)

    from scipy.stats import kurtosis as scipy_kurtosis

    ic_features = []
    for k in range(n_comp_actual):
        ic = sources[:, k]
        kurt = float(scipy_kurtosis(ic, fisher=True))
        var_exp = float(np.var(ic) / (np.var(X) + 1e-30))
        sf = spectral_flatness(ic)
        score = abs(kurt) * (1.0 - sf)
        ic_features.append({'ic': k, 'kurtosis': kurt, 'var_exp': var_exp, 'spectral_flatness': sf, 'score': score})

    ic_features_sorted = sorted(ic_features, key=lambda d: d['score'], reverse=True)
    cardiac_ic_indices = [d['ic'] for d in ic_features_sorted[:top_k]]

    # Feature summary table
    st.markdown("#### IC Feature Summary")
    import pandas as pd
    feat_df = pd.DataFrame(ic_features_sorted).rename(columns={
        'ic': 'IC', 'kurtosis': 'Kurtosis', 'var_exp': 'Var Explained',
        'spectral_flatness': 'Spectral Flatness', 'score': 'Cardiac Score'
    })
    feat_df['Cardiac?'] = feat_df['IC'].apply(lambda i: '*** Removed' if i in cardiac_ic_indices else '')
    st.dataframe(feat_df.style.format({
        'Kurtosis': '{:.3f}', 'Var Explained': '{:.4f}',
        'Spectral Flatness': '{:.4f}', 'Cardiac Score': '{:.4f}'
    }), use_container_width=True)

    # IC waveform grid
    st.markdown("#### IC Waveforms (group-level projection)")
    n_cols = min(4, n_comp_actual)
    n_rows_grid = int(np.ceil(n_comp_actual / n_cols))
    fig_ics, axes_ics = plt.subplots(n_rows_grid, n_cols, figsize=(4 * n_cols, 3 * n_rows_grid))
    axes_ics = np.array(axes_ics).reshape(-1)

    ic_times = np.arange(n_comp_actual)
    for k in range(n_comp_actual):
        ax = axes_ics[k]
        ic_waveform = sources[:, k]
        ax.plot(ic_waveform, color='crimson' if k in cardiac_ic_indices else 'steelblue',
                linewidth=0.8, alpha=0.7)
        kurt_val = ic_features[k]['kurtosis']
        sf_val = ic_features[k]['spectral_flatness']
        label = f"IC{k+1}\nKurt={kurt_val:.2f}\nSF={sf_val:.3f}"
        if k in cardiac_ic_indices:
            label += "\n[CARDIAC]"
            ax.set_facecolor('#fff0f0')
        ax.set_title(label, fontsize=7)
        ax.axhline(0, color='gray', linewidth=0.5)
        ax.tick_params(labelsize=6)

    for k in range(n_comp_actual, len(axes_ics)):
        axes_ics[k].set_visible(False)

    fig_ics.suptitle(
        f"IC Waveforms — {method_label} | {selected_group} / {selected_stage}\n"
        f"Red = Identified Cardiac ICs (top {top_k})",
        fontsize=10, fontweight='bold'
    )
    fig_ics.tight_layout()
    st.pyplot(fig_ics, use_container_width=True)
    plt.close(fig_ics)

    # Reconstruct cleaned signal
    st.markdown("#### Cleaned HEP per Channel (cardiac ICs removed)")
    mask = np.ones(n_comp_actual, dtype=bool)
    for k in cardiac_ic_indices:
        mask[k] = False

    if FASTICA_AVAILABLE and method_label == "FastICA":
        sources_clean = sources.copy()
        for k in cardiac_ic_indices:
            sources_clean[:, k] = 0.0
        X_clean = sources_clean @ ica_model.mixing_.T + ica_model.mean_
    else:
        sources_clean = sources.copy()
        for k in cardiac_ic_indices:
            sources_clean[:, k] = 0.0
        X_clean = sources_clean @ mixing.T + X.mean(axis=0)

    # Group by channel and plot
    ch_orig = {ch: [] for ch in common_channels}
    ch_clean = {ch: [] for ch in common_channels}

    for row_idx, (pid, ch) in enumerate(row_labels):
        if ch in ch_orig:
            ch_orig[ch].append(X[row_idx])
            ch_clean[ch].append(X_clean[row_idx])

    display_channels = ['Average'] + common_channels

    with st.expander("Average across all channels", expanded=True):
        all_orig = np.array([v for vals in ch_orig.values() for v in vals]) * 1e6
        all_clean = np.array([v for vals in ch_clean.values() for v in vals]) * 1e6
        if len(all_orig):
            fig_avg, ax_avg = plt.subplots(figsize=(12, 5))
            ax_avg.plot(times, all_orig.mean(axis=0), color='steelblue', linewidth=2.5,
                        label='Original avg')
            ax_avg.plot(times, all_clean.mean(axis=0), color='darkorange', linewidth=2.5,
                        linestyle='--', label='Cleaned avg (cardiac ICs removed)')
            ax_avg.axvline(0, color='gray', linewidth=0.8, linestyle='--')
            ax_avg.axhline(0, color='gray', linewidth=0.5)
            ax_avg.set_xlabel("Time (s)")
            ax_avg.set_ylabel("Amplitude (µV)")
            ax_avg.legend(fontsize=9)
            ax_avg.grid(alpha=0.2)
            ax_avg.set_title(
                f"ECG-Free ICA Cleaned HEP — Average | {selected_group} / {selected_stage}",
                fontsize=10, fontweight='bold'
            )
            fig_avg.tight_layout()
            st.pyplot(fig_avg, use_container_width=True)
            plt.close(fig_avg)

    for ch_name in common_channels:
        orig_list = ch_orig[ch_name]
        clean_list = ch_clean[ch_name]
        n = min(len(orig_list), len(clean_list))
        if n == 0:
            continue
        orig_arr = np.array(orig_list[:n]) * 1e6
        clean_arr = np.array(clean_list[:n]) * 1e6

        with st.expander(f"Channel: {ch_name}  (n={n})", expanded=False):
            fig, axes = plt.subplots(1, 2, figsize=(14, 5))
            cmap = plt.get_cmap('tab20' if n <= 20 else 'hsv', max(n, 1))
            colors = [cmap(i / max(n - 1, 1)) for i in range(n)]

            for ax, arr, title in zip(
                axes,
                [orig_arr, clean_arr],
                ["Original HEP", "Cleaned HEP (ECG-Free ICA)"]
            ):
                for i, trace in enumerate(arr):
                    ax.plot(times, trace, color=colors[i], alpha=0.35, linewidth=0.8)
                ax.plot(times, arr.mean(axis=0), color='black', linewidth=2, label='Group avg')
                ax.set_title(title, fontsize=10)
                ax.set_xlabel("Time (s)")
                ax.set_ylabel("Amplitude (µV)")
                ax.axvline(0, color='gray', linewidth=0.8, linestyle='--')
                ax.axhline(0, color='gray', linewidth=0.5)
                ax.grid(alpha=0.2)

            fig.suptitle(
                f"ECG-Free ICA — Channel: {ch_name} | {selected_group} / {selected_stage}",
                fontsize=11, fontweight='bold'
            )
            fig.tight_layout()
            st.pyplot(fig, use_container_width=True)
            plt.close(fig)


def run_single_group_analysis(base_path, selected_stage):
    """
    Logic for Single Group Analysis mode.
    """
    # Get available groups
    available_groups = [g for g in os.listdir(base_path) if os.path.isdir(os.path.join(base_path, g))]
    selected_group = st.selectbox("Select Group", available_groups, index=1)
    
    st.write(f"Processing individuals for {selected_group}...")
    if st.runtime.exists():
        col1, col2 = st.columns(2)
        with col1:
            test_run = st.checkbox("Test Run (first 5 files only)", value=False, key="test_run_single")
        with col2:
            recompute_cache = st.button("Recompute Cache", key="recompute_cache_single")
    else:
        test_run = True
        recompute_cache = False

    if recompute_cache:
        if hasattr(get_group_individuals, "clear"):
            get_group_individuals.clear()
    individuals = get_group_individuals(selected_group, selected_stage, base_path, test_run=test_run, recompute_cache=recompute_cache)
    
    # Collect globally skipped or removed logs
    all_logs = []
    if individuals:
        for ind in individuals:
            if len(ind) > 7 and ind[7] is not None:
                log_msg = ind[7]
                all_logs.append((ind[0], log_msg))

    repetitive_excluded_pids = {
        pid
        for pid, log_msg in all_logs
        if log_msg.get('perc', 0.0) > REPETITIVE_RPEAK_EXCLUDE_PERC
    }

    if all_logs:
        with st.expander(f"Repetitive R-Peak Artifact Details ({len(all_logs)} patients flagged)"):
            for pid, l in all_logs:
                st.markdown(f"**Patient: {pid}**")
                st.write(f"- Total initial R-peaks: {l['total']}")
                st.write(f"- Total removed: {l['removed']} ({l['perc']:.1f}%)")
                if l['info']:
                    for info in l['info']:
                        st.write(f"  - {info}")
                if pid in repetitive_excluded_pids:
                    st.warning(
                        f"Removed from individuals because more than {REPETITIVE_RPEAK_EXCLUDE_PERC:.0f}% "
                        "of R-peaks were marked as repetitive artifacts."
                    )
                st.markdown("---")

    if repetitive_excluded_pids and individuals:
        individuals = [ind for ind in individuals if ind[0] not in repetitive_excluded_pids]
        if not individuals:
            st.warning(
                "All loaded patients were removed because more than "
                f"{REPETITIVE_RPEAK_EXCLUDE_PERC:.0f}% of their R-peaks were repetitive artifacts."
            )
            return
    
    # make st number_input to ger from user jitter_sec defult .1
    jitter_sec = st.number_input("Jitter (seconds)", min_value=0.01, max_value=.5, value=0.1, step=0.05)
    # checkbox 'show ecg only plots'
    show_ecg_only = st.checkbox("Show ECG Only Plots", value=False)
    show_noise_analysis = st.checkbox("Show ECG Repetitive Noise Analysis", value=False)
    show_single_patient_all = st.checkbox("Show Single Patient All Channels", value=False)
    show_patients_comparison = st.checkbox("Show EEG-ECG Patients Comparison", value=False)
    show_group_ecg = st.checkbox("Show Group ECG Analysis", value=False)
    show_bhi = st.checkbox("Show Brain-Heart Interplay (BHI)", value=False)
    show_ica = st.checkbox("Show ICA ECG Cleaning of HEP", value=False)
    show_ecg_reduction = st.checkbox("Show ECG Signal Reduction (normalized subtraction)", value=False)
    show_nc_hep_cleaning = st.checkbox(
        "NC_HEP_CLEANING",
        value=False,
        help="Align the ECG HEP peak to the global EEG HEP peak, fit the ECG artifact amplitude per EEG channel, and subtract it.",
    )
    show_ecg_free_ica = st.checkbox("Show Automated ECG-Free ICA Component Identification", value=False)
    show_hep_t_peak_windows = st.checkbox("Show HEP T-peak Windows (EEG aligned to ECG T-wave)", value=True)
    show_snr_optimization = st.checkbox("Show SNR Optimization (compare cleaning methods)", value=False,
                                        help="Compare multiple signal cleaning methods and their effect on SNR.")
    show_fft_psd = st.checkbox("Show FFT / PSD Analysis", value=False,
                               help="Plot FFT amplitude spectrum and power spectral density for a selected patient.")
    if show_ica or show_snr_optimization:
        with st.expander("SNR Filter Range (post-cleaning)", expanded=False):
            _snr_col1, _snr_col2 = st.columns(2)
            with _snr_col1:
                snr_min_val = st.number_input("Min SNR (dB)", min_value=-20.0, max_value=20.0, value=1.5, step=0.5,
                                              key="single_snr_min",
                                              help="Only show patients with post-cleaning SNR ≥ this value.")
            with _snr_col2:
                snr_max_val = st.number_input("Max SNR (dB)", min_value=-20.0, max_value=50.0, value=6.0, step=0.5,
                                              key="single_snr_max",
                                              help="Only show patients with post-cleaning SNR ≤ this value.")
    else:
        snr_min_val, snr_max_val = 1.5, 6.0

    if individuals:
        if show_noise_analysis:
            handle_ecg_noise_detection(base_path, selected_group, selected_stage)

        if show_single_patient_all:
            handle_single_patient_view(
                individuals,
                selected_group,
                selected_stage,
                base_path,
                n_permutations=100,
                p_threshold=0.05,
                jitter_sec=jitter_sec,
            )

        if show_patients_comparison:
            n_compare = st.slider("Number of patients to compare", min_value=1, max_value=len(individuals), value=min(4, len(individuals)))
            plot_patients_butterfly_comparison(individuals[:n_compare], selected_group, selected_stage)

        if show_ecg_only:
            plot_ecg_hep_individuals(individuals, selected_group, selected_stage)

        if show_fft_psd:
            st.subheader("FFT / PSD Analysis")
            from scipy.fft import rfft, rfftfreq

            all_pids_fft = [ind[0] for ind in individuals]
            selected_pid_fft = st.selectbox("Select Patient", all_pids_fft, key="fft_psd_patient")
            ind_fft = next(i for i in individuals if i[0] == selected_pid_fft)

            pid_f, hep_data_f, times_f, ch_names_f = ind_fft[0], ind_fft[1], ind_fft[2], ind_fft[3]
            ecg_hep_f = ind_fft[5] if len(ind_fft) > 5 else None
            ecg_ch_f = ind_fft[6] if len(ind_fft) > 6 else []
            sfreq_f = 1.0 / np.mean(np.diff(times_f))

            all_channels_fft = list(ch_names_f) + (list(ecg_ch_f) if ecg_ch_f else [])
            selected_channels_fft = st.multiselect(
                "Select Channels", options=all_channels_fft,
                default=all_channels_fft[:min(4, len(all_channels_fft))],
                key="fft_psd_channels",
            )

            fft_col1, fft_col2 = st.columns(2)
            with fft_col1:
                freq_min = st.number_input("Min Frequency (Hz)", min_value=0.0, max_value=150.0,
                                           value=0.5, step=0.5, key="fft_psd_fmin")
            with fft_col2:
                freq_max = st.number_input("Max Frequency (Hz)", min_value=1.0, max_value=sfreq_f / 2,
                                           value=min(40.0, sfreq_f / 2), step=1.0, key="fft_psd_fmax")

            use_log_scale = st.checkbox("Log Y-axis", value=True, key="fft_psd_log")

            if selected_channels_fft:
                n_ch_fft = len(selected_channels_fft)
                fig_fft, axes_fft = plt.subplots(n_ch_fft, 2, figsize=(14, 3 * n_ch_fft), squeeze=False)
                fig_fft.suptitle(f"FFT & PSD — Patient {selected_pid_fft} ({selected_group} / {selected_stage})",
                                 fontsize=13, fontweight='bold')

                for ch_i, ch in enumerate(selected_channels_fft):
                    # Resolve channel data
                    if ch in ch_names_f:
                        ch_idx = list(ch_names_f).index(ch)
                        signal = hep_data_f[ch_idx]
                    elif ecg_hep_f is not None and ecg_ch_f and ch in ecg_ch_f:
                        ch_idx = list(ecg_ch_f).index(ch)
                        signal = ecg_hep_f[ch_idx]
                    else:
                        continue

                    n_samp = len(signal)
                    freqs_f = rfftfreq(n_samp, d=1.0 / sfreq_f)
                    mask = (freqs_f >= freq_min) & (freqs_f <= freq_max)

                    # FFT amplitude
                    fft_amp = np.abs(rfft(signal)) / n_samp * 2
                    ax_amp = axes_fft[ch_i, 0]
                    ax_amp.plot(freqs_f[mask], fft_amp[mask], color='steelblue', linewidth=1.2)
                    ax_amp.set_title(f"{ch} — FFT Amplitude")
                    ax_amp.set_xlabel("Frequency (Hz)")
                    ax_amp.set_ylabel("Amplitude (μV)")
                    if use_log_scale:
                        ax_amp.set_yscale('log')
                    ax_amp.grid(True, which='both', alpha=0.3)

                    # Power spectrum (|FFT|²)
                    psd = (np.abs(rfft(signal)) ** 2) / (n_samp * sfreq_f)
                    ax_psd = axes_fft[ch_i, 1]
                    ax_psd.plot(freqs_f[mask], psd[mask], color='darkorange', linewidth=1.2)
                    ax_psd.set_title(f"{ch} — Power Spectrum (μV²/Hz)")
                    ax_psd.set_xlabel("Frequency (Hz)")
                    ax_psd.set_ylabel("Power (μV²/Hz)")
                    if use_log_scale:
                        ax_psd.set_yscale('log')
                    ax_psd.grid(True, which='both', alpha=0.3)

                fig_fft.tight_layout()
                st.pyplot(fig_fft, use_container_width=True)
                plt.close(fig_fft)
            else:
                st.info("Select at least one channel to plot.")

        st.divider()
        st.subheader("Exclude Patients (Global)")
        
        # Get base patient IDs for the current group/stage
        all_pids = [ind[0] for ind in individuals]
        all_base_pids = list(dict.fromkeys([pid.split('_')[0] if '_' in pid else pid for pid in all_pids]))
        
        # Load globally excluded patients from CSV
        csv_path = os.path.join(base_path, "excluded_patients.csv")
        global_excluded = load_excluded_pids(base_path)

        # Filter the loaded exclusions to only default-select those that exist in this group/stage
        default_excluded = [pid for pid in global_excluded if pid in all_base_pids]

        excluded_pids = st.multiselect(
            "Select patient IDs to exclude from all group average analyses (saved across all sleep stages):", 
            options=all_base_pids,
            default=default_excluded,
            key="exclude_pids_multiselect"
        )
        
        # If selection changed relative to what was in the CSV for *these* patients, update the CSV
        if set(excluded_pids) != set(default_excluded):
            # Retain any exclusions from other groups/stages, remove existing exclusions for this group/stage, 
            # and append the newly selected ones.
            other_excluded = [pid for pid in global_excluded if pid not in all_base_pids]
            new_global_excluded = list(set(other_excluded + excluded_pids))
            
            try:
                pd.DataFrame({'patient_id': new_global_excluded}).to_csv(csv_path, index=False)
            except Exception as e:
                st.error(f"Failed to save excluded patients to CSV: {e}")

        # Filter individuals by checking if their base ID is in the excluded list
        filtered_individuals = []
        for ind in individuals:
            pid = str(ind[0])
            base_pid = pid.split('_')[0] if '_' in pid else pid
            if base_pid not in excluded_pids and pid not in excluded_pids:
                filtered_individuals.append(ind)
        
        if not filtered_individuals:
            st.warning("All patients have been excluded. No group average analysis can be performed.")
            return
        st.divider()

        if show_group_ecg:
            plot_group_ecg_analysis(filtered_individuals, selected_group, selected_stage)

        if show_bhi:
            plot_bhi_analysis(filtered_individuals, selected_group, selected_stage, base_path)

        if show_ica:
            handle_ica_ecg_cleaning(filtered_individuals, selected_group, selected_stage, base_path=base_path,
                                    snr_min=snr_min_val, snr_max=snr_max_val)

        if show_ecg_reduction:
            handle_ecg_reduction(filtered_individuals, selected_group, selected_stage)

        if show_nc_hep_cleaning:
            handle_nc_hep_cleaning(filtered_individuals, selected_group, selected_stage)

        if show_ecg_free_ica:
            handle_ecg_free_ica(filtered_individuals, selected_group, selected_stage)

        if show_hep_t_peak_windows:
            handle_hep_t_peak_windows(filtered_individuals, selected_group, selected_stage)

        if show_snr_optimization:
            handle_snr_optimization(filtered_individuals, selected_group, selected_stage,
                                    snr_min=snr_min_val, snr_max=snr_max_val)

        # Identify common channels across all filtered individuals
        common_channels = identify_common_eeg_channels(filtered_individuals)

        if not common_channels:
            st.error("No common EEG channels found across all individuals in this group.")
            return
        st.divider()
        st.title("Group HEP Analysis")

        # --- Plot Per Channel ---
        st.subheader("Per-Channel Analysis")
        channel_p_values = {}
        
        # Add pseudo-channels for Average, Median, ECG, Left, Right, Middle
        display_channels = ['Average', 'Median', 'ECG', 'Left', 'Right', 'Middle'] + common_channels
        
        for ch_name in display_channels:
            with st.expander(f"Channel: {ch_name}", expanded=False):
                fig, ax = plt.subplots(figsize=(16, 9))
                
                # Plot filtered individuals
                all_full_heps = []
                all_full_hep_pids = []

                n_subj = len(filtered_individuals)
                cmap = plt.get_cmap('tab20' if n_subj <= 20 else 'hsv', max(n_subj, 1))
                subj_colors = [cmap(i / max(n_subj - 1, 1)) for i in range(n_subj)]

                for i, ind in enumerate(filtered_individuals):
                    pid, hep_full, times, ch_names, rpeaks, ecg_hep, ecg_ch = ind[:7]
                    subj_color = subj_colors[i]

                    hep_clean = apply_ecg_regression(hep_full, ecg_hep)

                    if ch_name == 'Average' or ch_name == 'Median':
                        valid_ch_indices = [ch_names.index(ch) for ch in common_channels if ch in ch_names]
                        if valid_ch_indices:
                            if ch_name == 'Average':
                                hep = summarize_channels_without_reference_cancellation(
                                    hep_clean[valid_ch_indices, :]
                                )
                            else:  # Median
                                hep = np.nanmedian(hep_clean[valid_ch_indices, :], axis=0)
                            ax.plot(times, hep * 1e6, color=subj_color, alpha=0.4, linewidth=1, label=pid)
                            all_full_heps.append(hep)
                            all_full_hep_pids.append(pid)
                    elif ch_name == 'ECG':
                        if ecg_hep is not None:
                            hep = np.asarray(ecg_hep).squeeze()
                            if hep.ndim == 1 and len(hep) == len(times):
                                ax.plot(times, hep * 1e6, color=subj_color, alpha=0.4, linewidth=1, label=pid)
                                all_full_heps.append(hep)
                                all_full_hep_pids.append(pid)
                    elif ch_name in ('Left', 'Right', 'Middle'):
                        _avail_chs = [ch for ch in common_channels if ch in ch_names]
                        _left_h, _right_h, _mid_h = get_hemisphere_channels(_avail_chs)
                        if ch_name == 'Left':
                            side_chs = _left_h
                        elif ch_name == 'Right':
                            side_chs = _right_h
                        else:  # Middle
                            side_chs = _mid_h
                        side_indices = [ch_names.index(ch) for ch in side_chs]
                        if side_indices:
                            hep = np.nanmean(hep_clean[side_indices, :], axis=0)
                            ax.plot(times, hep * 1e6, color=subj_color, alpha=0.4, linewidth=1, label=pid)
                            all_full_heps.append(hep)
                            all_full_hep_pids.append(pid)
                    else:
                        if ch_name in ch_names:
                            ch_idx = ch_names.index(ch_name)
                            hep = hep_clean[ch_idx]
                            ax.plot(times, hep * 1e6, color=subj_color, alpha=0.4, linewidth=1, label=pid)
                            all_full_heps.append(hep)
                            all_full_hep_pids.append(pid)
                
                avg_hep = None
                mad_hep = None
                sig_windows = None
                min_p = 1.0
                if all_full_heps:
                    _heps_arr = np.array(all_full_heps)
                    avg_hep = np.nanmedian(_heps_arr, axis=0)
                    mad_hep = np.nanmedian(np.abs(_heps_arr - avg_hep), axis=0)
                    sig_windows, _, per_pt_info = permutation_cluster_jitter_test(
                        _heps_arr, times, jitter_sec=jitter_sec
                    )
                    if sig_windows:
                        min_p = min([w['p_value'] for w in sig_windows])
                    n_sig_pt  = per_pt_info.get('n_significant', 0)
                    n_total_pt = len(all_full_heps)
                    fisher_p  = per_pt_info.get('fisher_p', 1.0)

                channel_p_values[ch_name] = min_p

                # Finalize with Median ± MAD — no legend, include per-patient significance summary
                patient_summary = ""
                if all_full_heps:
                    patient_summary = f" | {n_sig_pt}/{n_total_pt} pts sig, Fisher p={fisher_p:.3f}"
                finalize_plot(
                    fig, ax,
                    f"Channel: {ch_name} - Group: {selected_group} - Stage: {selected_stage}{patient_summary}",
                    avg_hep=avg_hep,
                    times=times,
                    n_subjects=len(filtered_individuals),
                    significant_windows=sig_windows,
                    mad_hep=mad_hep,
                    show_legend=False,
                )

                if all_full_heps:
                    df_hep_csv = pd.DataFrame(
                        np.array(all_full_heps) * 1e6,
                        index=all_full_hep_pids,
                        columns=[f"{t:.4f}" for t in times]
                    )
                    df_hep_csv.index.name = "patient_id"
                    csv_bytes = df_hep_csv.to_csv().encode("utf-8")
                    st.download_button(
                        label=f"Download HEP data — {ch_name} (μV)",
                        data=csv_bytes,
                        file_name=f"HEP_{selected_group}_{selected_stage}_{ch_name}.csv",
                        mime="text/csv",
                        key=f"dl_hep_{ch_name}",
                    )

        # ── Pre-compute ICA+CSD p-values for topomap ─────────────────────────
        channel_p_values_csd = {}
        _csd_per_ch_topo = {ch: [] for ch in common_channels}
        _csd_times_topo = None
        _montage_topo = mne.channels.make_standard_montage('standard_1020')
        with st.spinner("Computing ICA+CSD data for topomap…"):
            for _ind_topo in filtered_individuals:
                _pid_topo, _hep_topo, _t_topo, _ch_topo, _, _ecg_topo, _ = _ind_topo[:7]
                _csd_times_topo = _t_topo
                # ICA regression (same as per-channel block above)
                _hep_ica_topo = _hep_topo.copy().astype(float)
                if _ecg_topo is not None:
                    _ev_topo = np.asarray(_ecg_topo).squeeze()
                    if _ev_topo.ndim == 1 and len(_ev_topo) == _hep_topo.shape[1]:
                        _den_topo = np.dot(_ev_topo, _ev_topo)
                        if _den_topo > 1e-20:
                            _eig_topo = np.dot(_hep_ica_topo, _ev_topo) / _den_topo
                            _hep_ica_topo = _hep_ica_topo - _eig_topo[:, None] * _ev_topo[None, :]
                # CSD requires ≥ 4 channels with known positions
                _valid_topo = [c for c in common_channels if c in _ch_topo]
                if len(_valid_topo) < 4:
                    continue
                _idx_topo = [_ch_topo.index(c) for c in _valid_topo]
                _mc_topo = _hep_ica_topo[_idx_topo, :]
                try:
                    _sfreq_topo = 1.0 / (_t_topo[1] - _t_topo[0]) if len(_t_topo) > 1 else 250.0
                    _info_topo = mne.create_info(ch_names=_valid_topo, sfreq=_sfreq_topo, ch_types='eeg')
                    _info_topo.set_montage(_montage_topo, on_missing='ignore')
                    _ev_mne_topo = mne.EvokedArray(_mc_topo, _info_topo, tmin=float(_t_topo[0]), verbose=False)
                    _ev_csd_topo = mne.preprocessing.compute_current_source_density(_ev_mne_topo)
                    for _ic_t, _cn_t in enumerate(_ev_csd_topo.ch_names):
                        if _cn_t in _csd_per_ch_topo:
                            _csd_per_ch_topo[_cn_t].append(_ev_csd_topo.data[_ic_t].copy())
                except Exception:
                    pass
        if _csd_times_topo is not None:
            for _ch_csd in common_channels:
                _traces_csd = _csd_per_ch_topo.get(_ch_csd, [])
                if len(_traces_csd) >= 3:
                    try:
                        _heps_csd_arr = np.array(_traces_csd)
                        _sig_csd_res, _, _ = permutation_cluster_jitter_test(
                            _heps_csd_arr, _csd_times_topo, jitter_sec=jitter_sec
                        )
                        channel_p_values_csd[_ch_csd] = (
                            min(w['p_value'] for w in _sig_csd_res) if _sig_csd_res else 1.0
                        )
                    except Exception:
                        channel_p_values_csd[_ch_csd] = 1.0

        # ── Helper: build and render a p-value topomap ───────────────────────
        def _render_pval_topomap(p_vals, ax_t, title):
            """Plot a p-value topomap into ax_t. Returns the AxesImage or None."""
            _mont = mne.channels.make_standard_montage('standard_1020')
            _mont_upper = [c.upper() for c in _mont.ch_names]
            _pch, _pd = [], []
            for _c, _p in p_vals.items():
                _cu = _c.upper()
                if _cu in _mont_upper:
                    _pch.append(_mont.ch_names[_mont_upper.index(_cu)])
                    _pd.append(_p)
            _std19 = ['Fp1', 'Fp2', 'F7', 'F3', 'Fz', 'F4', 'F8',
                      'C3', 'Cz', 'C4', 'P3', 'Pz', 'P4', 'O1', 'O2']
            _aliases = {'T7': 'T3', 'T8': 'T4', 'P7': 'T5', 'P8': 'T6'}
            for _bc in _std19:
                if not any(_bc.upper() == _x.upper() for _x in _pch):
                    _pch.append(_bc); _pd.append(0.05)
            for _nn, _on in _aliases.items():
                if not any(_c2.upper() in [_nn.upper(), _on.upper()] for _c2 in _pch):
                    _pch.append(_nn); _pd.append(0.05)
            if not _pch:
                return None
            _info2 = mne.create_info(ch_names=_pch, sfreq=250., ch_types='eeg')
            _info2.set_montage(_mont, on_missing='ignore')
            _valid2 = np.array([
                not np.any(np.isnan(_ch['loc'][:3])) and np.any(_ch['loc'][:3] != 0)
                for _ch in _info2['chs']
            ])
            if not np.any(_valid2):
                return None
            _pch_v = [_pch[i] for i in np.where(_valid2)[0]]
            _da = np.clip(np.array(_pd)[_valid2], 0, 0.05)
            _info2 = mne.pick_info(_info2, np.where(_valid2)[0])
            _res = mne.viz.plot_topomap(
                _da, _info2, axes=ax_t,
                cmap='Reds_r', names=_pch_v, vlim=(0, 0.05), extrapolate='head'
            )
            _im = _res[0] if isinstance(_res, tuple) else _res
            _cb = plt.colorbar(_im, ax=ax_t)
            _cb.set_label("p-value")
            ax_t.set_title(title)
            return _im

        # Plot Topomap of P-values
        if channel_p_values or channel_p_values_csd:
            st.divider()
            st.subheader("Significant Channels Topomap (Minimum cluster P-value)")
            topo_col1, topo_col2 = st.columns(2)

            with topo_col1:
                st.markdown("**ICA-cleaned**")
                if channel_p_values:
                    try:
                        fig_topo1, ax_topo1 = plt.subplots(figsize=(6, 5))
                        ok1 = _render_pval_topomap(
                            channel_p_values, ax_topo1,
                            f"ICA: {selected_group}, {selected_stage}"
                        )
                        if ok1 is not None:
                            st.pyplot(fig_topo1, use_container_width=False)
                        else:
                            st.warning("Could not match channels to standard 10-20 montage.")
                        plt.close(fig_topo1)
                    except Exception as e:
                        import traceback as _tb2
                        st.error(
                            f"**Error generating ICA topomap** ({type(e).__name__}): {e}"
                        )
                        with st.expander("Traceback — ICA Topomap Error"):
                            st.code(_tb2.format_exc())

            with topo_col2:
                st.markdown("**ICA + CSD-cleaned**")
                if channel_p_values_csd:
                    try:
                        fig_topo2, ax_topo2 = plt.subplots(figsize=(6, 5))
                        ok2 = _render_pval_topomap(
                            channel_p_values_csd, ax_topo2,
                            f"ICA+CSD: {selected_group}, {selected_stage}"
                        )
                        if ok2 is not None:
                            st.pyplot(fig_topo2, use_container_width=False)
                        else:
                            st.warning("Could not match channels to standard 10-20 montage.")
                        plt.close(fig_topo2)
                    except Exception as e:
                        import traceback as _tb3
                        st.error(
                            f"**Error generating ICA+CSD topomap** ({type(e).__name__}): {e}"
                        )
                        with st.expander("Traceback — ICA+CSD Topomap Error"):
                            st.code(_tb3.format_exc())
                else:
                    st.info("ICA+CSD topomap not available (insufficient data or CSD failed).")
    else:
        st.error(f"No data found for group {selected_group} in stage {selected_stage}")

def run_compare_groups_all_stages_analysis(base_path):
    """
    Logic for Compare Groups All Sleep Stages mode.
    Runs group comparison across all sleep stages ['W', 'light_sleep', 'N3', 'R'],
    showing results in a combined multi-stage plot and individual per-stage plots.
    """
    SLEEP_STAGES = ['W', 'light_sleep', 'N3', 'R']
    STAGE_COLORS = {'W': '#e74c3c', 'light_sleep': '#e67e22', 'N3': '#2980b9', 'R': '#9b59b6'}
    TOPO_JITTER_SEC = 0.1

    available_groups = sorted([g for g in os.listdir(base_path) if os.path.isdir(os.path.join(base_path, g))])
    if not available_groups:
        st.error("No groups found in the data directory.")
        return

    default_groups = default_hep_groups(available_groups)
    selected_groups = st.multiselect(
        "Select Groups to Compare",
        options=available_groups,
        default=default_groups,
        key="cmp_all_stages_selected_groups",
    )
    if not selected_groups:
        st.warning("Please select at least one group.")
        return
    groups_label = " vs ".join(selected_groups) if len(selected_groups) == 2 else ", ".join(selected_groups)

    with st.expander("⚙️ All-Stages Analysis Settings", expanded=True):
        col1, col2, col3, col4, col5, col6 = st.columns(6)
        with col1:
            if st.runtime.exists():
                test_run = st.checkbox("Test Run (5 files/group)", value=False, key="cmp_all_stages_test_run")
            else:
                test_run = True
        with col2:
            n_permutations = st.slider("Permutations", 50, 500, 200, 50, key="cmp_all_stages_n_perm")
        with col3:
            use_zscore = st.checkbox("Z-score subjects", value=True, key="cmp_all_stages_zscore")
        with col4:
            apply_ica = st.checkbox("Apply ICA", value=True, key="cmp_all_stages_apply_ica",
                                    help="Apply ICA to remove ECG artifact components from EEG channels.")
        with col5:
            show_combined = st.checkbox("Show Combined Plot", value=True, key="cmp_all_stages_combined",
                                        help="Show all stages side-by-side in a single combined figure.")
        with col6:
            recompute_all_cache = st.button("Recompute All Caches", key="cmp_all_stages_recompute_cache")

    # ── Per-stage/group cache reset buttons ─────────────────────────────────
    if 'cmp_all_stages_recompute_combos' not in st.session_state:
        st.session_state.cmp_all_stages_recompute_combos = set()

    with st.expander("🔄 Cache Reset (per Group × Sleep Stage)", expanded=False):
        st.caption("Press a button to recompute the cache for a specific group × sleep stage.")
        header_cols = st.columns([2] + [1] * len(SLEEP_STAGES))
        header_cols[0].markdown("**Group**")
        for i, stage in enumerate(SLEEP_STAGES):
            header_cols[i + 1].markdown(f"**{stage}**")
        for group in selected_groups:
            row_cols = st.columns([2] + [1] * len(SLEEP_STAGES))
            row_cols[0].markdown(f"_{group}_")
            for i, stage in enumerate(SLEEP_STAGES):
                if row_cols[i + 1].button(stage, key=f"reset_cache_{group}_{stage}"):
                    st.session_state.cmp_all_stages_recompute_combos.add((group, stage))
                    if hasattr(get_group_individuals, "clear"):
                        get_group_individuals.clear()
                    st.toast(f"Cache reset queued for {group} / {stage}", icon="🔄")

    if recompute_all_cache:
        for _g in selected_groups:
            for _s in SLEEP_STAGES:
                st.session_state.cmp_all_stages_recompute_combos.add((_g, _s))
        if hasattr(get_group_individuals, "clear"):
            get_group_individuals.clear()

    recompute_combos = frozenset(st.session_state.cmp_all_stages_recompute_combos)

    # Load globally excluded patients
    global_excluded_pids = load_excluded_pids(base_path)

    def _valid_hep_individuals(individuals, label):
        """Keep only individuals with usable HEP arrays for this all-stages view."""
        valid = []
        skipped = 0
        for ind in individuals or []:
            if ind is None or len(ind) < 4:
                skipped += 1
                continue
            hep_data, times, ch_names = ind[1], ind[2], ind[3]
            if hep_data is None or times is None or ch_names is None:
                skipped += 1
                continue
            try:
                hep_arr = np.asarray(hep_data)
            except Exception:
                skipped += 1
                continue
            if hep_arr.ndim < 2 or len(ch_names) == 0 or hep_arr.shape[0] < len(ch_names):
                skipped += 1
                continue
            valid.append(ind)
        if skipped:
            st.caption(f"Skipped {skipped} patient(s) without usable HEP data for {label}.")
        return valid

    # ── Load data for all stages and groups ─────────────────────────────────
    _load_status = st.empty()
    _load_status.info("Loading data for all sleep stages...")
    all_stage_data = {}  # stage -> group -> list of individual tuples

    progress = st.progress(0)
    for s_idx, stage in enumerate(SLEEP_STAGES):
        all_stage_data[stage] = {}
        for group in selected_groups:
            needs_recompute = (group, stage) in recompute_combos
            with st.spinner(f"Loading {group} / {stage}…"):
                inds = get_group_individuals(group, stage, base_path, test_run=test_run, apply_ica=apply_ica, recompute_cache=needs_recompute)
            inds_filtered = filter_excluded(
                _valid_hep_individuals(
                    inds,
                    f"{group} / {stage}",
                ),
                global_excluded_pids,
            )
            if inds_filtered:
                all_stage_data[stage][group] = inds_filtered
        progress.progress((s_idx + 1) / len(SLEEP_STAGES))
    progress.empty()
    _load_status.empty()
    st.session_state.cmp_all_stages_recompute_combos.clear()

    # ── Load raw + ICA datasets for cross-stage topomap summaries ───────────
    topomap_sources = {}
    topomap_sources['ica' if apply_ica else 'raw'] = all_stage_data

    for method_name, method_apply_ica in (('raw', False), ('ica', True)):
        if method_name in topomap_sources:
            continue
        method_stage_data = {}
        for stage in SLEEP_STAGES:
            method_stage_data[stage] = {}
            for group in selected_groups:
                needs_recompute = (group, stage) in recompute_combos
                with st.spinner(f"Loading {method_name.upper()} topomap data for {group} / {stage}…"):
                    inds = get_group_individuals(
                        group,
                        stage,
                        base_path,
                        test_run=test_run,
                        apply_ica=method_apply_ica,
                        recompute_cache=needs_recompute,
                    )
                inds_filtered = filter_excluded(
                    _valid_hep_individuals(
                        inds,
                        f"{method_name.upper()} {group} / {stage}",
                    ),
                    global_excluded_pids,
                )
                if inds_filtered:
                    method_stage_data[stage][group] = inds_filtered
        topomap_sources[method_name] = method_stage_data

    line_stat_caption = (
        "Solid colored lines show the group median HEP at each time point, not the mean/average. "
        "The shaded region is median absolute deviation (MAD) around the median, not standard deviation. "
        f"Amplitude values are {'per-subject z-scores' if use_zscore else 'microvolts (µV)'}."
    )

    def _prepare_all_stages_trace(waveform):
        waveform = np.asarray(waveform, dtype=float)
        if use_zscore and np.std(waveform) > 0:
            return (waveform - np.mean(waveform)) / np.std(waveform)
        if not use_zscore:
            return waveform * 1e6
        return waveform

    def _show_plot_count_table(rows, label):
        if rows:
            st.dataframe(pd.DataFrame(rows), use_container_width=True)
        else:
            st.caption(f"No patient-count rows available for {label}.")

    # ── Combined multi-stage plot ────────────────────────────────────────────
    if show_combined:
        st.subheader("Combined: All Sleep Stages × Groups")

        # Collect common channels across all loaded data
        all_ch_names = None
        for stage in SLEEP_STAGES:
            for group in selected_groups:
                inds = all_stage_data[stage].get(group, [])
                if inds:
                    ch_names_candidate = inds[0][3]
                    if all_ch_names is None:
                        all_ch_names = ch_names_candidate
                    else:
                        all_ch_names = [c for c in all_ch_names if c in ch_names_candidate]

        if all_ch_names is None:
            st.warning("No data loaded for any stage/group combination.")
        else:
            # Show a summary table of available data counts
            summary_rows = []
            for stage in SLEEP_STAGES:
                row = {'Stage': stage}
                for group in selected_groups:
                    row[group] = len(all_stage_data[stage].get(group, []))
                summary_rows.append(row)
            st.dataframe(pd.DataFrame(summary_rows).set_index('Stage'), use_container_width=True)

            # For the combined plot, pick up to 3 representative channels
            plot_channels = all_ch_names[:min(3, len(all_ch_names))]

            n_stages = len(SLEEP_STAGES)
            n_plot_ch = len(plot_channels)
            fig_combined, axes = plt.subplots(
                n_plot_ch, n_stages,
                figsize=(4 * n_stages, 3 * n_plot_ch),
                squeeze=False,
                sharey='row'
            )
            fig_combined.suptitle(
                f"HEP Group Comparison - Median +/- MAD | Stages: All Sleep Stages | Groups: {groups_label}",
                fontsize=14, fontweight='bold'
            )
            combined_count_rows = []

            for s_idx, stage in enumerate(SLEEP_STAGES):
                for ch_idx, ch_name in enumerate(plot_channels):
                    ax = axes[ch_idx][s_idx]
                    ax.axhline(0, color='gray', lw=0.5, ls='--')
                    ax.axvline(0, color='gray', lw=0.5, ls='--')

                    has_data = False
                    for g_idx, group in enumerate(selected_groups):
                        inds = all_stage_data[stage].get(group, [])
                        if not inds:
                            continue
                        waveforms = []
                        times_arr = None
                        for ind in inds:
                            pid, hep_data, times, ch_names_ind = ind[0], ind[1], ind[2], ind[3]
                            if ch_name not in ch_names_ind:
                                continue
                            ch_idx_ind = list(ch_names_ind).index(ch_name)
                            w = _prepare_all_stages_trace(hep_data[ch_idx_ind])
                            waveforms.append(w)
                            if times_arr is None:
                                times_arr = times
                        if not waveforms or times_arr is None:
                            continue
                        waveforms = np.array(waveforms)
                        median_wave, mad_wave = median_and_mad(waveforms, axis=0)
                        color = plt.cm.Set1(g_idx / max(len(selected_groups), 1))
                        ax.plot(times_arr * 1000, median_wave, color=color, lw=1.5, label=f"{group} (n={len(waveforms)})")
                        ax.fill_between(times_arr * 1000, median_wave - mad_wave, median_wave + mad_wave,
                                        color=color, alpha=0.2)
                        combined_count_rows.append({
                            "Stage": stage,
                            "Subplot channel": ch_name,
                            "Group": group,
                            "Patients plotted": len(waveforms),
                            "Line statistic": "Median",
                            "Shaded band": "MAD",
                        })
                        has_data = True

                    if ch_idx == 0:
                        ax.set_title(f"Stage: {stage}", fontsize=11, color=STAGE_COLORS.get(stage, 'black'), fontweight='bold')
                    if s_idx == 0:
                        ax.set_ylabel(ch_name, fontsize=9)
                    if ch_idx == n_plot_ch - 1:
                        ax.set_xlabel("Time (ms)", fontsize=8)
                    if has_data and ch_idx == 0 and s_idx == n_stages - 1:
                        ax.legend(fontsize=7, loc='upper right')
                    ax.tick_params(labelsize=7)

            fig_combined.tight_layout()
            st.pyplot(fig_combined, use_container_width=True)
            plt.close(fig_combined)
            st.caption(line_stat_caption)
            _show_plot_count_table(combined_count_rows, "the combined all-stages plot")

    # ── Per-stage individual plots ───────────────────────────────────────────
    st.subheader("Per-Stage Group Comparisons")
    for stage in SLEEP_STAGES:
        stage_data = all_stage_data[stage]
        if not any(stage_data.values()):
            st.warning(f"No data available for stage {stage}.")
            continue

        with st.expander(f"Stage: {stage}", expanded=(stage == 'light_sleep')):
            st.markdown(f"**Sleep Stage: {stage}** — Groups: {', '.join(selected_groups)}")

            ch_names_stage = None
            for group in selected_groups:
                inds = stage_data.get(group, [])
                if inds:
                    candidate = inds[0][3]
                    if ch_names_stage is None:
                        ch_names_stage = list(candidate)
                    else:
                        ch_names_stage = [c for c in ch_names_stage if c in candidate]

            if ch_names_stage is None:
                st.warning(f"No common channels for stage {stage}.")
                continue

            n_ch = len(ch_names_stage)
            n_cols = min(4, n_ch)
            n_rows = int(np.ceil(n_ch / n_cols))
            # unique patient ids per group for this stage (each individual tuple is (patient_id, ...))
            stage_groups_label = ", ".join(
                f"{group} (n={len(set(ind[0] for ind in stage_data.get(group, [])))})"
                for group in selected_groups
            )
            fig_stage, axes_stage = plt.subplots(n_rows, n_cols, figsize=(5 * n_cols, 3 * n_rows), squeeze=False)
            fig_stage.suptitle(
                f"HEP Group Comparison - Median +/- MAD | Stage: {stage} | Groups: {stage_groups_label}",
                fontsize=13, fontweight='bold',
                color=STAGE_COLORS.get(stage, 'black')
            )
            stage_count_rows = []

            for ch_i, ch_name in enumerate(ch_names_stage):
                row_i, col_i = divmod(ch_i, n_cols)
                ax = axes_stage[row_i][col_i]
                ax.axhline(0, color='gray', lw=0.5, ls='--')
                ax.axvline(0, color='gray', lw=0.5, ls='--')
                ax.set_title(ch_name, fontsize=8)

                for g_idx, group in enumerate(selected_groups):
                    inds = stage_data.get(group, [])
                    waveforms = []
                    times_arr = None
                    for ind in inds:
                        pid, hep_data, times, ch_names_ind = ind[0], ind[1], ind[2], ind[3]
                        if ch_name not in ch_names_ind:
                            continue
                        ch_idx_ind = list(ch_names_ind).index(ch_name)
                        w = _prepare_all_stages_trace(hep_data[ch_idx_ind])
                        waveforms.append(w)
                        if times_arr is None:
                            times_arr = times
                    if not waveforms or times_arr is None:
                        continue
                    waveforms = np.array(waveforms)
                    median_wave, mad_wave = median_and_mad(waveforms, axis=0)
                    color = plt.cm.Set1(g_idx / max(len(selected_groups), 1))
                    ax.plot(times_arr * 1000, median_wave, color=color, lw=1.5, label=f"{group} (n={len(waveforms)})")
                    ax.fill_between(times_arr * 1000, median_wave - mad_wave, median_wave + mad_wave,
                                    color=color, alpha=0.2)
                    stage_count_rows.append({
                        "Stage": stage,
                        "Subplot channel": ch_name,
                        "Group": group,
                        "Patients plotted": len(waveforms),
                        "Line statistic": "Median",
                        "Shaded band": "MAD",
                    })

                ax.tick_params(labelsize=7)
                if ch_i == 0:
                    ax.legend(fontsize=7)
                ax.set_xlabel("Time (ms)", fontsize=7)
                ylabel = "Z-score" if use_zscore else "Amplitude (µV)"
                ax.set_ylabel(ylabel, fontsize=7)

            for ch_i in range(n_ch, n_rows * n_cols):
                row_i, col_i = divmod(ch_i, n_cols)
                axes_stage[row_i][col_i].axis('off')

            fig_stage.tight_layout()
            st.pyplot(fig_stage, use_container_width=True)
            plt.close(fig_stage)
            st.caption(line_stat_caption)
            _show_plot_count_table(stage_count_rows, f"stage {stage}")

    # ── Significant-channel topomaps across all sleep stages ────────────────
    st.divider()
    st.subheader("Significant Channels Topomap (Minimum cluster P-value)")

    if len(selected_groups) != 2:
        st.info("Select exactly 2 groups to render the all-sleep-stage group-comparison topomaps.")
        return

    group_a, group_b = selected_groups
    topomap_methods = [
        ("raw", "Raw", False),
        ("ica", "ICA Cleaned", False),
        ("ica_csd", "ICA + CSD Cleaned", True),
    ]

    with st.spinner("Computing significant-channel topomaps across sleep stages…"):
        topomap_results = {}
        for method_key, method_label, use_csd in topomap_methods:
            source_key = 'ica' if method_key == 'ica_csd' else method_key
            method_data = topomap_sources.get(source_key, {})

            stage_diff_maps = {}
            pooled_group_inds = {group: [] for group in selected_groups}
            for stage in SLEEP_STAGES:
                stage_groups = method_data.get(stage, {})
                inds_a = stage_groups.get(group_a, [])
                inds_b = stage_groups.get(group_b, [])
                if inds_a and inds_b:
                    stage_diff_maps[stage] = _compute_two_group_channel_pvals(
                        inds_a,
                        inds_b,
                        n_permutations=n_permutations,
                        jitter_sec=TOPO_JITTER_SEC,
                        use_csd=use_csd,
                    )
                else:
                    stage_diff_maps[stage] = {}

                for group in selected_groups:
                    pooled_group_inds[group].extend(stage_groups.get(group, []))

            pooled_group_maps = {
                group: _compute_stage_or_group_pvals(
                    pooled_group_inds[group],
                    n_permutations=n_permutations,
                    jitter_sec=TOPO_JITTER_SEC,
                    use_csd=use_csd,
                )
                for group in selected_groups
            }

            pooled_diff_map = _compute_two_group_channel_pvals(
                pooled_group_inds[group_a],
                pooled_group_inds[group_b],
                n_permutations=n_permutations,
                jitter_sec=TOPO_JITTER_SEC,
                use_csd=use_csd,
            )

            topomap_results[method_key] = {
                "label": method_label,
                "stage_diff_maps": stage_diff_maps,
                "pooled_group_maps": pooled_group_maps,
                "pooled_diff_map": pooled_diff_map,
            }

    for method_key, method_label, _use_csd in topomap_methods:
        result = topomap_results.get(method_key, {})
        stage_diff_maps = result.get("stage_diff_maps", {})
        pooled_group_maps = result.get("pooled_group_maps", {})
        pooled_diff_map = result.get("pooled_diff_map", {})

        st.markdown(f"#### {method_label}")

        fig_stage_topo, axes_stage_topo = plt.subplots(1, len(SLEEP_STAGES), figsize=(4 * len(SLEEP_STAGES), 4))
        if len(SLEEP_STAGES) == 1:
            axes_stage_topo = [axes_stage_topo]
        fig_stage_topo.suptitle(
            f"{method_label}: Group Difference Topomaps | Stages: All Sleep Stages | Groups: {group_a} vs {group_b}",
            fontsize=13, fontweight='bold'
        )
        for idx, stage in enumerate(SLEEP_STAGES):
            p_map = stage_diff_maps.get(stage, {})
            if p_map:
                _render_pval_topomap(p_map, axes_stage_topo[idx], stage)
            else:
                axes_stage_topo[idx].axis('off')
                axes_stage_topo[idx].set_title(f"{stage}\nNo data")
        fig_stage_topo.tight_layout()
        st.pyplot(fig_stage_topo, use_container_width=True)
        plt.close(fig_stage_topo)

        fig_group_topo, axes_group_topo = plt.subplots(1, len(selected_groups), figsize=(5 * len(selected_groups), 4))
        if len(selected_groups) == 1:
            axes_group_topo = [axes_group_topo]
        fig_group_topo.suptitle(
            f"{method_label}: Average per Group | Stages: All Sleep Stages | Groups: {group_a} vs {group_b}",
            fontsize=13, fontweight='bold'
        )
        for idx, group in enumerate(selected_groups):
            p_map = pooled_group_maps.get(group, {})
            if p_map:
                _render_pval_topomap(p_map, axes_group_topo[idx], f"{group} | all stages")
            else:
                axes_group_topo[idx].axis('off')
                axes_group_topo[idx].set_title(f"{group}\nNo data")
        fig_group_topo.tight_layout()
        st.pyplot(fig_group_topo, use_container_width=True)
        plt.close(fig_group_topo)

        fig_diff_topo, ax_diff_topo = plt.subplots(figsize=(5, 4.5))
        fig_diff_topo.suptitle(
            f"{method_label}: Group Difference | Stages: All Sleep Stages | Groups: {group_a} vs {group_b}",
            fontsize=13, fontweight='bold'
        )
        if pooled_diff_map:
            _render_pval_topomap(pooled_diff_map, ax_diff_topo, "All sleep stages")
        else:
            ax_diff_topo.axis('off')
            ax_diff_topo.set_title("No data")
        fig_diff_topo.tight_layout()
        st.pyplot(fig_diff_topo, use_container_width=False)
        plt.close(fig_diff_topo)


def run_compare_sleep_stages_analysis(base_path):
    """
    Logic for Compare Sleep Stages mode.
    Plots ECG HEP and EEG HEP across different sleep stages for single patients or group averages
    (only for patients with valid data in ALL stages).
    """
    # 1. Select Group
    available_groups = [g for g in os.listdir(base_path) if os.path.isdir(os.path.join(base_path, g))]
    if not available_groups:
        st.error("No groups found.")
        return
    selected_group = st.selectbox("Select Group", available_groups, index=0)

    # 2. Select Patient and Filter for Completeness
    sleep_stages = ['W', 'light_sleep', 'N3', 'R']
    
    if st.runtime.exists():
        col1, col2, col3 = st.columns(3)
        with col1:
            test_run = st.checkbox("Test Run (first 5 files only)", value=False, key="test_run_compare_stages")
        with col2:
            recompute_cache = st.button("Recompute Cache", key="recompute_cache_compare_stages")
        with col3:
            apply_ica = st.checkbox("Apply ICA", value=True, key="apply_ica_compare_stages",
                                    help="Apply ICA to remove ECG artifact components from EEG channels.")
    else:
        test_run = True
        recompute_cache = False
        apply_ica = True
    st.info("Scanning for patients with valid data (passed R-peaks test) across ALL sleep stages... (this may take a moment if not cached)")
    if recompute_cache:
        if hasattr(get_group_individuals, "clear"):
            get_group_individuals.clear()
    progress_scan = st.progress(0)
    
    valid_patients_per_stage = []
    # Cache all loaded individuals per stage for quick access later
    all_stage_individuals = {}
    
    # ── Load globally excluded patients ─────────────────────────────────────
    global_excluded_pids = load_excluded_pids(base_path)

    unfiltered_patients_per_stage = []  # for display table (includes excluded)
    for idx, stage in enumerate(sleep_stages):
        # We use get_group_individuals to reliably find valid files that passed processing
        stage_individuals = get_group_individuals(selected_group, stage, base_path, test_run=test_run, recompute_cache=recompute_cache, apply_ica=apply_ica)

        # Track unfiltered patients (for the availability table)
        unfiltered_stage_patients = set()
        for ind in stage_individuals:
            pid_full = ind[0]
            pid_base = pid_full.split('_')[0]
            unfiltered_stage_patients.add(pid_base)
        unfiltered_patients_per_stage.append(unfiltered_stage_patients)

        # Filter out globally excluded patients
        stage_individuals = filter_excluded(stage_individuals, global_excluded_pids)
        all_stage_individuals[stage] = stage_individuals

        stage_patients = set()
        for ind in stage_individuals:
            pid_full = ind[0]
            pid_base = pid_full.split('_')[0]
            stage_patients.add(pid_base)
        valid_patients_per_stage.append(stage_patients)
        progress_scan.progress((idx + 1) / len(sleep_stages))
        
    progress_scan.empty()
    
    if valid_patients_per_stage:
        all_patients = set.intersection(*valid_patients_per_stage)
        union_all_patients = set.union(*valid_patients_per_stage)
    else:
        all_patients = set()
        union_all_patients = set()

    # Expandable container showing all patients and their available stages
    union_all_patients_unfiltered = set.union(*unfiltered_patients_per_stage) if unfiltered_patients_per_stage else set()
    if union_all_patients_unfiltered:
        excluded_set = set(global_excluded_pids)

        with st.expander("Patient Sleep Stage Availability", expanded=False):
            st.markdown("This table details which sleep stages are available for each patient in this group.")
            stage_avail_data = []
            for pid in sorted(list(union_all_patients_unfiltered)):
                row = {'Patient ID': pid}
                n_available = 0
                for stage, stage_patients in zip(sleep_stages, unfiltered_patients_per_stage):
                    has_stage = pid in stage_patients
                    row[stage] = "Y" if has_stage else "N"
                    if has_stage:
                        n_available += 1
                row['Total Stages'] = n_available
                row['Excluded'] = "Y" if pid in excluded_set else "N"
                stage_avail_data.append(row)

            df_availability = pd.DataFrame(stage_avail_data).set_index('Patient ID')

            def color_rows(row):
                total = row['Total Stages']
                if total == len(sleep_stages):
                    bg_color = 'background-color: rgba(0, 255, 0, 0.2)'  # Green for all
                elif total == 0:
                    bg_color = 'background-color: rgba(255, 0, 0, 0.4)'  # Red for none
                elif total <= 2:
                    bg_color = 'background-color: rgba(255, 100, 0, 0.3)' # Orange for few
                else:
                    bg_color = 'background-color: rgba(255, 255, 0, 0.2)' # Yellow for some
                styles = [bg_color] * len(row)
                if row.get('Excluded') == "Y":
                    excluded_col_idx = row.index.get_loc('Excluded')
                    styles[excluded_col_idx] = 'background-color: rgba(220, 38, 38, 0.5); color: white; font-weight: bold'
                return styles

            st.dataframe(df_availability.style.apply(color_rows, axis=1), use_container_width=True)

    if not all_patients:
        stage_counts = {
            stage: len(stage_patients)
            for stage, stage_patients in zip(sleep_stages, valid_patients_per_stage)
        }
        counts_text = ", ".join(f"{stage}: {count}" for stage, count in stage_counts.items())
        st.warning(
            f"No patients found in group {selected_group} with valid data in ALL stages. "
            f"Valid patient counts by stage: {counts_text}. "
            f"Unique patients with at least one valid stage: {len(union_all_patients)}."
        )
        return

    # 3. Choose Mode
    analysis_mode = st.radio("Analysis Type", ["Single Patient", "Group Average"], horizontal=True)

    if analysis_mode == "Single Patient":
        selected_pid = st.selectbox("Select Patient", sorted(list(all_patients)))
    else:
        selected_pid = "Group Average"

    st.subheader(f"Comparisons across Sleep Stages - {selected_pid} - {selected_group}")
    
    stage_data = {}
    
    # Define colors for stages
    stage_colors = {
        'W': 'orange',
        'light_sleep': 'lightblue',
        'N3': 'blue',
        'R': 'red'
    }

    def _build_patient_ecg_stage_traces(pid):
        """Per-stage (times, ecg_trace) for one patient, used for example plots."""
        traces = {}
        for stage in sleep_stages:
            for ind in all_stage_individuals[stage]:
                if ind[0].split('_')[0] == pid:
                    ecg_hep, times_p = ind[5], ind[2]
                    if ecg_hep is not None and times_p is not None:
                        ecg_trace = np.asarray(ecg_hep, dtype=float).squeeze()
                        min_len = min(len(times_p), len(ecg_trace))
                        traces[stage] = (np.asarray(times_p)[:min_len], ecg_trace[:min_len])
                    break
        return traces

    if analysis_mode == "Group Average":
        st.markdown("#### Example Individual Patients (ECG HEP across Sleep Stages)")
        example_pids = sorted(list(all_patients))[:3]
        example_cols = st.columns(len(example_pids))
        for col, pid in zip(example_cols, example_pids):
            with col:
                fig_ex, ax_ex = plt.subplots(figsize=(4, 3))
                for stage, (t_p, trace_p) in _build_patient_ecg_stage_traces(pid).items():
                    ax_ex.plot(t_p, trace_p * 1e6, label=stage, color=stage_colors.get(stage, 'gray'),
                               linewidth=1.5, alpha=0.85)
                ax_ex.axvline(0, color='black', linestyle='--', alpha=0.4)
                ax_ex.set_title(f"Patient {pid}", fontsize=10)
                ax_ex.set_xlabel("Time (s)", fontsize=8)
                ax_ex.set_ylabel(u"μV", fontsize=8)
                ax_ex.legend(fontsize=7)
                st.pyplot(fig_ex, use_container_width=True)
                plt.close(fig_ex)

    for stage in sleep_stages:
        stage_individuals = all_stage_individuals[stage]
        
        if analysis_mode == "Single Patient":
            for ind in stage_individuals:
                pid_full = ind[0]
                pid_base = pid_full.split('_')[0]
                if pid_base == selected_pid:
                    ecg_hep = ind[5]
                    eeg_hep = ind[1]
                    times = ind[2]
                    if times is not None:
                        common_len = len(times)
                        if ecg_hep is not None:
                            common_len = min(common_len, np.asarray(ecg_hep).shape[-1])
                        if eeg_hep is not None:
                            common_len = min(common_len, np.asarray(eeg_hep).shape[-1])
                        times = np.asarray(times)[:common_len]
                        if ecg_hep is not None:
                            ecg_hep = np.asarray(ecg_hep)[..., :common_len]
                        if eeg_hep is not None:
                            eeg_hep = np.asarray(eeg_hep)[..., :common_len]
                    stage_data[stage] = {
                        'ecg_hep': ecg_hep,
                        'eeg_hep': eeg_hep,
                        'times': times,
                        'ch_names': ind[3],
                        'n_epochs': len(ind[4]),
                        'n_subjects': 1
                    }
                    break
        else:
            # Group Average over all completely valid patients
            ecg_heps = []
            eeg_heps_list = []
            n_epochs_list = []
            times = None
            
            # Find common channels for this stage across all valid patients
            common_channels_stage = None
            valid_inds_dict = {}
            for ind in stage_individuals:
                pid_full = ind[0]
                pid_base = pid_full.split('_')[0]
                if pid_base in all_patients:
                    valid_inds_dict[pid_base] = ind
                    if common_channels_stage is None:
                        common_channels_stage = set(ind[3])
                    else:
                        common_channels_stage = common_channels_stage.intersection(set(ind[3]))
                        
            # Ensure valid_inds are ordered consistently by pid_base for pairwise tests
            valid_inds = [valid_inds_dict[p] for p in sorted(list(all_patients))]
                        
            if common_channels_stage:
                common_channels_stage = sorted(list(common_channels_stage))
                # Optional regex to filter EEG channels
                common_channels_stage = [ch for ch in common_channels_stage if re.match(r'^[a-zA-Z][0-9]*$', ch) or re.match(r'^[a-zA-Z]z$', ch)]
            else:
                common_channels_stage = []

            for ind in valid_inds:
                if ind[5] is not None:
                    ecg_heps.append(ind[5])
                if ind[1] is not None:
                    eeg_data = apply_ecg_regression(ind[1], ind[5] if len(ind) > 5 else None)
                    ch_names_ind = ind[3]
                    aligned_eeg = []
                    for ch in common_channels_stage:
                        ch_idx = ch_names_ind.index(ch)
                        aligned_eeg.append(eeg_data[ch_idx])
                    if aligned_eeg:
                        eeg_heps_list.append(aligned_eeg)
                n_epochs_list.append(len(ind[4]))
                if times is None:
                    times = ind[2]
                    
            if ecg_heps and eeg_heps_list and times is not None:
                # Ensure all arrays in ecg_heps have the same shape
                # Sometimes different sampling rates cause slight length mismatches
                min_len_ecg = min(len(x[0]) if len(x.shape) > 1 else len(x) for x in ecg_heps)
                ecg_heps_aligned = [x[:, :min_len_ecg] if len(x.shape) > 1 else x[:min_len_ecg] for x in ecg_heps]
                
                min_len_eeg = min(len(x[0]) for x in eeg_heps_list)
                eeg_heps_aligned = [[ch_arr[:min_len_eeg] for ch_arr in patient_eeg] for patient_eeg in eeg_heps_list]
                
                common_len = min(len(times), min_len_ecg, min_len_eeg)
                ecg_heps_aligned = [
                    x[:, :common_len] if len(x.shape) > 1 else x[:common_len]
                    for x in ecg_heps_aligned
                ]
                eeg_heps_aligned = [
                    [ch_arr[:common_len] for ch_arr in patient_eeg]
                    for patient_eeg in eeg_heps_aligned
                ]
                times = times[:common_len]

                avg_ecg = np.nanmedian(ecg_heps_aligned, axis=0)
                avg_eeg = np.nanmedian(eeg_heps_aligned, axis=0)
                avg_epochs = int(np.mean(n_epochs_list))
                
                stage_data[stage] = {
                    'ecg_hep': avg_ecg,
                    'eeg_hep': avg_eeg,
                    'times': times,
                    'ch_names': common_channels_stage,
                    'n_epochs': avg_epochs,
                    'n_subjects': len(ecg_heps),
                    'ecg_heps_aligned': np.array(ecg_heps_aligned),
                    'eeg_heps_aligned': np.array(eeg_heps_aligned)
                }
                
    if not stage_data:
        st.warning("No data found for any stage.")
        return

    # --- Group-level cluster permutation test: is HEP significantly != 0 per stage? ---
    stage_sig_windows = {}
    stage_permtest_info = {}
    if analysis_mode == "Group Average":
        for stage, data in stage_data.items():
            ecg_arr = data.get('ecg_heps_aligned')
            times_stage = data.get('times')
            if ecg_arr is None or times_stage is None:
                continue
            ecg_arr = np.asarray(ecg_arr, dtype=float)
            if ecg_arr.ndim == 3:
                ecg_arr = ecg_arr[:, 0, :]
            if ecg_arr.shape[0] < 3:
                continue
            sig_windows, _, per_patient_info = permutation_cluster_jitter_test(
                ecg_arr, np.asarray(times_stage)[:ecg_arr.shape[1]], n_permutations=100
            )
            stage_sig_windows[stage] = sig_windows
            stage_permtest_info[stage] = per_patient_info

    # --- Plot ECG ---
    fig_ecg, ax_ecg = plt.subplots(figsize=(10, 6))
    has_ecg_data = False

    for stage, data in stage_data.items():
        ecg_hep = data['ecg_hep']
        times = data['times']
        
        if ecg_hep is not None and times is not None:
             # ecg_hep is (1, n_times)
             n_epochs = data['n_epochs']
             label_str = f"{stage} (n={n_epochs})" if analysis_mode == "Single Patient" else f"{stage} (N={data['n_subjects']}, avg epochs={n_epochs})"
             ecg_trace = np.asarray(ecg_hep, dtype=float).squeeze()
             min_len = min(len(times), len(ecg_trace))
             ax_ecg.plot(np.asarray(times)[:min_len], ecg_trace[:min_len] * 1e6, label=label_str, color=stage_colors.get(stage, 'gray'), linewidth=2, alpha=0.8)
             for win in stage_sig_windows.get(stage, []):
                 ax_ecg.axvspan(win['start'], win['end'], color=stage_colors.get(stage, 'gray'), alpha=0.15)
             has_ecg_data = True

    if has_ecg_data:
        title_str = f"ECG HEP across Sleep Stages - {selected_pid}" if analysis_mode == "Single Patient" else f"ECG HEP across Sleep Stages - Group Average (N={len(all_patients)} patients)"
        ax_ecg.set_title(title_str)
        ax_ecg.set_xlabel("Time (s)")
        ax_ecg.set_ylabel("Amplitude (μV)")
        ax_ecg.grid(True)
        ax_ecg.legend()
        ax_ecg.axvline(0, color='black', linestyle='--', alpha=0.5)
        st.pyplot(fig_ecg, use_container_width=True)

        # --- Significant HEP windows per stage (cluster permutation test vs 0) ---
        if analysis_mode == "Group Average" and stage_sig_windows:
            st.markdown("#### Significant HEP Windows per Sleep Stage (cluster permutation test, HEP != 0)")
            rows = []
            for stage in stage_data.keys():
                info = stage_permtest_info.get(stage)
                if info is None:
                    continue
                windows = stage_sig_windows.get(stage, [])
                win_str = ", ".join(
                    f"{w['start']*1000:.0f}-{w['end']*1000:.0f}ms (p={w['p_value']:.4f})" for w in windows
                ) or "none"
                rows.append({
                    'Sleep Stage': stage,
                    'N patients': len(info['p_values']),
                    'N sig. patients': info['n_significant'],
                    'Fisher p': f"{info['fisher_p']:.4f}",
                    'Significant windows': win_str,
                })
            if rows:
                st.dataframe(pd.DataFrame(rows), use_container_width=True)

            fig_sig, ax_sig = plt.subplots(figsize=(10, max(2, len(stage_sig_windows) * 0.6)))
            y_ticks, y_labels = [], []
            for y, (stage, windows) in enumerate(stage_sig_windows.items(), start=1):
                y_ticks.append(y)
                y_labels.append(stage)
                for w in windows:
                    ax_sig.plot([w['start'], w['end']], [y, y], color=stage_colors.get(stage, 'gray'),
                                linewidth=6, solid_capstyle='butt')
            ax_sig.set_yticks(y_ticks)
            ax_sig.set_yticklabels(y_labels)
            ax_sig.set_xlabel("Time (s)")
            ax_sig.set_title("Cluster-Significant HEP Windows (Group Average, vs 0)")
            ax_sig.axvline(0, color='black', linestyle='--', alpha=0.5)
            ax_sig.grid(True, axis='x', alpha=0.3)
            ax_sig.set_xlim(ax_ecg.get_xlim())
            st.pyplot(fig_sig, use_container_width=True)
            plt.close(fig_sig)

        # --- Per-patient significance classification + significant-only group average ---
        if analysis_mode == "Group Average" and stage_permtest_info:
            st.markdown("#### Per-Patient Significance Classification per Sleep Stage")
            sorted_patients = sorted(list(all_patients))
            patient_rows = []
            for stage, info in stage_permtest_info.items():
                sig_flags = info.get('significant', [])
                for pid, is_sig in zip(sorted_patients, sig_flags):
                    patient_rows.append({
                        'Patient ID': pid,
                        'Sleep Stage': stage,
                        'Significant': bool(is_sig),
                    })
            if patient_rows:
                st.dataframe(pd.DataFrame(patient_rows), use_container_width=True)

            fig_sigonly, ax_sigonly = plt.subplots(figsize=(10, 6))
            has_sigonly_data = False
            for stage, data in stage_data.items():
                info = stage_permtest_info.get(stage)
                if info is None:
                    continue
                sig_flags = np.asarray(info.get('significant', []), dtype=bool)
                if sig_flags.sum() == 0:
                    st.info(f"No significant patients for stage '{stage}' - skipping significant-only average.")
                    continue
                ecg_arr = np.asarray(data.get('ecg_heps_aligned'), dtype=float)
                if ecg_arr.ndim == 3:
                    ecg_arr = ecg_arr[:, 0, :]
                sig_ecg = ecg_arr[sig_flags]
                avg_sig_ecg = np.nanmedian(sig_ecg, axis=0)
                times_stage = np.asarray(data.get('times'))
                min_len = min(len(times_stage), len(avg_sig_ecg))
                ax_sigonly.plot(times_stage[:min_len], avg_sig_ecg[:min_len] * 1e6,
                                 label=f"{stage} (N={int(sig_flags.sum())} significant)",
                                 color=stage_colors.get(stage, 'gray'), linewidth=2, alpha=0.8)
                has_sigonly_data = True

            if has_sigonly_data:
                ax_sigonly.set_title("Group Average ECG HEP - Significant Patients Only")
                ax_sigonly.set_xlabel("Time (s)")
                ax_sigonly.set_ylabel("Amplitude (μV)")
                ax_sigonly.grid(True)
                ax_sigonly.legend()
                ax_sigonly.axvline(0, color='black', linestyle='--', alpha=0.5)
                st.pyplot(fig_sigonly, use_container_width=True)
            plt.close(fig_sigonly)

        # Add Pairwise Comparison Matrix for Group Average
        if analysis_mode == "Group Average" and len(stage_data) > 1:
            st.markdown("#### Pairwise Significance (ECG HEP)")
            stages_list = list(stage_data.keys())
            
            # Use a standard T-test + FDR for pairwise matrix to get exact p-values for all pairs
            try:
                from mne.stats import fdr_correction
                HAS_FDR = True
            except ImportError:
                HAS_FDR = False

            # Build p-value lookup dict for all pairs
            pair_pvals = {}  # (stage_a, stage_b) -> (min_p, is_significant, windows_str)
            pair_sig_clusters = {} # (stage_a, stage_b) -> [list of sub-arrays of sig_times]
            for i in range(len(stages_list)):
                for j in range(i + 1, len(stages_list)):
                    stage_a = stages_list[i]
                    stage_b = stages_list[j]
                    heps_a = stage_data[stage_a].get('ecg_heps_aligned')
                    heps_b = stage_data[stage_b].get('ecg_heps_aligned')
                    
                    if heps_a is not None and heps_b is not None and len(heps_a) > 1 and len(heps_b) > 1:
                        min_len = min(heps_a.shape[-1], heps_b.shape[-1])
                        val_a = heps_a[:, 0, :min_len] if len(heps_a.shape) == 3 else heps_a[:, :min_len]
                        val_b = heps_b[:, 0, :min_len] if len(heps_b.shape) == 3 else heps_b[:, :min_len]
                        
                        # We use a simple t-test at each time point
                        t_stat, p_vals = stats.ttest_ind(val_a, val_b, axis=0, equal_var=False)
                        
                        if HAS_FDR:
                            reject, p_vals_corrected_fdr = fdr_correction(p_vals, alpha=0.05, method='indep')
                            p_vals = p_vals_corrected_fdr
                            
                        min_p = np.nanmin(p_vals)
                        is_sig = min_p < 0.05
                        
                        win_strs = ""
                        clusters = []
                        if is_sig:
                            sig_times = stage_data[stage_a]['times'][:min_len][p_vals < 0.05] * 1000
                            if len(sig_times) > 0:
                                # Find continuous clusters of significance
                                split_indices = np.where(np.diff(sig_times) > 5)[0] + 1 
                                clusters = np.split(sig_times, split_indices)
                                win_strs = ", ".join([f"{c[0]:.0f}–{c[-1]:.0f}ms" for c in clusters if len(c) > 1])

                        pair_pvals[(stage_a, stage_b)] = (min_p, is_sig, win_strs)
                        pair_sig_clusters[(stage_a, stage_b)] = clusters
                    else:
                        pair_pvals[(stage_a, stage_b)] = (None, None, "")  # no data
                        pair_sig_clusters[(stage_a, stage_b)] = []

            # ── Render as HTML matrix table ──────────────────────────────────
            n = len(stages_list)
            cmap = plt.get_cmap('Reds_r')

            # Color scheme
            COLOR_NONSIG = "#f0f0f0"   # gray   – not significant
            COLOR_NODATA = "#ffffff"   # white  – no data (missing)
            COLOR_DIAG   = "#d9d9d9"   # darker gray – diagonal

            html = "<style>table.pairwise{border-collapse:collapse;font-family:monospace;font-size:13px}" \
                   "table.pairwise td,table.pairwise th{border:1px solid #bbb;padding:6px 10px;text-align:center;min-width:90px}" \
                   "table.pairwise th{background:#444;color:white;font-weight:bold}" \
                   "</style><table class='pairwise'>"

            # Header row
            html += "<tr><th></th>"
            for s in stages_list:
                html += f"<th>{s}</th>"
            html += "</tr>"

            for row_s in stages_list:
                html += f"<tr><th>{row_s}</th>"
                for col_s in stages_list:
                    if row_s == col_s:
                        html += f"<td style='background:{COLOR_DIAG};color:#000000;font-weight:bold'>{row_s}</td>"
                    else:
                        # Canonical key order (always smaller index first)
                        key = (row_s, col_s) if (row_s, col_s) in pair_pvals else (col_s, row_s)
                        if key in pair_pvals:
                            min_p, is_sig, wins = pair_pvals[key]
                            if is_sig is None:
                                bg = COLOR_NODATA
                                text_c = "#000000"
                                cell_text = "n/a"
                                title = "No data"
                            else:
                                rgba = cmap(min_p / 0.05) if min_p <= 0.05 else mcolors.to_rgba(COLOR_NONSIG)
                                bg = mcolors.to_hex(rgba)
                                luminance = 0.299 * rgba[0] + 0.587 * rgba[1] + 0.114 * rgba[2]
                                text_c = "#ffffff" if luminance < 0.5 else "#000000"
                                
                                # Show Exact P-Value instead of 'n.s.'
                                p_text = f"p = {min_p:.4f}" if min_p >= 0.0001 else "p < 1e-4"
                                cell_text = f"<b>{p_text}</b><br><small>{wins}</small>" if is_sig else f"{p_text}"
                                title = f"p={min_p:.4f}"
                        else:
                            bg = COLOR_NODATA
                            text_c = "#000000"
                            cell_text = "n/a"
                            title = "No data"
                        html += f"<td style='background:{bg};color:{text_c}' title='{title}'>{cell_text}</td>"
                html += "</tr>"
            html += "</table>"
            st.markdown(html, unsafe_allow_html=True)

            # --- Plot Significant Windows ---
            sig_pairs_to_plot = {k: v for k, v in pair_sig_clusters.items() if len(v) > 0 and any(len(c) > 1 for c in v)}
            if sig_pairs_to_plot:
                st.markdown("#### Significant Difference Windows (ECG HEP)")
                fig_win, ax_win = plt.subplots(figsize=(10, max(3, len(sig_pairs_to_plot) * 0.5)))
                
                y_labels = []
                y_ticks = []
                
                for idx, (pair, clusters) in enumerate(sig_pairs_to_plot.items()):
                    y_pos = len(sig_pairs_to_plot) - idx
                    y_labels.append(f"{pair[0]} vs {pair[1]}")
                    y_ticks.append(y_pos)
                    
                    for cluster in clusters:
                        if len(cluster) > 1:
                            start_time = cluster[0] / 1000.0  # convert back to seconds
                            end_time = cluster[-1] / 1000.0
                            ax_win.plot([start_time, end_time], [y_pos, y_pos], color='black', linewidth=4, solid_capstyle='butt')
                            
                ax_win.set_yticks(y_ticks)
                ax_win.set_yticklabels(y_labels)
                ax_win.set_xlabel("Time (s)")
                ax_win.set_title("Significant Difference Windows (p < 0.05)")
                ax_win.grid(True, axis='x', alpha=0.3)
                ax_win.axvline(0, color='r', linestyle='--', alpha=0.5)
                ax_win.set_xlim(ax_ecg.get_xlim())
                fig_win.tight_layout()
                st.pyplot(fig_win, use_container_width=True)
                plt.close(fig_win)
    else:
        st.write("No ECG data available.")

    # --- Plot EEG using the same channel flow as Single Group Analysis ---
    st.divider()
    st.title("Group HEP Analysis Across Sleep Stages")
    st.subheader("Per-Channel Analysis")

    available_stage_channels = [
        set(data.get('ch_names', []))
        for data in stage_data.values()
        if data.get('ch_names')
    ]
    if not available_stage_channels:
        st.warning("No EEG channels found.")
        return

    common_channels_all_stages = sorted(set.intersection(*available_stage_channels))
    common_channels_all_stages = [
        ch for ch in common_channels_all_stages
        if re.match(r'^[A-Za-z]{1,3}[0-9]+$', ch) or re.match(r'^[A-Za-z]{1,2}z$', ch, re.IGNORECASE)
    ]
    if not common_channels_all_stages:
        st.warning("No common EEG channels found across the displayed sleep stages.")
        return

    display_channels = ['Average', 'Median', 'ECG', 'Left', 'Right', 'Middle'] + common_channels_all_stages

    def _stage_channel_matrix(stage, data, ch_name):
        times = data.get('times')
        if times is None:
            return None, None

        if ch_name == 'ECG':
            if analysis_mode == "Group Average" and data.get('ecg_heps_aligned') is not None:
                ecg_arr = np.asarray(data['ecg_heps_aligned'], dtype=float)
                if ecg_arr.ndim == 3:
                    ecg_arr = ecg_arr[:, 0, :]
                elif ecg_arr.ndim == 2:
                    pass
                else:
                    return None, None
                return ecg_arr, times[:ecg_arr.shape[1]]

            ecg_hep = data.get('ecg_hep')
            if ecg_hep is None:
                return None, None
            ecg_trace = np.asarray(ecg_hep, dtype=float).squeeze()
            if ecg_trace.ndim != 1:
                return None, None
            return ecg_trace[None, :], times[:len(ecg_trace)]

        ch_names = list(data.get('ch_names', []))
        if analysis_mode == "Group Average" and data.get('eeg_heps_aligned') is not None:
            eeg_arr = np.asarray(data['eeg_heps_aligned'], dtype=float)
            if eeg_arr.ndim != 3:
                return None, None

            if ch_name in ('Average', 'Median', 'Left', 'Right', 'Middle'):
                if ch_name in ('Average', 'Median'):
                    selected_chs = [ch for ch in common_channels_all_stages if ch in ch_names]
                else:
                    left_chs, right_chs, mid_chs = get_hemisphere_channels(
                        [ch for ch in common_channels_all_stages if ch in ch_names]
                    )
                    selected_chs = {'Left': left_chs, 'Right': right_chs, 'Middle': mid_chs}[ch_name]
                idx = [ch_names.index(ch) for ch in selected_chs if ch in ch_names]
                if not idx:
                    return None, None
                if ch_name == 'Median':
                    mat = np.nanmedian(eeg_arr[:, idx, :], axis=1)
                else:
                    mat = np.array([
                        summarize_channels_without_reference_cancellation(row[idx, :])
                        for row in eeg_arr
                    ])
                return mat, times[:mat.shape[1]]

            if ch_name not in ch_names:
                return None, None
            ch_idx = ch_names.index(ch_name)
            mat = eeg_arr[:, ch_idx, :]
            return mat, times[:mat.shape[1]]

        eeg_hep = data.get('eeg_hep')
        ecg_hep = data.get('ecg_hep')
        if eeg_hep is None or not ch_names:
            return None, None
        eeg_clean = apply_ecg_regression(np.asarray(eeg_hep, dtype=float), ecg_hep)

        if ch_name in ('Average', 'Median', 'Left', 'Right', 'Middle'):
            if ch_name in ('Average', 'Median'):
                selected_chs = [ch for ch in common_channels_all_stages if ch in ch_names]
            else:
                left_chs, right_chs, mid_chs = get_hemisphere_channels(
                    [ch for ch in common_channels_all_stages if ch in ch_names]
                )
                selected_chs = {'Left': left_chs, 'Right': right_chs, 'Middle': mid_chs}[ch_name]
            idx = [ch_names.index(ch) for ch in selected_chs if ch in ch_names]
            if not idx:
                return None, None
            if ch_name == 'Median':
                trace = np.nanmedian(eeg_clean[idx, :], axis=0)
            else:
                trace = summarize_channels_without_reference_cancellation(eeg_clean[idx, :])
            return trace[None, :], times[:len(trace)]

        if ch_name not in ch_names:
            return None, None
        trace = eeg_clean[ch_names.index(ch_name)]
        return trace[None, :], times[:len(trace)]

    channel_stage_p_values = {}
    for ch_name in display_channels:
        with st.expander(f"Channel: {ch_name}", expanded=(ch_name in ('Average', 'Median'))):
            fig, ax = plt.subplots(figsize=(16, 9))
            plotted_any = False
            channel_stage_p_values[ch_name] = {}

            for stage in sleep_stages:
                data = stage_data.get(stage)
                if not data:
                    continue
                mat, times = _stage_channel_matrix(stage, data, ch_name)
                if mat is None or times is None or mat.size == 0:
                    continue

                min_len = min(mat.shape[1], len(times))
                mat = mat[:, :min_len]
                times = np.asarray(times[:min_len], dtype=float)
                mat_uv = mat * 1e6
                median_trace, mad_trace = median_and_mad(mat_uv, axis=0)
                color = stage_colors.get(stage, 'gray')

                stage_label = (
                    f"{stage} (N={mat.shape[0]})"
                    if analysis_mode == "Group Average"
                    else f"{stage} (epochs={data.get('n_epochs', 'n/a')})"
                )
                ax.plot(times, median_trace, color=color, linewidth=2.2, label=stage_label)
                if mat.shape[0] > 1:
                    ax.fill_between(times, median_trace - mad_trace, median_trace + mad_trace,
                                    color=color, alpha=0.18)

                if analysis_mode == "Group Average" and mat.shape[0] >= 3:
                    try:
                        sig_windows, _, per_pt_info = permutation_cluster_jitter_test(
                            mat,
                            times,
                            n_permutations=100,
                            p_threshold=0.05,
                            jitter_sec=0.1,
                        )
                    except Exception:
                        sig_windows, per_pt_info = [], {}
                    min_p = min((w['p_value'] for w in sig_windows), default=1.0)
                    channel_stage_p_values[ch_name][stage] = min_p
                    if sig_windows:
                        for win in sig_windows:
                            ax.axvspan(win['start'], win['end'], color=color, alpha=0.10)
                    n_sig_pt = per_pt_info.get('n_significant', 0)
                    fisher_p = per_pt_info.get('fisher_p', 1.0)
                    ax.text(
                        0.01,
                        0.98 - 0.05 * len(channel_stage_p_values[ch_name]),
                        f"{stage}: p={min_p:.3f}, {n_sig_pt}/{mat.shape[0]} pts sig, Fisher={fisher_p:.3f}",
                        transform=ax.transAxes,
                        color=color,
                        fontsize=8,
                        va='top',
                    )
                plotted_any = True

            ax.axvline(0, color='red', linestyle='--', alpha=0.5)
            ax.axhline(0, color='black', linewidth=0.5, alpha=0.35)
            ax.set_xlabel("Time (s)")
            ax.set_ylabel("Amplitude (µV)")
            ax.grid(True, alpha=0.25)
            ax.legend(fontsize=9, loc='upper right')
            ax.set_title(
                f"Channel: {ch_name} - Median +/- MAD - Group: {selected_group} - All Sleep Stages - {analysis_mode}",
                fontsize=12,
                fontweight='bold',
            )

            if plotted_any:
                fig.tight_layout()
                st.pyplot(fig, use_container_width=True)
                plt.close(fig)
            else:
                plt.close(fig)
                st.info(f"No data available for {ch_name}.")

    if analysis_mode == "Group Average" and channel_stage_p_values:
        pval_rows = []
        for ch_name, stage_map in channel_stage_p_values.items():
            if not stage_map:
                continue
            row = {"Channel": ch_name}
            row.update({stage: stage_map.get(stage, np.nan) for stage in sleep_stages})
            pval_rows.append(row)
        if pval_rows:
            with st.expander("Per-channel minimum p-values by sleep stage", expanded=False):
                st.dataframe(pd.DataFrame(pval_rows).style.format(precision=4), use_container_width=True)

    # --- Plot Spatial Topomaps ---
    st.divider()
    st.subheader("Spatial Topomaps")
    st.caption("Select a time window to visualize the average HEP amplitude distribution, PSD, and (for Group Average) spatial significance across the scalp.")
    
    col_t1, col_t2 = st.columns(2)
    with col_t1:
        topo_tmin = st.number_input("Start Time (ms)", min_value=-500, max_value=1000, value=200, step=10, key="topo_tmin_compare")
    with col_t2:
        topo_tmax = st.number_input("End Time (ms)", min_value=-500, max_value=1000, value=400, step=10, key="topo_tmax_compare")
        
    if topo_tmin >= topo_tmax:
        st.warning("Start time must be less than end time.")
    else:
        from mne.time_frequency import psd_array_welch
        montage = mne.channels.make_standard_montage('standard_1020')
        montage_ch_names_upper = [ch.upper() for ch in montage.ch_names]
        
        # Pre-compute global min/max for the same color scale
        global_max_abs = 0.0
        global_psd_max = 0.0
        stage_plot_data = {}
        valid_times = False
        
        for stage, data in stage_data.items():
            if data['eeg_hep'] is None:
                continue
            times_ms = data['times'] * 1000
            t_mask = (times_ms >= topo_tmin) & (times_ms <= topo_tmax)
            if not np.any(t_mask):
                continue
            valid_times = True
            
            # 1. Median Amplitude
            mean_amps = np.nanmedian(data['eeg_hep'][:, t_mask], axis=1) * 1e6
            
            # 2. PSD (Total Power)
            hep_windowed = data['eeg_hep'][:, t_mask]
            sfreq = 1000.0 / np.mean(np.diff(times_ms)) if len(times_ms) > 1 else 250.0
            try:
                n_fft = min(256, hep_windowed.shape[1])
                if n_fft > 0:
                    psds, freqs = psd_array_welch(hep_windowed, sfreq=sfreq, fmin=0.5, fmax=40.0, n_fft=n_fft, verbose=False)
                    psd_total = np.sum(psds, axis=1)
                else:
                    psd_total = np.zeros(hep_windowed.shape[0])
            except Exception:
                psd_total = np.zeros(hep_windowed.shape[0])
                
            # 3. Jitter P-value (Group Average only)
            p_values = None
            if analysis_mode == "Group Average" and data.get('eeg_heps_aligned') is not None:
                p_values = []
                eeg_aligned = data.get('eeg_heps_aligned')
                eeg_win = eeg_aligned[:, :, t_mask]
                times_win = data['times'][t_mask]
                
                for ch_idx in range(eeg_win.shape[1]):
                    channel_data = eeg_win[:, ch_idx, :]
                    sig_windows, _, _pt = permutation_cluster_jitter_test(channel_data, times_win, n_permutations=100, p_threshold=0.05, jitter_sec=0.1)
                    if sig_windows:
                        min_p = min([w['p_value'] for w in sig_windows])
                        p_values.append(min_p)
                    else:
                        p_values.append(1.0)  # Non-significant (will be plotted as white)

            plot_ch_names = []
            plot_data_amp = []
            plot_data_psd = []
            plot_data_pval = [] if p_values is not None else None
            
            for i, ch in enumerate(data['ch_names']):
                ch_upper = ch.upper()
                if ch_upper in montage_ch_names_upper:
                    m_idx = montage_ch_names_upper.index(ch_upper)
                    standard_name = montage.ch_names[m_idx]
                    plot_ch_names.append(standard_name)
                    plot_data_amp.append(mean_amps[i])
                    plot_data_psd.append(psd_total[i])
                    if plot_data_pval is not None:
                        plot_data_pval.append(p_values[i])
                    
            # Pad missing 10-20 standard channels to keep head shape normal
            standard_19_base = ['Fp1', 'Fp2', 'F7', 'F3', 'Fz', 'F4', 'F8', 'C3', 'Cz', 'C4', 'P3', 'Pz', 'P4', 'O1', 'O2']
            aliases = {'T7': 'T3', 'T8': 'T4', 'P7': 'T5', 'P8': 'T6'}
            
            def pad_channel_if_missing(ch_name, pad_amp, pad_psd, pad_pval):
                if not any(ch_name.upper() == p_ch.upper() for p_ch in plot_ch_names):
                    plot_ch_names.append(ch_name)
                    plot_data_amp.append(pad_amp)
                    plot_data_psd.append(pad_psd)
                    if plot_data_pval is not None:
                        plot_data_pval.append(pad_pval)
            
            for base_ch in standard_19_base:
                pad_channel_if_missing(base_ch, 0.0, 0.0, 0.05)
                    
            for new_name, old_name in aliases.items():
                if not any(ch.upper() in [new_name.upper(), old_name.upper()] for ch in plot_ch_names):
                    pad_channel_if_missing(new_name, 0.0, 0.0, 0.05)

            if plot_data_amp:
                curr_max = np.max(np.abs(plot_data_amp))
                if curr_max > global_max_abs:
                    global_max_abs = curr_max
            if plot_data_psd:
                curr_psd_max = np.max(plot_data_psd)
                if curr_psd_max > global_psd_max:
                    global_psd_max = curr_psd_max
            
            stage_plot_data[stage] = {
                'names': plot_ch_names,
                'amp': plot_data_amp,
                'psd': plot_data_psd,
                'pval': plot_data_pval
            }
            
        if not valid_times:
            st.warning("Invalid time window selected (no data points).")
        elif not stage_plot_data:
            st.warning("No channels matched the standard montage.")
        else:
            n_stages = len(stage_plot_data)
            stages_list = list(stage_plot_data.keys())
            
            vmin_amp, vmax_amp = -global_max_abs, global_max_abs
            if vmax_amp == 0: vmax_amp, vmin_amp = 1, -1
            
            vmin_psd, vmax_psd = 0, global_psd_max
            if vmax_psd == 0: vmax_psd = 1

            # Helper: build info object for a channel list
            def make_info(names):
                info = mne.create_info(ch_names=list(names), sfreq=250., ch_types='eeg')
                info.set_montage(montage, on_missing='ignore')
                valid = np.array([
                    not np.any(np.isnan(ch['loc'][:3])) and np.any(ch['loc'][:3] != 0)
                    for ch in info['chs']
                ])
                info = mne.pick_info(info, np.where(valid)[0])
                return info, valid

            # ── Figure 1: Median Amplitude per stage ─────────────────────
            st.markdown("#### Median HEP Amplitude")
            fig_amp, axes_amp = plt.subplots(1, n_stages, figsize=(4 * n_stages + 1, 4))
            if n_stages == 1: axes_amp = [axes_amp]
            im_amp = None
            
            # Prepare summary stats
            amp_stats = []
            
            for idx, stage in enumerate(stages_list):
                p = stage_plot_data[stage]
                info_, _valid_mask = make_info(p['names'])
                r = mne.viz.plot_topomap(np.array(p['amp'])[_valid_mask], info_, axes=axes_amp[idx],
                                         cmap='RdBu_r', names=np.array(p['names'])[_valid_mask].tolist(),
                                         vlim=(vmin_amp, vmax_amp), extrapolate='head', show=False)
                im_amp = r[0] if isinstance(r, tuple) else r
                axes_amp[idx].set_title(stage)
                
                # Compute stats
                stage_n = stage_data[stage]['n_subjects']
                amp_vals = np.array(p['amp'])
                amp_stats.append({
                    'Stage': stage,
                    'N': stage_n,
                    'Min (µV)': np.min(amp_vals),
                    'Max (µV)': np.max(amp_vals),
                    'Std (µV)': np.std(amp_vals)
                })
                
            fig_amp.subplots_adjust(right=0.88)
            if im_amp is not None:
                cbar_ax = fig_amp.add_axes([0.91, 0.15, 0.02, 0.7])
                fig_amp.colorbar(im_amp, cax=cbar_ax).set_label("Median Amplitude (µV)")
            st.pyplot(fig_amp, use_container_width=False)
            plt.close(fig_amp)
            
            # Display Stats & Expandable Table for Amp
            st.markdown("**Summary Statistics (Amplitude)**")
            for stat in amp_stats:
                st.write(f"- **{stat['Stage']}** (N={stat['N']}): Min={stat['Min (µV)']:.3f}, Max={stat['Max (µV)']:.3f}, SD={stat['Std (µV)']:.3f}")
            with st.expander("Raw Amplitude Values by Channel"):
                amp_df = pd.DataFrame({stage: stage_plot_data[stage]['amp'] for stage in stages_list}, index=stage_plot_data[stages_list[0]]['names'])
                st.dataframe(amp_df.style.format("{:.3f}"))

            # ── Figure 2: P-value per stage (Group Average only) ─────────
            if analysis_mode == "Group Average":
                has_pvals = any(stage_plot_data[s]['pval'] is not None for s in stages_list)
                if has_pvals:
                    st.markdown("#### Spatial Significance (Min. Cluster P-value per Stage)")
                    fig_pval, axes_pval = plt.subplots(1, n_stages, figsize=(4 * n_stages + 1, 4))
                    if n_stages == 1: axes_pval = [axes_pval]
                    im_pval = None
                    
                    pval_stats = []
                    
                    for idx, stage in enumerate(stages_list):
                        p = stage_plot_data[stage]
                        if p['pval'] is not None:
                            info_, _valid_mask = make_info(p['names'])
                            data_pv = np.clip(np.array(p['pval']), 0, 0.05)[_valid_mask]
                            r = mne.viz.plot_topomap(data_pv, info_, axes=axes_pval[idx],
                                                     cmap='Reds_r', names=np.array(p['names'])[_valid_mask].tolist(),
                                                     vlim=(0, 0.05), extrapolate='head', show=False)
                            im_pval = r[0] if isinstance(r, tuple) else r
                            
                            pval_vals = np.array(p['pval'])
                            pval_stats.append({
                                'Stage': stage,
                                'N': stage_data[stage]['n_subjects'],
                                'Min P': np.min(pval_vals),
                                'Max P': np.max(pval_vals),
                                'Sig Channels Count': np.sum(pval_vals < 0.05)
                            })
                        else:
                            axes_pval[idx].axis('off')
                        axes_pval[idx].set_title(stage)
                    fig_pval.subplots_adjust(right=0.88)
                    if im_pval is not None:
                        cbar_ax = fig_pval.add_axes([0.91, 0.15, 0.02, 0.7])
                        fig_pval.colorbar(im_pval, cax=cbar_ax).set_label("p-value")
                    st.pyplot(fig_pval, use_container_width=False)
                    plt.close(fig_pval)
                    
                    # Display Stats & Expandable Table for P-vals
                    if pval_stats:
                        st.markdown("**Summary Statistics (P-values)**")
                        for stat in pval_stats:
                            st.write(f"- **{stat['Stage']}** (N={stat['N']}): Min={stat['Min P']:.4f}, Max={stat['Max P']:.4f}, Significant Channels (p<0.05)={stat['Sig Channels Count']}")
                        with st.expander("Raw P-Values by Channel"):
                            p_columns = {stage: stage_plot_data[stage]['pval'] for stage in stages_list if stage_plot_data[stage]['pval'] is not None}
                            if p_columns:
                                pval_df = pd.DataFrame(p_columns, index=stage_plot_data[stages_list[0]]['names'])
                                st.dataframe(pval_df.style.format("{:.4f}"))

            # ── Figure 3: Pairwise stage differences ─────────────────────
            pairs = [(stages_list[i], stages_list[j])
                     for i in range(len(stages_list))
                     for j in range(i + 1, len(stages_list))]

            if pairs:
                st.markdown("#### Pairwise Stage Difference (Median Amplitude, µV)")
                n_pairs = len(pairs)
                n_cols = min(3, n_pairs)
                n_rows = int(np.ceil(n_pairs / n_cols))
                fig_diff, axes_diff = plt.subplots(n_rows, n_cols, figsize=(4 * n_cols + 1, 4 * n_rows))
                if n_pairs == 1: axes_diff = [axes_diff]
                else: axes_diff = axes_diff.flatten()

                # Compute all diff arrays to find global max for a shared scale
                diff_data_list = []
                pair_names_list = []
                for (s_a, s_b) in pairs:
                    p_a = stage_plot_data[s_a]
                    p_b = stage_plot_data[s_b]
                    # Match channels — use the union of names (pad missing with 0)
                    all_names_pair = list(dict.fromkeys(p_a['names'] + p_b['names']))
                    amp_a_map = dict(zip(p_a['names'], p_a['amp']))
                    amp_b_map = dict(zip(p_b['names'], p_b['amp']))
                    diff_arr = np.array([
                        amp_a_map.get(ch, 0.0) - amp_b_map.get(ch, 0.0)
                        for ch in all_names_pair
                    ])
                    diff_data_list.append((all_names_pair, diff_arr))
                    pair_names_list.append((s_a, s_b))

                global_diff_max = max(np.max(np.abs(d)) for _, d in diff_data_list) or 1.0

                im_diff = None
                diff_stats = []
                diff_df_dict = {}
                
                for idx, ((s_a, s_b), (pair_names, diff_arr)) in enumerate(
                        zip(pair_names_list, diff_data_list)):
                    info_, _valid_mask = make_info(pair_names)
                    r = mne.viz.plot_topomap(diff_arr[_valid_mask], info_, axes=axes_diff[idx],
                                             cmap='RdBu_r', names=np.array(pair_names)[_valid_mask].tolist(),
                                             vlim=(-global_diff_max, global_diff_max),
                                             extrapolate='head', show=False)
                    im_diff = r[0] if isinstance(r, tuple) else r
                    label_ab = f"{s_a} − {s_b}"
                    axes_diff[idx].set_title(label_ab)
                    
                    diff_stats.append({
                        'Pair': label_ab,
                        'Min Diff': np.min(diff_arr),
                        'Max Diff': np.max(diff_arr),
                        'Std Diff': np.std(diff_arr)
                    })
                    diff_df_dict[label_ab] = dict(zip(pair_names, diff_arr))

                # Hide unused axes
                for j in range(len(pair_names_list), len(axes_diff)):
                    axes_diff[j].axis('off')

                fig_diff.subplots_adjust(right=0.88)
                if im_diff is not None:
                    cbar_ax = fig_diff.add_axes([0.91, 0.15, 0.02, 0.7])
                    fig_diff.colorbar(im_diff, cax=cbar_ax).set_label("Difference (µV)")
                st.pyplot(fig_diff, use_container_width=False)
                plt.close(fig_diff)
                
                # Display Stats & Expandable Table for Pairwise Diff
                st.markdown("**Summary Statistics (Pairwise Differences)**")
                for stat in diff_stats:
                    st.write(f"- **{stat['Pair']}**: Min={stat['Min Diff']:.3f}, Max={stat['Max Diff']:.3f}, SD={stat['Std Diff']:.3f}")
                with st.expander("Raw Pairwise Differences by Channel"):
                    # Ensure we build a dataframe across all possible channel names
                    diff_df = pd.DataFrame(diff_df_dict)
                st.dataframe(diff_df.style.format("{:.3f}"))

def run_compare_groups_non_eeg_analysis(base_path, selected_stage):
    """
    Compare ECG-aligned non-EEG channels between two user-selected groups.
    Uses the same ECG cleaning / R-peak alignment path as the HEP workflow,
    with optional ICA applied before extracting the aligned averages.
    """
    available_groups = sorted([g for g in os.listdir(base_path) if os.path.isdir(os.path.join(base_path, g))])
    if len(available_groups) < 2:
        st.error("At least two groups are required for non-EEG comparison.")
        return

    default_groups = default_hep_groups(available_groups)
    if len(default_groups) < 2:
        default_groups = available_groups[:2]

    selected_groups = st.multiselect(
        "Select Two Groups to Compare",
        options=available_groups,
        default=default_groups[:2],
        key="cmp_non_eeg_selected_groups",
    )
    if len(selected_groups) != 2:
        st.warning("Please select exactly two groups for the non-EEG comparison.")
        return

    group_a, group_b = selected_groups

    with st.expander("⚙️ Non-EEG Analysis Settings", expanded=True):
        col1, col2, col3, col4, col5 = st.columns(5)
        with col1:
            test_run = st.checkbox("Test Run (10 files/group)", value=False, key="cmp_non_eeg_test_run")
        with col2:
            n_permutations = st.slider("Permutations", 50, 500, 200, 50, key="cmp_non_eeg_n_perm")
        with col3:
            jitter_sec = st.number_input("Jitter (s)", 0.01, 0.5, 0.1, 0.05, key="cmp_non_eeg_jitter")
        with col4:
            use_zscore = st.checkbox(
                "Z-score per subject/channel",
                value=False,
                key="cmp_non_eeg_zscore",
                help="Useful when channels have very different raw scales across patients.",
            )
        with col5:
            apply_ica = st.checkbox(
                "Apply ICA before alignment",
                value=True,
                key="cmp_non_eeg_apply_ica",
                help="Runs the same ICA cleaning path used in the HEP comparison workflow.",
            )

    recompute_cache = st.button(
        "Recompute Non-EEG Cache",
        key="cmp_non_eeg_recompute_cache",
        help="Force reprocessing of all patient data for the non-EEG comparison.",
    )
    if recompute_cache:
        if hasattr(get_group_non_eeg_individuals, "clear"):
            get_group_non_eeg_individuals.clear()

    group_individuals = {}
    for group in selected_groups:
        with st.spinner(f"Loading non-EEG aligned data for {group}…"):
            inds = get_group_non_eeg_individuals(
                group,
                selected_stage,
                base_path,
                test_run=test_run,
                recompute_cache=recompute_cache,
                apply_ica=apply_ica,
            )
        if inds:
            group_individuals[group] = inds
        else:
            st.warning(f"No valid non-EEG aligned data for group **{group}** in stage {selected_stage}.")

    if len(group_individuals) < 2:
        st.error("Both groups need valid non-EEG aligned data.")
        return

    group_channel_counts = {}
    for group, inds in group_individuals.items():
        counts = Counter()
        for _, _, _, ch_names, _, _ in inds:
            counts.update(ch_names)
        group_channel_counts[group] = counts

    common_channels = sorted(
        ch for ch in set(group_channel_counts[group_a]) & set(group_channel_counts[group_b])
        if group_channel_counts[group_a][ch] >= max(1, math.ceil(0.3 * len(group_individuals[group_a])))
        and group_channel_counts[group_b][ch] >= max(1, math.ceil(0.3 * len(group_individuals[group_b])))
    )
    if not common_channels:
        st.error("No common non-EEG channels were found across the two selected groups.")
        return

    selected_channels = st.multiselect(
        "Select Non-EEG Channels",
        options=common_channels,
        default=common_channels[:min(8, len(common_channels))],
        key="cmp_non_eeg_channels",
    )
    if not selected_channels:
        st.warning("Please select at least one non-EEG channel.")
        return

    group_channel_data = {group: {} for group in selected_groups}
    times = None
    for group, inds in group_individuals.items():
        for ch in selected_channels:
            traces = []
            pids = []
            for pid, aligned_data, ind_times, ch_names, _, _ in inds:
                if ch in ch_names:
                    ch_idx = ch_names.index(ch)
                    trace = np.asarray(aligned_data[ch_idx], dtype=float)
                    if trace.ndim == 1 and np.isfinite(trace).any():
                        traces.append(trace)
                        pids.append(pid)
                        if times is None and ind_times is not None:
                            times = np.asarray(ind_times)
            if traces:
                mat = np.vstack(traces)
                mat = scale_matrix(mat, use_zscore)
                group_channel_data[group][ch] = {"matrix": mat, "pids": pids}

    valid_channels = [
        ch for ch in selected_channels
        if ch in group_channel_data[group_a] and ch in group_channel_data[group_b]
        and group_channel_data[group_a][ch]["matrix"].shape[0] >= 2
        and group_channel_data[group_b][ch]["matrix"].shape[0] >= 2
    ]
    if not valid_channels:
        st.error("No selected non-EEG channels had at least two patients in each group.")
        return

    st.markdown("---")
    st.subheader("Interactive Non-EEG Channel Comparison")
    st.caption(
        "Each panel shows ECG-aligned averages for a non-EEG channel after the same ECG cleaning and optional ICA path used in the HEP analysis."
    )

    try:
        from plotly.subplots import make_subplots
        import plotly.graph_objects as go
    except ImportError:
        st.error("Plotly is required for this mode. Please install `plotly`.")
        return

    amp_unit = "Z-score" if use_zscore else "µV"
    summary_rows = []

    for ch in valid_channels:
        mat_a = group_channel_data[group_a][ch]["matrix"]
        mat_b = group_channel_data[group_b][ch]["matrix"]
        n_a = mat_a.shape[0]
        n_b = mat_b.shape[0]

        sig_windows, t_obs, cohens_d = permutation_two_group_cluster_test(
            mat_a,
            mat_b,
            times,
            n_permutations=n_permutations,
            p_threshold=0.05,
            jitter_sec=jitter_sec,
            label_a=group_a,
            label_b=group_b,
            channel_label=ch,
            button_key=f"non_eeg_explain_{selected_stage}_{ch}",
        )

        med_a, mad_a = median_and_mad(mat_a, axis=0)
        med_b, mad_b = median_and_mad(mat_b, axis=0)
        diff = med_a - med_b
        _, point_p = stats.ttest_ind(mat_a, mat_b, axis=0, equal_var=False, nan_policy='omit')
        peak_idx = int(np.nanargmax(np.abs(diff)))
        min_cluster_p = min((w['p_value'] for w in sig_windows), default=np.nan)

        summary_rows.append({
            "channel": ch,
            f"{group_a}_n": n_a,
            f"{group_b}_n": n_b,
            "peak_diff_time_ms": times[peak_idx] * 1000.0,
            f"{group_a}_median_at_peak": med_a[peak_idx],
            f"{group_b}_median_at_peak": med_b[peak_idx],
            "peak_diff": diff[peak_idx],
            "peak_t": t_obs[peak_idx],
            "peak_uncorrected_p": point_p[peak_idx],
            "min_cluster_p": min_cluster_p,
            "significant_clusters": len(sig_windows),
        })

        fig = make_subplots(
            rows=2,
            cols=1,
            shared_xaxes=True,
            vertical_spacing=0.08,
            row_heights=[0.7, 0.3],
            subplot_titles=(
                f"{ch}: {group_a} vs {group_b}",
                "Difference and Welch t-statistic",
            ),
        )

        color_a = "#1f77b4"
        color_b = "#d62728"
        fill_a = "rgba(31,119,180,0.18)"
        fill_b = "rgba(214,39,40,0.18)"

        fig.add_trace(
            go.Scatter(
                x=np.concatenate([times, times[::-1]]),
                y=np.concatenate([med_a + mad_a, (med_a - mad_a)[::-1]]),
                fill="toself",
                fillcolor=fill_a,
                line=dict(color="rgba(0,0,0,0)"),
                hoverinfo="skip",
                name=f"{group_a} ± MAD",
            ),
            row=1, col=1,
        )
        fig.add_trace(
            go.Scatter(
                x=times,
                y=med_a,
                mode="lines",
                line=dict(color=color_a, width=3),
                name=f"{group_a} median (n={n_a})",
            ),
            row=1, col=1,
        )
        fig.add_trace(
            go.Scatter(
                x=np.concatenate([times, times[::-1]]),
                y=np.concatenate([med_b + mad_b, (med_b - mad_b)[::-1]]),
                fill="toself",
                fillcolor=fill_b,
                line=dict(color="rgba(0,0,0,0)"),
                hoverinfo="skip",
                name=f"{group_b} ± MAD",
            ),
            row=1, col=1,
        )
        fig.add_trace(
            go.Scatter(
                x=times,
                y=med_b,
                mode="lines",
                line=dict(color=color_b, width=3),
                name=f"{group_b} median (n={n_b})",
            ),
            row=1, col=1,
        )
        fig.add_trace(
            go.Scatter(
                x=times,
                y=diff,
                mode="lines",
                line=dict(color="#6a3d9a", width=2.5),
                name=f"{group_a} - {group_b}",
            ),
            row=2, col=1,
        )
        fig.add_trace(
            go.Scatter(
                x=times,
                y=t_obs,
                mode="lines",
                line=dict(color="#444444", width=1.5, dash="dot"),
                name="Welch t",
            ),
            row=2, col=1,
        )

        for win in sig_windows:
            fig.add_vrect(
                x0=win['start'],
                x1=win['end'],
                fillcolor="rgba(255,165,0,0.18)",
                line_width=0,
                annotation_text=f"p={win['p_value']:.3f}",
                annotation_position="top left",
                row=1,
                col=1,
            )
            fig.add_vrect(
                x0=win['start'],
                x1=win['end'],
                fillcolor="rgba(255,165,0,0.18)",
                line_width=0,
                row=2,
                col=1,
            )

        fig.add_vline(x=0, line_dash="dash", line_color="gray", row=1, col=1)
        fig.add_vline(x=0, line_dash="dash", line_color="gray", row=2, col=1)
        fig.add_hline(y=0, line_color="gray", line_width=1, row=1, col=1)
        fig.add_hline(y=0, line_color="gray", line_width=1, row=2, col=1)

        title_text = (
            f"ECG-Aligned Non-EEG Comparison - Median +/- MAD | {ch} | {group_a} vs {group_b} | {selected_stage} | "
            f"min cluster p={min_cluster_p:.3g}"
            if np.isfinite(min_cluster_p) else
            f"ECG-Aligned Non-EEG Comparison - Median +/- MAD | {ch} | {group_a} vs {group_b} | {selected_stage} | no significant clusters"
        )
        fig.update_layout(
            height=650,
            hovermode="x unified",
            title=dict(text=title_text, font=dict(size=14)),
            legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
        )
        fig.update_xaxes(title_text="Time relative to R-peak (s)", row=2, col=1)
        fig.update_yaxes(title_text=f"Amplitude ({amp_unit})", row=1, col=1)
        fig.update_yaxes(title_text=f"Difference ({amp_unit})", row=2, col=1)

        st.plotly_chart(fig, use_container_width=True)

        if sig_windows:
            st.success(
                ", ".join(
                    [
                        f"{ch}: {w['start']*1000:.0f}-{w['end']*1000:.0f} ms, p={w['p_value']:.4f}, {w['direction']}"
                        for w in sig_windows
                    ]
                )
            )
        else:
            st.info(f"{ch}: no significant cluster-corrected group difference at p < 0.05.")

    if summary_rows:
        st.markdown("---")
        st.subheader("Statistical Summary")
        summary_df = pd.DataFrame(summary_rows).sort_values(
            by=["min_cluster_p", "peak_uncorrected_p"],
            na_position="last",
        )
        formatters = {
            "peak_diff_time_ms": "{:.1f}",
            f"{group_a}_median_at_peak": "{:.3f}",
            f"{group_b}_median_at_peak": "{:.3f}",
            "peak_diff": "{:.3f}",
            "peak_t": "{:.3f}",
            "peak_uncorrected_p": "{:.4g}",
            "min_cluster_p": "{:.4g}",
        }
        st.dataframe(summary_df.style.format(formatters), use_container_width=True)
        st.caption(
            "Cluster p-values come from the same jittered cluster-permutation test used elsewhere in this dashboard. "
            "Peak uncorrected p-values are provided as descriptive point-wise Welch tests."
        )


def _run_waveform_autoencoder_embedding(matrix, latent_dim=2, max_epochs=120, random_state=42):
    """Return 2D latent coordinates, reconstruction error, and method label for waveform rows."""
    X = np.asarray(matrix, dtype=float)
    if X.ndim != 2 or X.shape[0] < 2 or X.shape[1] < 2:
        return None, None, "Unavailable"

    X = np.nan_to_num(X, nan=np.nanmedian(X), posinf=0.0, neginf=0.0)
    row_mu = np.mean(X, axis=1, keepdims=True)
    row_sd = np.std(X, axis=1, keepdims=True)
    Xn = (X - row_mu) / np.where(row_sd <= 1e-12, 1.0, row_sd)

    latent_dim = int(max(1, min(latent_dim, Xn.shape[0] - 1, Xn.shape[1])))
    if TORCH_AVAILABLE and Xn.shape[0] >= 4:
        try:
            torch.manual_seed(random_state)
            device = torch.device("cpu")
            xt = torch.tensor(Xn, dtype=torch.float32, device=device)
            hidden = int(min(64, max(8, Xn.shape[1] // 4)))
            model = nn.Sequential(
                nn.Linear(Xn.shape[1], hidden),
                nn.ReLU(),
                nn.Linear(hidden, latent_dim),
                nn.Linear(latent_dim, hidden),
                nn.ReLU(),
                nn.Linear(hidden, Xn.shape[1]),
            ).to(device)
            opt = torch.optim.Adam(model.parameters(), lr=0.01, weight_decay=1e-4)
            for _ in range(int(max_epochs)):
                opt.zero_grad()
                recon = model(xt)
                loss = torch.mean((recon - xt) ** 2)
                loss.backward()
                opt.step()

            encoder = model[:3]
            with torch.no_grad():
                z = encoder(xt).cpu().numpy()
                recon = model(xt).cpu().numpy()
            err = np.mean((recon - Xn) ** 2, axis=1)
            if z.shape[1] == 1:
                z = np.column_stack([z[:, 0], np.zeros(z.shape[0])])
            return z[:, :2], err, "Neural autoencoder"
        except Exception:
            pass

    # Linear autoencoder fallback: rank-2 SVD reconstruction.
    Xc = Xn - np.mean(Xn, axis=0, keepdims=True)
    u, s, vt = np.linalg.svd(Xc, full_matrices=False)
    z = u[:, :latent_dim] * s[:latent_dim]
    recon = (u[:, :latent_dim] * s[:latent_dim]) @ vt[:latent_dim, :] + np.mean(Xn, axis=0, keepdims=True)
    err = np.mean((recon - Xn) ** 2, axis=1)
    if z.shape[1] == 1:
        z = np.column_stack([z[:, 0], np.zeros(z.shape[0])])
    return z[:, :2], err, "Linear autoencoder (SVD)"


def _latent_group_permutation_test(z, labels, n_permutations=1000, random_state=42):
    """Permutation test for group separation in autoencoder latent space."""
    z = np.asarray(z, dtype=float)
    labels = np.asarray(labels)
    if z.ndim != 2 or z.shape[0] != len(labels) or z.shape[0] < 4:
        return np.nan, np.nan

    finite = np.all(np.isfinite(z), axis=1)
    z = z[finite]
    labels = labels[finite]
    groups = np.unique(labels)
    if len(groups) < 2 or any(np.sum(labels == g) < 2 for g in groups):
        return np.nan, np.nan

    def _between_group_stat(curr_labels):
        grand = np.mean(z, axis=0)
        stat = 0.0
        for group in groups:
            group_z = z[curr_labels == group]
            if len(group_z) == 0:
                continue
            diff = np.mean(group_z, axis=0) - grand
            stat += len(group_z) * float(np.dot(diff, diff))
        return stat

    obs = _between_group_stat(labels)
    rng = np.random.default_rng(random_state)
    null = np.zeros(int(n_permutations), dtype=float)
    for i in range(int(n_permutations)):
        null[i] = _between_group_stat(rng.permutation(labels))
    p_val = (np.sum(null >= obs) + 1) / (len(null) + 1)
    return float(obs), float(p_val)


def _permanova_latent_test(z, labels, n_permutations=1000, random_state=42):
    """One-way PERMANOVA on Euclidean distances in autoencoder latent space."""
    z = np.asarray(z, dtype=float)
    labels = np.asarray(labels)
    if z.ndim != 2 or z.shape[0] != len(labels) or z.shape[0] < 4:
        return np.nan, np.nan, {}

    finite = np.all(np.isfinite(z), axis=1)
    z = z[finite]
    labels = labels[finite]
    groups = np.unique(labels)
    n = len(labels)
    k = len(groups)
    if k < 2 or n <= k or any(np.sum(labels == g) < 2 for g in groups):
        return np.nan, np.nan, {}

    def _pseudo_f(curr_labels):
        grand = np.mean(z, axis=0)
        ss_between = 0.0
        ss_within = 0.0
        for group in groups:
            group_z = z[curr_labels == group]
            if len(group_z) == 0:
                continue
            centroid = np.mean(group_z, axis=0)
            ss_between += len(group_z) * float(np.sum((centroid - grand) ** 2))
            ss_within += float(np.sum((group_z - centroid) ** 2))
        if ss_within <= 1e-18:
            return np.inf
        return (ss_between / (k - 1)) / (ss_within / (n - k))

    obs_f = _pseudo_f(labels)
    rng = np.random.default_rng(random_state)
    null = np.zeros(int(n_permutations), dtype=float)
    for i in range(int(n_permutations)):
        null[i] = _pseudo_f(rng.permutation(labels))
    p_val = (np.sum(null >= obs_f) + 1) / (len(null) + 1)

    centroids = {group: np.mean(z[labels == group], axis=0) for group in groups}
    dist_to_centroid = np.array([
        float(np.linalg.norm(point - centroids[label]))
        for point, label in zip(z, labels)
    ])
    info = {
        "labels": labels,
        "z": z,
        "centroids": centroids,
        "dist_to_centroid": dist_to_centroid,
    }
    return float(obs_f), float(p_val), info


def _permdisp_latent_test(z, labels, n_permutations=1000, random_state=42):
    """Permutation test for equal group dispersion in autoencoder latent space."""
    z = np.asarray(z, dtype=float)
    labels = np.asarray(labels)
    if z.ndim != 2 or z.shape[0] != len(labels) or z.shape[0] < 4:
        return np.nan, np.nan

    finite = np.all(np.isfinite(z), axis=1)
    z = z[finite]
    labels = labels[finite]
    groups = np.unique(labels)
    n = len(labels)
    k = len(groups)
    if k < 2 or n <= k or any(np.sum(labels == g) < 2 for g in groups):
        return np.nan, np.nan

    def _dispersion_f(curr_labels):
        dists = []
        dist_labels = []
        for group in groups:
            group_z = z[curr_labels == group]
            if len(group_z) < 2:
                continue
            centroid = np.mean(group_z, axis=0)
            group_dist = np.linalg.norm(group_z - centroid, axis=1)
            dists.append(group_dist)
            dist_labels.extend([group] * len(group_dist))
        if len(dists) < 2:
            return np.nan

        dists = np.concatenate(dists)
        dist_labels = np.asarray(dist_labels)
        grand = np.mean(dists)
        ss_between = 0.0
        ss_within = 0.0
        for group in groups:
            group_dist = dists[dist_labels == group]
            if len(group_dist) == 0:
                continue
            group_mean = np.mean(group_dist)
            ss_between += len(group_dist) * float((group_mean - grand) ** 2)
            ss_within += float(np.sum((group_dist - group_mean) ** 2))
        if ss_within <= 1e-18:
            return np.inf if ss_between > 1e-18 else 0.0
        return (ss_between / (k - 1)) / (ss_within / (n - k))

    obs_f = _dispersion_f(labels)
    if not np.isfinite(obs_f):
        return float(obs_f), np.nan

    rng = np.random.default_rng(random_state)
    null = np.zeros(int(n_permutations), dtype=float)
    for i in range(int(n_permutations)):
        null[i] = _dispersion_f(rng.permutation(labels))
    null = null[np.isfinite(null)]
    if len(null) == 0:
        return float(obs_f), np.nan
    p_val = (np.sum(null >= obs_f) + 1) / (len(null) + 1)
    return float(obs_f), float(p_val)


def _format_p_value(p_val):
    if p_val is None or not np.isfinite(p_val):
        return "n/a"
    return "p < 0.001" if p_val < 0.001 else f"p = {p_val:.3f}"


def _format_sig_windows(sig_windows):
    if not sig_windows:
        return "No cluster-corrected significant time window was found."
    return "Significant time window(s): " + ", ".join(
        f"{win['start'] * 1000:.0f}-{win['end'] * 1000:.0f} ms ({_format_p_value(win['p_value'])}, {win.get('direction', '')})"
        for win in sig_windows
    )


def _pcgc_figure_explanation(
    *,
    groups,
    unit,
    latent_p,
    min_cluster_p=None,
    sig_windows=None,
    include_ecg=False,
    all_channels=False,
):
    line_text = (
        "Solid colored lines show the group median HEP waveform; the translucent band around each line is median absolute deviation (MAD). "
        "The vertical dashed line marks the R-peak at t=0."
    )
    if include_ecg:
        line_text += " In the HEP + ECG panel, dotted colored lines show the median ECG waveform on the right y-axis."

    stat_text = (
        f"Autoencoder latent feature separation across {', '.join(groups)}: {_format_p_value(latent_p)}. "
        "This p-value is from a permutation test of between-group centroid separation in the 2D latent space."
    )
    if all_channels:
        return f"{line_text} {stat_text}"

    cluster_text = (
        f"Waveform group-difference cluster p-value: {_format_p_value(min_cluster_p)}. "
        f"{_format_sig_windows(sig_windows)}"
    )
    return f"{line_text} The purple line is the median waveform difference for the first two groups, and the dotted gray line is the Welch t-statistic. {cluster_text} {stat_text}"


def _pcgc_permanova_explanation(permanova_f, permanova_p, permdisp_f=None, permdisp_p=None):
    permdisp_text = ""
    if permdisp_f is not None:
        permdisp_text = (
            f" PERMDISP: F = {permdisp_f:.3f}, {_format_p_value(permdisp_p)}; "
            "this tests whether groups differ in spread around their own centroid."
            if np.isfinite(permdisp_f)
            else " PERMDISP could not be computed for this channel."
        )
    return (
        f"PERMANOVA on autoencoder latent features: pseudo-F = {permanova_f:.3f}, {_format_p_value(permanova_p)}. "
        "This tests whether the groups occupy different regions of the latent feature space using Euclidean distances. "
        "Each point is one subject-channel waveform; the y-axis shows that point's distance from its own group centroid."
        f"{permdisp_text}"
        if np.isfinite(permanova_f)
        else "PERMANOVA could not be computed for this channel because the latent feature data were insufficient."
    )


def _plotly_group_trace_with_mad(fig, times, matrix, group, color, row=1, col=1, secondary_y=False, name_suffix=""):
    import plotly.graph_objects as go

    med, mad = median_and_mad(matrix, axis=0)
    fig.add_trace(
        go.Scatter(
            x=np.concatenate([times, times[::-1]]),
            y=np.concatenate([med + mad, (med - mad)[::-1]]),
            fill="toself",
            fillcolor=_hex_to_rgba(color, 0.16),
            line=dict(color="rgba(0,0,0,0)"),
            hoverinfo="skip",
            name=f"{group} MAD{name_suffix}",
            showlegend=False,
        ),
        row=row,
        col=col,
        secondary_y=secondary_y,
    )
    fig.add_trace(
        go.Scatter(
            x=times,
            y=med,
            mode="lines",
            line=dict(color=color, width=2.7),
            name=f"{group} median (n={matrix.shape[0]}){name_suffix}",
        ),
        row=row,
        col=col,
        secondary_y=secondary_y,
    )


def _hex_to_rgba(hex_color, alpha):
    rgba = mcolors.to_rgba(hex_color, alpha=alpha)
    return f"rgba({int(rgba[0]*255)},{int(rgba[1]*255)},{int(rgba[2]*255)},{rgba[3]:.3f})"


def _is_pcgc_scalp_channel(ch_name):
    """Keep EEG channel labels for per-channel group comparison."""
    ch = str(ch_name).strip()
    if not ch or _is_ecg_channel(ch):
        return False
    return _is_eeg_channel(ch)


def _compute_nn_distances(z, labels):
    """Per-sample within-group and cross-group nearest-neighbor distances in latent space.

    For each sample computes:
      - within-group NN distance: distance to nearest neighbour in the same group
      - cross-group NN distance: distance to nearest neighbour in the other group
    Returns (nn_within, nn_cross, p_value) where p_value is a one-sided Mann-Whitney U
    test that within-group distances are less than cross-group distances.
    """
    if z is None or len(z) < 4:
        return None, None, np.nan
    from scipy.spatial.distance import cdist
    unique_groups = np.unique(labels)
    if len(unique_groups) < 2:
        return None, None, np.nan
    nn_within = np.full(len(z), np.nan)
    nn_cross = np.full(len(z), np.nan)
    dist_matrix = cdist(z, z, metric='euclidean')
    np.fill_diagonal(dist_matrix, np.inf)
    for i, lbl in enumerate(labels):
        same_mask = labels == lbl
        same_mask = same_mask.copy()
        same_mask[i] = False
        diff_mask = labels != lbl
        if np.any(same_mask):
            nn_within[i] = float(np.min(dist_matrix[i][same_mask]))
        if np.any(diff_mask):
            nn_cross[i] = float(np.min(dist_matrix[i][diff_mask]))
    valid = ~(np.isnan(nn_within) | np.isnan(nn_cross))
    if valid.sum() < 4:
        return nn_within, nn_cross, np.nan
    try:
        _, p = stats.mannwhitneyu(nn_within[valid], nn_cross[valid], alternative='less')
    except Exception:
        p = np.nan
    return nn_within, nn_cross, float(p)


def run_per_channel_group_comparison_analysis(base_path, selected_stage):
    """Plotly per-channel EEG HEP group comparison with ECG overlay and autoencoder features."""
    available_groups = sorted([g for g in os.listdir(base_path) if os.path.isdir(os.path.join(base_path, g))])
    if len(available_groups) < 2:
        st.error("At least two groups are required for per-channel group comparison.")
        return

    default_groups = default_hep_groups(available_groups)
    selected_groups = st.multiselect(
        "Select Groups to Compare",
        options=available_groups,
        default=default_groups,
        key="pcgc_selected_groups",
    )
    if len(selected_groups) < 2:
        st.warning("Please select at least two groups.")
        return

    with st.expander("Per-Channel Group Comparison Settings", expanded=True):
        col1, col2, col3, col4, col5, col6, col7 = st.columns(7)
        with col1:
            test_run = st.checkbox("Test Run (5 files/group)", value=False, key="pcgc_test_run")
        with col2:
            n_permutations = st.slider("Permutations", 50, 500, 200, 50, key="pcgc_n_perm")
        with col3:
            jitter_sec = st.number_input("Jitter (s)", 0.01, 0.5, 0.1, 0.05, key="pcgc_jitter")
        with col4:
            use_zscore = st.checkbox("Z-score subjects", value=True, key="pcgc_zscore")
        with col5:
            apply_ica = st.checkbox("Apply ICA", value=True, key="pcgc_apply_ica")
        with col6:
            clean_ecg = st.checkbox("ECG-regress EEG", value=True, key="pcgc_clean_ecg")
        with col7:
            parallel_workers = st.number_input(
                "Patient workers",
                min_value=1,
                max_value=max(1, os.cpu_count() or 1),
                value=DEFAULT_PATIENT_WORKERS,
                step=1,
                key="pcgc_patient_workers",
            )

    recompute_cache = st.button("Recompute Cache", key="pcgc_recompute_cache")
    if recompute_cache and hasattr(get_group_individuals, "clear"):
        get_group_individuals.clear()

    global_excluded_pids = load_excluded_pids(base_path)
    group_individuals = {}
    for group in selected_groups:
        with st.spinner(f"Loading {group}…"):
            inds = get_group_individuals(
                group,
                selected_stage,
                base_path,
                test_run=test_run,
                apply_ica=apply_ica,
                recompute_cache=recompute_cache,
                parallel_workers=parallel_workers,
            )
        inds = filter_excluded(inds, global_excluded_pids) if inds else []
        if inds:
            group_individuals[group] = inds
        else:
            st.warning(f"No valid HEP data for **{group}** in {selected_stage}.")

    groups_with_data = [g for g in selected_groups if g in group_individuals]
    if len(groups_with_data) < 2:
        st.error("At least two selected groups need valid data.")
        return

    group_channel_sets = []
    for group in groups_with_data:
        counts = Counter()
        inds = group_individuals[group]
        for ind in inds:
            ch_names = ind[3] if len(ind) > 3 else None
            if ch_names is not None:
                counts.update(set(ch_names))
        min_count = max(1, math.ceil(0.5 * len(inds)))
        group_channel_sets.append({
            ch for ch, count in counts.items()
            if count >= min_count and _is_pcgc_scalp_channel(ch)
        })
    common_channels = sorted(set.intersection(*group_channel_sets)) if group_channel_sets else []
    has_common_ecg = all(
        sum(
            len(ind) > 5
            and ind[5] is not None
            and np.asarray(ind[5]).squeeze().ndim == 1
            and len(np.asarray(ind[5]).squeeze()) > 1
            for ind in group_individuals[g]
        ) >= max(1, math.ceil(0.5 * len(group_individuals[g])))
        for g in groups_with_data
    )
    selectable_channels = common_channels + (["ECG"] if has_common_ecg else [])
    if not selectable_channels:
        st.error("No EEG channels or ECG/EKG HEP data exist in every selected group.")
        return

    selected_channels = st.multiselect(
        "Select Common EEG / ECG Channels",
        options=selectable_channels,
        default=selectable_channels[:min(12, len(selectable_channels))],
        key="pcgc_channels",
    )
    if not selected_channels:
        st.warning("Please select at least one channel.")
        return
    if len(groups_with_data) > 2:
        st.info(
            "Waveform and autoencoder plots include all selected groups. "
            f"Cluster-permutation statistics are shown for the first two groups: {groups_with_data[0]} vs {groups_with_data[1]}."
        )

    try:
        from plotly.subplots import make_subplots
        import plotly.graph_objects as go
    except ImportError:
        st.error("Plotly is required for this mode. Please install `plotly`.")
        return

    palette = ['#1f77b4', '#d62728', '#2ca02c', '#9467bd', '#ff7f0e', '#17becf']
    group_color = {g: palette[i % len(palette)] for i, g in enumerate(groups_with_data)}
    unit = "Z-score" if use_zscore else "µV"
    group_channel_data = {g: {} for g in groups_with_data}
    group_ecg_data = {g: [] for g in groups_with_data}
    patient_channel_data = {g: {} for g in groups_with_data}
    times_ref = None

    for group in groups_with_data:
        for ind in group_individuals[group]:
            pid, hep_data, times, ch_names, _, ecg_hep, _ = ind[:7]
            if hep_data is None or times is None or ch_names is None:
                continue
            if times_ref is None:
                times_ref = np.asarray(times, dtype=float)
            eeg_data = apply_ecg_regression(np.asarray(hep_data, dtype=float), ecg_hep) if clean_ecg else np.asarray(hep_data, dtype=float)
            for ch in selected_channels:
                if ch == "ECG":
                    continue
                if ch in ch_names:
                    trace_1d = eeg_data[ch_names.index(ch)]
                    group_channel_data[group].setdefault(ch, []).append(trace_1d)
                    patient_channel_data[group].setdefault(pid, {})[ch] = trace_1d
            if ecg_hep is not None:
                ecg_trace = np.asarray(ecg_hep, dtype=float).squeeze()
                if ecg_trace.ndim == 1 and len(ecg_trace) > 1:
                    group_ecg_data[group].append(ecg_trace * 1e6)

    if times_ref is None:
        st.error("No valid time axis found.")
        return

    for group in groups_with_data:
        if "ECG" in selected_channels and group_ecg_data[group]:
            ecg_stacked, _, _ = _stack_traces_with_common_length(group_ecg_data[group], times=times_ref)
            if ecg_stacked is not None:
                group_channel_data[group]["ECG"] = scale_matrix(ecg_stacked, use_zscore)
        for ch in list(group_channel_data[group]):
            traces = group_channel_data[group][ch]
            if ch == "ECG" and isinstance(traces, np.ndarray):
                continue
            stacked, _, times_use = _stack_traces_with_common_length(traces, times=times_ref)
            if stacked is None:
                del group_channel_data[group][ch]
                continue
            group_channel_data[group][ch] = scale_matrix(stacked, use_zscore)

    valid_channels = [
        ch for ch in selected_channels
        if all(ch in group_channel_data[g] and group_channel_data[g][ch].shape[0] >= 2 for g in groups_with_data)
    ]
    if not valid_channels:
        st.error("No selected channel had at least two subjects in every selected group.")
        return

    st.subheader("All Selected Channels Together")

    # Per-patient encoding: one row per patient with all channels hstacked.
    min_ch_len = min(
        group_channel_data[g][ch].shape[1]
        for g in groups_with_data
        for ch in valid_channels
    )
    aggregate_by_group = {}       # (n_patients, n_times) channel-averaged — for waveform
    multivariate_by_group = {}    # (n_patients, n_channels * n_times) — for autoencoder
    patient_ids_by_group = {}

    for group in groups_with_data:
        common_pids = sorted([
            pid for pid, ch_dict in patient_channel_data[group].items()
            if all(ch in ch_dict for ch in valid_channels)
        ])
        if not common_pids:
            mats = [group_channel_data[group][ch][:, :min_ch_len] for ch in valid_channels]
            n_pat = min(m.shape[0] for m in mats)
            aggregate_by_group[group] = np.mean(
                np.stack([m[:n_pat] for m in mats], axis=0), axis=0
            )
            multivariate_by_group[group] = np.hstack([m[:n_pat] for m in mats])
            patient_ids_by_group[group] = list(range(n_pat))
        else:
            ch_mats = [
                np.vstack([patient_channel_data[group][pid][ch][:min_ch_len] for pid in common_pids])
                for ch in valid_channels
            ]
            aggregate_by_group[group] = np.mean(np.stack(ch_mats, axis=0), axis=0)
            multivariate_by_group[group] = np.hstack(ch_mats)
            patient_ids_by_group[group] = common_pids

    times_all = times_ref[:min_ch_len]
    for group in groups_with_data:
        aggregate_by_group[group] = aggregate_by_group[group][:, :len(times_all)]

    fig_all = make_subplots(
        rows=1,
        cols=2,
        subplot_titles=("All selected channels: HEP median +/- MAD", "Autoencoder latent features"),
    )
    for group in groups_with_data:
        _plotly_group_trace_with_mad(fig_all, times_all, aggregate_by_group[group], group, group_color[group], row=1, col=1)

    X_all = np.vstack([multivariate_by_group[g] for g in groups_with_data])
    labels_all = np.concatenate([[g] * multivariate_by_group[g].shape[0] for g in groups_with_data])
    z_all, err_all, ae_method_all = _run_waveform_autoencoder_embedding(X_all)
    _, latent_p_all = _latent_group_permutation_test(
        z_all,
        labels_all,
        n_permutations=max(200, n_permutations),
        random_state=42,
    ) if z_all is not None else (np.nan, np.nan)
    if z_all is not None:
        for group in groups_with_data:
            mask = labels_all == group
            fig_all.add_trace(
                go.Scatter(
                    x=z_all[mask, 0],
                    y=z_all[mask, 1],
                    mode="markers",
                    marker=dict(color=group_color[group], size=8, opacity=0.75),
                    name=f"{group} latent",
                    text=[f"{group}<br>recon MSE={e:.4f}" for e in err_all[mask]],
                ),
                row=1,
                col=2,
            )
    fig_all.add_vline(x=0, line_dash="dash", line_color="gray", row=1, col=1)
    fig_all.add_hline(y=0, line_color="gray", line_width=1, row=1, col=1)
    fig_all.update_layout(
        height=520,
        hovermode="x unified",
        title=(
            f"Per-Channel Group Comparison | {selected_stage} | {len(valid_channels)} channels | "
            f"{ae_method_all} | latent {_format_p_value(latent_p_all)}"
        ),
    )
    fig_all.update_xaxes(title_text="Time relative to R-peak (s)", row=1, col=1)
    fig_all.update_yaxes(title_text=f"Amplitude ({unit})", row=1, col=1)
    fig_all.update_xaxes(title_text="Latent 1", row=1, col=2)
    fig_all.update_yaxes(title_text="Latent 2", row=1, col=2)
    st.plotly_chart(fig_all, use_container_width=True)
    st.caption(
        _pcgc_figure_explanation(
            groups=groups_with_data,
            unit=unit,
            latent_p=latent_p_all,
            all_channels=True,
        )
    )

    summary_rows = []
    group_a, group_b = groups_with_data[:2]
    st.subheader("Per-Channel Plots")
    for ch in valid_channels:
        min_len = min(group_channel_data[g][ch].shape[1] for g in groups_with_data)
        times = times_ref[:min_len]
        channel_mats = {g: group_channel_data[g][ch][:, :min_len] for g in groups_with_data}

        X = np.vstack([channel_mats[g] for g in groups_with_data])
        labels = np.concatenate([[g] * channel_mats[g].shape[0] for g in groups_with_data])
        z, recon_err, ae_method = _run_waveform_autoencoder_embedding(X)
        _, latent_p = _latent_group_permutation_test(
            z,
            labels,
            n_permutations=max(200, n_permutations),
            random_state=42,
        ) if z is not None else (np.nan, np.nan)
        permanova_f, permanova_p, permanova_info = _permanova_latent_test(
            z,
            labels,
            n_permutations=max(200, n_permutations),
            random_state=42,
        ) if z is not None else (np.nan, np.nan, {})
        permdisp_f, permdisp_p = _permdisp_latent_test(
            z,
            labels,
            n_permutations=max(200, n_permutations),
            random_state=42,
        ) if z is not None else (np.nan, np.nan)

        sig_windows, t_obs, cohens_d = permutation_two_group_cluster_test(
            channel_mats[group_a],
            channel_mats[group_b],
            times,
            n_permutations=n_permutations,
            p_threshold=0.05,
            jitter_sec=jitter_sec,
            label_a=group_a,
            label_b=group_b,
            channel_label=ch,
            button_key=f"pcgc_perm_{selected_stage}_{ch}",
        )
        med_a, _ = median_and_mad(channel_mats[group_a], axis=0)
        med_b, _ = median_and_mad(channel_mats[group_b], axis=0)
        diff = med_a - med_b
        point_t, point_p = stats.ttest_ind(channel_mats[group_a], channel_mats[group_b], axis=0, equal_var=False, nan_policy="omit")
        peak_idx = int(np.nanargmax(np.abs(diff)))
        min_cluster_p = min((w["p_value"] for w in sig_windows), default=np.nan)
        summary_rows.append({
            "channel": ch,
            f"{group_a}_n": channel_mats[group_a].shape[0],
            f"{group_b}_n": channel_mats[group_b].shape[0],
            "peak_diff_time_ms": times[peak_idx] * 1000.0,
            "peak_diff": diff[peak_idx],
            "peak_t": point_t[peak_idx],
            "peak_uncorrected_p": point_p[peak_idx],
            "min_cluster_p": min_cluster_p,
            "autoencoder_method": ae_method,
            "autoencoder_latent_p": latent_p,
            "autoencoder_permanova_F": permanova_f,
            "autoencoder_permanova_p": permanova_p,
            "autoencoder_permdisp_F": permdisp_f,
            "autoencoder_permdisp_p": permdisp_p,
            "autoencoder_error_p": (
                stats.ttest_ind(
                    recon_err[labels == group_a],
                    recon_err[labels == group_b],
                    equal_var=False,
                    nan_policy="omit",
                ).pvalue
                if recon_err is not None and group_a in labels and group_b in labels else np.nan
            ),
        })

        fig = make_subplots(
            rows=2,
            cols=2,
            specs=[[{}, {}], [{"secondary_y": True}, {}]],
            subplot_titles=(
                f"{ch}: HEP median +/- MAD",
                f"{ch}: autoencoder latent features",
                f"{ch}: HEP + ECG median",
                "Difference, t-statistic, and significant windows",
            ),
            vertical_spacing=0.12,
            horizontal_spacing=0.08,
        )
        for group in groups_with_data:
            _plotly_group_trace_with_mad(fig, times, channel_mats[group], group, group_color[group], row=1, col=1)
            _plotly_group_trace_with_mad(fig, times, channel_mats[group], group, group_color[group], row=2, col=1)
            ecg_stack, _, _ = _stack_traces_with_common_length(group_ecg_data.get(group, []), times=times_ref)
            if ecg_stack is not None:
                ecg_len = min(min_len, ecg_stack.shape[1], len(times))
                ecg_stack = ecg_stack[:, :ecg_len]
                ecg_med, _ = median_and_mad(ecg_stack, axis=0)
                fig.add_trace(
                    go.Scatter(
                        x=times[:ecg_len],
                        y=ecg_med,
                        mode="lines",
                        line=dict(color=group_color[group], width=1.6, dash="dot"),
                        name=f"{group} ECG median",
                    ),
                    row=2,
                    col=1,
                    secondary_y=True,
                )
        if z is not None:
            for group in groups_with_data:
                mask = labels == group
                fig.add_trace(
                    go.Scatter(
                        x=z[mask, 0],
                        y=z[mask, 1],
                        mode="markers",
                        marker=dict(color=group_color[group], size=9, opacity=0.78),
                        name=f"{group} AE latent",
                        text=[f"{group}<br>{ch}<br>recon MSE={e:.4f}" for e in recon_err[mask]],
                    ),
                    row=1,
                    col=2,
                )

        fig.add_trace(
            go.Scatter(x=times, y=diff, mode="lines", line=dict(color="#6a3d9a", width=2.5), name=f"{group_a} - {group_b}"),
            row=2,
            col=2,
        )
        fig.add_trace(
            go.Scatter(x=times, y=t_obs, mode="lines", line=dict(color="#444", width=1.4, dash="dot"), name="Welch t"),
            row=2,
            col=2,
        )
        for win in sig_windows:
            for row, col in [(1, 1), (2, 1), (2, 2)]:
                fig.add_vrect(
                    x0=win["start"],
                    x1=win["end"],
                    fillcolor="rgba(255,165,0,0.18)",
                    line_width=0,
                    row=row,
                    col=col,
                )
        for row, col in [(1, 1), (2, 1), (2, 2)]:
            fig.add_vline(x=0, line_dash="dash", line_color="gray", row=row, col=col)
            fig.add_hline(y=0, line_color="gray", line_width=1, row=row, col=col)

        fig.update_layout(
            height=820,
            hovermode="x unified",
            title=f"Per-Channel Group Comparison | {ch} | {selected_stage} | {ae_method} | latent {_format_p_value(latent_p)}",
            legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
        )
        fig.update_xaxes(title_text="Time relative to R-peak (s)", row=2, col=1)
        fig.update_xaxes(title_text="Time relative to R-peak (s)", row=2, col=2)
        fig.update_yaxes(title_text=f"HEP ({unit})", row=1, col=1)
        fig.update_yaxes(title_text=f"HEP ({unit})", row=2, col=1, secondary_y=False)
        fig.update_yaxes(title_text="ECG (µV)", row=2, col=1, secondary_y=True)
        fig.update_yaxes(title_text=f"Difference ({unit}) / t", row=2, col=2)
        fig.update_xaxes(title_text="Latent 1", row=1, col=2)
        fig.update_yaxes(title_text="Latent 2", row=1, col=2)
        st.plotly_chart(fig, use_container_width=True)
        st.caption(
            _pcgc_figure_explanation(
                groups=groups_with_data,
                unit=unit,
                latent_p=latent_p,
                min_cluster_p=min_cluster_p,
                sig_windows=sig_windows,
                include_ecg=True,
            )
        )

        if permanova_info:
            fig_perm = go.Figure()
            perm_labels = permanova_info["labels"]
            perm_dist = permanova_info["dist_to_centroid"]
            perm_z = permanova_info["z"]
            for group in groups_with_data:
                mask = perm_labels == group
                if not np.any(mask):
                    continue
                fig_perm.add_trace(
                    go.Box(
                        y=perm_dist[mask],
                        name=group,
                        marker_color=group_color[group],
                        boxpoints="all",
                        jitter=0.35,
                        pointpos=0,
                        customdata=np.column_stack([perm_z[mask, 0], perm_z[mask, 1]]),
                        hovertemplate=(
                            f"{group}<br>"
                            "distance to centroid=%{y:.4f}<br>"
                            "latent 1=%{customdata[0]:.4f}<br>"
                            "latent 2=%{customdata[1]:.4f}<extra></extra>"
                        ),
                    )
                )
            fig_perm.update_layout(
                height=420,
                title=(
                    f"{ch}: Autoencoder Latent Features PERMANOVA / PERMDISP | "
                    f"PERMANOVA pseudo-F={permanova_f:.3f}, {_format_p_value(permanova_p)} | "
                    f"PERMDISP F={permdisp_f:.3f}, {_format_p_value(permdisp_p)}"
                ),
                yaxis_title="Distance to group centroid in latent space",
                xaxis_title="Group",
                showlegend=False,
            )
            st.plotly_chart(fig_perm, use_container_width=True)
            st.caption(_pcgc_permanova_explanation(permanova_f, permanova_p, permdisp_f, permdisp_p))

        # ── Nearest Neighbor Distance Distribution ───────────────────────────
        if z is not None and len(groups_with_data) >= 2:
            nn_within, nn_cross, nn_p = _compute_nn_distances(z, labels)
            if nn_within is not None and nn_cross is not None:
                fig_nn = go.Figure()
                for grp in groups_with_data:
                    grp_mask = labels == grp
                    if not np.any(grp_mask):
                        continue
                    color = group_color[grp]
                    fig_nn.add_trace(
                        go.Violin(
                            y=nn_within[grp_mask],
                            name=f"{grp} within-group",
                            box_visible=True,
                            meanline_visible=True,
                            points="all",
                            jitter=0.3,
                            marker_color=color,
                            line_color=color,
                            opacity=0.85,
                            legendgroup=grp,
                            hovertemplate=f"{grp} within-group<br>distance=%{{y:.4f}}<extra></extra>",
                        )
                    )
                    fig_nn.add_trace(
                        go.Violin(
                            y=nn_cross[grp_mask],
                            name=f"{grp} cross-group",
                            box_visible=True,
                            meanline_visible=True,
                            points="all",
                            jitter=0.3,
                            marker_color=color,
                            line_color=color,
                            fillcolor=_hex_to_rgba(color, 0.12),
                            opacity=0.55,
                            legendgroup=grp,
                            hovertemplate=f"{grp} cross-group<br>distance=%{{y:.4f}}<extra></extra>",
                        )
                    )
                nn_p_str = _format_p_value(nn_p) if not (isinstance(nn_p, float) and np.isnan(nn_p)) else "n/a"
                fig_nn.update_layout(
                    height=420,
                    title=(
                        f"{ch}: Nearest Neighbor Distance Distribution | "
                        f"Within vs Cross-Group | Mann-Whitney {nn_p_str}"
                    ),
                    yaxis_title="NN distance in latent space",
                    xaxis_title="Group / type",
                    violinmode="group",
                )
                st.plotly_chart(fig_nn, use_container_width=True)
                st.caption(
                    "Within-group: each subject's distance to its nearest neighbour in the same group "
                    "in autoencoder latent space. "
                    "Cross-group: distance to nearest neighbour in the other group. "
                    "If groups are well-separated, within-group distances should be smaller than cross-group distances. "
                    f"Mann-Whitney U p-value ({nn_p_str}): one-sided test that within-group distances < cross-group distances."
                )

    if summary_rows:
        st.subheader("Per-Channel Statistical Summary")
        summary_df = pd.DataFrame(summary_rows).sort_values(["min_cluster_p", "peak_uncorrected_p"], na_position="last")
        st.dataframe(
            summary_df.style.format({
                "peak_diff_time_ms": "{:.1f}",
                "peak_diff": "{:.3f}",
                "peak_t": "{:.3f}",
                "peak_uncorrected_p": "{:.4g}",
                "min_cluster_p": "{:.4g}",
                "autoencoder_latent_p": "{:.4g}",
                "autoencoder_permanova_F": "{:.3f}",
                "autoencoder_permanova_p": "{:.4g}",
                "autoencoder_permdisp_F": "{:.3f}",
                "autoencoder_permdisp_p": "{:.4g}",
                "autoencoder_error_p": "{:.4g}",
            }),
            use_container_width=True,
        )
        st.caption(
            "autoencoder_latent_p tests whether groups separate in the 2D autoencoder feature space. "
            "autoencoder_permanova_p is a PERMANOVA test on Euclidean distances in that same latent space. "
            "autoencoder_permdisp_p tests whether the groups have different dispersion around their latent-space centroids. "
            "autoencoder_error_p tests whether reconstruction error differs between the first two groups."
        )

# ═══════════════════════════════════════════════════════════════════════════
# DIAGNOSIS GROUP COMPARISON — helpers
# ═══════════════════════════════════════════════════════════════════════════

_DIAG_BASE = "/storage/pblab_shared_data2/Nir/Cobrad/EDF_Format/Harvard_Electroencephalography/ECG/I0001"

def _load_12sl_diagnosis_map():
    """Load 12SL ECG diagnoses → {patient_id_str: category_name}."""
    import os as _os2
    diag_acq_path = _os2.path.join(_DIAG_BASE, "12SL_diagnoses", "diagnoses_acquisition.csv")
    diag_dict_path = _os2.path.join(_DIAG_BASE, "12SL_diagnoses", "diagnoses_dictionary.csv")
    if not _os2.path.exists(diag_acq_path):
        return {}, {}
    code_to_name = {}
    if _os2.path.exists(diag_dict_path):
        dd = pd.read_csv(diag_dict_path)
        for _, row in dd.iterrows():
            try:
                code_to_name[int(row['codes'])] = str(row['diagnoses']).strip()
            except Exception:
                pass

    def _code_to_category(codes_list):
        afib_codes = set(range(180, 230))
        if set(codes_list) & afib_codes:
            return "Atrial Fibrillation/Flutter"
        if any(300 <= c <= 399 for c in codes_list):
            return "Bundle Branch Block / WPW"
        if any(400 <= c <= 499 for c in codes_list):
            return "ST-T Changes / Ischemia"
        if any(500 <= c <= 599 for c in codes_list):
            return "Hypertrophy"
        if any(700 <= c <= 799 for c in codes_list):
            return "Pacemaker Rhythm"
        if any(101 <= c <= 179 for c in codes_list):
            return "Conduction Disorder"
        if any(30 <= c <= 99 for c in codes_list):
            return "Rhythm Disorder"
        if any(c >= 1670 for c in codes_list):
            return "Normal Sinus Rhythm"
        return "Other"

    patient_categories = {}
    try:
        chunks = pd.read_csv(diag_acq_path, chunksize=200000, low_memory=False)
        for chunk in chunks:
            chunk['bdsp_id'] = chunk['FileName'].str.extract(r'/de_(\d+)_')
            chunk = chunk.dropna(subset=['bdsp_id'])
            for _, row in chunk.iterrows():
                bdsp_str = str(row['bdsp_id'])
                if bdsp_str in patient_categories:
                    continue
                codes_raw = row.get('codes_physician', '') or row.get('codes_software', '')
                if not codes_raw or str(codes_raw) == 'nan':
                    codes_raw = row.get('codes_software', '') or ''
                try:
                    codes_list = [int(c.strip()) for c in str(codes_raw).split(',') if c.strip().isdigit()]
                except Exception:
                    codes_list = []
                if codes_list:
                    patient_categories[bdsp_str] = _code_to_category(codes_list)
    except Exception as e:
        pass
    return patient_categories, code_to_name


def _match_pid_to_bdsp(pid_str):
    """Extract candidate BDSPPatientID strings from a pickle patient ID string."""
    import re as _re2
    candidates = []
    nums = _re2.findall(r'\d+', pid_str)
    for n in nums:
        if len(n) >= 9:
            candidates.append(n[-9:])
            candidates.append(n)
            stripped = n.lstrip('0')
            if len(stripped) >= 9:
                candidates.append(stripped[-9:])
    return list(dict.fromkeys(candidates))


_EHR_BASE = "/storage/pblab_shared_data2/Nir/Cobrad/EDF_Format/Harvard_Electroencephalography/EHR"
_EHR_DIAG_BASE = os.path.join(
    _EHR_BASE,
    "data_Structured",
    "Parquet",
    "Parquet",
    "bdsp_i0006_diagnosis",
)
_EHR_I0002_ICD10_BASE = os.path.join(
    _EHR_BASE,
    "I0002-EHR",
    "I0002_structured",
    "icd10_codes_nax_2024_parquet",
)
_DXGC_EHR_SOURCE_VERSION = "i0006_i0002_exclusive_auto_v2"

_DIAG_CATEGORIES = {
    "Obstructive Sleep Apnea": lambda n: "sleep apnea" in n.lower(),
    "Hypertension": lambda n: "hypertension" in n.lower(),
    "Diabetes": lambda n: "diabetes" in n.lower(),
    "Heart Failure": lambda n: "heart failure" in n.lower(),
    "Coronary Artery Disease": lambda n: "coronary" in n.lower() or "atherosclerotic heart" in n.lower(),
    "Heart Transplant": lambda n: "heart transplant" in n.lower() or "transplant status" in n.lower(),
    "Atrial Fibrillation": lambda n: "atrial fibrillation" in n.lower() or "atrial flutter" in n.lower(),
    "COPD": lambda n: "chronic obstructive" in n.lower(),
    "Obesity": lambda n: "obesity" in n.lower() or "obese" in n.lower(),
    "Depression": lambda n: "depressive" in n.lower() or "depression" in n.lower(),
    "Anxiety": lambda n: "anxiety" in n.lower(),
    "Cognitive Impairment / Dementia": lambda n: "cognitive" in n.lower() or "alzheimer" in n.lower() or "dementia" in n.lower(),
    "Stroke / Cerebrovascular": lambda n: "infarction" in n.lower() or "cerebrovascular" in n.lower() or "stroke" in n.lower() or "hemorrhage" in n.lower(),
    "Kidney Disease": lambda n: "kidney" in n.lower() or "renal" in n.lower(),
    "Anemia": lambda n: "anemia" in n.lower(),
}


def _add_dx_rows_to_patient_map(patient_map, df, pid_col, dx_col, target_ids=None):
    """Merge diagnosis rows into {bdsp_patient_id: set(dx_name, ...)}."""
    if df is None or df.empty or pid_col not in df.columns or dx_col not in df.columns:
        return
    work = df[[pid_col, dx_col]].copy()
    work = work[work[pid_col].notna() & work[dx_col].notna()]
    if work.empty:
        return
    work[pid_col] = pd.to_numeric(work[pid_col], errors="coerce")
    work = work.dropna(subset=[pid_col])
    work[pid_col] = work[pid_col].astype(int)
    if target_ids is not None:
        work = work[work[pid_col].isin(target_ids)]
    work[dx_col] = work[dx_col].astype(str).str.strip()
    work = work[(work[dx_col] != "") & (work[dx_col].str.lower() != "not recorded")]
    for pid_val, grp in work.groupby(pid_col):
        patient_map.setdefault(int(pid_val), set()).update(grp[dx_col].dropna().unique().tolist())


def _read_ehr_parquet(path, columns, pid_col, target_ids=None):
    """Read an EHR parquet chunk, using patient-ID filters when the parquet engine supports them."""
    if target_ids:
        try:
            return pd.read_parquet(
                path,
                columns=columns,
                filters=[(pid_col, "in", list(target_ids))],
            )
        except Exception:
            pass
    return pd.read_parquet(path, columns=columns)


def _load_ehr_diagnosis_map(target_ids=None):
    """Load I0006 and I0002 diagnosis parquet chunks -> {bdsp_patient_id: [dx_name, ...]}."""
    import glob as _glob
    target_ids = set(int(x) for x in target_ids) if target_ids is not None else None
    patient_map = {}

    chunks = sorted(_glob.glob(os.path.join(_EHR_DIAG_BASE, "*.parquet")))
    for ch in chunks:
        try:
            df_ch = _read_ehr_parquet(
                ch,
                columns=['bdsp_patient_id', 'dx_name', 'dx_code_icd_10'],
                pid_col='bdsp_patient_id',
                target_ids=target_ids,
            )
            _add_dx_rows_to_patient_map(patient_map, df_ch, 'bdsp_patient_id', 'dx_name', target_ids)
        except Exception:
            pass

    i0002_chunks = sorted(_glob.glob(os.path.join(_EHR_I0002_ICD10_BASE, "*.parquet")))
    for ch in i0002_chunks:
        try:
            df_ch = _read_ehr_parquet(
                ch,
                columns=['BDSPPatientID', 'LongDescription'],
                pid_col='BDSPPatientID',
                target_ids=target_ids,
            )
            _add_dx_rows_to_patient_map(patient_map, df_ch, 'BDSPPatientID', 'LongDescription', target_ids)
        except Exception:
            pass

    return {pid: sorted(dx_names) for pid, dx_names in patient_map.items()}


def _pid_to_bdsp_ehr(pid_str):
    """SUB-I0006179000017_... or SUB-I0002150000096_... -> 9-digit BDSPPatientID."""
    import re as _re_dxg
    m = _re_dxg.search(r'I\d{4}(\d{9})', str(pid_str))
    if m:
        return int(m.group(1))
    return None


def _assign_categories(dx_names_list):
    """Map list of dx_names to set of broad category names."""
    cats = set()
    for dx in dx_names_list:
        for cat_name, test_fn in _DIAG_CATEGORIES.items():
            if test_fn(dx):
                cats.add(cat_name)
    return cats


def _dxgc_cache_key(hep_groups, stage, grouping_mode, min_group_size, apply_ecg_reg, extra=None):
    """Compute cache key from analysis settings."""
    import hashlib as _hl, json as _js
    key = {"g": sorted(hep_groups), "s": stage, "gm": grouping_mode,
           "mg": int(min_group_size), "ecg": bool(apply_ecg_reg),
           "ehr": _DXGC_EHR_SOURCE_VERSION}
    if extra:
        key["extra"] = extra
    return _hl.md5(_js.dumps(key, sort_keys=True).encode()).hexdigest()[:12]


def _dxgc_cache_path(base_path, key):
    """Return path to cache pickle for given key."""
    cache_dir = os.path.join(base_path, ".dxgc_cache")
    os.makedirs(cache_dir, exist_ok=True)
    return os.path.join(cache_dir, f"dxgc_{key}.pkl")


def _dxgc_count_on_disk(base_path, hep_groups, stage):
    """Count patient pkl files on disk (fast, no loading) for cache invalidation."""
    total = 0
    for grp in hep_groups:
        d = os.path.join(base_path, grp, stage)
        if os.path.isdir(d):
            total += len([f for f in os.listdir(d)
                          if f.endswith('.pkl')
                          and not f.startswith('individuals_cache')
                          and not f.startswith('non_eeg_individuals_cache')])
    return total


def run_diagnosis_group_comparison_analysis(base_path, selected_stage):
    """Cluster patients by EHR diagnosis from I0006 and I0002 parquet data."""
    EEG_DISPLAY_SCALE = 1e6
    EEG_DISPLAY_UNIT = "µV"

    def _eeg_display_matrix(matrix):
        return np.asarray(matrix, dtype=float) * EEG_DISPLAY_SCALE

    def _show_non_ica_figure(fig, **kwargs):
        fig._skip_ica_cleaned_title = True
        st.pyplot(fig, **kwargs)

    def _filter_diagnosis_category_groups(group_data, cat_counts):
        if grouping_mode != "By Diagnosis Category" or not group_data:
            return group_data, None

        options = sorted(
            group_data.keys(),
            key=lambda g: (-int(cat_counts.get(g, group_data[g].get("n", 0))), str(g)),
        )
        default = options[:min(4, len(options))]
        selected = st.multiselect(
            "Select diagnosis categories to plot",
            options=options,
            default=default,
            key="dxgc_plot_cats_v2",
        )
        if not selected:
            st.warning("Select at least one diagnosis category to plot.")
            return {}, None
        return {g: group_data[g] for g in selected if g in group_data}, selected

    def _show_diagnosis_category_venn(patient_info, cat_counts, key_suffix):
        """Show exact patient intersections for two or three broad diagnosis categories."""
        available = [
            cat for cat, _ in sorted(cat_counts.items(), key=lambda item: (-item[1], item[0]))
            if int(cat_counts.get(cat, 0)) > 0
        ]
        if len(available) < 2:
            st.info("At least two populated diagnosis categories are needed for an overlap diagram.")
            return

        st.markdown("**Diagnosis Category Overlap (Venn Diagram)**")
        selected = st.multiselect(
            "Choose 2 or 3 diagnosis categories",
            options=available,
            default=available[:min(3, len(available))],
            max_selections=3,
            key=f"dxgc_table_venn_categories_{key_suffix}",
        )
        if len(selected) not in (2, 3):
            st.caption("Select exactly 2 or 3 categories to draw the Venn diagram.")
            return

        pid_sets = {
            cat: {
                str(pid) for pid, info in patient_info.items()
                if cat in set(info.get("categories", []))
            }
            for cat in selected
        }
        colors = ['#4C78A8', '#F58518', '#54A24B']
        fig, ax = plt.subplots(figsize=(9, 6))
        ax.set_aspect('equal')
        ax.axis('off')

        from matplotlib.patches import Circle as _VennCircle
        if len(selected) == 2:
            a, b = selected
            only_a = len(pid_sets[a] - pid_sets[b])
            only_b = len(pid_sets[b] - pid_sets[a])
            both = len(pid_sets[a] & pid_sets[b])
            for center, color in [((-0.7, 0), colors[0]), ((0.7, 0), colors[1])]:
                ax.add_patch(_VennCircle(center, 1.45, facecolor=color, edgecolor=color,
                                         alpha=0.35, linewidth=2))
            ax.text(-1.25, 0, str(only_a), ha='center', va='center', fontsize=15, fontweight='bold')
            ax.text(1.25, 0, str(only_b), ha='center', va='center', fontsize=15, fontweight='bold')
            ax.text(0, 0, str(both), ha='center', va='center', fontsize=15, fontweight='bold')
            ax.text(-1.0, 1.65, f"{a}\nN={len(pid_sets[a])}", ha='center', va='bottom', fontsize=10)
            ax.text(1.0, 1.65, f"{b}\nN={len(pid_sets[b])}", ha='center', va='bottom', fontsize=10)
            ax.set_xlim(-2.4, 2.4); ax.set_ylim(-1.7, 2.25)
        else:
            a, b, c = selected
            sa, sb, sc = pid_sets[a], pid_sets[b], pid_sets[c]
            counts = {
                "a": len(sa - sb - sc), "b": len(sb - sa - sc), "c": len(sc - sa - sb),
                "ab": len((sa & sb) - sc), "ac": len((sa & sc) - sb),
                "bc": len((sb & sc) - sa), "abc": len(sa & sb & sc),
            }
            centers = [(-0.75, 0.45), (0.75, 0.45), (0, -0.65)]
            for center, color in zip(centers, colors):
                ax.add_patch(_VennCircle(center, 1.5, facecolor=color, edgecolor=color,
                                         alpha=0.30, linewidth=2))
            positions = {
                "a": (-1.25, 0.65), "b": (1.25, 0.65), "c": (0, -1.25),
                "ab": (0, 0.85), "ac": (-0.62, -0.42), "bc": (0.62, -0.42),
                "abc": (0, 0.0),
            }
            for region, pos in positions.items():
                ax.text(*pos, str(counts[region]), ha='center', va='center',
                        fontsize=14, fontweight='bold')
            label_positions = [(-1.55, 1.85), (1.55, 1.85), (0, -2.25)]
            for cat, pos in zip(selected, label_positions):
                ax.text(*pos, f"{cat}\nN={len(pid_sets[cat])}", ha='center', va='center', fontsize=10)
            ax.set_xlim(-2.6, 2.6); ax.set_ylim(-2.6, 2.4)

        ax.set_title("Patient Intersections Between Diagnosis Categories",
                     fontsize=13, fontweight='bold', pad=12)
        fig.tight_layout()
        _show_non_ica_figure(fig, use_container_width=True)
        plt.close(fig)
        st.caption(
            "Numbers in non-overlapping regions are category-only counts; overlapping regions show "
            "patients assigned to every category covering that region. Circle areas are illustrative."
        )

    st.header("\U0001fac0\U0001f9e0 Diagnosis Group Comparison")
    st.markdown(
        """
        Groups patients by their clinical diagnoses from the **BDSP EHR databases**
        (`bdsp_i0006_diagnosis` and I0002 ICD10 parquet files).

        **Grouping modes include:**
        - **By Diagnosis Category**: broad ICD categories (OSA, Hypertension, Diabetes, etc.).
          Each patient may appear in multiple categories. Compare patients *with* vs *without* a diagnosis,
          or compare multiple diagnosis groups head-to-head.
        - **By Specific Diagnosis**: search and select exact diagnosis names (e.g. "Obstructive sleep apnea").
        - **Diagnosis Comorbidity Subgroups**: pick one index diagnosis, then compare large mutually exclusive
          subgroups such as OSA only, OSA + hypertension, OSA + hypertension + obesity, and so on.

        Only groups with >= min_patients shown.
        """
    )

    available_groups = sorted([g for g in os.listdir(base_path) if os.path.isdir(os.path.join(base_path, g))])
    selected_hep_groups = st.multiselect(
        "HEP groups to load",
        options=available_groups,
        default=default_hep_groups(available_groups),
        key="dxgc_hep_groups",
    )
    if not selected_hep_groups:
        st.warning("Select at least one HEP group.")
        return

    grouping_mode = st.radio(
        "Grouping mode",
        [
            "By Diagnosis Category",
            "By Specific Diagnosis",
            "Diagnosis Comorbidity Subgroups",
            "By Age Group",
            "By Sex",
            "Exclusive Diagnosis Category",
        ],
        horizontal=True,
        index=5,
        key="dxgc_grouping_mode",
    )

    comorbidity_base_category = None
    comorbidity_categories = []
    comorbidity_max_subgroups = 12
    comorbidity_min_disease_count = 10
    if grouping_mode == "Diagnosis Comorbidity Subgroups":
        category_options = list(_DIAG_CATEGORIES.keys())
        default_base = "Obstructive Sleep Apnea" if "Obstructive Sleep Apnea" in category_options else category_options[0]
        st.markdown("**Comorbidity subgroup setup**")
        cbase, cmax = st.columns(2)
        with cbase:
            comorbidity_base_category = st.selectbox(
                "Index diagnosis",
                options=category_options,
                index=category_options.index(default_base),
                key="dxgc_comorb_base_cat",
                help="Patients must have this broad diagnosis category before comorbidity subgroups are formed.",
            )
        preferred_comorbidities = [
            "Hypertension",
            "Obesity",
            "Diabetes",
            "Heart Failure",
            "Atrial Fibrillation",
            "Coronary Artery Disease",
            "COPD",
            "Kidney Disease",
        ]
        default_comorbidities = [
            cat for cat in preferred_comorbidities
            if cat in category_options and cat != comorbidity_base_category
        ]
        available_comorbidities = [cat for cat in category_options if cat != comorbidity_base_category]
        with cmax:
            comorbidity_max_subgroups = st.number_input(
                "Max displayed subgroups",
                min_value=2,
                max_value=50,
                value=12,
                step=1,
                key="dxgc_comorb_max_subgroups",
                help="Keeps plots and pairwise tests readable by using the largest valid subgroups.",
            )
        comorbidity_categories = st.multiselect(
            "Major comorbidity categories to combine",
            options=available_comorbidities,
            default=default_comorbidities,
            key="dxgc_comorb_cats",
            help="Exact combinations of these categories define the subgroups. Categories not selected here are ignored for subgroup labels.",
        )
        if not comorbidity_categories:
            st.warning("Select at least one comorbidity category.")
            return
        comorbidity_min_disease_count = st.number_input(
            "Ignore comorbidities below N patients",
            min_value=1,
            value=10,
            step=1,
            key="dxgc_comorb_min_disease_count",
            help="Rare comorbidity categories are removed from subgroup labels, but their patients remain in the broader subgroup.",
        )

    col1, col2, col3 = st.columns(3)
    with col1:
        if grouping_mode == "Diagnosis Comorbidity Subgroups":
            min_group_size = st.number_input(
                "Min patients per subgroup",
                min_value=1,
                value=500,
                step=1,
                key="dxgc_mingroup_comorb",
                help="Minimum subgroup size to include in plots and statistics.",
            )
        else:
            min_group_size = st.number_input(
                "Min patients per group",
                min_value=2,
                value=3,
                step=1,
                key="dxgc_mingroup",
            )
    with col2:
        n_permutations = st.number_input("Permutations", 50, 500, 200, 50, key="dxgc_nperms")
    with col3:
        apply_ecg_reg = st.checkbox("Regress out ECG from EEG", True, key="dxgc_ecgreg")

    # ── Cache management ────────────────────────────────────────────────────
    import pickle as _pkl_dxgc
    from datetime import datetime as _dt_dxgc

    cache_extra = None
    if grouping_mode == "Diagnosis Comorbidity Subgroups":
        cache_extra = {
            "base_category": comorbidity_base_category,
            "comorbidity_categories": list(comorbidity_categories),
            "max_subgroups": int(comorbidity_max_subgroups),
            "min_disease_count": int(comorbidity_min_disease_count),
        }
    cache_key = _dxgc_cache_key(
        selected_hep_groups,
        selected_stage,
        grouping_mode,
        min_group_size,
        apply_ecg_reg,
        extra=cache_extra,
    )
    cache_path = _dxgc_cache_path(base_path, cache_key)

    cached_data = None
    cache_status = "missing"
    if os.path.exists(cache_path):
        try:
            with open(cache_path, "rb") as _cf:
                cached_data = _pkl_dxgc.load(_cf)
            current_n = _dxgc_count_on_disk(base_path, selected_hep_groups, selected_stage)
            cached_n = cached_data.get("n_individuals_on_disk", -1)
            if current_n != cached_n:
                cache_status = "stale"
                st.warning(
                    f"Patient count changed on disk ({cached_n} → {current_n}). "
                    "Cache will be recomputed automatically."
                )
            else:
                cache_status = "fresh"
                st.success(
                    f"✅ Cached results from {cached_data.get('created_at', '?')} | "
                    f"N={cached_data.get('n_individuals', '?')} patients | "
                    f"Grouping: {grouping_mode}. Showing cached data."
                )
        except Exception:
            cached_data = None
            cache_status = "missing"

    btn_label = "🔄 Reload & Recompute" if cache_status == "fresh" else "🚀 Load & Analyse"
    force_recompute = st.button(btn_label, key="dxgc_load")

    use_cache = (cache_status == "fresh") and not force_recompute
    run_full = force_recompute or cache_status != "fresh"

    if use_cache and cached_data is not None:
        group_data = cached_data["group_data"]
        for _gd in group_data.values():
            _gd.setdefault("eeg_per_ch", {})
            _gd.setdefault("eeg_mv", None)
        patient_info = cached_data["patient_info"]
        cat_counts = cached_data["cat_counts"]
        all_individuals_n = cached_data["n_individuals"]
        st.caption(f"Loaded from cache: {len(group_data)} groups, {all_individuals_n} total patients.")
        # Show patient table from cache
        st.subheader("\U0001f4cb Patient Diagnosis Table (Sample)")
        if cached_data.get("patient_table"):
            st.dataframe(pd.DataFrame(cached_data["patient_table"]), use_container_width=True)
        _show_diagnosis_category_venn(patient_info, cat_counts, "cached")
        # Show prevalence table from cache
        st.subheader("\U0001f4ca Diagnosis Category Prevalence")
        if cat_counts:
            all_ind_n = cached_data.get("n_individuals", 1)
            st.dataframe(pd.DataFrame([
                {"Diagnosis Category": c, "N Patients with diagnosis": n,
                 "% of loaded patients": f"{100*n/max(all_ind_n,1):.1f}%"}
                for c, n in sorted(cat_counts.items(), key=lambda x: -x[1])
            ]), use_container_width=True)
        # Group size table from cache
        st.subheader("\U0001f4ca Group Distribution")
        if cached_data.get("group_sizes"):
            st.dataframe(pd.DataFrame(cached_data["group_sizes"]), use_container_width=True)
        valid_groups = {k: [] for k in group_data.keys()}  # empty in cache path — no patient-level iteration
        palette = ['#1f77b4', '#d62728', '#2ca02c', '#9467bd', '#ff7f0e', '#17becf', '#8c564b', '#e377c2', '#bcbd22', '#7f7f7f']
        group_color = {g: palette[i % len(palette)] for i, g in enumerate(sorted(group_data))}
    elif not run_full:
        return
    else:
        pass  # fall through to full computation below

    if not use_cache:
        # Load psg_metadata for age/sex
        meta_paths = [
            "/storage/pblab_shared_data2/Nir/Cobrad/EDF_Format/The_Human_Sleep_Project/psg_metadata.csv",
            "/storage/pblab_shared_data2/Nir/Cobrad/EDF_Format/Young_The_Human_Sleep_Project/psg_metadata.csv",
        ]
        meta_frames = []
        for mp in meta_paths:
            if os.path.exists(mp):
                try:
                    meta_frames.append(pd.read_csv(mp, low_memory=False))
                except Exception:
                    pass
        meta_df = None
        if meta_frames:
            meta_df = pd.concat(meta_frames, ignore_index=True)
            meta_df['BDSPPatientID'] = pd.to_numeric(meta_df['BDSPPatientID'], errors='coerce')
            meta_df = meta_df.dropna(subset=['BDSPPatientID'])
            meta_df['BDSPPatientID'] = meta_df['BDSPPatientID'].astype(int)
            meta_df = meta_df.drop_duplicates(subset=['BDSPPatientID'])

        # Load HEP individuals
        all_individuals = []
        for grp in selected_hep_groups:
            with st.spinner(f"Loading {grp} / {selected_stage}..."):
                inds = get_group_individuals(grp, selected_stage, base_path)
                all_individuals.extend(inds or [])
        st.caption(f"Total patient records loaded: {len(all_individuals)}")
        if not all_individuals:
            st.error("No patient data loaded.")
            return

        selected_bdsp_ids = sorted({
            bdsp_id for bdsp_id in (_pid_to_bdsp_ehr(str(ind[0])) for ind in all_individuals)
            if bdsp_id is not None
        })
        with st.spinner("Loading EHR diagnosis data for selected patients..."):
            ehr_map = _load_ehr_diagnosis_map(target_ids=selected_bdsp_ids)
        st.caption(
            f"EHR diagnosis map loaded from I0006 + I0002: "
            f"{len(ehr_map)}/{len(selected_bdsp_ids)} selected patients with diagnosis records."
        )

        # Match each patient to EHR and meta
        patient_info = {}  # pid_str -> {bdsp_id, dx_names, categories, age, sex, age_group}
        for ind in all_individuals:
            pid = str(ind[0])
            bdsp_id = _pid_to_bdsp_ehr(pid)
            dx_names = ehr_map.get(bdsp_id, []) if bdsp_id is not None else []
            cats = _assign_categories(dx_names)
            age, sex, age_group = None, None, "Unknown"
            if meta_df is not None and bdsp_id is not None and bdsp_id in meta_df['BDSPPatientID'].values:
                row = meta_df[meta_df['BDSPPatientID'] == bdsp_id].iloc[0]
                age = row.get('AgeAtVisit', None)
                sex = row.get('SexDSC', None)
                age_group = ("<40" if pd.notna(age) and age < 40 else
                             "40-59" if pd.notna(age) and age < 60 else
                             "60-79" if pd.notna(age) and age < 80 else
                             "80+" if pd.notna(age) else "Unknown")
            patient_info[pid] = {
                "bdsp_id": bdsp_id,
                "dx_names": dx_names,
                "categories": cats,
                "age": age,
                "sex": sex,
                "age_group": age_group,
            }

        ehr_matched = sum(1 for v in patient_info.values() if v["dx_names"])
        st.caption(f"EHR-matched: {ehr_matched}/{len(all_individuals)} patients have diagnosis records.")

        # Patient-diagnosis sample table
        st.subheader("\U0001f4cb Patient Diagnosis Table (Sample)")
        st.caption("Shows each patient's EHR diagnoses. Categories are broad ICD groupings.")
        table_rows = []
        for pid, info in list(patient_info.items())[:50]:
            table_rows.append({
                "PID": pid,
                "BDSPPatientID": str(info["bdsp_id"]) if info["bdsp_id"] else "—",
                "Age": info["age"] if info["age"] is not None else "—",
                "Sex": info["sex"] or "—",
                "Diagnosis Categories": ", ".join(sorted(info["categories"])) or "None matched",
                "N specific diagnoses": len(info["dx_names"]),
                "Sample diagnoses": "; ".join(info["dx_names"][:3]) + ("..." if len(info["dx_names"]) > 3 else ""),
            })
        st.dataframe(pd.DataFrame(table_rows), use_container_width=True)

        # Diagnosis category counts table
        cat_counts = Counter()
        for info in patient_info.values():
            for c in info["categories"]:
                cat_counts[c] += 1
        _show_diagnosis_category_venn(patient_info, cat_counts, "loaded")
        cat_df = pd.DataFrame([
            {"Diagnosis Category": c, "N Patients with diagnosis": n,
             "% of loaded patients": f"{100*n/len(all_individuals):.1f}%"}
            for c, n in sorted(cat_counts.items(), key=lambda x: -x[1])
        ])
        st.subheader("\U0001f4ca Diagnosis Category Prevalence")
        st.dataframe(cat_df, use_container_width=True)

        # Build groups based on mode
        group_to_inds = {}

        if grouping_mode == "By Diagnosis Category":
            all_cats = sorted(cat_counts.keys(), key=lambda c: -cat_counts[c])
            if not all_cats:
                st.warning("No diagnosis categories found for loaded patients.")
                return
            for cat in all_cats:
                group_to_inds[cat] = [
                    ind for ind in all_individuals
                    if cat in patient_info[str(ind[0])]["categories"]
                ]
            no_diag = [ind for ind in all_individuals if not patient_info[str(ind[0])]["categories"]]
            if no_diag:
                group_to_inds["No Matched Diagnosis"] = no_diag

        elif grouping_mode == "By Specific Diagnosis":
            all_dx_names = sorted(set(dx for info in patient_info.values() for dx in info["dx_names"]))
            all_cats_for_filter = sorted(cat_counts.keys(), key=lambda c: -cat_counts[c])
            cat_filter = st.multiselect(
                "Filter by diagnosis category (optional):",
                options=all_cats_for_filter,
                default=[],
                key="dxgc_spec_cat_filter",
                help="Narrow the list below to diagnoses associated with these categories.",
            )
            if cat_filter:
                cat_filter_set = set(cat_filter)
                filtered_by_cat = sorted(set(
                    dx for info in patient_info.values()
                    if info["categories"] & cat_filter_set
                    for dx in info["dx_names"]
                ))
            else:
                filtered_by_cat = all_dx_names
            search_term = st.text_input("Search diagnosis names (filter):", key="dxgc_search")
            filtered_dx = [d for d in filtered_by_cat if search_term.lower() in d.lower()] if search_term else filtered_by_cat[:100]
            selected_dx = st.multiselect(
                "Select specific diagnoses (each = one group: patients WITH that diagnosis)",
                options=filtered_dx,
                default=[],
                key="dxgc_spec_dx",
            )
            if not selected_dx:
                st.info("Select specific diagnoses above to start comparison.")
                return
            exclusive_spec = st.checkbox(
                "Exclusive groups — exclude patients belonging to multiple selected groups",
                value=True,
                key="dxgc_spec_exclusive",
                help="When checked, a patient who has more than one of the selected diagnoses is excluded from all groups.",
            )
            selected_dx_set = set(selected_dx)
            for dx in selected_dx:
                if exclusive_spec:
                    group_to_inds[dx] = [
                        ind for ind in all_individuals
                        if dx in patient_info[str(ind[0])]["dx_names"]
                        and not (set(patient_info[str(ind[0])]["dx_names"]) & (selected_dx_set - {dx}))
                    ]
                else:
                    group_to_inds[dx] = [
                        ind for ind in all_individuals
                        if dx in patient_info[str(ind[0])]["dx_names"]
                    ]

        elif grouping_mode == "Diagnosis Comorbidity Subgroups":
            base_cat = comorbidity_base_category
            comorb_cats = [cat for cat in comorbidity_categories if cat != base_cat]
            base_inds = [
                ind for ind in all_individuals
                if base_cat in patient_info[str(ind[0])]["categories"]
            ]
            if not base_inds:
                st.warning(f"No loaded patients have the index diagnosis category: {base_cat}.")
                return

            comorb_counts = {
                cat: sum(
                    1 for ind in base_inds
                    if cat in patient_info[str(ind[0])]["categories"]
                )
                for cat in comorb_cats
            }
            active_comorb_cats = [
                cat for cat in comorb_cats
                if comorb_counts.get(cat, 0) >= int(comorbidity_min_disease_count)
            ]
            ignored_comorb_cats = [
                cat for cat in comorb_cats
                if comorb_counts.get(cat, 0) < int(comorbidity_min_disease_count)
            ]
            if ignored_comorb_cats:
                st.caption(
                    f"Ignoring {len(ignored_comorb_cats)} rare comorbidity categories with "
                    f"< {int(comorbidity_min_disease_count)} patients. Patients are retained and "
                    "assigned using the remaining larger comorbidity categories."
                )
                st.dataframe(
                    pd.DataFrame([
                        {"Ignored comorbidity": cat, "N patients within index diagnosis": comorb_counts.get(cat, 0)}
                        for cat in sorted(ignored_comorb_cats, key=lambda c: (comorb_counts.get(c, 0), c))
                    ]),
                    use_container_width=True,
                )
            if not active_comorb_cats:
                st.warning(
                    f"No selected comorbidity categories have >= {int(comorbidity_min_disease_count)} "
                    f"patients among {base_cat} patients. All patients will be grouped as {base_cat} only."
                )

            for ind in base_inds:
                pid = str(ind[0])
                patient_cats = set(patient_info[pid]["categories"])
                combo = tuple(cat for cat in active_comorb_cats if cat in patient_cats)
                if combo:
                    label = f"{base_cat} + " + " + ".join(combo)
                else:
                    label = f"{base_cat} only"
                group_to_inds.setdefault(label, []).append(ind)

            subgroup_rows = []
            for label, inds in sorted(group_to_inds.items(), key=lambda item: (-len(item[1]), item[0])):
                subgroup_rows.append({
                    "Subgroup": label,
                    "N Patients": len(inds),
                    "% of index diagnosis": f"{100*len(inds)/max(len(base_inds), 1):.1f}%",
                    "Included": "Yes" if len(inds) >= int(min_group_size) else f"No (< {int(min_group_size)})",
                })
            st.caption(
                f"Comorbidity subgroups are mutually exclusive exact combinations among "
                f"{len(base_inds)} patients with {base_cat}. "
                "Unselected and rare diagnosis categories are ignored when assigning subgroup labels, "
                "but their patients remain in the broader subgroup."
            )
            st.dataframe(pd.DataFrame(subgroup_rows), use_container_width=True)

        elif grouping_mode == "By Age Group":
            for ind in all_individuals:
                ag = patient_info[str(ind[0])]["age_group"]
                group_to_inds.setdefault(ag, []).append(ind)

        elif grouping_mode == "Exclusive Diagnosis Category":
            all_cats_sorted = sorted(cat_counts.keys(), key=lambda c: -cat_counts[c])
            if not all_cats_sorted:
                st.warning("No diagnosis categories found for loaded patients.")
                return
            exclusive_strategy = st.radio(
                "Exclusive grouping strategy",
                ["Auto-find largest non-overlapping groups", "Only pure single-category overlap removal"],
                horizontal=True,
                key="dxgc_excl_strategy",
            )

            if exclusive_strategy == "Auto-find largest non-overlapping groups":
                import itertools as _itertools
                c1, c2 = st.columns(2)
                with c1:
                    max_auto_groups = st.number_input(
                        "Max non-overlapping groups", 2, 20, 6, 1, key="dxgc_excl_auto_max_groups"
                    )
                with c2:
                    max_candidate_cats = st.number_input(
                        "Candidate categories to search", 2, 50, min(20, max(2, len(all_cats_sorted))), 1,
                        key="dxgc_excl_auto_candidates",
                    )
                k = int(max_auto_groups)
                # Cap search space to avoid combinatorial explosion (C(20,8)~125k is fine)
                search_cats = all_cats_sorted[:min(int(max_candidate_cats), 20)]

                # Pre-compute patient-id sets per category
                cat_pid_sets = {
                    cat: frozenset(str(ind[0]) for ind in all_individuals
                                   if cat in patient_info[str(ind[0])]["categories"])
                    for cat in search_cats
                }

                # Find combo of exactly k_try categories maximising min(exclusive_count).
                # Exclusive count = patients in that category and NO other selected category.
                best_combo = None
                best_score = -1
                best_excl_counts = None
                found_k = None

                for k_try in range(min(k, len(search_cats)), 1, -1):
                    for combo in _itertools.combinations(search_cats, k_try):
                        combo_sets = [cat_pid_sets[c] for c in combo]
                        counts = []
                        for i, c in enumerate(combo):
                            others = frozenset().union(*(combo_sets[j] for j in range(k_try) if j != i))
                            counts.append(len(combo_sets[i] - others))
                        score = min(counts)
                        if score >= int(min_group_size) and score > best_score:
                            best_score = score
                            best_combo = combo
                            best_excl_counts = counts
                    if best_combo is not None:
                        found_k = k_try
                        break

                if best_combo is None:
                    st.warning(
                        "Could not find non-overlapping diagnosis groups that meet the minimum size. "
                        "Lower Min patients per group or increase Candidate categories to search."
                    )
                    return

                excl_info = []
                combo_sets = [cat_pid_sets[c] for c in best_combo]
                for i, cat in enumerate(best_combo):
                    others = frozenset().union(*(combo_sets[j] for j in range(found_k) if j != i))
                    excl_pids = combo_sets[i] - others
                    excl_inds = [ind for ind in all_individuals if str(ind[0]) in excl_pids]
                    group_to_inds[cat] = excl_inds
                    excl_info.append({
                        "Category": cat,
                        "Total with diagnosis": len(combo_sets[i]),
                        "Exclusive patients": len(excl_inds),
                        "Excluded (overlap with others)": len(combo_sets[i]) - len(excl_inds),
                    })

                st.caption(
                    f"**Auto exclusive mode:** searched top-{len(search_cats)} categories, "
                    f"found best {found_k}-group combo maximising the smallest group. "
                    "Patients belonging to multiple selected categories are excluded."
                )
                st.dataframe(pd.DataFrame(excl_info), use_container_width=True)
            else:
                top_n_excl = st.number_input(
                    "Top N categories to consider", 2, 20, 5, 1, key="dxgc_excl_topn"
                )
                top_cats_excl = all_cats_sorted[:int(top_n_excl)]
                selected_excl_cats = st.multiselect(
                    "Select diagnosis categories — each group = patients with ONLY that category (no overlap with other selected categories)",
                    options=all_cats_sorted,
                    default=top_cats_excl,
                    key="dxgc_excl_cats",
                )
                if not selected_excl_cats:
                    st.warning("Select at least one category.")
                    return
                selected_excl_set = set(selected_excl_cats)
                excl_info = []
                for cat in selected_excl_cats:
                    total_with_cat = sum(
                        1 for ind in all_individuals
                        if cat in patient_info[str(ind[0])]["categories"]
                    )
                    exclusive_inds = [
                        ind for ind in all_individuals
                        if cat in patient_info[str(ind[0])]["categories"]
                        and not (set(patient_info[str(ind[0])]["categories"]) & (selected_excl_set - {cat}))
                    ]
                    group_to_inds[cat] = exclusive_inds
                    excl_info.append({
                        "Category": cat,
                        "Total with diagnosis": total_with_cat,
                        "Exclusive (no overlap)": len(exclusive_inds),
                        "Excluded (overlap)": total_with_cat - len(exclusive_inds),
                    })
                st.caption(
                    "**Pure exclusive mode:** each group contains only patients with that diagnosis "
                    "and none of the other selected diagnoses."
                )
                st.dataframe(pd.DataFrame(excl_info), use_container_width=True)

        else:  # By Sex
            for ind in all_individuals:
                sex = str(patient_info[str(ind[0])]["sex"] or "Unknown")
                group_to_inds.setdefault(sex, []).append(ind)

        # Filter by min size
        all_sizes = sorted(group_to_inds.items(), key=lambda x: -len(x[1]))
        valid_groups = {k: v for k, v in group_to_inds.items() if len(v) >= int(min_group_size)}
        if grouping_mode == "Diagnosis Comorbidity Subgroups":
            valid_groups = dict(
                sorted(valid_groups.items(), key=lambda item: (-len(item[1]), item[0]))[
                    :int(comorbidity_max_subgroups)
                ]
            )

        st.subheader("\U0001f4ca Group Distribution")
        def _included_label(group_name, inds):
            if group_name in valid_groups:
                return "Yes"
            if len(inds) < int(min_group_size):
                return f"No (< {int(min_group_size)})"
            if grouping_mode == "Diagnosis Comorbidity Subgroups":
                return f"No (outside top {int(comorbidity_max_subgroups)})"
            return "No"

        size_df = pd.DataFrame([
            {"Group": k, "N Patients": len(v),
             "Included": _included_label(k, v)}
            for k, v in all_sizes
        ])
        st.dataframe(size_df, use_container_width=True)

        if not valid_groups:
            st.info(f"No groups meet minimum size {int(min_group_size)}. Reduce threshold or select more groups.")
            return

        palette = ['#1f77b4', '#d62728', '#2ca02c', '#9467bd', '#ff7f0e', '#17becf', '#8c564b', '#e377c2', '#bcbd22', '#7f7f7f']
        group_color = {g: palette[i % len(palette)] for i, g in enumerate(sorted(valid_groups))}

        def _extract_matrices(inds, clean_ecg=True):
            eeg_by_ch = {}
            ecg_traces = []
            times_ref = None
            ch_cnt = Counter()
            for ind in inds:
                try:
                    unpacked = list(ind) + [None] * 7
                    pid, hep_data, times, ch_names, _, ecg_hep = unpacked[:6]
                except Exception:
                    continue
                if hep_data is None or times is None or ch_names is None:
                    continue
                if times_ref is None:
                    times_ref = np.asarray(times, dtype=float)
                eeg_arr = apply_ecg_regression(np.asarray(hep_data, dtype=float), ecg_hep) if (clean_ecg and ecg_hep is not None) else np.asarray(hep_data, dtype=float)
                for ch in ch_names:
                    if _is_eeg_channel(ch):
                        idx = ch_names.index(ch)
                        if idx < eeg_arr.shape[0]:
                            eeg_by_ch.setdefault(ch, []).append(eeg_arr[idx])
                            ch_cnt[ch] += 1
                if ecg_hep is not None:
                    ecg_t = np.asarray(ecg_hep, dtype=float).squeeze()
                    if ecg_t.ndim == 1:
                        ecg_traces.append(ecg_t * 1e6)
            if times_ref is None:
                return None, None, None, [], {}, None
            n_p = len(inds)
            common = [c for c, cnt in ch_cnt.items() if cnt >= max(2, n_p * 0.5)]
            if not common:
                return None, None, times_ref, [], {}, None
            min_len = min(len(eeg_by_ch[c][0]) for c in common if eeg_by_ch.get(c))
            pat_eeg = []
            for i in range(n_p):
                tr = [eeg_by_ch[c][i][:min_len] for c in common if i < len(eeg_by_ch.get(c, []))]
                if tr:
                    pat_eeg.append(np.mean(tr, axis=0))
            if not pat_eeg:
                return None, None, times_ref[:min_len], common, {}, None
            eeg_mat = np.vstack(pat_eeg)
            ecg_mat = None
            if ecg_traces:
                ve = [t[:min_len] for t in ecg_traces if len(t) >= min_len]
                if ve:
                    ecg_mat = np.vstack(ve)
            # Per-channel matrices: {ch: (n_patients_with_ch, min_len)}
            eeg_per_ch = {}
            for c in common:
                rows = [r[:min_len] for r in eeg_by_ch.get(c, [])]
                if rows:
                    eeg_per_ch[c] = np.vstack(rows)
            # Multivariate matrix: each patient row = concatenation of all common channels
            mv_rows = []
            for i in range(len(pat_eeg)):
                row_parts = [eeg_by_ch[c][i][:min_len] for c in common if i < len(eeg_by_ch.get(c, []))]
                if len(row_parts) == len(common):
                    mv_rows.append(np.concatenate(row_parts))
            eeg_mv = np.vstack(mv_rows) if mv_rows else None
            return eeg_mat, ecg_mat, times_ref[:min_len], common, eeg_per_ch, eeg_mv

        group_data = {}
        for cat, inds in valid_groups.items():
            with st.spinner(f"Building matrices for {cat} (N={len(inds)})..."):
                em, ecm, t, chs, per_ch, mv = _extract_matrices(inds, clean_ecg=apply_ecg_reg)
                if em is not None:
                    per_ch_raw = per_ch
                    if apply_ecg_reg:
                        _, _, _, _, per_ch_raw, _ = _extract_matrices(inds, clean_ecg=False)
                    group_data[cat] = {
                        "eeg": em, "ecg": ecm, "times": t, "channels": chs, "n": len(inds),
                        "eeg_per_ch": per_ch,
                        "eeg_per_ch_raw": per_ch_raw,
                        "eeg_mv": mv,
                    }

        if not group_data:
            st.error("No valid EEG data for any group.")
            return

        # Save to cache
        try:
            import pickle as _pkl_dxgc2
            from datetime import datetime as _dt_dxgc2
            patient_table_cache = []
            for pid, info in list(patient_info.items())[:50]:
                patient_table_cache.append({
                    "PID": pid,
                    "BDSPPatientID": str(info["bdsp_id"]) if info["bdsp_id"] else "—",
                    "Age": info["age"] if info["age"] is not None else "—",
                    "Sex": info["sex"] or "—",
                    "Diagnosis Categories": ", ".join(sorted(info["categories"])) or "None matched",
                    "N specific diagnoses": len(info["dx_names"]),
                    "Sample diagnoses": "; ".join(info["dx_names"][:3]) + ("..." if len(info["dx_names"]) > 3 else ""),
                })
            group_sizes_cache = [
                {"Group": k, "N Patients": len(v),
                 "Included": _included_label(k, v)}
                for k, v in sorted(group_to_inds.items(), key=lambda x: -len(x[1]))
            ]
            cache_payload = {
                "group_data": group_data,
                "patient_info": patient_info,
                "cat_counts": cat_counts,
                "n_individuals": len(all_individuals),
                "n_individuals_on_disk": _dxgc_count_on_disk(base_path, selected_hep_groups, selected_stage),
                "created_at": _dt_dxgc2.now().strftime("%Y-%m-%d %H:%M"),
                "patient_table": patient_table_cache,
                "group_sizes": group_sizes_cache,
                "cache_key": cache_key,
            }
            with open(cache_path, "wb") as _cf2:
                _pkl_dxgc2.dump(cache_payload, _cf2)
            st.caption(f"✅ Results cached to `{os.path.basename(cache_path)}`.")
        except Exception as _ce:
            st.caption(f"Cache save failed: {_ce}")

    group_data, _ = _filter_diagnosis_category_groups(group_data, cat_counts)
    if not group_data:
        return
    palette = ['#1f77b4', '#d62728', '#2ca02c', '#9467bd', '#ff7f0e', '#17becf', '#8c564b', '#e377c2', '#bcbd22', '#7f7f7f']
    group_color = {g: palette[i % len(palette)] for i, g in enumerate(group_data)}

    # PLOT A: EEG HEP
    st.markdown("---")
    st.subheader("\U0001f4c8 A. EEG HEP by Diagnosis Group (Median +/- MAD)")
    st.markdown(
        "**Interpretation:** EEG response time-locked to each heartbeat R-peak (t=0). "
        "Shaded = MAD (robust spread). Differences between diagnosis groups reflect "
        "how cardiac diagnosis modulates brain-heart coupling."
    )
    eeg_sum = []
    fig_eeg, ax_eeg = plt.subplots(figsize=(14, 5))
    ax_eeg.axvline(0, color='red', ls='--', alpha=0.6, label='R-peak')
    ax_eeg.axhline(0, color='black', lw=0.5, alpha=0.3)
    for cat, gd in group_data.items():
        mat = gd["eeg"]; t = gd["times"]
        ml = min(mat.shape[1], len(t))
        mat_uv = _eeg_display_matrix(mat[:, :ml])
        med, mad_v = median_and_mad(mat_uv, axis=0)
        ax_eeg.plot(t[:ml], med, color=group_color[cat], lw=2, label=f"{cat} (N={gd['n']})")
        ax_eeg.fill_between(t[:ml], med - mad_v, med + mad_v, color=group_color[cat], alpha=0.18)
        eeg_sum.append({"Group": cat, "N": mat.shape[0],
            f"Peak |Amp| ({EEG_DISPLAY_UNIT})": f"{float(np.max(np.abs(med))):.4f}",
            "Peak Time (s)": f"{float(t[np.argmax(np.abs(med))]):.3f}",
            f"Mean MAD ({EEG_DISPLAY_UNIT})": f"{float(np.mean(mad_v)):.4f}",
            "N EEG ch": len(gd["channels"])})
    ax_eeg.set_xlabel("Time relative to R-peak (s)")
    ax_eeg.set_ylabel(f"EEG amplitude ({EEG_DISPLAY_UNIT})")
    ax_eeg.set_title(f"EEG HEP — {grouping_mode}", fontsize=13, fontweight='bold')
    ax_eeg.legend(fontsize=9, loc='best')
    ax_eeg.grid(True, alpha=0.25)
    fig_eeg.tight_layout()
    _show_non_ica_figure(fig_eeg, use_container_width=True)
    plt.close(fig_eeg)
    st.markdown("**EEG HEP summary table:**")
    st.caption("Amplitude and MAD are computed from the raw EEG HEP after converting stored volts to µV; no z-score scaling is applied here.")
    st.dataframe(pd.DataFrame(eeg_sum), use_container_width=True)

    # PLOT A2: Per-Channel EEG HEP
    st.markdown("---")
    st.subheader("📊 A2. Per-Channel EEG HEP by Diagnosis Group")
    st.markdown(
        "**Interpretation:** Same data as plot A broken down per EEG channel. "
        "Each subplot shows median ± MAD for each diagnosis group. Colors match plot A."
    )
    _per_ch_groups = {cat: gd["eeg_per_ch"] for cat, gd in group_data.items() if gd.get("eeg_per_ch")}
    if _per_ch_groups:
        _common_chs = sorted(set.intersection(*[set(d.keys()) for d in _per_ch_groups.values()])) if len(_per_ch_groups) > 1 else sorted(next(iter(_per_ch_groups.values())).keys())
        _common_chs = [ch for ch in _common_chs if str(ch).strip().lower() != "sao2"]
        if not _common_chs:
            st.info("No channels common to all diagnosis groups.")
        else:
            _n_chs = len(_common_chs)
            _ncols = min(4, _n_chs)
            _nrows = math.ceil(_n_chs / _ncols)
            _fig_pc, _axes_pc = plt.subplots(_nrows, _ncols, figsize=(5 * _ncols, 3.5 * _nrows), squeeze=False)
            _axes_flat = _axes_pc.flatten()
            for _i_ch, _ch in enumerate(_common_chs):
                _ax = _axes_flat[_i_ch]
                _ax.axvline(0, color='red', ls='--', alpha=0.5, lw=0.8)
                _ax.axhline(0, color='black', lw=0.4, alpha=0.3)
                for _cat, _per_ch in _per_ch_groups.items():
                    if _ch not in _per_ch:
                        continue
                    _mat = _per_ch[_ch]
                    _t = group_data[_cat]["times"]
                    _ml = min(_mat.shape[1], len(_t))
                    _mat_uv = _eeg_display_matrix(_mat[:, :_ml])
                    _med, _mad_v = median_and_mad(_mat_uv, axis=0)
                    _col = group_color.get(_cat, 'gray')
                    _ax.plot(_t[:_ml], _med, color=_col, lw=1.5, label=f"{_cat} (N={_mat.shape[0]})")
                    _ax.fill_between(_t[:_ml], _med - _mad_v, _med + _mad_v, color=_col, alpha=0.15)
                _ax.set_title(_ch, fontsize=8.5, fontweight='bold', pad=4)
                _ax.grid(True, alpha=0.2)
                _ax.tick_params(labelsize=6)
            for _j_ch in range(_n_chs, len(_axes_flat)):
                _axes_flat[_j_ch].set_visible(False)
            _handles, _labels_leg = _axes_flat[0].get_legend_handles_labels()
            if _handles:
                _fig_pc.legend(_handles, _labels_leg, loc='lower center',
                               bbox_to_anchor=(0.5, -0.02), ncol=len(_per_ch_groups),
                               fontsize=9, frameon=False)
            _cleaned_label = " | ECG regression applied" if apply_ecg_reg else ""
            _fig_pc.suptitle(
                f"Per-Channel EEG HEP — {grouping_mode}  ({_n_chs} channels){_cleaned_label}",
                fontsize=12, fontweight='bold',
            )
            _fig_pc.tight_layout()
            _show_non_ica_figure(_fig_pc, use_container_width=True)
            plt.close(_fig_pc)

            # Raw per-channel (no ECG regression) — shown only when ECG regression is active
            _raw_per_ch_groups = {
                cat: gd["eeg_per_ch_raw"]
                for cat, gd in group_data.items()
                if gd.get("eeg_per_ch_raw")
            }
            if apply_ecg_reg and _raw_per_ch_groups:
                _r_common = (
                    sorted(set.intersection(*[set(d.keys()) for d in _raw_per_ch_groups.values()]))
                    if len(_raw_per_ch_groups) > 1
                    else sorted(next(iter(_raw_per_ch_groups.values())).keys())
                )
                _r_common = [ch for ch in _r_common if str(ch).strip().lower() != "sao2"]
                if _r_common:
                    _rn = len(_r_common)
                    _rc = min(4, _rn)
                    _rr = math.ceil(_rn / _rc)
                    _fig_raw, _axes_raw = plt.subplots(_rr, _rc, figsize=(5 * _rc, 3.5 * _rr), squeeze=False)
                    _axes_raw_flat = _axes_raw.flatten()
                    for _ri, _rch in enumerate(_r_common):
                        _ax_r = _axes_raw_flat[_ri]
                        _ax_r.axvline(0, color='red', ls='--', alpha=0.5, lw=0.8)
                        _ax_r.axhline(0, color='black', lw=0.4, alpha=0.3)
                        for _rcat, _rpc in _raw_per_ch_groups.items():
                            if _rch not in _rpc:
                                continue
                            _rm = _rpc[_rch]
                            _rt = group_data[_rcat]["times"]
                            _rml = min(_rm.shape[1], len(_rt))
                            _rm_uv = _eeg_display_matrix(_rm[:, :_rml])
                            _rmed, _rmad = median_and_mad(_rm_uv, axis=0)
                            _rcol = group_color.get(_rcat, 'gray')
                            _ax_r.plot(_rt[:_rml], _rmed, color=_rcol, lw=1.5, label=f"{_rcat} (N={_rm.shape[0]})")
                            _ax_r.fill_between(_rt[:_rml], _rmed - _rmad, _rmed + _rmad, color=_rcol, alpha=0.15)
                        _ax_r.set_title(_rch, fontsize=8.5, fontweight='bold', pad=4)
                        _ax_r.grid(True, alpha=0.2)
                        _ax_r.tick_params(labelsize=6)
                    for _rj in range(_rn, len(_axes_raw_flat)):
                        _axes_raw_flat[_rj].set_visible(False)
                    _rh, _rl = _axes_raw_flat[0].get_legend_handles_labels()
                    if _rh:
                        _fig_raw.legend(_rh, _rl, loc='lower center',
                                        bbox_to_anchor=(0.5, -0.02), ncol=len(_raw_per_ch_groups),
                                        fontsize=9, frameon=False)
                    _fig_raw.suptitle(
                        f"Per-Channel EEG HEP — {grouping_mode}  ({_rn} channels) | Raw (no ECG regression)",
                        fontsize=12, fontweight='bold',
                    )
                    _fig_raw.tight_layout()
                    _show_non_ica_figure(_fig_raw, use_container_width=True)
                    plt.close(_fig_raw)
            elif apply_ecg_reg:
                st.info("Raw per-channel data not in cache. Rebuild the cache to enable raw comparison.")
    else:
        st.info("Per-channel data not available. Please rebuild the cache.")

    # PLOT B: ECG HEP
    st.markdown("---")
    st.subheader("\U0001fac0 B. ECG HEP by Diagnosis Group (Median +/- MAD)")
    st.markdown("**Interpretation:** Average ECG around each R-peak. Cardiac morphology differences between diagnosis groups (e.g. broader QRS in heart failure, ST changes in CAD).")
    ecg_groups = {c: gd for c, gd in group_data.items() if gd.get("ecg") is not None}
    if ecg_groups:
        ecg_sum = []
        fig_ecg, ax_ecg = plt.subplots(figsize=(14, 5))
        ax_ecg.axvline(0, color='red', ls='--', alpha=0.6, label='R-peak')
        ax_ecg.axhline(0, color='black', lw=0.5, alpha=0.3)
        for cat, gd in ecg_groups.items():
            mat = gd["ecg"]; t = gd["times"]
            ml = min(mat.shape[1], len(t))
            med, mad_v = median_and_mad(mat[:, :ml], axis=0)
            ax_ecg.plot(t[:ml], med, color=group_color[cat], lw=2, label=f"{cat} (N={mat.shape[0]})")
            ax_ecg.fill_between(t[:ml], med - mad_v, med + mad_v, color=group_color[cat], alpha=0.18)
            ecg_sum.append({"Group": cat, "N ECG": mat.shape[0],
                "Peak ECG Amp (uV)": f"{float(np.max(np.abs(med))):.1f}",
                "Peak Time (s)": f"{float(t[np.argmax(np.abs(med))]):.3f}"})
        ax_ecg.set_xlabel("Time (s)"); ax_ecg.set_ylabel("ECG amplitude (uV)")
        ax_ecg.set_title(f"ECG HEP — {grouping_mode}", fontsize=13, fontweight='bold')
        ax_ecg.legend(fontsize=9); ax_ecg.grid(True, alpha=0.25)
        fig_ecg.tight_layout()
        _show_non_ica_figure(fig_ecg, use_container_width=True); plt.close(fig_ecg)
        st.dataframe(pd.DataFrame(ecg_sum), use_container_width=True)
    else:
        st.info("No ECG data available for any group.")

    # PLOT C: EEG vs 0
    st.markdown("---")
    st.subheader("\U0001f4ca C. EEG HEP vs Baseline (0) — Each Group")
    st.markdown("**Interpretation:** Shaded time points = significantly non-zero EEG (p<0.05, one-sample t-test, uncorrected). Wider/stronger shading = stronger heartbeat-evoked potential.")
    vs0_rows = []
    n_grps = len(group_data)
    fig_vs0, axes_vs0 = plt.subplots(1, n_grps, figsize=(max(14, 4*n_grps), 4), sharey=True)
    if n_grps == 1:
        axes_vs0 = [axes_vs0]
    for ax_v, (cat, gd) in zip(axes_vs0, group_data.items()):
        mat = gd["eeg"]; t = gd["times"]
        ml = min(mat.shape[1], len(t))
        t_st, p_v = stats.ttest_1samp(mat[:, :ml], 0, axis=0)
        mat_uv = _eeg_display_matrix(mat[:, :ml])
        med, _ = median_and_mad(mat_uv, axis=0)
        sig = p_v < 0.05
        ax_v.axhline(0, color='black', lw=0.5, alpha=0.3); ax_v.axvline(0, color='red', ls='--', alpha=0.5)
        ax_v.plot(t[:ml], med, color=group_color[cat], lw=2)
        ax_v.fill_between(t[:ml], med, 0, where=sig, color=group_color[cat], alpha=0.4, label='p<0.05')
        ax_v.set_title(f"{cat}\nN={gd['n']}", fontsize=9, fontweight='bold')
        ax_v.set_xlabel("Time (s)"); ax_v.grid(True, alpha=0.2); ax_v.legend(fontsize=7)
        ns = int(np.sum(sig))
        st_f = float(t[np.where(sig)[0][0]]) if ns > 0 else float('nan')
        vs0_rows.append({"Group": cat, "N": mat.shape[0],
            "N sig. tp (p<0.05)": ns,
            "First sig. time (s)": f"{st_f:.3f}" if ns > 0 else "none",
            "Max |t|": f"{float(np.max(np.abs(t_st))):.2f}",
            "Min p": f"{float(np.min(p_v)):.4f}"})
    axes_vs0[0].set_ylabel(f"EEG amplitude ({EEG_DISPLAY_UNIT})")
    fig_vs0.suptitle("EEG HEP vs 0 (one-sample t-test, shaded = p<0.05 uncorrected)", fontsize=11, fontweight='bold')
    fig_vs0.tight_layout()
    _show_non_ica_figure(fig_vs0, use_container_width=True); plt.close(fig_vs0)
    st.markdown("**vs-0 statistics table:**")
    st.dataframe(pd.DataFrame(vs0_rows), use_container_width=True)

    # PLOT D: EEG-ECG Correlation
    st.markdown("---")
    st.subheader("\U0001f517 D. EEG–ECG Cross-Correlation per Patient")
    st.markdown(
        "**Interpretation:** Pearson r between EEG and ECG HEP traces per patient. "
        "High r = EEG mirrors ECG shape (cardiac artifact or tight cardiac-neural coupling). "
        "After ECG regression (if enabled), residual r reflects true cardiac-neural coupling. "
        "Kruskal-Wallis: overall difference; Mann-Whitney: pairwise."
    )
    corr_res = {}
    for cat, gd in group_data.items():
        em = gd["eeg"]; ecm = gd.get("ecg")
        if ecm is None: continue
        ml = min(em.shape[1], ecm.shape[1])
        np_ = min(em.shape[0], ecm.shape[0])
        corrs = []
        for i in range(np_):
            e, c = em[i, :ml], ecm[i, :ml]
            if np.std(e) > 0 and np.std(c) > 0:
                corrs.append(float(np.corrcoef(e, c)[0, 1]))
        if corrs: corr_res[cat] = corrs
    if corr_res:
        fig_cor, ax_cor = plt.subplots(figsize=(max(6, len(corr_res)*2.5), 5))
        cats_c = list(corr_res.keys())
        bp = ax_cor.boxplot([corr_res[c] for c in cats_c], positions=list(range(len(cats_c))), patch_artist=True)
        for patch, cat in zip(bp['boxes'], cats_c):
            patch.set_facecolor(group_color.get(cat, '#aaa')); patch.set_alpha(0.7)
        ax_cor.axhline(0, color='black', lw=0.5, ls='--', alpha=0.5)
        ax_cor.set_xticks(list(range(len(cats_c))))
        ax_cor.set_xticklabels([f"{c}\n(N={len(corr_res[c])})" for c in cats_c], fontsize=8)
        ax_cor.set_ylabel("Pearson r"); ax_cor.grid(True, alpha=0.2)
        if len(corr_res) >= 2:
            kw_h, kw_p = stats.kruskal(*corr_res.values())
            ax_cor.set_title(f"EEG–ECG Pearson r | KW H={kw_h:.2f}, p={kw_p:.4f}", fontsize=11, fontweight='bold')
        fig_cor.tight_layout()
        _show_non_ica_figure(fig_cor, use_container_width=True); plt.close(fig_cor)
        c_sum = [{"Group": c, "N": len(corr_res[c]),
            "Median r": f"{float(np.median(corr_res[c])):.3f}",
            "Mean r": f"{float(np.mean(corr_res[c])):.3f}",
            "Std r": f"{float(np.std(corr_res[c])):.3f}"} for c in cats_c]
        st.markdown("**Correlation summary:**"); st.dataframe(pd.DataFrame(c_sum), use_container_width=True)
        if len(cats_c) >= 2:
            pw = []
            for i2 in range(len(cats_c)):
                for j2 in range(i2+1, len(cats_c)):
                    ca, cb = cats_c[i2], cats_c[j2]
                    u, p_mw = stats.mannwhitneyu(corr_res[ca], corr_res[cb], alternative='two-sided')
                    pw.append({"Group A": ca, "Group B": cb, "Mann-Whitney U": f"{u:.1f}",
                        "p-value": f"{p_mw:.4f}", "Sig (p<0.05)": "Yes" if p_mw < 0.05 else "No"})
            st.markdown("**Pairwise correlation comparison:**"); st.dataframe(pd.DataFrame(pw), use_container_width=True)
    else:
        st.info("EEG–ECG correlation not available (no ECG data).")

    # PLOT E: ECG Modulation
    st.markdown("---")
    st.subheader("\U0001f4e1 E. ECG->EEG Modulation (beta over Time)")
    st.markdown("**Interpretation:** Regression beta predicting EEG from ECG amplitude at each time point across patients. Shaded = p<0.05. Larger |beta| = stronger ECG influence on EEG.")
    mod_rows = []
    fig_mod, ax_mod = plt.subplots(figsize=(14, 5))
    ax_mod.axhline(0, color='black', lw=0.5, alpha=0.3); ax_mod.axvline(0, color='red', ls='--', alpha=0.5)
    for cat, gd in group_data.items():
        em = gd["eeg"]; ecm = gd.get("ecg")
        if ecm is None: continue
        ml = min(em.shape[1], ecm.shape[1]); np_ = min(em.shape[0], ecm.shape[0])
        t = gd["times"][:ml]
        betas = np.zeros(ml); pvals_m = np.ones(ml)
        ec_c = ecm[:np_, :ml] - ecm[:np_, :ml].mean(axis=0, keepdims=True)
        ee_c = em[:np_, :ml] - em[:np_, :ml].mean(axis=0, keepdims=True)
        for ti in range(ml):
            x_t = ec_c[:, ti]; y_t = ee_c[:, ti]; ss = np.sum(x_t**2)
            if ss > 0:
                b = np.dot(x_t, y_t) / ss; betas[ti] = b
                res = y_t - b*x_t; se = np.sqrt(np.sum(res**2) / max(1, np_-2) / ss)
                if se > 0: pvals_m[ti] = float(2*(1-stats.t.cdf(abs(b/se), df=np_-2)))
        color = group_color[cat]
        ax_mod.plot(t, betas, color=color, lw=2, label=f"{cat} (N={np_})")
        ax_mod.fill_between(t, betas, 0, where=pvals_m < 0.05, color=color, alpha=0.25)
        mod_rows.append({"Group": cat, "N": np_,
            "Peak |beta|": f"{float(np.max(np.abs(betas))):.4f}",
            "Peak beta Time (s)": f"{float(t[np.argmax(np.abs(betas))]):.3f}",
            "N sig. tp (p<0.05)": int(np.sum(pvals_m < 0.05))})
    ax_mod.set_xlabel("Time relative to R-peak (s)"); ax_mod.set_ylabel("ECG->EEG beta")
    ax_mod.set_title(f"ECG->EEG Modulation — {grouping_mode}", fontsize=13, fontweight='bold')
    ax_mod.legend(fontsize=9); ax_mod.grid(True, alpha=0.25)
    fig_mod.tight_layout()
    _show_non_ica_figure(fig_mod, use_container_width=True); plt.close(fig_mod)
    if mod_rows:
        st.markdown("**Modulation summary:**"); st.dataframe(pd.DataFrame(mod_rows), use_container_width=True)

    # PLOT F: Pairwise EEG comparison
    st.markdown("---")
    st.subheader("F. Pairwise Between-Group EEG Comparison (Cluster Permutation)")
    st.markdown(
        "**Interpretation:** Cluster permutation test comparing EEG HEP waveforms between each pair of groups. "
        "Top panel: group medians ± MAD and A−B difference (dashed black). "
        "Bottom panel: T-statistic — dashed lines = ±clustering threshold (p=0.05). "
        "Orange shading = significant cluster windows (cluster-mass permutation, p<0.05). Tables show windows."
    )
    grp_list = [c for c in group_data if group_data[c]["eeg"] is not None]
    pw_summary = []
    if len(grp_list) < 2:
        st.info("Need >= 2 groups for pairwise comparison.")
    else:
        for i_p in range(len(grp_list)):
            for j_p in range(i_p+1, len(grp_list)):
                ca, cb = grp_list[i_p], grp_list[j_p]
                ea = group_data[ca]["eeg"]; eb = group_data[cb]["eeg"]
                ta = group_data[ca]["times"]; tb = group_data[cb]["times"]
                ml = min(ea.shape[1], eb.shape[1], len(ta), len(tb))
                ea = ea[:, :ml]; eb = eb[:, :ml]; t_c = ta[:ml]
                st.markdown(f"#### {ca} vs {cb}")
                with st.spinner(f"Cluster permutation: {ca} vs {cb}..."):
                    try:
                        sig_wins, t_st_p, cd_p = permutation_two_group_cluster_test(
                            ea, eb, t_c, n_permutations=int(n_permutations),
                            p_threshold=0.05, jitter_sec=0.1, label_a=ca, label_b=cb)
                    except Exception as e:
                        st.warning(f"Permutation test failed: {e}"); continue
                ea_uv = _eeg_display_matrix(ea)
                eb_uv = _eeg_display_matrix(eb)
                ma, mda = median_and_mad(ea_uv, axis=0); mb, mdb = median_and_mad(eb_uv, axis=0)
                fig_p, (ax1, ax2) = plt.subplots(2, 1, figsize=(14, 7), sharex=True,
                                                   gridspec_kw={'height_ratios': [3, 1.5]})
                ax1.axhline(0, color='black', lw=0.4, alpha=0.3); ax1.axvline(0, color='red', ls='--', alpha=0.5)
                ax1.plot(t_c, ma, color=group_color.get(ca,'blue'), lw=2, label=f"{ca} (N={ea.shape[0]})")
                ax1.fill_between(t_c, ma-mda, ma+mda, color=group_color.get(ca,'blue'), alpha=0.15)
                ax1.plot(t_c, mb, color=group_color.get(cb,'red'), lw=2, label=f"{cb} (N={eb.shape[0]})")
                ax1.fill_between(t_c, mb-mdb, mb+mdb, color=group_color.get(cb,'red'), alpha=0.15)
                ax1.plot(t_c, ma-mb, color='black', lw=1.5, ls='--', label=f"{ca}-{cb}")
                for win in sig_wins:
                    ax1.axvspan(win['start'], win['end'], color='orange', alpha=0.25)
                    ax2.axvspan(win['start'], win['end'], color='orange', alpha=0.25)
                ax1.legend(fontsize=9); ax1.set_ylabel(f"EEG amplitude ({EEG_DISPLAY_UNIT})"); ax1.grid(True, alpha=0.2)
                if t_st_p is not None:
                    _df_p = ea.shape[0] + eb.shape[0] - 2
                    _t_thresh_p = stats.t.ppf(1 - 0.05 / 2, df=_df_p)
                    ax2.plot(t_c, t_st_p, color='#5c35cc', lw=2, label=f"T-statistic ({ca}−{cb})")
                    ax2.axhline(0, color='black', lw=0.4)
                    ax2.axhline(_t_thresh_p, color='gray', lw=0.8, ls='--', alpha=0.7,
                                label=f'±t thresh={_t_thresh_p:.2f} (p=0.05)')
                    ax2.axhline(-_t_thresh_p, color='gray', lw=0.8, ls='--', alpha=0.7)
                    ax2.fill_between(t_c, 0, t_st_p, where=(t_st_p > 0),
                                     color='#2196F3', alpha=0.2, interpolate=True)
                    ax2.fill_between(t_c, 0, t_st_p, where=(t_st_p < 0),
                                     color='#FF5722', alpha=0.2, interpolate=True)
                    ax2.set_ylabel("T-statistic"); ax2.legend(fontsize=8, loc='upper right')
                ax2.set_xlabel("Time relative to R-peak (s)"); ax2.grid(True, alpha=0.2)
                ns = len(sig_wins)
                mp = min((w['p_value'] for w in sig_wins), default=float('nan'))
                ps = f"p={'<0.001' if mp<0.001 else f'{mp:.3f}'}" if sig_wins else "no sig. windows"
                fig_p.suptitle(f"EEG HEP: {ca} vs {cb} | {ns} sig. window(s) | {ps}", fontsize=12, fontweight='bold')
                fig_p.tight_layout()
                _show_non_ica_figure(fig_p, use_container_width=True); plt.close(fig_p)
                if sig_wins:
                    st.markdown(f"**Significant windows:**")
                    st.dataframe(pd.DataFrame([
                        {"Start (s)": f"{w['start']:.3f}", "End (s)": f"{w['end']:.3f}",
                         "Duration (ms)": f"{(w['end']-w['start'])*1000:.0f}", "p-value": f"{w['p_value']:.4f}"}
                        for w in sig_wins]), use_container_width=True)
                else:
                    st.caption(f"No significant EEG differences between {ca} and {cb}.")

                # Per-channel pairwise breakdown
                ch_a = group_data[ca].get("eeg_per_ch", {})
                ch_b = group_data[cb].get("eeg_per_ch", {})
                shared_chs = [
                    ch for ch in sorted(set(ch_a) & set(ch_b))
                    if str(ch).strip().lower() != "sao2"
                ]
                if shared_chs:
                    st.markdown(f"**Per-electrode significance: {ca} vs {cb}**")
                    st.caption(
                        "Each electrode is tested separately with the same two-group cluster-permutation approach. "
                        "'Significant at any time' means at least one cluster-corrected window with p < 0.05."
                    )
                    show_only_sig_channels = st.toggle(
                        "Show only channels with significant windows",
                        value=False,
                        key=f"dxgc_only_sig_channel_plots_{i_p}_{j_p}",
                        help="Turn on to hide channel plots with no significant cluster-corrected window.",
                    )
                    with st.spinner(f"Testing each electrode: {ca} vs {cb}..."):
                        ch_pw_rows = []
                        ch_pw_plot_data = {}
                        for ch_k in shared_chs:
                            ea_k = ch_a[ch_k]; eb_k = ch_b[ch_k]
                            t_k = group_data[ca]["times"]
                            ml_k = min(ea_k.shape[1], eb_k.shape[1], len(t_k))
                            ea_k = ea_k[:, :ml_k]; eb_k = eb_k[:, :ml_k]; tk = t_k[:ml_k]
                            try:
                                sw_k, _, cd_k = permutation_two_group_cluster_test(
                                    ea_k, eb_k, tk, n_permutations=int(n_permutations),
                                    p_threshold=0.05, jitter_sec=0.1,
                                    label_a=ca, label_b=cb, channel_label=ch_k,
                                    show_explanation=False)
                            except Exception:
                                sw_k, cd_k = [], None
                            ch_pw_plot_data[ch_k] = (ea_k, eb_k, tk, sw_k)
                            ns_k = len(sw_k)
                            mp_k = min((w['p_value'] for w in sw_k), default=float('nan'))
                            if sw_k:
                                windows_k = "; ".join(
                                    f"{w['start']:.3f}-{w['end']:.3f}s (p={w['p_value']:.4f}, {w.get('direction', '')})"
                                    for w in sw_k
                                )
                            else:
                                windows_k = "None"
                            ch_pw_rows.append({
                                "Electrode": ch_k,
                                f"N({ca})": ea_k.shape[0],
                                f"N({cb})": eb_k.shape[0],
                                "Significant at any time": "Yes" if sw_k else "No",
                                "Significant windows": ns_k,
                                "Min cluster p-value": mp_k if sw_k else np.nan,
                                "Window(s)": windows_k,
                            })

                    visible_ch_pw_plot_data = {
                        ch: plot_values
                        for ch, plot_values in ch_pw_plot_data.items()
                        if not show_only_sig_channels or plot_values[3]
                    }
                    if not visible_ch_pw_plot_data:
                        st.info(f"No channels have significant windows for {ca} vs {cb}.")
                    else:
                        n_ch_plot = len(visible_ch_pw_plot_data)
                        ncols_ch = min(4, n_ch_plot)
                        nrows_ch = math.ceil(n_ch_plot / ncols_ch)
                        fig_ch_pw, axes_ch_pw = plt.subplots(
                            nrows_ch, ncols_ch,
                            figsize=(5 * ncols_ch, 3.6 * nrows_ch),
                            sharex=True,
                            squeeze=False,
                        )
                        axes_ch_pw_flat = axes_ch_pw.flatten()
                        for ch_idx, (ch_k, (ea_k, eb_k, tk, sw_k)) in enumerate(visible_ch_pw_plot_data.items()):
                            ax_ch = axes_ch_pw_flat[ch_idx]
                            ea_k_uv = _eeg_display_matrix(ea_k)
                            eb_k_uv = _eeg_display_matrix(eb_k)
                            med_a, mad_a = median_and_mad(ea_k_uv, axis=0)
                            med_b, mad_b = median_and_mad(eb_k_uv, axis=0)
                            color_a = group_color.get(ca, "#1f77b4")
                            color_b = group_color.get(cb, "#d62728")
                            ax_ch.plot(tk, med_a, color=color_a, lw=1.6, label=f"{ca} (N={ea_k.shape[0]})")
                            ax_ch.fill_between(tk, med_a - mad_a, med_a + mad_a, color=color_a, alpha=0.15)
                            ax_ch.plot(tk, med_b, color=color_b, lw=1.6, label=f"{cb} (N={eb_k.shape[0]})")
                            ax_ch.fill_between(tk, med_b - mad_b, med_b + mad_b, color=color_b, alpha=0.15)
                            for win_k in sw_k:
                                ax_ch.axvspan(win_k["start"], win_k["end"], color="orange", alpha=0.25)
                            ax_ch.axhline(0, color="black", lw=0.4, alpha=0.35)
                            ax_ch.axvline(0, color="red", lw=0.8, ls="--", alpha=0.5)
                            min_p_k = min((win["p_value"] for win in sw_k), default=np.nan)
                            sig_text = (
                                f"{len(sw_k)} sig. window(s), p={'<0.001' if min_p_k < 0.001 else f'{min_p_k:.3f}'}"
                                if sw_k else "no sig. windows"
                            )
                            ax_ch.set_title(f"{ch_k} | {sig_text}", fontsize=9, fontweight="bold")
                            ax_ch.set_xlabel("Time from R-peak (s)")
                            ax_ch.set_ylabel(f"Amplitude ({EEG_DISPLAY_UNIT})")
                            ax_ch.grid(True, alpha=0.2)
                            ax_ch.tick_params(labelsize=7)
                        for unused_idx in range(n_ch_plot, len(axes_ch_pw_flat)):
                            axes_ch_pw_flat[unused_idx].set_visible(False)
                        handles_ch, labels_ch = axes_ch_pw_flat[0].get_legend_handles_labels()
                        if handles_ch:
                            fig_ch_pw.legend(
                                handles_ch, labels_ch, loc="lower center",
                                bbox_to_anchor=(0.5, -0.01), ncol=2, fontsize=8, frameon=False,
                            )
                        plot_scope = "Significant channels only" if show_only_sig_channels else "All channels"
                        fig_ch_pw.suptitle(
                            f"Per-Channel EEG HEP: {ca} vs {cb} | {plot_scope} | Orange = significant windows",
                            fontsize=12, fontweight="bold",
                        )
                        fig_ch_pw.tight_layout(rect=[0, 0.04, 1, 0.96])
                        _show_non_ica_figure(fig_ch_pw, use_container_width=True)
                        plt.close(fig_ch_pw)

                    ch_pw_df = pd.DataFrame(ch_pw_rows)
                    if not ch_pw_df.empty:
                        ch_pw_df = ch_pw_df.sort_values(
                            ["Significant windows", "Min cluster p-value", "Electrode"],
                            ascending=[False, True, True],
                            na_position="last",
                        )
                        st.dataframe(
                            ch_pw_df,
                            use_container_width=True,
                            column_config={
                                "Min cluster p-value": st.column_config.NumberColumn(
                                    "Min cluster p-value", format="%.4f"
                                )
                            },
                        )
                        st.download_button(
                            label="Download per-electrode significance CSV",
                            data=ch_pw_df.to_csv(index=False).encode("utf-8"),
                            file_name=(
                                f"per_electrode_pairwise_eeg_{ca}_vs_{cb}.csv"
                                .replace(" ", "_")
                                .replace("/", "-")
                            ),
                            mime="text/csv",
                            key=f"download_pairwise_eeg_electrodes_{i_p}_{j_p}",
                        )
                    if not any(row["Significant at any time"] == "Yes" for row in ch_pw_rows):
                        st.caption(f"No individual electrode had a significant cluster-corrected window for {ca} vs {cb}.")
                else:
                    st.info(f"No shared per-electrode EEG data available for {ca} vs {cb}.")

                pw_summary.append({"Comparison": f"{ca} vs {cb}", "N sig. windows": ns,
                    "Min p-value": f"{mp:.4f}" if sig_wins else "n/a"})
        if pw_summary:
            st.markdown("---")
            st.subheader("\U0001f4cb G. Pairwise Comparison Summary")
            st.dataframe(pd.DataFrame(pw_summary), use_container_width=True)

    # ── PLOT H: Autoencoder latent space ──────────────────────────────────────
    st.markdown("---")
    st.subheader("\U0001f9e0 H. Autoencoder Latent Space — Group Separation")
    st.markdown(
        "**Interpretation:** Each dot = one patient. 2-D latent representation of the HEP waveform "
        "computed by a neural (or linear fallback) autoencoder. Tighter within-group clusters and "
        "wider between-group separation = stronger diagnosis-linked brain–heart coupling pattern. "
        "H.1 = mean EEG (all channels averaged). H.2 = per EEG channel. H.3 = all channels stacked."
    )

    ae_groups = [g for g in group_data if group_data[g]["eeg"] is not None]
    if len(ae_groups) < 2:
        st.info("Need >= 2 groups with EEG data for autoencoder comparison.")
    else:
        # ── H.1 All-channels-mean ─────────────────────────────────────────────
        X_all_list, lbl_all = [], []
        for g in ae_groups:
            m = group_data[g]["eeg"]
            X_all_list.append(m)
            lbl_all.extend([g] * m.shape[0])
        try:
            X_all = np.vstack(X_all_list)
            lbl_all = np.array(lbl_all)
            with st.spinner("Running autoencoder (mean EEG, all groups)..."):
                z_all, err_all, ae_method_all = _run_waveform_autoencoder_embedding(X_all)
            if z_all is not None:
                fig_ae_all, ax_ae_all = plt.subplots(figsize=(8, 6))
                for g in ae_groups:
                    mask = lbl_all == g
                    n_g = int(mask.sum())
                    ax_ae_all.scatter(z_all[mask, 0], z_all[mask, 1],
                                      color=group_color.get(g, 'gray'), label=f"{g} (N={n_g})",
                                      alpha=0.7, s=60, edgecolors='k', linewidths=0.3)
                ax_ae_all.set_xlabel("Latent dim 1"); ax_ae_all.set_ylabel("Latent dim 2")
                ax_ae_all.set_title(f"H.1 Mean EEG HEP — {ae_method_all}", fontsize=11, fontweight='bold')
                ax_ae_all.legend(fontsize=9, bbox_to_anchor=(1.01, 1), loc='upper left')
                ax_ae_all.grid(True, alpha=0.25)
                fig_ae_all.tight_layout()
                _show_non_ica_figure(fig_ae_all, use_container_width=True)
                plt.close(fig_ae_all)
        except Exception as _ae_e:
            st.warning(f"Mean-EEG autoencoder failed: {_ae_e}")

        # ── H.2 Per-channel ───────────────────────────────────────────────────
        all_chs_ae = sorted(set(ch for g in ae_groups for ch in group_data[g].get("eeg_per_ch", {})))
        if all_chs_ae:
            st.markdown("**H.2 Per-channel latent spaces:**")
            max_ch_plots = st.slider("Max channels to plot", 1, min(len(all_chs_ae), 20),
                                     min(len(all_chs_ae), 6), key="dxgc_ae_maxch")
            chs_to_plot = all_chs_ae[:max_ch_plots]
            n_cols_ae = min(3, len(chs_to_plot))
            n_rows_ae = int(np.ceil(len(chs_to_plot) / n_cols_ae))
            fig_ae_ch, axes_ae_ch = plt.subplots(n_rows_ae, n_cols_ae,
                                                   figsize=(6 * n_cols_ae, 5 * n_rows_ae))
            axes_flat = np.array(axes_ae_ch).ravel() if hasattr(np.array(axes_ae_ch), 'ravel') else [axes_ae_ch]
            with st.spinner("Running per-channel autoencoders..."):
                for ax_i, ch_k in enumerate(chs_to_plot):
                    ax_k = axes_flat[ax_i]
                    X_ch_list, lbl_ch = [], []
                    for g in ae_groups:
                        m_ch = group_data[g].get("eeg_per_ch", {}).get(ch_k)
                        if m_ch is not None and m_ch.shape[0] >= 2:
                            X_ch_list.append(m_ch)
                            lbl_ch.extend([g] * m_ch.shape[0])
                    if len(X_ch_list) < 2:
                        ax_k.set_visible(False); continue
                    try:
                        X_ch = np.vstack(X_ch_list)
                        lbl_ch_arr = np.array(lbl_ch)
                        z_ch, _, ae_m_ch = _run_waveform_autoencoder_embedding(X_ch)
                        if z_ch is None:
                            ax_k.set_visible(False); continue
                        for g in ae_groups:
                            mask_k = lbl_ch_arr == g
                            n_gk = int(mask_k.sum())
                            if n_gk == 0:
                                continue
                            ax_k.scatter(z_ch[mask_k, 0], z_ch[mask_k, 1],
                                         color=group_color.get(g, 'gray'),
                                         label=f"{g} (N={n_gk})",
                                         alpha=0.7, s=40, edgecolors='k', linewidths=0.2)
                        ax_k.set_title(f"{ch_k} — {ae_m_ch}", fontsize=8, fontweight='bold')
                        ax_k.set_xlabel("L1", fontsize=7); ax_k.set_ylabel("L2", fontsize=7)
                        ax_k.legend(fontsize=6, loc='best')
                        ax_k.grid(True, alpha=0.2)
                        ax_k.tick_params(labelsize=7)
                    except Exception:
                        ax_k.set_visible(False)
            for ax_rem in axes_flat[len(chs_to_plot):]:
                ax_rem.set_visible(False)
            fig_ae_ch.tight_layout()
            _show_non_ica_figure(fig_ae_ch, use_container_width=True)
            plt.close(fig_ae_ch)
        else:
            st.info("Per-channel data not available. Click '\U0001f504 Reload & Recompute' to rebuild cache with per-channel data.")

        # ── H.3 All-channels multivariate ─────────────────────────────────────
        mv_list, lbl_mv = [], []
        for g in ae_groups:
            mv_m = group_data[g].get("eeg_mv")
            if mv_m is not None and mv_m.shape[0] >= 2:
                mv_list.append(mv_m)
                lbl_mv.extend([g] * mv_m.shape[0])
        if len(mv_list) >= 2:
            try:
                # Trim columns to common width
                min_mv_cols = min(m.shape[1] for m in mv_list)
                X_mv = np.vstack([m[:, :min_mv_cols] for m in mv_list])
                lbl_mv_arr = np.array(lbl_mv)
                with st.spinner("Running autoencoder (all channels stacked)..."):
                    z_mv, _, ae_m_mv = _run_waveform_autoencoder_embedding(X_mv)
                if z_mv is not None:
                    st.markdown("**H.3 All channels stacked (multivariate):**")
                    fig_mv, ax_mv = plt.subplots(figsize=(8, 6))
                    for g in ae_groups:
                        mask_mv = lbl_mv_arr == g
                        n_gm = int(mask_mv.sum())
                        ax_mv.scatter(z_mv[mask_mv, 0], z_mv[mask_mv, 1],
                                      color=group_color.get(g, 'gray'),
                                      label=f"{g} (N={n_gm})", alpha=0.7, s=60,
                                      edgecolors='k', linewidths=0.3)
                    ax_mv.set_xlabel("Latent dim 1"); ax_mv.set_ylabel("Latent dim 2")
                    ax_mv.set_title(f"H.3 All EEG Channels (stacked) — {ae_m_mv}", fontsize=11, fontweight='bold')
                    ax_mv.legend(fontsize=9, bbox_to_anchor=(1.01, 1), loc='upper left')
                    ax_mv.grid(True, alpha=0.25)
                    fig_mv.tight_layout()
                    _show_non_ica_figure(fig_mv, use_container_width=True)
                    plt.close(fig_mv)
            except Exception as _mv_e:
                st.warning(f"Multivariate autoencoder failed: {_mv_e}")
        else:
            st.info("Multivariate (all channels stacked) data not available. Rebuild cache to enable.")

    # ── PLOT I: Venn / Overlap diagram ────────────────────────────────────────
    if grouping_mode in ("By Diagnosis Category", "By Specific Diagnosis") and len(group_data) >= 2:
        st.markdown("---")
        st.subheader("\U0001f9e9 I. Patient Overlap Between Groups")
        st.markdown(
            "**Interpretation:** Patients can have multiple diagnoses, so groups may overlap. "
            "Each cell shows the number of patients in both groups simultaneously."
        )
        try:
            pid_sets = {}
            for g in group_data:
                inds_for_g = valid_groups.get(g) if valid_groups else []
                if inds_for_g:
                    pid_sets[g] = set(str(ind[0]) for ind in inds_for_g)
                else:
                    pid_sets[g] = set()  # cache path: no patient-level data available
            groups_ov = sorted(g for g in pid_sets if pid_sets[g])
            n_ov = len(groups_ov)
            if n_ov < 2:
                st.info("Overlap data not available from cache. Rebuild cache to enable Venn diagram.")
            else:
                any_overlap = any(
                    len(pid_sets[groups_ov[i]] & pid_sets[groups_ov[j]]) > 0
                    for i in range(n_ov) for j in range(i+1, n_ov)
                )
                if not any_overlap:
                    st.info("Groups are non-overlapping (mutually exclusive).")
                else:
                    ov_rows = []
                    for i, ga in enumerate(groups_ov):
                        for j, gb in enumerate(groups_ov):
                            if i < j:
                                ov = len(pid_sets[ga] & pid_sets[gb])
                                ov_rows.append({
                                    "Group A": ga, "N(A)": len(pid_sets[ga]),
                                    "Group B": gb, "N(B)": len(pid_sets[gb]),
                                    "Overlap N": ov,
                                    "Overlap %A": f"{100*ov/max(len(pid_sets[ga]),1):.1f}%",
                                    "Overlap %B": f"{100*ov/max(len(pid_sets[gb]),1):.1f}%",
                                })
                    st.dataframe(pd.DataFrame(ov_rows), use_container_width=True)

                    ov_mat = np.zeros((n_ov, n_ov), dtype=int)
                    for i in range(n_ov):
                        for j in range(n_ov):
                            ov_mat[i, j] = len(pid_sets[groups_ov[i]] & pid_sets[groups_ov[j]])
                    fig_ov, ax_ov = plt.subplots(figsize=(max(6, n_ov*1.4), max(5, n_ov*1.2)))
                    im = ax_ov.imshow(ov_mat, cmap='Blues', aspect='auto')
                    ax_ov.set_xticks(range(n_ov)); ax_ov.set_yticks(range(n_ov))
                    ax_ov.set_xticklabels(
                        [f"{g}\n(N={len(pid_sets[g])})" for g in groups_ov],
                        rotation=35, ha='right', fontsize=8)
                    ax_ov.set_yticklabels([f"{g} (N={len(pid_sets[g])})" for g in groups_ov], fontsize=8)
                    for i in range(n_ov):
                        for j in range(n_ov):
                            ax_ov.text(j, i, str(ov_mat[i, j]), ha='center', va='center',
                                       fontsize=9,
                                       color='white' if ov_mat[i, j] > ov_mat.max()*0.6 else 'black')
                    plt.colorbar(im, ax=ax_ov, label='N patients in common')
                    ax_ov.set_title(f"Patient Overlap Matrix — {grouping_mode}", fontsize=11, fontweight='bold')
                    fig_ov.tight_layout()
                    _show_non_ica_figure(fig_ov, use_container_width=True)
                    plt.close(fig_ov)

                    if n_ov in (2, 3):
                        try:
                            import matplotlib_venn as _mvenn
                            fig_venn, ax_venn = plt.subplots(figsize=(7, 5))
                            sets_v = [pid_sets[g] for g in groups_ov]
                            labels_v = [f"{g}\nN={len(pid_sets[g])}" for g in groups_ov]
                            if n_ov == 2:
                                _mvenn.venn2(sets_v, set_labels=labels_v, ax=ax_venn)
                            else:
                                _mvenn.venn3(sets_v, set_labels=labels_v, ax=ax_venn)
                            ax_venn.set_title(f"Venn Diagram — {grouping_mode}", fontsize=11, fontweight='bold')
                            fig_venn.tight_layout()
                            _show_non_ica_figure(fig_venn, use_container_width=True)
                            plt.close(fig_venn)
                        except ImportError:
                            st.caption("Install `matplotlib-venn` for a Venn diagram. Overlap table shown above.")
                        except Exception as _venn_e:
                            st.caption(f"Venn diagram skipped: {_venn_e}")
        except Exception as _ov_e:
            st.warning(f"Overlap analysis failed: {_ov_e}")


def main():
    st.title("HEP Group Comparison Dashboard")
    st.write("Comparing Amplitude vs Time (Heartbeat Evoked Potential).")

    st.sidebar.header("Export Tools")
    if 'pptx_figures_data' in st.session_state and len(st.session_state.pptx_figures_data) > 0:
        st.sidebar.write(f"Figures saved: {len(st.session_state.pptx_figures_data)}")
        if st.sidebar.button("Prepare PowerPoint Report"):
            with st.spinner("Generating PowerPoint..."):
                pptx_io = generate_pptx()
                if pptx_io:
                    st.sidebar.download_button(
                        label="Download PPTX",
                        data=pptx_io,
                        file_name="HEP_Report.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                else:
                    st.sidebar.error("python-pptx is not installed.")
        if st.sidebar.button("Clear Saved Figures"):
            st.session_state.pptx_figures_data = []
            st.rerun()
    else:
        st.sidebar.write("No figures saved yet.")


    base_path = "/storage/pblab_shared_data2/Nir/Cobrad/pickles_sleep_stage"

    # Select Sleep Stage
    sleep_stages = ['light_sleep', 'N3', 'R', 'W']
    selected_stage = st.selectbox("Select Sleep Stage", sleep_stages, index=0)
    
    # Analysis Mode Selection
    _all_modes = [
        "Single Group Analysis",
        "Compare Groups",
        "Compare Sleep Stages",
        "Compare Groups All Sleep Stages",
        "Compare Groups Non-EEG Channels",
        "Per Channel Group Comparison",
        "Diagnosis Group Comparison",
    ]
    _cli_mode = _parse_mode_arg()
    _default_index = _all_modes.index(_cli_mode) if _cli_mode in _all_modes else 1
    mode = st.radio("Analysis Mode", _all_modes, index=_default_index)

    if mode == "Compare Groups":
        run_compare_groups_analysis(base_path, selected_stage)
    elif mode == "Compare Sleep Stages":
        run_compare_sleep_stages_analysis(base_path)
    elif mode == "Compare Groups All Sleep Stages":
        run_compare_groups_all_stages_analysis(base_path)
    elif mode == "Compare Groups Non-EEG Channels":
        run_compare_groups_non_eeg_analysis(base_path, selected_stage)
    elif mode == "Per Channel Group Comparison":
        run_per_channel_group_comparison_analysis(base_path, selected_stage)
    elif mode == "Diagnosis Group Comparison":
        run_diagnosis_group_comparison_analysis(base_path, selected_stage)
    else: # Single Group Analysis
        run_single_group_analysis(base_path, selected_stage)

if __name__ == "__main__":
    main()

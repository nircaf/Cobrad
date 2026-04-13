import os
import glob
import numpy as np
import pandas as pd
from sklearn.preprocessing import StandardScaler
from scipy.stats import iqr
import mne
import matplotlib.pyplot as plt
import seaborn as sns
from scipy import stats
from statsmodels.stats.multitest import multipletests
import statsmodels.stats.multitest as smm
from mne_connectivity import SpectralConnectivity as spectral_connectivity
import statsmodels.api as sm
import streamlit as st
import json
import pickle
import uuid
import re
import io
from tqdm import tqdm
from concurrent.futures import ProcessPoolExecutor
from scipy.signal import welch
import dabest
from statsmodels.stats.power import TTestIndPower
from scipy.stats import spearmanr
from statsmodels.stats.multitest import fdrcorrection
from scipy.stats import mannwhitneyu
import sys
import mne.io.array
sys.modules['mne.io.array.array'] = mne.io.array
# Optional imports for plotly functionality
try:
    import plotly.graph_objects as go  # noqa: F401
    import plotly.express as px  # noqa: F401
    from plotly.subplots import make_subplots  # noqa: F401
    PLOTLY_AVAILABLE = True
except ImportError:
    PLOTLY_AVAILABLE = False

_config = json.load(open(os.path.join(os.path.dirname(__file__), 'config.json')))
eeg_channels = _config['eeg_channels']
eeg_dict_convertion = _config['eeg_dict_convertion']


# Base color palette for consistent plotting across functions
BASE_PALETTE = ["#1F77B4", "#E41A1C", "#2CA02C", "#FF7F0E", "#9467BD",
               "#8C564B", "#17BECF", "#D62728", "#BCBD22", "#7F7F7F"]


def st_progress_updater(total: int, label: str = "Progress"):
    """Return a progress updater for Streamlit.

    Usage:
        updater = st_progress_updater(100, "Processing files")
        for i, item in enumerate(items, start=1):
            # do work
            updater(i, f"Item {i} / {len(items)}")

    The returned updater accepts two args: current (int) and optional message (str).
    If Streamlit (`st`) is not available (e.g., running outside Streamlit), this
    becomes a no-op but still callable.
    """
    try:
        import streamlit as st
    except Exception:
        # st not available — return a no-op updater
        def _noop(current: int, message: str | None = None):
            return

        return _noop

    progress = st.progress(0)
    placeholder = st.empty()

    def _updater(current: int, message: str | None = None):
        # clamp current to [0, total]
        pct = int(max(0, min(total and int(current * 100 / total) or 0, 100)))
        try:
            progress.progress(pct)
            if message:
                placeholder.text(f"{label}: {message} ({pct}%)")
        except Exception:
            # Streamlit may raise if session ended — swallow safely
            pass

    return _updater


class st_progress_context:
    """Context manager convenience wrapper around `st_progress_updater`.

    Example:
        with st_progress_context(total=len(items), label="Processing") as update:
            for i, item in enumerate(items, start=1):
                # do work
                update(i, f"Item {i}")
    """
    # Quick usage example (also works outside of Streamlit gracefully):
    #
    # >>> from utils.eeg_utils import st_progress_context
    # >>> items = range(10)
    # >>> with st_progress_context(total=len(items), label="Processing") as update:
    # ...     for i, it in enumerate(items, start=1):
    # ...         # do work
    # ...         update(i, f"Processed {i}")
    def __init__(self, total: int, label: str = "Progress"):
        self.total = total
        self.label = label
        self._updater = None

    def __enter__(self):
        self._updater = st_progress_updater(self.total, self.label)
        return self._updater

    def __exit__(self, exc_type, exc, tb):
        # on exit, set progress to 100%
        if self._updater is not None:
            try:
                self._updater(self.total, "Done")
            except Exception:
                pass
power_bands = {
    "delta": [.5, 4],
    "theta": [4, 8],
    "alpha": [8, 12],
    "beta": [12, 30],
    "gamma": [30, 100]
}

import mne  # noqa: E402

def slice_stratified_windows(
    raw: mne.io.BaseRaw,
    n_bins: int = 6,
    window_s: float = 600.0,
    windows_per_bin: int = 3,
    min_gap_s: float | None = None,
    grid_s: float = 1.0,
    seed: int = 0,
    exclude_desc: tuple[str, ...] = ("BAD", "ARTIFACT", "EDGE", "REJECT")
):
    """
    Split the recording duration into N equal bins, and from EACH bin select
    random short windows totaling approximately M minutes (minutes_per_bin).
    Selection is random within each bin, respects excluded annotations, and
    enforces a minimum gap between chosen windows. If not enough candidates
    exist, a relaxed fallback sampling is applied to reach the target count.

    Parameters
    ----------
    raw : mne.io.Raw
        Continuous EEG data (preloaded).
    n_bins : int
        Number of equal bins for the recording duration.
    window_s : float
        Window length in seconds (e.g., 600 for 10 minutes).
    windows_per_bin : int
        Number of non-overlapping windows to sample from each bin.
    min_gap_s : float | None
        Minimum time between window starts. Defaults to window_s.
    grid_s : float
        Candidate start grid resolution (seconds).
    seed : int
        RNG seed for reproducibility.
    exclude_desc : tuple[str, ...]
        Annotation descriptions to exclude (case-insensitive substring match).

    Returns
    -------
    X : np.ndarray, shape (n_windows, n_channels, n_times)
        Stacked windows from all bins.
    meta : dict
        Metadata including per-window bin index and counts per bin:
          - starts_s, stops_s
          - bin_index (per window)
          - per_bin_counts (length n_bins)
          - bin_ranges_s: list[(start, stop)] for each bin
          - sfreq, ch_names, window_s, etc.
    """
    sfreq = float(raw.info["sfreq"])
    n_ch = len(raw.ch_names)
    # Duration in seconds (raw.times starts at 0.0)
    dur_s = float(raw.times[-1])
    if dur_s <= 0:
        raise RuntimeError("Recording duration is non-positive.")
    bin_len = dur_s / float(n_bins)
    rng = np.random.default_rng(seed)
    gap = window_s if min_gap_s is None else float(min_gap_s)

    # ------------------ helpers ------------------ #
    def clean_intervals():
        # Start with the full interval, subtract excluded annotations.
        clean = [(0.0, dur_s)]
        ex = []
        # mne.Annotations elements have attributes .description, .onset, .duration
        for ann in getattr(raw, "annotations", []):
            desc = str(getattr(ann, "description", ""))
            if any(tag.lower() in desc.lower() for tag in exclude_desc):
                onset = float(getattr(ann, "onset", 0.0))
                duration = float(getattr(ann, "duration", 0.0))
                ex.append((onset, onset + duration))

        def subtract(intervals, rm):
            out = []
            s, e = rm
            for a, b in intervals:
                if e <= a or s >= b:
                    out.append((a, b))
                else:
                    if a < s:
                        out.append((a, s))
                    if e < b:
                        out.append((e, b))
            return out

        for rm in ex:
            clean = subtract(clean, rm)
        # Keep only intervals that can hold at least one window
        return [(a, b) for a, b in clean if (b - a) >= window_s]

    def candidate_starts(intervals, start, stop):
        # Generate candidate starts on a grid, clipped to [start, stop - window_s]
        cands = []
        for a, b in intervals:
            a_clip, b_clip = max(a, start), min(b, stop)
            end = b_clip - window_s
            if end <= a_clip:
                continue
            # inclusive start, exclusive end
            starts = np.arange(a_clip, end + 1e-9, grid_s, dtype=float)
            if starts.size:
                cands.append(starts)
        if not cands:
            return np.array([], dtype=float)
        return np.unique(np.concatenate(cands))

    def choose_with_gap(cands, n_target, gap_s):
        # Shuffle candidates and pick greedily with min gap constraint
        c = np.array(cands, dtype=float)
        rng.shuffle(c)
        chosen = []
        for t in c:
            if all(abs(t - u) >= gap_s for u in chosen):
                chosen.append(float(t))
                if len(chosen) >= n_target:
                    break
        return np.array(sorted(chosen), dtype=float)

    # --------------------------------------------- #

    clean_ints = clean_intervals()
    if not clean_ints:
        raise RuntimeError("No clean intervals of at least window_s were found after excluding annotations.")

    # Target windows per bin
    n_per_bin = int(windows_per_bin)
    n_times = int(round(window_s * sfreq))

    # Collect per-bin selections
    starts_all = []
    stops_all = []
    bin_index_all = []
    per_bin_counts = []
    bin_ranges = []

    for i in range(n_bins):
        b0, b1 = i * bin_len, (i + 1) * bin_len
        bin_ranges.append((b0, b1))
        cands = candidate_starts(clean_ints, b0, b1)

        selected = np.array([], dtype=float)
        if cands.size > 0 and n_per_bin > 0:
            # First pass: enforce gap
            selected = choose_with_gap(cands, n_per_bin, gap)
            # If not enough, relax gap (half gap)
            if selected.size < n_per_bin and cands.size > 0:
                need = n_per_bin - selected.size
                # Prefer starts far from already chosen to reduce heavy overlap
                remaining = np.setdiff1d(cands, selected)
                if remaining.size == 0:
                    remaining = cands
                # Fallback: random choice from remaining (may violate gap if necessary)
                extra = rng.choice(remaining, size=need, replace=(remaining.size < need))
                selected = np.sort(np.concatenate([selected, extra]))

        per_bin_counts.append(int(selected.size))
        if selected.size:
            starts_all.extend(selected.tolist())
            stops_all.extend((selected + window_s).tolist())
            bin_index_all.extend([i] * selected.size)

    # Convert to arrays
    starts = np.array(starts_all, dtype=float)
    stops = np.array(stops_all, dtype=float)
    bin_index = np.array(bin_index_all, dtype=int)

    # Extract data safely (skip any window that would exceed data length)
    X_list = []
    valid_starts = []
    valid_stops = []
    valid_bins = []
    n_samples_total = raw.n_times
    for t0, t1, bidx in zip(starts, stops, bin_index):
        s0 = int(round(t0 * sfreq))
        s1 = s0 + n_times
        if s0 < 0 or s1 > n_samples_total:
            continue
        X_list.append(raw.get_data(start=s0, stop=s1))
        valid_starts.append(t0)
        valid_stops.append(t1)
        valid_bins.append(bidx)

    if len(X_list) == 0:
        raise RuntimeError("No windows could be extracted. Consider reducing minutes_per_bin or constraints.")

    X = np.stack(X_list, axis=0)  # (n_windows, n_channels, n_times)
    meta = dict(
        starts_s=np.array(valid_starts, dtype=float),
        stops_s=np.array(valid_stops, dtype=float),
        bin_index=np.array(valid_bins, dtype=int),
        per_bin_counts=np.array(per_bin_counts, dtype=int),
        bin_ranges_s=bin_ranges,
        sfreq=sfreq,
        ch_names=list(raw.ch_names),
        window_s=window_s,
        n_windows=X.shape[0],
        n_channels=n_ch,
        n_times=n_times,
        n_bins=n_bins,
        windows_per_bin=windows_per_bin,
        grid_s=grid_s,
        gap_s=gap,
    )
    return X, meta

# ---------------- Example usage ----------------
# X, meta = slice_stratified_windows(raw, n_bins=6, window_s=4.0, minutes_per_bin=10, seed=42)
# -> About 6 * (10 min / 4 s) ≈ 900 windows total

def mean_of_resized_arrays(arrays):
    # Get the shapes of all arrays
    shapes = np.array([arr.shape for arr in arrays])
    
    # Compute median dimensions
    median_shape = tuple(np.median(shapes, axis=0).astype(int))
    
    # Resize all arrays to the median shape
    resized_arrays = np.array([np.resize(arr, median_shape) for arr in arrays])
    
    # Compute the mean
    return np.mean(resized_arrays, axis=0)

def rename_channels(raw):
    try:
        channel_names = raw.ch_names
        # remove channels that don't have EEG, ECG, or EOG in their name from raw
        # raw.drop_channels([channel for channel in raw.ch_names if 'EEG' not in channel])
        # remove EEG from channel name
        raw.rename_channels({channel: channel.replace('EEG', '').strip() for channel in raw.ch_names})
        # all channels that have EEG, their split(' ')[-1] needs to be uppercase first letter, all rest lower case
        raw.rename_channels({channel: ' '.join([part.capitalize() if i == len(channel.split(' ')) - 1 else part for i, part in enumerate(channel.split(' '))]) if 'EEG' in channel else channel for channel in raw.ch_names})
        # rename channels based on eeg_utils.eeg_dict_convertion
        valid_channels = set(raw.info['ch_names'])
        valid_rename_dict = {k: v for k, v in eeg_dict_convertion.items() if k in valid_channels}
        # if valid_rename_dict == {}, use the first .split('-')[0] as the channel name if it already exists, use .split('-')[1]
        if len(valid_rename_dict) == 0:
            channels_to_remove = []
            for channel in raw.ch_names:
                if '-' in channel:
                    split_0 = channel.split('-')[0]
                    split_1 = channel.split('-')[1]
                    existing_values = valid_rename_dict.values()
                    
                    # If both exist, remove the channel
                    if split_0 in existing_values and split_1 in existing_values:
                        channels_to_remove.append(channel)
                    # If split[0] exists, use split[1]
                    elif split_0 in existing_values:
                        valid_rename_dict[channel] = split_1
                    # Default: use split[0]
                    else:
                        valid_rename_dict[channel] = split_0
            # Remove channels where both splits already exist
            if channels_to_remove:
                raw.drop_channels(channels_to_remove)
        # Drop duplicate targets: if multiple source channels map to the same name, keep only the first
        seen_targets = {}
        dupes_to_drop = []
        for src, tgt in valid_rename_dict.items():
            if tgt in seen_targets:
                dupes_to_drop.append(src)
            else:
                seen_targets[tgt] = src
        # Also drop sources whose target already exists as an unrenamed channel
        existing_channels = set(raw.ch_names)
        for src, tgt in list(valid_rename_dict.items()):
            if tgt in existing_channels and src != tgt and src not in dupes_to_drop:
                dupes_to_drop.append(src)
        if dupes_to_drop:
            valid_rename_dict = {k: v for k, v in valid_rename_dict.items() if k not in dupes_to_drop}
            raw.drop_channels([ch for ch in dupes_to_drop if ch in raw.ch_names])
        raw.rename_channels(valid_rename_dict)
    except Exception:
        print("Error in rename_channels, proceeding without renaming. Error details:")
        pass
    for channel in raw.ch_names:
        if channel.lower() in [ch.lower() for ch in eeg_channels]:
            raw.set_channel_types({channel: 'eeg'})
            channel_name = channel
            if channel.lower().startswith('fp'):
                # if channel starts with fp (ignore case) capitalize it
                channel_name = channel_name.capitalize()
                raw.rename_channels({channel: channel_name})
            # if channel ends with z (ignore case) lowercase it
            if channel_name.lower().endswith('z'):
                # lower the last letter
                raw.rename_channels({channel_name: channel_name[:-1] + channel_name[-1].lower()})
        elif channel.lower() in ('eog', 'ecg', 'emg'):
            if channel.lower() == 'eog':
                ch_type = 'eog'
            elif channel.lower() == 'ecg':
                ch_type = 'ecg'
            else:
                ch_type = 'emg'
            print(f'Setting channel type to {ch_type} for {channel}')
            raw.set_channel_types({channel: ch_type})
        else:
            raw.set_channel_types({channel: 'misc'})
    return raw

def stat_text_get(group_data, col=None):
    lower_bound = -1e-30
    upper_bound = 1e30
    if col is None:
        q1 = np.clip(group_data.quantile(0.25), lower_bound, upper_bound)
        q3 = np.clip(group_data.quantile(0.75), lower_bound, upper_bound)
        iqr = np.clip(q3 - q1, lower_bound, upper_bound)
        stats_dict = {
            "N": len(group_data),
            "Mean": np.clip(group_data.mean(), lower_bound, upper_bound),
            "Median": np.clip(group_data.median(), lower_bound, upper_bound),
            "Max": np.clip(group_data.max(), lower_bound, upper_bound),
            "Min": np.clip(group_data.min(), lower_bound, upper_bound),
            "Std": np.clip(group_data.std(), lower_bound, upper_bound),
            "IQR": iqr,
        }
    else:
        q1 = np.clip(group_data[col].quantile(0.25), lower_bound, upper_bound)
        q3 = np.clip(group_data[col].quantile(0.75), lower_bound, upper_bound)
        iqr = np.clip(q3 - q1, lower_bound, upper_bound)
        stats_dict = {
            "N": len(group_data),
            "Mean": np.clip(group_data[col].mean(), lower_bound, upper_bound),
            "Median": np.clip(group_data[col].median(), lower_bound, upper_bound),
            "Max": np.clip(group_data[col].max(), lower_bound, upper_bound),
            "Min": np.clip(group_data[col].min(), lower_bound, upper_bound),
            "Std": np.clip(group_data[col].std(), lower_bound, upper_bound),
            "IQR": iqr,
        }

    stats_text = (
        f"N = {stats_dict['N']}, "
        f"Mean = {stats_dict['Mean']:.2e}, "
        f"Median = {stats_dict['Median']:.2e}, "
        f"Max = {stats_dict['Max']:.2e}, "
        f"Min = {stats_dict['Min']:.2e}, "
        f"Std = {stats_dict['Std']:.2e}, "
        f"IQR = {stats_dict['IQR']:.2e}"
    )
    return stats_text, stats_dict


def plot_tsne_by_group(
    df,
    features=None,
    group_col="Group",
    perplexity=30,
    random_state=42,
    is_streamlit=True,
    dbscan_eps=3,
    dbscan_min_samples=5
):
    """
    Plot a t-SNE of the DataFrame, coloring by group and drawing a convex hull polygon around clusters of points
    within each group if the cluster has more than dbscan_min_samples points.

    Parameters
    ----------
    df : pd.DataFrame
        DataFrame containing features and group column.
    features : list of str, optional
        List of columns to use as features. If None, use all numeric columns except group_col.
    group_col : str
        Name of the column to use for coloring.
    perplexity : int
        t-SNE perplexity parameter.
    random_state : int
        t-SNE random state.
    is_streamlit : bool
        Whether to use Streamlit plotting.
    dbscan_eps : float
        The maximum distance between two samples for one to be considered as in the neighborhood of the other (DBSCAN).
    dbscan_min_samples : int
        The number of samples in a neighborhood for a point to be considered as a core point (DBSCAN).
    """
    from sklearn.manifold import TSNE
    from scipy.spatial import ConvexHull
    from matplotlib.patches import Polygon
    from sklearn.cluster import DBSCAN

    if features is None:
        features = [col for col in df.select_dtypes(include=[np.number]).columns if col != group_col]
    X = df[features].values
    # Impute missing values with the mean of each column
    from sklearn.impute import SimpleImputer
    imputer = SimpleImputer(strategy="mean")
    X = imputer.fit_transform(X)
    groups = df[group_col].values

    tsne = TSNE(n_components=2, perplexity=perplexity, random_state=random_state)
    X_embedded = tsne.fit_transform(X)

    plt.figure(figsize=(16, 9))
    unique_groups = np.unique(groups)
    colors = plt.cm.get_cmap('Set1', len(unique_groups))
    ax = plt.gca()

    for i, group in enumerate(unique_groups):
        idx = groups == group
        points = X_embedded[idx]
        # Cluster points to find dense regions
        if points.shape[0] >= dbscan_min_samples:
            db = DBSCAN(eps=dbscan_eps, min_samples=dbscan_min_samples).fit(points)
            labels = db.labels_
            unique_labels = set(labels)
            for cluster_label in unique_labels:
                if cluster_label == -1:
                    continue  # Skip noise
                cluster_points = points[labels == cluster_label]
                if cluster_points.shape[0] >= dbscan_min_samples:
                    try:
                        hull = ConvexHull(cluster_points)
                        polygon = Polygon(cluster_points[hull.vertices], closed=True, facecolor=colors(i), alpha=0.2, edgecolor=colors(i), linewidth=2)
                        ax.add_patch(polygon)
                    except Exception:
                        pass  # Skip degenerate clusters
        plt.scatter(points[:, 0], points[:, 1], label=str(group), alpha=0.8, s=60, color=colors(i))
    plt.xlabel("t-SNE 1")
    plt.ylabel("t-SNE 2")
    plt.title("t-SNE by Group")
    plt.legend()
    plt.tight_layout()

    if is_streamlit:
        st_pyplot_func(plt, filename="tsne_by_group")
    else:
        plt.show()

def remove_outliers(df, col, group_col=None, threshold=5, by_group=True):
    if by_group and group_col is not None:
        # process each group separately, preserve original order
        def filter_group(group):
            mean = group[col].mean()
            std = group[col].std()
            return group[np.abs(group[col] - mean) <= threshold * std]
        
        # use groupby with group_keys=False to preserve order
        return (
            df.groupby(group_col, group_keys=False, sort=False)
              .apply(filter_group)
              .reset_index(drop=True)
        )
    else:
        # global filtering
        mean = df[col].mean()
        std = df[col].std()
        return df[np.abs(df[col] - mean) <= threshold * std].reset_index(drop=True)
    
def boxplot_plot_dabest(results_df, combined_df, col, output_dir,figures_dir=None,is_streamlit=False,analysis_type=None, show_histograms=False, selected_feature=None):
    # Remove outliers from each group
    cleaned_df = remove_outliers(combined_df, col, 'Group')
    unique_groups = cleaned_df['Group'].unique()
    # DABEST estimation plot only if exactly 2 groups
    if cleaned_df.groupby('Group')[col].count().min() > 0:
        # Ensure 'Group' is categorical and idx matches
        cleaned_df['Group'] = pd.Categorical(cleaned_df['Group'], categories=unique_groups.tolist(), ordered=True)
        dabest_obj = dabest.load(cleaned_df, x='Group', y=col, idx=unique_groups.tolist())

        # Extract p-value and Cohen's d from dabest results
        dabest_results = dabest_obj.mean_diff.results
        dabest_results['pvalue_mann_whitney'].values[0] if 'pvalue_mann_whitney' in dabest_results else np.nan
        dabest_results.columns

        def add_sig_lines(is_streamlit=False, est=None):
            # --- Add significance lines and asterisks (including '****') for ALL pairwise comparisons ---
            if est is not None:
                fig = est  # est is the Figure object
                ax = fig.axes[0]  # Get the main Axes
            else:
                ax = plt.gca()
            y_min, y_max = ax.get_ylim()
            ax.get_xticks()
            xticklabels = [tick.get_text() for tick in ax.get_xticklabels()]
            if not any(xticklabels):
                xticklabels = [str(g) for g in unique_groups]
            group_to_x = {str(g): i for i, g in enumerate(unique_groups)}
            group_y_max = {}
            for i, group in enumerate(unique_groups):
                ys = cleaned_df[cleaned_df['Group'] == group][col].values
                if len(ys) > 0:
                    group_y_max[group] = np.max(ys)
                else:
                    group_y_max[group] = y_max * 0.95
            base_height = max(group_y_max.values()) if group_y_max else y_max * 0.95
            y_stack = base_height + 0.05 * (y_max - y_min)
            y_increment = 0.05 * (y_max - y_min)
            # For all unique pairs of groups, compute p-value and annotate
            from itertools import combinations
            pval_lines = []
            for group1, group2 in combinations(unique_groups, 2):
                data1 = cleaned_df[cleaned_df['Group'] == group1][col].dropna()
                data2 = cleaned_df[cleaned_df['Group'] == group2][col].dropna()
                if len(data1) < 2 or len(data2) < 2:
                    continue
                try:
                    stat, pval = mannwhitneyu(data1, data2, alternative='two-sided')
                except Exception:
                    continue

                # Determine significance symbol
                if pval < 0.0001:
                    stars = '****'
                elif pval < 0.001:
                    stars = '***'
                elif pval < 0.01:
                    stars = '**'
                elif pval < 0.05:
                    stars = '*'
                else:
                    stars = 'ns'
                # Collect tuple for Streamlit display and logic
                if is_streamlit and len(unique_groups) > 1:
                    pval_lines.append((group1, group2, pval))
                x1 = group_to_x.get(str(group1), None)
                x2 = group_to_x.get(str(group2), None)
                if x1 is None or x2 is None:
                    continue
                # Draw the line
                x_coords = [x1, x1, x2, x2]
                y_coords = [y_stack, y_stack + y_increment, y_stack + y_increment, y_stack]
                ax.plot(x_coords, y_coords, lw=1, c='k')
                ax.text((x1 + x2) / 2, y_stack + y_increment + 0.01 * (y_max - y_min), stars,
                        ha='center', va='bottom', color='k', fontsize=12, fontweight='bold')
                y_stack += y_increment * 1.2
            return pval_lines
        pval_lines = add_sig_lines(is_streamlit=is_streamlit)

        # Compute r^2 (eta squared for ANOVA)
        group_means = cleaned_df.groupby('Group')[col].mean()
        grand_mean = cleaned_df[col].mean()
        ss_between = sum(cleaned_df.groupby('Group').size() * (group_means - grand_mean) ** 2)
        ss_total = sum((cleaned_df[col] - grand_mean) ** 2)
        ss_between / ss_total if ss_total > 0 else np.nan
        def cohen_d_power():
            # Format Cohen's d
            try:
                cd = dabest_obj.cohens_d.results
                cd['difference'].values[0] if 'difference' in cd.columns else np.nan
            except Exception as e:
                st.write(f"Could not compute Cohen's d: {e}")
                return 
            analysis = TTestIndPower()
            alpha = 0.05
            alternative = 'two-sided'  # or 'larger'/'smaller' if you had a directional hypothesis

            # Compute observed power per contrast
            cd = cd.copy()
            observed_power = [
                analysis.power(effect_size=abs(d),          # power depends on |d|
                            nobs1=int(n1),
                            ratio=float(n2)/float(n1),
                            alpha=alpha,
                            alternative=alternative)
                for d, n1, n2 in zip(cd["difference"], cd["control_N"], cd["test_N"])
            ]
            # If observed_power is nan, set to 1.0 if p < 0.05, else 0.0
            pval_col = "pvalue_mann_whitney" if "pvalue_mann_whitney" in dabest_results else "pvalue_permutation"
            pvals = dabest_results[pval_col].values if pval_col in dabest_results else [np.nan] * len(observed_power)
            for i, (power, pval) in enumerate(zip(observed_power, pvals)):
                if not np.isfinite(power):
                    if pval < 0.05:
                        observed_power[i] = 1.0
                    else:
                        observed_power[i] = 0.0
                else:
                    observed_power[i] = min(max(power, 0.0), 1.0)
        # cohen_d_power()
        # plt_title = f"""p = {p_value:.3e}, r² = {r_squared:.2f}
        # Cohen's d = {cohen_d:.2f} with 95% CI [{cd['bca_low'].values[0]:.2f}, {cd['bca_high'].values[0]:.2f}]
        # Observed power = {100*observed_power[0]:.2f}%\n\n"""
        # Use selected_feature if provided, otherwise use col
        title_text = selected_feature if selected_feature is not None else col
        # Build p-value string for title
        if pval_lines:
            min_pval = min(pval for _, _, pval in pval_lines)
            pval_str = f"p = {min_pval:.3e}"
        else:
            pval_str = ""
        plt_title = f"{title_text} - Comparison between Groups\n{pval_str}"
        # suptitle
        est = dabest_obj.mean_diff.plot(raw_marker_size=3, raw_label=col,float_contrast=True,fig_size=(8,8),dpi=300)
        plt.suptitle(plt_title, fontsize=14)
        ax = est.axes[0]
        annotate_pvals_with_jitter(ax, pval_lines=pval_lines, cleaned_df=cleaned_df, col=col, unique_groups=unique_groups)
        if figures_dir:
            os.makedirs(f'{figures_dir}/dabest_plots/{output_dir}', exist_ok=True)
            plt.tight_layout()
            plt.savefig(f"{figures_dir}/dabest_plots/{output_dir}/{col}_dabest_plot.png")
            plt.close()
        if is_streamlit:
            # Show p-value lines in Streamlit: first line, rest in expander
            if len(pval_lines) > 0 and any(pval < 0.05 for _, _, pval in pval_lines) or analysis_type =='Full':
                st.divider()
                # Use selected_feature if provided, otherwise use col
                title_text = selected_feature if selected_feature is not None else col
                st.subheader(f"Estimation Plot: {title_text}")
                # Build DataFrame from pval_lines
                pval_df = pd.DataFrame(pval_lines, columns=["Group 1", "Group 2", "P-value"])
                pval_df["P-value"] = pval_df["P-value"].map(lambda x: f"{x:.3e}")  # format in scientific notation
                output_df = pd.DataFrame()
                if len(pval_lines) > 0:
                    with st.expander("Show numeric results"):
                        st.write(plt_title)
                        # write mean ± std for each group
                        for group in combined_df['Group'].unique():
                            stat_dict = {}
                            stat_dict['Group'] = group
                            group_data = combined_df[combined_df['Group'] == group][col]
                            stats_text, stat_dict2 = stat_text_get(group_data)
                            stat_dict = {**stat_dict, **stat_dict2}
                            st.write(f"{group} mean {stat_dict['Mean']:.2e} ± {stat_dict['Std']:.2e}")
                            # add stat_dict to pval_df
                            # Find p-value for this group (if it exists in any comparison)
                            group_pvals = pval_df[(pval_df['Group 1'] == group) | (pval_df['Group 2'] == group)]['P-value']
                            if len(group_pvals) > 0:
                                stat_dict['P-value'] = group_pvals.iloc[0]  # Take first p-value found
                            else:
                                stat_dict['P-value'] = 'N/A'
                            stat_row = pd.DataFrame([stat_dict])
                            output_df = pd.concat([output_df, stat_row], ignore_index=True)
                        # Show as table instead of separate lines
                        st.dataframe(output_df, use_container_width=True)
                st_pyplot_func(plt, filename=f"{col}_dabest_plot", pval_df=output_df) 
    return # results_df
    figures_dir = 'figures'

from collections import defaultdict  # noqa: E402

def annotate_pvals_with_jitter(ax, pval_lines, cleaned_df, col, unique_groups):
    """Annotate pairwise p-values with U-bars + stars while:
       - collapsing duplicate (g1,g2) pairs (keep smallest p),
       - stacking overlapping bars,
       - adding tiny deterministic jitter to avoid exact overlap.
    """
    # ---------- collapse duplicate pairs (keep smallest p) ----------
    collapsed = {}
    for g1, g2, p in pval_lines:
        key = tuple(sorted((g1, g2)))
        if key not in collapsed or p < collapsed[key][0]:
            collapsed[key] = (p, key[0], key[1])
    pval_lines_collapsed = [(v[1], v[2], v[0]) for v in collapsed.values()]
    if not pval_lines_collapsed:
        return

    # ---------- helper: p -> star string ----------
    def pval_to_stars(p):
        if p < 0.0001:
            return '****'
        if p < 0.001:
            return '***'
        if p < 0.01:
            return '**'
        if p < 0.05:
            return '*'
        return None

    # ---------- helper: get label positions (fallback to indices) ----------
    def get_label_positions(ax, unique_groups):
        xticks = ax.get_xticks()
        xticklabels = [t.get_text().strip() for t in ax.get_xticklabels()]
        if any(lbl for lbl in xticklabels):
            mapping = {}
            # map first occurrence of label -> tick
            for lbl, x in zip(xticklabels, xticks):
                if lbl and lbl not in mapping:
                    mapping[lbl] = x
            # fallback for groups not found among ticklabels
            for i, g in enumerate(unique_groups):
                mapping.setdefault(g, xticks[i] if i < len(xticks) else i)
            return mapping
        else:
            return {g: i for i, g in enumerate(unique_groups)}

    label_pos = get_label_positions(ax, unique_groups)

    # ---------- plotting params ----------
    yrange = ax.get_ylim()[1] - ax.get_ylim()[0]
    spacer = 0.07 * yrange
    stack_counter = defaultdict(int)

    # sort so larger p (less significant) are drawn first; smallest p plotted last (on top)
    pval_sorted = sorted(pval_lines_collapsed, key=lambda t: t[2], reverse=True)

    for idx, (g1, g2, p) in enumerate(pval_sorted):
        stars = pval_to_stars(p)
        if not stars:
            continue

        # get numeric x positions (fallback to index positions)
        x1 = label_pos.get(g1)
        x2 = label_pos.get(g2)
        if x1 is None or x2 is None:
            try:
                i1, i2 = unique_groups.index(g1), unique_groups.index(g2)
            except ValueError:
                continue
            x1, x2 = i1, i2

        # ensure order and tiny horizontal nudge if identical
        if x2 < x1:
            x1, x2 = x2, x1
        if abs(x2 - x1) < 1e-6:
            x1 -= 0.02
            x2 += 0.02

        # stacking key for near-identical pairs
        key = (round(float(x1), 3), round(float(x2), 3))
        stack_index = stack_counter[key]

        # compute base y from data maxima of the two groups (safe fallback)
        try:
            y1 = cleaned_df[cleaned_df['Group'] == g1][col].max()
            y2 = cleaned_df[cleaned_df['Group'] == g2][col].max()
            base_y = max(float(y1) if not np.isnan(y1) else ax.get_ylim()[0],
                         float(y2) if not np.isnan(y2) else ax.get_ylim()[0])
        except Exception:
            base_y = ax.get_ylim()[1] * 0.9

        # stacking + tiny deterministic jitter to avoid exact overlap
        jitter = ((idx % 5) * 0.05) * yrange
        y = base_y + 0.05 * yrange + stack_index * spacer + jitter

        # expand ylim if needed so the annotation is not clipped
        top = ax.get_ylim()[1]
        if y + 0.12 * yrange > top:
            ax.set_ylim(ax.get_ylim()[0], y + 0.18 * yrange)
            yrange = ax.get_ylim()[1] - ax.get_ylim()[0]
            spacer = 0.07 * yrange

        cap = 0.02 * yrange

        # draw a thin U-shaped bar + stars
        ax.plot([x1, x1, x2, x2], [y, y + cap, y + cap, y],
                lw=1.0, color='k', clip_on=False, zorder=5, solid_capstyle='projecting')
        ax.text((x1 + x2) / 2, y + cap + 0.01 * yrange, stars,
                ha='center', va='bottom', fontsize=12, zorder=6, clip_on=False, alpha=0.95)

        stack_counter[key] += 1

from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

# initialize session state pptx only once
if "pptx" not in st.session_state:
    st.session_state.pptx = Presentation()
    st.session_state.pptx_sections = []  # Track sections for TOC
    st.session_state.current_section = None  # Track current section

def _add_title_slide(pptx, project_name=None):
    """Add a title slide to the presentation."""
    title_slide_layout = pptx.slide_layouts[0]  # Title slide layout
    slide = pptx.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "EEG Analysis Report"
    project_display = project_name if project_name else "EEG Analysis"
    subtitle.text = "Statistical Analysis and Visualization\n" + \
                   f"Generated from {project_display} Pipeline"

def _add_toc_slide(pptx, sections):
    """Add a table of contents slide."""
    blank_slide_layout = pptx.slide_layouts[6]  # Blank layout
    slide = pptx.slides.add_slide(blank_slide_layout)
    
    # Add title
    left = Inches(1)
    top = Inches(0.5)
    width = Inches(8)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Table of Contents"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Inches(0.4)
    title_para.font.bold = True
    
    # Add section list
    toc_top = Inches(1.5)
    toc_left = Inches(1)
    toc_width = Inches(8)
    
    for i, section in enumerate(sections, 1):
        text_box = slide.shapes.add_textbox(toc_left, toc_top + (i-1) * Inches(0.4), toc_width, Inches(0.35))
        text_frame = text_box.text_frame
        text_frame.text = f"{i}. {section}"
        para = text_frame.paragraphs[0]
        para.font.size = Inches(0.25)

def add_table_to_pptx(df, title="Group Statistics and Comparisons"):
    """Add a DataFrame as a table slide to the PPTX presentation."""
    if "pptx" not in st.session_state:
        st.session_state.pptx = Presentation()
        st.session_state.pptx_sections = []
    
    pptx = st.session_state.pptx
    blank_slide_layout = pptx.slide_layouts[6]  # Blank layout
    slide = pptx.slides.add_slide(blank_slide_layout)
    
    # Add title
    title_left = Inches(0.5)
    title_top = Inches(0.3)
    title_width = Inches(9)
    title_height = Inches(0.5)
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Inches(0.25)
    title_para.font.bold = True
    
    # Calculate table dimensions
    num_rows = len(df) + 1  # +1 for header
    num_cols = len(df.columns)
    
    # Position and size the table
    table_left = Inches(0.3)
    table_top = Inches(1.0)
    table_width = Inches(9.4)
    table_height = Inches(0.3 * num_rows)  # Dynamic height based on rows
    
    # Add table to slide
    table = slide.shapes.add_table(
        num_rows, num_cols,
        table_left, table_top, table_width, table_height
    ).table
    
    # Set column widths (distribute evenly)
    col_width = int(table_width / num_cols)
    for col_idx in range(num_cols):
        table.columns[col_idx].width = col_width
    
    # Populate header row
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Inches(0.12)
    
    # Populate data rows
    for row_idx, (_, row) in enumerate(df.iterrows(), start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Inches(0.10)
    
    # Track section
    if title not in st.session_state.pptx_sections:
        st.session_state.pptx_sections.append(title)

def _get_section_name_from_filename(filename):
    """Determine section name from filename."""
    filename_lower = filename.lower()
    if 'pairplot' in filename_lower:
        return 'Pair Plot Analysis'
    elif 'scatterplot' in filename_lower or 'regression' in filename_lower:
        return 'Scatter Plot with Regression Analysis'
    elif 'tsne' in filename_lower:
        return 't-SNE Dimensionality Reduction'
    elif 'dabest' in filename_lower or 'boxplot' in filename_lower:
        return 'DABEST Boxplot Comparison'
    elif 'topomap' in filename_lower:
        return 'Topographic Map (Topomap) Analysis'
    elif 'forest_plot' in filename_lower:
        return 'Forest Plot Meta-Analysis'
    elif 'average_trajectory' in filename_lower:
        return 'Average Trajectory Analysis'
    elif 'longitudinal' in filename_lower:
        return 'Longitudinal Analysis'
    elif any(x in filename_lower for x in ['eeg', 'raw', 'channel', 'delta']):
        return 'EEG Raw Data Visualization'
    elif 'spectrogram' in filename_lower:
        return 'Spectrogram Analysis'
    else:
        return 'General Analysis'

def _extract_column_name_from_filename(filename):
    """Extract column name(s) from filename patterns."""
    import re
    # Remove file extension
    name_without_ext = filename.replace('.png', '').replace('.svg', '')
    
    # Pattern 1: average_trajectory_{project}_{column} or longitudinal_{project}_{column}
    # Extract the last part after the last underscore (assuming it's the column)
    if 'average_trajectory' in name_without_ext.lower() or 'longitudinal' in name_without_ext.lower():
        parts = name_without_ext.split('_')
        # Find the index of 'trajectory' or 'longitudinal'
        if 'average_trajectory' in name_without_ext.lower():
            idx = name_without_ext.lower().find('average_trajectory')
            remaining = name_without_ext[idx + len('average_trajectory'):].lstrip('_')
        else:
            idx = name_without_ext.lower().find('longitudinal')
            remaining = name_without_ext[idx + len('longitudinal'):].lstrip('_')
        
        # Split remaining and take the last part as column name
        if remaining:
            parts = remaining.split('_')
            # Skip project name (usually first part), column is last
            if len(parts) > 1:
                return parts[-1]  # Last part is the column
            elif len(parts) == 1:
                return parts[0]
    
    # Pattern 2: scatterplot_{column}_vs_all_y_zscore
    if 'scatterplot' in name_without_ext.lower():
        match = re.search(r'scatterplot_([^_]+)_vs', name_without_ext.lower())
        if match:
            return match.group(1)
        # Alternative: scatterplot_{column}.png
        match = re.search(r'scatterplot_([^.]+)', name_without_ext.lower())
        if match:
            return match.group(1)
    
    # Pattern 3: pairplot_{columns_separated_by_underscores}
    if 'pairplot' in name_without_ext.lower():
        match = re.search(r'pairplot_(.+)', name_without_ext.lower())
        if match:
            cols = match.group(1)
            # If multiple columns, return them joined
            return cols.replace('_', ', ')
    
    # Pattern 4: {column}_dabest_plot
    if 'dabest' in name_without_ext.lower():
        match = re.search(r'^([^_]+)_dabest', name_without_ext.lower())
        if match:
            return match.group(1)
    
    # Pattern 5: {x_col}_vs_{y_col}_regression
    if 'regression' in name_without_ext.lower() and '_vs_' in name_without_ext:
        match = re.search(r'([^_]+)_vs_([^_]+)_regression', name_without_ext.lower())
        if match:
            return f"{match.group(1)} vs {match.group(2)}"
    
    # Pattern 6: Try to extract any column-like pattern (last meaningful segment)
    # This is a fallback for other patterns
    parts = name_without_ext.split('_')
    if len(parts) > 1:
        # Return the last part as it's likely the column name
        return parts[-1]
    
    return None

def _get_section_explanation(section_name):
    """Get explanation text for each section."""
    explanations = {
        'Pair Plot Analysis': 
            "Pair plots display pairwise relationships between multiple variables. "
            "Each plot shows scatterplots (lower triangle), histograms (diagonal), "
            "and optionally correlation coefficients. This helps identify relationships, "
            "distributions, and potential correlations between clinical and EEG features.",
        
        'Scatter Plot with Regression Analysis':
            "Scatter plots show the relationship between two continuous variables with a regression line. "
            "R² (R-squared) indicates the proportion of variance explained (0-1, higher is better). "
            "P-value indicates statistical significance of the relationship. "
            "Regression lines show the trend: positive slope indicates positive correlation, negative slope indicates negative correlation.",
        
        't-SNE Dimensionality Reduction':
            "t-SNE (t-Distributed Stochastic Neighbor Embedding) reduces high-dimensional data to 2D for visualization. "
            "Points closer together in the 2D plot are more similar in the original high-dimensional space. "
            "This helps identify clusters and patterns in complex EEG data. "
            "Different colors represent different groups for comparison.",
        
        'DABEST Boxplot Comparison':
            "DABEST (Data Analysis with Bootstrap ESTimation) plots show effect sizes between groups using bootstrap confidence intervals. "
            "Boxplots display median (line), quartiles (box), and outliers. "
            "The estimation plot shows mean difference with 95% confidence intervals. "
            "P-values indicate statistical significance (*** p<0.001, ** p<0.01, * p<0.05, ns=not significant). "
            "Cohen's d effect size: |d|<0.2 (negligible), 0.2-0.5 (small), 0.5-0.8 (medium), >0.8 (large).",
        
        'Topographic Map (Topomap) Analysis':
            "Topomaps show spatial distribution of EEG measures across the scalp. "
            "Color intensity represents the magnitude of the measure (e.g., power, p-values). "
            "P-value topomaps show statistical significance: darker colors indicate lower p-values (more significant). "
            "This visualizes which brain regions show significant differences between groups.",
        
        'Forest Plot Meta-Analysis':
            "Forest plots summarize effect sizes (e.g., Cohen's d) across multiple features. "
            "Each horizontal line represents a feature's effect size with confidence interval. "
            "Diamonds show overall effect. Features on the right favor one group, on the left favor the other. "
            "P-values and corrected p-values (accounting for multiple comparisons) are provided.",
        
        'Average Trajectory Analysis':
            "Shows the average progression of a measure over time across groups. "
            "Lines represent mean values with confidence intervals. "
            "R² indicates how well time explains the variance. "
            "P-values test whether slopes differ significantly between groups.",
        
        'Longitudinal Analysis':
            "Longitudinal plots show individual trajectories over time. "
            "Each line represents a subject's progression. "
            "This helps identify individual patterns, variability, and group differences in temporal trends.",
        
        'EEG Raw Data Visualization':
            "Raw EEG data displays show time-series electrical activity from different brain regions (channels). "
            "Amplitude (y-axis) represents voltage in microvolts, time (x-axis) in seconds. "
            "Delta power analysis identifies channels/regions with highest/lowest slow-wave activity (0.5-4 Hz), "
            "which is important for detecting brain states like sleep or pathological slowing.",
        
        'Spectrogram Analysis':
            "Spectrograms show frequency content over time using color intensity. "
            "X-axis: time, Y-axis: frequency, Color: power/intensity. "
            "This reveals temporal changes in brain rhythms (delta, theta, alpha, beta, gamma bands).",
        
        'General Analysis':
            "General statistical analyses and visualizations from the EEG processing pipeline."
    }
    return explanations.get(section_name, explanations['General Analysis'])

def _add_section_explanation_slide(pptx, section_name):
    """Add an explanation slide for a section."""
    blank_slide_layout = pptx.slide_layouts[6]  # Blank layout
    slide = pptx.slides.add_slide(blank_slide_layout)
    
    # Add section title
    left = Inches(1)
    top = Inches(0.5)
    width = Inches(8)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = section_name
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Inches(0.4)
    title_para.font.bold = True
    
    # Add explanation text
    explanation_top = Inches(1.5)
    explanation_left = Inches(1)
    explanation_width = Inches(8)
    explanation_height = Inches(5)
    text_box = slide.shapes.add_textbox(explanation_left, explanation_top, explanation_width, explanation_height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = _get_section_explanation(section_name)
    
    # Format paragraphs
    for para in text_frame.paragraphs:
        para.font.size = Inches(0.2)
        para.space_after = Inches(0.1)

def st_pyplot_func(fig, filename="plot",pval_df=None):
    """Show matplotlib figure in Streamlit, add SVG download,
       and append to in-memory PPTX stored in session_state."""
    
    # show plot
    st.pyplot(fig)

    # --- Add download for SVG ---
    svg_buffer = io.BytesIO()
    fig.savefig(svg_buffer, format="svg", bbox_inches="tight")
    svg_buffer.seek(0)
    st.download_button(
        label="Download this plot as SVG",
        data=svg_buffer,
        file_name=f"{filename}.svg",
        mime="image/svg+xml",
        key=str(uuid.uuid4())
    )

    # --- Initialize pptx in session_state if it doesn't exist ---
    if "pptx" not in st.session_state:
        st.session_state.pptx = Presentation()
        st.session_state.pptx_sections = []
        st.session_state.current_section = None

    # --- Detect section change and add explanation slide if needed ---
    section_name = _get_section_name_from_filename(filename)
    if section_name != st.session_state.current_section:
        # New section detected
        if section_name not in st.session_state.pptx_sections:
            st.session_state.pptx_sections.append(section_name)
        # Add explanation slide before the plot
        _add_section_explanation_slide(st.session_state.pptx, section_name)
        st.session_state.current_section = section_name

    # --- Append PNG to pptx in session_state ---
    png_buffer = io.BytesIO()
    fig.savefig(png_buffer, format="png", dpi=300, bbox_inches="tight")
    png_buffer.seek(0)

    slide = st.session_state.pptx.slides.add_slide(
        st.session_state.pptx.slide_layouts[6]  # blank layout
    )
    
    # Extract column name from filename
    column_name = _extract_column_name_from_filename(filename)
    
    # Build slide title with section name and column name
    if column_name:
        slide_title = f"{section_name}: {column_name}"
    else:
        slide_title = section_name
    
    # Add section label at top of slide
    left_label = Inches(0.5)
    top_label = Inches(0.2)
    width_label = Inches(9)
    height_label = Inches(0.3)
    label_box = slide.shapes.add_textbox(left_label, top_label, width_label, height_label)
    label_frame = label_box.text_frame
    label_frame.text = slide_title
    label_para = label_frame.paragraphs[0]
    label_para.font.size = Inches(0.18)
    label_para.font.italic = True
    # Use theme default color by not explicitly setting RGB
    
    # Add the plot image (adjusted position to account for label)
    slide.shapes.add_picture(png_buffer, Inches(1), Inches(0.6), height=Inches(5))
    
    # --- Add pval_df as table to the same slide ---
    if pval_df is not None and not pval_df.empty:
        # Position the table below the figure
        table_left = Inches(1)
        table_top = Inches(6.3)  # Adjusted for label
        table_width = Inches(8)
        table_height = Inches(1.5)
        
        # Convert DataFrame to table format for PowerPoint
        table_data = [list(pval_df.columns)]  # Header row
        for _, row in pval_df.iterrows():
            table_data.append(list(row))
        
        # Add table to slide
        table = slide.shapes.add_table(
            len(table_data), len(table_data[0]),
            table_left, table_top, table_width, table_height
        ).table
        
        # Populate table with data
        for i, row_data in enumerate(table_data):
            for j, cell_data in enumerate(row_data):
                table.cell(i, j).text = str(cell_data)
                # Style the header row
                if i == 0:
                    table.cell(i, j).text_frame.paragraphs[0].font.bold = True


def _copy_slide_shapes(source_slide, target_slide):
    """Copy all shapes from source slide to target slide."""
    try:
        from pptx.shapes.picture import Picture
        PictureClass = Picture
    except ImportError:
        PictureClass = None
    
    for shape in source_slide.shapes:
        try:
            # Check if it's a picture shape
            is_picture = (PictureClass and isinstance(shape, PictureClass)) or \
                        (hasattr(shape, 'image') and hasattr(shape.image, 'blob'))
            if is_picture:
                # Copy image
                img_stream = io.BytesIO(shape.image.blob)
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                target_slide.shapes.add_picture(img_stream, left, top, width, height)
            elif hasattr(shape, 'table'):
                # Copy table
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                table_data = []
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text)
                    table_data.append(row_data)
                if table_data:
                    new_table = target_slide.shapes.add_table(
                        len(table_data), len(table_data[0]),
                        left, top, width, height
                    ).table
                    for i, row_data in enumerate(table_data):
                        for j, cell_text in enumerate(row_data):
                            new_table.cell(i, j).text = str(cell_text)
                            if i == 0:
                                new_table.cell(i, j).text_frame.paragraphs[0].font.bold = True
            elif hasattr(shape, 'text_frame'):
                # Copy text box
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                text_box = target_slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
                text_frame.text = shape.text_frame.text
                # Copy formatting for first paragraph
                if shape.text_frame.paragraphs:
                    source_para = shape.text_frame.paragraphs[0]
                    if text_frame.paragraphs:
                        text_frame.paragraphs[0].font.size = source_para.font.size
                        text_frame.paragraphs[0].font.bold = source_para.font.bold
                        text_frame.paragraphs[0].font.italic = source_para.font.italic
        except Exception:
            # Skip shapes that can't be copied
            continue

def download_pptx_button(label="Download all plots as PPTX"):
    """Provide a Streamlit button to download the aggregated PPTX."""
    # Ensure pptx and sections are initialized
    if "pptx" not in st.session_state:
        st.session_state.pptx = Presentation()
        st.session_state.pptx_sections = []
    
    # Create a new presentation with title, TOC, and then all existing slides
    final_pptx = Presentation()
    
    # Add title slide with project name from session state
    project_name = st.session_state.get('project_name', None)
    _add_title_slide(final_pptx, project_name)
    
    # Add TOC slide if we have sections
    if st.session_state.pptx_sections:
        _add_toc_slide(final_pptx, st.session_state.pptx_sections)
    
    # Copy all slides from original presentation
    for source_slide in st.session_state.pptx.slides:
        # Create blank slide in target presentation
        blank_slide_layout = final_pptx.slide_layouts[6]  # Blank layout
        target_slide = final_pptx.slides.add_slide(blank_slide_layout)
        
        # Copy all shapes from source to target
        _copy_slide_shapes(source_slide, target_slide)
    
    # Save to buffer
    pptx_buffer = io.BytesIO()
    final_pptx.save(pptx_buffer)
    pptx_buffer.seek(0)

    st.sidebar.download_button(
        label=label,
        data=pptx_buffer,
        file_name="plots.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        key="pptx-download",
        use_container_width=False
    )

def download_excel_button(label="Download analysis results as Excel", patient_df=None):
    """Provide a Streamlit button to download patient-level analysis results as Excel."""
    # Try to find patient-level data in session state if not provided
    if patient_df is None:
        # Check common session state keys for patient data (prioritize df_wnv3)
        possible_keys = [
            'df_wnv3', 'df_wnv2', 'combined_df', 'patient_results_df', 'analysis_results_df', 
            'patient_df', 'results_df', 'df_wnv', 'controls_df'
        ]
        
        patient_df = None
        for key in possible_keys:
            if key in st.session_state and isinstance(st.session_state[key], pd.DataFrame):
                patient_df = st.session_state[key]
                break
        
        # If still not found, try to combine common dataframes
        if patient_df is None:
            dfs_to_combine = []
            if 'df_wnv3' in st.session_state and isinstance(st.session_state['df_wnv3'], pd.DataFrame):
                dfs_to_combine.append(st.session_state['df_wnv3'])
            elif 'df_wnv2' in st.session_state and isinstance(st.session_state['df_wnv2'], pd.DataFrame):
                dfs_to_combine.append(st.session_state['df_wnv2'])
            if 'df_wnv' in st.session_state and isinstance(st.session_state['df_wnv'], pd.DataFrame):
                dfs_to_combine.append(st.session_state['df_wnv'])
            if 'controls_df' in st.session_state and isinstance(st.session_state['controls_df'], pd.DataFrame):
                dfs_to_combine.append(st.session_state['controls_df'])
            if 'combined_df' in st.session_state and isinstance(st.session_state['combined_df'], pd.DataFrame):
                dfs_to_combine.append(st.session_state['combined_df'])
            
            if dfs_to_combine:
                patient_df = dfs_to_combine[0]  # Use the first (most relevant) dataframe
    
    # If no data found, show warning and return
    if patient_df is None or patient_df.empty:
        st.sidebar.warning("No patient data available for export.")
        return
    
    # Aggregate by patient if patient_id column exists
    if 'patient_id' in patient_df.columns:
        # Get numeric and non-numeric columns
        numeric_cols = patient_df.select_dtypes(include=[np.number]).columns.tolist()
        non_numeric_cols = [col for col in patient_df.columns 
                           if col not in numeric_cols and col not in ['patient_id']]
        
        # Group by patient_id and aggregate
        agg_dict = {col: 'mean' for col in numeric_cols}
        agg_dict.update({col: 'first' for col in non_numeric_cols})
        
        patient_df_agg = patient_df.groupby('patient_id').agg(agg_dict).reset_index()
    elif 'patient_number' in patient_df.columns:
        # Alternative patient identifier
        numeric_cols = patient_df.select_dtypes(include=[np.number]).columns.tolist()
        non_numeric_cols = [col for col in patient_df.columns 
                           if col not in numeric_cols and col not in ['patient_number']]
        
        agg_dict = {col: 'mean' for col in numeric_cols}
        agg_dict.update({col: 'first' for col in non_numeric_cols})
        
        patient_df_agg = patient_df.groupby('patient_number').agg(agg_dict).reset_index()
    else:
        # No patient identifier found, use data as-is
        patient_df_agg = patient_df.copy()
    
    # Convert timezone-aware datetime columns to timezone-naive for Excel compatibility
    def make_timezone_naive(df):
        """Convert timezone-aware datetime columns to timezone-naive."""
        df_copy = df.copy()
        for col in df_copy.columns:
            if pd.api.types.is_datetime64_any_dtype(df_copy[col]):
                # Check if timezone-aware
                if df_copy[col].dt.tz is not None:
                    # Convert to UTC then remove timezone
                    df_copy[col] = df_copy[col].dt.tz_convert('UTC').dt.tz_localize(None)
        return df_copy
    
    # Apply timezone conversion
    patient_df_agg = make_timezone_naive(patient_df_agg)
    
    # Create Excel file in memory
    excel_buffer = io.BytesIO()
    try:
        with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
            # Write main patient data
            patient_df_agg.to_excel(writer, index=False, sheet_name='Patient Results')
            
            # If results_df exists in session state, add it as a separate sheet
            if 'results_df' in st.session_state and isinstance(st.session_state['results_df'], pd.DataFrame):
                if not st.session_state['results_df'].empty:
                    results_df_clean = make_timezone_naive(st.session_state['results_df'])
                    results_df_clean.to_excel(writer, index=False, sheet_name='Statistical Results')
        
        excel_buffer.seek(0)
        
        st.sidebar.download_button(
            label=label,
            data=excel_buffer.getvalue(),
            file_name="analysis_results.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            key="excel-download",
            use_container_width=False
        )
    except ImportError:
        st.sidebar.error("openpyxl is required for Excel export. Please install it with: pip install openpyxl")
    except Exception as e:
        st.sidebar.error(f"Error creating Excel file: {str(e)}")

def scatter_plot_with_regression(results_df, combined_df, x_col, y_col, output_dir, figures_dir=None, is_streamlit=False, analysis_type=None, show_histograms=False):
    plt.figure(figsize=(10, 6))
    sns.scatterplot(x=x_col, y=y_col, data=combined_df, alpha=0.5, color='black')
    sns.regplot(x=x_col, y=y_col, data=combined_df, scatter=False, color='blue')
    # Perform linear regression
    X = sm.add_constant(combined_df[x_col])
    y = combined_df[y_col]
    model = sm.OLS(y, X).fit()
    p_value = model.pvalues[1]  # p-value for the slope
    r_squared = model.rsquared  # R-squared value
    # Determine significance symbol
    if (p_value < 0.001):
        sig_symbol = '***'
    elif (p_value < 0.01):
        sig_symbol = '**'
    elif (p_value < 0.05):
        sig_symbol = '*'
    else:
        sig_symbol = 'ns'
    if sig_symbol != 'ns':
        # Add stats results to the title
        plt.title(
            f"{x_col} vs {y_col} Regression\n"
            f"Slope p = {p_value:.3e} ({sig_symbol}), R^2 = {r_squared:.2f}, n = {len(combined_df)}",
            ha='center'
        )
        plt.xlabel(x_col)
        plt.ylabel(y_col)
        plt.tight_layout()
        # Make folder {figures_dir}
        if is_streamlit:
            st.divider()
            st.subheader(f"Scatterplot of {x_col} vs {y_col}")
            # write {x_col} mean ± std
            st.write(f"Mean: {x_col} = {combined_df[x_col].mean():.2e} ± {combined_df[x_col].std():.2e}")
            st.write(f"Mean: {y_col} = {combined_df[y_col].mean():.2e} ± {combined_df[y_col].std():.2e}")
            st.write(f"P-value {p_value:.3e}, R^2 {r_squared:.2f}")
            st_pyplot_func(plt,filename=f"{x_col}_vs_{y_col}_regression")
        else:
            os.makedirs(f'{figures_dir}/scatterplots/{output_dir}', exist_ok=True)
            plt.savefig(f"{figures_dir}/scatterplots/{output_dir}/{y_col}_regression.png")
            plt.close()
    
def analyze_and_correct(combined_df, columns_to_analyze, groups=['Control', 'WNV']):
    def analyze_groups(combined_df, col, groups):
        results = []
        for i, group1 in enumerate(groups):
            for group2 in groups[i+1:]:
                control_data = combined_df[combined_df['Group'] == group1][col].dropna()
                case_data = combined_df[combined_df['Group'] == group2][col].dropna()
                if len(control_data) < 2 or len(case_data) < 2:
                    continue
                # Normality test
                _, normal_p = stats.normaltest(combined_df[col].dropna())
                # Choose appropriate test
                if normal_p < 0.05:  # Non-parametric
                    stat, p = stats.mannwhitneyu(control_data, case_data)
                    test_used = "Mann-Whitney U"
                else:  # Parametric
                    stat, p = stats.ttest_ind(control_data, case_data)
                    test_used = "T-test"
                # Effect size
                cohen_d = (case_data.mean() - control_data.mean()) / np.sqrt((
                    (len(case_data)-1)*case_data.std()**2 + 
                    (len(control_data)-1)*control_data.std()**2) / 
                    (len(case_data) + len(control_data) - 2))
                results.append({
                    'Variable': col,
                    'Group1': group1,
                    'Group2': group2,
                    'Test': test_used,
                    'Statistic': stat,
                    'p_value': p,
                    'Cohen_d': cohen_d,
                    'Mean_Group1': control_data.mean(),
                    'Mean_Group2': case_data.mean(),
                    'Std_Group1': control_data.std(),
                    'Std_Group2': case_data.std(),
                    'median_Group1': control_data.median(),
                    'median_Group2': case_data.median(),
                    'MAD_Group1': stats.median_abs_deviation(control_data),
                    'MAD_Group2': stats.median_abs_deviation(case_data),
                })
        return results

    all_results = []    
    # Statistical analysis
    for col in columns_to_analyze:
        results = analyze_groups(combined_df, col, groups)
        if results:
            all_results.extend(results)

    # Create results DataFrame
    results_df = pd.DataFrame(all_results)

    # Multiple testing correction
    if not results_df.empty:
        rej, adj_p, _, _ = multipletests(results_df['p_value'], method='fdr_bh')
        results_df['adj_p_value'] = adj_p
        results_df['Significant'] = adj_p < 0.05

    return results_df

def topomap_group_data( band, montage,control_data,wnv_data,output_dir,figures_dir,is_streamlit=False):
    # Calculate p-values for each channel
    common_channels = control_data.columns.intersection(wnv_data.columns)
    dict_p_values = {}
    for common_channel in common_channels:
        control_channel = control_data[common_channel].dropna()
        wnv_channel = wnv_data[common_channel].dropna()
        if len(control_channel) < 2 or len(wnv_channel) < 2:
            continue
        _, p = stats.mannwhitneyu(control_channel, wnv_channel)
        dict_p_values[common_channel] = p
    df_p_values = pd.DataFrame(dict_p_values, index=[0]).T
    reject, pvals_corrected, _, _ = smm.multipletests(df_p_values[0].values, alpha=0.05, method='fdr_bh')
    df_p_values['pvals_corrected'] = pvals_corrected
    if any(df_p_values['pvals_corrected'] < 0.05):
        ch_names = df_p_values.index.tolist()
        # Create an info object
        info = mne.create_info(ch_names=ch_names, sfreq=256, ch_types='eeg')
        info.set_montage(montage)
        # Create an EvokedArray object for p-values
        p_evoked = mne.EvokedArray(df_p_values['pvals_corrected'].values.reshape(-1, 1), info)
        # Plot the topomap of p-values
        fig, ax = plt.subplots()
        vlim_max = min(0.05, df_p_values['pvals_corrected'].max())
        im, cm = mne.viz.plot_topomap(p_evoked.data[:, 0], p_evoked.info, axes=ax, show=False, cmap='jet_r', vlim=[0, vlim_max])
        fig.colorbar(im, ax=ax)
        plt.title(f"{band} {output_dir} P-Value Topomap")
        # if any value in pvals_corrected is less than 0.05
        # make folder {figures_dir}/boxplots/{output_dir}
        # Save the figure
        if is_streamlit:
            st.write(f"Topomap for {band} {output_dir} P-Value")
            st_pyplot_func(plt)
        else:
            os.makedirs(f'{figures_dir}/topomaps_p_values/{output_dir}', exist_ok=True)
            plt.savefig(f"{figures_dir}/topomaps_p_values/{output_dir}/p_values_{band}_topomap.png")
            plt.close()

def process_group_data(group, run_df, frequency_bands, eeg_dict_convertion, eeg_channels, montage,group_data):
    # get only columns that say EEG
    eeg_df = run_df.filter(like='EEG')
    group_data[group] = {}
    for band in frequency_bands:
        # get the columns which have the band name
        power_df = eeg_df.filter(like=band)
        # column name split ' '[-1]
        power_df.columns = power_df.columns.str.split(' ').str[-1]
        power_df.columns = [eeg_dict_convertion[col] if col in eeg_dict_convertion else col for col in power_df.columns]
        # leave only the eeg channels
        # Get the channels that exist in the DataFrame
        existing_channels = [ch for ch in eeg_channels if ch in power_df.columns]
        # Filter the DataFrame to include only the existing channels
        power_df = power_df[existing_channels]
        # drop columns more than half of the values are NaN
        power_df = power_df.dropna(axis=1, thresh=power_df.shape[0]//2)
        # drop duplicates columns
        power_df.columns
        # change column names to only be the channel. if not
        # Extract the power values for the current band
        # power_values = power_df[band].values
        ch_names = power_df.columns.tolist()
        # Create an info object
        info = mne.create_info(ch_names=ch_names, sfreq=256, ch_types='eeg')
        info.set_montage(montage)
        power_values = power_df.T.values
        group_data[group][band] = power_df

        # Create an EvokedArray object
        evoked = mne.EvokedArray(power_values, info)
        # Plot the topomap
        fig, ax = plt.subplots()
        im, cm = mne.viz.plot_topomap(evoked.data[:, 0], evoked.info, axes=ax, show=False)
        fig.colorbar(im, ax=ax)
        # plt.title(f"{band} Topomap")
        # # Save the figure
        # os.makedirs(f'{figures_dir}/topomaps', exist_ok=True)
        # plt.savefig(f"{figures_dir}/topomaps/{group}_{band}_topomap.png")
        # plt.close()
    return group_data 

def weighted_avg(df, weight_col, numeric_cols):
    # Calculate weighted average for numeric columns, ignoring NaNs
    weighted_df = df[numeric_cols].multiply(df[weight_col], axis=0).sum(skipna=True) / df[weight_col].sum(skipna=True)
    # Preserve non-numeric columns by taking the first non-null value
    non_numeric_cols = df.columns.difference(numeric_cols)
    preserved_df = df[non_numeric_cols].iloc[0]
    # Combine the results
    return pd.concat([preserved_df, weighted_df])
#%% WNV
def wnv_get_files():
    # load clinical data from WNV_merged_291224_KP.xlsx
    df_wnv = pd.read_excel('WNV_merged_291224_KP.xlsx')
    # df_wnv replace '.' in column names with '_'
    df_wnv.columns = df_wnv.columns.str.replace('.', '_')
    # df_wnv replace 'NA' with np.nan
    df_wnv = df_wnv.replace('NA', np.nan)
    # Configuration
    patients_folder = "west_nile_virus"
    case_file = f"{patients_folder}.csv"
    # Read and prepare data
    controls = pd.read_csv(f'{patients_folder}_controls.csv')
    cases = pd.read_csv(case_file)
    wnv_ids = [file.split('/')[-2] for file in cases['file_path']]
    # to int
    wnv_ids = [int(id) for id in wnv_ids]
    cases['ID'] = wnv_ids
    # merge the dataframes
    df_merged_outer = pd.merge(df_wnv, cases, on='ID', how='outer',indicator=True)
    df_merged = df_merged_outer[df_merged_outer['_merge'] == 'both']
    wnv_files = os.listdir('west_nile_virus')
    # remove .DS_Store
    wnv_files = [file for file in wnv_files if 'DS_Store' not in file]
    wnv_files = [file.split('.edf')[0] for file in wnv_files]
    # also split '-'
    wnv_files = [file.split('-')[0] for file in wnv_files]
    # remove duplicates
    wnv_files = list(set(wnv_files))
    # to int
    wnv_files = [int(file) for file in wnv_files]
    # get df in column ID matches with wnv_files
    # print what wnv_files are not in df
    print([file for file in wnv_files if file not in df_wnv['ID'].values])
    df_wnv2 = df_wnv[df_wnv['ID'].isin(wnv_files)]
    # avg lines with same ID
    numeric_cols = df_merged.select_dtypes(include=[np.number]).columns
    # Group by ID and calculate the mean of each numeric column
    df_wnv2 = df_merged.groupby('ID').apply(weighted_avg, weight_col='duration_min', numeric_cols=numeric_cols).reset_index(drop=True)
    # delta = MRS_FOLLOW_UP - MRS_prior
    df_wnv2['MRS_delta_follow_up'] = df_wnv2['MRS_FOLLOW_UP'] - df_wnv2['MRS_prior']
    df_wnv2['MRS_delta_at_peak_illness'] = df_wnv2['MRS_at_peak_illness'] - df_wnv2['MRS_prior']
    df_wnv2['MRS_delta_peak_minus_follow_up'] = df_wnv2['MRS_at_peak_illness'] - df_wnv2['MRS_FOLLOW_UP']
    cases_group_name = 'WNV'
    return df_wnv,patients_folder,controls,df_wnv2,cases_group_name 
    # The following lines were unreachable after return and have been removed:
    # df_wnv2.to_csv('wnv_grouped.csv', index=False)
    # df_wnv2['ENCEPHALITIS'].sum()
    # df_wnv2['MENINGITIS'].sum()
    # [col for col in df_wnv2.columns if 'MRS' in col]
    # df_wnv2.columns.tolist()
    # df_merged['ID'].to_list()

def find_consecutive_sequences(events, min_duration_sec=5):
    sequences = []
    temp_seq = [events[0]]

    for i in range(1, len(events)):
        if events[i] == events[i - 1] + min_duration_sec:
            temp_seq.append(events[i])
        else:
            # get sum diff
            temp_seq_sum_diff = sum(np.diff(temp_seq))
            if min_duration_sec < 5:
                if temp_seq_sum_diff >= 5:
                    sequences.append(temp_seq)
            else:
                sequences.append(temp_seq)
            temp_seq = [events[i]]

    if len(temp_seq) >= 5:  # Check the last sequence
        sequences.append(temp_seq)

    return sequences

def read_edf_mne(file_path):
    """
    Read EDF file using MNE with error handling for malformed files.
    Tries multiple approaches if the default fails.
    """
    # Try different approaches to read the EDF file
    read_attempts = [
        # Attempt 1: Default with latin1 encoding
        {'preload': True, 'encoding': 'latin1', 'verbose': False},
        # Attempt 2: Without preload first, then load
        {'preload': False, 'encoding': 'latin1', 'verbose': False},
        # Attempt 3: Try without encoding specification
        {'preload': True, 'verbose': False},
        # Attempt 4: Try with preload=False and no encoding
        {'preload': False, 'verbose': False},
        # Attempt 5: Try with exclude parameter to skip problematic channels
        {'preload': True, 'encoding': 'latin1', 'exclude': [], 'verbose': False},
        # Attempt 6: Try with stim_channel=None (for EDF files without stim channels)
        {'preload': True, 'encoding': 'latin1', 'stim_channel': None, 'verbose': False},
        # Attempt 7: Try with exclude and stim_channel=None
        {'preload': True, 'encoding': 'latin1', 'exclude': [], 'stim_channel': None, 'verbose': False},
    ]
    
    raw = None
    last_error = None
    
    for i, params in enumerate(read_attempts):
        try:
            raw = mne.io.read_raw_edf(file_path, **params)
            # If preload was False, load the data now
            if not params.get('preload', True):
                raw.load_data()
            # Success - break out of loop
            break
        except Exception as e:
            last_error = e
            error_str = str(e)
            error_repr = repr(e)
            type(e).__name__
            # Check for the specific empty string to int error in multiple ways
            is_int_error = (
                "invalid literal for int()" in error_str or 
                "base 10" in error_str or 
                "''" in error_str or
                "invalid literal for int()" in error_repr or
                "base 10" in error_repr or
                "''" in error_repr or
                (isinstance(e, ValueError) and ("int()" in error_str or "base 10" in error_str))
            )
            
            if is_int_error:
                print(f"Warning: MNE attempt {i+1} failed with int conversion error for {os.path.basename(file_path)}: {error_str}")
                print("  Trying pyedflib as alternative reading method...")
                # Try using pyedflib as fallback immediately when we detect this error
                try:
                    import pyedflib
                    # Read with pyedflib and convert to MNE format
                    edf = pyedflib.EdfReader(file_path)
                    n_channels = edf.signals_in_file
                    ch_names = [edf.getSignalLabel(i) for i in range(n_channels)]
                    sfreqs = [edf.getSampleFrequency(i) for i in range(n_channels)]
                    # Use the most common sampling frequency, or first one if all same
                    sfreq = sfreqs[0] if sfreqs else 256
                    if len(set(sfreqs)) > 1:
                        # If different sampling rates, use the most common one
                        from collections import Counter
                        sfreq = Counter(sfreqs).most_common(1)[0][0]
                    
                    # Get data
                    data = np.array([edf.readSignal(i) for i in range(n_channels)])
                    # Create MNE info - use 'misc' for unknown channel types
                    ch_types = ['misc'] * n_channels
                    info = mne.create_info(ch_names=ch_names, sfreq=sfreq, ch_types=ch_types)
                    # Create Raw object
                    raw = mne.io.RawArray(data, info)
                    edf.close()
                    print("Successfully read file using pyedflib fallback")
                    break
                except ImportError:
                    print("pyedflib not available for fallback reading")
                    # Continue with next MNE attempt if pyedflib not available
                    if i < len(read_attempts) - 1:
                        continue
                except Exception as fallback_error:
                    print(f"pyedflib fallback also failed: {fallback_error}")
                    # Continue with next MNE attempt if pyedflib fails
                    if i < len(read_attempts) - 1:
                        continue
            else:
                # For non-int errors, continue with next attempt or re-raise on last attempt
                if i < len(read_attempts) - 1:
                    print(f"Warning: Attempt {i+1} failed for {os.path.basename(file_path)}: {error_str}")
                    continue
                else:
                    # Last attempt failed, try pyedflib as final fallback
                    print("Warning: All MNE attempts failed. Trying pyedflib as final fallback...")
                    try:
                        import pyedflib
                        # Read with pyedflib and convert to MNE format
                        edf = pyedflib.EdfReader(file_path)
                        n_channels = edf.signals_in_file
                        ch_names = [edf.getSignalLabel(i) for i in range(n_channels)]
                        sfreqs = [edf.getSampleFrequency(i) for i in range(n_channels)]
                        # Use the most common sampling frequency, or first one if all same
                        sfreq = sfreqs[0] if sfreqs else 256
                        if len(set(sfreqs)) > 1:
                            # If different sampling rates, use the most common one
                            from collections import Counter
                            sfreq = Counter(sfreqs).most_common(1)[0][0]
                        
                        # Get data
                        data = np.array([edf.readSignal(i) for i in range(n_channels)])
                        # Create MNE info - use 'misc' for unknown channel types
                        ch_types = ['misc'] * n_channels
                        info = mne.create_info(ch_names=ch_names, sfreq=sfreq, ch_types=ch_types)
                        # Create Raw object
                        raw = mne.io.RawArray(data, info)
                        edf.close()
                        print("Successfully read file using pyedflib fallback")
                        break
                    except ImportError:
                        print("pyedflib not available for fallback reading")
                    except Exception as fallback_error:
                        print(f"pyedflib fallback also failed: {fallback_error}")
                    # Re-raise the original error if pyedflib also fails
                    raise
    
    # If all attempts failed, raise the last error with context
    if raw is None:
        error_msg = f"Failed to read EDF file {file_path} after multiple attempts."
        if last_error:
            error_msg += f" Last error: {last_error}"
        raise ValueError(error_msg)
    
    # Try to create metadata from raw object
    try:
        metadata = {
            'file_name': os.path.basename(file_path),
            'start_date': raw.info['meas_date'],
            'duration_sec': raw.times[-1],
            'duration_min': raw.times[-1] / 60,
            'number_of_signals': len(raw.ch_names),
            'signal_labels': raw.ch_names,
            'sampling_frequency': raw.info['sfreq'],
            'highpass': raw.info['highpass'],
            'lowpass': raw.info['lowpass'],
            'annotations': raw.annotations if raw.annotations else None,
            'n_samples': raw.n_times,
            'bad_channels': raw.info['bads']
        }
    except Exception as e:
        # If metadata creation fails, create minimal metadata and warn
        print(f"Warning: Could not extract full metadata from {os.path.basename(file_path)}: {e}")
        print("  Creating minimal metadata...")
        try:
            metadata = {
                'file_name': os.path.basename(file_path),
                'number_of_signals': len(raw.ch_names) if hasattr(raw, 'ch_names') else 0,
                'signal_labels': raw.ch_names if hasattr(raw, 'ch_names') else [],
                'sampling_frequency': raw.info['sfreq'] if hasattr(raw, 'info') and 'sfreq' in raw.info else None,
                'n_samples': raw.n_times if hasattr(raw, 'n_times') else None,
            }
            # Try to get other fields with defaults
            if hasattr(raw, 'times') and len(raw.times) > 0:
                metadata['duration_sec'] = raw.times[-1]
                metadata['duration_min'] = raw.times[-1] / 60
            else:
                metadata['duration_sec'] = None
                metadata['duration_min'] = None
            
            if hasattr(raw, 'info'):
                metadata['start_date'] = raw.info.get('meas_date', None)
                metadata['highpass'] = raw.info.get('highpass', None)
                metadata['lowpass'] = raw.info.get('lowpass', None)
                metadata['bad_channels'] = raw.info.get('bads', [])
            else:
                metadata['start_date'] = None
                metadata['highpass'] = None
                metadata['lowpass'] = None
                metadata['bad_channels'] = []
            
            metadata['annotations'] = raw.annotations if hasattr(raw, 'annotations') and raw.annotations else None
        except Exception as e2:
            # If even minimal metadata fails, raise an error
            error_msg = f"Failed to extract metadata from EDF file {file_path} after reading."
            error_msg += f" Reading succeeded but metadata extraction failed: {e2}"
            raise ValueError(error_msg) from e2
    
    return metadata, raw

def eeg_data_to_features(raw, window_size_sec=5, min_duration_sec=5):
    raw_eeg = raw.copy().pick_types(eeg=True)
    eeg_data = raw_eeg.get_data()
    sf = raw_eeg.info['sfreq']
    window_size = int(window_size_sec * sf)

    # Initialize storage
    psd_list = []
    fft_list = []
    mpf_list = []
    df_list = []
    dfv_list = []
    pswe_mpf_avg = []
    pswe_events_per_channel = []

    # --- NEW: Initialize network feature arrays per band ---
    band_feature_arrays = {band: {'efficiency': [], 'clustering': [], 'assortativity': [], 'modularity': []}
                          for band in power_bands}

    # Loop channels
    for i, channel_data in enumerate(eeg_data):
        ch_psd, ch_fft, ch_mpf, ch_df, ch_pswe,pswe_mpf_list_temp = [], [], [], [], [],[]
        # Slide windows
        for start in range(0, len(channel_data) - window_size + 1, window_size):
            segment = eeg_data[:,start:start + window_size]
            win = channel_data[start:start + window_size]
            # PSD
            psd_vals, freqs = mne.time_frequency.psd_array_welch(
                win, sf, fmin=1, fmax=40, n_fft=int(sf)
            )
            psd_vals = psd_vals.squeeze()

            # FFT
            fft_vals = np.fft.fft(win)

            # MPF and dominant freq
            mpf = np.sum(psd_vals * freqs) / np.sum(psd_vals)
            dfreq = freqs[np.argmax(psd_vals)]

            ch_psd.append(psd_vals)
            ch_fft.append(fft_vals)
            ch_mpf.append(mpf)
            ch_df.append(dfreq)

            # PSWE detection
            if mpf < 6.0:
                ch_pswe.append(start / sf)
                pswe_mpf_list_temp.append(mpf)
            # --- MODIFIED: Compute network features for each window and band ---
            if i == 0:
                for band_name, (fmin, fmax) in power_bands.items():
                    efficiency, clustering, assortativity, modularity = compute_network_features(segment, sf, [fmin, fmax])
                    band_feature_arrays[band_name]['efficiency'].append(efficiency)
                    band_feature_arrays[band_name]['clustering'].append(clustering)
                    band_feature_arrays[band_name]['assortativity'].append(assortativity)
                    band_feature_arrays[band_name]['modularity'].append(modularity)
        # Store per-channel
        pswe_mpf_avg.append(np.mean(pswe_mpf_list_temp) if pswe_mpf_list_temp else np.nan)
        psd_list.append(ch_psd)
        fft_list.append(ch_fft)
        mpf_list.append(ch_mpf)
        df_list.append(ch_df)
        dfv_list.append(np.std(ch_df))
        pswe_events_per_channel.append(ch_pswe)

    # Convert to arrays
    psd_arr = np.array(psd_list)
    fft_arr = np.array(fft_list)
    mpf_arr = np.array(mpf_list)
    df_arr = np.array(df_list)
    dfv_arr = np.array(dfv_list)

    # PSWE stats
    total_duration = eeg_data.shape[1] / sf
    overall_pswe_durations = []
    pswe_stats = []

    for channel in pswe_events_per_channel:
        if channel:
            sequences = find_consecutive_sequences(np.array(channel), min_duration_sec)
            durations = [len(seq)*min_duration_sec for seq in sequences]
            # Bin edges: 0, 60, 120, ..., up to total_duration
            n_bins = int(np.ceil(total_duration / 60))
            bin_edges = np.arange(0, n_bins + 1) * 60
            # Count events in each bin
            event_times = np.array(channel)
            event_counts, _ = np.histogram(event_times, bins=bin_edges)
            # Bin index is minute number
            events_per_minute = {minute: int(event_counts[minute]) for minute in range(n_bins)}
        else:
            durations = []
            events_per_minute = {}
            pswe_stats.append({
                'pswe_percentage': 0,
                'pswe_events_per_minute': 0,
                'pswe_avg_length': 0,
            })
            continue

        pswe_total = sum(durations)
        pct = (pswe_total / total_duration) * 100
        avg_len = float(np.mean(durations) if durations else 0)
        # mean events_per_minute
        events_p_minute = float(np.mean(list(events_per_minute.values())) if events_per_minute else 0)
        pswe_stats.append({
            'pswe_percentage': pct,
            'pswe_events_per_minute': events_p_minute,
            'pswe_avg_length': avg_len,
            'pswe_mpf_avg': pswe_mpf_avg
        })
        overall_pswe_durations.extend(durations)

    # Overall PSWE
    overall_pswe_median = np.median(overall_pswe_durations) if overall_pswe_durations else 0
    overall_pct = (overall_pswe_median / total_duration) * 100
    overall_evpm = len(overall_pswe_durations) / ((total_duration * len(pswe_events_per_channel)) / 60)
    overall_avg = np.mean(overall_pswe_durations) if overall_pswe_durations else 0

    # Bandpower
    def bandpower(data, sf, band, window_sec=None, relative=False):
        low, high = band
        nperseg = int((window_sec or (2/low)) * sf)
        freqs, psd = welch(data, sf, nperseg=nperseg)
        idx = np.logical_and(freqs >= low, freqs <= high)
        bp = np.trapz(psd[:, idx], freqs[idx], axis=1)
        if relative:
            bp /= np.trapz(psd, freqs, axis=1)
        return bp

    band_powers = {name: bandpower(eeg_data, sf, rng,window_sec=window_size_sec) for name, rng in power_bands.items()}

    # Time-domain stats
    skew = np.apply_along_axis(lambda x: pd.Series(x).skew(), 1, eeg_data)
    kurt = np.apply_along_axis(lambda x: pd.Series(x).kurt(), 1, eeg_data)

    # Compute connectivity matrix (coherence) using MNE spectral_connectivity
    try:
        # spectral_connectivity expects shape (n_epochs, n_signals, n_times) -> wrap our data as a single epoch
        data_for_conn = eeg_data[np.newaxis, :, :]  # shape (1, n_channels, n_times)
        conn_res = spectral_connectivity(
            data_for_conn,
            method='coh',
            sfreq=sf,
            fmin=1.0,
            fmax=40.0,
            faverage=True,
            verbose=False,
        )
        conn = conn_res[0]  # connectivity array
        # conn can be (n_channels, n_channels, n_freqs) or (n_connections, n_freqs)
        if conn.ndim == 3:
            conn_mat = np.mean(conn, axis=2)
        elif conn.ndim == 2:
            # convert upper-tri vector to symmetric matrix
            n_ch = eeg_data.shape[0]
            conn_mean = np.mean(conn, axis=1)
            conn_mat = np.zeros((n_ch, n_ch))
            triu_idx = np.triu_indices(n_ch, k=1)
            conn_mat[triu_idx] = conn_mean
            conn_mat[(triu_idx[1], triu_idx[0])] = conn_mean
        else:
            conn_mat = np.full((eeg_data.shape[0], eeg_data.shape[0]), np.nan)
        conn_df = pd.DataFrame(conn_mat, index=raw_eeg.ch_names, columns=raw_eeg.ch_names)
    except Exception as e:
        # If connectivity fails, return an empty dict and continue
        conn_df = pd.DataFrame()
        print(f"Connectivity computation failed: {e}")

    # Compile metadata
    meta = {
        'overall_pswe_median_percentage': overall_pct,
        'overall_pswe_events_per_minute': overall_evpm,
        'overall_pswe_avg_length': overall_avg,
        'overall_min': eeg_data.min(),
        'overall_max': eeg_data.max(),
        'overall_mean': eeg_data.mean(),
        'overall_median': np.median(eeg_data),
        'overall_std': eeg_data.std(),
        'overall_psd_mean': psd_arr.mean(),
        'overall_psd_std': psd_arr.std(),
        'overall_fft_mean': np.abs(fft_arr).mean(),
        'overall_fft_std': np.abs(fft_arr).std(),
        'overall_delta_power': band_powers['delta'].mean(),
        'overall_theta_power': band_powers['theta'].mean(),
        'overall_alpha_power': band_powers['alpha'].mean(),
        'overall_beta_power': band_powers['beta'].mean(),
        'overall_gamma_power': band_powers['gamma'].mean(),
        'overall_skewness': skew.mean(),
        'overall_kurtosis': kurt.mean(),
        'overall_mpf_mean': mpf_arr.mean(),
        'overall_mpf_median': np.median(mpf_arr),
        'overall_df_mean': df_arr.mean(),
        'overall_dfv_std': dfv_arr.mean(),
    }


    # Channel-specific metadata
    for i, ch in enumerate(raw_eeg.ch_names):
        meta.update({
            f'min_EEG {ch}': eeg_data[i].min(),
            f'max_EEG {ch}': eeg_data[i].max(),
            f'mean_EEG {ch}': eeg_data[i].mean(),
            f'median_EEG {ch}': np.median(eeg_data[i]),
            f'std_EEG {ch}': eeg_data[i].std(),
            f'psd_mean_EEG {ch}': psd_arr[i].mean(),
            f'psd_std_EEG {ch}': psd_arr[i].std(),
            f'fft_mean_EEG {ch}': np.abs(fft_arr[i]).mean(),
            f'fft_std_EEG {ch}': np.abs(fft_arr[i]).std(),
            f'delta_power_EEG {ch}': band_powers['delta'][i],
            f'theta_power_EEG {ch}': band_powers['theta'][i],
            f'alpha_power_EEG {ch}': band_powers['alpha'][i],
            f'beta_power_EEG {ch}': band_powers['beta'][i],
            f'gamma_power_EEG {ch}': band_powers['gamma'][i],
            f'skewness_EEG {ch}': skew[i],
            f'kurtosis_EEG {ch}': kurt[i],
            f'mean_mpf_EEG {ch}': mpf_arr[i].mean(),
            f'median_mpf_EEG {ch}': np.median(mpf_arr[i]),
            f'dfv_mean_EEG {ch}': df_arr[i].mean(),
            f'dfv_std_EEG {ch}': dfv_arr[i]
        })
        if pswe_stats:
            meta.update({
                f'pswe_percentage_EEG {ch}': pswe_stats[i]['pswe_percentage'],
                f'pswe_events_per_minute_EEG {ch}': pswe_stats[i]['pswe_events_per_minute'],
                f'pswe_avg_length_EEG {ch}': pswe_stats[i]['pswe_avg_length'],
            })
    return meta, conn_df

def process_file_2(file, pickles_location, sample_window_size):
    """
    Process a single file to extract metadata.
    """
    print(f"Processing file: {file}")
    # Extract patient ID from the file name
    patient_id = re.search(r'(\d{4})-\d{3}', file).group(1)
    print(f"Patient ID: {patient_id}")
    # Load the raw EEG data from the pickle file
    raw = pickle.load(open(os.path.join(pickles_location, file), 'rb'))
    metadata_window = eeg_data_to_features(raw, window_size_sec=sample_window_size)
    return patient_id, metadata_window

#%% TGA
def _parse_series_mixed(s: pd.Series, *, primary_dayfirst: bool) -> pd.Series:
    """Try parsing with primary_dayfirst, then flip for any NaT."""
    p1 = pd.to_datetime(s, errors="coerce", dayfirst=primary_dayfirst, infer_datetime_format=True)
    missing = p1.isna()
    if missing.any():
        p2 = pd.to_datetime(s[missing], errors="coerce", dayfirst=not primary_dayfirst, infer_datetime_format=True)
        p1.loc[missing] = p2
    return p1

def _standardize_columns(df: pd.DataFrame) -> pd.DataFrame:
    mapper = {
        "tga time onset": "onset_raw",
        "tga_time_onset": "onset_raw",
        "tga onset": "onset_raw",
        "tga time final": "final_raw",
        "tga_time_final": "final_raw",
        "tga final": "final_raw",
        "final diagnosis? 0=no tga; 1=yes; 2=maybe2": "diagnosis",
        "final diagnosis": "diagnosis",
        "eeg_date": "eeg_raw",
        "eeg time": "eeg_raw",
        "eeg": "eeg_raw",
    }
    cols_norm = {c: c.strip().lower().replace("\u200f", "").replace("\u202c", "") for c in df.columns}
    df2 = df.copy()
    df2.columns = [cols_norm[c] for c in df.columns]
    for k, v in mapper.items():
        if k in df2.columns and v not in df2.columns:
            df2.rename(columns={k: v}, inplace=True)
    # Keep original names alongside standardized ones when possible
    return df2

def _derive_phase(eeg, onset, final):
    if pd.isna(eeg) or pd.isna(onset) or pd.isna(final):
        return "unknown"
    if eeg < onset:
        return "pre-onset"
    if onset <= eeg <= final:
        return "during"
    return "post-final"

def enrich_tga_dataframe(df: pd.DataFrame) -> pd.DataFrame:
    df0 = _standardize_columns(df)

    # Parse datetimes with sensible defaults:
    # TGA times are often logged day-first (e.g., 20/9/2014),
    # EEG entries are often month-first (e.g., 9/8/14).
    onset_dt = _parse_series_mixed(df0.get("onset_raw", pd.Series([], dtype=object)), primary_dayfirst=True)
    final_dt = _parse_series_mixed(df0.get("final_raw", pd.Series([], dtype=object)), primary_dayfirst=True)
    eeg_dt   = _parse_series_mixed(df0.get("eeg_raw", pd.Series([], dtype=object)), primary_dayfirst=False)

    out = df.copy()  # preserve original columns/ordering
    out["onset_dt"] = onset_dt
    out["final_dt"] = final_dt
    out["eeg_dt"]   = eeg_dt

    # Core durations (hours, float with 3 decimals)
    tga_total = (final_dt - onset_dt).dt.total_seconds() / 3600.0
    eeg_minus_onset = (eeg_dt - onset_dt).dt.total_seconds() / 3600.0
    eeg_minus_final = (eeg_dt - final_dt).dt.total_seconds() / 3600.0

    out["tga_total_hours"] = np.round(tga_total, 3)
    out["eeg_minus_onset_hours"] = np.round(eeg_minus_onset, 3)
    out["eeg_minus_final_hours"] = np.round(eeg_minus_final, 3)

    # Valid if duration between 0 and 48 hours
    out["duration_valid"] = (tga_total >= 0) & (tga_total <= 48)

    # Days between EEG and onset date (absolute) to catch suspicious spans
    # Use .normalize() to compare by date only
    onset_date = onset_dt.dt.normalize()
    eeg_date   = eeg_dt.dt.normalize()
    span_days = (eeg_date - onset_date).dt.days.abs()
    out["suspect_span_days"] = span_days
    out["suspect_span_gt30d"] = span_days > 30

    # Missingness
    out["missing_any"] = out[["onset_dt", "final_dt", "eeg_dt"]].isna().any(axis=1)

    # Notes column for quick QA
    notes = []
    for i, row in out.iterrows():
        msgs = []
        if row.get("missing_any", False):
            msgs.append("missing datetimes")
        if not pd.isna(row.get("tga_total_hours")):
            if row["tga_total_hours"] < 0:
                msgs.append("final < onset")
            elif row["tga_total_hours"] > 48:
                msgs.append("duration > 48h (check dates)")
        if bool(row.get("suspect_span_gt30d", False)):
            msgs.append("EEG far from onset (>30d)")
        notes.append("; ".join(msgs) if msgs else "")
    out["notes"] = notes

    # Optional: normalize/clean diagnosis to int category if present
    diag_col = None
    for c in out.columns:
        if c.strip().lower().startswith("final diagnosis"):
            diag_col = c
            break
    if diag_col is not None:
        out["diagnosis_int"] = pd.to_numeric(out[diag_col], errors="coerce").astype("Int64")

    return out

def tga_get_files():
    def get_df_clinical():
        """Finds and reads the clinical Excel file from EDF_Format/TGA directory."""
        tga_dir = os.path.join('EDF_Format', 'TGA')
        excel_files = []
        
        if not os.path.isdir(tga_dir):
            raise ValueError(f"Directory {tga_dir} does not exist")
        
        for fname in os.listdir(tga_dir):
            if fname.lower().endswith(('.xlsx', '.xls')):
                excel_files.append(os.path.join(tga_dir, fname))
        
        if len(excel_files) == 0:
            raise ValueError(f"No Excel file found in {tga_dir}")
        elif len(excel_files) > 1:
            print(f"Warning: Multiple Excel files found in {tga_dir}, using first one: {excel_files[0]}")
        
        filepath = excel_files[0]
        try:
            df = pd.read_excel(filepath)
            # remove rows and col all nan
            df = df.dropna(how='all')
            df = df.dropna(axis=1, how='all')
            return df
        except Exception as e:
            raise ValueError(f"Error reading clinical file {filepath}: {e}")


    def get_tga_controls():
        """Get TGA controls from EDF_Format/TGA_control/control and merge with sex and age from Excel file."""
        # Read controls from parquet_results/TGA_control directory
        parquet_dir = os.path.join('parquet_results', 'TGA_control')
        controls_list = []
        
        if os.path.isdir(parquet_dir):
            for fname in sorted(os.listdir(parquet_dir)):
                if not fname.lower().endswith('.parquet'):
                    continue
                fpath = os.path.join(parquet_dir, fname)
                try:
                    df = pd.read_parquet(fpath)
                    # Ensure there's a file_name column
                    if 'file_name' not in df.columns:
                        base = fname.replace('.parquet', '').replace('.edf', '')
                        df['file_name'] = base
                    controls_list.append(df)
                except Exception as e:
                    print(f"Warning: failed reading {fpath}: {e}")
                    continue
        
        if len(controls_list) == 0:
            # If no parquet files found, try reading from TGA_controls.parquet in root directory
            root_parquet_path = 'TGA_control.parquet'
            if os.path.exists(root_parquet_path):
                try:
                    controls = pd.read_parquet(root_parquet_path)
                    # Ensure there's a file_name column
                    if 'file_name' not in controls.columns:
                        # Try to extract from index or create a default
                        controls['file_name'] = controls.index.astype(str)
                except Exception as e:
                    print(f"Warning: failed reading {root_parquet_path}: {e}")
                    controls = None
        else:
            controls = pd.concat(controls_list, ignore_index=True)
        
        # Process controls similar to get_cobrad_controls
        # Ensure file_name exists for ID extraction
        if 'file_name' not in controls.columns:
            raise ValueError("Controls dataframe must have 'file_name' column")
        
        controls['ID'] = controls['file_name'].apply(lambda x: x.split('_')[0] if '_' in str(x) else str(x).replace('.edf', '').replace('.EEG', '')).astype(str)
        numeric_cols = controls.select_dtypes(include=[np.number]).columns
        
        # Keep file_name for merging later - use first file_name per ID
        if 'duration_min' in controls.columns:
            # Group by ID and aggregate, keeping file_name
            def aggregate_group(x):
                result = {}
                for col in numeric_cols:
                    if col != 'duration_min':
                        result[col] = (x[col].multiply(x['duration_min'], axis=0)).sum(skipna=False) / x['duration_min'].sum(skipna=False)
                result['file_name'] = x['file_name'].iloc[0]
                return pd.Series(result)
            
            controls = controls.groupby('ID').apply(aggregate_group).reset_index()
        else:
            # If no duration_min, just take mean but keep file_name
            agg_dict = {col: 'mean' for col in numeric_cols}
            agg_dict['file_name'] = 'first'
            controls = controls.groupby('ID').agg(agg_dict).reset_index()
        
        # Read Excel file to get sex and age
        excel_path = 'EDF_Format/TGA_control/control/control_10JUL25.xlsx'
        if os.path.exists(excel_path):
            df_demo = pd.read_excel(excel_path)
            # Convert gender column (Hebrew: נקבה = female, זכר = male)
            df_demo['gender'] = df_demo['gender'].astype(str).str.strip()
            df_demo['sex'] = df_demo['gender'].apply(lambda x: 'f' if 'נקבה' in x or 'female' in x.lower() or x == '2' else ('m' if 'זכר' in x or 'male' in x.lower() or x == '1' else np.nan))
            df_demo = convert_sex_column(df_demo, 'sex')
            
            # Prepare for merge - use EEG FILE NAME to match with controls
            df_demo['file_name_clean'] = df_demo['EEG FILE NAME'].astype(str).str.strip().str.upper()
            # Remove .edf or .EEG extensions if present
            df_demo['file_name_clean'] = df_demo['file_name_clean'].str.replace('.EDF', '', regex=False).str.replace('.EEG', '', regex=False)
            # Create ID from file_name (similar to controls) - use the file_name itself as ID if no underscore
            df_demo['ID'] = df_demo['file_name_clean'].apply(lambda x: x.split('_')[0] if '_' in str(x) else str(x)).astype(str)
            
            # Clean controls file_name for matching
            if 'file_name' in controls.columns:
                controls['file_name_clean'] = controls['file_name'].astype(str).str.strip().str.upper().str.replace('.EDF', '', regex=False).str.replace('.EEG', '', regex=False)
                # Merge on file_name_clean first (most direct match)
                controls = controls.merge(
                    df_demo[['file_name_clean', 'sex', 'age']].drop_duplicates(subset='file_name_clean'),
                    on='file_name_clean',
                    how='left'
                )
                # If some rows don't have sex/age, try merging on ID as fallback
                missing_mask = controls['sex'].isna() | controls['age'].isna()
                if missing_mask.any():
                    # Use existing ID column for controls that are missing data
                    controls_missing = controls.loc[missing_mask, ['ID']].copy()
                    # Merge on ID
                    controls_missing = controls_missing.merge(
                        df_demo[['ID', 'sex', 'age']].drop_duplicates(subset='ID'),
                        on='ID',
                        how='left'
                    )
                    # Update missing values
                    controls.loc[missing_mask, 'sex'] = controls_missing['sex'].values
                    controls.loc[missing_mask, 'age'] = controls_missing['age'].values
                # Remove temporary column
                controls = controls.drop(columns=['file_name_clean'], errors='ignore')
            else:
                # If no file_name in controls, just merge on ID
                controls = controls.merge(
                    df_demo[['ID', 'sex', 'age']].drop_duplicates(subset='ID'),
                    on='ID',
                    how='left'
                )
        else:
            print(f"Warning: Excel file not found at {excel_path}, sex and age columns will not be added")
        # rename sex to gender
        controls = controls.rename(columns={'sex': 'gender'})
        return controls

    # Example usage inside tga_get_files
    clinical_df = get_df_clinical()
    clinical_df = convert_sex_column(clinical_df,'Gender_Text')
    # read TGA.oarquet
    df_wnv = pd.read_parquet('TGA.parquet')
    # get clinical features
    clinical_df = enrich_tga_dataframe(clinical_df)
    # remove from file_name .edf
    df_wnv['file_name'] = df_wnv['file_name'].str.replace('.edf', '', regex=False)
    patients_folder = "TGA"
    controls = get_tga_controls()
    # remove controls that are missing age or gender
    controls = controls[controls['age'].notna() & controls['gender'].notna()]
    # df_wnv2 is merged df_wnv and clinical_df right on EEG FILE NAME, left on PATIENT_ID
    df_wnv2 = df_wnv.merge(clinical_df, left_on='file_name', right_on='EEG FILE NAME', how='inner')
    cases_group_name = "TGA"
    return df_wnv, patients_folder, controls, df_wnv2, cases_group_name

from pathlib import Path  # noqa: E402
def get_clinical_data(project_name):
    # Attempt to find a clinical Excel file under EDF_Format/{project_name}
    clinical_df = pd.DataFrame()
    edf_proj_dir = os.path.join('EDF_Format', project_name)
    if os.path.isdir(edf_proj_dir):
        # look for any .xls/.xlsx files in the directory and subdirs
        excel_files = []
        for root, dirs, files in os.walk(edf_proj_dir):
            for f in files:
                if f.lower().endswith(('.xls', '.xlsx')):
                    excel_files.append(os.path.join(root, f))
        if excel_files:
            clinical_path = excel_files[0]
            try:
                clinical_df = pd.read_excel(clinical_path)
            except Exception as e:
                raise ValueError(f"Failed reading clinical file {clinical_path}: {e}")
                print(f"Warning: failed reading clinical file {clinical_path}: {e}")
    else:
        print(f"Warning: EDF_Format project directory not found: {edf_proj_dir}")

    # Normalize clinical df: prefer column 'record_id' -> 'ID'
    if not clinical_df.empty:
        if 'record_id' in clinical_df.columns:
            clinical_df = clinical_df.rename(columns={'record_id': 'ID'})
        # Ensure ID is string-like and strip whitespace / extensions
        if 'ID' in clinical_df.columns:
            clinical_df['ID'] = clinical_df['ID'].astype(str).str.strip().str.replace('\.edf$', '', regex=True).str.lower()
    # find the col of sex or gender ignore case
    gender_col = next((col for col in clinical_df.columns if col.lower() in ['sex', 'gender']), None)
    if gender_col:
        # clnical df column gender or sex ignore case to  convert_sex_column
        clinical_df = convert_sex_column(clinical_df, gender_col)
    return clinical_df

def generic_get_files(project_name: str):
    """
    Generic loader for projects that store per-record parquet files under
    `parquet_results/{project_name}` and a clinical Excel table under
    `EDF_Format/{project_name}`. The clinical table is expected to contain
    a column named 'record_id' (per the user's description) which will be
    matched to the EDF/parquet file names (without the .edf extension).

    Returns the same 5-tuple as the specialized loaders used in the app:
      (df_parquet_concat, patients_folder, controls_df, df_wnv2_grouped, cases_group_name)

    If no clinical file or parquet files are found, the function will
    return empty DataFrames where applicable but will not raise.
    """
    patients_folder = project_name
    parquet_dir = os.path.join('parquet_results', project_name)
    df_list = []
    # Read parquet files
    if os.path.isdir(parquet_dir):
        for fname in sorted(os.listdir(parquet_dir)):
            if not fname.lower().endswith('.parquet'):
                continue
            fpath = os.path.join(parquet_dir, fname)
            try:
                df = pd.read_parquet(fpath)
            except Exception as e:
                # If reading fails, try to continue
                print(f"Warning: failed reading {fpath}: {e}")
                continue
            # Ensure there's a file_name column - use the parquet filename as fallback
            if 'file_name' not in df.columns:
                base = fname.replace('.parquet', '')
                # remove trailing .edf if present
                if base.lower().endswith('.edf'):
                    base = base[:-4]
                df['file_name'] = base
            filebase = fname.split('.edf')[0]
            # find file path in 'EDF_Format' folder and subfolders with glob
            pattern = os.path.join('EDF_Format', patients_folder, '**', filebase+'*')
            matches = glob.glob(pattern, recursive=True)
            if matches:
                filepath = matches[0]  # first match
            else:
                print(f"Warning: file not found: {pattern}")
                continue
            # Find the least common (most specific) folder in the path
            path_parts = Path(filepath).parts
            # Remove 'EDF_Format' and the project folder to get the relevant path
            relevant_parts = path_parts[path_parts.index('EDF_Format')+2:]  # Skip 'EDF_Format' and project folder
            if relevant_parts:
                # The least common folder is the first part after removing EDF_Format and project folder
                df['mother_folder'] = relevant_parts[0]
            df_list.append(df)
    else:
        raise ValueError(f"No parquet files found for {project_name}")

    if len(df_list) == 0:
        raise ValueError(f"No parquet files found for {project_name}")
    else:
        df_parquet = pd.concat(df_list, ignore_index=True, sort=False)
    # df_parquet ID = file_name.split'.'[0]
    df_parquet['ID'] = df_parquet['file_name'].astype(str).apply(lambda x: os.path.basename(x).replace('.edf', '').lower().strip())


    clinical_df = get_clinical_data(project_name)
    # Determine which column in df_parquet best matches clinical_df['ID']
    def find_best_matching_column(clinical_df, df_parquet):
        """Determine which column in df_parquet best matches clinical_df['ID'] values"""
        if clinical_df.empty or df_parquet.empty or 'ID' not in clinical_df.columns:
            return None
        
        clinical_ids = clinical_df['ID'].astype(str).str.lower().str.strip()
        best_column = None
        best_match_count = 0
        
        # Check 'ID' column in df_parquet
        if 'ID' in df_parquet.columns:
            parquet_ids = df_parquet['ID'].astype(str).str.lower().str.strip()
            matches = 0
            for clinical_id in clinical_ids:
                for parquet_id in parquet_ids:
                    if str(clinical_id) in str(parquet_id) or str(parquet_id) in str(clinical_id):
                        matches += 1
                        break
            if matches > best_match_count:
                best_match_count = matches
                best_column = 'ID'
        
        # Check 'mother_folder' column in df_parquet
        if 'mother_folder' in df_parquet.columns:
            parquet_folders = df_parquet['mother_folder'].astype(str).str.lower().str.strip()
            matches = 0
            for clinical_id in clinical_ids:
                for parquet_folder in parquet_folders:
                    if str(clinical_id) in str(parquet_folder) or str(parquet_folder) in str(clinical_id):
                        matches += 1
                        break
            if matches > best_match_count:
                best_match_count = matches
                best_column = 'mother_folder'
        
        print(f"Best matching column: {best_column} with {best_match_count} matches out of {len(clinical_ids)} clinical IDs")
        return best_column
    
    # Find the best matching column and create ID column if needed
    best_match_col = find_best_matching_column(clinical_df, df_parquet)
    if best_match_col and best_match_col != 'ID':
        # Create ID column from the best matching column
        if best_match_col == 'mother_folder':
            df_parquet['ID'] = df_parquet['mother_folder']
        elif best_match_col == 'file_name':
            df_parquet['ID'] = df_parquet['file_name'].astype(str).apply(lambda x: os.path.basename(x).replace('.edf', '').lower().strip())
    
    # Merge clinical and parquet data
    df_wnv2 = pd.DataFrame()
    if (not clinical_df.empty) and (not df_parquet.empty) and ('ID' in clinical_df.columns) and ('ID' in df_parquet.columns):
        try:
            # Create a mapping where df_parquet ID contains clinical_df ID
            def find_matching_id(clinical_id):
                """Find df_parquet ID that contains the clinical_id"""
                for parquet_id in df_parquet['ID']:
                    if str(clinical_id).lower() in str(parquet_id).lower():
                        return parquet_id
                return None
            
            # Add matching parquet ID to clinical_df
            clinical_df['matched_parquet_id'] = clinical_df['ID'].apply(find_matching_id)
            
            # Filter out rows where no match was found
            clinical_df_matched = clinical_df[clinical_df['matched_parquet_id'].notna()].copy()
            
            if len(clinical_df_matched) == 0:
                print("Warning: No matching IDs found between clinical and parquet data")
                df_merged = pd.DataFrame()
            else:
                # Merge using the matched parquet ID
                df_merged = pd.merge(clinical_df_matched, df_parquet, 
                                   left_on='matched_parquet_id', right_on='ID', 
                                   how='inner', suffixes=('_clinical', '_parquet'))
                
                # Merge ID columns into one ID column
                if 'ID_clinical' in df_merged.columns and 'ID_parquet' in df_merged.columns:
                    # Use parquet ID as primary, fallback to clinical ID
                    df_merged['ID'] = df_merged['ID_parquet'].fillna(df_merged['ID_clinical'])
                    # Drop the separate ID columns
                    df_merged = df_merged.drop(columns=['ID_clinical', 'ID_parquet'], errors='ignore')
                elif 'ID_clinical' in df_merged.columns:
                    df_merged['ID'] = df_merged['ID_clinical']
                    df_merged = df_merged.drop(columns=['ID_clinical'], errors='ignore')
                elif 'ID_parquet' in df_merged.columns:
                    df_merged['ID'] = df_merged['ID_parquet']
                    df_merged = df_merged.drop(columns=['ID_parquet'], errors='ignore')
        except Exception as e:
            print(f"Warning: merge failed: {e}")
            df_merged = pd.DataFrame()
    else:
        # If we couldn't merge (missing pieces), try a looser merge by checking substring matches
        df_merged = pd.DataFrame()
        if (not clinical_df.empty) and (not df_parquet.empty):
            # Try matching clinical 'ID' values as substrings of df_parquet['file_name']
            if 'file_name' in df_parquet.columns and 'ID' in clinical_df.columns:
                # Build a reverse index for quick lookup
                mapping = {}
                for idx, row in df_parquet.iterrows():
                    fname = str(row['file_name'])
                    key = os.path.basename(fname).replace('.edf', '')
                    mapping.setdefault(key, []).append(idx)
                matched_rows = []
                for _, c in clinical_df.iterrows():
                    cid = str(c.get('ID', ''))
                    if cid in mapping:
                        matched_rows.extend(mapping[cid])
                if matched_rows:
                    df_merged = pd.concat([clinical_df.reset_index(drop=True), df_parquet.loc[matched_rows].reset_index(drop=True)], axis=1, ignore_index=False)

    if project_name == "reading_epilepsy_cut":
        pre_df, post_df, pairing_info = load_pre_post_data(project_name)
        df_wnv2 = pd.concat([post_df], ignore_index=True)
    # Controls: generic loader doesn't have controls; return empty DataFrame
    controls = get_cobrad_controls()
    cases_group_name = project_name
    # assess if df_merged is empty
    if df_merged.empty:
        return df_parquet, patients_folder, controls, df_wnv2, cases_group_name
        raise ValueError("No matching IDs found between clinical and parquet data")
    df_wnv2 = aggregate_dataframe(df_merged, weighted_avg)
    # Create grouped/aggregated df (df_wnv2) similar to other loaders
    if not df_wnv2.empty:
        # Drop helper columns if present
        if '_merge' in df_wnv2.columns:
            df_wnv2 = df_wnv2.drop(columns=['_merge'], errors='ignore')
        # Clean up common unwanted columns and coerce numeric strings
        cols_to_drop = ['highpass', 'lowpass', 'n_samples', 'size', 'patient_number', 'duration_sec']
        df_wnv2 = df_wnv2.drop(columns=[col for col in df_wnv2.columns if any(drop in col for drop in cols_to_drop)], errors='ignore')
        df_wnv2 = df_wnv2.applymap(lambda x: np.nan if isinstance(x, str) and 'nan' in x.lower() else x)
        for col in df_wnv2.columns:
            try:
                df_wnv2[col] = pd.to_numeric(df_wnv2[col])
            except Exception:
                pass
    return df_parquet, patients_folder, controls, df_wnv2, cases_group_name
    df_merged.columns.unique().tolist()

def get_patient_groups(dataset="reading_epilepsy_cut"):
    """Get patient group assignments from EDF_Format folder structure."""
    patient_groups = {}
    
    if dataset == "reading_epilepsy_cut":
        edf_dir = "EDF_Format/reading_epilepsy_cut"
        if os.path.exists(edf_dir):
            # Walk through patient folders
            for root, dirs, files in os.walk(edf_dir):
                # Get patient group name from folder
                if root != edf_dir:  # Skip the root directory
                    patient_group = os.path.basename(root)
                    
                    # Process EDF files in this patient group
                    for file in files:
                        if file.endswith('.edf'):
                            # Extract patient ID from filename (before first underscore)
                            patient_id = file.split('_')[0]
                            patient_groups[patient_id] = patient_group
    
    return patient_groups

def load_pre_post_data(dataset="VNS_PRE_POST_25", analysis_level="Per Patient (averaged)"):
    """Load pre and post data from parquet files and determine group membership from EDF directory structure."""
    
    # Load all parquet files
    parquet_dir = f"parquet_results/{dataset}"
    parquet_files = glob.glob(os.path.join(parquet_dir, "*.parquet"))
    
    pre_data = []
    post_data = []
    
    if dataset == "reading_epilepsy_cut":
        # Get patient group assignments from EDF_Format folder structure
        patient_groups = get_patient_groups(dataset)
        print(f"Patient groups: {patient_groups}")
        
        # Special handling for reading_epilepsy_cut dataset
        for parquet_file in parquet_files:
            try:
                # Extract filename without extension
                filename = os.path.basename(parquet_file).replace('.edf.parquet', '')
                
                # Load parquet data
                df = pd.read_parquet(parquet_file)
                
                # Determine group and subject from filename
                # Format: BA0012UZ_reading.edf.parquet or BA0012UZ_non_reading.edf.parquet
                if '_non_reading' in filename:
                    group = "PRE"   # non_reading = pre (before reading)
                    subject_id = filename.replace('_non_reading', '')
                elif '_reading' in filename:
                    group = "POST"  # reading = post (after reading)
                    subject_id = filename.replace('_reading', '')
                else:
                    continue  # Skip files that don't match the pattern
                
                # Extract patient ID and get patient group
                patient_id = subject_id.split('_')[0]
                patient_group = patient_groups.get(patient_id, 'Unknown')
                
                # Add group label and patient information
                df['Group'] = group
                df['Subject_ID'] = subject_id
                df['Patient_ID'] = patient_id
                df['Patient_Group'] = patient_group
                
                if group == "POST":
                    post_data.append(df)
                else:
                    pre_data.append(df)
                    
            except Exception as e:
                st.warning(f"Could not load {parquet_file}: {e}")
                continue
    else:
        # Original logic for VNS_PRE_POST_25 dataset
        file_map = {}
        # EDF directory structure for determining pre/post
        edf_base_dir = f"EDF_Format/{dataset}"
        
        for parquet_file in parquet_files:
            try:
                # Extract filename without extension
                filename = os.path.basename(parquet_file).replace('.edf.parquet', '')
                
                # Load parquet data
                df = pd.read_parquet(parquet_file)
                
                # Determine if this is pre or post based on EDF directory structure
                # Build a mapping from filename to (group, Subject_ID) for all files in edf_base_dir
                    
                for root, dirs, files in os.walk(edf_base_dir):
                    for file in files:
                        if filename in file:
                            key = file.split('.')[0]  # Remove extension for matching
                            group = root.split(os.sep)[-1].upper()
                            subject_id = root.split(os.sep)[-2]
                            file_map[key] = (group, subject_id)
                            print(f"Mapping file {key} to group {group}, subject {subject_id}. Filename: {filename}, file: {file}")
                            break
                    # if key in file_map:
                    if filename in file_map:
                        break
                load_pre_post_data._file_map = file_map

                # Try to match filename in the file_map
                group_subject = file_map.get(filename, (False, None))
                is_post, Subject_ID = group_subject
                
                # Add group label
                df['Group'] = is_post
                df['Subject_ID'] = Subject_ID
                
                if is_post == "POST":
                    post_data.append(df)
                else:
                    pre_data.append(df)
                # print how many in pre and post
            except Exception as e:
                st.warning(f"Could not load {parquet_file}: {e}")
                continue
    print(f"Loaded {len(pre_data)} PRE and {len(post_data)} POST files so far.")
    if not pre_data and not post_data:
        st.error("No data found!")
        return None, None, None
    
    pre_df = pd.concat(pre_data, ignore_index=True) if pre_data else pd.DataFrame()
    post_df = pd.concat(post_data, ignore_index=True) if post_data else pd.DataFrame()
    
    # Apply averaging based on analysis level
    if analysis_level == "Per Patient (averaged)" and dataset == "reading_epilepsy_cut":
        # For per-patient analysis, average by Patient_Group instead of Subject_ID
        if len(pre_df) > 0:
            numeric_cols = pre_df.select_dtypes(include=np.number).columns.tolist()
            non_numeric_cols = [col for col in pre_df.columns if col not in numeric_cols and col not in ['Subject_ID', 'Patient_ID', 'Patient_Group']]
            pre_df = pre_df.groupby('Patient_Group').agg({**{col: 'mean' for col in numeric_cols}, **{col: 'first' for col in non_numeric_cols}}).reset_index()
            # Use Patient_Group as Subject_ID for consistency
            pre_df['Subject_ID'] = pre_df['Patient_Group']
            
        if len(post_df) > 0:
            numeric_cols = post_df.select_dtypes(include=np.number).columns.tolist()
            non_numeric_cols = [col for col in post_df.columns if col not in numeric_cols and col not in ['Subject_ID', 'Patient_ID', 'Patient_Group']]
            post_df = post_df.groupby('Patient_Group').agg({**{col: 'mean' for col in numeric_cols}, **{col: 'first' for col in non_numeric_cols}}).reset_index()
            # Use Patient_Group as Subject_ID for consistency
            post_df['Subject_ID'] = post_df['Patient_Group']
    elif analysis_level == "All Samples (mixed paired/non-paired)":
        # For mixed analysis, average by Subject_ID but keep all samples
        # Multiple scans per patient will be averaged per Subject_ID
        if len(pre_df) > 0:
            numeric_cols = pre_df.select_dtypes(include=np.number).columns.tolist()
            non_numeric_cols = [col for col in pre_df.columns if col not in numeric_cols and col != 'Subject_ID']
            pre_df = pre_df.groupby('Subject_ID').agg({**{col: 'mean' for col in numeric_cols}, **{col: 'first' for col in non_numeric_cols}}).reset_index()
        if len(post_df) > 0:
            numeric_cols = post_df.select_dtypes(include=np.number).columns.tolist()
            non_numeric_cols = [col for col in post_df.columns if col not in numeric_cols and col != 'Subject_ID']
            post_df = post_df.groupby('Subject_ID').agg({**{col: 'mean' for col in numeric_cols}, **{col: 'first' for col in non_numeric_cols}}).reset_index()
    else:
        # Original logic: mean over Subject_ID for numeric columns, first for non-numeric columns
        if len(pre_df) > 0:
            numeric_cols = pre_df.select_dtypes(include=np.number).columns.tolist()
            non_numeric_cols = [col for col in pre_df.columns if col not in numeric_cols and col != 'Subject_ID']
            pre_df = pre_df.groupby('Subject_ID').agg({**{col: 'mean' for col in numeric_cols}, **{col: 'first' for col in non_numeric_cols}}).reset_index()
        if len(post_df) > 0:
            numeric_cols = post_df.select_dtypes(include=np.number).columns.tolist()
            non_numeric_cols = [col for col in post_df.columns if col not in numeric_cols and col != 'Subject_ID']
            post_df = post_df.groupby('Subject_ID').agg({**{col: 'mean' for col in numeric_cols}, **{col: 'first' for col in non_numeric_cols}}).reset_index()
    # Find patients with both PRE and POST data
    pre_subjects = set(pre_df['Subject_ID'].unique()) if len(pre_df) > 0 else set()
    post_subjects = set(post_df['Subject_ID'].unique()) if len(post_df) > 0 else set()
    paired_subjects = pre_subjects.intersection(post_subjects)
    unpaired_pre_subjects = pre_subjects - post_subjects
    unpaired_post_subjects = post_subjects - pre_subjects
    
    # Handle different analysis levels
    if analysis_level == "All Samples (mixed paired/non-paired)":
        # Return all data: paired and unpaired
        # Sort by Subject_ID for consistency
        pre_df = pre_df.sort_values('Subject_ID').reset_index(drop=True)
        post_df = post_df.sort_values('Subject_ID').reset_index(drop=True)
        
        # Create a combined info dict to track paired/unpaired status
        pairing_info = {
            'paired': paired_subjects,
            'unpaired_pre': unpaired_pre_subjects,
            'unpaired_post': unpaired_post_subjects,
            'total_pre': len(pre_df),
            'total_post': len(post_df),
            'paired_count': len(paired_subjects),
            'unpaired_pre_count': len(unpaired_pre_subjects),
            'unpaired_post_count': len(unpaired_post_subjects)
        }
        
        print(f'Mixed analysis - Paired: {len(paired_subjects)}, Unpaired PRE: {len(unpaired_pre_subjects)}, Unpaired POST: {len(unpaired_post_subjects)}')
        return pre_df, post_df, pairing_info
    else:
        # Original behavior: filter to only include paired subjects
        if len(paired_subjects) > 0:
            pre_paired = pre_df[pre_df['Subject_ID'].isin(paired_subjects)].copy()
            post_paired = post_df[post_df['Subject_ID'].isin(paired_subjects)].copy()
            
            # Sort by Subject_ID to ensure proper pairing
            pre_paired = pre_paired.sort_values('Subject_ID').reset_index(drop=True)
            post_paired = post_paired.sort_values('Subject_ID').reset_index(drop=True)
            print(f'len pre_paired: {len(pre_paired)}, len post_paired: {len(post_paired)}')
            return pre_paired, post_paired, paired_subjects
        else:
            st.warning("No patients found with both PRE and POST data!")
            return None, None, None

def aggregate_dataframe(df_merged, weighted_avg):
    """
    Aggregate a dataframe by mother_folder (if present) and ID.
    Uses weighted average if 'duration_min' exists, otherwise falls back to mean.
    Non-numeric columns preserve the first value.

    Parameters
    ----------
    df_merged : pd.DataFrame
        Input dataframe.
    weighted_avg : function
        Function signature must be:
        weighted_avg(group, weight_col, numeric_cols)

    Returns
    -------
    df_merged : pd.DataFrame
        Aggregated dataframe (mother_folder level if exists).
    df_wnv2 : pd.DataFrame
        Aggregated dataframe (ID level).
    """
    # --- Handle mother_folder level aggregation ---
    if 'mother_folder' in df_merged.columns and df_merged['mother_folder'].nunique() > 1:
        count_mother = df_merged['mother_folder'].value_counts()

        if 'duration_min' in df_merged.columns:
            numeric_cols = df_merged.select_dtypes(include=[np.number]).columns
            try:
                df_merged = (
                    df_merged.groupby('mother_folder')
                    .apply(weighted_avg, weight_col='duration_min', numeric_cols=numeric_cols)
                    .reset_index(drop=True)
                )
            except Exception:
                print("Warning: weighted average failed, falling back to simple mean")
                df_merged = (
                    df_merged.groupby('mother_folder')
                    .mean(numeric_only=True)
                    .reset_index()
                )
        else:
            print("Warning: 'duration_min' column not found, falling back to simple mean")
            df_merged = (
                df_merged.groupby('mother_folder')
                .agg(lambda x: x.mean() if pd.api.types.is_numeric_dtype(x) else x.iloc[0])
                .reset_index()
            )
    # --- Handle ID level aggregation ---
    elif not df_merged.empty:
        if '_merge' in df_merged.columns:
            df_merged = df_merged.drop(columns=['_merge'], errors='ignore')

        if 'duration_min' in df_merged.columns:
            numeric_cols = df_merged.select_dtypes(include=[np.number]).columns
            try:
                (
                    df_merged.groupby('ID')
                    .apply(weighted_avg, weight_col='duration_min', numeric_cols=numeric_cols)
                    .reset_index(drop=True)
                )
            except Exception:
                print("Warning: weighted average failed at ID level, falling back to simple mean")
                (
                    df_merged.groupby('ID')
                    .mean(numeric_only=True)
                    .reset_index()
                )
        else:
            (
                df_merged.groupby('ID')
                .mean(numeric_only=True)
                .reset_index()
            )
    # Add counts
    df_merged = df_merged.merge(
        count_mother.rename('num_records'),
        left_on='mother_folder',
        right_index=True,
        how='left'
    )
    return df_merged
#%% COBRAD
def get_cobrad_controls(patients_folder='EDF'):
    # read Controls_controls.parquet
    controls = pd.read_parquet('Controls_controls.parquet')
    controls['ID'] = controls['file_name'].apply(lambda x: x.split('_')[0]).astype(str)
    numeric_cols = controls.select_dtypes(include=[np.number]).columns
    controls = controls.groupby('ID').apply(
        lambda x: (x[numeric_cols].multiply(x['duration_min'], axis=0)).sum(skipna=False) / x['duration_min'].sum(skipna=False)
    ).reset_index()

    controls_dir = 'parquet_results/Controls'
    age_gender_df = get_controls_ages_genders(controls_dir)
    # Merge on ID
    controls = controls.merge(age_gender_df[['ID', 'Age', 'Gender']], on='ID', how='left')
    # Rename Age and Gender to lowercase
    controls = controls.rename(columns={'Age': 'age', 'Gender': 'gender'})
    return controls

def cobrad_get_files(sample_window_size=0,only_awake=False,sleep_only=False):
    patients_folder = "EDF"
    sheets_to_read = ['clinical', 'medications', 'npi-q', 'epworth', 'isi', 'ecog_12','Sheet4','seizures']
    dfs = pd.read_excel('COBRAD_clinical_24022025.xlsx', sheet_name=sheets_to_read)
    def get_df_wnv():
        # read sheets clinical, medications, npi-q, epworth,isi, ecpg_12 from COBRAD_clinical_24022025.xlsx
        sheets_to_sum_vals = ['epworth', 'isi', 'ecog_12','Sheet4','npi-q','seizures', 'medications']
        # Rename 'record_id' to 'ID' in each DataFrame and convert to string
        for sheet in sheets_to_read:
            dfs[sheet] = dfs[sheet].rename(columns={'record_id': 'ID'}).astype(str)
            # drop col contain has_eeg or has eeg . ignore case
            dfs[sheet] = dfs[sheet].drop(columns=[col for col in dfs[sheet].columns if 'has eeg' in col.lower() or 'has_eeg' in col.lower()])
            # One-hot encode drugs in the medications sheet
            if sheet == 'medications':
                # Extract the drug names no 'nan'
                drug_names = dfs[sheet]['name_drug_1'].dropna().unique()
                # remove 'nan'
                drug_names = [drug for drug in drug_names if 'nan' not in drug]
                # Create one-hot encoded columns for each drug
                for drug in drug_names:
                    dfs[sheet][f'{drug}'] = dfs[sheet]['name_drug_1'].apply(lambda x: 1 if x == drug else 0)
                # Keep only the ID and one-hot encoded drug columns
                drug_columns = [f'{drug}' for drug in drug_names]
                dfs[sheet] = dfs[sheet][['ID'] + drug_columns]  
                # merge columns per ID
                dfs[sheet] = dfs[sheet].groupby('ID').sum().reset_index()   
                # dict how many ID take a drug
                drug_counts = dfs[sheet].drop(columns='ID').sum().to_dict() 
                # only above 3
                drug_counts = {k: v for k, v in drug_counts.items() if v > 3}  
                # read utils/drug_groups.json
                with open('utils/drug_groups.json', 'r') as f:
                    drug_groups = json.load(f)
                for key_number_groups in drug_groups:
                    for key_group, value_group in drug_groups[key_number_groups].items():
                        # value_group remove /n and stip
                        value_group = [group.strip() for group in value_group]
                        # get the columns that contain the key_group
                        columns = [col for col in dfs[sheet].columns if any(group in col for group in value_group)]
                        # sum the columns and create a new column with the name of the group
                        dfs[sheet][f'{key_number_groups}_groups_{key_group}'] = dfs[sheet][columns].sum(axis=1)
            if sheet in sheets_to_sum_vals:
                # to numeric all columns but ID
                dfs[sheet] = pd.concat([dfs[sheet]['ID'], dfs[sheet].drop(columns='ID').apply(pd.to_numeric, errors='coerce')], axis=1)
                dfs[sheet][f'{sheet}_sum'] = dfs[sheet].drop(columns='ID').sum(axis=1)
                if sheet =='seizures':
                    # replace -9 with np.nan
                    dfs[sheet] = dfs[sheet].replace(-9, np.nan)
            # add the name of sheet at beggining of column all cols but ID
            dfs[sheet].columns = [f'{sheet}_{col}' if col != 'ID' else col for col in dfs[sheet].columns]
        # Merge all DataFrames on 'ID'
        df_wnv = dfs[sheets_to_read[0]]
        for sheet in sheets_to_read[1:]:
            df_wnv = pd.merge(df_wnv, dfs[sheet], on='ID', how='outer')
        # Rename any column that contains 'clinical_sex' (case-insensitive) to 'sex'
        for col in df_wnv.columns:
            if isinstance(col, str) and 'clinical_sex' in col.lower():
                # If 'sex' already exists, prefer the existing 'sex' column; otherwise rename
                if 'sex' not in df_wnv.columns:
                    df_wnv = df_wnv.rename(columns={col: 'sex'})
                else:
                    # If both exist, keep both but standardize name by creating 'sex_orig' for clarity
                    df_wnv = df_wnv.rename(columns={col: f"{col}_orig"})
                break
        return df_wnv

    controls = get_cobrad_controls(patients_folder)

    def get_cases_cobrad():
        if sample_window_size == 0:
            if only_awake:
                case_file = f"{patients_folder}_awake.csv"
            else:
                case_file = f"{patients_folder}.csv"
            cases = pd.read_csv(case_file)
            cases['ID'] = cases['csv_file_name'].apply(id_from_csv_filename).astype(str)
            # Sort by ID
            cases = cases.sort_values(by='ID')
        else:
            # Process pickles to calculate average values for each patient
            case_folder = 'wake_EDF' if only_awake else 'EDF'
            pickles_location = f'pickles/{case_folder}'
            files = [file for file in os.listdir(pickles_location) if file.endswith('.pkl')]

            # Dictionary to store metadata for each patient
            patient_metadata = {}
            # Use ProcessPoolExecutor for parallel processing
            with ProcessPoolExecutor() as executor:
                list(tqdm(executor.map(process_file_2, files, [pickles_location] * len(files), [sample_window_size] * len(files)), 
                                    total=len(files), desc="Processing files"))

            # Aggregate metadata for each patient
            for patient_id, metadata_list in patient_metadata.items():
                # Convert list of metadata dictionaries to a DataFrame
                patient_df = pd.DataFrame(metadata_list)
                # Calculate mean for numeric columns
                patient_avg = patient_df.mean(numeric_only=True).to_dict()
                patient_avg['ID'] = patient_id
                cases.append(patient_avg)

            # Convert the list of dictionaries to a DataFrame
            cases = pd.DataFrame(cases)

        return cases
    cases = get_cases_cobrad()
    df_wnv = get_df_wnv()
    df_wnv.columns = df_wnv.columns.str.replace('.', '_').str.replace('<', '').str.replace('>', '')

    df_merged = pd.merge(df_wnv, cases, on='ID', how='outer',indicator=True)
    # Get all files that end with .edf from EDF folder and subfolders
    eeg_files = []
    for root, dirs, files in os.walk('EDF'):
        for file in files:
            if file.endswith('.edf'):
                eeg_files.append(os.path.join(root, file))
    # print outer
    failed_ids = df_merged[df_merged['_merge'] == 'left_only']['ID'].unique()
    # check if they exist in eeg_files
    for id_to_check in failed_ids:
        # check if id exists contains in eeg_files
        if any(id_to_check in file for file in eeg_files):
            # print(f'{id_to_check}')
            pass
    df_merged = df_merged[df_merged['_merge'] == 'both'].drop(columns='_merge')
    numeric_cols = df_merged.select_dtypes(include=[np.number]).columns
    # Group by ID and apply the weighted average function
    df_wnv2 = df_merged.groupby('ID').apply(weighted_avg, weight_col='duration_min', numeric_cols=numeric_cols).reset_index(drop=True)
    # remove highpass, lowpass, n_samples, size, patient_number
    cols_to_drop = ['highpass', 'lowpass', 'n_samples', 'size', 'patient_number','duration_sec']
    # drop if conatins any of the cols_to_drop
    df_wnv2 = df_wnv2.drop(columns=[col for col in df_wnv2.columns if any([drop in col for drop in cols_to_drop])])
    # replace all nan with np.nan ignore case
    df_wnv2 = df_wnv2.applymap(lambda x: np.nan if isinstance(x, str) and 'nan' in x.lower() else x)
    # numeric strings to float or int
    for col in df_wnv2.columns:
        try:
            df_wnv2[col] = pd.to_numeric(df_wnv2[col])
        except Exception:
            # print(f'Could not convert {col} to numeric')
            pass
    # if 'COBRAD_descriptive.xlsx' doesnt exist
    if not os.path.exists('COBRAD_descriptive.xlsx'):
        # Create a dictionary to store descriptive statistics for each sheet
        desc_stats = {}
        for sheet in sheets_to_read:
            # get only the ids that are in df_wnv2
            dfs[sheet] = dfs[sheet][dfs[sheet]['ID'].isin(df_wnv2['ID'])]
            df_desc = custom_describe(dfs[sheet])
            desc_stats[sheet] = df_desc
        # Save all descriptive statistics to one Excel file
        with pd.ExcelWriter('COBRAD_descriptive.xlsx') as writer:
            for sheet_name, df_desc in desc_stats.items():
                df_desc.to_excel(writer, sheet_name=sheet_name)
    # df_wnv save to csv
    cases_group_name = 'EDF' #'COBRAD'
    return df_wnv,patients_folder,controls,df_wnv2,cases_group_name
    df_merged['ID'].unique()
    df_merged['_merge'].unique()
    df_merged[df_merged['_merge'] == 'right_only']['ID'].unique()
    df_wnv2['ID'].unique()

def get_clinical_and_boxplot_cols(df_wnv2):
       boxplot_columns = [col for col in df_wnv2.columns if 'overall' in col.lower()]
       # split file name .[0] and then '-'[0] to get the ID
       clinical_columns_all = df_wnv2.columns.tolist()
       # remove boxplot_columns from clinical_columns
       clinical_columns = [col for col in clinical_columns_all if col not in boxplot_columns]
       # Remove columns that contain 'EEG'
       clinical_columns = [col for col in clinical_columns if 'EEG' not in col]
       # remove 'ID' from clinical_columns
       clinical_columns = [col for col in clinical_columns if col != 'ID']
       return clinical_columns,boxplot_columns

def save_raw_data(df, figures_dir):
    # Ensure the directory exists
    os.makedirs(f'{figures_dir}/raw_data', exist_ok=True)
    # Save the DataFrame to a CSV file
    df.to_csv(f'{figures_dir}/raw_data/raw_data.csv', index=False)
   
# Custom descriptive statistics function
def custom_describe(df):
    # Convert columns with object dtype to more relevant types
    for col in df.select_dtypes(include=['object']).columns:
        try:
            df[col] = pd.to_datetime(df[col])
        except (ValueError, TypeError):
            try:
                df[col] = pd.to_numeric(df[col])
            except ValueError:
                pass  # If conversion fails, keep the column as object
    stats = {}
    numeric_columns = df.select_dtypes(include=[np.number]).columns
    for col in df.columns:
        # if col values are numeric cases.select_dtypes(include=[np.number]).columns
        if col in numeric_columns:
            stats[col] = {
                'count': df[col].count(),
                'mean': df[col].mean(),
                'std': df[col].std(),
                'median': df[col].median(),
                'unique': df[col].nunique(),
                'min': df[col].min(),
                'max': df[col].max(),
                'iqr': iqr(df[col].dropna())  # IQR requires non-null values
            }
        else:  # Non-numeric columns (e.g., record_id)
            stats[col] = {
                'count': df[col].count(),
                'unique': df[col].nunique()
            }
    df_ret = pd.DataFrame(stats)
    # round 2
    return df_ret.round(2)

import mne  # noqa: E402

def detect_sleep_stages(raw, plot=True):
    """
    Detect sleep stages using YASA on a selected EEG channel.

    Parameters:
    -----------
    raw : mne.io.Raw
        MNE Raw object containing the EEG data.
    channel : str, optional
        The channel name to use for sleep staging (default is 'Fpz').
    epoch_sec : int, optional
        The length (in seconds) of each epoch (default is 30).
    plot : bool, optional
        If True, plots the hypnogram of the sleep stages.

    Returns:
    --------
    predicted_stages : numpy.ndarray
        Array of predicted sleep stage labels for each epoch.
    """
    import yasa

    # Check if the chosen channel is available in the raw data
    channels = raw.info['ch_names']
    # Fz if exists if not C4 if not C3
    eeg_name = 'Fz' if 'Fz' in channels else 'C4' if 'C4' in channels else 'C3'
    # Initialize and run YASA's sleep staging
    ss = yasa.SleepStaging(raw,eeg_name=eeg_name)
    predicted_stages = ss.predict()

    # Plot hypnogram if requested
    if plot:
        plt.figure(figsize=(10, 4))
        plt.plot(predicted_stages, drawstyle='steps-post')
        plt.xlabel('Epochs (sec segments)')
        plt.ylabel('Sleep Stage')
        plt.title(f'Sleep Staging Hypnogram from {eeg_name}')
        # Map stages to labels (0 = Wake, 1 = N1, 2 = N2, 3 = N3, 5 = REM)
        plt.yticks([0, 1, 2, 3, 5], ['W', 'N1', 'N2', 'N3', 'REM'])
        plt.grid(True)
        plt.show(block=False)
    return predicted_stages

def fix_predicted_stages(predicted_stages, min_non_w=10):
    """
    Fix the predicted_stages array by turning sequences of fewer than `min_non_w` 'W' values
    into the nearest neighboring stage.

    Parameters:
    -----------
    predicted_stages : np.ndarray
        Array of predicted sleep stage labels.
    min_non_w : int
        Minimum number of consecutive 'W' values to keep as 'W'.

    Returns:
    --------
    np.ndarray
        Fixed predicted_stages array.
    """
    fixed_stages = predicted_stages.copy()
    w_indices = np.where(predicted_stages == 'W')[0]  # Indices of 'W' values

    # Group consecutive indices of 'W' values
    groups = np.split(w_indices, np.where(np.diff(w_indices) != 1)[0] + 1)

    for group in groups:
        if len(group) < min_non_w and len(group) > 0:
            # Replace with the nearest neighboring stage
            start_idx = group[0]
            end_idx = group[-1]

            # Get the value before the group (if it exists)
            before_value = fixed_stages[start_idx - 1] if start_idx > 0 else None
            # Get the value after the group (if it exists)
            after_value = fixed_stages[end_idx + 1] if end_idx + 1 < len(fixed_stages) else None

            # Decide the replacement value
            if before_value == after_value:
                replacement_value = before_value  # Use the same value if both neighbors are the same
            elif before_value is not None:
                replacement_value = before_value  # Prefer the value before the group
            elif after_value is not None:
                replacement_value = after_value  # Otherwise, use the value after the group
            else:
                replacement_value = 'W'  # Default to 'W' if no neighbors exist

            # Replace the group with the determined value
            fixed_stages[group] = replacement_value

    return fixed_stages

def detect_sleep(cases_group_name, save_sleep_only=False):
    # Get pickle files from pickles/{cases_group_name}/*.pkl
    case_files = []
    for root, dirs, files in os.walk(f'pickles/{cases_group_name}'):
        for file in files:
            if file.endswith('.pkl'):
                case_files.append(os.path.join(root, file))

    # Create output directory for wake data
    wake_dir = f'pickles/wake_{cases_group_name}'
    os.makedirs(wake_dir, exist_ok=True)
    if save_sleep_only:
        sleep_dir = f'pickles/sleep_{cases_group_name}'
        os.makedirs(sleep_dir, exist_ok=True)

    for case_file in case_files:
               # Load the file
        try:
            with open(case_file, 'rb') as f:
                raw = pickle.load(f)
        except Exception as e:
            print(f"Failed to load {case_file}: {e}")
            continue
        # if raw time less than 5 minutes continue
        if raw.times[-1] - raw.times[0] < 5 * 60:
            print(f"File {case_file} has less than 5 minutes of data. Skipping...")
            continue
        # Get the ID from the file name
        ID = os.path.basename(case_file).split('.')[0].split(' ')[0]

        # Detect sleep stages
        predicted_stages = detect_sleep_stages(raw, plot=False)
        duration_sec = raw.times[-1] - raw.times[0]
        epoch_length_sec = int(duration_sec / len(predicted_stages))
        # min_non_w to be 10 minutes. based on epoch_length_sec
        min_non_w = int(10 * 60 / epoch_length_sec)
        # Fix the predicted_stages array
        predicted_stages = fix_predicted_stages(predicted_stages, min_non_w=min_non_w)

        # If there are no awake ('W') segments, skip this file
        if not any(predicted_stages == 'W'):
            print(f"No awake segments found for ID {ID}. Skipping...")
            continue

        # Get the indices of the awake ('W') segments
        awake_indices = np.where(predicted_stages == 'W')[0]
        # Group consecutive awake indices
        awake_groups = np.split(awake_indices, np.where(np.diff(awake_indices) != 1)[0] + 1)
        # Initialize an empty list to store the awake segments
        awake_segments = []
        # Convert awake groups to time in seconds and crop the raw data
        for group in awake_groups:
            start_time = group[0] * epoch_length_sec
            end_time = (group[-1] + 1) * epoch_length_sec
            awake_segments.append(raw.copy().crop(tmin=start_time, tmax=end_time, include_tmax=False))

        # Concatenate all awake segments into a single raw object
        awake_raw = mne.concatenate_raws(awake_segments)

        # Save the clipped raw data to a new pickle file
        wake_file_path = os.path.join(f'{case_file}')
        # remove .edf and .csv
        wake_file_path = wake_file_path.replace('.edf', '').replace('.csv', '')
        # replace 'EDF' with 'wake_EDF'
        wake_file_path = wake_file_path.replace('EDF', f'wake_{cases_group_name}')
        with open(wake_file_path, 'wb') as f:
            pickle.dump(awake_raw, f)

        print(f"Saved awake data for ID {ID} to {wake_file_path}")

        # Save sleep segments if requested
        if save_sleep_only:
            # Find indices where predicted_stages != 'W'
            sleep_indices = np.where(predicted_stages != 'W')[0]
            if len(sleep_indices) == 0:
                print(f"No sleep segments found for ID {ID}.")
                continue
            sleep_groups = np.split(sleep_indices, np.where(np.diff(sleep_indices) != 1)[0] + 1)
            sleep_segments = []
            for group in sleep_groups:
                start_time = group[0] * epoch_length_sec
                end_time = (group[-1] + 1) * epoch_length_sec
                sleep_segments.append(raw.copy().crop(tmin=start_time, tmax=end_time, include_tmax=False))
            if sleep_segments:
                sleep_raw = mne.concatenate_raws(sleep_segments)
                sleep_file_path = case_file.replace('.edf', '').replace('.csv', '')
                sleep_file_path = sleep_file_path.replace('EDF', f'sleep_{cases_group_name}')
                with open(sleep_file_path, 'wb') as f:
                    pickle.dump(sleep_raw, f)
                print(f"Saved sleep data for ID {ID} to {sleep_file_path}")
    return

def read_pkl_files(path):
    """
    Read all pickle files from a given path, process them per ID, and save individual metadata as CSV files.
    Merge all individual CSV files into one consolidated CSV file.

    Parameters:
    -----------
    path : str
        Path to the directory containing the pickle files.

    Returns:
    --------
    pd.DataFrame
        DataFrame containing merged metadata for all file IDs.
    """
    pkl_files = [os.path.join(path, f) for f in os.listdir(path) if f.endswith('.pkl')]
    metadata_list = []
    temp_dir = 'temps_awake_EDF'
    os.makedirs(temp_dir, exist_ok=True)

    for file in pkl_files:
        with open(file, 'rb') as f:
            raw = pickle.load(f)
            # Extract ID from the filename (assuming ID is part of the filename)
            file_id = os.path.basename(file).split(' ')[0]
            metadata = eeg_data_to_features(raw)
            metadata['ID'] = file_id  # Add ID to metadata
            metadata['duration_sec'] = raw.times[-1] - raw.times[0]
            metadata['duration_min'] = metadata['duration_sec'] / 60
            metadata_list.append(metadata)
            # replace .pkl with .csv
            file_name = os.path.basename(file).replace('.pkl', '')
            # Save individual metadata to a CSV file
            individual_csv_path = os.path.join(temp_dir, f'{file_name}.csv')
            pd.DataFrame([metadata]).to_csv(individual_csv_path, index=False)
    # Merge all individual CSV files into one consolidated CSV file
    merge_csv_files(temp_dir)

    
def id_from_csv_filename(filename):
    # re get '0345-{3%d}' from filename
    match = re.search(r'0345-\d{3}', filename)
    if match:
        return match.group(0)[1:]
    else:
        print(f"ID not found in filename: {filename}")
        return None
    
def merge_csv_files(temp_dir):
    # Merge all individual CSV files into one consolidated CSV file
    all_metadata_df = pd.DataFrame()
    for f in os.listdir(temp_dir):
        if f.endswith('.csv'):
            temp_df = pd.read_csv(os.path.join(temp_dir, f))

            temp_df['file_name'] = f.split(' ')[0]  # Extract file name from filename
            temp_df['csv_file_name'] = f  # Extract file name from filename
            # with regex find 0345-%d
            temp_df['ID'] = temp_df['csv_file_name'].apply(id_from_csv_filename)
            # remove ID from file name
            all_metadata_df = pd.concat([all_metadata_df, temp_df], ignore_index=True)
    consolidated_csv_path = os.path.join('EDF_awake.csv')
    # move ID to the first column
    cols = all_metadata_df.columns.tolist()
    cols.insert(0, cols.pop(cols.index('ID')))
    all_metadata_df = all_metadata_df[cols]
    # Save the consolidated DataFrame to a CSV file
    all_metadata_df.to_csv(consolidated_csv_path, index=False)
    return
    # print how many unique ID
    f"Total unique IDs in the consolidated CSV: {all_metadata_df['ID'].nunique()}"
    all_metadata_df['ID'].unique()

def raw_run(cases_group_name='EDF'):
    group = 'west_nile_virus' if cases_group_name == 'WNV' else 'EDF'
    pickle_files = [f for f in os.listdir(f'pickles/{group}') if f.endswith('.pkl')]
    pickle_file = st.selectbox('Select a file', pickle_files)
    
    with open(f'pickles/{group}/{pickle_file}', 'rb') as f:
        raw = pickle.load(f)
    
    total_duration = raw.times[-1]
    st.write(f"Total duration: {total_duration:.2f} seconds ({total_duration / 60:.2f} minutes)")

    start_time = st.number_input("Select start time (seconds)", min_value=0.0, max_value=max(0.0, total_duration), value=0.0, step=10.0, format="%.1f")
    end_time = min(start_time + 10, total_duration)
    cropped_raw = raw.copy().crop(tmin=start_time, tmax=end_time, include_tmax=False)

    st.write(f"Displaying data from {start_time:.2f} to {end_time:.2f} seconds")
    fig = cropped_raw.plot(show=False, block=False)
    st_pyplot_func(fig)
    
    # Compute power spectral density (PSD) for delta band (0.5-4 Hz)
    psd, freqs = raw.compute_psd(fmin=0.5, fmax=4.0).get_data(return_freqs=True)
    delta_power = psd.mean(axis=1)  # Average across frequencies for each channel

    # Find the channel with the highest delta power
    max_channel_idx = delta_power.argmax()
    max_channel_name = raw.info['ch_names'][max_channel_idx]
    st.write(f"Channel with most delta power: {max_channel_name}")

    window_size = st.sidebar.slider("Select window size in seconds", 10, int(total_duration//2), 5)
    # Plot the channel with the highest delta power
    fig_max_channel = raw.plot(start=start_time, duration=window_size, picks=[max_channel_idx], show=False, block=False)
    st.write(f"Plot of {max_channel_name} with most delta power")
    st_pyplot_func(fig_max_channel)

    # Find the channel with the lowest delta power
    min_channel_idx = delta_power.argmin()
    min_channel_name = raw.info['ch_names'][min_channel_idx]
    st.write(f"Channel with least delta power: {min_channel_name}")

    # Plot the channel with the lowest delta power
    fig_min_channel = raw.plot(start=start_time, duration=window_size, picks=[min_channel_idx], show=False, block=False)
    st.write(f"Plot of {min_channel_name} with least delta power")
    st_pyplot_func(fig_min_channel)

    # Find region with the highest delta power across all channels
    max_power_idx = delta_power.argmax()
    start_slow = max_power_idx * window_size
    end_slow = min(start_slow + window_size, total_duration)
    slowing_raw = raw.copy().crop(tmin=start_slow, tmax=end_slow, include_tmax=False)
    st.write(f"Most delta power: {start_slow:.2f} to {end_slow:.2f} seconds")
    fig_slow = slowing_raw.plot(show=False, block=False)
    st_pyplot_func(fig_slow)

    # Find region with the lowest delta power across all channels
    min_power_idx = delta_power.argmin()
    start_fast = min_power_idx * window_size
    end_fast = min(start_fast + window_size, total_duration)
    speeding_raw = raw.copy().crop(tmin=start_fast, tmax=end_fast, include_tmax=False)
    st.write(f"Least delta power: {start_fast:.2f} to {end_fast:.2f} seconds")
    fig_fast = speeding_raw.plot(show=False, block=False)
    st_pyplot_func(fig_fast)

import yasa  # noqa: E402
import pandas as pd  # noqa: E402
import streamlit as st  # noqa: E402
from mne.time_frequency import psd_array_welch  # for PSD  # noqa: E402

def spectrogram_run(group, win_sec=5,is_streamlit=False):

    # --- Otherwise, generate it from pickles ---
    pickle_files = [f for f in os.listdir(f'pickles/{group}') if f.endswith('.pkl')]
    default_sample_k = 25
    load_all = st.button(f"Load all pickles ({len(pickle_files)})")
    if not load_all:
        pickle_files = pickle_files[:default_sample_k]

    arr = []
    progress = st.progress(0)
    total_files = len(pickle_files)

    for i, pickle_file in enumerate(pickle_files, start=1):
        raw = pd.read_pickle(f'pickles/{group}/{pickle_file}')
        data = raw.get_data()
        arr.append(data)
        raw.info['ch_names']
        progress.progress(i / total_files)

    arr_mean = mean_of_resized_arrays(arr)
    n_samples = len(arr)

    # --- PSD on mean data ---
    sf = raw.info['sfreq']
    psds, freqs = psd_array_welch(
        arr_mean,
        sfreq=sf,
        fmin=0.5,
        fmax=40.0,
        n_fft=int(sf * win_sec)
    )
    ch_names = raw.info['ch_names']
    eeg_channels_lower = [ch.lower() for ch in eeg_channels]
    ch_names = sorted(
        ch_names,
        key=lambda x: eeg_channels_lower.index(x.lower())
        if x.lower() in eeg_channels_lower else len(eeg_channels_lower)
    )
    psd_df = pd.DataFrame(psds.T, index=freqs, columns=ch_names)
    st.write(f"Mean PSD over {n_samples} samples")
    st.line_chart(psd_df)

    # --- Generate spectrogram figure ---
    st.write(f"Mean Spectrogram over {n_samples} samples")
    plot_spectrogram_mean(group, arr_mean, win_sec=win_sec, sf=sf)

def plot_spectrogram_mean(group,arr_mean,figures_dir=None,win_sec=5,sf=256):
    # mean over all channels
    fig = yasa.plot_spectrogram(arr_mean.mean(axis=0), sf,win_sec=win_sec, ch_names=['mean'],cmap='jet')
    st_pyplot_func(fig)
    return fig

# Utility function to convert sex column
def convert_sex_column(clinical_df, sex_column):
    """
    Converts 'M'/'m'/'male' to 0 and 'F'/'f'/'female' to 1 in the specified sex column of clinical_df.
    Returns the modified DataFrame.
    """
    mapping = {
        'm': 0, 'male': 0,
        'f': 1, 'female': 1
    }
    clinical_df[sex_column] = clinical_df[sex_column].astype(str).str.strip().str.lower().map(lambda x: mapping.get(x, np.nan))
    return clinical_df
        
def clean_df_demographics(df_demographics,patient_names):
    # strip all values
    df_demographics = df_demographics.map(lambda x: x.strip() if isinstance(x, str) else x)
    # find what column has values 'נקבה'
    sex_col_num = df_demographics.columns[df_demographics.isin(['נקבה']).any()][0]
    # move to be first col
    df_demographics = pd.concat([df_demographics.loc[:, [sex_col_num]], df_demographics.drop(columns=[sex_col_num])], axis=1)
    # rename to gender
    df_demographics.rename(columns={sex_col_num: 'Gender'}, inplace=True)
    # what column contains patient_names at least 1
    id_col_num = df_demographics.columns[df_demographics.isin(patient_names).any()][0]
    # move to be first col
    df_demographics = pd.concat([df_demographics.loc[:, [id_col_num]], df_demographics.drop(columns=[id_col_num])], axis=1)
    # rename to ID
    df_demographics.rename(columns={id_col_num: 'ID'}, inplace=True)
    return df_demographics
    df_demographics.iloc[:,0]


def get_controls_ages_genders(selected_folder, controls_dir='EDF_Format/Controls/Controls'):
    # Define age groups and their ranges
    age_groups = [
        '18-30 data', '31-40 data', '41-50 data', '51-60 data',
        '61-70 data', '71-80 data', '81-90 data', '90+ data'
    ]
    age_groups_int = [
        (18, 30), (31, 40), (41, 50), (51, 60),
        (61, 70), (71, 80), (81, 90), (90, 120)
    ]

    # Get control IDs from the selected folder (.parquet files)
    control_ids = [f.split('_')[0] for f in os.listdir(selected_folder) if f.endswith('.parquet')]
    results = []

    for group, (start, end) in zip(age_groups, age_groups_int):
        group_name = group.split(' ')[0]
        group_path = os.path.join(controls_dir, group)

        # Read demographics
        demo_path = os.path.join(group_path, f'{group_name}.xlsx')
        if not os.path.exists(demo_path):
            print(f"Warning: {demo_path} does not exist")
            continue  # Skip if file doesn't exist
        df_demo = pd.read_excel(demo_path, header=None)
        # Find which column contains any of the control_ids
        col_with_ids = None
        for col in df_demo.columns:
            if df_demo[col].isin(control_ids).any():
                col_with_ids = col
                break
        if col_with_ids is not None:
            # Filter rows where the column matches control_ids
            df_demo2 = df_demo[df_demo[col_with_ids].isin(control_ids)].copy()
            df_demo2 = clean_df_demographics(df_demo2, control_ids)
        else:
            continue

        # For age groups 61-90, also read "final data"
        if start >= 61 and end <= 90:
            demo_path2 = os.path.join(group_path, f'{group_name} final data.xlsx')
            if os.path.exists(demo_path2):
                # rename column 5 to 'Age'
                df_demo2 = pd.read_excel(demo_path2, header=None)
                df_demo2.rename(columns={5: 'Age'}, inplace=True)
                df_demo2 = clean_df_demographics(df_demo2, control_ids)
            else:
                print(f"Warning: {demo_path2} does not exist")
                # If final data doesn't exist, use regular file like other age groups
                df_demo2.rename(columns={7: 'Age'}, inplace=True)
        else:
            # rename column 7 to Age
            df_demo2.rename(columns={7: 'Age'}, inplace=True)

        # Remove duplicate columns and NaN IDs
        sex_col_num = 1
        df_demo2.iloc[:, sex_col_num] = df_demo2.iloc[:, sex_col_num].apply(lambda x: 'f' if x == 'נקבה' else 'm')

        # Filter for selected controls
        df_selected = df_demo2[df_demo2['ID'].isin(control_ids)]
        results.append(df_selected[['ID', df_selected.columns[sex_col_num], 'Age']])  # Assuming age is col 2

    # Concatenate all results
    final_df = pd.concat(results, ignore_index=True)
    final_df.columns = ['ID', 'Gender', 'Age']
    # final_df remove ID coplicates
    final_df = final_df.loc[:, ~final_df.columns.duplicated()].dropna(subset=['ID'])
    # convert_sex_column
    final_df = convert_sex_column(final_df,'Gender')
    return final_df

from itertools import combinations  # noqa: E402
from scipy.signal import coherence, windows  # noqa: E402
import networkx as nx  # noqa: E402
def compute_network_features(eeg_data, sfreq, freq_band, threshold_density=0.2):
    """
    Compute network features from EEG data using coherence in a given frequency band.

    Parameters:
        eeg_data (ndarray): EEG data (channels x samples)
        sfreq (float): Sampling frequency
        freq_band (tuple): Frequency band (low, high) in Hz
        threshold_density (float): Proportion of top connections to retain in binarized graph

    Returns:
        dict: Network metrics (clustering, efficiency, assortativity, modularity)
    """
    n_channels = eeg_data.shape[0]
    coh_matrix = np.zeros((n_channels, n_channels))
    # sfreq to int
    sfreq = int(sfreq)
    
    # Validate input data
    if n_channels < 2:
        return 0.0, 0.0, 0.0, 0.0
    
    # Check for valid data
    if np.any(~np.isfinite(eeg_data)):
        # Replace non-finite values with zeros
        eeg_data = np.nan_to_num(eeg_data, nan=0.0, posinf=0.0, neginf=0.0)
    # Compute coherence between all pairs
    for i, j in combinations(range(n_channels), 2):
        try:
            # Check for valid data before computing coherence
            if (np.std(eeg_data[i]) < 1e-10 or np.std(eeg_data[j]) < 1e-10 or 
                np.allclose(eeg_data[i], eeg_data[j], atol=1e-10)):
                # Skip if signals are too similar or have no variance
                coh_matrix[i, j] = coh_matrix[j, i] = 0.0
                continue
                
            # Suppress the RuntimeWarning from scipy coherence
            import warnings
            with warnings.catch_warnings():
                warnings.simplefilter("ignore", RuntimeWarning)
                # Use a smaller window size to avoid division by zero issues
                window_size = min(sfreq, 512)  # Use smaller window size
                try:
                    f, Cxy = coherence(eeg_data[i], eeg_data[j], fs=sfreq, 
                                       window=windows.hann(window_size),
                                       nperseg=window_size, noverlap=window_size//2)
                except Exception:
                    # If coherence fails, use correlation as fallback
                    corr = np.corrcoef(eeg_data[i], eeg_data[j])[0, 1]
                    if np.isfinite(corr):
                        mean_coh = abs(corr)
                    else:
                        mean_coh = 0.0
                    coh_matrix[i, j] = coh_matrix[j, i] = mean_coh
                    continue
            freq_mask = (f >= freq_band[0]) & (f <= freq_band[1])
            
            # Handle potential NaN or infinite values
            Cxy_clean = Cxy[freq_mask]
            Cxy_clean = Cxy_clean[np.isfinite(Cxy_clean)]
            
            if len(Cxy_clean) > 0:
                mean_coh = np.mean(Cxy_clean)
                # Ensure coherence is between 0 and 1
                mean_coh = np.clip(mean_coh, 0.0, 1.0)
            else:
                mean_coh = 0.0
                
            coh_matrix[i, j] = coh_matrix[j, i] = mean_coh
            
        except Exception:
            # If coherence calculation fails, set to 0
            coh_matrix[i, j] = coh_matrix[j, i] = 0.0

    # Binarize: keep top X% of connections (thresholding by density)
    n_possible = n_channels * (n_channels - 1) / 2
    n_edges = int(threshold_density * n_possible)
    triu_vals = coh_matrix[np.triu_indices(n_channels, k=1)]
    thresh_val = np.sort(triu_vals)[-n_edges] if n_edges > 0 else 0
    binarized = (coh_matrix >= thresh_val).astype(int)
    np.fill_diagonal(binarized, 0)

    # Build graph and compute features
    G = nx.from_numpy_array(binarized)

    try:
        clustering = nx.transitivity(G)
        if not np.isfinite(clustering):
            clustering = 0.0
    except Exception:
        clustering = 0.0

    try:
        efficiency = nx.global_efficiency(G)
        if not np.isfinite(efficiency):
            efficiency = 0.0
    except Exception:
        efficiency = 0.0

    try:
        assortativity = nx.degree_pearson_correlation_coefficient(G)
        if not np.isfinite(assortativity):
            assortativity = 0.0
    except Exception:
        assortativity = 0.0

    try:
        modularity = nx.algorithms.community.modularity(G, nx.community.label_propagation_communities(G))
        if not np.isfinite(modularity):
            modularity = 0.0
    except Exception:
        modularity = 0.0

    return efficiency, clustering, assortativity, modularity

def _create_interactive_significant_pairs_plot(df_pair, var_names, pair_to_stats, name, save_plot, save_dir, is_streamlit, hue, size_values, palette_dict):
    """Create an interactive scatter plot showing only significant pairs using Plotly."""
    import plotly.graph_objects as go
    from plotly.subplots import make_subplots
    import numpy as np
    from scipy.stats import spearmanr
    
    # Find significant pairs (p < 0.05)
    significant_pairs = []
    for (i, j), (r, p_fdr) in pair_to_stats.items():
        if not np.isnan(r) and not np.isnan(p_fdr) and p_fdr < 0.05:
            significant_pairs.append((i, j, r, p_fdr, var_names[i], var_names[j]))
    
    if not significant_pairs:
        print("No significant pairs found for interactive plot.")
        return
    
    # Create subplots layout - 1 column, rows = number of pairs
    n_pairs = len(significant_pairs)
    rows, cols = n_pairs, 1
    
    # Create subplots
    fig = make_subplots(
        rows=rows, cols=cols,
        subplot_titles=[f"{var_names[i]} vs {var_names[j]}" for i, j, r, p_fdr, x_name, y_name in significant_pairs],
        specs=[[{"secondary_y": False}] for _ in range(rows)]
    )
    
    # Plot each significant pair
    for idx, (i, j, r, p_fdr, x_name, y_name) in enumerate(significant_pairs):
        row = idx + 1  # Since we have 1 column, row = index + 1
        col = 1
        
        x_data = df_pair[var_names[i]]
        y_data = df_pair[var_names[j]]
        
        # Clean data
        mask = ~np.isnan(x_data) & ~np.isnan(y_data)
        x_clean = np.array(x_data)[mask]
        y_clean = np.array(y_data)[mask]
        
        if '__hue__' in df_pair.columns and hue is not None:
            hue_clean = np.array(df_pair['__hue__'])[mask]
            unique_hues = np.unique(hue_clean)
            
            # Plot with hue
            for hue_val in unique_hues:
                hue_mask = hue_clean == hue_val
                if np.sum(hue_mask) > 0:
                    x_hue = x_clean[hue_mask]
                    y_hue = y_clean[hue_mask]
                    color = palette_dict.get(hue_val, 'gray') if palette_dict else 'gray'
                    
                    # Add size variation if size_values provided
                    if size_values is not None:
                        size_clean = np.array(size_values)[mask]
                        size_hue = size_clean[hue_mask]
                        size_normalized = (size_hue - np.min(size_hue)) / (np.max(size_hue) - np.min(size_hue) + 1e-8)
                        marker_size = 8 + 12 * size_normalized  # Size between 8 and 20
                    else:
                        marker_size = 10
                    
                    # Create hover text
                    hover_text = [f"{x_name}: {x:.3f}<br>{y_name}: {y:.3f}<br>Group: {hue_val}" 
                                 for x, y in zip(x_hue, y_hue)]
                    
                    # Add scatter plot
                    fig.add_trace(
                        go.Scatter(
                            x=x_hue, y=y_hue,
                            mode='markers',
                            name=str(hue_val),
                            marker=dict(
                                color=color,
                                size=marker_size,
                                opacity=0.7
                            ),
                            text=hover_text,
                            hovertemplate='%{text}<extra></extra>',
                            showlegend=True
                        ),
                        row=row, col=col
                    )
                    
                    # Add regression line for this hue group
                    if len(x_hue) > 1:
                        z = np.polyfit(x_hue, y_hue, 1)
                        x_line = np.linspace(x_hue.min(), x_hue.max(), 100)
                        y_line = z[0] * x_line + z[1]
                        
                        # Calculate correlation for this group
                        r_hue, p_hue = spearmanr(x_hue, y_hue)
                        
                        fig.add_trace(
                            go.Scatter(
                                x=x_line, y=y_line,
                                mode='lines',
                                name=f"{hue_val} trend",
                                line=dict(color=color, width=2),
                                showlegend=False,
                                hovertemplate=f"{hue_val} trend<br>R²={r_hue**2:.3f}, p={p_hue:.3g}<extra></extra>"
                            ),
                            row=row, col=col
                        )
        else:
            # Plot without hue
            if size_values is not None:
                size_clean = np.array(size_values)[mask]
                size_normalized = (size_clean - np.min(size_clean)) / (np.max(size_clean) - np.min(size_clean) + 1e-8)
                marker_size = 8 + 12 * size_normalized
            else:
                marker_size = 10
            
            # Use base palette color for this pair
            marker_color = BASE_PALETTE[idx % len(BASE_PALETTE)]
            
            # Create hover text
            hover_text = [f"{x_name}: {x:.3f}<br>{y_name}: {y:.3f}" for x, y in zip(x_clean, y_clean)]
            
            # Add scatter plot
            fig.add_trace(
                go.Scatter(
                    x=x_clean, y=y_clean,
                    mode='markers',
                    name='Data',
                    marker=dict(
                        color=marker_color,
                        size=marker_size,
                        opacity=0.7
                    ),
                    text=hover_text,
                    hovertemplate='%{text}<extra></extra>',
                    showlegend=False
                ),
                row=row, col=col
            )
            
            # Add regression line
            if len(x_clean) > 1:
                z = np.polyfit(x_clean, y_clean, 1)
                x_line = np.linspace(x_clean.min(), x_clean.max(), 100)
                y_line = z[0] * x_line + z[1]
                
                fig.add_trace(
                    go.Scatter(
                        x=x_line, y=y_line,
                        mode='lines',
                        name='Trend',
                        line=dict(color=marker_color, width=2),
                        showlegend=False,
                        hovertemplate=f"Trend line<br>R²={r**2:.3f}, p={p_fdr:.3g}<extra></extra>"
                    ),
                    row=row, col=col
                )
        
        # Update axes labels
        fig.update_xaxes(title_text=x_name, row=row, col=col)
        fig.update_yaxes(title_text=y_name, row=row, col=col)
    
    # Update layout
    fig.update_layout(
        title=f'Significant Pairs - {name.replace("_", " ")} (N={len(df_pair)})',
        height=400 * rows,  # Each subplot gets 400px height
        width=800,  # Fixed width for single column
        showlegend=True,
        hovermode='closest'
    )
    
    # Add grid
    fig.update_xaxes(showgrid=True, gridwidth=1, gridcolor='lightgray')
    fig.update_yaxes(showgrid=True, gridwidth=1, gridcolor='lightgray')
    
    if save_plot:
        import os
        os.makedirs(save_dir, exist_ok=True)
        fname = f"{save_dir}/interactive_significant_pairs_{name}.html"
        fig.write_html(fname)
        print(f"Interactive plot saved to: {fname}")
    else:
        if is_streamlit:
            # For streamlit, show the plotly figure
            st.plotly_chart(fig, use_container_width=True)
        else:
            fig.show()

def _plot_metric_vs_hrv(t, arrs, labels, save_plot, name, save_dir, is_streamlit=False, hue=None, size_values=None):
    if len(arrs) == 2:
        import dabest
        df = pd.DataFrame({
            'value': np.concatenate([arrs[0], arrs[1]]),
            'group': [labels[0]] * len(arrs[0]) + [labels[1]] * len(arrs[1]),
            'pair_id': list(range(len(arrs[0]))) + list(range(len(arrs[1])))
        })
        dabest_obj = dabest.load(data=df, x='group', y='value', id_col='pair_id', paired=True)
        fig = dabest_obj.plot(
            swarm_label="Value",
            contrast_label="Mean difference",
            custom_palette=None,
            show_pairs=True,
            raw_marker_size=4,
            halfviolin_width=0.6,
            fig_size=(6, 4)
        )
        if save_plot:
            os.makedirs(save_dir, exist_ok=True)
            fname = f"{save_dir}/dabest_{labels[0]}_vs_{labels[1]}_{name}.png"
            fig.savefig(fname, dpi=300, bbox_inches='tight')
            plt.close(fig)
        else:
            if is_streamlit:
                # Ensure the figure is properly managed by pyplot
                plt.figure(fig.number)
                st.pyplot(fig)
            else:
                plt.show()
    elif len(arrs) > 2:
        n = len(arrs[0])
        # assert are all the arrs same length
        assert all(len(arr) == n for arr in arrs), "All input arrays must"
        pd.DataFrame({label: arr for label, arr in zip(labels, arrs)})
        df = pd.DataFrame({
            'value': np.concatenate(arrs),
            'group': np.concatenate([[label]*n for label in labels]),
            'pair_id': np.tile(np.arange(n), len(arrs))
        })
        df_pair = pd.DataFrame({label: arr for label, arr in zip(labels, arrs)})
        if hue is not None:
            try:
                if len(hue) == len(df_pair):
                    # ignore index
                    df_pair['__hue__'] = hue.reset_index(drop=True)
                    hue_arr = np.array(hue)
                    # Create a custom palette of distinct RGB colors
                    import matplotlib.colors as mcolors
                    unique_groups = list(pd.Series(hue_arr).dropna().unique())
                    n_groups = len(unique_groups)
                    # If more groups than palette, extend with distinct colors
                    if n_groups > len(BASE_PALETTE):
                        extra_colors = list(mcolors.CSS4_COLORS.values())
                        extra_colors = [c for c in extra_colors if c not in BASE_PALETTE]
                        custom_palette = BASE_PALETTE + extra_colors[:n_groups-len(BASE_PALETTE)]
                    else:
                        custom_palette = BASE_PALETTE[:n_groups]
                    palette_dict = {group: custom_palette[i % len(custom_palette)] for i, group in enumerate(unique_groups)}
            except Exception as e:
                print(f"Error processing hue: {e}")

        g = sns.PairGrid(df_pair, diag_sharey=False)
        def calculate_robust_correlations(df_pair, var_names):
            """Calculate robust correlations for all variable pairs with FDR correction."""
            n_vars = len(var_names)
            pair_indices, r_vals, p_vals = [], [], []
            
            for i in range(n_vars):
                for j in range(i+1, n_vars):
                    x = df_pair[var_names[i]]
                    y = df_pair[var_names[j]]
                    
                    # Convert to numeric, coercing errors to NaN
                    try:
                        x_numeric = pd.to_numeric(x, errors='coerce')
                        y_numeric = pd.to_numeric(y, errors='coerce')
                    except Exception:
                        r, p = np.nan, np.nan
                        pair_indices.append((i, j))
                        r_vals.append(r)
                        p_vals.append(p)
                        continue
                    
                    # Create mask for valid numeric values
                    mask = ~np.isnan(x_numeric) & ~np.isnan(y_numeric)
                    x_clean = np.array(x_numeric)[mask]
                    y_clean = np.array(y_numeric)[mask]
                    
                    if len(x_clean) > 1 and len(y_clean) > 1:
                        # Use normal linear regression
                        r, p = spearmanr(x_clean, y_clean)
                    else:
                        r, p = np.nan, np.nan
                    
                    pair_indices.append((i, j))
                    r_vals.append(r)
                    p_vals.append(p)
            
            # Apply FDR correction
            p_vals_array = np.array([p if not np.isnan(p) else 1.0 for p in p_vals])
            reject, pvals_fdr = fdrcorrection(p_vals_array, alpha=0.05, method='indep')
            pair_to_stats = {(i, j): (r_vals[idx], pvals_fdr[idx]) for idx, (i, j) in enumerate(pair_indices)}
            
            return pair_to_stats
        
        var_names = [c for c in df_pair.columns if c != '__hue__']
        pair_to_stats = calculate_robust_correlations(df_pair, var_names)

        def scatter_with_stats(x, y, hue=None, palette_dict=None, **kwargs):
            ax = kwargs.get('ax', plt.gca())
            i = j = None
            for idx, col in enumerate(var_names):
                if np.all(x == df_pair[col]):
                    i = idx
                if np.all(y == df_pair[col]):
                    j = idx
            if i is not None and j is not None and i != j:
                key = (min(i, j), max(i, j))
                r, p_fdr = pair_to_stats.get(key, (np.nan, np.nan))
            else:
                r, p_fdr = np.nan, np.nan

            # Convert to numeric, coercing errors to NaN
            try:
                x_numeric = pd.to_numeric(x, errors='coerce')
                y_numeric = pd.to_numeric(y, errors='coerce')
            except Exception:
                # If conversion fails, skip this plot
                return
            
            mask = ~np.isnan(x_numeric) & ~np.isnan(y_numeric)
            x_clean = np.array(x_numeric)[mask]
            y_clean = np.array(y_numeric)[mask]
            hue_clean = None
            if hue is not None:
                try:
                    hue_arr = np.array(hue)
                    hue_clean = hue_arr[mask]
                except Exception:
                    hue_clean = None
            if hue_clean is not None:
                # Normalize size_values to 0-1 range for alpha if size_values is available
                if size_values is not None:
                    size_normalized = (size_values - np.min(size_values)) / (np.max(size_values) - np.min(size_values))
                    sns.scatterplot(x=x_clean, y=y_clean, hue=hue_clean, ax=ax, alpha=size_normalized, palette=palette_dict, legend=False)
                else:
                    sns.scatterplot(x=x_clean, y=y_clean, hue=hue_clean, ax=ax, alpha=0.7, palette=palette_dict, legend=False)
                # Create normal regression lines and calculate stats for each hue group
                unique_hues = np.unique(hue_clean)
                hue_stats = []
                for hue_val in unique_hues:
                    hue_mask = hue_clean == hue_val
                    if np.sum(hue_mask) > 1:  # Need at least 2 points for regression
                        x_hue = x_clean[hue_mask]
                        y_hue = y_clean[hue_mask]
                        if len(x_hue) > 1 and len(y_hue) > 1:
                            color = palette_dict.get(hue_val, 'gray') if palette_dict else 'gray'
                            sns.regplot(x=x_hue, y=y_hue, ax=ax, scatter=False, color=color, 
                                      line_kws={'alpha':0.8})
                            
                            # Calculate correlation for this hue group
                            r_hue, p_hue = spearmanr(x_hue, y_hue)
                            hue_stats.append((hue_val, r_hue, p_hue))
                
                # Display statistics for each hue group with colored text
                if hue_stats:
                    y_pos = 0.85
                    for i, (hue_val, r_hue, p_hue) in enumerate(hue_stats):
                        if not np.isnan(r_hue) and not np.isnan(p_hue):
                            color = palette_dict.get(hue_val, 'gray') if palette_dict else 'gray'
                            stats_text = f"R²={r_hue**2:.2f}, p={p_hue:.3g}"
                            ax.annotate(stats_text,
                                      xy=(0.05, y_pos - i*0.08), xycoords='axes fraction', 
                                      fontsize=8, color=color, ha='left', va='center',
                                      bbox=dict(boxstyle="round,pad=0.2", fc="white", ec=color, alpha=0.3))
            else:
                # Normalize size_values to 0-1 range for alpha if size_values is available
                if size_values is not None:
                    size_normalized = (size_values - np.min(size_values)) / (np.max(size_values) - np.min(size_values))
                    sns.scatterplot(x=x_clean, y=y_clean, ax=ax, alpha=size_normalized)
                else:
                    sns.scatterplot(x=x_clean, y=y_clean, ax=ax, alpha=0.7)
                if len(x_clean) > 1 and len(y_clean) > 1:
                    sns.regplot(x=x_clean, y=y_clean, ax=ax, scatter=False, color='gray', line_kws={'alpha':0.8})
                # Display overall statistics when no hue grouping
                if not np.isnan(r) and not np.isnan(p_fdr) and p_fdr < 0.05:
                    ax.annotate(f"$R^2$={r**2:.2f}\np={p_fdr:.3g}",
                                xy=(0.05, 0.85), xycoords='axes fraction', fontsize=9,
                                bbox=dict(boxstyle="round,pad=0.2", fc="white", ec="gray", alpha=0.7))

        def lineplot_with_time(x, y,palette_dict=None, **kwargs):
            ax = kwargs.get('ax', plt.gca())
            
            # Convert to numeric, coercing errors to NaN
            try:
                x_numeric = pd.to_numeric(x, errors='coerce')
                y_numeric = pd.to_numeric(y, errors='coerce')
            except Exception:
                # If conversion fails, skip this plot
                return
            
            # Filter out NaN values
            mask = ~np.isnan(x_numeric) & ~np.isnan(y_numeric)
            if mask.sum() < 2:  # Need at least 2 points
                return
                
            x_clean = x_numeric[mask]
            y_clean = y_numeric[mask]
            
            t_temp = np.linspace(min(min(x_clean), min(y_clean)), max(max(x_clean), max(y_clean)), num=len(x_clean))
            i = j = None
            for idx, col in enumerate(var_names):
                if np.all(x == df_pair[col]):
                    i = idx
                if np.all(y == df_pair[col]):
                    j = idx
            if i is not None and j is not None and i != j:
                key = (min(i, j), max(i, j))
                r, p_fdr = pair_to_stats.get(key, (np.nan, np.nan))
            else:
                r, p_fdr = np.nan, np.nan
            
            ax.plot(t_temp, x_clean, label='x', alpha=0.5,color= 'blue')
            ax.plot(t_temp, y_clean, label='y', alpha=0.5,color='red')
            if not np.isnan(r) and not np.isnan(p_fdr) and p_fdr < 0.05:
                ax.annotate(f"$R^2$={r**2:.2f}\np={p_fdr:.3g}",
                            xy=(0.05, 0.85), xycoords='axes fraction', fontsize=9,
                            bbox=dict(boxstyle="round,pad=0.2", fc="white", ec="gray", alpha=0.7))
            ax.legend(fontsize=8)
            ax.set_xlabel('t')
            ax.set_ylabel('Value')

        g.map_diag(sns.histplot, kde=True, bins=10)
        if '__hue__' in df_pair.columns:
            g.map_upper(scatter_with_stats, hue=df_pair['__hue__'],palette_dict=palette_dict)
            g.map_lower(lineplot_with_time, palette_dict=palette_dict)
        else:
            g.map_upper(scatter_with_stats)
            g.map_lower(lineplot_with_time)
        g.fig.suptitle(f'{name.replace("_"," ")}\nN={len(df_pair)}')
        # 1) Nuke any legends Seaborn created
        if hasattr(g, "_legend") and g._legend is not None:
            g._legend.remove()
        for ax in g.axes.flat:
            leg = ax.get_legend()
            if leg is not None:
                leg.remove()

        # 2) Build one legend and put it on the first row, last col (only if hue exists)
        if '__hue__' in df_pair.columns:
            first_ax = g.axes[0, -1]  # top-right subplot
            levels = pd.Index(df_pair['__hue__'].dropna().unique())
            import matplotlib as mpl
            handles = []
            for lvl in levels:
                color = palette_dict.get(lvl) if palette_dict is not None else None
                handles.append(
                    mpl.lines.Line2D(
                        [], [], marker='o', linestyle='',
                        markerfacecolor=color, markeredgecolor='none',
                        markersize=6, label=str(lvl)
                    )
                )

            leg = first_ax.legend(
                handles=handles,
                title="Group",          # <- new title
                loc="best",
                frameon=True,
                fontsize=8,             # <- smaller labels
                title_fontsize=9
            )
        plt.tight_layout()
        if save_plot:
            os.makedirs(save_dir, exist_ok=True)
            fname2 = f"{save_dir}/pairgrid_all_{name}.png"
            g.savefig(fname2, dpi=300, bbox_inches='tight')
            plt.close(g.fig)
        else:
            if is_streamlit:
                # Ensure the figure is properly managed by pyplot
                plt.figure(g.fig.number)
                st.pyplot(g.fig)
            else:
                plt.show()
        
        # Create interactive scatter plot for significant pairs only
        # _create_interactive_significant_pairs_plot(df_pair, var_names, pair_to_stats, name, save_plot, save_dir, is_streamlit, hue, size_values, palette_dict)
    else:
        fig, ax = plt.subplots(figsize=(6, 4))
        ax.plot(t, arrs[0], label=labels[0])
        ax.set_title(labels[0])
        ax.set_xlabel('Time (s)')
        ax.legend()
        plt.tight_layout()
        if save_plot:
            os.makedirs(save_dir, exist_ok=True)
            fname = f"{save_dir}/{labels[0]}_{name}.png"
            fig.savefig(fname, dpi=300, bbox_inches='tight')
            plt.close(fig)
        else:
            if is_streamlit:
                # Ensure the figure is properly managed by pyplot
                plt.figure(fig.number)
                st.pyplot(fig)
            else:
                plt.show()

def only_plots(results_df, save_plot, save_dir, edf_pickle_name="plot", band="band", step_sec=5,is_streamlit=False,hue=None,size_feature=None):

    labels = [
        'Vagal_SD1',
        'Sympathetic_SD2',
        'Efficiency',
        'Clustering',
        'Modularity',
        'Assortativity'
    ]
    arrs_zscore = [results_df[col].values for col in labels if col in results_df]
    t = np.arange(len(results_df)) * step_sec
    
    # Get size_feature data if specified
    size_values = None
    if size_feature and size_feature in results_df.columns:
        size_values = results_df[size_feature].values
    elif size_feature:
        # If size_feature is specified but not in results_df, try to find it in the original data
        # This might happen if the feature was in the original data but not in the processed results_df
        st.warning(f"Size feature '{size_feature}' not found in results data. Using default size.")
        size_values = None
    
    _plot_metric_vs_hrv(t, arrs_zscore, labels, save_plot, f"{edf_pickle_name}_{band}", save_dir, is_streamlit=is_streamlit,hue=hue,size_values=size_values)
    return
    hue.unique()

def enhanced_clustering_plots(results_df, df_wnv3_clean, save_plot, save_dir, edf_pickle_name="plot", band="band", n_clusters=3, is_streamlit=False, size_feature=None):
    """
    Enhanced clustering plots function that performs K-means clustering and classification analysis.
    
    Parameters:
    -----------
    results_df : pd.DataFrame
        DataFrame with clustering features (Vagal_SD1, Sympathetic_SD2, etc.)
    df_wnv3_clean : pd.DataFrame
        DataFrame with clinical features for classifier training
    n_clusters : int
        Number of clusters for K-means (default: 3)
    """
    from sklearn.cluster import KMeans
    from sklearn.ensemble import RandomForestClassifier
    
    # Define the clustering features
    cluster_features = [
        'Vagal_SD1',
        'Sympathetic_SD2', 
        'Efficiency',
        'Clustering',
        'Modularity',
        'Assortativity'
    ]
    
    # Filter to only include features that exist in results_df
    available_features = [col for col in cluster_features if col in results_df.columns]
    if len(available_features) < 2:
        if is_streamlit:
            st.error(f"Not enough clustering features available. Found: {available_features}")
        return
    
    # Prepare data for clustering
    X_cluster = results_df[available_features].dropna()
    if len(X_cluster) < n_clusters:
        if is_streamlit:
            st.error(f"Not enough data points for clustering. Need at least {n_clusters}, got {len(X_cluster)}")
        return
    
    # Perform K-means clustering
    scaler = StandardScaler()
    X_scaled = scaler.fit_transform(X_cluster)
    
    kmeans = KMeans(n_clusters=n_clusters, random_state=42, n_init=10)
    cluster_labels = kmeans.fit_predict(X_scaled)
    # print unique from cluster_labels
    np.unique(cluster_labels)
    # Add cluster labels to results_df
    results_df_with_clusters = results_df.copy()
    results_df_with_clusters['cluster'] = cluster_labels
    
    if is_streamlit:
        st.write(f"K-means clustering completed with {n_clusters} clusters")
        st.write(f"Cluster distribution: {pd.Series(cluster_labels).value_counts().to_dict()}")
    
    # Train classifier if df_wnv3_clean is provided
    feature_importance = None
    if df_wnv3_clean is not None and not df_wnv3_clean.empty:
        # Find common indices between results_df and df_wnv3_clean
        common_indices = results_df_with_clusters.index.intersection(df_wnv3_clean.index)
        if len(common_indices) > 10:  # Need sufficient data for training
            # Prepare features for classification
            # Get numeric columns from df_wnv3_clean
            numeric_cols = df_wnv3_clean.select_dtypes(include=[np.number]).columns
            X_classifier = df_wnv3_clean.loc[common_indices, numeric_cols].dropna(axis=1)
            
            # Remove columns with too many missing values or constant values
            X_classifier = X_classifier.loc[:, X_classifier.nunique() > 1]
            X_classifier = X_classifier.loc[:, X_classifier.isnull().mean() < 0.5]
            
            if len(X_classifier.columns) > 0:
                # Get cluster labels for common indices - ensure we only use indices that exist in both
                y_classifier = results_df_with_clusters.loc[common_indices, 'cluster'].dropna()
                
                # Find the intersection of indices that exist in both X and y after all filtering
                final_common_indices = X_classifier.index.intersection(y_classifier.index)
                
                if len(final_common_indices) > 10:
                    X_final = X_classifier.loc[final_common_indices]
                    y_final = y_classifier.loc[final_common_indices]
                    
                    # Ensure X_final and y_final have the same number of samples
                    if len(X_final) == len(y_final):
                        # Train Random Forest classifier
                        rf = RandomForestClassifier(n_estimators=100, random_state=42)
                        rf.fit(X_final, y_final)
                        
                        # Get feature importance
                        feature_importance = pd.DataFrame({
                            'feature': X_final.columns,
                            'importance': rf.feature_importances_
                        }).sort_values('importance', ascending=False)
                        
                        if is_streamlit:
                            st.write("Classifier trained successfully")
                            st.write(f"Training data: {len(X_final)} samples, {len(X_final.columns)} features")
                            st.write("Top 10 most important features:")
                            st.dataframe(feature_importance.head(10))
                    else:
                        if is_streamlit:
                            st.warning(f"Data dimension mismatch: X has {len(X_final)} samples, y has {len(y_final)} samples. Skipping classifier training.")
                else:
                    if is_streamlit:
                        st.warning(f"Not enough common data for classifier training. Found {len(final_common_indices)} samples, need at least 10.")
    
    # Create PairGrid with enhanced functionality
    _create_enhanced_pairgrid(results_df_with_clusters, available_features, cluster_labels, 
                             feature_importance, edf_pickle_name, band, save_plot, save_dir, 
                             is_streamlit, size_feature)
    
    return results_df_with_clusters, feature_importance

def _create_enhanced_pairgrid(results_df, features, cluster_labels, feature_importance, 
                            name, band, save_plot, save_dir, is_streamlit, size_feature):
    """Create enhanced PairGrid with clustering and feature importance."""
    import matplotlib.pyplot as plt
    import seaborn as sns
    from scipy.stats import spearmanr
    
    # Prepare data for plotting
    df_plot = results_df[features + ['cluster']].dropna()
    
    # Create color palette for clusters
    n_clusters = len(np.unique(cluster_labels))
    colors = plt.cm.Set1(np.linspace(0, 1, n_clusters))
    cluster_colors = {i: colors[i] for i in range(n_clusters)}
    
    # Create PairGrid
    g = sns.PairGrid(df_plot, diag_sharey=False)
    
    def plot_upper(x, y, **kwargs):
        """Plot upper triangle with cluster-colored scatter plots."""
        ax = kwargs.get('ax', plt.gca())
        
        # Get cluster labels for this plot
        
        # Create scatter plot with cluster colors
        for cluster_id in range(n_clusters):
            mask = df_plot['cluster'] == cluster_id
            if mask.sum() > 0:
                ax.scatter(x[mask], y[mask], 
                          c=cluster_colors[cluster_id], 
                          label=f'Cluster {cluster_id}',
                          alpha=0.7, s=50)
        
        # Add regression line for each cluster
        for cluster_id in range(n_clusters):
            mask = df_plot['cluster'] == cluster_id
            if mask.sum() > 2:  # Need at least 3 points for regression
                x_cluster = x[mask]
                y_cluster = y[mask]
                if len(x_cluster) > 2:
                    # Clean data: remove NaN and infinite values
                    valid_mask = np.isfinite(x_cluster) & np.isfinite(y_cluster)
                    x_clean = x_cluster[valid_mask]
                    y_clean = y_cluster[valid_mask]
                    
                    # Check if we have enough valid data points and variance
                    if (len(x_clean) > 2 and 
                        np.var(x_clean) > 1e-10 and 
                        np.var(y_clean) > 1e-10):
                        try:
                            z = np.polyfit(x_clean, y_clean, 1)
                            p = np.poly1d(z)
                            x_line = np.linspace(x_clean.min(), x_clean.max(), 100)
                            ax.plot(x_line, p(x_line), color=cluster_colors[cluster_id], 
                                   linestyle='--', alpha=0.8, linewidth=2)
                        except (np.linalg.LinAlgError, ValueError):
                            # Skip this cluster if polyfit fails
                            continue
        
        ax.legend(bbox_to_anchor=(1.05, 1), loc='upper left', fontsize=8)
    
    def plot_lower(x, y, **kwargs):
        """Plot lower triangle with feature importance or correlation matrix."""
        ax = kwargs.get('ax', plt.gca())
        
        if feature_importance is not None:
            # Show feature importance as text
            x_name = x.name
            y_name = y.name
            
            # Find importance for these features if they exist
            x_imp = feature_importance[feature_importance['feature'] == x_name]['importance'].values
            y_imp = feature_importance[feature_importance['feature'] == y_name]['importance'].values
            
            if len(x_imp) > 0 and len(y_imp) > 0:
                ax.text(0.5, 0.7, f'{x_name} importance:\n{x_imp[0]:.3f}', 
                       transform=ax.transAxes, ha='center', va='center',
                       bbox=dict(boxstyle="round,pad=0.3", facecolor='lightblue', alpha=0.7))
                ax.text(0.5, 0.3, f'{y_name} importance:\n{y_imp[0]:.3f}', 
                       transform=ax.transAxes, ha='center', va='center',
                       bbox=dict(boxstyle="round,pad=0.3", facecolor='lightgreen', alpha=0.7))
            else:
                # Fallback to correlation
                r, p = spearmanr(x, y)
                ax.text(0.5, 0.5, f'R = {r:.3f}\np = {p:.3f}', 
                       transform=ax.transAxes, ha='center', va='center',
                       bbox=dict(boxstyle="round,pad=0.3", facecolor='lightgray', alpha=0.7))
        else:
            # Show correlation
            r, p = spearmanr(x, y)
            ax.text(0.5, 0.5, f'R = {r:.3f}\np = {p:.3f}', 
                   transform=ax.transAxes, ha='center', va='center',
                   bbox=dict(boxstyle="round,pad=0.3", facecolor='lightgray', alpha=0.7))
        
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.set_xticks([])
        ax.set_yticks([])
    
    def plot_diag(x, **kwargs):
        """Plot diagonal with histograms for each cluster."""
        ax = kwargs.get('ax', plt.gca())
        
        for cluster_id in range(n_clusters):
            mask = df_plot['cluster'] == cluster_id
            if mask.sum() > 0:
                ax.hist(x[mask], alpha=0.6, color=cluster_colors[cluster_id], 
                       label=f'Cluster {cluster_id}', bins=10)
        
        ax.legend(fontsize=8)
    
    # Apply the plotting functions
    g.map_upper(plot_upper)
    g.map_lower(plot_lower)
    g.map_diag(plot_diag)
    
    # Set title
    g.fig.suptitle(f'{name.replace("_"," ")} - Clustering Analysis\nN={len(df_plot)}, Clusters={n_clusters}', 
                   fontsize=14, y=1.02)
    
    plt.tight_layout()
    
    # Save or display
    if save_plot:
        os.makedirs(save_dir, exist_ok=True)
        fname = f"{save_dir}/enhanced_pairgrid_{name}_{band}.png"
        g.savefig(fname, dpi=300, bbox_inches='tight')
        plt.close(g.fig)
    else:
        if is_streamlit:
            st.pyplot(g.fig)
        else:
            plt.show()

def detect_hep_bands(temp_dir='temps_EDF_HEP', candidates=('delta','theta','alpha','beta','gamma')):
    """
    Detect available HEP bands from filenames in temp_dir by searching for candidate tokens.
    Falls back to the provided candidates if none detected or dir missing.
    """
    if not os.path.isdir(temp_dir):
        return list(candidates)
    files = [f.lower() for f in os.listdir(temp_dir)]
    detected = [b for b in candidates if any(b in f for f in files)]
    return detected or list(candidates)

def hep_run(temp_dir='temps_EDF_HEP', bands=None, step_sec=5, save_plot=False, save_dir='figures_HEP', is_streamlit=False):
    """
    Iterate over bands, read .parquet files from temp_dir whose filename contains the band,
    z-score key columns, and plot using only_plots.

    Parameters
    ----------
    temp_dir : str
        Directory with parquet files (default 'temps_EDF_HEP').
    bands : list[str] | None
        Bands to process. If None, will attempt to auto-detect from filenames
        among ['delta','theta','alpha','beta','gamma'].
    step_sec : int
        Step size (seconds) for the time axis used by only_plots.
    save_plot : bool
        If True, save plots to save_dir.
    save_dir : str
        Directory to save plots (default 'figures_HEP').
    is_streamlit : bool
        If True, emit Streamlit UI messages and use Streamlit plotting.

    Behavior
    --------
    For each band:
      - Collect parquet files whose name contains the band token (case-insensitive)
      - Concatenate to a single DataFrame (row-wise)
      - Compute z-score per column for:
          Vagal_SD1, Sympathetic_SD2, Efficiency, Clustering, Modularity, Assortativity
        storing them in columns with suffix '_zscore'
      - Call only_plots(results_df, ...) to render/save plots
    """
    if not os.path.isdir(temp_dir):
        if is_streamlit:
            st.warning(f"Directory not found: {temp_dir}")
        return

    all_files = [f for f in os.listdir(temp_dir) if f.lower().endswith('.parquet')]
    if len(all_files) == 0:
        if is_streamlit:
            st.info(f"No .parquet files found in {temp_dir}")
        return

    default_bands = ['delta','theta','alpha','beta','gamma']
    if bands is None or len(bands) == 0:
        bands = [b for b in default_bands if any(b in f.lower() for f in all_files)] or default_bands

    for band in bands:
        band_files = [f for f in all_files if band.lower() in f.lower()]
        if len(band_files) == 0:
            if is_streamlit:
                st.write(f"No files found for band '{band}'.")
            continue

        dfs = []
        for fname in band_files:
            fpath = os.path.join(temp_dir, fname)
            try:
                df = pd.read_parquet(fpath)
            except Exception as e:
                if is_streamlit:
                    st.write(f"Failed reading {fpath}: {e}")
                continue
            dfs.append(df)

        if not dfs:
            continue

        results_df = pd.concat(dfs, ignore_index=True)

        # Compute z-scores per column for key metrics
        base_cols = ['Vagal_SD1','Sympathetic_SD2','Efficiency','Clustering','Modularity','Assortativity']
        for base in base_cols:
            # find matching column name ignoring case
            match = next((c for c in results_df.columns if isinstance(c, str) and c.strip().lower() == base.lower()), None)
            if match is None:
                continue
            results_df[match] = pd.to_numeric(results_df[match], errors='coerce')
            mu = results_df[match].mean()
            sigma = results_df[match].std()
            if sigma is not None and np.isfinite(sigma) and sigma > 0:
                results_df[f"{base}_zscore"] = (results_df[match] - mu) / sigma
            else:
                results_df[f"{base}_zscore"] = np.nan

        os.makedirs(save_dir, exist_ok=True)
        only_plots(
            results_df=results_df,
            save_plot=save_plot,
            save_dir=save_dir,
            edf_pickle_name=str(band),
            band=str(band),
            step_sec=step_sec,
            is_streamlit=is_streamlit
        )
        
def raise_error(error, context=""):
    """
    Helper function to raise errors with optional context information.
    
    Parameters
    ----------
    error : Exception
        The exception to raise
    context : str
        Additional context information for debugging
    """
    if context:
        print(f"Error context: {context}")
    raise error

# if name == main
if __name__ == '__main__':
    # Example usage
    pass
    # tga_get_files()
    # get_controls_ages_genders('temps_Controls_EDF')
    # detect_sleep('EDF', save_sleep_only=True)
    # read_pkl_files('pickles/wake_EDF')

    # merge_csv_files('temps_awake_EDF')

    # df_wnv,patients_folder,controls,df_wnv2,cases_group_name = cobrad_get_files(sample_window_size=600,only_awake=True)
    
    # df_wnv,patients_folder,controls,df_wnv2,cases_group_name = wnv_get_files()

# %%